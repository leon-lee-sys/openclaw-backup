# -*- coding: utf-8 -*-
"""
道德经第30天 课件生成器
章节：第91-93章（道者万物之奥、勇于不敢、不吾知）
风格：商务蓝/深色系，与第24天模板一致
内容详尽，每章至少4页
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# 配色方案 - 商务蓝/深色系
PRIMARY = RGBColor(30, 60, 114)       # 深蓝
SECONDARY = RGBColor(70, 130, 180)    # 中蓝
ACCENT = RGBColor(204, 120, 50)       # 橙色强调
TEXT = RGBColor(51, 51, 51)           # 深灰文字
BG = RGBColor(245, 247, 250)          # 浅灰蓝背景
WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(20, 35, 60)        # 深蓝背景

# 字体配置（比标准模板大2号）
TITLE_SIZE = Pt(38)
HEADER_SIZE = Pt(30)
YUANWEN_SIZE = Pt(22)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(16)
TOTAL_SLIDES = 17

def set_dark_title_bar(slide, title):
    """深色标题栏"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = HEADER_SIZE
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_slide_number(slide, num, total=TOTAL_SLIDES):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def add_yuanwen_section(slide, lines, y=1.5, height=3.2):
    """原文区域"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = SECONDARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line if line else "　"
        pp.font.size = YUANWEN_SIZE
        pp.font.color.rgb = WHITE
        pp.space_before = Pt(3)

def add_commentary_section(slide, title, content, y):
    """释义区域"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    box.line.color.rgb = SECONDARY
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"【{title}】"
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = BODY_SIZE
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=SECONDARY):
    """核心观点框"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(15)
    p2.font.color.rgb = WHITE

def add_points_list(slide, points, y_start=3.2, box_height=1.25):
    """要点列表"""
    y = y_start
    for icon_title, explanation in points:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = icon_title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = explanation
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)
        y += box_height + 0.1
    return y

def add_managed_insight_box(slide, title, items, y):
    """管理启示框"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(2.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    box.line.color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"📊 {title}"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    for item in items:
        p2 = tf.add_paragraph()
        p2.text = f"• {item}"
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)

def add_real_connection_box(slide, title, items, y):
    """现实联想框"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    box.line.color.rgb = ACCENT
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"🌍 {title}"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT
    for item in items:
        p2 = tf.add_paragraph()
        p2.text = f"• {item}"
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)

# ============================================================
# 第1页：封面
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.6))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.9))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第30课"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = WHITE
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.6))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第91-93章"
p2.font.size = Pt(52)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.9))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "道者万物之奥 · 勇于不敢 · 不吾知"
p3.font.size = Pt(26)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.8), Inches(7), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = DARK_BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "道者万物之奥。善人之宝，不善人之所保。"
pq1.font.size = Pt(22)
pq1.font.italic = True
pq1.font.color.rgb = WHITE
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第六十二章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = RGBColor(180, 180, 180)
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1)

# ============================================================
# 第2页：今日课程概览
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s2, "今日课程 · 三章概览")
chapters = [
    ("第六十二章", "道者万物之奥", "道者万物之奥。善人之宝，不善人之所保。道者，卑谦处下，不可须臾离也。"),
    ("第六十三章", "勇于不敢", "为无为，事无事，味无味。大小多少，报怨以德。为而不争，则莫能与之争。"),
    ("第六十四章", "不吾知", "其安易持，其未兆易谋。慎终如始，则无败事。知足者富，强行者有志。")
]
y2 = 1.5
for num, title, desc in chapters:
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y2), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf2c = box.text_frame
    tf2c.word_wrap = True
    p2c1 = tf2c.paragraphs[0]
    p2c1.text = f"📖 {num}：{title}"
    p2c1.font.size = Pt(20)
    p2c1.font.bold = True
    p2c1.font.color.rgb = PRIMARY
    p2c2 = tf2c.add_paragraph()
    p2c2.text = desc
    p2c2.font.size = Pt(15)
    p2c2.font.color.rgb = TEXT
    p2c2.space_before = Pt(5)
    y2 += 1.65
add_slide_number(s2, 2)

# ============================================================
# 第3页：第六十二章 原文
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s3, "第六十二章 · 道者万物之奥")
add_yuanwen_section(s3, [
    "道者万物之奥。",
    "善人之宝，不善人之所保。",
    "",
    "美言可以市，尊行可以贺人。",
    "人之不善，何弃之有？",
    "",
    "古之所以贵此道者何？",
    "不曰求以得，有罪以免耶？",
    "",
    "故为天下贵。"
], y=1.5, height=4.2)
add_commentary_section(s3, "白话解读",
    "道是万物的奥妙所在。善良的人把它当作珍宝，不善良的人也因它而得到保全。美好的言语可以换来尊敬，良好的行为可以使人得到庆贺。即使有不善的人，道又怎会舍弃他们呢？自古以来人们重视道的原因是什么？不就是因为：有所求就能得到，有罪过也能免除吗？所以道被天下人看得很珍贵。",
    5.8)
add_slide_number(s3, 3)

# ============================================================
# 第4页：第六十二章 逐句释义
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s4, "第六十二章 · 逐句释义")
add_points_list(s4, [
    ("① 道者万物之奥", "道是万物的奥妙所在，深藏于万物之后，主宰着一切——道是宇宙间最根本的存在。"),
    ("② 善人之宝，不善人之所保", "善良的人把道当作珍宝，不善良的人也因此得到道的庇护——道不分善恶，普遍慈悲。"),
    ("③ 美言可以市，尊行可以贺人", "美好的言语可以换来尊敬，良好的行为可以用于庆贺——道可以转化为社会资本。"),
    ("④ 不曰求以得，有罪以免耶", "有所求就能得到，有罪过也能免除——道有赦免的力量，给人改过自新的机会。"),
    ("⑤ 故为天下贵", "所以道被天下人看得很珍贵——道是人间最可宝贵的资源。")
], y_start=1.5, box_height=1.15)
add_slide_number(s4, 4)

# ============================================================
# 第5页：第六十二章 核心要义
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s5, "第六十二章 · 核心要义")
add_core_box(s5, "🌟 核心观点", "道者万物之奥——道是万物的根本所在，善人以此为宝，不善人亦可得保全。道之珍贵，在于其包容与慈悲。")
add_points_list(s5, [
    ("1. 道是万物之奥", "道深藏于万物之后，主宰天地——万事万物莫不遵循道的法则。"),
    ("2. 善人视道为宝", "善良的人珍惜道，把道当作最宝贵的财富——道是修身正心的根本。"),
    ("3. 不善人亦得保", "即使行为不善的人，道也不会弃之不顾——道的慈悲无有拣择。"),
    ("4. 求以得，免以赦", "追求道就能得到，有罪过也能得到赦免——道给人改过自新的机会。"),
    ("5. 天下贵", "道被天下人视为最珍贵的东西——因为道能带来真正的安宁与救赎。")
], y_start=3.0, box_height=1.1)
add_slide_number(s5, 5)

# ============================================================
# 第6页：第六十二章 管理启示
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s6, "第六十二章 · 管理启示")
add_managed_insight_box(s6, "管理智慧", [
    "包容文化：善人与不善人皆可用，道的包容精神启示管理者要容纳不同性格的员工。",
    "价值投资：道是'善人之宝'，启示企业要把核心价值观当作最珍贵的资产来守护。",
    "救赎机制：'有罪以免'启示建立容错机制，让犯错的员工有改正的机会。",
    "慈悲领导：道不弃不善者，启示领导者在管理中保持慈悲，给予每个人成长空间。",
    "无差别对待：道对善人与不善人同样慈悲，启示管理中一视同仁，不搞小圈子。"
], y=1.5)
add_slide_number(s6, 6)

# ============================================================
# 第7页：第六十三章 原文
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s7, "第六十三章 · 勇于不敢")
add_yuanwen_section(s7, [
    "为无为，事无事，味无味。",
    "",
    "大小多少，报怨以德。",
    "",
    "图难于其易，为大于其细。",
    "天下难事必作于易，",
    "天下大事必作于细。",
    "",
    "夫轻诺必寡信，多易必多难。",
    "",
    "是以圣人犹难之，故终无难。"
], y=1.5, height=3.8)
add_commentary_section(s7, "白话解读",
    "以'无为'的态度去作为，以'无事'的方式去做事，以'无味'的心态去品味。无论大事小事、多的少的，都用德来报答怨恨。谋划困难的事要从容易的地方入手，做大事要从细小的地方开始。天下的难事一定从容易的开始，天下的大事一定从细小的开始。轻易许诺的人一定缺少信用，把事情看得太容易一定会遇到很多困难。因此圣人把容易的事也当作困难的事，所以始终没有困难。",
    5.4)
add_slide_number(s7, 7)

# ============================================================
# 第8页：第六十三章 逐句释义
# ============================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s8, "第六十三章 · 逐句释义")
add_points_list(s8, [
    ("① 为无为，事无事，味无味", "以'无为'的态度去作为，以'无事'的方式去做事——有为而似无为，自然而然。"),
    ("② 大小多少，报怨以德", "无论大事小事、多与少，都以德报怨——用宽容和善良回应怨恨。"),
    ("③ 图难于其易，为大于其细", "谋划困难的事要从容易处入手，做大事要从细小处开始——从易到难，从细到大。"),
    ("④ 天下难事必作于易", "天下的难事一定从容易的开始——难是易的积累，大是小的发展。"),
    ("⑤ 轻诺必寡信，多易必多难", "轻易许诺一定少信，把事情看得太容易一定多难——谨慎行事，敬畏困难。"),
    ("⑥ 圣人犹难之，故终无难", "圣人把容易的事也当作困难，所以始终没有困难——防患于未然。")
], y_start=1.5, box_height=1.05)
add_slide_number(s8, 8)

# ============================================================
# 第9页：第六十三章 核心要义
# ============================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s9, "第六十三章 · 核心要义")
add_core_box(s9, "🌟 核心观点", "勇于不敢——真正的勇敢不是贸然行事，而是懂得敬畏、谨慎处之。圣人把容易的事也当困难，所以始终没有困难。")
add_real_connection_box(s9, "现实联想", [
    "易者，难之原：容易的事不做，难事就越积越多，最终无法解决。",
    "轻诺寡信：随意承诺而不兑现，就会失去他人的信任。",
    "多易多难：把事情看得太简单，准备不足，就会遇到更多困难。",
    "报怨以德：用善良回应怨恨，不是软弱，而是最有力量的态度。",
    "从细到大：任何伟大的成就都是从小处积累而来。"
], y=3.0)
add_slide_number(s9, 9)

# ============================================================
# 第10页：第六十三章 管理启示
# ============================================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s10, "第六十三章 · 管理启示")
add_managed_insight_box(s10, "管理智慧", [
    "无为而治：'为无为'启示管理者不要事事亲力亲为，要发挥团队的主动性和创造力。",
    "以德报怨：面对团队内部的矛盾和抱怨，应用宽容和善良来回应，而非针锋相对。",
    "从易着手：面对复杂项目，先从简单的部分切入，逐个击破，而非一口吃成胖子。",
    "重视细节：'为大于其细'启示在项目管理中要关注细节，小问题不及时处理会变成大麻烦。",
    "谨慎承诺：不轻易许诺，一旦许诺就一定要兑现，这样才能建立个人的信誉和威信。",
    "防患未然：把容易的事也当困难来做，提前识别风险，避免问题发生。"
], y=1.5)
add_slide_number(s10, 10)

# ============================================================
# 第11页：第六十四章 原文
# ============================================================
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s11, "第六十四章 · 不吾知")
add_yuanwen_section(s11, [
    "其安易持，其未兆易谋。",
    "其脆易泮，其微易散。",
    "",
    "为之于未有，治之于未乱。",
    "",
    "慎终如始，则无败事。",
    "",
    "知足者富，强行者有志。",
    "",
    "不失其所者久，死而不亡者寿。"
], y=1.5, height=3.8)
add_commentary_section(s11, "白话解读",
    "事物在安定的时候容易保持，在没有征兆的时候容易谋划。在脆弱的时候容易分开，在细微的时候容易消散。在事情还没有发生的时候就有所作为，在混乱还没有出现的时候就进行治理。如果在结束的时候也能像开始那样谨慎，就不会失败。知道满足的人才是富有的，坚持力行的人才是有志向的。不失去根本的人才能长久，死了但精神不消亡的人才是长寿的。",
    5.4)
add_slide_number(s11, 11)

# ============================================================
# 第12页：第六十四章 逐句释义
# ============================================================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s12, "第六十四章 · 逐句释义")
add_points_list(s12, [
    ("① 其安易持，其未兆易谋", "事物在安定时容易保持，在未显露迹象时容易谋划——防患于未然。"),
    ("② 其脆易泮，其微易散", "在脆弱时容易分开，在细微时容易消散——小问题容易解决。"),
    ("③ 为之于未有，治之于未乱", "在事情还未发生时就处理，在混乱还未出现时就预防——预防胜于补救。"),
    ("④ 慎终如始，则无败事", "在结束时也能像开始那样谨慎，就不会失败——善始善终最重要。"),
    ("⑤ 知足者富，强行者有志", "知道满足的人才是富有的，坚持力行的人是有志向的——知足与坚持是成功双翼。"),
    ("⑥ 不失其所者久，死而不亡者寿", "不失去根本的人才能长久，死了但精神不消亡才是长寿——肉体会死，但精神可以永存。")
], y_start=1.5, box_height=1.05)
add_slide_number(s12, 12)

# ============================================================
# 第13页：第六十四章 核心要义
# ============================================================
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s13, "第六十四章 · 核心要义")
add_core_box(s13, "🌟 核心观点", "不吾知——知足者富，强行者有志。在安定时要想到危险，在未乱时要预防混乱。慎终如始，方能无败。")
add_real_connection_box(s13, "现实联想", [
    "未雨绸缪：在事情还未发生时就做好准备，是智者的行为方式。",
    "知足常乐：知道满足的人，内心永远是充实的，不会被欲望驱使。",
    "善始善终：开始容易，坚持难；在结束时保持谨慎，才能确保成功。",
    "死而不亡：肉体会消亡，但留下的思想和精神可以流传千古。",
    "未兆易谋：事物在萌芽阶段最容易引导和控制，错过时机就会困难重重。"
], y=3.0)
add_slide_number(s13, 13)

# ============================================================
# 第14页：第六十四章 管理启示
# ============================================================
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s14, "第六十四章 · 管理启示")
add_managed_insight_box(s14, "管理智慧", [
    "风险管理：'为之于未有'启示管理者要有预见性，在风险还未显露时就做好准备。",
    "保持警觉：'其安易持'启示在企业顺境时更要保持警惕，防止问题积累。",
    "知足文化：'知足者富'启示企业要设定合理的目标，避免过度扩张和盲目攀比。",
    "持续奋斗：'强行者有志'启示无论取得多大成就，都要保持奋斗精神。",
    "基业常青：'不失其所者久'启示企业要守住核心价值观和使命，才能长盛不衰。",
    "精神传承：'死而不亡者寿'启示企业要建立可以传承的文化和品牌，精神永存。"
], y=1.5)
add_slide_number(s14, 14)

# ============================================================
# 第15页：三章精华总览
# ============================================================
s15 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s15, "三章精华 · 一以贯之")
sums = [
    ("第六十二章", "道者万物之奥", "道是万物的奥妙所在。善人视其为珍宝，不善人亦可得保全。道之珍贵在于其普遍的包容与慈悲。"),
    ("第六十三章", "勇于不敢", "为无为，事无事，味无味。图难于其易，为大于其细。谨慎行事，善始善终，终无败事。"),
    ("第六十四章", "不吾知", "知足者富，强行者有志。为之于未有，治之于未乱。慎终如始，死而不亡者寿。")
]
y15 = 1.5
for ch, title, content in sums:
    box = s15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y15), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf15 = box.text_frame
    tf15.word_wrap = True
    p15_1 = tf15.paragraphs[0]
    p15_1.text = f"📖 {ch}：{title}"
    p15_1.font.size = Pt(18)
    p15_1.font.bold = True
    p15_1.font.color.rgb = PRIMARY
    p15_2 = tf15.add_paragraph()
    p15_2.text = content
    p15_2.font.size = Pt(15)
    p15_2.font.color.rgb = TEXT
    p15_2.space_before = Pt(5)
    y15 += 1.65
add_slide_number(s15, 15)

# ============================================================
# 第16页：现代应用
# ============================================================
s16 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s16, "📌 现代应用")
apps = [
    ("🏢 企业治理", "道者万物之奥——企业要有核心价值观作为指导，如同道一样成为组织的根本。"),
    ("⚖️ 冲突处理", "报怨以德——面对客户投诉或员工矛盾，以善良和宽容回应，而非对抗。"),
    ("📈 项目管理", "图难于其易——大项目要从简单部分开始，小问题要及时处理，否则积重难返。"),
    ("🎯 目标设定", "知足者富——设定合理目标，不要被欲望驱使，知足才是真正的富足。"),
    ("⏰ 时间管理", "慎终如始——无论做什么事，开始和结束都要同样谨慎，才能确保完美收官。"),
    ("🌱 个人成长", "为之于未有——个人成长要提前规划，在问题还未出现时就做好准备。")
]
add_points_list(s16, apps, y_start=1.5, box_height=1.1)
add_slide_number(s16, 16)

# ============================================================
# 第17页：结语
# ============================================================
s17 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s17, "第30课 · 结语")
final_box = s17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(2))
final_box.fill.solid()
final_box.fill.fore_color.rgb = DARK_BG
tff = final_box.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "道者万物之奥"
pf1.font.size = Pt(36)
pf1.font.bold = True
pf1.font.color.rgb = WHITE
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "知足者富，强行者有志"
pf2.font.size = Pt(18)
pf2.font.italic = True
pf2.font.color.rgb = SECONDARY
pf2.alignment = PP_ALIGN.CENTER
pf2.space_before = Pt(10)
pf3 = tff.add_paragraph()
pf3.text = "为之于未有，治之于未乱"
pf3.font.size = Pt(16)
pf3.font.color.rgb = RGBColor(180, 180, 180)
pf3.alignment = PP_ALIGN.CENTER
footer_note = s17.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
tffn = footer_note.text_frame
tffn.word_wrap = True
pfn1 = tffn.paragraphs[0]
pfn1.text = "道家智慧：道是万物的奥妙，包容一切，善人与不善人皆可得保全。"
pfn1.font.size = Pt(16)
pfn1.font.color.rgb = TEXT
pfn1.alignment = PP_ALIGN.CENTER
pfn2 = tffn.add_paragraph()
pfn2.text = "知足不辱，知止不殆——这才是人生的大智慧。"
pfn2.font.size = Pt(15)
pfn2.font.color.rgb = SECONDARY
pfn2.alignment = PP_ALIGN.CENTER
pfn2.space_before = Pt(8)
add_slide_number(s17, 17)

# 保存文件
output_path = '/Users/mac/.openclaw/workspace/courses/道德经_第30天.pptx'
prs.save(output_path)
print(f"✅ 第30天课件已生成：{output_path}")
print(f"   共 {TOTAL_SLIDES} 页")