# -*- coding: utf-8 -*-
"""
道德经第31天 课件生成器
章节：第94-96章（知足者富、企者不立、有物混成）
风格：商务蓝/深色系，与第24天模板一致
内容详尽，每章至少4页
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# 配色方案 - 商务蓝/深色系
PRIMARY = RGBColor(30, 60, 114)       # 深蓝
SECONDARY = RGBColor(70, 130, 180)    # 中蓝
ACCENT = RGBColor(204, 120, 50)       # 橙色强调
TEXT = RGBColor(51, 51, 51)           # 深灰文字
BG = RGBColor(245, 247, 250)          # 浅灰蓝背景
WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(20, 35, 60)        # 深蓝背景

# 字体配置（比标准模板大2号）
TITLE_SIZE = Pt(38)
HEADER_SIZE = Pt(30)
YUANWEN_SIZE = Pt(22)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(16)
TOTAL_SLIDES = 17

def set_dark_title_bar(slide, title):
    """深色标题栏"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = HEADER_SIZE
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_slide_number(slide, num, total=TOTAL_SLIDES):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def add_yuanwen_section(slide, lines, y=1.5, height=3.2):
    """原文区域"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = SECONDARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line if line else "　"
        pp.font.size = YUANWEN_SIZE
        pp.font.color.rgb = WHITE
        pp.space_before = Pt(3)

def add_commentary_section(slide, title, content, y):
    """释义区域"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    box.line.color.rgb = SECONDARY
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"【{title}】"
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = BODY_SIZE
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=SECONDARY):
    """核心观点框"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(15)
    p2.font.color.rgb = WHITE

def add_points_list(slide, points, y_start=3.2, box_height=1.25):
    """要点列表"""
    y = y_start
    for icon_title, explanation in points:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = icon_title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = explanation
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)
        y += box_height + 0.1
    return y

def add_managed_insight_box(slide, title, items, y):
    """管理启示框"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(2.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    box.line.color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"📊 {title}"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    for item in items:
        p2 = tf.add_paragraph()
        p2.text = f"• {item}"
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)

def add_real_connection_box(slide, title, items, y):
    """现实联想框"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    box.line.color.rgb = ACCENT
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"🌍 {title}"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT
    for item in items:
        p2 = tf.add_paragraph()
        p2.text = f"• {item}"
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)

# ============================================================
# 第1页：封面
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.6))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.9))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第31课"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = WHITE
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.6))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第94-96章"
p2.font.size = Pt(52)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.9))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "知足者富 · 企者不立 · 有物混成"
p3.font.size = Pt(26)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.8), Inches(7), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = DARK_BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "知足者富，强行者有志。不失其所者久，死而不亡者寿。"
pq1.font.size = Pt(20)
pq1.font.italic = True
pq1.font.color.rgb = WHITE
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第六十四章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = RGBColor(180, 180, 180)
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1)

# ============================================================
# 第2页：今日课程概览
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s2, "今日课程 · 三章概览")
chapters = [
    ("第六十五章", "知足者富", "知足者富，强行者有志。不失其所者久，死而不亡者寿。人之迷，其日固久。知足不辱，知止不殆。"),
    ("第六十六章", "企者不立", "企者不立，跨者不行。自见者不明，自是者不彰，自伐者无功，自矜者不长。物或恶之，故有道者不处。"),
    ("第六十七章", "有物混成", "有物混成，先天地生。寂兮寥兮，独立不改，周行而不殆。万物归焉而不为主。可为天下母。")
]
y2 = 1.5
for num, title, desc in chapters:
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y2), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf2c = box.text_frame
    tf2c.word_wrap = True
    p2c1 = tf2c.paragraphs[0]
    p2c1.text = f"📖 {num}：{title}"
    p2c1.font.size = Pt(20)
    p2c1.font.bold = True
    p2c1.font.color.rgb = PRIMARY
    p2c2 = tf2c.add_paragraph()
    p2c2.text = desc
    p2c2.font.size = Pt(15)
    p2c2.font.color.rgb = TEXT
    p2c2.space_before = Pt(5)
    y2 += 1.65
add_slide_number(s2, 2)

# ============================================================
# 第3页：第六十五章 原文
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s3, "第六十五章 · 知足者富")
add_yuanwen_section(s3, [
    "知足者富，强行者有志。",
    "",
    "不失其所者久，",
    "死而不亡者寿。",
    "",
    "人之迷，其日固久。",
    "",
    "知足不辱，知止不殆，",
    "可以长久。",
    "",
    "归根曰静，是谓复命。",
    "复命曰常，知常曰明。"
], y=1.5, height=4.2)
add_commentary_section(s3, "白话解读",
    "知道满足的人才是富有的，坚持力行的人才是有志向的。不失去根本的人才能长久，死了但精神不消亡的人才是长寿的。人们对道的迷惑，由来已经很久了。知道满足就不会受到屈辱，知道停止就不会有危险，这样才可以保持长久。回归根本叫做静，这就是复归于天命。复归于天命叫做永恒，知道永恒叫做明智。",
    5.8)
add_slide_number(s3, 3)

# ============================================================
# 第4页：第六十五章 逐句释义
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s4, "第六十五章 · 逐句释义")
add_points_list(s4, [
    ("① 知足者富", "知道满足的人才是富有的——真正的富足不在于拥有多少，而在于内心的满足。"),
    ("② 强行者有志", "坚持力行的人才是有志向的——意志坚定、持之以恒，才能成就事业。"),
    ("③ 不失其所者久", "不失去根本的人才能长久——守住本心，不被外物所诱，才能长盛不衰。"),
    ("④ 死而不亡者寿", "死了但精神不消亡的人才是长寿——肉体会死，但留下的人和思想可以永存。"),
    ("⑤ 人之迷，其日固久", "人们对道的迷惑，由来已经很久了——人们往往在追求中迷失，不知止足。"),
    ("⑥ 知足不辱，知止不殆", "知道满足就不会受屈辱，知道停止就不会有危险——知足知止，方可长久。")
], y_start=1.5, box_height=1.05)
add_slide_number(s4, 4)

# ============================================================
# 第5页：第六十五章 核心要义
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s5, "第六十五章 · 核心要义")
add_core_box(s5, "🌟 核心观点", "知足者富——知道满足才是真正的富足。强行者有志，持之以恒方能成就。知足不辱，知止不殆，方可长久。")
add_points_list(s5, [
    ("1. 知足是真正的富", "外在的财富可能被夺走，但内心的知足永不失——这是真正的精神富足。"),
    ("2. 强行者有志", "坚持力行、有志向的人，才能克服困难、成就大事——意志坚定是成功的基石。"),
    ("3. 不失其所", "不失去根本、不偏离本心，是保持长久的秘密——守住初心，方得始终。"),
    ("4. 死而不亡", "肉体可死，但精神不朽——屈原、文天祥，他们的正气永远流传。"),
    ("5. 知止不殆", "知道何时该停止，就不会陷入危险——贪得无厌者最终必然失败。")
], y_start=3.0, box_height=1.1)
add_slide_number(s5, 5)

# ============================================================
# 第6页：第六十五章 管理启示
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s6, "第六十五章 · 管理启示")
add_managed_insight_box(s6, "管理智慧", [
    "知足文化：'知足者富'启示企业要设定合理的发展目标，不要过度扩张贪多嚼不烂。",
    "持续奋斗：'强行者有志'启示管理者要保持奋斗精神，在成功后依然持续精进。",
    "基业常青：'不失其所者久'启示企业要坚守核心价值观和社会责任，才能长期健康发展。",
    "精神传承：'死而不亡者寿'启示企业要建立可以传承的文化，让组织精神永存。",
    "风险管控：'知止不殆'启示在经营中要知道何时收手，避免盲目扩张带来的风险。",
    "防止迷失：'人之迷，其日固久'启示管理者要时刻保持清醒，不被成功冲昏头脑。"
], y=1.5)
add_slide_number(s6, 6)

# ============================================================
# 第7页：第六十六章 原文
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s7, "第六十六章 · 企者不立")
add_yuanwen_section(s7, [
    "企者不立，跨者不行。",
    "",
    "自见者不明，自是者不彰，",
    "自伐者无功，自矜者不长。",
    "",
    "物或恶之，故有道者不处。",
    "",
    "企者必刺[欲字]，跨者必蹶，",
    "余食赘行，物或恶之。",
    "",
    "故有道者不处。"
], y=1.5, height=3.8)
add_commentary_section(s7, "白话解读",
    "踮起脚跟的人站不稳，迈开大步的人走不远。自我表现的人不能明白事理，自以为是的人不能明辨是非，自我夸耀的人不能有功，自我矜持的人不能长久。事物都厌恶这样的行为，所以有道的人不会这样做。踮脚的人必有欲望过度的危险，迈步的人必有跌倒的危险。吃剩的食物和多余的行为，事物都厌恶。所以有道的人不会这样做。",
    5.4)
add_slide_number(s7, 7)

# ============================================================
# 第8页：第六十六章 逐句释义
# ============================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s8, "第六十六章 · 逐句释义")
add_points_list(s8, [
    ("① 企者不立，跨者不行", "踮起脚跟的人站不稳，迈开大步的人走不远——急于求成反而无法成功。"),
    ("② 自见者不明", "自我表现的人不能明白事理——总想展示自己的人，反而显得不明智。"),
    ("③ 自是者不彰", "自以为是的人不能明辨是非——认为自己一切都对，就无法看清真相。"),
    ("④ 自伐者无功", "自我夸耀的人不能有功——总夸耀自己功劳的人，反而没人承认其功。"),
    ("⑤ 自矜者不长", "自我矜持的人不能长久——骄傲自大的人，最终必然失败。"),
    ("⑥ 物或恶之，故有道者不处", "事物都厌恶这些行为，所以有道的人不会这样做——谦卑是道的本质。")
], y_start=1.5, box_height=1.05)
add_slide_number(s8, 8)

# ============================================================
# 第9页：第六十六章 核心要义
# ============================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s9, "第六十六章 · 核心要义")
add_core_box(s9, "🌟 核心观点", "企者不立——踮脚站立不稳，迈步行走不远。自见、自是、自伐、自矜，皆是败事之因。有道者不处这些行为。")
add_real_connection_box(s9, "现实联想", [
    "急于求成：踮脚站不稳、迈步走不远——越是急于成功，越容易失败。",
    "自见之蔽：总想表现自己的人，反而让人看不清其真实能力。",
    "自是之蔽：认为自己对的人，拒绝接受意见，最终判断失误。",
    "自伐之弊：总夸耀自己功劳，反而让人反感，没人愿意合作。",
    "自矜之害：骄傲自大的人，最终众叛亲离，无人跟随。",
    "谦卑之道：有道的人谦卑处下，所以能得人心、聚人才。"
], y=3.0)
add_slide_number(s9, 9)

# ============================================================
# 第10页：第六十六章 管理启示
# ============================================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s10, "第六十六章 · 管理启示")
add_managed_insight_box(s10, "管理智慧", [
    "戒急用缓：'企者不立'启示管理者不要急于求成，要稳扎稳打、步步为营。",
    "谦逊领导：'自见者不明'启示领导不要总表现自己，而要让团队成员展示能力。",
    "海纳百川：'自是者不彰'启示管理者要听取不同意见，兼听则明，偏信则暗。",
    "功成不居：'自伐者无功'启示管理者做了大成绩不要自夸，让团队来分享荣耀。",
    "戒骄戒躁：'自矜者不长'启示成功后更要谦虚，不能骄傲自满。",
    "处下为上：'有道者不处'启示真正的领导力来自谦卑处下，而非高高在上。"
], y=1.5)
add_slide_number(s10, 10)

# ============================================================
# 第11页：第六十七章 原文
# ============================================================
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s11, "第六十七章 · 有物混成")
add_yuanwen_section(s11, [
    "有物混成，先天地生。",
    "",
    "寂兮寥兮，独立不改，",
    "周行而不殆。",
    "",
    "万物归焉而不为主。",
    "",
    "可名为大，大曰逝，",
    "逝曰远，远曰反。",
    "",
    "故能成其大。"
], y=1.5, height=3.8)
add_commentary_section(s11, "白话解读",
    "有一个东西在天地之前就混合而成了。它寂静啊空虚啊，独自存在而永不改变，循环运行而从不懈怠。万物归附于它而它却不以为是万物的主人。可以称它为'大'，大而且运行不息，运行不息而愈加深远，愈加深远而又返回本源。所以它能成就其伟大。",
    5.4)
add_slide_number(s11, 11)

# ============================================================
# 第12页：第六十七章 逐句释义
# ============================================================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s12, "第六十七章 · 逐句释义")
add_points_list(s12, [
    ("① 有物混成，先天地生", "有一个东西在天地之前就混合而成了——道先于天地，是宇宙的根源。"),
    ("② 寂兮寥兮，独立不改", "寂静啊空虚啊，独自存在而永不改变——道无形无相，却永恒不变。"),
    ("③ 周行而不殆", "循环运行而从不懈怠——道永不停息地运行，生生不息。"),
    ("④ 万物归焉而不为主", "万物归附于它而它却不以为是万物的主人——道生养万物却不主宰。"),
    ("⑤ 大曰逝，逝曰远，远曰反", "大而且运行，运行而愈加深远，愈加深远而又返回——道循环往复。"),
    ("⑥ 故能成其大", "所以它能成就其伟大——因为道不争、不居，所以能成就其大。")
], y_start=1.5, box_height=1.05)
add_slide_number(s12, 12)

# ============================================================
# 第13页：第六十七章 核心要义
# ============================================================
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s13, "第六十七章 · 核心要义")
add_core_box(s13, "🌟 核心观点", "有物混成——道先天地而生，寂寥独立，周行不殆。生养万物而不为主宰，这才是真正的伟大。")
add_real_connection_box(s13, "现实联想", [
    "道法自然：道先于天地，是宇宙的根本法则——人要顺应自然规律，不可逆天而行。",
    "独立不改：道永不改变，有自己的原则——人也要有自己的核心价值观，不随风摇摆。",
    "周行不殆：道循环运行，永不停止——自然循环、社会更替，皆有其规律。",
    "生而不主：道生养万物但不主宰——领导者要成就员工但不控制，给空间但不放任。",
    "大曰逝：道伟大而运行不息——真正伟大的人永远在前进，不会满足于现状。",
    "远曰反：运行到极致又返回——月满则亏，水满则溢，盛世必有衰败，要懂进退。"
], y=3.0)
add_slide_number(s13, 13)

# ============================================================
# 第14页：第六十七章 管理启示
# ============================================================
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s14, "第六十七章 · 管理启示")
add_managed_insight_box(s14, "管理智慧", [
    "遵循规律：'道法自然'启示管理者要顺应市场规律和企业发展阶段，不可逆势而行。",
    "坚守本心：'独立不改'启示企业要有核心使命和价值观，不因短期利益而动摇。",
    "生生不息：'周行而不殆'启示组织要保持活力，不断创新，不断进步。",
    "成就员工：'万物归焉而不为主'启示领导者要帮助员工成长，而不是控制员工。",
    "不争而成：'不为主'启示管理者不要争功争名，让团队成功就是自己的成功。",
    "循环发展：'远曰反'启示企业发展到一定阶段要懂得回归本原，不忘初心。"
], y=1.5)
add_slide_number(s14, 14)

# ============================================================
# 第15页：三章精华总览
# ============================================================
s15 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s15, "三章精华 · 一以贯之")
sums = [
    ("第六十五章", "知足者富", "知足者富，强行者有志。不失其所者久，死而不亡者寿。知足不辱，知止不殆，方可长久。"),
    ("第六十六章", "企者不立", "企者不立，跨者不行。自见者不明，自是者不彰，自伐者无功，自矜者不长。有道者谦卑处下。"),
    ("第六十七章", "有物混成", "有物混成，先天地生。寂兮寥兮，独立不改，周行不殆。生养万物而不为主，可为天下母。")
]
y15 = 1.5
for ch, title, content in sums:
    box = s15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y15), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf15 = box.text_frame
    tf15.word_wrap = True
    p15_1 = tf15.paragraphs[0]
    p15_1.text = f"📖 {ch}：{title}"
    p15_1.font.size = Pt(18)
    p15_1.font.bold = True
    p15_1.font.color.rgb = PRIMARY
    p15_2 = tf15.add_paragraph()
    p15_2.text = content
    p15_2.font.size = Pt(15)
    p15_2.font.color.rgb = TEXT
    p15_2.space_before = Pt(5)
    y15 += 1.65
add_slide_number(s15, 15)

# ============================================================
# 第16页：现代应用
# ============================================================
s16 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s16, "📌 现代应用")
apps = [
    ("💰 理财观念", "知足者富——不要被欲望驱使，设定合理的目标，稳健投资才是真正的富足。"),
    ("🏃 身心健康", "企者不立——不要透支身体，踮脚站立不稳，迈步走不远——健康是长久之道。"),
    ("🤝 人际关系", "自见者不明——不要总表现自己，多倾听他人，反而更能获得尊重。"),
    ("💼 职业发展", "强行者有志——持之以恒，不断精进，职业生涯才能走得更远。"),
    ("🏛️ 企业文化", "生而不主——优秀的领导者成就员工但不控制，给予空间但不放任。"),
    ("🌱 终身成长", "有物混成——人要不断学习、不断进步，终身成长才能不被时代淘汰。")
]
add_points_list(s16, apps, y_start=1.5, box_height=1.1)
add_slide_number(s16, 16)

# ============================================================
# 第17页：结语
# ============================================================
s17 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s17, "第31课 · 结语")
final_box = s17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(2))
final_box.fill.solid()
final_box.fill.fore_color.rgb = DARK_BG
tff = final_box.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "知足者富，企者不立"
pf1.font.size = Pt(36)
pf1.font.bold = True
pf1.font.color.rgb = WHITE
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "有物混成，先天地生"
pf2.font.size = Pt(18)
pf2.font.italic = True
pf2.font.color.rgb = SECONDARY
pf2.alignment = PP_ALIGN.CENTER
pf2.space_before = Pt(10)
pf3 = tff.add_paragraph()
pf3.text = "万物归焉而不为主，可为天下母"
pf3.font.size = Pt(16)
pf3.font.color.rgb = RGBColor(180, 180, 180)
pf3.alignment = PP_ALIGN.CENTER
footer_note = s17.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
tffn = footer_note.text_frame
tffn.word_wrap = True
pfn1 = tffn.paragraphs[0]
pfn1.text = "道家智慧：知足不辱，知止不殆。不自见故明，不自是故彰，不自伐故有功。"
pfn1.font.size = Pt(16)
pfn1.font.color.rgb = TEXT
pfn1.alignment = PP_ALIGN.CENTER
pfn2 = tffn.add_paragraph()
pfn2.text = "谦卑处下，无为而成——这是老子留给我们的人生真谛。"
pfn2.font.size = Pt(15)
pfn2.font.color.rgb = SECONDARY
pfn2.alignment = PP_ALIGN.CENTER
pfn2.space_before = Pt(8)
add_slide_number(s17, 17)

# 保存文件
output_path = '/Users/mac/.openclaw/workspace/courses/道德经_第31天.pptx'
prs.save(output_path)
print(f"✅ 第31天课件已生成：{output_path}")
print(f"   共 {TOTAL_SLIDES} 页")