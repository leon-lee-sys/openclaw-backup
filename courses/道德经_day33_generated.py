# -*- coding: utf-8 -*-
"""
道德经第33天 课件生成器
章节：第100-101章 + 第81章总结（道德经完结）
风格：商务蓝/深色系，与第29/30/31/32天模板一致
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# 配色方案 - 商务蓝/深色系
PRIMARY = RGBColor(30, 60, 114)       # 深蓝
SECONDARY = RGBColor(70, 130, 180)    # 中蓝
ACCENT = RGBColor(204, 120, 50)       # 橙色强调
TEXT = RGBColor(51, 51, 51)           # 深灰文字
BG = RGBColor(245, 247, 250)          # 浅灰蓝背景
WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(20, 35, 60)        # 深蓝背景

# 字体配置
TITLE_SIZE = Pt(38)
HEADER_SIZE = Pt(30)
YUANWEN_SIZE = Pt(22)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(16)

def set_dark_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = HEADER_SIZE
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def add_yuanwen_section(slide, lines, y=1.5, height=3.2):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = SECONDARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line if line else "　"
        pp.font.size = YUANWEN_SIZE
        pp.font.color.rgb = WHITE
        pp.space_before = Pt(3)

def add_commentary_section(slide, title, content, y):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    box.line.color.rgb = SECONDARY
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"【{title}】"
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = BODY_SIZE
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=SECONDARY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(15)
    p2.font.color.rgb = WHITE

def add_points_list(slide, points, y_start=3.2):
    y = y_start
    for icon_title, explanation in points:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.25))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = icon_title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = explanation
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)
        y += 1.35
    return y

TOTAL = 16

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.6))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.9))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典33章 · 终章特辑"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = WHITE
main = s1.shapes.add_textbox(Inches(1), Inches(2.6), Inches(8), Inches(1.6))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第100-101章 · 第81章"
p2.font.size = Pt(48)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.4), Inches(8), Inches(0.9))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "为而不争 · 圣人之道 · 道德经完结篇"
p3.font.size = Pt(26)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.8), Inches(7), Inches(1.4))
quote.fill.solid()
quote.fill.fore_color.rgb = DARK_BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "信言不美，美言不信。善者不辩，辩者不善。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = WHITE
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "知者不博，博者不知。—— 第八十一章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = RGBColor(180, 180, 180)
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, TOTAL)

# ===== 第2页：终章概览 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s2, "终章课程 · 三章收官")
chapters = [
    ("第一百章", "为而不争", "为而不争，与而不取。善为道者不以知争，功成不居。玄德深矣，与物反矣。"),
    ("第一百零一章", "圣人之道", "圣人之道，为而不争，功成不居。不自见故明，不自是故彰。"),
    ("第八十一章", "总结·信言不美", "信言不美，美言不信。善者不辩，辩者不善。知者不博，博者不知——全篇总结。")
]
y2 = 1.5
for num, title, desc in chapters:
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y2), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf2c = box.text_frame
    tf2c.word_wrap = True
    p2c1 = tf2c.paragraphs[0]
    p2c1.text = f"📖 {num}：{title}"
    p2c1.font.size = Pt(20)
    p2c1.font.bold = True
    p2c1.font.color.rgb = PRIMARY
    p2c2 = tf2c.add_paragraph()
    p2c2.text = desc
    p2c2.font.size = Pt(16)
    p2c2.font.color.rgb = TEXT
    p2c2.space_before = Pt(5)
    y2 += 1.65
add_slide_number(s2, 2, TOTAL)

# ===== 第3页：第一百章 原文 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s3, "第一百章 · 为而不争")
add_yuanwen_section(s3, [
    "为而不争，与而不取。",
    "",
    "善为道者，不以知争。",
    "",
    "以其不争，",
    "故天下莫能与之争。",
    "",
    "为而不争，功成不居。",
    "夫唯不居，是以不去。",
    "",
    "玄德深矣，远矣，",
    "与物反矣，",
    "然后乃至大顺。"
], y=1.5, height=4.2)
add_commentary_section(s3, "白话解读",
    "有所作为但不争夺，给予但不索取。善于行道的人，不以智慧争夺。正因为他不争，所以天下没有人能与他争夺。有作为但不争夺，功业成就却不居功。正因为不居功，所以功业不会离去。玄德深奥啊，遥远啊，与万物同返于道，然后才能达到最大的顺应。",
    5.8)
add_slide_number(s3, 3, TOTAL)

# ===== 第4页：第一百章 逐句释义 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s4, "第一百章 · 逐句释义")
add_points_list(s4, [
    ("① 为而不争，与而不取",
     "有所作为但不争夺，给予但不索取——为道者无私无欲，只管付出，不问收获。"),
    ("② 善为道者，不以知争",
     "善于行道的人，不以智慧争夺——不用智巧争抢，不以才华压人。"),
    ("③ 以其不争，故天下莫能与之争",
     "正因为他有不争的美德，所以天下没有人能与他争夺——不争者无敌于天下。"),
    ("④ 功成不居，夫唯不居，是以不去",
     "功业成就却不居功，正因为不居功，所以功业不会离去——退让反而保全。"),
    ("⑤ 玄德深矣，远矣，与物反矣",
     "玄德深奥遥远，与万物同返于道——返回道的本源，达到大顺境界。")
], y_start=1.5)
add_slide_number(s4, 4, TOTAL)

# ===== 第5页：第一百章 本章小结 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s5, "第一百章 · 本章小结")
add_core_box(s5, "🌟 核心观点", "为而不争——行道者有所作为但不争夺，功成不居。不争之争，是为玄德。夫唯不争，天下莫能与之争。")
add_points_list(s5, [
    ("1. 为而不争", "做事尽心尽力，但不与人争夺——只问付出，不问回报。"),
    ("2. 与而不取", "给予但不索取——慷慨奉献，不求报偿。"),
    ("3. 不以知争", "不用智巧争夺——不耍小聪明，不争一时之快。"),
    ("4. 功成不居", "成就功业却不居功——谦退让贤，反而保全。"),
    ("5. 玄德大顺", "玄德深远，与物反矣——返回道的本源，达到大顺境界。")
], y_start=3.1)
add_slide_number(s5, 5, TOTAL)

# ===== 第6页：第一百零一章 原文 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s6, "第一百零一章 · 圣人之道")
add_yuanwen_section(s6, [
    "圣人之道，为而不争，功成不居。",
    "",
    "不自见故明，不自是故彰，",
    "不自伐故有功，不自矜故长。",
    "",
    "夫唯不争，故天下莫能与之争。",
    "",
    "古之所谓曲则全者，岂虚言哉！",
    "诚全而归之。"
], y=1.5, height=3.5)
add_commentary_section(s6, "白话解读",
    "圣人的道，有所作为但不争夺，功业成就却不居功。不自我显示所以明智，不自以为是所以彰显，不自我夸耀所以有功，不自我骄傲所以长久。正因为不争夺，所以天下没有人能与他争夺。古人所说的委屈反而能保全，难道是空话吗！确实能够保全而归附于道。",
    5.1)
add_slide_number(s6, 6, TOTAL)

# ===== 第7页：第一百零一章 逐句释义 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s7, "第一百零一章 · 逐句释义")
add_points_list(s7, [
    ("① 圣人之道，为而不争，功成不居",
     "圣人的道——有所作为但不争夺，功成不居，是圣人的行为准则。"),
    ("② 不自见故明，不自是故彰",
     "不自我显示所以明智，不自以为是所以彰显——谦逊的人反而能看清全局。"),
    ("③ 不自伐故有功，不自矜故长",
     "不自我夸耀所以有功，不自我骄傲所以长久——不争功反而功成名就。"),
    ("④ 夫唯不争，故天下莫能与之争",
     "正因为不争夺，所以天下没有人能与他争夺——不争是最高明的争。"),
    ("⑤ 曲则全——诚全而归之",
     "委屈反而能保全——真正的保全不是争来的，而是让出来的。")
], y_start=1.5)
add_slide_number(s7, 7, TOTAL)

# ===== 第8页：第一百零一章 本章小结 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s8, "第一百零一章 · 本章小结")
add_core_box(s8, "🌟 核心观点", "圣人之道——为而不争，功成不居。不自见、不自是、不自伐、不自矜。曲则全，诚全而归之。")
add_points_list(s8, [
    ("1. 四不精神", "不自见、不自是、不自伐、不自矜——不自我彰显，反而明智、彰显、有功、长久。"),
    ("2. 不争之争", "夫唯不争，故天下莫能与之争——不争是最高明的策略。"),
    ("3. 曲则全", "委屈反而能保全——退一步海阔天空，让三分皆大欢喜。"),
    ("4. 诚全归之", "确实能够保全而归附于道——真诚地退让，最终能保全自己。")
], y_start=3.1)
add_slide_number(s8, 8, TOTAL)

# ===== 第9页：第八十一章 原文（全书总结）=====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s9, "第八十一章 · 信言不美（全书总结）")
add_yuanwen_section(s9, [
    "信言不美，美言不信。",
    "",
    "善者不辩，辩者不善。",
    "",
    "知者不博，博者不知。",
    "",
    "圣人不积，既以为人己愈有，",
    "既以与人己愈多。",
    "",
    "天之道，利而不害。",
    "圣人之道，为而不争。",
    "",
    "—— 终章 ——"
], y=1.5, height=4.5)
add_commentary_section(s9, "白话解读",
    "真实的话不华美，华美的话不真实。善良的人不巧辩，巧辩的人不善良。有智慧的人不广博，广博的人没有智慧。圣人不积蓄，尽力帮助别人自己反而更富有，尽量给予别人自己反而更充足。天的道，利益万物而不伤害。圣人的道，有所作为但不争夺。",
    6.1)
add_slide_number(s9, 9, TOTAL)

# ===== 第10页：第八十一章 逐句释义 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s10, "第八十一章 · 逐句释义")
add_points_list(s10, [
    ("① 信言不美，美言不信",
     "真实的话不华美，华美的话不真实——花言巧语往往不可信，真话往往不好听。"),
    ("② 善者不辩，辩者不善",
     "善良的人不巧辩，巧辩的人不善良——真正善良的人不需要争辩，真理自在人心。"),
    ("③ 知者不博，博者不知",
     "有智慧的人不广博，广博的人没有智慧——专精与广博的矛盾，真正智者有所取舍。"),
    ("④ 圣人不积，为人愈有，与人愈多",
     "圣人不积蓄，帮助别人自己更富有，给予别人自己更充足——越分享越富有。"),
    ("⑤ 天之道，利而不害；圣人之道，为而不争",
     "天道利物无害，圣道为而不争——全书核心宗旨的最终总结。")
], y_start=1.5)
add_slide_number(s10, 10, TOTAL)

# ===== 第11页：第八十一章 本章小结 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s11, "第八十一章 · 本章小结")
add_core_box(s11, "🌟 核心观点", "信言不美，美言不信——终章总结：真实胜于华美，利他胜于利己，为而不争是天之道、圣人之道。")
add_points_list(s11, [
    ("1. 信与美", "信言不美，美言不信——真实的话往往不华美，华美的话往往不真实。"),
    ("2. 善与辩", "善者不辩，辩者不善——真正善良的人不需要巧辩，真理自在人心。"),
    ("3. 知与博", "知者不博，博者不知——专精与广博的矛盾，有智慧的人有所取舍。"),
    ("4. 圣人不积", "既以为人己愈有，既以与人己愈多——越分享越富有，越给予越充足。"),
    ("5. 道之总结", "天之道，利而不害；圣人之道，为而不争——全书核心宗旨的最终总结。")
], y_start=3.1)
add_slide_number(s11, 11, TOTAL)

# ===== 第12页：三章精华总览 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s12, "终章三篇 · 一以贯之")
sums = [
    ("第一百章", "为而不争", "为而不争，与而不取，功成不居。玄德深矣，与物反矣——不争是最高明的争。"),
    ("第一百零一章", "圣人之道", "圣人之道：为而不争，功成不居。不自见故明，不自是故彰。不自伐故有功，不自矜故长。"),
    ("第八十一章", "全书总结", "信言不美，美言不信。善者不辩，辩者不善。圣人不积，既以为人己愈有，既以与人己愈多。天之道，利而不害；圣人之道，为而不争。")
]
y12 = 1.5
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.55))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    p12_2.space_before = Pt(5)
    y12 += 1.65
add_slide_number(s12, 12, TOTAL)

# ===== 第13页：道德经81章核心脉络 =====
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s13, "📌 道德经81章核心脉络")
core_lines = [
    ("第一章", "道可道，非常道——开篇明义，道不可言说"),
    ("第十五章", "微妙玄通，深不可识——古之善为道者"),
    ("第三十三章", "知人者智，自知者明——认识自我"),
    ("第五十七章", "以正治国，以奇用兵——治国之道"),
    ("第六十三章", "为无为，事无事，味无味——无为而为"),
    ("第六十六章", "江海所以能为百谷王者——处下不争"),
    ("第七十八章", "天下莫柔弱于水——柔弱胜刚强"),
    ("第八十一章", "圣人之道，为而不争——全书宗旨")
]
y13 = 1.5
for i, (ch, content) in enumerate(core_lines):
    box = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y13), Inches(9), Inches(0.88))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf13 = box.text_frame
    tf13.word_wrap = True
    p13_1 = tf13.paragraphs[0]
    p13_1.text = f"{ch}：{content}"
    p13_1.font.size = Pt(16)
    p13_1.font.bold = True
    p13_1.font.color.rgb = PRIMARY
    y13 += 0.93
add_slide_number(s13, 13, TOTAL)

# ===== 第14页：现实联想 =====
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s14, "📌 现实联想")
apps = [
    ("💎 信言不美", "真实的话往往不中听——忠言逆耳利于行，多听真实反馈，少被美言迷惑。"),
    ("🤝 为而不争", "做事但不争抢——职场中做事尽心、不争功劳的人，反而容易获得晋升机会。"),
    ("🎁 与人愈多", "给予反而更富有——分享知识、经验、资源，帮助他人的同时自己也在成长。"),
    ("📊 不自伐", "不自我夸耀——有功劳不说，反而让人觉得可靠，功劳自然不会被埋没。"),
    ("🌊 曲则全", "委屈反而能保全——职场中适当让步、顾全大局的人，往往走得更远。")
]
add_points_list(s14, apps, y_start=1.5)
add_slide_number(s14, 14, TOTAL)

# ===== 第15页：管理启示 =====
s15 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s15, "💼 管理启示")
insights = [
    ("🎯 为而不争的领导力", "领导者做事但不为争夺功劳，让团队成员获得荣誉。不争的领导者反而能凝聚人心、获得追随。"),
    ("📣 不自伐的谦逊管理", "有功不自我夸耀——这样的管理者让团队成员感到被尊重，团队更加团结，功劳自然保全。"),
    ("🤝 与人愈多的领导格局", "领导者越分享、越给予团队成长机会，自己收获的尊重和能力也越多。"),
    ("⚖️ 曲则全的管理艺术", "懂得在适当时候让步、顾全大局的管理者——看似吃亏，实则保全了更大的利益。"),
    ("🏆 信言不美的决策智慧", "决策时多听真实但不好听的声音，而不是只听爱听的话——真实信息是正确决策的基础。")
]
add_points_list(s15, insights, y_start=1.5)
add_slide_number(s15, 15, TOTAL)

# ===== 第16页：全书结语 =====
s16 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s16, "全书结语 · 道德经81章圆满")
# 大字核心句
core_box = s16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(2.2))
core_box.fill.solid()
core_box.fill.fore_color.rgb = DARK_BG
tfc = core_box.text_frame
tfc.word_wrap = True
pc1 = tfc.paragraphs[0]
pc1.text = "天之道，利而不害"
pc1.font.size = Pt(36)
pc1.font.bold = True
pc1.font.color.rgb = WHITE
pc1.alignment = PP_ALIGN.CENTER
pc2 = tfc.add_paragraph()
pc2.text = "圣人之道，为而不争"
pc2.font.size = Pt(36)
pc2.font.bold = True
pc2.font.color.rgb = SECONDARY
pc2.alignment = PP_ALIGN.CENTER
pc2.space_before = Pt(8)
pc3 = tfc.add_paragraph()
pc3.text = "——《道德经》终章"
pc3.font.size = Pt(16)
pc3.font.color.rgb = RGBColor(180, 180, 180)
pc3.alignment = PP_ALIGN.CENTER
# 底部总结
footer_box = s16.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.4))
tff = footer_box.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "81章道德经，核心在于：道法自然，为无为，事无事，味无味。"
pf1.font.size = Pt(16)
pf1.font.color.rgb = TEXT
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "上善若水，柔弱胜刚强，夫唯不争，天下莫能与之争。"
pf2.font.size = Pt(16)
pf2.font.color.rgb = SECONDARY
pf2.alignment = PP_ALIGN.CENTER
pf2.space_before = Pt(6)
pf3 = tff.add_paragraph()
pf3.text = "知足者富，强行者有志，死而不亡者寿。"
pf3.font.size = Pt(16)
pf3.font.color.rgb = TEXT
pf3.alignment = PP_ALIGN.CENTER
pf3.space_before = Pt(6)
pf4 = tff.add_paragraph()
pf4.text = "信言不美，美言不信；圣人不积，既以为人己愈有，既以与人己愈多。"
pf4.font.size = Pt(15)
pf4.font.color.rgb = RGBColor(100, 100, 100)
pf4.alignment = PP_ALIGN.CENTER
pf4.space_before = Pt(6)
add_slide_number(s16, 16, TOTAL)

# 保存文件
output_path = '/Users/mac/.openclaw/workspace/courses/道德经_第33天.pptx'
prs.save(output_path)
print(f"Done: {output_path}")