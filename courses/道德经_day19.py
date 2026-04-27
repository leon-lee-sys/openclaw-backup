# -*- coding: utf-8 -*-
"""
道德经第19天 课件生成器
章节：第61-63章
风格：精读深度讲解版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
ACCENT = RGBColor(204, 0, 0)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_yuanwen_box(slide, lines, y=1.4, height=2.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(20)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=BLUE_LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(255, 255, 255)

def add_points(slide, points, y_start=3.1):
    y = y_start
    for title, content in points:
        pt_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT
        y += 1.1
    return y

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.5))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第19课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第61-63章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "大邦者下流 · 出生入死 · 道者万物之奥"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "大邦者下流，天下之牝。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = PRIMARY
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第六十一章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = SECONDARY
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, 12)

# ===== 第2页：第六十一章详解 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s2, "第六十一章 · 大邦者下流")
add_yuanwen_box(s2, [
    "大邦者下流，天下之牝，",
    "天下之交也。",
    "牝常以静胜牡，",
    "以静为下。",
    "",
    "故大邦以下小邦，则取小邦；",
    "小邦以下大邦，则取大邦。",
    "",
    "故或下以取，",
    "或下而取。",
    "大邦不过欲兼畜人，",
    "小邦不过欲入事人。",
    "夫两者各得其所欲，",
    "大者宜为下。"
], y=1.4, height=4.8)
content2 = s2.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(2))
tf_c2 = content2.text_frame
tf_c2.word_wrap = True
pc2 = tf_c2.paragraphs[0]
pc2.text = "【白话解读】大国应当像处于江河的下游，甘居天下雌柔的位置。雌柔常以安静制服雄强，以静为下。大国对小国谦下，就能取得小国的信任；小国对大国谦下，就能被大国容纳。"
pc2.font.size = Pt(14)
pc2.font.color.rgb = TEXT
add_slide_number(s2, 2, 12)

# ===== 第3页：第六十一章核心 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s3, "第六十一章 · 核心思想")
add_core_box(s3, "🌟 核心观点", "大国谦下，则天下归心。柔弱胜刚强，安静为下，是道的法则。")
add_points(s3, [
    ("1. 大邦下流", "大国像江河下流，甘居雌柔位置，以静制动。"),
    ("2. 牝静胜牡", "雌柔常以安静制服雄强——柔弱胜刚强。"),
    ("3. 各得其欲", "大国想兼畜人，小国想入事人——谦下才能各得所需。"),
    ("4. 大者宜为下", "越是强大的，越应当谦下——这是领导的艺术。")
])
add_slide_number(s3, 3, 12)

# ===== 第4页：第六十二章详解 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s4, "第六十二章 · 道者万物之奥")
add_yuanwen_box(s4, [
    "道者，万物之奥，",
    "善人之宝，",
    "不善人之所保。",
    "",
    "美言可以市，",
    "尊行可以加人。",
    "人之不善，何弃之有？",
    "",
    "故立天子，置三公，",
    "虽有拱璧以先驷马，",
    "不如坐进此道。",
    "",
    "古之所以贵此道者何？",
    "不曰求以得，有罪以免邪？",
    "故为天下贵。"
], y=1.4, height=4.8)
content4 = s4.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(2))
tf_c4 = content4.text_frame
tf_c4.word_wrap = True
pc4 = tf_c4.paragraphs[0]
pc4.text = "【白话解读】道是万物的保护者，善人珍视它，不善的人也依靠它。即使有不善的人，道也不会抛弃他们。所以即使有美玉驷马的尊贵，不如坐下来进献此道。"
pc4.font.size = Pt(14)
pc4.font.color.rgb = TEXT
add_slide_number(s4, 4, 12)

# ===== 第5页：第六十二章核心 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s5, "第六十二章 · 核心思想")
add_core_box(s5, "🌟 核心观点", "道是天下最尊贵的东西——有求就能得到，有罪也能免除。")
add_points(s5, [
    ("1. 道者万物之奥", "道是万物的保护者和藏庇——无处不在，无所不包。"),
    ("2. 善人之宝", "道是善人的珍宝——追求道就是追求善。"),
    ("3. 不善人之所保", "道也是不善人的保护——即使不善，道也不弃。"),
    ("4. 有罪以免", "道有求就能得到，有罪也能免除——这是天下最可贵的。")
])
add_slide_number(s5, 5, 12)

# ===== 第6页：现代反思 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s6, "第六十一-六十二章 · 现代反思")
items6 = [
    ("🤝 外交智慧", "大邦者下流 → 真正的强国不称霸，而是谦下待人，这样才能赢得尊重。"),
    ("💼 商业谈判", "各得其欲 → 好的谈判是让双方都得到想要的，而不是一方全赢。"),
    ("👔 职场领导", "大者宜为下 → 领导谦下，员工才愿意追随；高高在上只会让人远离。"),
    ("🌍 国际关系", "不善人之所保 → 即使立场不同，对话与合作的大门也不应关闭。")
]
y6 = 1.4
for icon_title, explanation in items6:
    box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y6), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf6 = box.text_frame
    tf6.word_wrap = True
    p6_1 = tf6.paragraphs[0]
    p6_1.text = icon_title
    p6_1.font.size = Pt(16)
    p6_1.font.bold = True
    p6_1.font.color.rgb = PRIMARY
    p6_2 = tf6.add_paragraph()
    p6_2.text = explanation
    p6_2.font.size = Pt(14)
    p6_2.font.color.rgb = TEXT
    y6 += 1.4
add_slide_number(s6, 6, 12)

# ===== 第7页：第六十三章详解 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s7, "第六十三章 · 为无为，事无事")
add_yuanwen_box(s7, [
    "为无为，",
    "事无事，",
    "味无味。",
    "",
    "大小多少，报怨以德。",
    "",
    "图难于其易，",
    "为大于其细。",
    "天下难事，必作于易；",
    "天下大事，必作于细。",
    "",
    "是以圣人终不为大，",
    "故能成其大。",
    "",
    "夫轻诺必寡信，",
    "多易必多难。",
    "是以圣人犹难之，",
    "故终无难。"
], y=1.4, height=5.2)
content7 = s7.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(2))
tf_c7 = content7.text_frame
tf_c7.word_wrap = True
pc7 = tf_c7.paragraphs[0]
pc7.text = "【白话解读】以无为的态度有所作为，以不滋事的方式处理事务，以无味为味。处理难事从容易的地方入手，做大事从小处着手。轻易许诺的人必定缺乏信用，把事情看得太容易必定会遇到困难。"
pc7.font.size = Pt(14)
pc7.font.color.rgb = TEXT
add_slide_number(s7, 7, 12)

# ===== 第8页：第六十三章核心 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s8, "第六十三章 · 核心思想")
add_core_box(s8, "🌟 核心观点", "为无为，事无事——从小处着手，从容易处入手，看似无为实则无不为。")
add_points(s8, [
    ("1. 为无为", "以无为的态度有所作为——不妄为，顺其自然。"),
    ("2. 图难于易", "处理难事从容易处入手——防患于未然。"),
    ("3. 为大于细", "做大事从小处着手——积少成多，聚沙成塔。"),
    ("4. 终无难", "把容易的事当难题认真对待，最终就没有困难了。")
])
add_slide_number(s8, 8, 12)

# ===== 第9页：西方思想 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s9, "🌍 西方思想对照")
phils = [
    ("修昔底德 · 强国困境", "雅典的崛起让斯巴达恐惧——强者谦下才能避免冲突。约前400年"),
    ("老子与马基雅维利", "无为而治 vs 铁腕统治——一柔一刚，治理风格迥异。约前400年/1513年"),
    ("杰文斯 · 边际效用", "价值取决于最后一单位的效用——小事不小，细处见真章。1867年"),
    ("孙子兵法 · 形势", "善战者，求之于势——顺势而为，事半功倍。约前500年")
]
y9 = 1.4
for name_theory, desc_year in phils:
    box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y9), Inches(9), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf9 = box.text_frame
    tf9.word_wrap = True
    p9_1 = tf9.paragraphs[0]
    p9_1.text = name_theory
    p9_1.font.size = Pt(16)
    p9_1.font.bold = True
    p9_1.font.color.rgb = PRIMARY
    p9_2 = tf9.add_paragraph()
    p9_2.text = desc_year
    p9_2.font.size = Pt(14)
    p9_2.font.color.rgb = TEXT
    y9 += 1.3
add_slide_number(s9, 9, 12)

# ===== 第10页：现代应用 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s10, "📌 现代应用")
add_points(s10, [
    ("项目管理", "图难于易 → 项目管理从小问题抓起，不要等到问题变大。"),
    ("目标设定", "为大于细 → 大目标分解成小步骤，每天进步一点点。"),
    ("信用积累", "轻诺必寡信 → 不要轻易许诺，说到做到才能建立信任。"),
    ("危机处理", "终无难 → 把每个小问题当大问题对待，才能避免大危机。")
])
add_slide_number(s10, 10, 12)

# ===== 第11页：中国典故 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s11, "📜 中国典故：张良的\"圯下拾履\"")
box11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.5))
box11.fill.solid()
box11.fill.fore_color.rgb = BG
tf11 = box11.text_frame
tf11.word_wrap = True
p11_1 = tf11.paragraphs[0]
p11_1.text = "【史记·留侯世家】"
p11_1.font.size = Pt(16)
p11_1.font.bold = True
p11_1.font.color.rgb = PRIMARY
p11_2 = tf11.add_paragraph()
p11_2.text = "张良刺杀秦始皇失败后，逃亡到下邳。有一天在桥上散步，一个老人故意把鞋甩到桥下，对张良说：'年轻人，下去帮我捡鞋！'张良惊讶，想打他，但看他年老，就下去捡了上来。老人又让他帮忙穿上，张良都照做了。老人说：'孺子可教也！'五日后，老人授予张良《太公兵法》。张良靠这部兵法辅佐刘邦，成就了汉朝大业。"
p11_2.font.size = Pt(15)
p11_2.font.color.rgb = TEXT
p11_2.space_before = Pt(10)
insight11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.1), Inches(9), Inches(1))
insight11.fill.solid()
insight11.fill.fore_color.rgb = BLUE_LIGHT
tfi11 = insight11.text_frame
tfi11.word_wrap = True
pi11 = tfi11.paragraphs[0]
pi11.text = "💡 启示"
pi11.font.size = Pt(18)
pi11.font.bold = True
pi11.font.color.rgb = RGBColor(255, 255, 255)
pi11_2 = tfi11.add_paragraph()
pi11_2.text = "张良能忍辱下人，最终成就大业——\"大者宜为下\"的最好例证。谦下不是软弱，是内在强大的表现。"
pi11_2.font.size = Pt(14)
pi11_2.font.color.rgb = RGBColor(255, 255, 255)
app11 = s11.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2))
tfa11 = app11.text_frame
tfa11.word_wrap = True
pa11 = tfa11.paragraphs[0]
pa11.text = "【现代应用】职场：尊重前辈和下属，看似吃亏实则积累人脉。创业：从小事做起，谦下待人，赢得信任和合作机会。"
pa11.font.size = Pt(14)
pa11.font.color.rgb = TEXT
add_slide_number(s11, 11, 12)

# ===== 第12页：总结 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s12, "第19课 · 三章精华总结")
sums = [
    ("第六十一章", "大邦者下流", "大国谦下则取小国，柔弱胜刚强，大者宜为下。"),
    ("第六十二章", "道者万物之奥", "道是天下贵重之物，有求以得，有罪以免。"),
    ("第六十三章", "为无为事无事", "图难于易，为大于细，终无难。轻诺必寡信。")
]
y12 = 1.4
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    y12 += 1.6
final = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(9), Inches(1))
final.fill.solid()
final.fill.fore_color.rgb = PRIMARY
tff = final.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "为无为，事无事，味无味。"
pf1.font.size = Pt(20)
pf1.font.bold = True
pf1.font.color.rgb = RGBColor(255, 255, 255)
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "—— 道德经第六十三章"
pf2.font.size = Pt(14)
pf2.font.color.rgb = RGBColor(200, 200, 200)
pf2.alignment = PP_ALIGN.CENTER
add_slide_number(s12, 12, 12)

prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第19天.pptx')
print("Done: 道德经_第19天.pptx (12 slides)")
