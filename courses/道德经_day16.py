# -*- coding: utf-8 -*-
"""
道德经第16天 课件生成器
章节：第52-54章
风格：精读深度讲解版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
ACCENT = RGBColor(204, 0, 0)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_yuanwen_box(slide, lines, y=1.4, height=2.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(20)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=BLUE_LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    return box

def add_points(slide, points, y_start=3.1):
    y = y_start
    for title, content in points:
        pt_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT
        y += 1.1
    return y

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.5))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第16课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第52-54章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "天下有始 · 盗夸非道 · 善建不拔"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "天下有始，以为天下母。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = PRIMARY
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第五十二章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = SECONDARY
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, 12)

# ===== 第2页：第五十二章详解 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s2, "第五十二章 · 天下有始")
add_yuanwen_box(s2, [
    "天下有始，以为天下母。",
    "既得其母，以知其子；",
    "既知其子，复守其母，没身不殆。",
    "",
    "塞其兑，闭其门，终身不勤；",
    "开其兑，济其事，终身不救。",
    "",
    "见小曰明，守柔曰强；",
    "用其光，复归其明，无遗身殃，是为袭常。"
], y=1.4, height=4.2)
content2 = s2.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(2))
tf_c2 = content2.text_frame
tf_c2.word_wrap = True
pc2 = tf_c2.paragraphs[0]
pc2.text = "【白话解读】道是宇宙起源，知其来源守其根本，便能终身不殆。"
pc2.font.size = Pt(14)
pc2.font.color.rgb = TEXT
add_slide_number(s2, 2, 12)

# ===== 第3页：第五十二章核心 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s3, "第五十二章 · 核心思想")
add_core_box(s3, "🌟 核心观点", "知其来源，守其根本，便能终身不殆。")
add_points(s3, [
    ("1. 天下有始", "道是宇宙的开始，是天下万物的根源。"),
    ("2. 知子守母", "知道万物的运作，仍要回归道的根本。"),
    ("3. 见小曰明", "能觉察细微变化，才是真正的明智。"),
    ("4. 守柔曰强", "保持柔弱谦下，才是真正的强大。")
])
add_slide_number(s3, 3, 12)

# ===== 第4页：第五十三章详解 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s4, "第五十三章 · 盗夸非道")
add_yuanwen_box(s4, [
    "使我介然有知，行于大道，唯施是畏。",
    "大道甚夷，而民好径。",
    "",
    "朝甚除，田甚芜，仓甚虚；",
    "服文彩，带利剑，厌饮食，",
    "财货有余，是谓盗夸。非道也哉！"
], y=1.4, height=3.8)
content4 = s4.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2.5))
tf_c4 = content4.text_frame
tf_c4.word_wrap = True
pc4 = tf_c4.paragraphs[0]
pc4.text = "【白话解读】走正道虽慢但稳，走捷径虽快但危险。表面的繁华掩盖不了内在的空虚。"
pc4.font.size = Pt(14)
pc4.font.color.rgb = TEXT
add_slide_number(s4, 4, 12)

# ===== 第5页：第五十三章核心 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s5, "第五十三章 · 核心思想")
add_core_box(s5, "⚠️ 警示", "盗夸非道——靠炫耀掩饰空虚，这不是道。", ACCENT)
add_points(s5, [
    ("1. 大道甚夷", "平坦的大道就在那里，但人们偏要走捷径。"),
    ("2. 民好径", "人心浮躁，总想找捷径成功。"),
    ("3. 表里不一", "表面光鲜亮丽，内在一无所有——这是社会的病态。"),
    ("4. 非道也哉", "强盗式的炫耀，不是真正的道。")
])
add_slide_number(s5, 5, 12)

# ===== 第6页：现代反思 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s6, "第五十三章 · 现代反思")
items6 = [
    ("🏭 社会现象", "朝甚除，田甚芜，仓甚虚 → 城市光鲜，农村凋敝"),
    ("💰 炫富心理", "服文彩，带利剑，厌饮食 → 用名牌掩饰内心不安"),
    ("📈 捷径心态", "而民好径 → 都想一夜暴富，不愿踏实做事"),
    ("🎭 盗夸心态", "财货有余，是谓盗夸 → 虚荣攀比的恶性循环")
]
y6 = 1.4
for icon_title, explanation in items6:
    box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y6), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf6 = box.text_frame
    tf6.word_wrap = True
    p6_1 = tf6.paragraphs[0]
    p6_1.text = icon_title
    p6_1.font.size = Pt(16)
    p6_1.font.bold = True
    p6_1.font.color.rgb = PRIMARY
    p6_2 = tf6.add_paragraph()
    p6_2.text = explanation
    p6_2.font.size = Pt(14)
    p6_2.font.color.rgb = TEXT
    y6 += 1.4
add_slide_number(s6, 6, 12)

# ===== 第7页：第五十四章详解 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s7, "第五十四章 · 善建不拔")
add_yuanwen_box(s7, [
    "善建者不拔，",
    "善抱者不脱，",
    "子孙以祭祀不辍。",
    "",
    "修之于身，其德乃真；",
    "修之于家，其德乃余；",
    "修之于乡，其德乃长；",
    "修之于国，其德乃丰；",
    "修之于天下，其德乃普。"
], y=1.4, height=4.0)
content7 = s7.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(2))
tf_c7 = content7.text_frame
tf_c7.word_wrap = True
pc7 = tf_c7.paragraphs[0]
pc7.text = "【白话解读】德行由近及远，从个人到天下，层层扩展，终身受益。"
pc7.font.size = Pt(14)
pc7.font.color.rgb = TEXT
add_slide_number(s7, 7, 12)

# ===== 第8页：第五十四章核心 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s8, "第五十四章 · 核心思想")
add_core_box(s8, "🌟 核心观点", "德行由近及远，从个人修养扩展到天下治理。")
add_points(s8, [
    ("1. 修身为本", "一切从个人修养开始——自身不正，难以服人。"),
    ("2. 由内而外", "从个人到家庭到乡里到国家到天下，德行逐步放大。"),
    ("3. 基业长青", "真正的事业传承，靠的是德行而非金钱或权力。"),
    ("4. 祭祀不辍", "有德的家族，子孙后代会永远纪念和传承。")
])
add_slide_number(s8, 8, 12)

# ===== 第9页：西方思想 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s9, "🌍 西方思想对照")
phils = [
    ("亚里士多德 · 美德伦理学", "美德是一种习惯，通过实践培养，形成品格。约前350年"),
    ("演化生物学 · 适应性", "真正适应性强的特质，才能在漫长历史中传承下去。1859年"),
    ("罗马斯多葛 · 四主德", "智慧、勇气、正义、节制——个人修养的四个支柱。约0年"),
    ("基因传递 · 遗传与文化", "真正传递的不是财富，而是智慧和价值观。现代")
]
y9 = 1.4
for name_theory, desc_year in phils:
    box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y9), Inches(9), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf9 = box.text_frame
    tf9.word_wrap = True
    p9_1 = tf9.paragraphs[0]
    p9_1.text = name_theory
    p9_1.font.size = Pt(16)
    p9_1.font.bold = True
    p9_1.font.color.rgb = PRIMARY
    p9_2 = tf9.add_paragraph()
    p9_2.text = desc_year
    p9_2.font.size = Pt(14)
    p9_2.font.color.rgb = TEXT
    y9 += 1.3
add_slide_number(s9, 9, 12)

# ===== 第10页：现代应用 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s10, "📌 现代应用")
add_points(s10, [
    ("企业管理", "企业文化的建立如同修道——需要从上到下的德行传递，才能基业长青。"),
    ("家庭教育", "父母修身，子女自然模仿。身教重于言教，德行是最好的传承。"),
    ("个人成长", "每天反省自己，修身养性，让自己的德行一点点进步。"),
    ("社会治理", "领导人有德，人民自然归附。以德服人者，天下归心。")
])
add_slide_number(s10, 10, 12)

# ===== 第11页：中国典故 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s11, "📜 中国典故：颜回的不二法门")
box11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.5))
box11.fill.solid()
box11.fill.fore_color.rgb = BG
tf11 = box11.text_frame
tf11.word_wrap = True
p11_1 = tf11.paragraphs[0]
p11_1.text = "【论语·雍也】"
p11_1.font.size = Pt(16)
p11_1.font.bold = True
p11_1.font.color.rgb = PRIMARY
p11_2 = tf11.add_paragraph()
p11_2.text = "子曰：'贤哉，回也！一箪食，一瓢饮，在陋巷，人不堪其忧，回也不改其乐。贤哉，回也！'"
p11_2.font.size = Pt(17)
p11_2.font.color.rgb = TEXT
p11_2.space_before = Pt(10)
insight11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.1), Inches(9), Inches(1))
insight11.fill.solid()
insight11.fill.fore_color.rgb = BLUE_LIGHT
tfi11 = insight11.text_frame
tfi11.word_wrap = True
pi11 = tfi11.paragraphs[0]
pi11.text = "💡 启示"
pi11.font.size = Pt(18)
pi11.font.bold = True
pi11.font.color.rgb = RGBColor(255, 255, 255)
pi11_2 = tfi11.add_paragraph()
pi11_2.text = "颜回身处简陋，却不改其乐——这就是'德'的体现。内在的德行才是根本，外在的物质不重要。"
pi11_2.font.size = Pt(14)
pi11_2.font.color.rgb = RGBColor(255, 255, 255)
app11 = s11.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2))
tfa11 = app11.text_frame
tfa11.word_wrap = True
pa11 = tfa11.paragraphs[0]
pa11.text = "【现代应用】职场不追求奢华，生活知足常乐，理财德本财末。"
pa11.font.size = Pt(14)
pa11.font.color.rgb = TEXT
add_slide_number(s11, 11, 12)

# ===== 第12页：总结 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s12, "第16课 · 三章精华总结")
sums = [
    ("第五十二章", "天下有始", "知其母，守其子，终身不殆。见小曰明，守柔曰强。"),
    ("第五十三章", "盗夸非道", "走正道不行捷径，盗夸非道，表里如一才是真。"),
    ("第五十四章", "善建不拔", "修德于身家乡国天下，由近及远，子孙传承。")
]
y12 = 1.4
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    y12 += 1.6
final = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(9), Inches(1))
final.fill.solid()
final.fill.fore_color.rgb = PRIMARY
tff = final.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "上善若水，水善利万物而不争。"
pf1.font.size = Pt(20)
pf1.font.bold = True
pf1.font.color.rgb = RGBColor(255, 255, 255)
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "—— 道德经第八章"
pf2.font.size = Pt(14)
pf2.font.color.rgb = RGBColor(200, 200, 200)
pf2.alignment = PP_ALIGN.CENTER
add_slide_number(s12, 12, 12)

prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第16天.pptx')
print("Done: 道德经_第16天.pptx (12 slides)")
