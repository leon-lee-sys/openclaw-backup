# -*- coding: utf-8 -*-
"""
道德经第22天 课件生成器
章节：第70-72章
风格：精读深度讲解版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
ACCENT = RGBColor(204, 0, 0)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_yuanwen_box(slide, lines, y=1.4, height=2.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(20)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=BLUE_LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(255, 255, 255)

def add_points(slide, points, y_start=3.1):
    y = y_start
    for title, content in points:
        pt_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT
        y += 1.1
    return y

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.5))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第22课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第70-72章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "知人者智 · 自知者明"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "知人者智，自知者明。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = PRIMARY
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第七十章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = SECONDARY
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, 12)

# ===== 第2页：第七十章详解 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s2, "第七十章 · 知人者智")
add_yuanwen_box(s2, [
    "知人者智，",
    "自知者明。",
    "",
    "胜人者有力，",
    "自胜者强。",
    "",
    "知足者富，",
    "强行者有志，",
    "",
    "不失其所者久，",
    "死而不亡者寿。"
], y=1.4, height=4.0)
content2 = s2.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(2))
tf_c2 = content2.text_frame
tf_c2.word_wrap = True
pc2 = tf_c2.paragraphs[0]
pc2.text = "【白话解读】能了解别人的人是智慧的，能认识自己的人是明白的。能够战胜别人是有力量的，能够战胜自己的人是强大的。知道满足的人是富有的，坚持力行的人是意志坚定的。不失掉根基的人才能长久，身死而精神不灭的人才是真正的长寿。"
pc2.font.size = Pt(13)
pc2.font.color.rgb = TEXT
add_slide_number(s2, 2, 12)

# ===== 第3页：第七十章核心 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s3, "第七十章 · 核心思想")
add_core_box(s3, "🌟 核心观点", "知人者智，自知者明——了解别人是聪明，认识自己才是真正的智慧。")
add_points(s3, [
    ("1. 知人者智", "能了解别人是智慧的——但这只是小智慧。"),
    ("2. 自知者明", "能认识自己才是明白的——这是真正的大智慧。"),
    ("3. 自胜者强", "能战胜自己才是真正的强大——最大的敌人是自己。"),
    ("4. 死而不亡者寿", "身死而精神不灭才是真正的长寿——精神永恒才是永生。")
])
add_slide_number(s3, 3, 12)

# ===== 第4页：第七十一章详解 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s4, "第七十一章 · 知不知")
add_yuanwen_box(s4, [
    "知不知，上；",
    "不知知，病。",
    "",
    "夫唯病病，",
    "是以不病。",
    "",
    "圣人不病，",
    "以其病病，",
    "是以不病。"
], y=1.4, height=3.5)
content4 = s4.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(2))
tf_c4 = content4.text_frame
tf_c4.word_wrap = True
pc4 = tf_c4.paragraphs[0]
pc4.text = "【白话解读】知道自己还有所不知，这是最好的；不知道却以为自己知道，这是缺点。把缺点当作缺点来对待，就不会有缺点。圣人没有缺点，因为他把缺点当作缺点来对待，所以没有缺点。"
pc4.font.size = Pt(14)
pc4.font.color.rgb = TEXT
add_slide_number(s4, 4, 12)

# ===== 第5页：第七十一章核心 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s5, "第七十一章 · 核心思想")
add_core_box(s5, "🌟 核心观点", "知不知，上——知道自己不知道，才是真正的上等智慧。不知却自以为知，是最大的缺点。")
add_points(s5, [
    ("1. 知不知，上", "知道自己还有所不知——这是最完美的智慧状态。"),
    ("2. 不知知，病", "不知道却以为自己知道——这是最大的毛病。"),
    ("3. 夫唯病病", "把缺点当作缺点来对待——承认问题才能解决问题。"),
    ("4. 圣人不病", "圣人没有缺点——因为他正视每一个缺点。")
])
add_slide_number(s5, 5, 12)

# ===== 第6页：第七十二章详解 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s6, "第七十二章 · 民不畏威")
add_yuanwen_box(s6, [
    "民不畏威，",
    "则大威至。",
    "",
    "无狎其所居，",
    "无厌其所生。",
    "",
    "夫唯不厌，",
    "是以不厌。",
    "",
    "是以圣人自知不自见，",
    "自爱不自贵。",
    "故去彼取此。"
], y=1.4, height=4.0)
content6 = s6.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(2))
tf_c6 = content6.text_frame
tf_c6.word_wrap = True
pc6 = tf_c6.paragraphs[0]
pc6.text = "【白话解读】当人民不害怕威压时，更大的祸患就要来了。不要逼迫人民的居所，不要厌恶人民的生活。只要不压迫人民，人民就不会厌恶统治者。所以圣人自知而不自我表现，自爱而不自以为高贵。因此舍弃后者而采取前者。"
pc6.font.size = Pt(13)
pc6.font.color.rgb = TEXT
add_slide_number(s6, 6, 12)

# ===== 第7页：第七十二章核心 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s7, "第七十二章 · 核心思想")
add_core_box(s7, "🌟 核心观点", "民不畏威，则大威至——当人民不再害怕压迫时，更大的反抗就要来了。统治者要懂得敬畏。")
add_points(s7, [
    ("1. 民不畏威，大威至", "当人民不怕压迫时，更大的灾难就要来了——革命的预兆。"),
    ("2. 无狎其所居", "不要逼得人民没有安身之处——给人民留生存空间。"),
    ("3. 自知不自见", "自知而不自我表现——有智慧但不张扬。"),
    ("4. 自爱不自贵", "自爱但不自以为高贵——谦卑是圣人品质。")
])
add_slide_number(s7, 7, 12)

# ===== 第8页：现代反思 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s8, "第七十-七十二章 · 现代反思")
items8 = [
    ("🧠 自我认知", "自知者明 → 认识自己是人生最重要的事，了解自己的优势和局限。"),
    ("💼 职场心态", "知不知，上 → 承认自己不知道，比假装知道更赢得尊重。"),
    ("🤝 管理哲学", "民不畏威 → 好的领导让员工敬畏但不恐惧，强制管理适得其反。"),
    ("🌱 人生智慧", "死而不亡者寿 → 留下有价值的思想和贡献，才是真正的永生。")
]
y8 = 1.4
for icon_title, explanation in items8:
    box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y8), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf8 = box.text_frame
    tf8.word_wrap = True
    p8_1 = tf8.paragraphs[0]
    p8_1.text = icon_title
    p8_1.font.size = Pt(16)
    p8_1.font.bold = True
    p8_1.font.color.rgb = PRIMARY
    p8_2 = tf8.add_paragraph()
    p8_2.text = explanation
    p8_2.font.size = Pt(14)
    p8_2.font.color.rgb = TEXT
    y8 += 1.4
add_slide_number(s8, 8, 12)

# ===== 第9页：西方思想 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s9, "🌍 西方思想对照")
phils = [
    ("苏格拉底 · 认识你自己", "苏格拉底说：认识你自己。知道自己无知才是真正的智慧——与自知者明相通。约前400年"),
    ("尼采 · 永恒轮回", "死而不亡者寿 → 精神可以超越肉体，达到永恒——尼采的永恒轮回思想。1889年"),
    ("马斯洛 · 自我实现", "了解自己的需求层次，才能真正成长——自我实现是最高需求。1943年"),
    ("乔布斯 · Stay Hungry, Stay Foolish", "保持饥饿，保持愚蠢——永远觉得自己不够好，才是前进的动力。2005年")
]
y9 = 1.4
for name_theory, desc_year in phils:
    box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y9), Inches(9), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf9 = box.text_frame
    tf9.word_wrap = True
    p9_1 = tf9.paragraphs[0]
    p9_1.text = name_theory
    p9_1.font.size = Pt(15)
    p9_1.font.bold = True
    p9_1.font.color.rgb = PRIMARY
    p9_2 = tf9.add_paragraph()
    p9_2.text = desc_year
    p9_2.font.size = Pt(13)
    p9_2.font.color.rgb = TEXT
    y9 += 1.3
add_slide_number(s9, 9, 12)

# ===== 第10页：现代应用 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s10, "📌 现代应用")
add_points(s10, [
    ("自我成长", "自知者明 → 定期反思，了解自己的优点和不足，持续改进。"),
    ("职场沟通", "知不知，上 → 承认自己不知道，比不懂装懂更能赢得信任。"),
    ("领导力", "自爱不自贵 → 好的领导谦卑待人，不摆架子。"),
    ("人生价值", "死而不亡者寿 → 思考你能留下什么，什么能被后人记住。")
])
add_slide_number(s10, 10, 12)

# ===== 第11页：中国典故 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s11, "📜 中国典故：秦穆公与偷马者")
box11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.5))
box11.fill.solid()
box11.fill.fore_color.rgb = BG
tf11 = box11.text_frame
tf11.word_wrap = True
p11_1 = tf11.paragraphs[0]
p11_1.text = "【史记·秦本纪】"
p11_1.font.size = Pt(16)
p11_1.font.bold = True
p11_1.font.color.rgb = PRIMARY
p11_2 = tf11.add_paragraph()
p11_2.text = "秦穆公丢失了一匹良马，他亲自去找，发现有人杀了马正在吃。左右的人要惩罚吃马者，穆公却说：有义士杀马而食之，是马也。并赐酒给他们。后来秦晋交战，穆公被困，幸得那些吃马的人相救，最终击败晋军，活了下来。"
p11_2.font.size = Pt(15)
p11_2.font.color.rgb = TEXT
p11_2.space_before = Pt(10)
insight11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.1), Inches(9), Inches(1))
insight11.fill.solid()
insight11.fill.fore_color.rgb = BLUE_LIGHT
tfi11 = insight11.text_frame
tfi11.word_wrap = True
pi11 = tfi11.paragraphs[0]
pi11.text = "💡 启示"
pi11.font.size = Pt(18)
pi11.font.bold = True
pi11.font.color.rgb = RGBColor(255, 255, 255)
pi11_2 = tfi11.add_paragraph()
pi11_2.text = "自爱不自贵——秦穆公身为君主，却不以尊贵自居，最终赢得了部下拼死相报。谦卑者得人心，这是永恒的道理。"
pi11_2.font.size = Pt(14)
pi11_2.font.color.rgb = RGBColor(255, 255, 255)
add_slide_number(s11, 11, 12)

# ===== 第12页：总结 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s12, "第22课 · 三章精华总结")
sums = [
    ("第七十章", "知人者智", "知人者智，自知者明。胜人者有力，自胜者强。"),
    ("第七十一章", "知不知，上", "知不知，上；不知知，病。知道自己不知道才是真正的智慧。"),
    ("第七十二章", "民不畏威", "自知不自见，自爱不自贵。去彼取此，谦卑待人。")
]
y12 = 1.4
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    y12 += 1.6
final = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(9), Inches(1))
final.fill.solid()
final.fill.fore_color.rgb = PRIMARY
tff = final.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "知人者智，自知者明。"
pf1.font.size = Pt(20)
pf1.font.bold = True
pf1.font.color.rgb = RGBColor(255, 255, 255)
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "—— 道德经第七十章"
pf2.font.size = Pt(14)
pf2.font.color.rgb = RGBColor(200, 200, 200)
pf2.alignment = PP_ALIGN.CENTER
add_slide_number(s12, 12, 12)

prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第22天.pptx')
print("Done: 道德经_第22天.pptx (12 slides, 70-72章)")
