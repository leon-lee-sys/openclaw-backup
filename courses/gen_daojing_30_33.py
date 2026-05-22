#!/usr/bin/env python3
"""生成道德经第30-33天课件"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

CHAPTERS = {
    30: [
        (91, "道者万物之奥", "道是万物的根本。善人的法宝，不善人的保护伞。美好的言辞可以换来尊重，良好的行为可以讨人欢喜。给予别人好处，没有比这更有福的了。"),
        (92, "勇于不敢", "勇于柔弱就能活，勇于坚强就会死。天道所厌恶的，谁能知道是什么缘故？因此连圣人也觉得难于解释。"),
        (93, "不吾知", "不被了解而不争辩，不得认可而不强求。始终不和颜悦色，不以声音压人。成长而不夸耀，看见小的就想着大的。")
    ],
    31: [
        (94, "知足者富", "知道满足的人就是富有。坚持力行的人就是有志。不丧失根本的人就能长久。死了而精神不灭的人就是长寿。"),
        (95, "企者不立", "踮起脚跟的人站不稳，跨大步走的人走不远。自我表现的人不能显明，自以为是的人不能彰显，自我夸耀的人没有功劳。"),
        (96, "有物混成", "有一个东西混然而成，在天地之前就产生了。寂静啊无声啊，独自存在而不改变，循环运行而不衰竭，可以作为天下的根本。人法地，地法天，天法道，道法自然。")
    ],
    32: [
        (97, "致知", "获得智慧的方法，古代有大圣人，知道天下最真实的情况，所以圣人的治理，关键在于知足。"),
        (98, "为而不争", "做事但不争夺。正因为不争夺，天下才没有人能与他争。古代所说的委屈才能保全，岂是空话！"),
        (99, "以其不争", "正因为他不争夺，所以天下没有人能与他争夺。")
    ],
    33: [
        (100, "小国寡民", "国家小人口少。即使有各种器具也不使用，使人民重视死亡而不向远处迁移。甘其食，美其服，安其居，乐其俗。"),
        (101, "甘其食", "以自己的食物为香甜，以自己的衣服为美好，以自己的居所为安定，以自己的习俗为快乐。邻国相望，鸡犬之声相闻，民至老死不相往来。"),
        (81, "信言不美", "真实的言语不华丽，华丽的言语不真实。善良的人不巧辩，巧辩的人不善良。圣人不积累，尽量帮助别人自己反而更充足。")
    ]
}

ORIGINALS = {
    91: "道者万物之奥。善人之宝，不善人之所保。",
    92: "勇于敢则杀，勇于不敢则活。此两者，或利或害。",
    93: "不吾知其否？不见是名？终不大声以色？",
    94: "知足者富也。强行者有志也。不失其所者久也。",
    95: "企者不立，跨者不行。自见者不明，自是者不彰。",
    96: "有物混成，先天地生。寂兮寥兮，独立不改。",
    97: "致知也，古代有大圣者，知天下之至情。",
    98: "为而不争。夫唯不争，天下莫能与之争。",
    99: "以其不争，故天下莫能与之争。",
    100: "小国寡民，使有什伯之器而不用，使民重死而不远徒。",
    101: "甘其食，美其服，安其居，乐其俗。",
    81: "信言不美，美言不信。善者不辩，辩者不善。"
}

ASSOCIATIONS = {
    91: "管理者要懂得'道'是根本，不要只看表面成绩",
    92: "柔弱不是懦弱，而是顺势而为的智慧",
    93: "不争是一种更高明的竞争策略",
    94: "知足不是躺平，而是懂得适可而止",
    95: "自我膨胀必然导致失败",
    96: "遵循自然规律，不要强行干预",
    97: "学习的关键是知足，不是贪多",
    98: "为而不争，方能成其私",
    99: "不争是最高明的领导艺术",
    100: "规模不是目的，幸福感才是",
    101: "简单的快乐最真实",
    81: "真实比华丽更有价值"
}

DARK_BLUE = RGBColor(0, 51, 102)
GRAY = RGBColor(80, 80, 80)
LIGHT_GRAY = RGBColor(100, 100, 100)

def add_title_slide(prs, day_num, ch_range):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_box.text_frame.text = f"道德经 · 第{day_num}天"
    title_box.text_frame.paragraphs[0].font.size = Pt(48)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
    sub_box.text_frame.text = f"第{ch_range}"
    sub_box.text_frame.paragraphs[0].font.size = Pt(24)
    sub_box.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    sub_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_chapter_slide(prs, ch_num, title, orig_text, note):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    ch_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1))
    ch_title.text_frame.text = f"第{ch_num}章 · {title}"
    ch_title.text_frame.paragraphs[0].font.size = Pt(36)
    ch_title.text_frame.paragraphs[0].font.bold = True
    ch_title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    ch_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    orig_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.2))
    orig_box.text_frame.text = orig_text
    orig_box.text_frame.paragraphs[0].font.size = Pt(20)
    orig_box.text_frame.paragraphs[0].font.color.rgb = GRAY
    orig_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    note_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.333), Inches(2))
    note_box.text_frame.text = note
    note_box.text_frame.paragraphs[0].font.size = Pt(18)
    note_box.text_frame.paragraphs[0].font.color.rgb = GRAY

def add_keypoint_slide(prs, ch_num, note, association):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    key_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.8))
    key_title.text_frame.text = f"第{ch_num}章 · 核心要义"
    key_title.text_frame.paragraphs[0].font.size = Pt(32)
    key_title.text_frame.paragraphs[0].font.bold = True
    key_title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

    key_point = note.split("。")[0] + "。"
    key_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
    key_box.text_frame.text = key_point
    key_box.text_frame.paragraphs[0].font.size = Pt(22)
    key_box.text_frame.paragraphs[0].font.color.rgb = GRAY

    assoc_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(1.5))
    assoc_box.text_frame.text = f"💡 {association}"
    assoc_box.text_frame.paragraphs[0].font.size = Pt(22)
    assoc_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 153)

def add_end_slide(prs, day_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    end_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    end_box.text_frame.text = f"第{day_num}天 · 完"
    end_box.text_frame.paragraphs[0].font.size = Pt(36)
    end_box.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    end_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_ppt(day_num, chapters):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ch_nums = [c[0] for c in chapters]
    ch_range = f"{min(ch_nums)}-{max(ch_nums)}章"

    add_title_slide(prs, day_num, ch_range)

    for ch_num, title, note in chapters:
        orig_text = ORIGINALS.get(ch_num, "")
        association = ASSOCIATIONS.get(ch_num, "道法自然，顺势而为")
        add_chapter_slide(prs, ch_num, title, orig_text, note)
        add_keypoint_slide(prs, ch_num, note, association)

    add_end_slide(prs, day_num)

    out_path = f'/Users/mac/.openclaw/workspace/courses/道德经_第{day_num}天.pptx'
    prs.save(out_path)
    print(f"✅ 生成: 道德经_第{day_num}天.pptx ({len(chapters)*2+2}页)")
    return out_path

def main():
    for day_num in [30, 31, 32, 33]:
        chapters = CHAPTERS[day_num]
        create_ppt(day_num, chapters)
    print("\n全部完成！")

if __name__ == '__main__':
    main()