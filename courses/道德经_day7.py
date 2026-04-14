from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

# ===== 第1页：封面 =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar.line.fill.background()
txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第七课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

main_title = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main_title.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第13-15章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER

subtitle = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = subtitle.text_frame
p3 = tf3.paragraphs[0]
p3.text = "宠辱若惊 · 致虚守静 · 微妙玄通"
p3.font.size = Pt(28)
p3.font.color.rgb = RGBColor(70, 130, 180)
p3.alignment = PP_ALIGN.CENTER

quote_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
txBox2 = slide1.shapes.add_textbox(Inches(2), Inches(5.9), Inches(6), Inches(1))
tf4 = txBox2.text_frame
p4 = tf4.paragraphs[0]
p4.text = "「致虚极，守静笃。万物并作，吾以观复。」"
p4.font.size = Pt(20)
p4.font.color.rgb = RGBColor(100, 100, 100)
p4.font.italic = True
p4.alignment = PP_ALIGN.CENTER

# ===== 第2页：第十三章 宠辱若惊 =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
bg2.fill.solid()
bg2.fill.fore_color.rgb = RGBColor(30, 60, 114)
bg2.line.fill.background()
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = title2.text_frame
tf.paragraphs[0].text = "第十三章 · 宠辱若惊"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

section_title = slide2.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
tf = section_title.text_frame
tf.paragraphs[0].text = "【原文】"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(30, 60, 114)

original = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
tf = original.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "宠辱若惊，贵大患若身。何谓宠辱若惊？宠为下，得之若惊，失之若惊，是谓宠辱若惊。何谓贵大患若身？吾所以有大患者，为吾有身。及吾无身，或无患。"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(60, 60, 60)
p.line_spacing = 1.5

section_title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.5))
tf = section_title2.text_frame
tf.paragraphs[0].text = "【白话解读】"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(30, 60, 114)

modern = slide2.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1.5))
tf = modern.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '本章论述"宠辱若惊"与"贵身"之道。得宠与受辱都使人惊惧，因为两者都暗示自身处于被评判的下位。真正的忧患源于过度执着于自我身形——把身体看得太重。反之，若能超越私身之累，便能无患。核心在于破除自我中心，不以荣辱为念，方可超然物外。'
p.font.size = Pt(15)
p.font.color.rgb = RGBColor(80, 80, 80)
p.line_spacing = 1.5

# ===== 第3页：第十三章 典故 =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
bg3.fill.solid()
bg3.fill.fore_color.rgb = RGBColor(30, 60, 114)
bg3.line.fill.background()
title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = title3.text_frame
tf.paragraphs[0].text = "第十三章 · 典故案例"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

story_title = slide3.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
tf = story_title.text_frame
tf.paragraphs[0].text = "【庄子·逍遥游】"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)

story = slide3.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(2))
tf = story.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '尧帝想让天下给许由，许由却说："你治理天下已经治理得很好，我若再来代替你，不过是名不符实。名是实的宾位，我何苦求名呢？"许由拒不接受天下，因为他明白：执着于"有天下"的名位，恰是"大患"的根源。真正的逍遥者，不为宠辱所动，不以天下为贵。'
p.font.size = Pt(15)
p.font.color.rgb = RGBColor(80, 80, 80)
p.line_spacing = 1.5

foreign_title = slide3.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
tf = foreign_title.text_frame
tf.paragraphs[0].text = "【斯多葛学派】"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)

foreign = slide3.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(2))
tf = foreign.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '古罗马哲学家爱比克泰德说："人不是被事物本身困扰，而是被他们对事物的看法困扰。"这与老子"贵大患若身"的智慧相通。斯多葛派主张区分可控与不可控之事，不为外在的宠辱得失而扰乱内心，与道家"无身无患"异曲同工。'
p.font.size = Pt(15)
p.font.color.rgb = RGBColor(80, 80, 80)
p.line_spacing = 1.5

# ===== 第4页：第十四章 致虚守静 =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
bg4.fill.solid()
bg4.fill.fore_color.rgb = RGBColor(30, 60, 114)
bg4.line.fill.background()
title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = title4.text_frame
tf.paragraphs[0].text = "第十四章 · 致虚守静"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

section_title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
tf = section_title4.text_frame
tf.paragraphs[0].text = "【原文】"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(30, 60, 114)

original4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
tf = original4.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "致虚极，守静笃。万物并作，吾以观复。夫物芸芸，各复归其根。归根曰静，是谓复命。复命曰常，知常曰明。不知常，妄作凶。知常容，容乃公，公乃全，全乃天，天乃道，道乃久，没身不殆。"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(60, 60, 60)
p.line_spacing = 1.5

section_title42 = slide4.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.5))
tf = section_title42.text_frame
tf.paragraphs[0].text = "【白话解读】"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(30, 60, 114)

modern4 = slide4.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(2))
tf = modern4.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '本章为全篇最精要之一。"致虚极，守静笃"——使心灵空明到极致坚守清静。万物蓬勃生长，我从中观察其循环往复的规律。一切纷繁复杂的事物，终将回归根本。回归根本叫"静"，这便是"复命"——复归本性的法则。认识这个法则叫"明"，不认识就会妄为招凶。认识法则则能包容，包容则公正，公正则周全，周全则合乎自然，合乎自然便合乎道，合道便能长久，终身免于危难。'
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(80, 80, 80)
p.line_spacing = 1.4

# ===== 第5页：第十四章 典故 =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
bg5.fill.solid()
bg5.fill.fore_color.rgb = RGBColor(30, 60, 114)
bg5.line.fill.background()
title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = title5.text_frame
tf.paragraphs[0].text = "第十四章 · 典故案例"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

story5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(2.5))
tf = story5.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = """【陶渊明·归去来兮辞】

陶渊明任彭泽县令八十余日后，叹道："归去来兮，田园将芜胡不归？"他不为五斗米折腰，毅然辞官归隐。正是参透了"万物归根"的道理——官场荣辱不过是身外之物，回归本性、顺应自然才是正道。他在田园中找到了生命的归宿，达到"复命"的境界。

【曾国藩·静字诀】

曾国藩以"静"字为修身第一要诀。他在日记中写道："心静则体察精，克制亦省力。"纷乱的乱世中，正是内心的清静使他做出正确决策。这正是"致虚极，守静笃"的现实应用。"""
p.font.size = Pt(15)
p.font.color.rgb = RGBColor(80, 80, 80)
p.line_spacing = 1.4

# ===== 第6页：第十五章 微妙玄通 =====
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
bg6.fill.solid()
bg6.fill.fore_color.rgb = RGBColor(30, 60, 114)
bg6.line.fill.background()
title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = title6.text_frame
tf.paragraphs[0].text = "第十五章 · 微妙玄通"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

section_title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
tf = section_title6.text_frame
tf.paragraphs[0].text = "【原文】"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(30, 60, 114)

original6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
tf = original6.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "古之善为士者，微妙玄通，深不可识。夫唯不可识，故强为之容：豫兮若冬涉川，犹兮若畏四邻，俨兮其若客，涣兮若冰之将释，敦兮其若朴，旷兮其若谷，混兮其若浊。孰能浊以静之徐清？孰能安以久动之徐生？保此道者不欲盈。夫唯不盈，故能蔽不新成。"
p.font.size = Pt(15)
p.font.color.rgb = RGBColor(60, 60, 60)
p.line_spacing = 1.4

section_title62 = slide6.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.5))
tf = section_title62.text_frame
tf.paragraphs[0].text = "【白话解读】"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(30, 60, 114)

modern6 = slide6.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(2))
tf = modern6.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '古代得道之人，思想微妙通达，深不可识。七种意象描绘其品格：谨慎如冬天涉水过河，警觉如畏惧四邻，恭敬如做客，洒脱如冰块消融，敦厚如未经雕琢的木料，空旷如深谷，浑厚如浊水。如何使浑水慢慢澄清？如何使安静转为徐徐生发？保持此道者不求盈满。正因不盈满，才能在守成中保有新的可能。'
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(80, 80, 80)
p.line_spacing = 1.4

# ===== 第7页：第十五章 典故 =====
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
bg7.fill.solid()
bg7.fill.fore_color.rgb = RGBColor(30, 60, 114)
bg7.line.fill.background()
title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = title7.text_frame
tf.paragraphs[0].text = "第十五章 · 典故案例"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

story7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(2.5))
tf = story7.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = """【张良·圮上进履】

张良在桥下为老人拾鞋、穿鞋，老人知其可教，遂传《太公兵法》。张良一生谨慎小心，"状貌若妇人"，却在风云变幻中保全自身，助刘邦定天下，正是"豫兮若冬涉川"的典范。

【曾国藩的用人之道】

曾国藩选才用人，不求全责备，喜欢用"乡气重"之人——纯朴而有活力。正是"敦兮其若朴"的智慧。他深知：过于完美的人反而难以成事，不求盈满才能持续进步。"""
p.font.size = Pt(15)
p.font.color.rgb = RGBColor(80, 80, 80)
p.line_spacing = 1.4

# ===== 第8页：精华提炼 =====
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
bg8.fill.solid()
bg8.fill.fore_color.rgb = RGBColor(30, 60, 114)
bg8.line.fill.background()
title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = title8.text_frame
tf.paragraphs[0].text = "精华提炼"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

point1 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.1), Inches(4.3), Inches(2))
point1.fill.solid()
point1.fill.fore_color.rgb = RGBColor(240, 248, 255)
p1 = slide8.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(4), Inches(0.5))
p1.tf = p1.text_frame
p1.tf.paragraphs[0].text = "第十三章：宠辱不惊"
p1.tf.paragraphs[0].font.size = Pt(16)
p1.tf.paragraphs[0].font.bold = True
p1.tf.paragraphs[0].font.color.rgb = RGBColor(30, 60, 114)
p1b = slide8.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(4), Inches(1.3))
p1b.tf = p1b.text_frame
p1b.tf.word_wrap = True
p1b.tf.paragraphs[0].text = "放下自我执着，不为宠辱得失所动。真正的内心安宁，来自于超越对身体的执念。"
p1b.tf.paragraphs[0].font.size = Pt(13)
p1b.tf.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)

point2 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.1), Inches(4.3), Inches(2))
point2.fill.solid()
point2.fill.fore_color.rgb = RGBColor(240, 248, 255)
p2 = slide8.shapes.add_textbox(Inches(5.4), Inches(1.2), Inches(4), Inches(0.5))
p2.tf = p2.text_frame
p2.tf.paragraphs[0].text = "第十四章：归根复命"
p2.tf.paragraphs[0].font.size = Pt(16)
p2.tf.paragraphs[0].font.bold = True
p2.tf.paragraphs[0].font.color.rgb = RGBColor(30, 60, 114)
p2b = slide8.shapes.add_textbox(Inches(5.4), Inches(1.7), Inches(4), Inches(1.3))
p2b.tf = p2b.text_frame
p2b.tf.word_wrap = True
p2b.tf.paragraphs[0].text = "万事万物终归根本，静是回归的途径。致虚守静，方能洞察事物发展规律。知常曰明。"
p2b.tf.paragraphs[0].font.size = Pt(13)
p2b.tf.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)

point3 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.3), Inches(4.3), Inches(2))
point3.fill.solid()
point3.fill.fore_color.rgb = RGBColor(240, 248, 255)
p3 = slide8.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(4), Inches(0.5))
p3.tf = p3.text_frame
p3.tf.paragraphs[0].text = "第十五章：微妙玄通"
p3.tf.paragraphs[0].font.size = Pt(16)
p3.tf.paragraphs[0].font.bold = True
p3.tf.paragraphs[0].font.color.rgb = RGBColor(30, 60, 114)
p3b = slide8.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(4), Inches(1.3))
p3b.tf = p3b.text_frame
p3b.tf.word_wrap = True
p3b.tf.paragraphs[0].text = "真正的智者深藏不露，谨慎而通达。不求盈满，方能保持持续成长的空间，守成而有新成。"
p3b.tf.paragraphs[0].font.size = Pt(13)
p3b.tf.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)

core = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(3.3), Inches(4.3), Inches(2))
core.fill.solid()
core.fill.fore_color.rgb = RGBColor(30, 60, 114)
p_core = slide8.shapes.add_textbox(Inches(5.4), Inches(3.5), Inches(4), Inches(0.5))
p_core.tf = p_core.text_frame
p_core.tf.paragraphs[0].text = "💡 本课核心"
p_core.tf.paragraphs[0].font.size = Pt(16)
p_core.tf.paragraphs[0].font.bold = True
p_core.tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
p_core2 = slide8.shapes.add_textbox(Inches(5.4), Inches(4.0), Inches(4), Inches(1.2))
p_core2.tf = p_core2.text_frame
p_core2.tf.word_wrap = True
p_core2.tf.paragraphs[0].text = "静心致虚，复归根本\n不争不盈，玄德长存"
p_core2.tf.paragraphs[0].font.size = Pt(14)
p_core2.tf.paragraphs[0].font.color.rgb = RGBColor(200, 220, 255)
p_core2.tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_slide_number(slide8, 8, 8)

prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第7天.pptx')
print("done")
