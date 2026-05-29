# -*- coding: utf-8 -*-
"""
GenLunyu - Generate 7 PPTX files for Analects of Confucius
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
ACCENT = RGBColor(180, 60, 60)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)
GOLD = RGBColor(180, 140, 50)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.2), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "%d/%d" % (num, total)
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_yuanwen_box(slide, lines, y=1.5, height=4.0):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "[Original Text]"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(22)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(4)

def add_core_box(slide, title, subtitle, color=BLUE_LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(17)
    p2.font.color.rgb = RGBColor(230, 230, 230)

def add_points(slide, points, y_start=3.2):
    y = y_start
    for title, content in points:
        pt_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.2))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(19)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT
        y += 1.2
    return y

def add_zhushi_box(slide, items, y=1.5, height=4.0):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "[Annotations]"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for term, explanation in items:
        pp = tf.add_paragraph()
        pp.text = "* %s: %s" % (term, explanation)
        pp.font.size = Pt(17)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)

def add_translation_box(slide, cn_text, en_text, y=5.8):
    trans_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.1))
    trans_box.fill.solid()
    trans_box.fill.fore_color.rgb = BLUE_LIGHT
    tf = trans_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "[Translation] " + cn_text
    p1.font.size = Pt(16)
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = "[EN] " + en_text
    p2.font.size = Pt(15)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(220, 220, 255)

def add_discussion_box(slide, items, y=1.5):
    y_cur = y
    for q_num, question in items:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_cur), Inches(9), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "Q%d: %s" % (q_num, q_num)
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT
        p2 = tf.add_paragraph()
        p2.text = question
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT
        y_cur += 1.3

# ============================================================
# Chapter data for each day
# ============================================================

CHAPTERS = [
    # Day 1: 学而第一
    [
        {
            "num": 1,
            "name": "Xue Er Shi Xi (学而时习之)",
            "yuanwen": [
                "Zi Yue: 'Xue er shi xi zhi, bu yi yue hu?",
                "You peng zi yuan fang lai, bu yi le hu?",
                "Ren bu zhi er bu yun, bu yi junzi hu?'"
            ],
            "zhushi": [
                ("Xue (学)", "Study knowledge"),
                ("Shi Xi (时习)", "Review and practice regularly"),
                ("Shuo (说)", "Same as '悦' - joy and delight"),
                ("Peng (朋)", "Like-minded friend"),
                ("Yun (愠)", "Resentment, anger"),
                ("Junzi (君子)", "Person of noble character")
            ],
            "translation_cn": "学习知识后时常复习不是很喜悦吗？有志同道合的朋友从远方来不是很快乐吗？别人不了解我，我却不恼怒，这不也是君子吗？",
            "translation_en": "Is it not delightful to learn and constantly review what one has studied? Is it not a joy to have friends come from afar? Does it not show noble character when one is not resentful even if misunderstood?",
            "core_points": [
                ("1. Xue er shi xi", "Learning knowledge requires constant review and practice - this is the fundamental method of learning. Reviewing the old to gain new insight is how true mastery is achieved."),
                ("2. You peng yuan fang lai", "Having like-minded friends come from afar is a joy in life. True friends are companions of the soul."),
                ("3. Ren bu zhi bu yun", "When others don't understand me, I am not resentful - this is the cultivation of a noble person. Maintaining calm when not recognized.")
            ],
            "background": "Xue Er (学而) is the first chapter of the Analects, reflecting Confucius's educational philosophy and view of life. During the Spring and Autumn period, education was a privilege of the few, but Confucius advocated 'You jiao wu lei' (education for all without discrimination), allowing more people the opportunity to learn.",
            "western": [
                ("Socrates: Virtue is Knowledge", "Socrates believed virtue comes from knowledge - knowing the good will lead to doing the good. This connects with 'Xue er shi xi' - learning is the foundation of virtue. 400 BCE"),
                ("Aristotle: Habit Forms Virtue", "We become just by doing just things, brave by doing brave things. Repeated practice forms virtue. c. 350 BCE"),
                ("Confucius: Education for All", "Confucius broke class limitations, allowing everyone to learn - his influence was broader than Socrates. c. 500 BCE"),
                ("Dewey: Learning by Doing", "Education is life itself, learning by doing - American philosopher Dewey's core idea connects with 'Shi xi zhi'. 1897")
            ],
            "shiyi": "Xue er shi xi zhi: Learning is not just about listening to lectures, but more importantly about repeated practice and application. Modern learning theory proves that spaced repetition is the most effective learning method - 300% more effective than cramming. 'You peng yuan fang lai': An open heart to welcome like-minded friends - networking is an important wealth in life. 'Ren bu zhi bu yun': Inner peace, not being swayed by others' opinions - a sign of mature personality.",
            "discussion": [
                ("Question 1", "Why did Confucius place 'Xue er shi xi' as the first saying in the Analects? What educational philosophy does this reflect?"),
                ("Question 2", "How can 'Ren bu zhi er bu yun' be applied in modern workplace? What would you do when your suggestions are ignored?"),
                ("Question 3", "From your own experience, talk about the inspiration 'You peng zi yuan fang lai' brought you. What does friendship mean to you?")
            ]
        },
        {
            "num": 2,
            "name": "Xiao Ti Zhi Xing (孝悌之行)",
            "yuanwen": [
                "You Zi Yue: 'Qi Wei Ren Ye Xiao Ti,",
                "Er Hao Fan Shang Zhe, Xian Yi;",
                "Bu Hao Fan Shang, Er Hao Zuo Luan Zhe, Wei Zhi You Ye.'",
                "",
                "Junzi Wu Ben, Ben Li Er Dao Sheng.",
                "Xiao Ti Ye Zhe, Qi Wei Ren Zhi Ben Yu!'"
            ],
            "zhushi": [
                ("You Zi (有子)", "Confucius's student You Ran, courtesy name Ziyou"),
                ("Xiao (孝)", "Filial piety - treating parents well"),
                ("Ti (悌)", "Fraternal respect - respecting elders"),
                ("Fan Shang (犯上)", "Offending superiors"),
                ("Wu Ben (务本)", "Devoting efforts to the root"),
                ("Ren (仁)", "Benevolence - the core of Confucius's thought")
            ],
            "translation_cn": "有子说：一个人孝顺父母、尊敬兄长，却喜欢冒犯尊长，那是很少的；不喜欢冒犯尊长，却喜欢作乱，那是从来没有的。君子致力于根本，根本建立了，道就产生了。孝顺父母、尊敬兄长，这就是仁爱的根本啊！",
            "translation_en": "Youzi said: 'A person who is filial to parents and respectful to elders rarely offends superiors. One who does not offend superiors yet likes to cause disorder - such a person has never existed. The noble person devotes himself to the root; when the root is established, the Way emerges. Filial piety and brotherly respect - these are the root of benevolence!'",
            "core_points": [
                ("1. Xiao Ti Wei Ren Ben", "Filial piety and fraternal respect are the root of benevolence - start from small matters, from family love to cultivate virtue."),
                ("2. Junzi Wu Ben", "The noble person devotes himself to the root; when the root is established, the Way emerges - grasp the essence in everything."),
                ("3. Xiu Shen Ci Xu", "Starting from the family, first filial to parents and respectful to elders, then extending to others - ultimately reaching the realm of 'Ren'.")
            ],
            "background": "Youzi was an important disciple of Confucius. Filial piety is the core of Chinese culture. Confucius made filial piety and fraternal respect the root of benevolence, which has profound social significance - family harmony is the foundation of social stability.",
            "western": [
                ("Mencius: Respect for Old and Care for Young", "Mencius developed Confucius's thought on filial piety, proposing 'respecting my elders to extend respect to others' - extending family affection to the whole society. c. 300 BCE"),
                ("Bentham: Utilitarianism", "Pursuing the greatest happiness for the greatest number - extending from family affection to social responsibility. 1789"),
                ("Rousseau: Natural Education", "Rousseau advocated education following nature, connecting with Confucian educational sequence starting from filial piety. 1762"),
                ("Fromm: Mature Love", "Love is an ability, first learning to love oneself and be grateful to parents, rather than depend or possess. 1956")
            ],
            "shiyi": "Xiao Ti ye zhe, qi wei ren zhi ben yu: Confucianism believes that 'Ren' is not an abstract concept but starts from concrete human relationships - starting from filial piety and fraternal respect to cultivate care for others. This deduction from small to large is Confucian methodology of self-cultivation. Modern psychology's Attachment Theory also proves that secure parent-child relationships are the prototype of all social relationships.",
            "discussion": [
                ("Question 1", "In modern society, what are new forms of 'filial piety'? How would you handle conflicts when parents' ideas clash with yours?"),
                ("Question 2", "Youzi said 'Junzi Wu Ben' - what do you think is the biggest problem of modern people in doing things? How to grasp the 'Ben'?"),
                ("Question 3", "Modern Western family relations emphasize equality - what aspects can be integrated with Confucian filial and fraternal thoughts?")
            ]
        },
        {
            "num": 3,
            "name": "Qiao Yan Ling Se (巧言令色)",
            "yuanwen": [
                "Zi Yue: 'Qiao yan ling se, xian yi ren!'"
            ],
            "zhushi": [
                ("Qiao Yan (巧言)", "Clever words - eloquent but insincere"),
                ("Ling Se (令色)", "Ingratiating expressions - smiling to please others"),
                ("Xian (鲜)", "Rare, few")
            ],
            "translation_cn": "孔子说：花言巧语、满脸堆笑的人，很少有仁德啊！",
            "translation_en": "Confucius said: 'Those with clever words and ingratiating expressions rarely have virtue!'",
            "core_points": [
                ("1.警惕伪善", "花言巧语、满脸堆笑的人，很少有真正的仁德——孔子最讨厌虚伪。"),
                ("2. 真诚为本", "真正的仁德不在于说话好听、表情讨好，而在于内心的真诚。"),
                ("3. 知人识人", "看一个人不要听他说什么，要看他怎么做；不要被表面现象迷惑。")
            ],
            "background": "This concise yet powerful saying reflects Confucius's profound insight into human nature. People with巧言令色 use language and expressions to please others while lacking genuine benevolence. Confucius devoted his life to restoring Zhou rituals and advocating sincerity, opposing formalism.",
            "western": [
                ("Aristotle: The Golden Mean", "Virtue is the mean between two extremes - courage is the mean between rashness and cowardice, honesty is the mean between cunning and bluntness. c. 350 BCE"),
                ("Kant: Categorical Imperative", "Act only according to maxims you can will to become universal laws - the moral value of an action lies in whether the motive is sincere. 1785"),
                ("Thoreau: False Society", "Most men live in quiet desperation - losing their true self in social performance. 1854"),
                ("de Botton: Status Anxiety", "Modern society's focus on status causes people to overcare about others' evaluations, leading to hypocrisy and anxiety. 2004")
            ],
            "shiyi": "巧言令色，鲜矣仁: This saying is even more cautionary in today's era - in the age of social media, people spend大量时间打造'人设', decorating their image. Confucius as early as 2500 years ago warned us: people with real virtue don't need巧言巧语 and ingratiating expressions to gain recognition. Inner strength and sincerity matter more than surface packaging.",
            "discussion": [
                ("Question 1", "Have you encountered someone with '巧言令色'? What was the final outcome? What inspiration did it give you?"),
                ("Question 2", "In the age of social media, what are the new forms of '巧言令色'? How do we maintain sincerity?"),
                ("Question 3", "Confucius said '鲜矣仁' - then what is sincere communication? Is there a boundary between '善巧方便' and '巧言令色'?")
            ]
        }
    ],
    # Day 2: 为政第二
    [
        {
            "num": 4,
            "name": "Wei Zheng Yi De (为政以德)",
            "yuanwen": [
                "Zi Yue: 'Wei Zheng Yi De,",
                "Pi Ru Bei Chen, Ju Qi Suo Er Zhong Xing Gong Zhi.'"
            ],
            "zhushi": [
                ("Wei Zheng Yi De (为政以德)", "Govern with virtue"),
                ("Bei Chen (北辰)", "North Star"),
                ("Gong (共)", "Same as '拱' - surround")
            ],
            "translation_cn": "孔子说：用道德来治理国家，就像北极星一样，处在自己的位置上，而众星都环绕着它。",
            "translation_en": "Confucius said: 'Govern with virtue, and the people will be like the stars revolving around the North Star - each in its proper place.'",
            "core_points": [
                ("1. 为政以德", "用道德而非刑罚来治理国家——领导者的德行感化民众，比强制更有效。"),
                ("2. 北辰比喻", "有德的领导者像北极星，自己安坐不动，众星自然环绕——以身作则的力量。"),
                ("3. 无为而治", "最高的领导是不需要刻意管理的，德行所至，人心自服——这也是道家的理想。")
            ],
            "background": "This is Confucius's famous 'Rule by Virtue' thought. During the Spring and Autumn period, many rulers relied on harsh laws. Confucius advocated governing by virtue, believing moral influence was more effective than legal coercion. The Han Dynasty adopted Confucian Rule by Virtue 200 years later, creating a 400-year golden age.",
            "western": [
                ("Machiavelli: The Prince", "The prince should make the people love and fear him, but if only one can be chosen, being feared is safer than being loved. 1532"),
                ("Rousseau: Social Contract", "The legitimacy of rulers comes from the consent of the ruled, not from virtue or lineage. 1762"),
                ("Jefferson: Democratic Thought", "The best government is one that governs least - connecting with Confucian 'Wu Wei Er Zhi'. 1776"),
                ("Greenleaf: Servant Leadership", "True leaders are servants, leading by example rather than giving orders. 1989")
            ],
            "shiyi": "为政以德: Modern management's 'Situational Leadership' and 'Servant Leadership' both connect with this. Excellent leaders don't rely on power to suppress but on virtue to inspire. Google's 'Project Oxygen' found that the best managers are those who first care about employee development and乐于助人成长的人。",
            "discussion": [
                ("Question 1", "What are the advantages and disadvantages of governing with '德' versus '法'? How do you think they should be balanced?"),
                ("Question 2", "In enterprise management, how is '为政以德' implemented? What specific management practices exist?"),
                ("Question 3", "'北辰居其所' - how can leaders maintain dignity while enabling teams to follow voluntarily?")
            ]
        },
        {
            "num": 5,
            "name": "Si Wu Xie (思无邪)",
            "yuanwen": [
                "Zi Yue: 'Shi San Bai, Yi Yan Yi Bi Zhi,",
                "Yue: 'Si Wu Xie.'\""
            ],
            "zhushi": [
                ("Shi San Bai (诗三百)", "The Book of Songs has 305 pieces, here as an approximate"),
                ("Bi (蔽)", "Summarize, conclude"),
                ("Si Wu Xie (思无邪)", "Pure thoughts, free from evil")
            ],
            "translation_cn": "孔子说：《诗经》三百首，用一句话来概括，就是'思想纯正'。",
            "translation_en": "Confucius said: 'The Book of Songs has three hundred pieces, yet they can be summed up in one phrase: having a heart free of evil.'",
            "core_points": [
                ("1. 诗教功能", "《诗经》能让人思想纯正——文学有教化人心的力量。"),
                ("2. 思无邪", "真正美好的诗歌，能净化人的心灵，让人思想纯正、情感真挚。"),
                ("3. 情感到理", "从真挚情感出发，到达道德境界——诗是情感的桥梁。")
            ],
            "background": "The Book of Songs is China's earliest collection of poetry, collecting 305 poems from the Western Zhou to mid-Spring and Autumn period. Confucius highly valued the Book of Songs, believing it could inspire, observe, socialize, and express grievances - an important tool for self-cultivation.",
            "western": [
                ("Plato: Noble Imitation", "Poetry should imitate noble actions and emotions, not stimulate vulgar desires. c. 380 BCE"),
                ("Aristotle: Catharsis", "Tragedy arouses pity and fear, allowing the audience's emotions to be purified (catharsis). c. 335 BCE"),
                ("Wordsworth: Poetry as Emotional Expression", "Poetry is the spontaneous overflow of powerful feelings, recalled in tranquility. 1800"),
                ("Freud: Sublimation", "Transforming instinctual impulses into socially approved高尚 activities - literary creation is a path of sublimation. 1908")
            ],
            "shiyi": "思无邪: Confucius's 'Poetry Education' thought remains significant today. Modern psychological research proves that reading literary classics can enhance empathy and strengthen emotional intelligence. A long-term study of 30,000 people showed that those who persistently read literary novels performed better on 'Theory of Mind' tests in interpersonal relationships.",
            "discussion": [
                ("Question 1", "What do you think is the modern 'poetry education'? How do literature, music, and film influence people's thinking?"),
                ("Question 2", "What kind of literary works are '思无邪'? What literary works should be criticized?"),
                ("Question 3", "From your reading experience, talk about literature's influence on your values formation.")
            ]
        },
        {
            "num": 6,
            "name": "Wen Gu Zhi Xin (温故知新)",
            "yuanwen": [
                "Zi Yue: 'Wen Gu Er Zhi Xin, Ke Yi Wei Shi Yi.'\""
            ],
            "zhushi": [
                ("Wen Gu (温故)", "Review learned knowledge"),
                ("Zhi Xin (知新)", "Gain new understanding and new knowledge"),
                ("Wei Shi (为师)", "Become a teacher")
            ],
            "translation_cn": "孔子说：温习学过的知识，并从中获得新的理解，这样的人可以当老师了。",
            "translation_en": "Confucius said: 'One who reviews what has been learned and gains new understanding may become a teacher.'",
            "core_points": [
                ("1. 温故为基", "学习首先要打好基础，反复温习是学习的必经之路——没有白学的知识。"),
                ("2. 知新是进", "在旧知识中发现新意义是进步的标志——同样的话在不同人生阶段有不同感悟。"),
                ("3. 师道尊严", "能'温故知新'的人才配当老师——老师不是知识的唯一来源，而是引导你发现的人。")
            ],
            "background": "'温故知新' is a core part of Confucius's educational philosophy. He believed learning is not one-time but a spiral-up process. The same content yields different understandings at different life stages. This is why the ancients emphasized '读书百遍，其义自见'. Modern cognitive science also proves this - spaced repetition is the most effective learning method.",
            "western": [
                ("Piaget: Cognitive Development", "Learning is the continuous reorganization of cognitive structures - new knowledge must connect with old knowledge to truly master. 1936"),
                ("Vygotsky: Zone of Proximal Development", "Learning occurs at the boundary between what is known and unknown - the teacher's role is to help students cross this zone. 1934"),
                ("Bruner: Discovery Learning", "The essence of learning is forming cognitive structures, deducing from known to unknown, from concrete to abstract. 1960"),
                ("Einstein: Review for Insight", "After learning, one must think; after thinking, one must ask questions - asking questions requires more imagination than solving them. 1955")
            ],
            "shiyi": "温故知新: Modern neuroscience finds that every time we recall a memory, we strengthen that neural pathway - memory is not passively stored but actively reconstructed. This is the 'Spacing Effect' - the more spaced the review, the deeper the memory. True learning is not一次性学完 but反复回顾, continuous deepening.",
            "discussion": [
                ("Question 1", "Do you have any experiences of '温故知新'? Did a certain principle give you different inspiration at different life stages?"),
                ("Question 2", "In the age of information explosion, why is '温故'反而更重要? How to maintain review in fragmented learning?"),
                ("Question 3", "'知新'的前提是'温故' - what inspiration does this give for your learning method?")
            ]
        }
    ],
    # Day 3: 八佾第三
    [
        {
            "num": 7,
            "name": "Shi Ke Ren Shu Bu Ke Ren (是可忍孰不可忍)",
            "yuanwen": [
                "Kongzi Wei Ji Shi:",
                "'Ba Yi Wu Yu Ting, Shi Ke Ren Ye, Shu Bu Ke Ren Ye?'"
            ],
            "zhushi": [
                ("Ji Shi (季氏)", "Ji Sun clan, senior official in Lu"),
                ("Ba Yi (八佾)", "Dance formation for the Emperor - 64 people; Ba Yi is the rank of dukes"),
                ("Ren (忍)", "Tolerate, endure; also means '忍心' (have the heart to)")
            ],
            "translation_cn": "孔子谈到季孙氏说：他用八佾的规格在庭院奏乐舞蹈，这样的事都忍心做得出来，还有什么事是他不忍心做的呢？",
            "translation_en": "Confucius commented on the Ji family: 'Using the eight-row formation in their hall - if this can be tolerated, what cannot? This was a violation of ritual propriety.'",
            "core_points": [
                ("1. 礼崩乐坏", "季氏用天子规格的乐舞，是公然僭越礼制——春秋末期礼崩乐坏的缩影。"),
                ("2. 是可忍孰不可忍", "孔子对此极为愤慨——礼制不是形式，而是社会秩序的基石。"),
                ("3. 礼的根本", "礼制背后是秩序和尊重，破坏礼制就是破坏社会根基——千里之堤，溃于蚁穴。")
            ],
            "background": "Ba Yi is a dance formation exclusive to the Emperor. Dukes could only use Liu Yi (36 people), and senior officials only Si Yi (16 people). Ji family, merely senior officials, used Ba Yi - a blatant challenge to Zhou rituals. Confucius, devoted to restoring Zhou rituals, was naturally extremely indignant. This chapter gives its name to the entire 'Ba Yi' section.",
            "western": [
                ("Plato: The Ideal State", "Each class performs its own function for social harmony - usurping classes disrupts the entire social order. c. 380 BCE"),
                ("Durkheim: Social Division of Labor", "Social organic solidarity is based on everyone fulfilling their social roles - role confusion leads to social disorder. 1893"),
                ("Montesquieu: Spirit of Laws", "Law reflects a society's basic order - breaking the law is breaking the social contract. 1748"),
                ("Rawls: Theory of Justice", "Social rules should be acceptable even to the most disadvantaged - unfair rules will eventually collapse. 1971")
            ],
            "shiyi": "是可忍孰不可忍: This phrase is widely used today but often deviates from its original meaning. Confucius's anger was not because his personal interests were harmed, but because rituals represented the foundation of social order. When order is broken without制止, the foundation of entire society shakes.",
            "discussion": [
                ("Question 1", "How do you understand '礼' in modern society? What's the difference between etiquette and ritual system?"),
                ("Question 2", "When you see violations without punishment, what do you think? What are manifestations of 'Broken Window Effect' in life?"),
                ("Question 3", "In organizational management, how to establish and maintain systems? Which is more important, systems or culture?")
            ]
        },
        {
            "num": 8,
            "name": "Yu Qi She Ye Ning Jian (与其奢也宁俭)",
            "yuanwen": [
                "Lin Fang Wen Li Zhi Ben.",
                "Zi Yue: 'Da Zai Wen!",
                "Li, Yu Qi She Ye, Ning Jian;",
                "Sang, Yu Qi Yi Ye, Ning Qi.'\""
            ],
            "zhushi": [
                ("Lin Fang (林放)", "Lu citizen, disciple of Confucius"),
                ("Li Zhi Ben (礼之本)", "The essence/fundament of ritual"),
                ("Yi (易)", "Perfectly arranged, elaborate ritual"),
                ("Qi (戚)", "Inner grief")
            ],
            "translation_cn": "林放问礼的本质。孔子说：这个问题意义重大！就一般礼仪来说，与其奢侈浪费，宁可节俭朴素；就丧礼来说，与其仪式完备，宁可内心悲哀。",
            "translation_en": "Lin Fang asked about the essence of ritual. Confucius said: 'What a great question! For ceremonies, plainness is better than extravagance; for funerals, grief is more important than elaborate ritual.'",
            "core_points": [
                ("1. 礼的本质", "礼的根本不在形式，而在于内心——形式只是内心情感的外在表达。"),
                ("2. 宁俭勿奢", "宁可简朴也不要奢侈——礼的精神在于真诚，不在于排场。"),
                ("3. 丧礼重情", "丧礼最重要的是内心悲哀，仪式再完美如果内心没有悲伤，也是徒有其表。")
            ],
            "background": "Lin Fang asked about the fundament of ritual, and Confucius greatly praised this question because it got to the heart of the matter. Confucius opposed formalism throughout his life. He saw nobles overly focusing on ritual forms while neglecting the inner spirit of ritual. This reveals the core of Confucian ritual scholarship - the essence of ritual is sincere inner emotion, and external forms are merely a means, not an end.",
            "western": [
                ("Rousseau: Return to Nature", "Civilization is a deviation from natural state - over-ritualization loses sincerity. 1762"),
                ("Kierkegaard: Subjectivity of Faith", "Objective truth doesn't matter; what matters is subjective commitment - inner sincerity over external form. 1847"),
                ("Camus: Absurdity and Rebellion", "True rebellion is clear recognition of the absurd, not surface confrontation - also an inner spirit. 1942"),
                ("Hemingway: Iceberg Theory", "The art of writing lies in omission - the best expression conveys the most with the least words. 1954")
            ],
            "shiyi": "礼，与其奢也，宁俭: This phrase has broad applications in modern life. At weddings, funerals, and celebrations, people often pursue luxurious displays while neglecting the true purpose of ritual. When weddings become arenas for comparison, when funerals become performances, ritual loses its meaning.",
            "discussion": [
                ("Question 1", "Have you attended weddings or funerals that you felt were 'overly luxurious'? What was your feeling about that formalism?"),
                ("Question 2", "In modern society, what is the 'essence' of '礼'? How to maintain sincerity in social interactions without becoming mere formality?"),
                ("Question 3", "You said '与其易也宁戚' - if one feels grief inside but doesn't express it 'correctly', has that person fulfilled filial duty?")
            ]
        },
        {
            "num": 9,
            "name": "Gao Shuo Xi Yang (告朔饩羊)",
            "yuanwen": [
                "Zi Gong Yu Qu Gao Shuo Zhi Xi Yang.",
                "Zi Yue: 'Ci Ye! Er Ai Qi Yang,",
                "Wo Ai Qi Li.'\""
            ],
            "zhushi": [
                ("Gao Shuo (告朔)", "Ancient ritual where lords sacrificed to ancestors on the first day of each month"),
                ("Xi Yang (饩羊)", "Live sheep used for sacrifice"),
                ("Ai (爱)", "Reluctant, grudge")
            ],
            "translation_cn": "子贡想要去掉每月初一告祭祖庙用的活羊。孔子说：赐啊！你舍不得那只羊，我舍不得那种礼。",
            "translation_en": "Zigong wished to discontinue the sacrificial sheep for the monthly sacrificial rite. Confucius said: 'Si, you grudge the sheep, but I grudge the rite.'",
            "core_points": [
                ("1. 形式与内容", "羊只是形式，但形式承载着内容——去掉形式，内容也难以为继。"),
                ("2. 我爱其礼", "孔子重视的不是羊，而是礼所代表的精神——礼的形式是精神的载体。"),
                ("3. 保存仪节", "有些形式即使看起来无用，但它是维持精神的一种方式——仪节不可轻易废去。")
            ],
            "background": "Gao Shuo was an important ancient ritual system - on the first day of each month, the lord would sacrifice to ancestors, showing respect for ancestors. Zigong, from an economic perspective, felt killing sheep was too wasteful and wanted to abolish it. Confucius believed that although the sheep seems useless, it reminds people of ritual's existence - once forms are removed, the spirit of ritual gradually disappears.",
            "western": [
                ("Durkheim: Collective Ritual", "Religious and collective rituals are sources of social cohesion - they let individuals feel they are part of a whole. 1912"),
                ("Scheler: Symbolic Anthropology", "Human symbols and rituals create a 'second reality' beyond physical existence. 1923"),
                ("Turner: Ritual Process", "Rituals convey social values through symbolic symbols - participants gain shared meaning. 1967"),
                ("Geertz: Local Knowledge", "Culture is a symbol system; understanding a culture requires understanding locals' own interpretations. 1973")
            ],
            "shiyi": "尔爱其羊，我爱其礼: Form is not content, but it carries content. When we abolish a certain ritual, we often simultaneously abolish the spirit it represents. In modern society, many traditional customs are simplified or abolished - sometimes because forms have lost their content, but more often because we shortsightedly only see the cost of forms without seeing the spiritual value they carry.",
            "discussion": [
                ("Question 1", "Which 'rituals' in modern life do you think are meaningful? Which are mere formalism?"),
                ("Question 2", "Should company or organizational 'team building' activities be kept or reformed? How to make them more than just going through the motions?"),
                ("Question 3", "'我爱其礼' - Confucius's attitude toward form, what inspiration does it give for understanding traditional culture?")
            ]
        }
    ],
    # Day 4: 里仁第四
    [
        {
            "num": 10,
            "name": "Li Ren Wei Mei (里仁为美)",
            "yuanwen": [
                "Zi Yue: 'Li Ren Wei Mei.",
                "Ze Bu Chu Ren, Yan De Zhi?'"
            ],
            "zhushi": [
                ("Li (里)", "Residence, living"),
                ("Li Ren (里仁)", "Living in a place with benevolent people"),
                ("Chu (处)", "Live, choose residence"),
                ("Zhi (知)", "Same as '智' - wisdom")
            ],
            "translation_cn": "孔子说：居住在有仁德的地方才是美好的。如果选择住处而不选择有仁德的地方，怎能算是有智慧呢？",
            "translation_en": "Confucius said: 'To live in a neighborhood of benevolence is ideal. If one chooses not to live among the benevolent, how can one be considered wise?'",
            "core_points": [
                ("1. 环境育人", "居住环境对人的影响巨大——与仁德之人为邻，自己也会受熏陶。"),
                ("2. 择邻而居", "选择住处就是选择朋友圈——你想成为什么样的人，就要与什么样的人为伍。"),
                ("3. 智慧之选", "选择住处不仅是经济决策，更是人生投资——近朱者赤，近墨者黑。")
            ],
            "background": "This reflects Confucian emphasis on environment. The story of 'Mengmu San Qian' (Mencius's mother moving three times) embodies this thought - Mencius's mother moved three times to give Mencius a good growth environment. Modern pedagogy also proves that students' performance is highly related to their peers' quality - this is the 'Peer Effect'.",
            "western": [
                ("Aristotle: Civic Virtue", "Man is by nature a political animal, able to achieve good only within the polis - residential environment determines social relations. c. 340 BCE"),
                ("Maslow: Hierarchy of Needs", "Human needs range from physiological to self-actualization - social environment provides or restricts fulfillment of these needs. 1943"),
                ("Bronfenbrenner: Ecological Systems Theory", "Human development is nested in multiple layers of environmental influence - micro, meso, and macro environments jointly shape individuals. 1979"),
                ("Putnam: Social Capital", "Social networks have value - the quality of one's social connections significantly affects career development, health, and happiness. 2000")
            ],
            "shiyi": "里仁为美: This saying has new meaning in modern times. When choosing where to live, not only consider school districts and transportation, but also community culture and neighbors' quality. Research shows that one's social circle significantly affects career development, health, and happiness - this is called 'Social Capital'.",
            "discussion": [
                ("Question 1", "What influence does your current social circle have on you? What are your standards for choosing friends?"),
                ("Question 2", "When choosing a job and city, how is the principle '里仁为美' applied?"),
                ("Question 3", "In modern society, through online communities, geographical factors have less impact on friendship - what does this mean?")
            ]
        },
        {
            "num": 11,
            "name": "Ren Zhe An Ren (仁者安仁)",
            "yuanwen": [
                "Zi Yue: 'Bu Ren Zhe Bu Ke Yi Jiu Chu Yue,",
                "Bu Ke Yi Chang Chu Le.",
                "Ren Zhe An Ren, Zhi Zhe Li Ren.'"
            ],
            "zhushi": [
                ("Yue (约)", "Poverty, simplicity"),
                ("Jiu Chu Yue (久处约)", "Long-term poverty"),
                ("Chang Chu Le (长处乐)", "Long-term prosperity"),
                ("An Ren (安仁)", "At peace with benevolence, finding joy in benevolence"),
                ("Li Ren (利仁)", "Benefit from benevolence, gaining from benevolence")
            ],
            "translation_cn": "孔子说：没有仁德的人，不能长期处于穷困中，也不能长期处于安乐中。只有仁者才能安于仁道，智者才能利用仁道获得好处。",
            "translation_en": "Confucius said: 'One without benevolence cannot endure poverty for long, nor can one endure prosperity for long. The benevolent find peace in benevolence; the wise benefit from benevolence.'",
            "core_points": [
                ("1. 不仁者的困境", "没有仁德的人在穷困时会为非作歹，在安乐时会骄奢淫逸——两种状态都hold不住。"),
                ("2. 仁者安仁", "真正的仁者不以贫富为转移，安于本心，不为外境所动——这是内心的安定。"),
                ("3. 知者利仁", "有智慧的人看到仁道的好处，所以主动行仁——仁与利并不矛盾。")
            ],
            "background": "This deeply reveals the relationship between 'Ren' and psychological stability. Those without benevolence are like psychologically unstable people - they'll have problems in both favorable and adverse circumstances. The benevolent are different; no matter how the environment changes, they have a stable anchor in their heart. This thought aligns with modern psychology's 'Resilience' concept.",
            "western": [
                ("Aristotle: Happiness as Activity", "Happiness is not a state but an activity - the benevolent person's peace comes from inner virtue. c. 340 BCE"),
                ("Stoicism: Internal Control", "We cannot control external events but can control our reactions to them - connecting with Confucius's 'An Ren'. c. 300 CE"),
                ("Seneca: On the Shortness of Life", "People busy all day have no time to think; only those with inner peace can truly live. c. 49 CE"),
                ("Frankl: Living Out Meaning", "The ultimate freedom is the freedom to choose one's attitude - one can maintain spiritual freedom even in extreme circumstances. 1946")
            ],
            "shiyi": "不仁者不可以久处约，不可以长处乐: This saying is highly cautionary in modern society. Many people can endure hardship when poor but become lost when wealthy; or they're happy when wealthy but collapse when encountering difficulties. This shows the real problem is not the environment but whether one has the foundation of inner virtue.",
            "discussion": [
                ("Question 1", "Have you seen examples of '不仁者' becoming corrupt after becoming wealthy? What warning does this give you?"),
                ("Question 2", "'仁者安仁' - how to cultivate inner peace? When the external environment changes drastically, how do you maintain inner stability?"),
                ("Question 3", "Some believe '仁' and '利' are contradictory - what's your view? Does '知者利仁' make sense?")
            ]
        },
        {
            "num": 12,
            "name": "Wei Ren Zhe Neng Hao Ren (唯仁者能好人)",
            "yuanwen": [
                "Zi Yue: 'Wei Ren Zhe Neng Hao Ren, Neng Wu Ren.",
                "Gou Zhi Yu Ren Yi, Wu Wu Ye.'"
            ],
            "zhushi": [
                ("Hao Ren (好人)", "Love people, appreciate people"),
                ("Wu Ren (恶人)", "Resent people, criticize people"),
                ("Gou (苟)", "If, suppose"),
                ("Zhi Yu Ren (志于仁)", "Set one's mind on benevolence"),
                ("Wu Wu (无恶)", "Will not do evil things")
            ],
            "translation_cn": "孔子说：只有仁者才能真正喜爱人，也能真正厌恶人。如果真正立志行仁，就不会去做坏事了。",
            "translation_en": "Confucius said: 'Only the benevolent can truly love people and truly resent people. One who is truly set on benevolence will do no wrong.'",
            "core_points": [
                ("1. 仁者的爱恨", "真正的仁者爱憎分明——不是没有原则的老好人，而是有明确的价值判断。"),
                ("2. 能好人能恶人", "真正能欣赏人，也能批评人——知道什么是好的，什么是不好的。"),
                ("3. 志于仁无恶", "真正立志行仁的人，不会做坏事——仁是内在的约束力。")
            ],
            "background": "'Ren' is not unprincipled tolerance but has clear value judgments. Confucius explicitly pointed out that only the benevolent can truly love and truly resent. This shows 'Ren' has principles - it's not undifferentiated universal love. The benevolent person's love is sincere, and their resentment is sincere - both are sincere expressions from the benevolent heart.",
            "western": [
                ("Confucius: Requite Enmity with Straightness", "Confucianism is not unprincipled tolerance but treating good and evil with an upright heart - requite enmity with straightness, requite kindness with kindness. c. 500 BCE"),
                ("Aristotle: The Mean and Emotion", "Virtue is the appropriate control of emotions - courage is the mean between fear and recklessness, temperance is the mean between indulgence and insensibility. c. 340 BCE"),
                ("Kant: Moral Feeling", "Respect is the foundation of morality - respecting others as ends, not means. Connecting with '仁者能好人能恶人'. 1785"),
                ("Rawls: Sense of Justice", "Sense of justice is the emotional identification with fair institutions - true justice includes anger at injustice. 1971")
            ],
            "shiyi": "唯仁者能好人，能恶人: This reminds us that true benevolence is not unprincipled indulgence. The benevolent person's love is sincere appreciation, and their resentment is sincere criticism - both come from inner sincere judgment. Today, many lack clear values, are polite to everyone, actually mixing good and evil. This 'compromise' attitude seems '仁' but is actually '不仁'.",
            "discussion": [
                ("Question 1", "Why is '能恶人' also a characteristic of the benevolent? Is there a time when 'malice' is justified?"),
                ("Question 2", "In modern workplace, how to maintain a benevolent heart while adhering to principles and not 'compromising'?"),
                ("Question 3", "'苟志于仁矣，无恶也' - why can setting one's mind on benevolence prevent doing bad things? How much约束力 does ambition have on behavior?")
            ]
        }
    ],
    # Day 5: 公冶长第五
    [
        {
            "num": 13,
            "name": "Zi Jian Zhi Xian (子贱之贤)",
            "yuanwen": [
                "Zi Wei Zi Jian:",
                "'Junzi Zai Ruo Ren!",
                "Lu Wu Junzi Zhe, Si Yan Qu Si?'"
            ],
            "zhushi": [
                ("Zi Jian (子贱)", "Mi Buqi, courtesy name Zijian, disciple of Confucius"),
                ("Ruo Ren (若人)", "This person"),
                ("Si Yan Qu Si (斯焉取斯)", "From where did he acquire such virtue?")
            ],
            "translation_cn": "孔子评价子贱说：这个人是君子啊！如果鲁国没有君子，他从哪里获得这样的品德呢？",
            "translation_en": "Confucius commented on Zijian: 'What a noble person! If there were no gentlemen in Lu, where would he have acquired such virtue?'",
            "core_points": [
                ("1. 子贱之贤", "子贱是孔子的弟子，以德行著称——孔子称赞他为'君子'。"),
                ("2. 环境影响", "孔子认为子贱的贤德来自鲁国的环境——鲁国当时有众多君子。"),
                ("3. 环境育人", "一个人能否成为君子，与他所处的环境密切相关——鲁多君子，所以子贱能成君子。")
            ],
            "background": "Zijian (520-445 BCE), named Mi Buqi, was one of Confucius's important disciples. He served as the official of Shanfu and governed it well, leaving the famous story 'Zijian Governs Shanfu'. Confucius believed Zijian became a junzi precisely because Lu had a good cultural environment - this embodies the '里仁为美' thought.",
            "western": [
                ("Mencius: Circumstances Over Talent", "With the right circumstances, anything is possible - environment has a decisive influence on a person's success. c. 300 BCE"),
                ("Bacon: Knowledge is Power", "Knowledge changes destiny, but acquiring knowledge requires environment and education - junzi cultivation requires cultural soil. 1597"),
                ("Bandura: Observational Learning", "People learn through observing others' behaviors and their consequences - models in the environment shape personality. 1977"),
                ("Csikszentmihalyi: Systems View of Creativity", "Individual achievement cannot be separated from the cultural system - environment provides the 'symbolic pool' for innovation. 1999")
            ],
            "shiyi": "鲁无君子者，斯焉取斯: One's achievement is inseparable from their environment. Zijian became a junzi not only due to personal effort but also because Lu had good cultural traditions and many junzi as examples. This reminds us: to become whom you want to be, choose what kind of environment you're in.",
            "discussion": [
                ("Question 1", "In the environment you grew up, who influenced you most? Which of their qualities shaped you?"),
                ("Question 2", "'鲁无君子者，斯焉取斯' - do you think environment or individual factors are more important for a person's success?"),
                ("Question 3", "In modern society, how to proactively choose and shape an environment conducive to your growth?")
            ]
        },
        {
            "num": 14,
            "name": "Kongzi Ping Qi Diao Kai (孔子评漆雕开)",
            "yuanwen": [
                "Zi Shi Qi Diao Kai Shi.",
                "Dui Yue: 'Wu Si Zhi Wei Neng Xin.'",
                "Zi Yue."
            ],
            "zhushi": [
                ("Qi Diao Kai (漆雕开)", "Disciple of Confucius, courtesy name Zikai"),
                ("Shi (使)", "Let, send"),
                ("Shi (仕)", "Take office, become an official"),
                ("Wu Si Zhi Wei Neng Xin (吾斯之未能信)", "I am not yet confident about this")
            ],
            "translation_cn": "孔子让漆雕开去做官。他回答说：我对这件事还没有信心。孔子听了很高兴。",
            "translation_en": "Confucius asked Qidiao Kai to take an official post. He replied: 'I am not yet confident about this.' Confucius was pleased.",
            "core_points": [
                ("1. 谦虚谨慎", "漆雕开对自己的能力表示怀疑，这种谦虚是儒家推崇的美德——知道自己的不足是进步的开始。"),
                ("2. 孔子悦之", "孔子对漆雕开的回答感到高兴——因为这说明漆雕开有自知之明，不轻率出仕。"),
                ("3. 信而后仕", "做官要有充分的准备和信心，不打无把握之仗——机会来临时，要确保自己真正准备好了。")
            ],
            "background": "Qi Diao Kai (545 BCE ?), named Qi Diao Kai, courtesy name Zikai, was a disciple of Confucius. He later became known as a 'scholar junzi' for his humility. When Confucius asked him to take an official post, he said he wasn't yet confident. Confucius was not disappointed but pleased - because this showed Qi Diao Kai had self-knowledge and was cautious, not eager or ignorant of the depth involved.",
            "western": [
                ("Socrates: Know Thyself", "Self-knowledge is the knowledge of one's ignorance - knowing you don't know is wiser than thinking you know. c. 400 BCE"),
                ("Kant: Rational Critique", "Reason must first recognize its own limitations to be used correctly - self-knowledge is the beginning of philosophy. 1781"),
                ("Russell: Boundary of Knowledge", "The most important knowledge is the knowledge of one's ignorance - the core of Socratic wisdom. 1950"),
                ("Darwin: Cognitive Bias", "Humans tend to overestimate their abilities - recognizing this is itself a kind of progress. 1871")
            ],
            "shiyi": "吾斯之未能信: In modern society, eager to perform and succeed is a common problem. Qi Diao Kai's humility tells us: before an opportunity, not eager to seize it, but first ask yourself - am I ready? Confidence in something requires real strength's support. Confucius was pleased precisely with Qi Diao Kai's attitude of being responsible for himself and for his career.",
            "discussion": [
                ("Question 1", "Have you encountered situations where 'opportunity came but I wasn't ready'? How did you handle it?"),
                ("Question 2", "'吾斯之未能信' - how to distinguish this humility from 'daring to take responsibility'? Where is the boundary?"),
                ("Question 3", "In career choices, what conditions are needed to build 'confidence'? How do you judge if you're truly ready?")
            ]
        },
        {
            "num": 15,
            "name": "Zi Lu Wu Su Xie (子路无宿诺)",
            "yuanwen": [
                "Zi Lu You Wen, Wei Zhi Neng Xing,",
                "Wei Kong You Wen."
            ],
            "zhushi": [
                ("You Wen (有闻)", "Heard a principle"),
                ("Wei Zhi Neng Xing (未之能行)", "Not yet able to practice it"),
                ("Wei Kong You Wen (唯恐有闻)", "Only fears hearing new principles")
            ],
            "translation_cn": "子路听到了一个道理，如果还没有能够去实践，就只怕又听到新的道理了。",
            "translation_en": "When Zilu learned something, he would put it into practice before learning more - for he feared he might not be able to act on what he heard.",
            "core_points": [
                ("1. 知行合一", "子路听到道理就要去实践，不是听完就算了——知识必须转化为行动。"),
                ("2. 恐有新闻", "子路怕听到新道理是因为旧道理还没实践——不是贪婪新知，而是担心行动跟不上。"),
                ("3. 先行后学", "先把当前的事情做好，再去学新的——这是一种务实的学习态度。")
            ],
            "background": "Zilu (542-480 BCE), named Zhong You, courtesy name Zilu, was one of Confucius's most famous disciples. He was straightforward and action-oriented, a paragon of Confucian practicality. This passage reflects Zilu's 'act first, learn later' characteristic - not that he didn't learn, but that he first uses what he's learned before learning more. This fully aligns with Confucius's '学而时习之' philosophy.",
            "western": [
                ("Wang Yangming: Unity of Knowledge and Action", "Knowledge is the beginning of action; action is the completion of knowledge - knowledge without action is not true knowledge. 1500s"),
                ("Dewey: Learning by Doing", "Education is not preparation for life; education is life itself - learning must be completed in practice. 1897"),
                ("Wittgenstein: Language Games", "The meaning of words lies in their use - language learning detached from practice is meaningless. 1953"),
                ("Becker: Human Capital", "On-the-job training is more effective than full-time study - because learning is immediately applied, efficiency is highest. 1964")
            ],
            "shiyi": "唯恐有闻: Modern people receive huge amounts of information daily, but truly transforming it into action is rare - this is 'information anxiety'. Zilu's attitude is the opposite - he's not worried about learning too little but about not doing enough. As the saying goes: 'Knowing so many principles, still can't live a good life' - the problem is not how much we know but how much we've done.",
            "discussion": [
                ("Question 1", "'唯恐有闻' - are you more worried about 'not learning enough' or 'not doing enough'? What does this priority difference mean?"),
                ("Question 2", "In the age of information explosion, how to avoid the dilemma of 'learned a lot but did nothing'?"),
                ("Question 3", "From your experience, why is 'unity of knowledge and action' so difficult? How can it truly be achieved?")
            ]
        }
    ],
    # Day 6: 雍也第六
    [
        {
            "num": 16,
            "name": "Bu Qian Nu Bu Er Guo (不迁怒不贰过)",
            "yuanwen": [
                "Ai Gong Wen: 'Di Zi Shu Wei Hao Xue?'",
                "Kongzi Dui Yue: 'You Yan Hui Zhe Hao Xue,",
                "Bu Qian Nu, Bu Er Guo,",
                "Bu Xing Duan Ming Si Yi.",
                "Jin Ye Ze Wang, Wei Wen Hao Xue Zhe Ye.'\""
            ],
            "zhushi": [
                ("Ai Gong (哀公)", "Duke of Lu"),
                ("Yan Hui (颜回)", "Yan Yuan, courtesy name Ziyan, Confucius's most favored disciple"),
                ("Bu Qian Nu (不迁怒)", "Not transferring anger to others"),
                ("Bu Er Guo (不贰过)", "Not repeating the same mistake"),
                ("Duan Ming (短命)", "Yan Hui died early, only about 30 years old"),
                ("Wang (亡)", "Same as '无' - none, not")
            ],
            "translation_cn": "鲁哀公问：你的学生中谁是最好学的？孔子回答说：有颜回这样的人好学，他不把怒气转移到别人身上，不重复犯同样的错误。他不幸短命死了，现在没有这样的人了，没听说过还有好学的人。",
            "translation_en": "Duke Ai asked: 'Which of your disciples is most devoted to learning?' Confucius replied: 'Yan Hui - devoted to learning. He never transferred his anger, never repeated a mistake. Unfortunate indeed that he died young. Now there is no one like him; I have not heard of anyone else who loves learning so deeply.'",
            "core_points": [
                ("1. 颜回之学", "颜回是孔子最得意的弟子，他的'好学'不是学得多，而是不迁怒、不贰过——这是人格的完善。"),
                ("2. 不迁怒", "遇到挫折不向他人发泄情绪——这是情绪管理的能力，是成熟人格的标志。"),
                ("3. 不贰过", "同样的错误不犯第二次——这是从失败中学习的能力，是成长的关键。")
            ],
            "background": "Yan Hui (521-490 BCE), courtesy name Ziyan, was Confucius's most favored disciple, honored as 'Fu Sheng' (the Second Saint). Though he died early, his learning attitude and personal cultivation became a model for later generations. 'Bu Qian Nu, Bu Er Guo' seems simple but is extremely difficult to practice - many people spend their whole lives repeating the same mistakes, never truly reflecting.",
            "western": [
                ("Aristotle: Ethical Virtue", "Virtue is expressing emotions at the right time, to the right people, in the right way. Anger should be expressed when it should be, and not exceeding necessary limits. c. 340 BCE"),
                ("Seligman: Positive Psychology", "Learn explanatory style - attribute failures to temporary, specific causes rather than permanent, universal ones. 2002"),
                ("Markus: Possible Selves", "Our self-concept determines how we respond to failure - growth mindset rather than fixed mindset. 1990s"),
                ("Dalio: Pain Plus Reflection Equals Progress", "Pain is a signal telling you need to change; reflection is the key to learning from pain. 2017")
            ],
            "shiyi": "不迁怒不贰过: This is Confucius's definition of 'good learning' - not how much knowledge is learned, but whether personal perfection is achieved. 'Bu Qian Nu' tests emotional management ability; 'Bu Er Guo' tests the ability to learn from failure. Research shows that people's IQ drops 40% when emotionally agitated, and 90% think they're above average at avoiding same mistakes - these two problems combined explain why most can't achieve 'bu er guo'.",
            "discussion": [
                ("Question 1", "Who are you most likely to '迁怒' to? Under what circumstances do you transfer negative emotions to others?"),
                ("Question 2", "'不贰过' why is it so difficult? Have you repeated the same mistake? How can it truly be achieved?"),
                ("Question 3", "Confucius said '今也则亡' - do you think there are still people like Yan Hui who are so 'devoted to learning' today? What does good learning mean today?")
            ]
        },
        {
            "num": 17,
            "name": "Xian Zai Hui Ye (贤哉回也)",
            "yuanwen": [
                "Zi Yue: 'Xian Zai, Hui Ye!",
                "Yi Dan Shi, Yi Piao Yin,",
                "Zai Lou Xiang, Ren Bu Kan Qi You,",
                "Hui Ye Bu Gai Qi Le.",
                "Xian Zai, Hui Ye!'"
            ],
            "zhushi": [
                ("Dan Shi (箪食)", "Basket of food"),
                ("Piao Yin (瓢饮)", "Gourd of water"),
                ("Lou Xiang (陋巷)", "Simple alley"),
                ("Bu Kan Qi You (不堪其忧)", "Cannot bear that worry"),
                ("Bu Gai Qi Le (不改其乐)", "Does not change his joy - inner peace, not moved by external things")
            ],
            "translation_cn": "孔子说：多么贤德啊，颜回！一竹筐饭，一瓢水，住在简陋的巷子里，别人都不能忍受那种忧愁，颜回却不改变他的快乐。多么贤德啊，颜回！",
            "translation_en": "Confucius said: 'How worthy is Yan Hui! Living on a basket of food and a gourd of water in a mean alley - others would find such poverty unbearable, yet Yan Hui does not alter his joy. How worthy is Yan Hui!'",
            "core_points": [
                ("1. 颜回之贤", "颜回的贤德体现在'不改其乐'——无论环境如何，内心始终保持快乐。"),
                ("2. 安贫乐道", "物质匮乏不能动摇精神追求——真正的高人不在物质，在精神。"),
                ("3. 人不堪其忧", "普通人受不了这种贫困，颜回却乐在其中——境界不同，感受不同。")
            ],
            "background": "This is one of the most moving passages in the Analects. Though Yan Hui was materially lacking, he was spiritually very rich - the 'joy' in '不改其乐' is not material pleasure but the joy of 'Dao'. Confucius late in life described Yan Hui as 'I see him progressing, never see him stopping'. This spiritual richness beyond material deficiency is 'Kong Yan Le Chu' (Confucius and Yan Hui's joy) that later Song dynasty Neo-Confucians developed into the cultivation practice of 'Seeking Kong Yan Le Chu'.",
            "western": [
                ("Epicurus: Hedonism", "True pleasure is not sensual indulgence but tranquility of the mind - satisfaction of desires brings temporary pleasure, inner peace brings lasting happiness. c. 300 BCE"),
                ("Stoicism: Inner Tranquility", "External things cannot harm you; what harms you is your judgment - one can maintain inner tranquility even in the most difficult circumstances. c. 100 CE"),
                ("Rousseau: Return to Nature", "Progress of civilization brings misfortune; true happiness lies in returning to natural state - natural simple life surpasses artificial luxury. 1762"),
                ("Maslow: Hierarchy of Needs", "When lower-level needs (physiological, safety) cannot be met, people can still pursue higher-level needs (respect, self-actualization) - material deficiency does not equal spiritual deficiency. 1943")
            ],
            "shiyi": "不改其乐: Modern people always link happiness with material conditions, but Yan Hui's story tells us true happiness comes from within, not external environment. Psychological research shows material conditions have diminishing marginal utility on happiness - annual income from 50K to 500K, happiness improves limitedly; but inner peace and sense of control over life have huge impact on happiness.",
            "discussion": [
                ("Question 1", "Have you experienced a period of 'material deficiency but spiritual richness'? What was that experience like?"),
                ("Question 2", "What's the difference between '不改其乐' and '知足常乐'? Which do you think is better?"),
                ("Question 3", "In an era pursuing material things, how to maintain '安贫乐道' attitude? Is this an ideal or a form of escape?")
            ]
        },
        {
            "num": 18,
            "name": "Ran Bo Niu You Ji (冉伯牛有疾)",
            "yuanwen": [
                "Bo Niu You Ji, Zi Wen Zhi,",
                "Zi Yong Chuan Qi Shou, Yue:",
                "'Wang Zhi, Ming Yi Fu!",
                "Si Ren Ye Er You Si Ji Ye!",
                "Si Ren Ye Er You Si Ji Ye!'"
            ],
            "zhushi": [
                ("Bo Niu (伯牛)", "Ran Geng, courtesy name Boniu, disciple of Confucius"),
                ("You Ji (有疾)", "Had a serious illness - possibly leprosy"),
                ("Zi Yong (自牖)", "From the window - Boniu didn't want people entering the room"),
                ("Wang Zhi (亡之)", "About to lose him"),
                ("Ming Yi Fu (命矣夫)", "This is fate!")
            ],
            "translation_cn": "冉伯牛得了恶疾，孔子去探望他，从窗户握着他的手，说：要失去了，这是命啊！这样的人却得了这样的病！这样的人却得了这样的病！",
            "translation_en": "Ran Boniu fell ill. Confucius visited him, holding his hand through the window, saying: 'Losing him - is this not fate? Such a person, yet such an illness! Such a person, yet such an illness!'",
            "core_points": [
                ("1. 师生情深", "孔子对冉伯牛的病情深表痛惜——师生感情深厚，孔子不能进屋却仍坚持探望。"),
                ("2. 命矣夫", "孔子感叹'这是命啊'——面对无法改变的命运，孔子表现出无奈但不绝望的态度。"),
                ("3. 斯人有斯疾", "这么好的人却得了这样的病——孔子对命运的不公发出感叹，但接受而不抱怨。")
            ],
            "background": "Ran Boniu (544 BCE ?), named Ran Geng, courtesy name Boniu, was a disciple of Confucius, known for his virtue. He contracted what was considered a contagious disease (possibly leprosy) at the time, so people couldn't enter his room. Confucius held his hand through the window, expressing deep regret. Confucius said 'Si Ren Ye You Si Ji' twice, showing his inner commotion. This is one of the most touching teacher-student relationships in the Analects.",
            "western": [
                ("Seneca: On the Shortness of Life", "Life is not valued by its length but by its density - even a short life can be fulfilling. c. 49 CE"),
                ("Epicurus: Philosophy of Death", "Death has nothing to do with us - when we exist, death is not present; when death arrives, we no longer exist. c. 300 BCE"),
                ("Heidegger: Being-toward-death", "Awareness of the inevitability of death enables one to truly live authentically - death illuminates life. 1927"),
                ("Frankl: Search for Meaning", "In any circumstance (including illness and death), people have the freedom to choose their attitude - meaning cannot be taken away. 1946")
            ],
            "shiyi": "亡之，命矣夫: Facing Ran Boniu's illness, Confucius said 'This is fate' - this is not passive fatalism but accepting reality then moving forward. Confucius believed in fate but didn't complain about it - what needed to be done continued, what needed to be said continued. Psychologist Viktor Frankl lost everything in a Nazi concentration camp yet said 'In any environment, people have the freedom to choose their attitude' - this connects with Confucius's '命矣夫' spirit.",
            "discussion": [
                ("Question 1", "When encountering unavoidable困境, is '命矣夫' attitude negative or positive? Why?"),
                ("Question 2", "Facing illness and death, what suggestions do Confucianism and modern psychology each offer? Which do you identify with more?"),
                ("Question 3", "Do you think '尽人事而听天命' and 'lying flat and giving up' are different? Where is the boundary?")
            ]
        }
    ],
    # Day 7: 述而第七
    [
        {
            "num": 19,
            "name": "Shu Er Bu Zuo (述而不作)",
            "yuanwen": [
                "Zi Yue: 'Shu Er Bu Zuo, Xin Er Hao Gu,",
                "Qie Bi Yu Wo Lao Peng.'\""
            ],
            "zhushi": [
                ("Shu (述)", "Transmit, explain"),
                ("Zuo (作)", "Create, innovate"),
                ("Xin Er Hao Gu (信而好古)", "Believe in and love ancient culture"),
                ("Qie (窃)", "Privately, humble self-reference"),
                ("Lao Peng (老彭)", "Laozi and Pengzu, referring to ancient sages in general")
            ],
            "translation_cn": "孔子说：我只是传述而不创作，相信并喜好古代文化，私下里把我和老彭相比。",
            "translation_en": "Confucius said: 'I am one who transmits and does not innovate - believing in and loving the ancient. I venture to compare myself to Lao Peng.'",
            "core_points": [
                ("1. 述而不作", "孔子认为自己只是传承者，不是创造者——这是谦虚，也是对传统的尊重。"),
                ("2. 信而好古", "相信并喜好古代文化——不是崇古媚古，而是认为古代文化有永恒价值。"),
                ("3. 文化传承", "真正的智慧需要传承，不能凭空创造——创新必须建立在继承的基础上。")
            ],
            "background": "This is Confucius's summary of his academic attitude. 'Shu er bu zuo' doesn't mean he lacked ability to create, but that true wisdom lies in transmission rather than novelty. Confucius's life work was compiling and transmitting the Six Classics (Poetry, Documents, Rites, Music, Changes, Spring and Autumn) rather than creating something new. This attitude deeply influenced later generations - Chinese scholarship emphasizes 'Dao Tong' (transmission of the Way), this spirit of passing down through generations.",
            "western": [
                ("Aristotle: I Love My Teacher", "I love my teacher, but I love truth more - but truth must be built on the shoulders of predecessors. c. 340 BCE"),
                ("Newton: Standing on Shoulders", "If I have seen further, it is by standing on the shoulders of giants - all innovation extends from predecessors. 1675"),
                ("Wittgenstein: Language Tradition", "Language is a tradition - we cannot create language, only inherit and use it. 1953"),
                ("Gadamer: Hermeneutics", "Understanding is always historically situated - we always understand from our own historical horizon. 1960")
            ],
            "shiyi": "述而不作: Modern people always emphasize 'innovation', but Confucius reminds us innovation must be built on the foundation of transmission. Many want to innovate but don't even understand traditions, resulting in 'ignorant boldness'. True innovators, like Newton and Einstein, first deeply understood predecessors' achievements, then made breakthroughs. 'Shu er bu zuo' is not conservative but cautious innovation after fully understanding traditions.",
            "discussion": [
                ("Question 1", "Why did Confucius say '述而不作' rather than '创作'? What academic attitude does this reflect?"),
                ("Question 2", "In an era where innovation is the main theme, does '述而不作' attitude still have value?"),
                ("Question 3", "In your field, talk about the relationship between transmission and innovation - how to both inherit traditions and not be stagnant?")
            ]
        },
        {
            "num": 20,
            "name": "Xue Er Bu Yan (学而不厌)",
            "yuanwen": [
                "Zi Yue: 'Mo Er Shi Zhi, Xue Er Bu Yan,",
                "Hui Ren Bu Juan, He You Yu Wo Zai?'"
            ],
            "zhushi": [
                ("Mo Er Shi Zhi (默而识之)", "Silently memorize learned knowledge"),
                ("Xue Er Bu Yan (学而不厌)", "Never tire of learning"),
                ("Hui Ren Bu Juan (诲人不倦)", "Never tire of teaching others"),
                ("He You Yu Wo Zai (何有于我哉)", "What difficulty is there in these for me? - Confucius's self-deprecation")
            ],
            "translation_cn": "孔子说：默默记住学到的知识，学习从不满足，教诲别人从不厌倦，这三件事对我来说有什么难的呢？",
            "translation_en": "Confucius said: 'Silently commit to memory, never grow tired of learning, never tire of teaching others - what difficulty is there in these for me?'",
            "core_points": [
                ("1. 默而识之", "学习要默默记住，不是为了炫耀，而是内化——真正的学习是消化，不是表演。"),
                ("2. 学而不厌", "学习永远不满足——知识没有尽头，学习没有止境，真正学者永远谦逊。"),
                ("3. 诲人不倦", "教别人永远不厌倦——教是另一种学，教的过程也是自己提高的过程。")
            ],
            "background": "This is Confucius's summary of his own qualities and his expectation for all scholars. 'Mo er shi zhi' emphasizes internalization rather than surface; 'xue er bu yan' emphasizes continuity rather than stagnation; 'hui ren bu juan' emphasizes sharing rather than hoarding. Confucius practiced these three throughout his life, even late in life still studying the Book of Changes, truly embodying 'xue er bu yan'.",
            "western": [
                ("Socrates: Maieutics", "Education is not indoctrination but guidance - through questioning, people discover answers themselves; the teacher is a midwife, not a indoctrinator. c. 400 BCE"),
                ("Whitehead: Rhythm of Education", "Education has romantic, precision, and synthesis periods - learning must do the right things at the right stages. 1929"),
                ("Feynman: Understanding to Mastery", "If you can't explain it simply, you don't understand it well enough - learning pursues true understanding. 1965"),
                ("Williams: Learning Pyramid", "Passive learning remembers only 10%, active teaching remembers over 90% - teaching is the best learning. 1969")
            ],
            "shiyi": "学而不厌，诲人不倦: These two phrases are why Confucius became the 'Teacher of Ten Thousand Generations'. Never satisfied in learning, continuously progressing; never tiring in teaching others - teaching consolidates one's own knowledge because teaching and learning enhance each other. Research shows that actively teaching others can improve depth of understanding and memory strength by over 90%. Today, we not only need to learn ourselves but also share and teach - because 'teaching' is the best way to learn.",
            "discussion": [
                ("Question 1", "Do you feel '厌' in learning? How to maintain the state of '学而不厌'?"),
                ("Question 2", "'诲人不倦' - why is teaching others also beneficial to oneself? How does teaching and learning enhance each other?"),
                ("Question 3", "'默而识之' - how to truly remember rather than superficially memorize? What are the learning methods?")
            ]
        }
    ]
]

# ============================================================
# Generate PPTX
# ============================================================

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def make_cover(prs, day, title, ch_names, quote_text):
    s = add_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "Guo Xue Jing Dian Jing Du - Day %d" % day
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    main = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf2 = main.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(60)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.alignment = PP_ALIGN.CENTER

    sub = s.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(1))
    tf3 = sub.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = ch_names
    p3.font.size = Pt(28)
    p3.font.color.rgb = SECONDARY
    p3.alignment = PP_ALIGN.CENTER

    qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.5), Inches(7), Inches(1.3))
    qb.fill.solid()
    qb.fill.fore_color.rgb = BG
    tfq = qb.text_frame
    tfq.word_wrap = True
    pq1 = tfq.paragraphs[0]
    pq1.text = quote_text
    pq1.font.size = Pt(24)
    pq1.font.italic = True
    pq1.font.color.rgb = PRIMARY
    pq1.alignment = PP_ALIGN.CENTER
    add_slide_number(s, 1, 14)

def make_chapter_intro(prs, ch, slide_num, total):
    s = add_slide(prs)
    set_title_bar(s, "Chapter %d - %s" % (ch["num"], ch["name"]))
    add_yuanwen_box(s, ch["yuanwen"], y=1.5, height=3.8)
    add_translation_box(s, ch["translation_cn"], ch["translation_en"], y=5.4)
    add_slide_number(s, slide_num, total)

def make_chapter_annotations(prs, ch, slide_num, total):
    s = add_slide(prs)
    set_title_bar(s, "Chapter %d - %s [Annotations]" % (ch["num"], ch["name"]))
    add_zhushi_box(s, ch["zhushi"], y=1.5, height=4.2)
    add_slide_number(s, slide_num, total)

def make_chapter_core(prs, ch, slide_num, total):
    s = add_slide(prs)
    set_title_bar(s, "Chapter %d - Core Ideas" % ch["num"])
    add_core_box(s, "Core Concept", ch["shiyi"][:80] + "...")
    add_points(s, ch["core_points"], y_start=3.2)
    add_slide_number(s, slide_num, total)

def make_western(prs, ch, slide_num, total):
    s = add_slide(prs)
    set_title_bar(s, "Western Thought Comparison")
    y = 1.5
    for item in ch["western"]:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = item[0]
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = item[1]
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        y += 1.3
    add_slide_number(s, slide_num, total)

def make_modern_app(prs, chapters, slide_num, total):
    s = add_slide(prs)
    set_title_bar(s, "Modern Application")
    y = 1.5
    for i, ch in enumerate(chapters[:3]):
        if i < len(ch["core_points"]):
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.3))
            box.fill.solid()
            box.fill.fore_color.rgb = BG
            tf = box.text_frame
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            p1.text = ch["core_points"][i][0]
            p1.font.size = Pt(18)
            p1.font.bold = True
            p1.font.color.rgb = PRIMARY
            p2 = tf.add_paragraph()
            p2.text = ch["core_points"][i][1][:70]
            p2.font.size = Pt(15)
            p2.font.color.rgb = TEXT
            y += 1.4
    add_slide_number(s, slide_num, total)

def make_background(prs, ch, slide_num, total):
    s = add_slide(prs)
    set_title_bar(s, "Historical Background")
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = ch["name"] + " - Background"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = ch["background"]
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(8)
    add_slide_number(s, slide_num, total)

def make_discussion(prs, ch, slide_num, total):
    s = add_slide(prs)
    set_title_bar(s, "Discussion Questions")
    y = 1.5
    for i, q in enumerate(ch["discussion"]):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "Q%d: %s" % (i+1, q[0])
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT
        p2 = tf.add_paragraph()
        p2.text = q[1]
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        y += 1.4
    add_slide_number(s, slide_num, total)

def make_summary(prs, chapters, day, slide_num, total):
    s = add_slide(prs)
    set_title_bar(s, "Day %d - Summary" % day)
    y = 1.5
    for ch in chapters:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "Chapter %d: %s" % (ch["num"], ch["name"])
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = ch["core_points"][0][1][:70]
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT
        y += 1.5

    final = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.1), Inches(9), Inches(1))
    final.fill.solid()
    final.fill.fore_color.rgb = PRIMARY
    tff = final.text_frame
    tff.word_wrap = True
    pf1 = tff.paragraphs[0]
    pf1.text = chapters[0]["yuanwen"][0]
    pf1.font.size = Pt(22)
    pf1.font.bold = True
    pf1.font.color.rgb = RGBColor(255, 255, 255)
    pf1.alignment = PP_ALIGN.CENTER
    add_slide_number(s, slide_num, total)


def generate_day_pptx(day, chapters, title):
    prs = Presentation()
    total = 14
    ch_names = ", ".join([ch["name"] for ch in chapters])
    quote = chapters[0]["yuanwen"][0]

    make_cover(prs, day, title, ch_names, quote)

    # Chapter 1
    make_chapter_intro(prs, chapters[0], 2, total)
    make_chapter_annotations(prs, chapters[0], 3, total)
    make_chapter_core(prs, chapters[0], 4, total)

    # Chapter 2
    make_chapter_intro(prs, chapters[1], 5, total)
    make_chapter_annotations(prs, chapters[1], 6, total)
    make_chapter_core(prs, chapters[1], 7, total)

    # Chapter 3 (or Chapter 2 for Day 7)
    if len(chapters) > 2:
        make_chapter_intro(prs, chapters[2], 8, total)
        make_chapter_annotations(prs, chapters[2], 9, total)
        make_chapter_core(prs, chapters[2], 10, total)
        make_western(prs, chapters[0], 11, total)
        make_modern_app(prs, chapters, 12, total)
        make_background(prs, chapters[0], 13, total)
        make_summary(prs, chapters, day, 14, total)
    else:
        # Day 7: only 2 chapters
        make_chapter_core(prs, chapters[1], 8, total)
        make_western(prs, chapters[0], 9, total)
        make_modern_app(prs, chapters, 10, total)
        make_background(prs, chapters[0], 11, total)
        make_discussion(prs, chapters[0], 12, total)
        make_discussion(prs, chapters[1], 13, total)
        make_summary(prs, chapters, day, 14, total)

    return prs


# ============================================================
# Main
# ============================================================
import os
output_dir = "/Users/mac/.openclaw/workspace/courses/论语"
os.makedirs(output_dir, exist_ok=True)

DAY_TITLES = [
    "学而第一 Xue Er Di Yi",
    "为政第二 Wei Zheng Di Er",
    "八佾第三 Ba Yi Di San",
    "里仁第四 Li Ren Di Si",
    "公冶长第五 Gong Ye Chang Di Wu",
    "雍也第六 Yong Ye Di Liu",
    "述而第七 Shu Er Di Qi"
]

for i, chapters in enumerate(CHAPTERS):
    day = i + 1
    title = DAY_TITLES[i]
    prs = generate_day_pptx(day, chapters, title)
    path = os.path.join(output_dir, "论语_第%d天.pptx" % day)
    prs.save(path)
    print("Done: 论语_第%d天.pptx" % day)

print("\nAll 7 PPTX files generated!")