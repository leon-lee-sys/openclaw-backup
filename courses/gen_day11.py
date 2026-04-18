from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案 - 学术白风格
COLORS = {
    'primary': RGBColor(0, 51, 102),      # 深蓝
    'secondary': RGBColor(102, 102, 102), # 灰
    'accent': RGBColor(204, 0, 0),        # 红
    'text': RGBColor(51, 51, 51),         # 深灰
    'bg': RGBColor(255, 255, 255),        # 白
    'light': RGBColor(240, 248, 255),     # 淡蓝
}

def add_slide_bg(slide, prs):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['bg']

def add_title_bar(slide, title, chapter_info=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()

    # 标题文字
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    if chapter_info:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(9), Inches(0.4))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = chapter_info
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(200, 220, 240)

def add_origin_block(slide, text, top=Inches(1.5)):
    """原文块"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['light']
    box.line.color.rgb = COLORS['primary']

    tb = slide.shapes.add_textbox(Inches(0.7), top + Inches(0.15), Inches(8.6), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "【原文】"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']

    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.size = Pt(15)
    p2.font.color.rgb = COLORS['primary']
    p2.font.name = "SimSun"

def add_bullets(slide, items, top, font_size=Pt(16), spacing=Pt(10)):
    tb = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = font_size
        p.font.color.rgb = COLORS['text']
        p.space_after = spacing

def add_quote_box(slide, text, top, color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(9), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = color or COLORS['secondary']
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.7), top + Inches(0.15), Inches(8.6), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"💡 {text}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============ 第1页：封面 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, prs)

# 顶部装饰条
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.3))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = COLORS['primary']
top_bar.line.fill.background()

# 主标题
tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "《道德经》每日课件"
p.font.size = Pt(42)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.alignment = PP_ALIGN.CENTER

# 副标题
p2 = tf.add_paragraph()
p2.text = "第11天 · 第34-36章"
p2.font.size = Pt(24)
p2.font.color.rgb = COLORS['secondary']
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

# 日期
p3 = tf.add_paragraph()
p3.text = "2026年4月18日"
p3.font.size = Pt(16)
p3.font.color.rgb = COLORS['secondary']
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(20)

# 底部引言
quote_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(10), Inches(1))
quote_bar.fill.solid()
quote_bar.fill.fore_color.rgb = COLORS['primary']
quote_bar.line.fill.background()

tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.65), Inches(9), Inches(0.7))
tf2 = tb2.text_frame
p_q = tf2.paragraphs[0]
p_q.text = "「上善若水，水善利万物而不争。」"
p_q.font.size = Pt(18)
p_q.font.italic = True
p_q.font.color.rgb = RGBColor(255, 255, 255)
p_q.alignment = PP_ALIGN.CENTER

# ============ 第2页：第三十四章 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, prs)
add_title_bar(slide, "第三十四章 · 大道泛兮", "Day 11 | 第34章")
add_origin_block(slide, "大道泛兮，其可左右。万物恃之而生而不辞，功成不名有。衣养万物而不为主，常无欲，可名于小；万物归焉而不为主，可名为大。以其终不自大，故能成其大。", Inches(1.5))
add_bullets(slide, [
    "大道如水，流布天地，无所不到、无所不能",
    "万物依赖道而生存，道却从不推辞、不居功",
    "道滋养万物却不做主，永远没有私欲",
    "道看似渺小——因为它从不自以为大",
    "道成就伟大——正因为它永远不自我膨胀"
], Inches(3.2), Pt(16))

# ============ 第3页：三十四章 · 管理智慧 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, prs)
add_title_bar(slide, "第三十四章 · 管理智慧", "Day 11 | 第34章")

add_quote_box(slide, "以其终不自大，故能成其大。——因为从不自居伟大，所以成就了伟大", Inches(1.3), COLORS['primary'])

add_bullets(slide, [
    "【领导力】真正的领导者，从不宣称自己是 leader",
    "【授权之道】培养人才却不替代，包融万物却不控制",
    "【无为之治】最高明的管理，是让被管理者感觉不到被管",
    "【谷歌案例】OKR取代KPI：目标透明却不施压，员工自驱",
    "【海尔的启示】张瑞敏"砸冰箱"后走向授权管理：让员工成为自己的CEO",
    "【日本经营之神】松下幸之助：企业是"社会公器"，不属于创始人",
    "【华为管理】任正非："让听见炮声的人呼唤炮火"——一线决策，赋能而非控制"
], Inches(2.5), Pt(14), Pt(9))

# ============ 第4页：三十五章 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, prs)
add_title_bar(slide, "第三十五章 · 执大象", "Day 11 | 第35章")
add_origin_block(slide, "执大象，天下往。往而不害，安平泰。乐与饵，过客止。道之出口，淡乎其无味，视之不足见，听之不足闻，用之不足既。", Inches(1.5))
add_bullets(slide, [
    "掌握大道者，天下人都会来归附",
    "归附者彼此不相害，因而安宁、太平、通泰",
    "世俗的乐声美食，能使过客驻足，却不能长久",
    "道的言语平淡无味——看似平常，实则妙用无穷",
    "看不见、听不到，却"用之不尽"——大道至简，功效无限"
], Inches(3.2), Pt(16))

# ============ 第5页：三十五章 · 生活智慧 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, prs)
add_title_bar(slide, "第三十五章 · 生活智慧", "Day 11 | 第35章")

add_quote_box(slide, "道之出口，淡乎其无味。——真正的好东西，往往看起来很普通", Inches(1.3), COLORS['secondary'])

add_bullets(slide, [
    "【处世】真正有魅力的人，不需要刻意表演",
    "【教育】最好的教育方式，往往平淡无奇，润物无声",
    "【健康】最简单的养生：规律作息、清淡饮食、平和心态",
    "【择业】选工作的标准：平台够大、发展可持续，而非一时高薪",
    "【投资】好公司往往业务简单、管理朴实（巴菲特"护城河"理念）",
    "【做人】真诚平淡却能打动人心——少套路，多真心",
    "【交友】真正的朋友不需要热闹场合维系，平日淡如水，久处不厌"
], Inches(2.5), Pt(14), Pt(9))

# ============ 第6页：三十六章 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, prs)
add_title_bar(slide, "第三十六章 · 将欲取之", "Day 11 | 第36章")
add_origin_block(slide, "将欲取天下而为之，吾见其不得已。天下神器，不可为也。为者败之，执者失之。物，或行或随，或呴或吹，或强或羸，或挫或隳。是以圣人去甚，去奢，去泰。", Inches(1.5))
add_bullets(slide, [
    "想靠强力夺取天下？我看他达不到目的",
    "天下是神圣之物——不可强取，不可硬控",
    "强为者必败，强执者必失",
    "万物各有本性：有的前行有的跟随，有的轻呼吸有的急吹",
    "圣人行事：去除极端、去除奢侈、去除过度"
], Inches(3.2), Pt(16))

# ============ 第7页：三十六章 · 现代启示 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, prs)
add_title_bar(slide, "第三十六章 · 现代启示", "Day 11 | 第36章")

add_quote_box(slide, "圣人去甚，去奢，去泰。——去除极端、奢侈、过度，方为正道", Inches(1.3), COLORS['accent'])

add_bullets(slide, [
    "【组织管理】过度管控 = 控制越多，失去越多（微管理之害）",
    "【创业者警句】马云："看不见对手在哪里，是最可怕的"——强为者败",
    "【个人成长】凡事不过度：努力但不极端，张弛有度方可持续",
    "【历史教训】秦朝强取天下，二世而亡——暴力可得一时，不能持久",
    "【适度原则】中庸之道：找到平衡点，而非走向任何一个极端",
    "【企业管理】GE杰克·韦尔奇：数一数二战略——去甚，不做末尾",
    "【人生哲学】曾国藩："既往不恋，当下不杂，未来不迎"——去奢，去泰"
], Inches(2.5), Pt(14), Pt(9))

# ============ 第8页：西方思想对照 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, prs)
add_title_bar(slide, "中西思想对照", "Day 11 | 第34-36章")

# 左栏
left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(4.6), Inches(5.5))
ltf = left_box.text_frame
ltf.word_wrap = True

items_left = [
    ("🇫🇷 卢梭（法国）", "\"最大的危险，不在于看清野心，而在于失去德性。\" — 强取者败"),
    ("🇩🇪 黑格尔（德国）", "\"存在即合理，但过度存在走向自我毁灭。\" — 物极必反"),
    ("🇬🇧 亚当·斯密（英国）", "\"市场有一只看不见的手，自动调节供需。\" — 无为而治的市场版"),
    ("🇺🇸 杰斐逊（美国）", "\"最少的政府是最好的政府。\" — 自由主义的道家声音"),
]
for i, (name, quote) in enumerate(items_left):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    p.text = name
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p2 = ltf.add_paragraph()
    p2.text = quote
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLORS['text']
    p2.space_after = Pt(12)

# 右栏标题
right_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.5), Inches(0.5))
rtf = right_title.text_frame
rp = rtf.paragraphs[0]
rp.text = "🌏 本章核心智慧"
rp.font.size = Pt(16)
rp.font.bold = True
rp.font.color.rgb = COLORS['primary']

right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.5), Inches(5))
rtf2 = right_box.text_frame
rtf2.word_wrap = True

insights = [
    "第34章：不自大，故成其大\n　　→ 谦卑是真正的力量",
    "第35章：道出口，淡无味\n　　→ 真味在平淡中",
    "第36章：去甚、去奢、去泰\n　　→ 极端是失败之源",
    "三章共贯：\n　　→ 无为、守中、返璞归真",
]
for i, ins in enumerate(insights):
    p = rtf2.paragraphs[0] if i == 0 else rtf2.add_paragraph()
    p.text = ins
    p.font.size = Pt(13)
    p.font.color.rgb = COLORS['text']
    p.space_after = Pt(14)

# ============ 第9页：精华提炼 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, prs)
add_title_bar(slide, "精华提炼 · 一句话总结", "Day 11 | 第34-36章")

# 三章核心
for i, (chapter, core, quote) in enumerate([
    ("第34章\n大道泛兮", "不自大，故成其大\n谦下不争，天下归之", "以其终不自大，\n故能成其大"),
    ("第35章\n执大象", "道之出口，淡无味\n用之不足既", "执大象，\n天下往"),
    ("第36章\n将欲取之", "去甚、去奢、去泰\n极端是败事之本", "为者败之，\n执者失之"),
]):
    top = Inches(1.5) + i * Inches(1.7)

    # 章节标签
    chapter_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(1.8), Inches(1.4))
    chapter_box.fill.solid()
    chapter_box.fill.fore_color.rgb = COLORS['primary']
    chapter_box.line.fill.background()

    ctb = slide.shapes.add_textbox(Inches(0.5), top + Inches(0.2), Inches(1.8), Inches(1.0))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = chapter
    cp.font.size = Pt(13)
    cp.font.bold = True
    cp.font.color.rgb = RGBColor(255, 255, 255)
    cp.alignment = PP_ALIGN.CENTER

    # 核心内容
    content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), top, Inches(4.5), Inches(1.4))
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = COLORS['light']
    content_box.line.color.rgb = COLORS['secondary']

    ctb2 = slide.shapes.add_textbox(Inches(2.7), top + Inches(0.2), Inches(4.1), Inches(1.0))
    ctf2 = ctb2.text_frame
    ctf2.word_wrap = True
    cp2 = ctf2.paragraphs[0]
    cp2.text = core
    cp2.font.size = Pt(14)
    cp2.font.bold = True
    cp2.font.color.rgb = COLORS['primary']

    # 金句
    quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), top, Inches(2.3), Inches(1.4))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = COLORS['secondary']
    quote_box.line.fill.background()

    qtb = slide.shapes.add_textbox(Inches(7.3), top + Inches(0.3), Inches(2.1), Inches(0.9))
    qtf = qtb.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.text = quote
    qp.font.size = Pt(11)
    qp.font.italic = True
    qp.font.color.rgb = RGBColor(255, 255, 255)
    qp.alignment = PP_ALIGN.CENTER

# 一句话总结
summary_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2), Inches(9), Inches(0.9))
summary_box.fill.solid()
summary_box.fill.fore_color.rgb = COLORS['primary']
summary_box.line.fill.background()

stb = slide.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(8.6), Inches(0.6))
stf = stb.text_frame
sp = stf.paragraphs[0]
sp.text = "🌟 三天三章，道之三味：谦下守中、守朴不争、返璞归真。"
sp.font.size = Pt(16)
sp.font.bold = True
sp.font.color.rgb = RGBColor(255, 255, 255)
sp.alignment = PP_ALIGN.CENTER

# ============ 第10页：结束页 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, prs)

top_bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.3))
top_bar2.fill.solid()
top_bar2.fill.fore_color.rgb = COLORS['primary']
top_bar2.line.fill.background()

tb_end = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
tf_end = tb_end.text_frame
ep = tf_end.paragraphs[0]
ep.text = "明天继续 · 第12天"
ep.font.size = Pt(36)
ep.font.bold = True
ep.font.color.rgb = COLORS['primary']
ep.alignment = PP_ALIGN.CENTER

ep2 = tf_end.add_paragraph()
ep2.text = "第37-39章"
ep2.font.size = Pt(22)
ep2.font.color.rgb = COLORS['secondary']
ep2.alignment = PP_ALIGN.CENTER
ep2.space_before = Pt(15)

ep3 = tf_end.add_paragraph()
ep3.text = "道常无为 · 始制有名 · 上德不德"
ep3.font.size = Pt(16)
ep3.font.color.rgb = COLORS['secondary']
ep3.alignment = PP_ALIGN.CENTER
ep3.space_before = Pt(10)

# 底部
bot_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(10), Inches(1))
bot_bar.fill.solid()
bot_bar.fill.fore_color.rgb = COLORS['primary']
bot_bar.line.fill.background()

btb = slide.shapes.add_textbox(Inches(0.5), Inches(6.65), Inches(9), Inches(0.7))
btf = btb.text_frame
bp = btf.paragraphs[0]
bp.text = "「上善若水，水善利万物而不争。」— 道德经第八章"
bp.font.size = Pt(16)
bp.font.italic = True
bp.font.color.rgb = RGBColor(200, 220, 240)
bp.alignment = PP_ALIGN.CENTER

# 保存
output_path = "/Users/mac/.openclaw/workspace/courses/道德经_第11天.pptx"
prs.save(output_path)
print(f"✅ 课件已保存：{output_path}")
