# -*- coding: utf-8 -*-
"""
道德经第20天 课件生成器
章节：第64-66章
风格：精读深度讲解版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
ACCENT = RGBColor(204, 0, 0)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_yuanwen_box(slide, lines, y=1.4, height=2.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(20)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=BLUE_LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(255, 255, 255)

def add_points(slide, points, y_start=3.1):
    y = y_start
    for title, content in points:
        pt_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT
        y += 1.1
    return y

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.5))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第20课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第64-66章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "其安易持 · 古之善为道者 · 江海能为百谷王"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "江海所以能为百谷王者，以其善下之。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = PRIMARY
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第六十六章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = SECONDARY
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, 12)

# ===== 第2页：第六十四章详解 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s2, "第六十四章 · 其安易持")
add_yuanwen_box(s2, [
    "其安易持，其未兆易谋。",
    "其脆易泮，其微易散。",
    "为之于未有，治之于未乱。",
    "",
    "合抱之木，生于毫末；",
    "九层之台，起于累土；",
    "千里之行，始于足下。",
    "",
    "为者败之，执者失之。",
    "是以圣人无为，故无败；",
    "无执，故无失。",
    "",
    "民之从事，常于几成而败之。",
    "慎终如始，则无败事。"
], y=1.4, height=4.8)
content2 = s2.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(2))
tf_c2 = content2.text_frame
tf_c2.word_wrap = True
pc2 = tf_c2.paragraphs[0]
pc2.text = "【白话解读】局面安稳时容易维持，事情还没露头时容易想办法。脆弱的东西容易破碎，细微的东西容易散开。在问题还没出现时就处理好，在混乱还没形成时就预防好。合抱的大树从细小的萌芽长成，九层高台从一筐土开始，千里远行从脚下开始。强为的人会失败，硬抓的人会失去。圣人顺其自然，所以不失败；不硬抓，所以不失去。"
pc2.font.size = Pt(13)
pc2.font.color.rgb = TEXT
add_slide_number(s2, 2, 12)

# ===== 第3页：第六十四章核心 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s3, "第六十四章 · 核心思想")
add_core_box(s3, "🌟 核心观点", "为之于未有，治之于未乱——最好的治理是预防问题发生，而非出了问题再补救。千里之行，始于足下——再大的事，也得从眼前的小事做起。")
add_points(s3, [
    ("1. 其安易持", "安稳时容易维持——居安思危，在顺境时保持警觉。"),
    ("2. 为之于未有", "在问题还没出现时就处理——预防永远胜于补救。"),
    ("3. 千里之行，始于足下", "再大的事业也是从眼前的小步开始——不要轻视每一天的积累。"),
    ("4. 慎终如始", "结束时像开始时一样谨慎——坚持到底才是真正的成功。")
])
add_slide_number(s3, 3, 12)

# ===== 第4页：第六十五章详解 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s4, "第六十五章 · 古之善为道者")
add_yuanwen_box(s4, [
    "古之善为道者，",
    "非以明民，将以愚之。",
    "",
    "民之难治，以其智多。",
    "故以智治国，国之贼；",
    "不以智治国，国之福。",
    "",
    "知此两者，亦稽式。",
    "常知稽式，是谓玄德。",
    "玄德深矣，远矣，",
    "与物反矣，",
    "然后乃至大顺。"
], y=1.4, height=4.2)
content4 = s4.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(2))
tf_c4 = content4.text_frame
tf_c4.word_wrap = True
pc4 = tf_c4.paragraphs[0]
pc4.text = "【白话解读】古代善于行道的人，不是让人民变得聪明精巧，而是让人民保持淳朴。百姓难以治理，是因为他们的巧智太多。所以用智巧来治理国家，是国家的祸害；不用智巧来治理，是国家的幸福。懂得这两个区别，就懂得了治国的法则。永远守住这个法则，就是深远的德行。这种深远的德行啊，又深又远，与万物一起复归道的怀抱，最终达到最大的顺应自然。"
pc4.font.size = Pt(13)
pc4.font.color.rgb = TEXT
add_slide_number(s4, 4, 12)

# ===== 第5页：第六十五章核心 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s5, "第六十五章 · 核心思想")
add_core_box(s5, "🌟 核心观点", "将以愚之——道家说的\"愚\"不是愚蠢，而是淳朴返璞。让百姓保持纯朴的心智，社会才能长治久安。")
add_points(s5, [
    ("1. 非以明民，将以愚之", "不让百姓变得巧智精巧，而是保持纯朴——巧智反而让社会更混乱。"),
    ("2. 以智治国，国之贼", "用智巧治理国家反而是祸害——算计越多，人心越乱。"),
    ("3. 玄德", "深远的德行——不求表面光鲜，但求内在深厚。"),
    ("4. 与物反矣", "与万物一起回归道的本质——最终达到大顺其自然的状态。")
])
add_slide_number(s5, 5, 12)

# ===== 第6页：第六十六章详解 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s6, "第六十六章 · 江海能为百谷王")
add_yuanwen_box(s6, [
    "江海所以能为百谷王者，",
    "以其善下之，",
    "故能为百谷王。",
    "",
    "是以圣人欲上民，",
    "必以言下之；",
    "欲先民，",
    "必以身后之。",
    "",
    "是以圣人处上而民不重，",
    "处前而民不害。",
    "是以天下乐推而不厌。",
    "",
    "以其不争，",
    "故天下莫能与之争。"
], y=1.4, height=4.2)
content6 = s6.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(2))
tf_c6 = content6.text_frame
tf_c6.word_wrap = True
pc6 = tf_c6.paragraphs[0]
pc6.text = "【白话解读】江海之所以能成为百川的领袖，是因为它善于处在低下的位置，所以能成为百川之王。所以圣人想要居于人民之上，必须先用言语表示谦下；想要站在人民前面，必须把自己放在后面。因此圣人虽然地位在上，但人民不觉得沉重；虽然站在前面，但人民不觉得受害。天下人都乐意推戴他而不讨厌他。因为他不争，所以天下没有人能和他争。"
pc6.font.size = Pt(13)
pc6.font.color.rgb = TEXT
add_slide_number(s6, 6, 12)

# ===== 第7页：第六十六章核心 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s7, "第六十六章 · 核心思想")
add_core_box(s7, "🌟 核心观点", "江海善下——能处于最低处，才能成大事。圣人不争，故天下莫能与之争——最强大的力量是不争。")
add_points(s7, [
    ("1. 善下之", "江海居下而成百谷王——谦下才是真正的尊贵。"),
    ("2. 欲上必下", "想居上必须先处下——想要领导，先学会服务。"),
    ("3. 处上而不重", "身居高位但百姓不觉沉重——好的领导让人轻松。"),
    ("4. 不争而莫能争", "因为不争，所以无人能争——最无敌的境界是不争。")
])
add_slide_number(s7, 7, 12)

# ===== 第8页：现代反思 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s8, "第六十四-六十六章 · 现代反思")
items8 = [
    ("🏛️ 领导力", "欲上必下 → 真正高明的领导者，是那些愿意服务别人的人。"),
    ("💰 理财智慧", "千里之行，始于足下 → 理财从每一笔小钱开始，复利需要时间。"),
    ("🔄 风险管理", "为之于未有 → 最好的投资风险控制是在亏损发生之前。"),
    ("🌱 长期主义", "慎终如始 → 不管做什么，坚持到最后才是真正的胜利。")
]
y8 = 1.4
for icon_title, explanation in items8:
    box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y8), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf8 = box.text_frame
    tf8.word_wrap = True
    p8_1 = tf8.paragraphs[0]
    p8_1.text = icon_title
    p8_1.font.size = Pt(16)
    p8_1.font.bold = True
    p8_1.font.color.rgb = PRIMARY
    p8_2 = tf8.add_paragraph()
    p8_2.text = explanation
    p8_2.font.size = Pt(14)
    p8_2.font.color.rgb = TEXT
    y8 += 1.4
add_slide_number(s8, 8, 12)

# ===== 第9页：西方思想 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s9, "🌍 西方思想对照")
phils = [
    ("老子与蒙田 · 怀疑主义", "真正的智慧在于知道自己的无知——中西方哲人都反对过度精巧的计算。约1580年"),
    ("亚当·斯密 · 无形的手", "最好的经济秩序不是政府设计出来的——与\"为道之愚\"的智慧相通。1776年"),
    ("海德格尔 · 老子影响", "西方存在主义大家也向东方哲学寻求突破——道法自然的智慧超越东西方。约1940年"),
    (" Buffett · 价值投资", "没有人愿意慢慢变富——但\"千里之行，始于足下\"正是复利思维的本质。1950年代至今")
]
y9 = 1.4
for name_theory, desc_year in phils:
    box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y9), Inches(9), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf9 = box.text_frame
    tf9.word_wrap = True
    p9_1 = tf9.paragraphs[0]
    p9_1.text = name_theory
    p9_1.font.size = Pt(15)
    p9_1.font.bold = True
    p9_1.font.color.rgb = PRIMARY
    p9_2 = tf9.add_paragraph()
    p9_2.text = desc_year
    p9_2.font.size = Pt(13)
    p9_2.font.color.rgb = TEXT
    y9 += 1.3
add_slide_number(s9, 9, 12)

# ===== 第10页：现代应用 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s10, "📌 现代应用")
add_points(s10, [
    ("企业管理", "为之于未有 → 建立风险预警系统，在问题爆发前就解决它。"),
    ("投资理财", "千里之行，始于足下 → 定投指数基金，坚持20年，你会超过大多数聪明人。"),
    ("人际交往", "善下之 → 想赢得尊重，先学会倾听和尊重别人。"),
    ("个人成长", "慎终如始 → 每天坚持一点点，比三天打鱼两天晒网强一万倍。")
])
add_slide_number(s10, 10, 12)

# ===== 第11页：中国典故 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s11, "📜 中国典故：萧何月下追韩信")
box11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.5))
box11.fill.solid()
box11.fill.fore_color.rgb = BG
tf11 = box11.text_frame
tf11.word_wrap = True
p11_1 = tf11.paragraphs[0]
p11_1.text = "【史记·淮阴侯列传】"
p11_1.font.size = Pt(16)
p11_1.font.bold = True
p11_1.font.color.rgb = PRIMARY
p11_2 = tf11.add_paragraph()
p11_2.text = "韩信原本只是项羽军中一个小小的执戟郎中，多次献策不被采纳，于是逃离楚营。萧何知道韩信是大才，连夜追赶，终于追回韩信，并向刘邦极力推荐。刘邦拜韩信为大将军，最终在楚汉争霸中击败项羽。萧何的伟大在于：他不看表面身份，只看内在才能；懂得在发现人才时果断行动。"
p11_2.font.size = Pt(15)
p11_2.font.color.rgb = TEXT
p11_2.space_before = Pt(10)
insight11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.1), Inches(9), Inches(1))
insight11.fill.solid()
insight11.fill.fore_color.rgb = BLUE_LIGHT
tfi11 = insight11.text_frame
tfi11.word_wrap = True
pi11 = tfi11.paragraphs[0]
pi11.text = "💡 启示"
pi11.font.size = Pt(18)
pi11.font.bold = True
pi11.font.color.rgb = RGBColor(255, 255, 255)
pi11_2 = tfi11.add_paragraph()
pi11_2.text = "\"千里之行，始于足下\"——韩信的成功始于那关键的一步：被萧何追回。但更关键的是萧何：他懂得识别和珍惜真正的人才。每个人都要找到自己的\"萧何\"。"
pi11_2.font.size = Pt(14)
pi11_2.font.color.rgb = RGBColor(255, 255, 255)
app11 = s11.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2))
tfa11 = app11.text_frame
tfa11.word_wrap = True
pa11 = tfa11.paragraphs[0]
pa11.text = "【现代应用】职场：识别并追随真正的领导者比职位高低更重要。人生：找到你的\"萧何\"——那个能发现你潜力的人。投资：找到适合自己的投资策略，然后坚持。"
pa11.font.size = Pt(14)
pa11.font.color.rgb = TEXT
add_slide_number(s11, 11, 12)

# ===== 第12页：总结 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s12, "第20课 · 三章精华总结")
sums = [
    ("第六十四章", "其安易持", "为之于未有，治之于未乱。千里之行，始于足下。"),
    ("第六十五章", "古之善为道者", "非以明民，将以愚之。以智治国，国之贼。"),
    ("第六十六章", "江海能为百谷王", "善下之，故能为百谷王。圣人不争，故天下莫能与之争。")
]
y12 = 1.4
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    y12 += 1.6
final = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(9), Inches(1))
final.fill.solid()
final.fill.fore_color.rgb = PRIMARY
tff = final.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "以其不争，故天下莫能与之争。"
pf1.font.size = Pt(20)
pf1.font.bold = True
pf1.font.color.rgb = RGBColor(255, 255, 255)
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "—— 道德经第六十六章"
pf2.font.size = Pt(14)
pf2.font.color.rgb = RGBColor(200, 200, 200)
pf2.alignment = PP_ALIGN.CENTER
add_slide_number(s12, 12, 12)

prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第20天.pptx')
print("Done: 道德经_第20天.pptx (12 slides, 64-66章)")
