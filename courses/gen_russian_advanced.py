#!/usr/bin/env python3
"""
俄语进阶7天课件生成器
在7天速成基础上深化，加入更多实用词汇、语法点和对话练习
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# 颜色配置
TITLE_BG = RGBColor(0x1A, 0x47, 0x8A)  # 深蓝
ACCENT = RGBColor(0xC0, 0x39, 0x2B)   # 俄罗斯红
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 7天进阶内容
ADVANCED_CONTENT = [
    {
        "day": "День 8",
        "title": "Приветствия и знакомство",
        "subtitle": "深化问候与介绍",
        "sections": [
            {
                "title": "进阶词汇",
                "words": [
                    ("Как дела?", "你好吗？"),
                    ("Очень хорошо!", "非常好！"),
                    ("Рад знакомству!", "很高兴认识你！"),
                    ("Как поживаете?", "您近来好吗？"),
                    ("Давно не виделись!", "好久不见了！"),
                    ("Что нового?", "有什么新鲜事？"),
                    ("Всё отлично, спасибо!", "一切都好，谢谢！"),
                    ("И у меня всё хорошо!", "我也一切都好！"),
                ]
            },
            {
                "title": "语法要点",
                "words": [
                    ("Вы → ты (亲密关系)", "您式转为你式"),
                    ("Откуда вы? (询问籍贯)", "您从哪里来？"),
                    ("Чем занимаетесь?", "您是做什么工作的？"),
                    ("Мне ... лет (年龄表达)", "我...岁"),
                ]
            },
            {
                "title": "实用对话",
                "dialogue": [
                    "— Здравствуйте! Как вас зовут?",
                    "— Меня зовут Ли. А как вас?",
                    "— Очень приятно! Откуда вы?",
                    "— Я из Шанхая. А вы?",
                    "— Я тоже из Китая! Как это здорово!",
                ]
            },
            {
                "title": "文化小贴士",
                "content": "俄罗斯人在正式场合习惯用Вы（您），关系熟悉后才用ты（你）。初次见面通常握手，关系亲密后可能拥抱贴面礼。"
            }
        ]
    },
    {
        "day": "День 9",
        "title": "Работа и карьера",
        "subtitle": "深化工作与职业",
        "sections": [
            {
                "title": "进阶词汇",
                "words": [
                    ("Я работаю в компании", "我在公司工作"),
                    ("Моя должность - менеджер", "我的职位是经理"),
                    ("Я отвечаю за проекты", "我负责项目"),
                    ("У меня важная встреча", "我有一个重要会议"),
                    ("Когда у вас перерыв?", "您什么时候休息？"),
                    ("Я занимаюсь маркетингом", "我从事市场营销"),
                    ("Мы команда", "我们是一个团队"),
                    ("Хорошая работа!", "工作出色！"),
                ]
            },
            {
                "title": "语法要点",
                "words": [
                    ("Кем вы работаете? (职业)", "您做什么工作？"),
                    ("Где вы работаете? (地点)", "您在哪儿工作？"),
                    ("Я работаю с 9 до 6", "我朝九晚六工作"),
                    ("У нас гибкий график", "我们的工作时间灵活"),
                ]
            },
            {
                "title": "实用对话",
                "dialogue": [
                    "— Где вы работаете?",
                    "— Я работаю в международной компании.",
                    "— Какой у вас график работы?",
                    "— Обычно с девяти утра до шести вечера.",
                    "— Это очень удобно!",
                ]
            },
            {
                "title": "文化小贴士",
                "content": "俄罗斯人很重视职业身份，见面常问工作。工作场合一般比较正式，注意使用Вы来称呼。"
            }
        ]
    },
    {
        "day": "День 10",
        "title": "В ресторане",
        "subtitle": "深化餐厅用语",
        "sections": [
            {
                "title": "进阶词汇",
                "words": [
                    ("Можно меню?", "可以给我菜单吗？"),
                    ("Я хотел бы заказать...", "我想点..."),
                    ("Какое у вас фирменное блюдо?", "你们有什么招牌菜？"),
                    ("Счёт, пожалуйста", "买单"),
                    ("Это очень вкусно!", "这个很好吃！"),
                    ("Мне, пожалуйста, воду", "请给我水"),
                    ("Без острого, пожалуйста", "请不要辣的"),
                    ("Можно оплатить картой?", "可以刷卡吗？"),
                ]
            },
            {
                "title": "语法要点",
                "words": [
                    ("Я бы хотел + AKK (愿望)", "我想要..."),
                    ("Мне нужно + N (需要)", "我需要..."),
                    ("Без + G (否定)", "不要..."),
                    ("С + I (一起)", "和...一起"),
                ]
            },
            {
                "title": "实用对话",
                "dialogue": [
                    "— Добрый вечер! Можно меню?",
                    "— Конечно! Вот, пожалуйста.",
                    "— Что вы рекомендуете?",
                    "— Наше фирменное блюдо - борщ!",
                    "— Отлично! Принесите, пожалуйста, борщ.",
                ]
            },
            {
                "title": "文化小贴士",
                "content": "在俄罗斯餐厅，给小费不是强制的，但如果你对服务满意，通常会给账单的10-15%。俄罗斯菜品口味偏重，多用油和盐。"
            }
        ]
    },
    {
        "day": "День 11",
        "title": "Путешествие",
        "subtitle": "深化旅行用语",
        "sections": [
            {
                "title": "进阶词汇",
                "words": [
                    ("Где находится вокзал?", "火车站在哪里？"),
                    ("Когда отправляется поезд?", "火车什么时候出发？"),
                    ("Мне нужен билет в один конец", "我需要单程票"),
                    ("Когда мы прибудем?", "我们什么时候到达？"),
                    ("Где можно взять такси?", "哪里可以打车？"),
                    ("Сколько стоит проезд?", "票价多少钱？"),
                    ("Я заблудился", "我迷路了"),
                    ("Вы говорите по-английски?", "您说英语吗？"),
                ]
            },
            {
                "title": "语法要点",
                "words": [
                    ("В/на + PR (地点)", "在...地方"),
                    ("Отправляться в/на + AKK", "出发前往..."),
                    ("Сколько стоит + N?", "...多少钱？"),
                    ("Мне нужен/нужна/нужно", "我需要（性数配合）"),
                ]
            },
            {
                "title": "实用对话",
                "dialogue": [
                    "— Извините, когда отправляется поезд до Москвы?",
                    "— В десять часов вечера.",
                    "— Мне нужно два билета. Сколько стоит?",
                    "— Два билета - четыре тысячи рублей.",
                    "— Хорошо, я беру!",
                ]
            },
            {
                "title": "文化小贴士",
                "content": "俄罗斯幅员辽阔，火车是重要的长途交通工具。俄罗斯有很多著名的火车站，莫斯科就有多个大型火车站，分别开往不同方向。"
            }
        ]
    },
    {
        "day": "День 12",
        "title": "Покупки",
        "subtitle": "深化购物用语",
        "sections": [
            {
                "title": "进阶词汇",
                "words": [
                    ("Сколько это стоит?", "这个多少钱？"),
                    ("Это слишком дорого", "这太贵了"),
                    ("У вас есть скидка?", "有折扣吗？"),
                    ("Можно примерить?", "可以试穿吗？"),
                    ("Где примерочная?", "试衣间在哪里？"),
                    ("Я возьму это", "我要这个"),
                    ("Дайте мне чек, пожалуйста", "请给我收据"),
                    ("Можно оплатить картой?", "可以刷卡吗？"),
                ]
            },
            {
                "title": "语法要点",
                "words": [
                    ("Этот/эта/это (指示代词)", "这个/这件/这辆"),
                    ("Слишком + PR/adv (程度)", "太..."),
                    ("У меня есть + N (拥有)", "我有..."),
                    ("Мне нужно + N (需要)", "我需要..."),
                ]
            },
            {
                "title": "实用对话",
                "dialogue": [
                    "— Здравствуйте! Сколько стоит эта куртка?",
                    "— Двадцать тысяч рублей.",
                    "— Это слишком дорого. Есть скидка?",
                    "— Сейчас скидка двадцать процентов!",
                    "— Хорошо, я возьму!",
                ]
            },
            {
                "title": "文化小贴士",
                "content": "在俄罗斯集市购物可以砍价，但在正规商场通常不讲价。俄罗斯人购物时很注重质量，喜欢买结实保暖的商品。"
            }
        ]
    },
    {
        "day": "День 13",
        "title": "Здоровье",
        "subtitle": "深化健康用语",
        "sections": [
            {
                "title": "进阶词汇",
                "words": [
                    ("Я плохо себя чувствую", "我感觉不舒服"),
                    ("У меня температура", "我发烧了"),
                    ("Где аптека?", "药店在哪里？"),
                    ("Мне нужен врач", "我需要看医生"),
                    ("У меня болит голова", "我头疼"),
                    ("Вызовите скорую помощь", "请叫救护车"),
                    ("Я аллергик", "我有过敏症"),
                    ("Мне нужны лекарства", "我需要买药"),
                ]
            },
            {
                "title": "语法要点",
                "words": [
                    ("У меня болит + N (身体部位)", "我...疼"),
                    ("Мне нужен/нужна/нужно + N", "我需要..."),
                    ("Я аллергик на + AKK", "我对...过敏"),
                    ("Вызвать + AKK (召唤)", "叫..."),
                ]
            },
            {
                "title": "实用对话",
                "dialogue": [
                    "— Алло! Скорая помощь?",
                    "— Да, что случилось?",
                    "— Мой друг упал и не может встать. Ему нужна помощь!",
                    "— Какой у вас адрес?",
                    "— Улица Пушкина, дом десять, квартира пять.",
                ]
            },
            {
                "title": "文化小贴士",
                "content": "在俄罗斯看病通常需要预约，急诊可以直接去医院。药店到处都是，购买处方药需要医生处方。俄罗斯有很好的医疗体系，但外语服务可能有限。"
            }
        ]
    },
    {
        "day": "День 14",
        "title": "Культура и искусство",
        "subtitle": "深化文化与艺术",
        "sections": [
            {
                "title": "进阶词汇",
                "words": [
                    ("Вы любите искусство?", "您喜欢艺术吗？"),
                    ("Мне нравится русская культура", "我喜欢俄罗斯文化"),
                    ("Когда открывается музей?", "博物馆什么时候开门？"),
                    ("Я хочу посетить театр", "我想去看话剧"),
                    ("Это очень красиво!", "这太美了！"),
                    ("У вас есть билеты на концерт?", "你们有音乐会票吗？"),
                    ("Кто ваш любимый писатель?", "您最喜欢的作家是谁？"),
                    ("Я люблю классическую музыку", "我喜欢古典音乐"),
                ]
            },
            {
                "title": "语法要点",
                "words": [
                    ("Мне нравится + N (我喜欢)", "......我喜欢"),
                    ("Кто ваш любимый + N?", "您最喜欢的......是谁？"),
                    ("Когда открывается + N?", "...什么时候开放？"),
                    ("Я хочу посетить + AKK", "我想参观..."),
                ]
            },
            {
                "title": "实用对话",
                "dialogue": [
                    "— Вы любите балет?",
                    "— Да, я очень люблю русский балет!",
                    "— Завтра вечером балет Лебединое озеро.",
                    "— Это классический спектакль! Я очень хочу пойти!",
                    "— Я могу купить билеты для нас обоих.",
                    "— Спасибо! Это было бы замечательно!",
                ]
            },
            {
                "title": "文化小贴士",
                "content": "俄罗斯文化底蕴深厚，芭蕾歌剧交响乐都是世界级的。莫斯科大剧院和马林斯基剧院是世界著名的艺术殿堂。俄罗斯人普遍受教育程度高，对文学诗歌和古典音乐有深厚兴趣。"
            }
        ]
    }
]

def create_slide(prs, day_info, output_path):
    """生成单日课件"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_tf = title_box.text_frame
    title_tf.paragraphs[0].text = "俄语进阶 " + day_info['day']
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = ACCENT
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(1))
    sub_tf = subtitle_box.text_frame
    sub_tf.paragraphs[0].text = day_info['title']
    sub_tf.paragraphs[0].font.size = Pt(28)
    sub_tf.paragraphs[0].font.color.rgb = TITLE_BG
    
    sub2_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(12.333), Inches(0.6))
    sub2_tf = sub2_box.text_frame
    sub2_tf.paragraphs[0].text = day_info['subtitle']
    sub2_tf.paragraphs[0].font.size = Pt(20)
    sub2_tf.paragraphs[0].font.color.rgb = DARK_TEXT
    
    # 内容页
    for section in day_info['sections']:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 标题栏背景
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
        bg.fill.solid()
        bg.fill.fore_color.rgb = TITLE_BG
        bg.line.fill.background()
        
        # 标题文字
        header = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.8))
        header_tf = header.text_frame
        header_tf.paragraphs[0].text = day_info['day'] + " - " + day_info['title'] + " | " + section['title']
        header_tf.paragraphs[0].font.size = Pt(24)
        header_tf.paragraphs[0].font.bold = True
        header_tf.paragraphs[0].font.color.rgb = WHITE
        header_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        y_start = Inches(1.2)
        
        if 'words' in section:
            # 词汇/语法页
            for i, (ru, cn) in enumerate(section['words']):
                y_pos = y_start + Inches(i * 0.65)
                
                ru_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(6), Inches(0.5))
                ru_tf = ru_box.text_frame
                ru_tf.paragraphs[0].text = ru
                ru_tf.paragraphs[0].font.size = Pt(18)
                ru_tf.paragraphs[0].font.bold = True
                ru_tf.paragraphs[0].font.color.rgb = DARK_TEXT
                
                cn_box = slide.shapes.add_textbox(Inches(6.5), y_pos, Inches(6), Inches(0.5))
                cn_tf = cn_box.text_frame
                cn_tf.paragraphs[0].text = cn
                cn_tf.paragraphs[0].font.size = Pt(18)
                cn_tf.paragraphs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)
                
                line = slide.shapes.add_shape(1, Inches(0.3), y_pos + Inches(0.5), Inches(12.733), Inches(0.01))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                line.line.fill.background()
        
        elif 'dialogue' in section:
            # 对话页
            for i, line_text in enumerate(section['dialogue']):
                y_pos = y_start + Inches(i * 0.75)
                
                is_question = line_text.strip().startswith("-")
                
                line_box = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(11.5), Inches(0.6))
                line_tf = line_box.text_frame
                line_tf.paragraphs[0].text = line_text
                line_tf.paragraphs[0].font.size = Pt(17)
                line_tf.paragraphs[0].font.color.rgb = TITLE_BG if is_question else DARK_TEXT
                if is_question:
                    line_tf.paragraphs[0].font.italic = True
        
        elif 'content' in section:
            # 文化贴士页
            content_box = slide.shapes.add_textbox(Inches(0.5), y_start, Inches(12.333), Inches(4))
            content_tf = content_box.text_frame
            content_tf.paragraphs[0].text = section['content']
            content_tf.paragraphs[0].font.size = Pt(20)
            content_tf.paragraphs[0].font.color.rgb = DARK_TEXT
            content_tf.paragraphs[0].line_spacing = 1.5
    
    prs.save(output_path)
    print("已生成: " + output_path)

def main():
    output_dir = Path.home() / ".openclaw/workspace/courses/俄语课件"
    output_dir.mkdir(parents=True, exist_ok=True)
    
    for i, day_info in enumerate(ADVANCED_CONTENT):
        day_num = i + 8  # Day 8-14
        output_path = output_dir / ("俄语进阶_Day" + str(day_num) + ".pptx")
        create_slide(None, day_info, str(output_path))
    
    print("\n全部7天进阶课件生成完毕！")

if __name__ == "__main__":
    main()