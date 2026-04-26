# -*- coding: utf-8 -*-
"""
道德经第17天 课件生成器
章节：第55-57章
风格：精读深度讲解版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
ACCENT = RGBColor(204, 0, 0)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_yuanwen_box(slide, lines, y=1.4, height=2.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(20)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=BLUE_LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(255, 255, 255)

def add_points(slide, points, y_start=3.1):
    y = y_start
    for title, content in points:
        pt_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT
        y += 1.1
    return y

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.5))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第17课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第55-57章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "含德之厚 · 希言自然 · 企者不立"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "含德之厚，比于赤子。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = PRIMARY
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第五十五章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = SECONDARY
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, 12)

# ===== 第2页：第五十五章详解 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s2, "第五十五章 · 含德之厚")
add_yuanwen_box(s2, [
    "含德之厚，比于赤子。",
    "蜂虿虺蛇不螫，猛兽不据，攫鸟不搏。",
    "骨弱筋柔而握固，未知牝牡之合而朘作，精之至也。",
    "",
    "终日号而不嗄，和之至也。",
    "",
    "知和曰常，知常曰明，",
    "益生曰祥，心使气曰强。",
    "物壮则老，谓之不道，不道早已。"
], y=1.4, height=4.0)
content2 = s2.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(2))
tf_c2 = content2.text_frame
tf_c2.word_wrap = True
pc2 = tf_c2.paragraphs[0]
pc2.text = "【白话解读】含德深厚的人，像初生婴儿——毒虫不咬他，猛兽不伤他。婴儿筋骨柔弱却握拳牢固，不知道男女之事却能勃起，这是精气充沛的表现。"
pc2.font.size = Pt(14)
pc2.font.color.rgb = TEXT
add_slide_number(s2, 2, 12)

# ===== 第3页：第五十五章核心 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s3, "第五十五章 · 核心思想")
add_core_box(s3, "🌟 核心观点", "最高的德行就像初生婴儿——纯真、柔弱、精气充沛，这才是道的境界。")
add_points(s3, [
    ("1. 含德之厚", "真正有德之人，像婴儿一样纯真无害，不具备攻击性。"),
    ("2. 精之至也", "精气充沛是道的表现——不是软弱，是内在充实。"),
    ("3. 和之至也", "阴阳调和是自然的状态——不是没有矛盾，是和谐平衡。"),
    ("4. 物壮则老", "过于强壮反而走向衰老——太过极端就不符合道了。")
])
add_slide_number(s3, 3, 12)

# ===== 第4页：第五十六章详解 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s4, "第五十六章 · 知者不言")
add_yuanwen_box(s4, [
    "知者不言，言者不知。",
    "",
    "塞其兑，闭其门，",
    "挫其锐，解其纷，",
    "和其光，同其尘，",
    "是谓玄同。",
    "",
    "故不可得而亲，不可得而疏；",
    "不可得而利，不可得而害；",
    "不可得而贵，不可得而贱。",
    "故为天下贵。"
], y=1.4, height=4.2)
content4 = s4.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(2))
tf_c4 = content4.text_frame
tf_c4.word_wrap = True
pc4 = tf_c4.paragraphs[0]
pc4.text = "【白话解读】真正知道的人不说话，说话的人不是真正知道。塞住嗜欲的孔窍，关闭欲念的心门，挫折锋芒，解散纷扰，涵敛光耀，混同尘世——这就是与道玄妙同一的境界。"
pc4.font.size = Pt(14)
pc4.font.color.rgb = TEXT
add_slide_number(s4, 4, 12)

# ===== 第5页：第五十六章核心 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s5, "第五十六章 · 核心思想")
add_core_box(s5, "🌟 核心观点", "与道玄同的人超越了亲疏、利害、贵贱的分别——这是天下最尊贵的境界。")
add_points(s5, [
    ("1. 知者不言", "真正明白道的人不夸夸其谈——道是不可言说的。"),
    ("2. 塞兑闭门", "关闭感官门窗，不被外物干扰——内观自心。"),
    ("3. 玄同", "与道合一的状态——不分彼此，没有分别心。"),
    ("4. 天下贵", "达到玄同境界的人，是天下最可贵的。")
])
add_slide_number(s5, 5, 12)

# ===== 第6页：现代反思 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s6, "第五十五-五十六章 · 现代反思")
items6 = [
    ("🏋️ 养生之道", "含德之厚比于赤子 → 养生不是吃补品，而是保养精气神，保持内心纯真。"),
    ("😤 情绪管理", "心使气曰强 → 控制情绪而不是被情绪控制，这才是真正的强大。"),
    ("🤐 言语之道", "知者不言 → 真正有智慧的人不急于表达，听比说更重要。"),
    ("⚖️ 平常心", "不可得而贵贱 → 对得失、荣辱保持平常心，不因外物而动摇。")
]
y6 = 1.4
for icon_title, explanation in items6:
    box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y6), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf6 = box.text_frame
    tf6.word_wrap = True
    p6_1 = tf6.paragraphs[0]
    p6_1.text = icon_title
    p6_1.font.size = Pt(16)
    p6_1.font.bold = True
    p6_1.font.color.rgb = PRIMARY
    p6_2 = tf6.add_paragraph()
    p6_2.text = explanation
    p6_2.font.size = Pt(14)
    p6_2.font.color.rgb = TEXT
    y6 += 1.4
add_slide_number(s6, 6, 12)

# ===== 第7页：第五十七章详解 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s7, "第五十七章 · 希言自然")
add_yuanwen_box(s7, [
    "希言自然。",
    "",
    "飘风不终朝，骤雨不终日。",
    "孰为此者？天地。",
    "天地尚不能久，而况于人乎？",
    "",
    "故从事于道者，道者同于道，",
    "德者同于德，失者同于失。",
    "同于道者，道亦乐得之；",
    "同于德者，德亦乐得之；",
    "同于失者，失亦乐得之。",
    "",
    "信不足焉，有不信焉。"
], y=1.4, height=4.4)
content7 = s7.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(9), Inches(2))
tf_c7 = content7.text_frame
tf_c7.word_wrap = True
pc7 = tf_c7.paragraphs[0]
pc7.text = "【白话解读】少说话才是自然的。狂风刮不了一早晨，暴雨下不了一整天。谁使它这样的？是天地。天地都不能持久，何况人呢？追求道的人与道合一，追求德的人与德合一。"
pc7.font.size = Pt(14)
pc7.font.color.rgb = TEXT
add_slide_number(s7, 7, 12)

# ===== 第8页：第五十七章核心 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s8, "第五十七章 · 核心思想")
add_core_box(s8, "🌟 核心观点", "狂风骤雨都不能持久，人更不能狂妄自大。顺应自然，少说多做，才能长久。")
add_points(s8, [
    ("1. 希言自然", "少说话才是道的运作方式——言多必失，不如行动。"),
    ("2. 飘风骤雨", "极端的手段不能持久——再大的狂风也有平息的时候。"),
    ("3. 从事于道", "追求什么就会得到什么——与道合一，道也乐于接纳你。"),
    ("4. 信不足", "诚信不足的人，人们自然不信任他——这是他自己的问题。")
])
add_slide_number(s8, 8, 12)

# ===== 第9页：西方思想 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s9, "🌍 西方思想对照")
phils = [
    ("卢梭 · 自然状态", "人类最初像婴儿一样纯真，社会腐化了自然状态。1754年"),
    ("赫拉克利特 · 万物流变", "一切都在变化中，没有什么是永恒的，包括狂风和暴雨。约前500年"),
    ("海德格尔 · 沉默", "真正的语言是沉默——言说要出于本真，不是空话。1927年"),
    ("怀特海 · 过程哲学", "实体是过程，而非静态物体——一切都在成为之中。1929年")
]
y9 = 1.4
for name_theory, desc_year in phils:
    box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y9), Inches(9), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf9 = box.text_frame
    tf9.word_wrap = True
    p9_1 = tf9.paragraphs[0]
    p9_1.text = name_theory
    p9_1.font.size = Pt(16)
    p9_1.font.bold = True
    p9_1.font.color.rgb = PRIMARY
    p9_2 = tf9.add_paragraph()
    p9_2.text = desc_year
    p9_2.font.size = Pt(14)
    p9_2.font.color.rgb = TEXT
    y9 += 1.3
add_slide_number(s9, 9, 12)

# ===== 第10页：现代应用 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s10, "📌 现代应用")
add_points(s10, [
    ("企业管理", "希言自然 → 领导少干预，给员工空间，让团队自我运作。"),
    ("个人成长", "从事于道者 → 找准方向，持续投入，时间会给你答案。"),
    ("危机处理", "飘风不终朝 → 再大的危机也会过去，保持冷静，不要过度反应。"),
    ("人际交往", "信不足焉 → 与人交往要守信用，一次失信可能毁掉所有信任。")
])
add_slide_number(s10, 10, 12)

# ===== 第11页：中国典故 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s11, "📜 中国典故：庄子的无用之树")
box11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.5))
box11.fill.solid()
box11.fill.fore_color.rgb = BG
tf11 = box11.text_frame
tf11.word_wrap = True
p11_1 = tf11.paragraphs[0]
p11_1.text = "【庄子·内篇·逍遥游】"
p11_1.font.size = Pt(16)
p11_1.font.bold = True
p11_1.font.color.rgb = PRIMARY
p11_2 = tf11.add_paragraph()
p11_2.text = "匠石之齐，至于曲辕，见栎社树。其大蔽数千牛，絜之百围。观者如市，匠伯不顾。弟子厌观之，走及匠石曰：'自吾执斧斤以随夫子，未尝见材如此其美也。先生不肯视，行不辍，何邪？'"
p11_2.font.size = Pt(16)
p11_2.font.color.rgb = TEXT
p11_2.space_before = Pt(10)
insight11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.1), Inches(9), Inches(1))
insight11.fill.solid()
insight11.fill.fore_color.rgb = BLUE_LIGHT
tfi11 = insight11.text_frame
tfi11.word_wrap = True
pi11 = tfi11.paragraphs[0]
pi11.text = "💡 启示"
pi11.font.size = Pt(18)
pi11.font.bold = True
pi11.font.color.rgb = RGBColor(255, 255, 255)
pi11_2 = tfi11.add_paragraph()
pi11_2.text = "这棵树因为\"没用\"所以没被砍伐——无用是大用。这与\"希言自然\"一样，都是道家顺其自然、不强求的智慧。"
pi11_2.font.size = Pt(14)
pi11_2.font.color.rgb = RGBColor(255, 255, 255)
app11 = s11.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2))
tfa11 = app11.text_frame
tfa11.word_wrap = True
pa11 = tfa11.paragraphs[0]
pa11.text = "【现代应用】职场：不必强求出人头地，做好本分，无用之地有时是最大的保护。生活中：知足常乐，不争才是最大的争。"
pa11.font.size = Pt(14)
pa11.font.color.rgb = TEXT
add_slide_number(s11, 11, 12)

# ===== 第12页：总结 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s12, "第17课 · 三章精华总结")
sums = [
    ("第五十五章", "含德之厚", "含德深厚如赤子，精气充沛，柔弱胜刚强，物壮则老。"),
    ("第五十六章", "知者不言", "知者不言，塞兑闭门，与道玄同，超越亲疏贵贱。"),
    ("第五十七章", "希言自然", "飘风骤雨不终日，少说多做，顺应自然才能长久。")
]
y12 = 1.4
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    y12 += 1.6
final = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(9), Inches(1))
final.fill.solid()
final.fill.fore_color.rgb = PRIMARY
tff = final.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "知者不言，言者不知。"
pf1.font.size = Pt(20)
pf1.font.bold = True
pf1.font.color.rgb = RGBColor(255, 255, 255)
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "—— 道德经第五十六章"
pf2.font.size = Pt(14)
pf2.font.color.rgb = RGBColor(200, 200, 200)
pf2.alignment = PP_ALIGN.CENTER
add_slide_number(s12, 12, 12)

prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第17天.pptx')
print("Done: 道德经_第17天.pptx (12 slides)")
