# -*- coding: utf-8 -*-
"""
道德经第32天 课件生成器
章节：第97-99章（致知、为而不争、以其不争）
风格：商务蓝/深色系，与第29/30/31天模板一致
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# 配色方案 - 商务蓝/深色系
PRIMARY = RGBColor(30, 60, 114)       # 深蓝
SECONDARY = RGBColor(70, 130, 180)    # 中蓝
ACCENT = RGBColor(204, 120, 50)       # 橙色强调
TEXT = RGBColor(51, 51, 51)           # 深灰文字
BG = RGBColor(245, 247, 250)          # 浅灰蓝背景
WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(20, 35, 60)        # 深蓝背景

# 字体配置
TITLE_SIZE = Pt(38)
HEADER_SIZE = Pt(30)
YUANWEN_SIZE = Pt(22)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(16)

def set_dark_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = HEADER_SIZE
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def add_yuanwen_section(slide, lines, y=1.5, height=3.2):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = SECONDARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line if line else "　"
        pp.font.size = YUANWEN_SIZE
        pp.font.color.rgb = WHITE
        pp.space_before = Pt(3)

def add_commentary_section(slide, title, content, y):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    box.line.color.rgb = SECONDARY
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"【{title}】"
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = BODY_SIZE
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=SECONDARY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(15)
    p2.font.color.rgb = WHITE

def add_points_list(slide, points, y_start=3.2):
    y = y_start
    for icon_title, explanation in points:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.25))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = icon_title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = explanation
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)
        y += 1.35
    return y

TOTAL = 14

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.6))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.9))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典32章 · 第32课"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = WHITE
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.6))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第97-99章"
p2.font.size = Pt(58)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.6), Inches(8), Inches(0.9))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "致知 · 为而不争 · 以其不争"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = DARK_BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "知人者智，自知者明。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = WHITE
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第七十三章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = RGBColor(180, 180, 180)
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, TOTAL)

# ===== 第2页：三章概览 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s2, "今日课程")
chapters = [
    ("第九十七章", "致知", "古之善为道者，微妙玄通，深不可识。真正的智者超越常人认知，难以描述。"),
    ("第九十八章", "为而不争", "为而不争，与而不取。不居功，不自是，以玄德感化天下，自然归附。"),
    ("第九十九章", "以其不争", "上善若水，水善利万物而不争。夫唯不争，天下莫能与之争。不争之争，是为玄德。")
]
y2 = 1.5
for num, title, desc in chapters:
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y2), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf2c = box.text_frame
    tf2c.word_wrap = True
    p2c1 = tf2c.paragraphs[0]
    p2c1.text = f"📖 {num}：{title}"
    p2c1.font.size = Pt(20)
    p2c1.font.bold = True
    p2c1.font.color.rgb = PRIMARY
    p2c2 = tf2c.add_paragraph()
    p2c2.text = desc
    p2c2.font.size = Pt(16)
    p2c2.font.color.rgb = TEXT
    p2c2.space_before = Pt(5)
    y2 += 1.65
add_slide_number(s2, 2, TOTAL)

# ===== 第3页：第九十七章 原文 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s3, "第九十七章 · 致知")
add_yuanwen_section(s3, [
    "古之善为道者，微妙玄通，深不可识。",
    "",
    "夫唯不可识，故强为之容：",
    "",
    "豫兮，若冬涉川；",
    "犹兮，若畏四邻；",
    "俨兮，其若客；",
    "涣兮，若冰之将释；",
    "敦兮，其若朴；",
    "旷兮，其若谷；",
    "混兮，其若浊。",
    "",
    "孰能浊以止？静之徐清。",
    "孰能安以久？动之徐生。",
    "",
    "保此道者，不欲盈。",
    "夫唯不盈，故能敝不新成。"
], y=1.5, height=4.8)
add_commentary_section(s3, "白话解读",
    "古代善于行道的人，微妙玄通，深不可识。正因为不可识，所以只能勉强形容：小心啊，像冬天涉水过河；迟疑啊，像害怕四周邻国；庄重啊，像在做客；流动啊，像冰将要融化；敦厚啊，像未经雕琢的木头；空旷啊，像幽深的山谷；浑厚啊，像浑浊的黄河。谁能让浑浊停止？安静下来就会慢慢清澈。谁能让安定长久？运动起来就会慢慢生长。保持这个道的人，不追求盈满。正因为不盈满，所以能够保持旧的状态而不追求新成。",
    6.4)
add_slide_number(s3, 3, TOTAL)

# ===== 第4页：第九十七章 逐句释义 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s4, "第九十七章 · 逐句释义")
add_points_list(s4, [
    ("① 古之善为道者，微妙玄通，深不可识",
     "古代善于行道的人，微妙玄通，深不可识——真正得道之人，其境界超越常人理解，无法用言语形容。"),
    ("② 豫兮，若冬涉川；犹兮，若畏四邻",
     "小心迟疑，如冬天涉水、害怕四邻——行事谨慎，不轻举妄动，敬畏天地人事。"),
    ("③ 俨兮，其若客；涣兮，若冰之将释",
     "庄重如客，流动如冰将融——既保持严肃恭敬，又自然流畅，不拘泥刻板。"),
    ("④ 敦兮，其若朴；旷兮，其若谷；混兮，其若浊",
     "敦厚如朴，空旷如谷，浑厚如浊——外表朴素、内心深邃、与世无争。"),
    ("⑤ 保此道者，不欲盈",
     "保持此道的人，不追求盈满——不满不溢，知止不贪，所以能保持长久。")
], y_start=1.5)
add_slide_number(s4, 4, TOTAL)

# ===== 第5页：第九十七章 本章小结 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s5, "第九十七章 · 本章小结")
add_core_box(s5, "🌟 核心观点", "致知——真正的智者微妙玄通，深不可识。他们谨慎若冬涉川，敦厚若朴，旷达若谷，与世无争，不欲盈满。")
add_points_list(s5, [
    ("1. 深不可识", "善为道者微妙玄通——真正的得道者境界深远，非常人所能识。"),
    ("2. 七种形容", "豫、犹、俨、涣、敦、旷、混——谨慎、迟疑、庄重、流动、敦厚、空旷、浑厚。"),
    ("3. 浊清安久", "静之徐清，动之徐生——动荡中方能清晰，安静中方能长久。"),
    ("4. 不欲盈", "保此道者不欲盈——不满不溢，知止不贪，方能持久。")
], y_start=3.1)
add_slide_number(s5, 5, TOTAL)

# ===== 第6页：第九十八章 原文 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s6, "第九十八章 · 为而不争")
add_yuanwen_section(s6, [
    "为而不争，与而不取。",
    "",
    "善为道者，不以知争。",
    "",
    "以其不争，",
    "故天下莫能与之争。",
    "",
    "为而不争，功成不居。",
    "夫唯不居，是以不去。",
    "",
    "玄德深矣，远矣，",
    "与物反矣，然后乃至大顺。"
], y=1.5, height=3.8)
add_commentary_section(s6, "白话解读",
    "有所作为但不争夺，给予但不索取。善于行道的人，不以智慧争夺。正因为他不争，所以天下没有人能与他争夺。有作为但不争夺，功业成就却不居功。正因为不居功，所以功业不会离去。玄德深奥啊，遥远啊，与万物同返于道，然后才能达到最大的顺应。",
    5.4)
add_slide_number(s6, 6, TOTAL)

# ===== 第7页：第九十八章 逐句释义 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s7, "第九十八章 · 逐句释义")
add_points_list(s7, [
    ("① 为而不争，与而不取",
     "有所作为但不争夺，给予但不索取——为道者无私无欲，只管付出，不问收获。"),
    ("② 善为道者，不以知争",
     "善于行道的人，不以智慧争夺——不用智巧争抢，不以才华压人。"),
    ("③ 以其不争，故天下莫能与之争",
     "正因为他有不争的美德，所以天下没有人能与他争夺——不争者无敌。"),
    ("④ 功成不居，夫唯不居，是以不去",
     "功业成就却不居功，正因为不居功，所以功业不会离去——退让反而保全。"),
    ("⑤ 玄德深矣，远矣",
     "玄德深奥遥远，与万物同返于道——然后达到最大的顺应境界。")
], y_start=1.5)
add_slide_number(s7, 7, TOTAL)

# ===== 第8页：第九十八章 本章小结 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s8, "第九十八章 · 本章小结")
add_core_box(s8, "🌟 核心观点", "为而不争——行道者有所作为但不争夺，功成不居。不争之争，是为玄德。夫唯不争，天下莫能与之争。")
add_points_list(s8, [
    ("1. 为而不争", "做事尽心尽力，但不与人争夺——只问付出，不问回报。"),
    ("2. 与而不取", "给予但不索取——慷慨奉献，不求报偿。"),
    ("3. 不以知争", "不用智巧争夺——不耍小聪明，不争一时之快。"),
    ("4. 功成不居", "成就功业却不居功——谦退让贤，反而保全。"),
    ("5. 玄德大顺", "玄德深远，与物反矣——返回道的本源，达到大顺境界。")
], y_start=3.1)
add_slide_number(s8, 8, TOTAL)

# ===== 第9页：第九十九章 原文 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s9, "第九十九章 · 以其不争")
add_yuanwen_section(s9, [
    "上善若水，水善利万物而不争。",
    "",
    "处众人之所恶，",
    "故几于道。",
    "",
    "居善地，心善渊，",
    "与善仁，言善信，",
    "正善治，事善能，动善时。",
    "",
    "夫唯不争，",
    "故无尤。",
    "",
    "江海所以能为百谷王者，",
    "以其善下之，是以能为百谷王。"
], y=1.5, height=4.5)
add_commentary_section(s9, "白话解读",
    "最高的善就像水一样，水善于利益万物而不争夺。停留在众人所厌恶的地方，所以接近于道。居住善于选择地方，心胸善于沉静，与人交往善于仁爱，说话善于诚信，执政善于治理，做事善于发挥才能，行动善于把握时机。正因为不争夺，所以没有过失。江海之所以能成为山谷之王，是因为它善于处在低下的位置，所以能成为山谷之王。",
    6.1)
add_slide_number(s9, 9, TOTAL)

# ===== 第10页：第九十九章 逐句释义 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s10, "第九十九章 · 逐句释义")
add_points_list(s10, [
    ("① 上善若水，水善利万物而不争",
     "最高的善就像水——水滋润万物而不争功，处下不争，近于道。"),
    ("② 处众人之所恶，故几于道",
     "停留在众人所厌恶的低洼处，所以接近于道——甘居人下，是道的品质。"),
    ("③ 七善：居、心、与、言、正、事、动",
     "居善地，心善渊，与善仁，言善信，正善治，事善能，动善时——处世的七种智慧。"),
    ("④ 夫唯不争，故无尤",
     "正因为不争夺，所以没有过失——不争则无矛盾，无尤则无忧。"),
    ("⑤ 江海善下，故能为百谷王",
     "江海善于处在低下位置，所以能成为山谷之王——处下是成就伟大的秘密。")
], y_start=1.5)
add_slide_number(s10, 10, TOTAL)

# ===== 第11页：第九十九章 本章小结 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s11, "第九十九章 · 本章小结")
add_core_box(s11, "🌟 核心观点", "上善若水——水善利万物而不争，处众人之所恶，故几于道。夫唯不争，天下莫能与之争。")
add_points_list(s11, [
    ("1. 上善若水", "最高的善像水一样——柔弱却能攻坚，利他而不争。"),
    ("2. 处下不争", "处众人之所恶——甘居人下，不与人争，是道的品质。"),
    ("3. 七善之道", "居善地，心善渊，与善仁，言善信，正善治，事善能，动善时——完整的处世智慧。"),
    ("4. 不争无尤", "不争夺就没有过失——无欲则刚，无争则无忧。"),
    ("5. 善下为王", "江海善下，故能为百谷王——处下是成就伟大的秘密。")
], y_start=3.1)
add_slide_number(s11, 11, TOTAL)

# ===== 第12页：三章总览与核心要义 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s12, "三章精华 · 一以贯之")
sums = [
    ("第九十七章", "致知", "古之善为道者，微妙玄通，深不可识。不欲盈，保此道者故能敝不新成——谨慎内敛，不求盈满。"),
    ("第九十八章", "为而不争", "为而不争，与而不取，功成不居。夫唯不争，天下莫能与之争——玄德深远，大顺于道。"),
    ("第九十九章", "以其不争", "上善若水，水善利万物而不争。夫唯不争，故无尤——处下不争，无尤无忧。")
]
y12 = 1.5
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.55))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    p12_2.space_before = Pt(5)
    y12 += 1.65
add_slide_number(s12, 12, TOTAL)

# ===== 第13页：现实联想 =====
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s13, "📌 现实联想")
apps = [
    ("🧠 认知谦逊", "致知 → 知道越多，越意识到自己无知。真正有智慧的人不会自以为是，而是保持开放心态。"),
    ("💼 为而不争", "为而不争 → 在职场中做事尽心，但不为名利争夺，容易获得信任和重用。"),
    ("🌊 上善若水", "水善利万物而不争 → 做生意也好，做管理也好，先让别人得利，自己自然受益。"),
    ("🏔️ 甘居人下", "处众人之所恶 → 年轻人主动承担脏活累活，看似吃亏，实则积累经验和信任。"),
    ("⏳ 不欲盈", "保此道者不欲盈 → 成功时不骄傲，不自满，保持空杯心态，才能持续进步。")
]
add_points_list(s13, apps, y_start=1.5)
add_slide_number(s13, 13, TOTAL)

# ===== 第14页：管理启示 =====
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s14, "💼 管理启示")
insights = [
    ("🎯 为而不争的领导者", "领导者做事但不为争夺功劳，让团队成员获得荣誉，自己退居幕后。这样的领导者更能凝聚团队，获得追随。"),
    ("💧 水的管理哲学", "像水一样滋养员工，而不是压迫员工。给予员工成长机会，关心员工福祉，自然赢得尊重。"),
    ("📊 不居功反而保功", "不居功的管理者反而能保全自己的地位。争功会失去人心，不争则众望所归。"),
    ("🌊 处下成就伟大", "江海善下故能为百谷王——管理者谦下，员工愿意归附。傲上的管理者最终孤立无援。"),
    ("⚖️ 七善的领导者素质", "居善地（定好位）、心善渊（有深度）、与善仁（有爱心）、言善信（有信用）、正善治（有方法）、事善能（有才能）、动善时（把握时机）——完整的领导力模型。")
]
add_points_list(s14, insights, y_start=1.5)
add_slide_number(s14, 14, TOTAL)

# 保存文件
output_path = '/Users/mac/.openclaw/workspace/courses/道德经_第32天.pptx'
prs.save(output_path)
print(f"Done: {output_path}")