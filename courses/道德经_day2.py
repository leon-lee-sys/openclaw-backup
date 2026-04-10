from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

# ===== 第1页：封面 =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar.line.fill.background()
txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第二课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

main_title = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main_title.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第4-6章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER

subtitle = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = subtitle.text_frame
p3 = tf3.paragraphs[0]
p3.text = "养身之道 · 天地之鉴 · 虚无之用"
p3.font.size = Pt(28)
p3.font.color.rgb = RGBColor(70, 130, 180)
p3.alignment = PP_ALIGN.CENTER

quote_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_q = quote_box.text_frame
tf_q.word_wrap = True
p_q1 = tf_q.paragraphs[0]
p_q1.text = "天地不仁，以万物为刍狗。"
p_q1.font.size = Pt(22)
p_q1.font.italic = True
p_q1.font.color.rgb = RGBColor(30, 60, 114)
p_q1.alignment = PP_ALIGN.CENTER
p_q2 = tf_q.add_paragraph()
p_q2.text = "—— 《道德经》第五章"
p_q2.font.size = Pt(14)
p_q2.font.color.rgb = RGBColor(150, 150, 150)
p_q2.alignment = PP_ALIGN.CENTER

add_slide_number(slide1, 1, 10)

# ===== 第2页：第四章详解 =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar2.fill.solid()
title_bar2.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar2.line.fill.background()
txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf5 = txBox2.text_frame
p5 = tf5.paragraphs[0]
p5.text = "第四章 · 养身之道"
p5.font.size = Pt(32)
p5.font.bold = True
p5.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.8))
yuanwen_box.fill.solid()
yuanwen_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw = yuanwen_box.text_frame
tf_yw.word_wrap = True
p_yw = tf_yw.paragraphs[0]
p_yw.text = "【原文】"
p_yw.font.size = Pt(18)
p_yw.font.bold = True
p_yw.font.color.rgb = RGBColor(30, 60, 114)
for line in ["道冲而用之或不盈，渊兮似万物之宗。", "挫其锐，解其纷，和其光，同其尘。"]:
    p = tf_yw.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(9), Inches(4))
tf_c = content_box.text_frame
tf_c.word_wrap = True
p_c1 = tf_c.paragraphs[0]
p_c1.text = "【白话解读】"
p_c1.font.size = Pt(18)
p_c1.font.bold = True
p_c1.font.color.rgb = RGBColor(30, 60, 114)
解读 = [
    ("道冲而不盈", "'道'是虚空的，却用之不竭——像深渊一样深邃，是万物的本源。"),
    ("挫锐解纷", "收敛锋芒，解决纷争——不露锋芒，才能避免冲突。"),
    ("和光同尘", "混合光彩，混同尘埃——与世无争，随俗而不染。"),
    ("核心：谦下", "保持虚空谦下，不自满，不显露，自然而然")
]
for title, content in 解读:
    p = tf_c.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(8)

add_slide_number(slide2, 2, 10)

# ===== 第3页：第四章生活应用 =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar3.fill.solid()
title_bar3.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar3.line.fill.background()
txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf6 = txBox3.text_frame
p6 = tf6.paragraphs[0]
p6.text = "第四章 · 生活智慧"
p6.font.size = Pt(32)
p6.font.bold = True
p6.font.color.rgb = RGBColor(255, 255, 255)

core_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.3))
core_box.fill.solid()
core_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_core = core_box.text_frame
tf_core.word_wrap = True
p_core = tf_core.paragraphs[0]
p_core.text = "🌟 核心智慧"
p_core.font.size = Pt(18)
p_core.font.bold = True
p_core.font.color.rgb = RGBColor(255, 255, 255)
p_core2 = tf_core.add_paragraph()
p_core2.text = "真正有修养的人，不显山露水，不锋芒毕露，而是谦下包容，与光同尘。"
p_core2.font.size = Pt(16)
p_core2.font.color.rgb = RGBColor(255, 255, 255)

points = [
    ("1. 职场启示", "有才能不炫耀，反而更受尊重。锋芒毕露容易招致嫉妒，谦下守柔才能长久。"),
    ("2. 处世哲学", "与人相处不争是非，不显自己高明，反而能化解矛盾，赢得人心。"),
    ("3. 自我修养", "内心保持虚空，像空杯一样，才能不断学习进步，自满则止步。")
]
y = 2.9
for title, content in points:
    pt_box = slide3.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
    tf_pt = pt_box.text_frame
    tf_pt.word_wrap = True
    p_pt1 = tf_pt.paragraphs[0]
    p_pt1.text = title
    p_pt1.font.size = Pt(18)
    p_pt1.font.bold = True
    p_pt1.font.color.rgb = RGBColor(30, 60, 114)
    p_pt2 = tf_pt.add_paragraph()
    p_pt2.text = content
    p_pt2.font.size = Pt(14)
    p_pt2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.2

add_slide_number(slide3, 3, 10)

# ===== 第4页：第五章详解 =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar4.fill.solid()
title_bar4.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar4.line.fill.background()
txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf7 = txBox4.text_frame
p7 = tf7.paragraphs[0]
p7.text = "第五章 · 天地之鉴"
p7.font.size = Pt(32)
p7.font.bold = True
p7.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.8))
yuanwen_box2.fill.solid()
yuanwen_box2.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw2 = yuanwen_box2.text_frame
tf_yw2.word_wrap = True
p_yw2 = tf_yw2.paragraphs[0]
p_yw2.text = "【原文】"
p_yw2.font.size = Pt(18)
p_yw2.font.bold = True
p_yw2.font.color.rgb = RGBColor(30, 60, 114)
for line in ["天地不仁，以万物为刍狗；", "圣人不仁，以百姓为刍狗。"]:
    p = tf_yw2.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box2 = slide4.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(9), Inches(4))
tf_c2 = content_box2.text_frame
tf_c2.word_wrap = True
p_c4 = tf_c2.paragraphs[0]
p_c4.text = "【白话解读】"
p_c4.font.size = Pt(18)
p_c4.font.bold = True
p_c4.font.color.rgb = RGBColor(30, 60, 114)
解读2 = [
    ("天地不仁", "天地没有偏爱，对万物一视同仁——不像人会有偏心私情。"),
    ("刍狗之喻", "刍狗是祭祀用的草狗，用时尊重，用完丢弃——天地对万物亦如此。"),
    ("圣人不仁", "圣人也像天地一样，没有偏爱，对百姓一视同仁。"),
    ("核心：平等", "真正的仁爱是不偏私的顺其自然，而非有差别的偏爱")
]
for title, content in 解读2:
    p = tf_c2.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide4, 4, 10)

# ===== 第5页：天地不仁典故 =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar5.fill.solid()
title_bar5.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar5.line.fill.background()
txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf8 = txBox5.text_frame
p8 = tf8.paragraphs[0]
p8.text = "📜 典故：天地不仁"
p8.font.size = Pt(32)
p8.font.bold = True
p8.font.color.rgb = RGBColor(255, 255, 255)

story_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
story_box.fill.solid()
story_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_story = story_box.text_frame
tf_story.word_wrap = True
p_s1 = tf_story.paragraphs[0]
p_s1.text = "【自然法则】"
p_s1.font.size = Pt(16)
p_s1.font.bold = True
p_s1.font.color.rgb = RGBColor(30, 60, 114)
p_s2 = tf_story.add_paragraph()
p_s2.text = "太阳照耀大地，不会因为你是好人就多照一小时；雨水滋润万物，不会因为你是坏人就少下一滴雨。天地按照规律运行，没有私心，没有偏袒。"
p_s2.font.size = Pt(15)
p_s2.font.color.rgb = RGBColor(51, 51, 51)
p_s2.space_before = Pt(8)
p_s3 = tf_story.add_paragraph()
p_s3.text = "这正是老子所说的\"天道无亲，常与善人\"——天道没有偏亲，但永远帮助顺应规律的人。"
p_s3.font.size = Pt(15)
p_s3.font.color.rgb = RGBColor(51, 51, 51)

change_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.6), Inches(9), Inches(1))
change_box.fill.solid()
change_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_ch = change_box.text_frame
tf_ch.word_wrap = True
p_ch = tf_ch.paragraphs[0]
p_ch.text = "⚡ 现代启示"
p_ch.font.size = Pt(18)
p_ch.font.bold = True
p_ch.font.color.rgb = RGBColor(255, 255, 255)
p_ch2 = tf_ch.add_paragraph()
p_ch2.text = "与其祈求天保佑，不如顺应天道，自己掌握命运"
p_ch2.font.size = Pt(16)
p_ch2.font.color.rgb = RGBColor(255, 255, 255)

app_box = slide5.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2.5))
tf_app = app_box.text_frame
tf_app.word_wrap = True
p_a1 = tf_app.paragraphs[0]
p_a1.text = "【处世智慧】"
p_a1.font.size = Pt(16)
p_a1.font.bold = True
p_a1.font.color.rgb = RGBColor(30, 60, 114)
apps = [
    "• 不怨天尤人：事情成败有自身原因，也有客观规律",
    "• 顺应时势：看清大势所趋，不逆势而行",
    "• 公平待人：对人对事一视同仁，不因私废公"
]
for app in apps:
    p = tf_app.add_paragraph()
    p.text = app
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide5, 5, 10)

# ===== 第6页：西方哲学对照 =====
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar6.fill.solid()
title_bar6.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar6.line.fill.background()
txBox6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf9 = txBox6.text_frame
p9 = tf9.paragraphs[0]
p9.text = "🌍 西方思想对照"
p9.font.size = Pt(32)
p9.font.bold = True
p9.font.color.rgb = RGBColor(255, 255, 255)

philosophers = [
    ("🇬🇷 斯多葛学派", "\"自然法\"", "宇宙有内在规律，人应顺应理性与自然法则生活，而非被情感左右。", "约前300年"),
    ("🇬🇧 休谟", "\"情感主义\"", "理性是情感的奴隶，人类的道德来自同情心，而非冷酷的理性计算。", "1751年"),
    ("🇩🇪 康德", "\"绝对命令\"", "道德法则应普遍适用——你的行为准则能否成为所有人的准则？", "1785年"),
    ("🇫🇷 卢梭", "\"自然状态\"", "人类本性善良，是社会腐化了人，应回归自然、顺应人性。", "1754年")
]
y = 1.4
for name, theory, desc, year in philosophers:
    phil_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.2))
    phil_box.fill.solid()
    phil_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_phil = phil_box.text_frame
    tf_phil.word_wrap = True
    p_n = tf_phil.paragraphs[0]
    p_n.text = f"{name} · {theory} · {year}"
    p_n.font.size = Pt(16)
    p_n.font.bold = True
    p_n.font.color.rgb = RGBColor(30, 60, 114)
    p_d = tf_phil.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.35

add_slide_number(slide6, 6, 10)

# ===== 第7页：第六章详解 =====
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar7.fill.solid()
title_bar7.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar7.line.fill.background()
txBox7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf10 = txBox7.text_frame
p10 = tf10.paragraphs[0]
p10.text = "第六章 · 虚无之用"
p10.font.size = Pt(32)
p10.font.bold = True
p10.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box3 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
yuanwen_box3.fill.solid()
yuanwen_box3.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw3 = yuanwen_box3.text_frame
tf_yw3.word_wrap = True
p_yw3 = tf_yw3.paragraphs[0]
p_yw3.text = "【原文】"
p_yw3.font.size = Pt(18)
p_yw3.font.bold = True
p_yw3.font.color.rgb = RGBColor(30, 60, 114)
for line in ["谷神不死，是谓玄牝。", "玄牝之门，是谓天地根。"]:
    p = tf_yw3.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box3 = slide7.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(4))
tf_c3 = content_box3.text_frame
tf_c3.word_wrap = True
p_c7 = tf_c3.paragraphs[0]
p_c7.text = "【白话解读】"
p_c7.font.size = Pt(18)
p_c7.font.bold = True
p_c7.font.color.rgb = RGBColor(30, 60, 114)
解读3 = [
    ("谷神不死", "'谷'象征虚空，'神'象征变化——虚空的变化永不停歇，生生不息。"),
    ("玄牝之门", "'玄牝'是微妙的母性——这是天地产生的根源所在。"),
    ("虚空的力量", "正因为虚空，才能容纳万物；正因为无用，才能无所不用。"),
    ("核心：守虚", "保持内心虚空，像山谷一样，才能获得无限的生命力")
]
for title, content in 解读3:
    p = tf_c3.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide7, 7, 10)

# ===== 第8页：道家养生智慧 =====
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar8.fill.solid()
title_bar8.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar8.line.fill.background()
txBox8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf11 = txBox8.text_frame
p11 = tf11.paragraphs[0]
p11.text = "📜 养生智慧：心斋坐忘"
p11.font.size = Pt(32)
p11.font.bold = True
p11.font.color.rgb = RGBColor(255, 255, 255)

history_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
history_box.fill.solid()
history_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_hist = history_box.text_frame
tf_hist.word_wrap = True
p_h1 = tf_hist.paragraphs[0]
p_h1.text = "【庄子·人间世】"
p_h1.font.size = Pt(16)
p_h1.font.bold = True
p_h1.font.color.rgb = RGBColor(30, 60, 114)
p_h2 = tf_hist.add_paragraph()
p_h2.text = "颜回问孔子如何\"心斋\"，孔子答曰：\"若一志，无听之以耳而听之以心，无听之以心而听之以气。听止于耳，心止于符。虚而待物，唯道集虚。\"——排除杂念，让心如虚空，才能与道合一。"
p_h2.font.size = Pt(15)
p_h2.font.color.rgb = RGBColor(51, 51, 51)

measures_box = slide8.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(4))
tf_m = measures_box.text_frame
tf_m.word_wrap = True
p_m1 = tf_m.paragraphs[0]
p_m1.text = "【现代养生启示】"
p_m1.font.size = Pt(18)
p_m1.font.bold = True
p_m1.font.color.rgb = RGBColor(30, 60, 114)
measures = [
    ("静坐冥想", "每天静心片刻，放空杂念，让身心回归虚空"),
    ("减少欲望", "欲海难填，知足常乐，减少物欲累赘"),
    ("顺应作息", "日出而作，日落而息，与自然节奏同步"),
    ("形神兼养", "不只是锻炼身体，更要涵养精神心态")
]
for title, content in measures:
    p = tf_m.add_paragraph()
    p.text = f"✓ {title}：{content}"
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(8)

add_slide_number(slide8, 8, 10)

# ===== 第9页：生活应用 =====
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar9.fill.solid()
title_bar9.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar9.line.fill.background()
txBox9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf12 = txBox9.text_frame
p12 = tf12.paragraphs[0]
p12.text = "💡 生活中的智慧"
p12.font.size = Pt(32)
p12.font.bold = True
p12.font.color.rgb = RGBColor(255, 255, 255)

app_items = [
    ("🏢 团队管理", "领导者的\"虚空\"——不放狂言，不自以为高明\n倾听下属建议，让团队发挥创造力\n案例：张良\"运筹帷幄\"而不亲临战场"),
    ("🧘 心态修养", "遇到挫折时保持虚空心态\n不被情绪填满，才能看清出路\n像山谷一样，空才能容"),
    ("💼 商业智慧", "最成功的商业是\"无形\"的——如平台、生态\n不直接干预，而是创造环境\n案例：微信\"连接一切\"而非\"控制一切\""),
    ("👨‍👩‍👧 家庭教育", "不把自己的期望强加给孩子\n给他们空间，让他们按自己的方式成长\n父母的\"虚空\"是给孩子最好的礼物")
]
y = 1.4
for title, content in app_items:
    app_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.3))
    app_box.fill.solid()
    app_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_app = app_box.text_frame
    tf_app.word_wrap = True
    p_a = tf_app.paragraphs[0]
    p_a.text = title
    p_a.font.size = Pt(16)
    p_a.font.bold = True
    p_a.font.color.rgb = RGBColor(30, 60, 114)
    p_b = tf_app.add_paragraph()
    p_b.text = content
    p_b.font.size = Pt(13)
    p_b.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.45

add_slide_number(slide9, 9, 10)

# ===== 第10页：精华总结 =====
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar10.fill.solid()
title_bar10.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar10.line.fill.background()
txBox10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf13 = txBox10.text_frame
p13 = tf13.paragraphs[0]
p13.text = "💎 精华提炼 · 第二课总结"
p13.font.size = Pt(32)
p13.font.bold = True
p13.font.color.rgb = RGBColor(255, 255, 255)

summary_items = [
    ("第四章·养身之道", "挫锐解纷，和光同尘——收敛锋芒，谦下包容"),
    ("第五章·天地之鉴", "天地不仁，公平无偏——顺应规律，不怨不尤"),
    ("第六章·虚无之用", "谷神不死，玄牝之门——保持虚空，生生不息")
]
y = 1.4
for title, content in summary_items:
    s_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.9))
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_s = s_box.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.text = f"{title}"
    p_s.font.size = Pt(16)
    p_s.font.bold = True
    p_s.font.color.rgb = RGBColor(30, 60, 114)
    p_s2 = tf_s.add_paragraph()
    p_s2.text = content
    p_s2.font.size = Pt(14)
    p_s2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.05

conclusion_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
conclusion_box.fill.solid()
conclusion_box.fill.fore_color.rgb = RGBColor(30, 60, 114)
tf_con = conclusion_box.text_frame
tf_con.word_wrap = True
p_con = tf_con.paragraphs[0]
p_con.text = "🌟 一句话总结"
p_con.font.size = Pt(18)
p_con.font.bold = True
p_con.font.color.rgb = RGBColor(255, 255, 255)
p_con2 = tf_con.add_paragraph()
p_con2.text = "上善若水，虚怀若谷"
p_con2.font.size = Pt(24)
p_con2.font.bold = True
p_con2.font.color.rgb = RGBColor(255, 255, 255)
p_con2.alignment = PP_ALIGN.CENTER
p_con3 = tf_con.add_paragraph()
p_con3.text = "最高的善就像水一样，谦下包容；最深的修养是内心如山谷般虚空。"
p_con3.font.size = Pt(16)
p_con3.font.color.rgb = RGBColor(200, 200, 200)

add_slide_number(slide10, 10, 10)

# 保存
prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第2天.pptx')
print("PPT已生成：道德经_第2天.pptx")
