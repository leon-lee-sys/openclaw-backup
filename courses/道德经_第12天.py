#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""道德经第12天课件 - 第三十七章至第三十九章"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def add_title(prs, text, top=0.3, size=32, color=RGBColor(30, 58, 95)):
    box = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.333), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return box

def add_content(prs, lines, top=1.3):
    box = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.333), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(10)

def add_page_num(prs, num, total):
    box = prs.slides[-1].shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(150, 150, 150)

def create_daodejing_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 10

    # ========== 封面 ==========
    s = create_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(30, 58, 95)
    bg.line.fill.background()
    add_title(prs, "国学经典24章 · 第12课", top=2.2, size=48, color=RGBColor(255,255,255))
    add_title(prs, "第三十七章 至 第三十九章", top=3.3, size=28, color=RGBColor(200,220,255))
    add_title(prs, "道常无为 · 始制有名 · 上德不德", top=4.1, size=20, color=RGBColor(150,180,220))
    add_title(prs, "2026年4月19日（周日）", top=5.5, size=16, color=RGBColor(150,150,150))

    # ========== 第三十七章 ==========
    s = create_slide(prs)
    add_title(prs, "第三十七章：道常无为", size=28)
    lines = [
        "【原文】",
        "道常无为而无不为。候王若能守之，万物将自化。",
        "化而欲作，吾将镇之以无名之朴。无名之朴，夫亦将无欲。",
        "不欲以静，天下将自定。",
        "",
        "【白话解读】",
        "  道永远是顺其自然、不妄为的，却没有什么做不到。",
        "  侯王若能遵循道，万物就会自然生化。",
        "  有人起贪欲时，我就用道的真朴来镇服。",
        "  保持无欲状态，社会自然安定。",
        "",
        "【要点】",
        "① 无为而治：顺应自然，不强行干预",
        "② 自化原理：万物按道运行，自我生化",
        "③ 无欲则静：克制私欲，社会自然安定"
    ]
    add_content(prs, lines)
    add_page_num(prs, 2, total)

    # ========== 第三十七章 生活智慧 ==========
    s = create_slide(prs)
    add_title(prs, "第三十七章：生活智慧", size=28)
    lines = [
        "核心智慧",
        "",
        "【企业管理】顺其自然，减少过度管控",
        "  给员工足够的自主空间，让他们自我驱动",
        "  出现问题时用道（价值观）来引导，而非强制",
        "",
        "【家庭教育】静待花开，不过度干预",
        "  让孩子自然成长，不过分管控",
        "  父母保持内心平静，家庭自然和谐",
        "",
        "【个人修养】少私寡欲，抱朴守真",
        "  减少物质欲望，追求精神充实",
        "  以静制动，内心平静则外事顺遂"
    ]
    add_content(prs, lines)
    add_page_num(prs, 3, total)

    # ========== 第三十八章 ==========
    s = create_slide(prs)
    add_title(prs, "第三十八章：始制有名", size=28)
    lines = [
        "【原文】",
        "上德不德，是以有德；下德不失德，是以无德。",
        "上德无为而无以为；下德为之而有以为。",
        "上仁为之而无以为；上义为之而有以为。",
        "上礼为之而莫之应，则攘臂而扔之。",
        "故失道而后德，失德而后仁，失仁而后义，失义而后礼。",
        "夫礼者，忠信之薄，而乱之首。",
        "前识者，道之华，而愚之始。",
        "",
        "【要点】",
        "① 德的层次：从道→德→仁→义→礼，层层递减",
        "② 上德境界：不刻意表现德，自然而有德",
        "③ 礼的警示：礼是社会混乱的开始，要回归道的本真"
    ]
    add_content(prs, lines, top=1.1)
    add_page_num(prs, 4, total)

    # ========== 第三十八章 典故与启示 ==========
    s = create_slide(prs)
    add_title(prs, "第三十八章：典故与西方思想", size=28)
    lines = [
        "古代典故",
        "",
        "【典故一：尧舜禹禅让】",
        "尧舜时代以德治国，老子认为这是上德境界——无所表现的德",
        "",
        "【典故二：周公制礼】",
        "周朝建立礼乐制度，老子却警告：过分强调礼是道德衰退的表现",
        "",
        "西方思想对照",
        "",
        "法国 伏尔泰",
        "最好的政府是最少干预的政府。与道常无为相通",
        "",
        "德国 黑格尔",
        "道德是精神的最低层次，与老子失道而后礼的层次退化思想相呼应"
    ]
    add_content(prs, lines, top=1.1)
    add_page_num(prs, 5, total)

    # ========== 第三十九章 ==========
    s = create_slide(prs)
    add_title(prs, "第三十九章：上德不德", size=28)
    lines = [
        "【原文】",
        "昔之得一者：天得一以清；地得一以宁；",
        "神得一以灵；谷得一以盈；万物得一以生；",
        "侯王得一以为天下贞。",
        "",
        "其致之：天无以清将恐裂；地无以宁将恐废；",
        "神无以灵将恐歇；谷无以盈将恐竭；",
        "万物无以生将恐灭；侯王无以贵高将恐蹶。",
        "",
        "【要点】",
        "① 一即道：万事万物都得道才能存在",
        "② 天清地宁：自然得道则和谐",
        "③ 侯王之贞：统治者得道才能治理天下",
        "④ 贵高必蹶：地位越高越要守道，否则危险"
    ]
    add_content(prs, lines, top=1.1)
    add_page_num(prs, 6, total)

    # ========== 第三十九章 现代启示 ==========
    s = create_slide(prs)
    add_title(prs, "第三十九章：管理启示", size=28)
    lines = [
        "核心智慧",
        "",
        "【企业经营管理】",
        "  核心价值观（一）是企业灵魂，没有则难以持久",
        "  越是成功的企业，越要回归本真，不能忘乎所以",
        "  保持警觉，冬天总会到来",
        "",
        "【国家治理】",
        "  老子强调侯王要守道，治国者更需道德",
        "  得道多助，失道寡助",
        "  与儒家仁政思想互补",
        "",
        "【个人修为】",
        "  成功时保持清醒（天清地宁）",
        "  得意时不忘形（神得一以灵）",
        "  越成功越要谨慎（贵高将恐蹶）"
    ]
    add_content(prs, lines)
    add_page_num(prs, 7, total)

    # ========== 西方思想对照 ==========
    s = create_slide(prs)
    add_title(prs, "三章综合：西方思想名言", size=28)
    lines = [
        "法国 伏尔泰",
        "最好的政府是最少干预的政府。与道常无为相通",
        "",
        "德国 康德",
        "自由不是你想做什么就做什么，而是你应该做什么就做什么。",
        "",
        "美国 爱默生",
        "宇宙万物都有内在的秩序，顺其自然才能和谐。",
        "",
        "英国 斯宾塞",
        "适者生存，而非强者的胜利。与道的柔弱胜刚强相通"
    ]
    add_content(prs, lines)
    add_page_num(prs, 8, total)

    # ========== 生活中的智慧 ==========
    s = create_slide(prs)
    add_title(prs, "三章精华：生活中的智慧", size=28)
    lines = [
        "内心修养",
        "每天留出时间静心，不过度焦虑结果。无欲以静，天下自定",
        "",
        "工作态度",
        "做好本分，不过度表现德行。上德不德，真正有德的人不刻意显示",
        "",
        "人际关系",
        "保持内心平静，不强求他人。道常无为，万物自化",
        "",
        "人生智慧",
        "越是成功时越要谨慎。侯王无以贵高将恐蹶"
    ]
    add_content(prs, lines)
    add_page_num(prs, 9, total)

    # ========== 总结页 ==========
    s = create_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(30, 58, 95)
    bg.line.fill.background()
    add_title(prs, "三章精华提炼", top=1.5, size=36, color=RGBColor(255,255,255))
    lines = [
        "第三十七章：道常无为——顺应自然，无为而治",
        "第三十八章：始制有名——德、仁、义、礼，回归本真",
        "第三十九章：上德不德——得道者不显德，真正有德",
        "",
        "一句话总结：",
        "真正的智慧在于顺其自然、不刻意表现，",
        "越是有所成就，越要回归道的本真，保持谦卑与警觉。"
    ]
    box = s.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(200, 220, 255)
        p.space_after = Pt(8)
    add_page_num(prs, 10, total)

    # 保存
    output_dir = os.path.expanduser("~/.openclaw/workspace/courses")
    os.makedirs(output_dir, exist_ok=True)
    output_path = f"{output_dir}/道德经_第12天_37至39章.pptx"
    prs.save(output_path)
    print(f"课件已保存：{output_path}")
    return output_path

if __name__ == "__main__":
    create_daodejing_ppt()
