# -*- coding: utf-8 -*-
"""
道德经第12天 课件生成器
章节：第37-39章
风格：学术白
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# 配色方案 - 学术白
PRIMARY   = RGBColor(0, 51, 102)    # 深蓝
SECONDARY = RGBColor(102, 102, 102) # 灰
ACCENT    = RGBColor(204, 0, 0)     # 深红
TEXT      = RGBColor(51, 51, 51)    # 深灰
BG        = RGBColor(255, 255, 255) # 白色
DARK_BG   = RGBColor(248, 250, 253) # 浅灰背景

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)

# ─── 工具函数 ───────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    p.alignment = align
    return tb

def add_title_bar(slide, title):
    add_rect(slide, 0, 0, Inches(10), Inches(1.1), PRIMARY)
    add_text(slide, title, Inches(0.4), Inches(0.15), Inches(9), Inches(0.8),
             size=26, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.LEFT)

def add_chapter_header(slide, chapter_num, chapter_name, quote):
    """章节标题页"""
    set_bg(slide, PRIMARY)
    # 章节编号
    add_text(slide, f"第 {chapter_num} 章", Inches(0.5), Inches(1.5),
             Inches(9), Inches(1), size=20, bold=False,
             color=RGBColor(180,210,240), align=PP_ALIGN.CENTER)
    # 章节名
    add_text(slide, chapter_name, Inches(0.5), Inches(2.5),
             Inches(9), Inches(1.2), size=40, bold=True,
             color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
    # 引文
    add_rect(slide, Inches(1), Inches(4.2), Inches(8), Inches(1.2), RGBColor(30,80,140))
    add_text(slide, f"「{quote}」", Inches(1.2), Inches(4.35), Inches(7.6), Inches(0.9),
             size=18, italic=True, color=RGBColor(220,235,255), align=PP_ALIGN.CENTER)

def content_page(slide, title, sections):
    """通用内容页"""
    set_bg(slide, BG)
    add_title_bar(slide, title)
    y = Inches(1.3)
    for sec_title, sec_body in sections:
        # 小标题
        add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.35), ACCENT)
        add_text(slide, sec_title, Inches(0.6), y - Pt(2), Inches(3), Inches(0.45),
                 size=15, bold=True, color=PRIMARY)
        y += Inches(0.42)
        # 正文
        tb = slide.shapes.add_textbox(Inches(0.6), y, Inches(8.8), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(sec_body.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"    {line.strip()}"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT
            p.space_after = Pt(4)
        y += Inches(0.9)

def two_col_page(slide, title, left_title, left_items, right_title, right_items):
    """双栏内容页"""
    set_bg(slide, BG)
    add_title_bar(slide, title)
    # 左栏标题
    add_rect(slide, Inches(0.3), Inches(1.3), Inches(4.5), Inches(0.45), PRIMARY)
    add_text(slide, left_title, Inches(0.4), Inches(1.32), Inches(4.3), Inches(0.42),
             size=14, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
    # 左栏内容
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(1.85), Inches(4.4), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)
    # 右栏标题
    add_rect(slide, Inches(5.2), Inches(1.3), Inches(4.5), Inches(0.45), ACCENT)
    add_text(slide, right_title, Inches(5.3), Inches(1.32), Inches(4.3), Inches(0.42),
             size=14, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
    # 右栏内容
    tb2 = slide.shapes.add_textbox(Inches(5.25), Inches(1.85), Inches(4.4), Inches(5.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)

def summary_page(slide, title, points, conclusion):
    """总结页"""
    set_bg(slide, BG)
    add_title_bar(slide, title)
    y = Inches(1.4)
    for pt in points:
        # 绿色勾选
        add_rect(slide, Inches(0.5), y + Inches(0.08), Inches(0.22), Inches(0.22), RGBColor(0,140,60))
        add_text(slide, "✓", Inches(0.5), y + Pt(1), Inches(0.22), Inches(0.3),
                 size=11, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
        add_text(slide, pt, Inches(0.85), y, Inches(8.5), Inches(0.45),
                 size=15, color=TEXT)
        y += Inches(0.55)
    # 结论框
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.4), Inches(5.8), Inches(9.2), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY
    add_text(slide, f"💡  核心结论", Inches(0.6), Inches(5.9), Inches(8.8), Inches(0.4),
             size=13, bold=True, color=RGBColor(180,210,240))
    add_text(slide, conclusion, Inches(0.6), Inches(6.25), Inches(8.8), Inches(0.7),
             size=15, color=RGBColor(255,255,255))

# ─── 封面 ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PRIMARY)
# 装饰条
add_rect(slide, 0, Inches(2.8), Inches(10), Inches(0.06), RGBColor(180,210,240))
add_rect(slide, 0, Inches(4.8), Inches(10), Inches(0.06), RGBColor(180,210,240))
add_text(slide, "《道德经》每日一课", Inches(0.5), Inches(2.0), Inches(9), Inches(0.8),
         size=18, color=RGBColor(180,210,240), align=PP_ALIGN.CENTER)
add_text(slide, "第 12 天", Inches(0.5), Inches(3.0), Inches(9), Inches(1.5),
         size=52, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
add_text(slide, "第37章 · 第38章 · 第39章", Inches(0.5), Inches(4.2), Inches(9), Inches(0.6),
         size=20, color=RGBColor(180,210,240), align=PP_ALIGN.CENTER)
add_text(slide, "道常无为 ｜ 始制有名 ｜ 上德不德", Inches(0.5), Inches(4.9), Inches(9), Inches(0.6),
         size=15, color=RGBColor(150,190,230), align=PP_ALIGN.CENTER)
add_text(slide, "制作：小燕子 🐦", Inches(0.5), Inches(6.5), Inches(9), Inches(0.5),
         size=12, color=RGBColor(120,170,210), align=PP_ALIGN.CENTER)

# ─── 第37章 ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chapter_header(slide, "37", "道常无为",
                   "道常无为而无不为。侯王若能守，万物将自化。")

content_page(slide, "第37章 · 原文与解读", [
    ("【原文】", "道常无为而无不为。侯王若能守之，万物将自化。\n化而欲作，吾将镇之以无名之朴。无名之朴，夫亦将无欲。\n不欲以静，天下将自定。"),
    ("【白话解读】", "大道始终顺应自然，不妄加干涉，却成就了一切。\n领导者若能守住这个原则，万物便会自然生化。\n当人心萌生贪欲时，便以道的质朴来镇伏，人若无私无欲，\n心灵归于宁静，天下自然安定。"),
    ("【核心思想】", "全章核心：一个「守」字。\n不是消极不作为，而是以不争之德，守住道的本体，\n让万物各归其位，自然而然地运转。"),
])

# ─── 第37章 · 典故 ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG)
add_title_bar(slide, "第37章 · 典故案例")
# 左侧案例
add_rect(slide, Inches(0.3), Inches(1.3), Inches(4.5), Inches(0.45), PRIMARY)
add_text(slide, "📜 案例一：汉初休养生息", Inches(0.4), Inches(1.33), Inches(4.3), Inches(0.4),
         size=13, bold=True, color=RGBColor(255,255,255))
tb = slide.shapes.add_textbox(Inches(0.35), Inches(1.85), Inches(4.4), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
lines = ["汉高祖刘邦建立汉朝后，面对战乱创伤，",
         "实行「萧规曹随」——不频繁更改政策，",
         "减轻赋税，让百姓休养生息。",
         "结果：文景之治，天下大定。",
         "这正是「道常无为而无不为」的实践。"]
for i, l in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"• {l}"; p.font.size = Pt(13); p.font.color.rgb = TEXT; p.space_after = Pt(8)

# 右侧案例
add_rect(slide, Inches(5.2), Inches(1.3), Inches(4.5), Inches(0.45), ACCENT)
add_text(slide, "📜 案例二：刘备「勿以善小」", Inches(5.3), Inches(1.33), Inches(4.3), Inches(0.4),
         size=13, bold=True, color=RGBColor(255,255,255))
tb2 = slide.shapes.add_textbox(Inches(5.25), Inches(1.85), Inches(4.4), Inches(2.5))
tf2 = tb2.text_frame; tf2.word_wrap = True
lines2 = ["刘备临终遗言：「勿以善小而不为，",
          "勿以恶小而为之。」",
          "他不主张用严厉刑法治理蜀国，",
          "而是以德感化官员与百姓。",
          "体现「天下将自定」的政治智慧。"]
for i, l in enumerate(lines2):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = f"• {l}"; p.font.size = Pt(13); p.font.color.rgb = TEXT; p.space_after = Pt(8)

# 管理启示框
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.3), Inches(5.0), Inches(9.4), Inches(2.1))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(240,248,255)
add_text(slide, "🏢 管理启示", Inches(0.5), Inches(5.1), Inches(3), Inches(0.4),
         size=14, bold=True, color=PRIMARY)
items = ["好的管理不是事必躬亲，而是设定正确的方向和原则",
         "「少即是多」——减少不必要的干预，组织自有活力",
         "领导者最大的智慧是「守」：守住核心价值观"]
y = Inches(5.55)
for item in items:
    add_text(slide, f"✓ {item}", Inches(0.5), y, Inches(9), Inches(0.42),
             size=13, color=TEXT)
    y += Inches(0.48)

# ─── 第38章 ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chapter_header(slide, "38", "上德不德",
                   "上德不德，是以有德；下德不失德，是以无德。")

content_page(slide, "第38章 · 原文与解读", [
    ("【原文】", "上德不德，是以有德；下德不失德，是以无德。\n上德无为而无以为；下德为之而有以为。\n上仁为之而无以为；上义为之而有以为。\n上礼为之而莫之应，则攘臂而扔之。"),
    ("【白话解读】", "真正有德之人，不刻意表现德行，所以真正有德；\n刻意不失德之人，往往失去了真正的德。\n上德者顺其自然，不带私心；下德者刻意作为，动机不纯。\n仁、义、礼三者，层层递减，至礼已是强拉硬拽。"),
    ("【层次图解】", "上德（不德）→ 上仁（无以为）→ 上义（有以为）→ 上礼（扔之）\n"
                     "    ↓有德     ↓    ↓     ↓无德\n"
                     "德行递减，人为递增，老子以此揭示失道而后德的历程。"),
])

# ─── 第38章 · 精华 ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
two_col_page(slide, "第38章 · 精华提炼",
    "🟦 上德 vs 下德",
    [
        "上德不德：真正有德，不自显",
        "顺应本性，无心而为",
        "「无不为」—— 无所不成",
        "德的最高境界：自然流露",
        "如阳光普照，不求回报",
        "代表：得道圣人",
    ],
    "🟥 下德不失德",
    [
        "刻意不失德：虚伪的道德",
        "人为造作，有心而为",
        "「为之而有以为」—— 有所图",
        "德的假象：形式大于实质",
        "如作秀表演，期待掌声",
        "代表：伪善之人",
    ]
)

# ─── 第39章 ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chapter_header(slide, "39", "上德不德（续）",
                   "昔之得一者：天得一以清，地得一以宁，神得一以灵。")

content_page(slide, "第39章 · 原文与解读", [
    ("【原文】", "昔之得一者：天得一以清，地得一以宁，\n神得一以灵，谷得一以盈，万物得一以生，\n侯王得一以为天下贞。"),
    ("【白话解读】", "自古以来得到「一」（即道）的事物：\n天得道而清明，地得道而安宁，\n人心得道而灵慧，河谷得道而充盈，\n万物得道而生长，领导者得道而为天下楷模。"),
    ("【反面论证】", "天无以清将恐裂，地无以宁将恐废，\n神无以灵将恐歇，谷无以盈将恐竭，\n万物无以生将恐灭，侯王无以贵将恐蹶。\n"
                     "——失去道，则一切崩溃，此为「贵高必以贱为本」之理。"),
])

# ─── 第39章 · 案例 ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
two_col_page(slide, "第39章 · 典故与启示",
    "📜 历史案例",
    [
        "秦朝统一六国后，秦始皇严刑峻法，",
        "焚书坑儒，以暴力治天下，",
        "结果：二世而亡，仅存15年。",
        "",
        "汉初刘邦以道家黄老之术治国，",
        "轻徭薄赋，与民休息，",
        "结果：开创四百年汉朝基业。",
        "",
        "对比：霸道 vs 王道",
    ],
    "🏢 现代管理启示",
    [
        "「贵以贱为本」—— 顾客是根本，",
        "     忽视用户需求的组织必倒。",
        "「高以下为基」—— 基层是基础，",
        "     不尊重员工的企业难长久。",
        "",
        "亚马逊的「Day 1」哲学：",
        "永远保持创业初期的谦逊与灵活，",
        "不因成功而傲慢，不因规模而僵化。",
        "这正是「得一」精神的现代诠释。",
    ]
)

# ─── 三章综合 ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PRIMARY)
add_text(slide, "第12天 · 三章精华总览", Inches(0.5), Inches(0.25), Inches(9), Inches(0.7),
         size=24, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)

boxes = [
    ("第37章", "道常无为", "无为而治\n顺势而为\n守朴去欲", RGBColor(30, 80, 140)),
    ("第38章", "上德不德", "真德不显\n伪德刻意\n层次分明", RGBColor(50, 100, 160)),
    ("第39章", "得一者成", "天清地宁\n贵以贱本\n高以下基", RGBColor(70, 120, 180)),
]
for i, (ch, name, pts, col) in enumerate(boxes):
    x = Inches(0.35 + i * 3.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.2), Inches(2.9), Inches(5.5))
    box.fill.solid(); box.fill.fore_color.rgb = col
    add_text(slide, ch, Inches(x.inches + 0.1), Inches(1.4), Inches(2.7), Inches(0.5),
             size=15, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
    add_text(slide, name, Inches(x.inches + 0.1), Inches(1.95), Inches(2.7), Inches(0.65),
             size=18, bold=True, color=RGBColor(220,235,255), align=PP_ALIGN.CENTER)
    add_rect(slide, x, Inches(2.7), Inches(2.9), Inches(0.05), RGBColor(180,210,240))
    y = Inches(2.9)
    for pt in pts.split('\n'):
        add_text(slide, f"✦ {pt}", Inches(x.inches + 0.15), y, Inches(2.6), Inches(0.55),
                 size=14, color=RGBColor(220,235,255), align=PP_ALIGN.CENTER)
        y += Inches(0.7)

# ─── 总结页 ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
summary_page(slide, "第12天 · 核心收获", [
    "「道常无为」—— 最高明的领导，是让被领导者感觉不到被管理",
    "「上德不德」—— 真正的德行是自然的，不是表演的",
    "「贵以贱为本」—— 越高贵，越要以卑微为根基",
    "「天下将自定」—— 给予事物空间和时间，它会自我调适",
], "无为不是躺平，而是「为」到极致之后的从容放下。以不争，成就无可争。")

# ─── 保存 ───────────────────────────────────────────────
output = "/Users/mac/.openclaw/workspace/courses/道德经_第12天.pptx"
prs.save(output)
print(f"✅ 已保存：{output}")
