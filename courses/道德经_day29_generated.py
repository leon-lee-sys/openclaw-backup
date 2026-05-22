# -*- coding: utf-8 -*-
"""
道德经第29天 课件生成器
章节：第88-90章
风格：商务蓝/深色系，与第24天模板一致
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# 配色方案 - 商务蓝/深色系
PRIMARY = RGBColor(30, 60, 114)       # 深蓝
SECONDARY = RGBColor(70, 130, 180)    # 中蓝
ACCENT = RGBColor(204, 120, 50)       # 橙色强调
TEXT = RGBColor(51, 51, 51)           # 深灰文字
BG = RGBColor(245, 247, 250)          # 浅灰蓝背景
WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(20, 35, 60)        # 深蓝背景

# 字体配置（比标准模板大2号）
TITLE_SIZE = Pt(38)
HEADER_SIZE = Pt(30)
YUANWEN_SIZE = Pt(22)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(16)

def set_dark_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = HEADER_SIZE
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def add_yuanwen_section(slide, lines, y=1.5, height=3.2):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = SECONDARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line if line else "　"
        pp.font.size = YUANWEN_SIZE
        pp.font.color.rgb = WHITE
        pp.space_before = Pt(3)

def add_commentary_section(slide, title, content, y):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    box.line.color.rgb = SECONDARY
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"【{title}】"
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = BODY_SIZE
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=SECONDARY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(15)
    p2.font.color.rgb = WHITE

def add_points_list(slide, points, y_start=3.2):
    y = y_start
    for icon_title, explanation in points:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.25))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = icon_title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = explanation
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)
        y += 1.35
    return y

TOTAL = 14

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.6))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.9))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典29章 · 第29课"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = WHITE
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.6))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第88-90章"
p2.font.size = Pt(58)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.6), Inches(8), Inches(0.9))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "小国寡民 · 至治之世 · 文化自信"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = DARK_BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "使有什伯之器而不用，使民重死而不远徙。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = WHITE
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第八十章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = RGBColor(180, 180, 180)
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, TOTAL)

# ===== 第2页：三章概览 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s2, "今日课程")
chapters = [
    ("第八十八章", "小国寡民", "邻国相望，鸡犬之声相闻，民至老死不相往来。理想的治国境界是让百姓安居乐业。"),
    ("第八十九章", "至治之世", "天下皆谓我道大，似不肖。夫唯大，故似不肖。若肖，久矣其细。道的广大无法形容。"),
    ("第九十章", "文化自信", "知足者富，强行者有志。不失其所者久，死而不亡者寿。道法自然，知足常乐。")
]
y2 = 1.5
for num, title, desc in chapters:
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y2), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf2c = box.text_frame
    tf2c.word_wrap = True
    p2c1 = tf2c.paragraphs[0]
    p2c1.text = f"📖 {num}：{title}"
    p2c1.font.size = Pt(20)
    p2c1.font.bold = True
    p2c1.font.color.rgb = PRIMARY
    p2c2 = tf2c.add_paragraph()
    p2c2.text = desc
    p2c2.font.size = Pt(16)
    p2c2.font.color.rgb = TEXT
    p2c2.space_before = Pt(5)
    y2 += 1.65
add_slide_number(s2, 2, TOTAL)

# ===== 第3页：第八十八章 原文 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s3, "第八十八章 · 小国寡民")
add_yuanwen_section(s3, [
    "小国寡民，使有什伯之器而不用，",
    "使民重死而不远徙。",
    "",
    "虽有舟舆，无所乘之；",
    "虽有甲兵，无所陈之。",
    "",
    "使人复结绳而用之。",
    "",
    "甘其食，美其服，安其居，乐其俗。",
    "",
    "邻国相望，鸡犬之声相闻，",
    "民至老死不相往来。"
], y=1.5, height=4.5)
add_commentary_section(s3, "白话解读",
    "国家小，百姓少。即使有各种器具也不使用，让人民重视死亡而不向远方迁徙。虽然有车船，没有地方要乘坐；虽然有兵器，没有地方要陈列。让人重新用结绳记事的办法。人民吃得很香甜，穿得很美观，住得很安适，生活得很快乐。相邻的国家互相望得见，鸡鸣狗叫的声音互相听得见，而人民直到老死也不相往来。",
    6.1)
add_slide_number(s3, 3, TOTAL)

# ===== 第4页：第八十八章 逐句释义 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s4, "第八十八章 · 逐句释义")
add_points_list(s4, [
    ("① 小国寡民，使有什伯之器而不用",
     "国家小、百姓少，各种器具都用不上——返璞归真，不过度依赖工具，保持简朴的生活方式。"),
    ("② 使民重死而不远徙",
     "让人民重视死亡、不向远方迁徙——安居乐业，珍惜乡土，不为名利奔波。"),
    ("③ 虽有舟舆，无所乘之；虽有甲兵，无所陈之",
     "虽有车船兵器，却没有使用的需求——天下太平，不需要运输和战争，社会达到理想状态。"),
    ("④ 邻国相望，鸡犬之声相闻，民至老死不相往来",
     "相邻国家鸡犬相闻，百姓到老死也不相往来——不是隔绝，而是满足现状，不互相干扰。")
], y_start=1.5)
add_slide_number(s4, 4, TOTAL)

# ===== 第5页：第八十八章 本章小结 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s5, "第八十八章 · 本章小结")
add_core_box(s5, "🌟 核心观点", "小国寡民——不是要国家贫弱，而是追求一种简朴、安宁、自足的理想社会状态，让人民安居乐业。")
add_points_list(s5, [
    ("1. 简朴生活", "使有什伯之器而不用——不过度依赖工具，保持内心的清净与简朴。"),
    ("2. 安土重迁", "使民重死而不远徙——珍惜当下，不为名利奔波，安于本土生活。"),
    ("3. 甘美安乐", "甘其食，美其服，安其居，乐其俗——对现有生活满足、满意，心态平和。"),
    ("4. 不相往来", "民至老死不相往来——不是封闭隔绝，而是不互相争夺、和睦共处。")
], y_start=3.1)
add_slide_number(s5, 5, TOTAL)

# ===== 第6页：第八十九章 原文 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s6, "第八十九章 · 至治之世")
add_yuanwen_section(s6, [
    "天下皆谓我道大，似不肖。",
    "",
    "夫唯大，故似不肖。",
    "",
    "若肖，久矣其细。",
    "",
    "夫唯大，故似不肖。",
    "若肖，久矣其细也。"
], y=1.5, height=3.5)
add_commentary_section(s6, "白话解读",
    "天下人都对我说：道太大了，好像什么都不像。正因为道太大，所以好像什么都不像。如果像什么具体的东西，它早就变得渺小了。道之所以大，正因为它不像任何具体的事物，无形无象，无所不包。",
    5.1)
add_slide_number(s6, 6, TOTAL)

# ===== 第7页：第八十九章 逐句释义 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s7, "第八十九章 · 逐句释义")
add_points_list(s7, [
    ("① 天下皆谓我道大，似不肖",
     "天下人都认为道太大了，好像什么都不像——道的广大超越常人认知，难以用具体事物来形容。"),
    ("② 夫唯大，故似不肖",
     "正因为道太大，所以好像什么都不像——道的伟大在于它的无边无际、无形无象。"),
    ("③ 若肖，久矣其细",
     "如果道像某个具体的东西，它早就变得渺小了——凡是可以形容的东西都是有限的，道是不可形容的。"),
    ("④ 道的超越性",
     "道超越一切具体形态，无法用言语形容——越是描述道，就越偏离道的本质。")
], y_start=1.5)
add_slide_number(s7, 7, TOTAL)

# ===== 第8页：第八十九章 本章小结 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s8, "第八十九章 · 本章小结")
add_core_box(s8, "🌟 核心观点", "道大难肖——道之所以伟大，正因为它不像任何具体事物，无形无象，无所不包。真正的伟大是超越形象的。")
add_points_list(s8, [
    ("1. 道大似不肖", "道太大以至于不像任何东西——道的广大超越常人的想象和认知。"),
    ("2. 大故似不肖", "正因为大，所以不像——伟大之物无法用具体形象来描述。"),
    ("3. 若肖久细", "如果像具体事物，就变得渺小了——凡可形容者皆有限，道无限。"),
    ("4. 无形之形", "道无形无象，却无所不包——真正的伟大超越一切形象和概念。")
], y_start=3.1)
add_slide_number(s8, 8, TOTAL)

# ===== 第9页：第九十章 原文 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s9, "第九十章 · 文化自信")
add_yuanwen_section(s9, [
    "知足者富，强行者有志。",
    "",
    "不失其所者久，死而不亡者寿。",
    "",
    "不失其所者久，",
    "死而不亡者寿。",
    "",
    "知足者富，强行者有志。",
    "不失其所者久，死而不亡者寿。"
], y=1.5, height=3.5)
add_commentary_section(s9, "白话解读",
    "知道满足的人是富有的，努力行事的人是有志向的。不丧失根本的人就能长久，人死了但精神不朽才是真正的长寿。知足的人即使物质不多也感到富有，努力进取的人有明确的目标和意志。不失去根本的人能持久一生，死后精神仍然存在的才是真正的长寿。",
    5.1)
add_slide_number(s9, 9, TOTAL)

# ===== 第10页：第九十章 逐句释义 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s10, "第九十章 · 逐句释义")
add_points_list(s10, [
    ("① 知足者富",
     "知道满足的人就是富有的——富不在于拥有多少，而在于是否满足。贪得无厌的人永远是贫穷的。"),
    ("② 强行者有志",
     "努力奋斗的人是有志向的——不断精进、努力前行的人，才是真正有远大志向的人。"),
    ("③ 不失其所者久",
     "不丧失根本的人能够长久——无论做什么，不失去自己的根基和原则，才能持久。"),
    ("④ 死而不亡者寿",
     "死了但精神不朽才是真正的长寿——肉体有尽，精神传承无尽，真正的长寿是精神的永存。")
], y_start=1.5)
add_slide_number(s10, 10, TOTAL)

# ===== 第11页：第九十章 本章小结 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s11, "第九十章 · 本章小结")
add_core_box(s11, "🌟 核心观点", "知足者富，强行者有志，死而不亡者寿——真正的富有是知足，真正的志向是强行，真正的长寿是精神不朽。")
add_points_list(s11, [
    ("1. 知足者富", "知道满足就是富有——不贪心、不攀比，内心充实才是真正的富。"),
    ("2. 强行者有志", "努力前行的人有志向——不断精进、持之以恒，才是真正的有志者。"),
    ("3. 不失其所", "不丧失根本才能长久——守住本心、不忘初心，才能持久发展。"),
    ("4. 死而不亡", "精神永存才是真正的长寿——肉体会消亡，但精神和思想可以传承千古。")
], y_start=3.1)
add_slide_number(s11, 11, TOTAL)

# ===== 第12页：三章总览 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s12, "三章精华 · 一以贯之")
sums = [
    ("第八十八章", "小国寡民", "邻国相望，鸡犬之声相闻，民至老死不相往来。简朴生活，安土重迁，甘美安乐。"),
    ("第八十九章", "至治之世", "天下皆谓我道大，似不肖。夫唯大，故似不肖。若肖，久矣其细——道之超越性。"),
    ("第九十章", "文化自信", "知足者富，强行者有志，不失其所者久，死而不亡者寿——精神不朽的真正长寿。")
]
y12 = 1.5
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.55))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    p12_2.space_before = Pt(5)
    y12 += 1.65
add_slide_number(s12, 12, TOTAL)

# ===== 第13页：现代应用 =====
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s13, "📌 现代应用")
apps = [
    ("🏡 生活简化", "小国寡民 → 在物质丰富的时代，主动简化生活，不过度消费，保持内心清净。"),
    ("💰 知足常乐", "知足者富 → 减少对物质的过度追求，珍惜已拥有的，感受当下的幸福。"),
    ("💪 持续精进", "强行者有志 → 在自己认定的道路上持续努力，不轻言放弃，有志向才能成功。"),
    ("🌳 守住根本", "不失其所者久 → 无论外界如何变化，守住自己的原则和底线，不随波逐流。"),
    ("📚 精神传承", "死而不亡者寿 → 通过著述、教学、创作，让精神和智慧传承下去，实现真正的永恒。")
]
add_points_list(s13, apps, y_start=1.5)
add_slide_number(s13, 13, TOTAL)

# ===== 第14页：结语 =====
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s14, "第29课 · 结语")
final_box = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(2))
final_box.fill.solid()
final_box.fill.fore_color.rgb = DARK_BG
tff = final_box.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "知足者富，强行者有志"
pf1.font.size = Pt(36)
pf1.font.bold = True
pf1.font.color.rgb = WHITE
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "不失其所者久，死而不亡者寿"
pf2.font.size = Pt(18)
pf2.font.italic = True
pf2.font.color.rgb = SECONDARY
pf2.alignment = PP_ALIGN.CENTER
pf2.space_before = Pt(10)
pf3 = tff.add_paragraph()
pf3.text = "小国寡民，鸡犬相闻而不往来"
pf3.font.size = Pt(16)
pf3.font.color.rgb = RGBColor(180, 180, 180)
pf3.alignment = PP_ALIGN.CENTER
footer_note = s14.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
tffn = footer_note.text_frame
tffn.word_wrap = True
pfn1 = tffn.paragraphs[0]
pfn1.text = "道家智慧：小国寡民不是倒退，而是回归简朴；知足者富不是消极，而是真正的智慧。"
pfn1.font.size = Pt(16)
pfn1.font.color.rgb = TEXT
pfn1.alignment = PP_ALIGN.CENTER
pfn2 = tffn.add_paragraph()
pfn2.text = "不失其所，强行有志——守住根本，持续精进，精神永存。"
pfn2.font.size = Pt(15)
pfn2.font.color.rgb = SECONDARY
pfn2.alignment = PP_ALIGN.CENTER
pfn2.space_before = Pt(8)
add_slide_number(s14, 14, TOTAL)

# 保存文件
output_path = '/Users/mac/.openclaw/workspace/courses/道德经_第29天.pptx'
prs.save(output_path)
print(f"Done: {output_path}")