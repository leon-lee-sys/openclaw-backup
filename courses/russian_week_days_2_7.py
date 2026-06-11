#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
俄语7天速成 Day 2-7 - 批量生成
"""
import os
import sys
sys.path.insert(0, '/opt/homebrew/lib/node_modules/openclaw/node_modules')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(30, 60, 114)
TEXT = RGBColor(51, 51, 51)
BLUE_LIGHT = RGBColor(70, 130, 180)
GOLD = RGBColor(180, 140, 50)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(60, 140, 90)

def make_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(200, 220, 255)

def add_footer(slide, day, theme=""):
    footer = slide.shapes.add_textbox(Inches(8.5), Inches(7.2), Inches(1.3), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "День %d/7" % day
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(160, 160, 160)
    p.alignment = PP_ALIGN.RIGHT

def add_box(slide, text, left, top, width, height, fill_color=None, font_size=16, bold=False, color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return box

def make_day2():
    prs = Presentation()
    day = 2
    theme = "На работе"
    
    # Cover
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语7天速成班"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    t2 = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = t2.text_frame
    p = tf.paragraphs[0]
    p.text = "День 2: На работе"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    t3 = s.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = t3.text_frame
    p = tf.paragraphs[0]
    p.text = "工作与职场"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER
    add_footer(s, day, theme)
    
    # Vocabulary
    s = make_slide(prs)
    set_title_bar(s, "核心词汇 · Слова", "Ключевые слова")
    words = [
        ("работа", "工作"),
        ("офис", "办公室"),
        ("начальник", "上司/老板"),
        ("коллега", "同事"),
        ("совещание", "会议"),
        ("отчёт", "报告"),
        ("дедлайн", "截止日期"),
        ("уволиться", "辞职"),
    ]
    y = 1.6
    for ru, cn in words:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(4), Inches(0.5), fill_color=BLUE_LIGHT, font_size=18, color=WHITE)
        add_box(s, cn, Inches(4.7), Inches(y), Inches(4.8), Inches(0.5), fill_color=GREEN, font_size=18, color=WHITE)
        y += 0.6
    add_footer(s, day, theme)
    
    # Phrases
    s = make_slide(prs)
    set_title_bar(s, "职场常用语 · Фразы", "На работе")
    phrases = [
        ("Где офис?", "办公室在哪里？"),
        ("Когда совещание?", "什么时候开会？"),
        ("У меня дедлайн завтра.", "我明天截止日期。"),
        ("Я отправлю отчёт.", "我会发送报告。"),
        ("Спасибо за помощь!", "谢谢帮忙！"),
        ("Мне нужно уволиться.", "我需要辞职。"),
    ]
    y = 1.6
    for ru, cn in phrases:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(5.5), Inches(0.55), fill_color=PRIMARY, font_size=17, color=WHITE)
        add_box(s, cn, Inches(6.2), Inches(y), Inches(3.3), Inches(0.55), font_size=17, color=TEXT)
        y += 0.65
    add_footer(s, day, theme)
    
    # Dialogue
    s = make_slide(prs)
    set_title_bar(s, "情景对话 · Диалог", "Ситуация")
    dlg = [
        "— Доброе утро! Как дела?",
        "— Доброе утро! Хорошо, спасибо. Когда совещание?",
        "— В два часа. Вы готовы с отчётом?",
        "— Да, я отправлю его сегодня.",
        "— Отлично! Спасибо.",
    ]
    y = 1.6
    for line in dlg:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        y += 0.5
    add_footer(s, day, theme)
    
    # Practice
    s = make_slide(prs)
    set_title_bar(s, "课后练习 · Упражнения", "Практика")
    exercises = [
        "1. 用俄语问同事：办公室在哪里？",
        "2. 翻译：我的截止日期是明天。",
        "3. 对话练习：询问会议时间",
        "4. 回答：Как дела на работе?（工作怎么样？）",
    ]
    y = 1.7
    for ex in exercises:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = ex
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT
        y += 0.55
    add_footer(s, day, theme)
    
    # Summary
    s = make_slide(prs)
    set_title_bar(s, "今日总结 · Итоги", "Что мы выучили")
    summary = ["✅ 职场核心词汇", "✅ 工作常用表达", "✅ 简单职场对话"]
    y = 1.7
    for item in summary:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GREEN
        y += 0.6
    add_box(s, "Завтра: День 3 - В ресторане (餐厅)", Inches(1), Inches(5.5), Inches(8), Inches(0.6), fill_color=PRIMARY, font_size=16, color=WHITE)
    add_footer(s, day, theme)
    
    path = "/Users/mac/.openclaw/workspace/courses/俄语课件/俄语速成_Day2.pptx"
    prs.save(path)
    print("Done: Day 2")
    return path

def make_day3():
    prs = Presentation()
    day = 3
    theme = "В ресторане"
    
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语7天速成班"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    t2 = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = t2.text_frame
    p = tf.paragraphs[0]
    p.text = "День 3: В ресторане"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    t3 = s.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = t3.text_frame
    p = tf.paragraphs[0]
    p.text = "餐厅与点餐"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "核心词汇 · Слова", "Ключевые слова")
    words = [
        ("ресторан", "餐厅"),
        ("меню", "菜单"),
        ("заказ", "点单/订单"),
        ("счёт", "账单"),
        ("вода", "水"),
        ("хлеб", "面包"),
        ("мясо", "肉"),
        ("овощи", "蔬菜"),
    ]
    y = 1.6
    for ru, cn in words:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(4), Inches(0.5), fill_color=BLUE_LIGHT, font_size=18, color=WHITE)
        add_box(s, cn, Inches(4.7), Inches(y), Inches(4.8), Inches(0.5), fill_color=GREEN, font_size=18, color=WHITE)
        y += 0.6
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "点餐常用语 · Фразы", "Как заказать")
    phrases = [
        ("Можно меню?", "请给我菜单？"),
        ("Что вы рекомендуете?", "您推荐什么？"),
        ("Я хочу заказать...", "我想点..."),
        ("Счёт, пожалуйста.", "买单。"),
        ("Это очень вкусно!", "这很好吃！"),
    ]
    y = 1.6
    for ru, cn in phrases:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(5.5), Inches(0.55), fill_color=PRIMARY, font_size=17, color=WHITE)
        add_box(s, cn, Inches(6.2), Inches(y), Inches(3.3), Inches(0.55), font_size=17, color=TEXT)
        y += 0.65
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "情景对话 · Диалог", "В ресторане")
    dlg = [
        "— Добрый вечер! Можно меню?",
        "— Пожалуйста! Что вы рекомендуете?",
        "— Сегодня мы рекомендуем мясо с овощами.",
        "— Хорошо, я возьму это. И воды, пожалуйста.",
        "— Счёт, пожалуйста. Спасибо за обед!",
    ]
    y = 1.6
    for line in dlg:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        y += 0.5
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "今日总结 · Итоги", "Что мы выучили")
    summary = ["✅ 餐厅词汇", "✅ 点餐用语", "✅ 买单表达"]
    y = 1.7
    for item in summary:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GREEN
        y += 0.6
    add_box(s, "Завтра: День 4 - Путешествие (旅行)", Inches(1), Inches(5.5), Inches(8), Inches(0.6), fill_color=PRIMARY, font_size=16, color=WHITE)
    add_footer(s, day, theme)
    
    path = "/Users/mac/.openclaw/workspace/courses/俄语课件/俄语速成_Day3.pptx"
    prs.save(path)
    print("Done: Day 3")
    return path

def make_day4():
    prs = Presentation()
    day = 4
    theme = "Путешествие"
    
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语7天速成班"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    t2 = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = t2.text_frame
    p = tf.paragraphs[0]
    p.text = "День 4: Путешествие"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    t3 = s.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = t3.text_frame
    p = tf.paragraphs[0]
    p.text = "旅行与交通"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "核心词汇 · Слова", "Ключевые слова")
    words = [
        ("аэропорт", "机场"),
        ("вокзал", "火车站"),
        ("билет", "票"),
        ("самолёт", "飞机"),
        ("поезд", "火车"),
        ("автобус", "公交车"),
        ("такси", "出租车"),
        ("гостиница", "酒店"),
    ]
    y = 1.6
    for ru, cn in words:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(4), Inches(0.5), fill_color=BLUE_LIGHT, font_size=18, color=WHITE)
        add_box(s, cn, Inches(4.7), Inches(y), Inches(4.8), Inches(0.5), fill_color=GREEN, font_size=18, color=WHITE)
        y += 0.6
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "交通常用语 · Фразы", "Как спросить")
    phrases = [
        ("Где аэропорт?", "机场在哪里？"),
        ("Один билет, пожалуйста.", "请给我一张票。"),
        ("Сколько стоит?", "多少钱？"),
        ("Когда отправление?", "什么时候出发？"),
        ("Мне нужно в гостиницу.", "我需要去酒店。"),
    ]
    y = 1.6
    for ru, cn in phrases:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(5.5), Inches(0.55), fill_color=PRIMARY, font_size=17, color=WHITE)
        add_box(s, cn, Inches(6.2), Inches(y), Inches(3.3), Inches(0.55), font_size=17, color=TEXT)
        y += 0.65
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "情景对话 · Диалог", "В аэропорту")
    dlg = [
        "— Здравствуйте! Один билет до Москвы, пожалуйста.",
        "— Когда вы хотите лететь?",
        "— Завтра утром.",
        "— Хорошо. Сколько стоит?",
        "— Двадцать тысяч рублей.",
    ]
    y = 1.6
    for line in dlg:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        y += 0.5
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "今日总结 · Итоги", "Что мы выучили")
    summary = ["✅ 交通词汇", "✅ 问路与购票", "✅ 机场对话"]
    y = 1.7
    for item in summary:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GREEN
        y += 0.6
    add_box(s, "Завтра: День 5 - Покупки (购物)", Inches(1), Inches(5.5), Inches(8), Inches(0.6), fill_color=PRIMARY, font_size=16, color=WHITE)
    add_footer(s, day, theme)
    
    path = "/Users/mac/.openclaw/workspace/courses/俄语课件/俄语速成_Day4.pptx"
    prs.save(path)
    print("Done: Day 4")
    return path

def make_day5():
    prs = Presentation()
    day = 5
    theme = "Покупки"
    
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语7天速成班"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    t2 = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = t2.text_frame
    p = tf.paragraphs[0]
    p.text = "День 5: Покупки"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    t3 = s.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = t3.text_frame
    p = tf.paragraphs[0]
    p.text = "购物与价格"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "核心词汇 · Слова", "Ключевые слова")
    words = [
        ("магазин", "商店"),
        ("цена", "价格"),
        ("скидка", "折扣"),
        ("оплата", "支付"),
        ("наличные", "现金"),
        ("карта", "卡"),
        ("размер", "尺寸"),
        ("цвет", "颜色"),
    ]
    y = 1.6
    for ru, cn in words:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(4), Inches(0.5), fill_color=BLUE_LIGHT, font_size=18, color=WHITE)
        add_box(s, cn, Inches(4.7), Inches(y), Inches(4.8), Inches(0.5), fill_color=GREEN, font_size=18, color=WHITE)
        y += 0.6
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "购物常用语 · Фразы", "Как купить")
    phrases = [
        ("Сколько стоит?", "多少钱？"),
        ("Это слишком дорого.", "这太贵了。"),
        ("Можно скидку?", "可以打折吗？"),
        ("Я беру это.", "我买了。"),
        ("Можно картой?", "可以刷卡吗？"),
    ]
    y = 1.6
    for ru, cn in phrases:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(5.5), Inches(0.55), fill_color=PRIMARY, font_size=17, color=WHITE)
        add_box(s, cn, Inches(6.2), Inches(y), Inches(3.3), Inches(0.55), font_size=17, color=TEXT)
        y += 0.65
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "情景对话 · Диалог", "В магазине")
    dlg = [
        "— Здравствуйте! Сколько стоит эта рубашка?",
        "— Две тысячи рублей.",
        "— О, это слишком дорого. Можно скидку?",
        "— Хорошо, будет одна тысяча пятьсот.",
        "— Хорошо, я беру это!",
    ]
    y = 1.6
    for line in dlg:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        y += 0.5
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "今日总结 · Итоги", "Что мы выучили")
    summary = ["✅ 购物词汇", "✅ 询价与砍价", "✅ 付款表达"]
    y = 1.7
    for item in summary:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GREEN
        y += 0.6
    add_box(s, "Завтра: День 6 - Здоровье (健康)", Inches(1), Inches(5.5), Inches(8), Inches(0.6), fill_color=PRIMARY, font_size=16, color=WHITE)
    add_footer(s, day, theme)
    
    path = "/Users/mac/.openclaw/workspace/courses/俄语课件/俄语速成_Day5.pptx"
    prs.save(path)
    print("Done: Day 5")
    return path

def make_day6():
    prs = Presentation()
    day = 6
    theme = "Здоровье"
    
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语7天速成班"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    t2 = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = t2.text_frame
    p = tf.paragraphs[0]
    p.text = "День 6: Здоровье"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    t3 = s.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = t3.text_frame
    p = tf.paragraphs[0]
    p.text = "健康与医疗"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "核心词汇 · Слова", "Ключевые слова")
    words = [
        ("больница", "医院"),
        ("врач", "医生"),
        ("аптека", "药房"),
        ("лекарство", "药"),
        ("температура", "发烧/温度"),
        ("голова", "头"),
        ("горло", "喉咙"),
        ("насморк", "流鼻涕"),
    ]
    y = 1.6
    for ru, cn in words:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(4), Inches(0.5), fill_color=BLUE_LIGHT, font_size=18, color=WHITE)
        add_box(s, cn, Inches(4.7), Inches(y), Inches(4.8), Inches(0.5), fill_color=GREEN, font_size=18, color=WHITE)
        y += 0.6
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "就医常用语 · Фразы", "Как сказать")
    phrases = [
        ("Где больница?", "医院在哪里？"),
        ("У меня болит голова.", "我头疼。"),
        ("Мне нужны лекарства.", "我需要药。"),
        ("Какая температура?", "多少度？"),
        ("Я себя плохо чувствую.", "我感觉不舒服。"),
    ]
    y = 1.6
    for ru, cn in phrases:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(5.5), Inches(0.55), fill_color=PRIMARY, font_size=17, color=WHITE)
        add_box(s, cn, Inches(6.2), Inches(y), Inches(3.3), Inches(0.55), font_size=17, color=TEXT)
        y += 0.65
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "情景对话 · Диалог", "У врача")
    dlg = [
        "— Здравствуйте! Что у вас болит?",
        "— У меня болит голова и горло.",
        "— Давайте измерим температуру.",
        "— Хорошо. Какая температура?",
        "— Тридцать семь и пять. Это простуда.",
    ]
    y = 1.6
    for line in dlg:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        y += 0.5
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "今日总结 · Итоги", "Что мы выучили")
    summary = ["✅ 医疗词汇", "✅ 描述症状", "✅ 药房与就医"]
    y = 1.7
    for item in summary:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GREEN
        y += 0.6
    add_box(s, "Завтра: День 7 - Культура (文化)", Inches(1), Inches(5.5), Inches(8), Inches(0.6), fill_color=PRIMARY, font_size=16, color=WHITE)
    add_footer(s, day, theme)
    
    path = "/Users/mac/.openclaw/workspace/courses/俄语课件/俄语速成_Day6.pptx"
    prs.save(path)
    print("Done: Day 6")
    return path

def make_day7():
    prs = Presentation()
    day = 7
    theme = "Культура"
    
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语7天速成班"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    t2 = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = t2.text_frame
    p = tf.paragraphs[0]
    p.text = "День 7: Культура"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    t3 = s.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = t3.text_frame
    p = tf.paragraphs[0]
    p.text = "文化与娱乐"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "核心词汇 · Слова", "Ключевые слова")
    words = [
        ("кино", "电影"),
        ("театр", "剧院"),
        ("музей", "博物馆"),
        ("музыка", "音乐"),
        ("искусство", "艺术"),
        ("выставка", "展览"),
        ("футбол", "足球"),
        ("книга", "书"),
    ]
    y = 1.6
    for ru, cn in words:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(4), Inches(0.5), fill_color=BLUE_LIGHT, font_size=18, color=WHITE)
        add_box(s, cn, Inches(4.7), Inches(y), Inches(4.8), Inches(0.5), fill_color=GREEN, font_size=18, color=WHITE)
        y += 0.6
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "娱乐常用语 · Фразы", "Как пригласить")
    phrases = [
        ("Вы хотите пойти в кино?", "你想去看电影吗？"),
        ("Давайте пойдём в театр!", "我们去看剧吧！"),
        ("Мне нравится русская музыка.", "我喜欢俄罗斯音乐。"),
        ("Когда выставка?", "展览什么时候？"),
        ("Это очень интересно!", "这很有趣！"),
    ]
    y = 1.6
    for ru, cn in phrases:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(5.5), Inches(0.55), fill_color=PRIMARY, font_size=17, color=WHITE)
        add_box(s, cn, Inches(6.2), Inches(y), Inches(3.3), Inches(0.55), font_size=17, color=TEXT)
        y += 0.65
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "情景对话 · Диалог", "Приглашение")
    dlg = [
        "— Привет! Что ты хочешь делать сегодня?",
        "— Не знаю. А что?",
        "— Вы хотите пойти в музей?",
        "— Да, это хорошая идея!",
        "— Отлично! До встречи!",
    ]
    y = 1.6
    for line in dlg:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        y += 0.5
    add_footer(s, day, theme)
    
    s = make_slide(prs)
    set_title_bar(s, "7天总结 · Итоги недели", "Всё что мы выучили")
    summary = [
        "✅ Day 1: Приветствие — 问候",
        "✅ Day 2: На работе — 工作",
        "✅ Day 3: В ресторане — 餐厅",
        "✅ Day 4: Путешествие — 旅行",
        "✅ Day 5: Покупки — 购物",
        "✅ Day 6: Здоровье — 健康",
        "✅ Day 7: Культура — 文化",
    ]
    y = 1.5
    for item in summary:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.45))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = GREEN
        y += 0.45
    add_box(s, "🎉 恭喜！7天速成班完成！", Inches(1), Inches(6.2), Inches(8), Inches(0.6), fill_color=GOLD, font_size=18, color=WHITE)
    add_footer(s, day, theme)
    
    path = "/Users/mac/.openclaw/workspace/courses/俄语课件/俄语速成_Day7.pptx"
    prs.save(path)
    print("Done: Day 7")
    return path

if __name__ == "__main__":
    make_day2()
    make_day3()
    make_day4()
    make_day5()
    make_day6()
    make_day7()
    print("\nDay 2-7 完成！")
