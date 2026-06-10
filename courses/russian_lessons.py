#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
俄语课件生成器 - 第一月语法基础
第1-7天：名词第1格、第2格
"""
import os
import sys
sys.path.insert(0, '/opt/homebrew/lib/node_modules/openclaw/node_modules')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(30, 60, 114)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)
GOLD = RGBColor(180, 140, 50)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(40, 40, 40)

def make_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def set_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(200, 220, 255)

def add_footer(slide, day_num, lesson_num, total=12):
    footer = slide.shapes.add_textbox(Inches(8.5), Inches(7.2), Inches(1.3), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Day%d-L%d/%d" % (day_num, lesson_num, total)
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(160, 160, 160)
    p.alignment = PP_ALIGN.RIGHT

def add_box(slide, text, left, top, width, height, fill_color=None, font_size=18, bold=False, color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()
    box.line.fill.background()
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return box

def make_lesson(day_num, lesson_num, title_cn, title_ru, content_blocks, output_dir):
    prs = Presentation()
    
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    main = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = main.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语基础语法 · 第%d天" % day_num
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    sub = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    tf2 = sub.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title_cn
    p2.font.size = Pt(36)
    p2.font.color.rgb = TEXT
    p2.alignment = PP_ALIGN.CENTER
    
    ru_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5), Inches(6), Inches(0.8))
    ru_box.fill.solid()
    ru_box.fill.fore_color.rgb = BLUE_LIGHT
    ru_box.line.fill.background()
    tf = ru_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_ru
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    add_footer(s, day_num, lesson_num)
    
    for idx, block in enumerate(content_blocks[:10]):
        s = make_slide(prs)
        set_title_bar(s, block["title_cn"], block.get("title_ru", ""))
        
        y = 1.6
        for item in block["content"]:
            if isinstance(item, tuple):
                if item[0] == "header":
                    h = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
                    tf = h.text_frame
                    p = tf.paragraphs[0]
                    p.text = item[1]
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = PRIMARY
                    y += 0.4
                elif item[0] == "example":
                    box = add_box(s, item[1], Inches(0.5), Inches(y), Inches(9), Inches(0.6), 
                                 fill_color=BLUE_LIGHT, font_size=16, color=WHITE)
                    y += 0.75
                elif item[0] == "note":
                    box = add_box(s, item[1], Inches(0.5), Inches(y), Inches(9), Inches(0.5),
                                 fill_color=GOLD, font_size=14, color=WHITE)
                    y += 0.65
            else:
                t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
                tf = t.text_frame
                p = tf.paragraphs[0]
                p.text = item
                p.font.size = Pt(15)
                p.font.color.rgb = TEXT
                y += 0.45
        
        add_footer(s, day_num, lesson_num)
    
    s = make_slide(prs)
    set_title_bar(s, "课后练习", "Uprazhneniya")
    
    practice = [
        "1. 将下列名词变为第2格：",
        "   dom, kniga, okno, drug, student",
        "",
        "2. 用所给名词的适当形式填空：",
        "   Ya student > Ya iz ___ (Moskva)",
        "   Eto kniga > Eto kniga ___ (moy drug)",
        "",
        "3. 中译俄：",
        "   这是我朋友的书。",
        "   我来自上海。",
    ]
    
    y = 1.7
    for line in practice:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        y += 0.38
    
    add_footer(s, day_num, lesson_num)
    
    filename = "俄语_Day%d_L%d.pptx" % (day_num, lesson_num)
    path = os.path.join(output_dir, filename)
    prs.save(path)
    print("Done: %s" % filename)
    return path

def gen_day1(output_dir):
    lesson_num = 1
    title_cn = "名词的第1格和第2格"
    title_ru = "Imenitelnyy & Roditelny padesh"
    
    content = [
        {"title_cn": "一、什么是第1格？", "title_ru": "Imenitelny padezh",
         "content": [
             ("header", '第1格 = 主格，表示"谁/什么"'),
             "第1格是名词的原形，是所有名词的基础形式。",
             ("example", "Student (学生) - etot student uchitsya."),
             ("example", "Kniga (书) - eta kniga ochen interesnaya."),
             ("note", "记忆方法：第1格就是词典里单词的原形！"),
         ]},
        {"title_cn": "二、第2格的基本用法", "title_ru": "Roditelny padezh",
         "content": [
             ("header", '第2格主要表示"所属"关系，相当于中文的"的"：'),
             "Eto kniga studenta. = 这是学生的书。",
             "Eto dom brata. = 这是弟弟的房子。",
             ("header", "第2格还表示'来自'："),
             "Ya iz Kitaya. = 我来自中国。",
             "Ona iz Moskvy. = 她来自莫斯科。",
             ("note", "关键语法：iz + 第2格 = 来自某地"),
         ]},
        {"title_cn": "三、名词第2格单数变化规则", "title_ru": "Sklonenie sushchestvitelnykh",
         "content": [
             ("header", "第一变格法（阴性名词）："),
             "以 -a / -ya 结尾的阴性名词：",
             ("example", "kniga > knigi (书)"),
             ("example", "devushka > devushki (姑娘)"),
             ("header", "第二变格法（阳性/中性名词）："),
             "以硬辅音结尾的阳性名词：",
             ("example", "dom > doma (房子)"),
             ("example", "student > studenta (学生)"),
         ]},
        {"title_cn": "四、名词第2格变化示例", "title_ru": "Primery skloneniya",
         "content": [
             ("header", "阳性名词："),
             "student > studenta (学生的)",
             "drug > друга (朋友的)",
             "vrach > vracha (医生的)",
             ("header", "阴性名词："),
             "kniga > knigi (书的)",
             "devushka > devushki (姑娘的)",
             ("note", "规律：通常在词尾加 -a 或 -i"),
         ]},
        {"title_cn": "五、易错点：特殊结尾", "title_ru": "Osobye sluchai",
         "content": [
             ("header", "以 -t 结尾的动词名词："),
             "chitat > chteniya (阅读的) - 特殊变化",
             ("header", "以 -ch 结尾的名词："),
             "noch > nochi (夜晚的) - ch > sh",
             ("note", "注意：这类词的变化不规则，需要单独记忆！"),
         ]},
        {"title_cn": "六、数量词后的第2格", "title_ru": "Roditelny posle chiselnykh",
         "content": [
             ("header", "数字2-4后名词用单数第2格："),
             "dva studenta (两个学生)",
             "tri knigi (三本书)",
             "chetyre devushki (四个姑娘)",
             ("header", "数字5以上名词用复数第2格："),
             "pyat studentov (五个学生)",
             ("note", "关键：1之后单数，2-4单数第2格，5+复数第2格"),
         ]},
        {"title_cn": "七、第2格否定用法", "title_ru": "Otritsatelny roditelny",
         "content": [
             ("header", '否定"没有/不是"时使用第2格：'),
             "U menya net vremeni. = 我没有时间。",
             "Eto ne moya kniga. = 这不是我的书。",
             ("header", "net = 没有（现在时）"),
             "ne bylo = 没有（过去时）",
             "ne budet = 没有（将来时）",
             ("note", "结构：U + 人称 + net + 物"),
         ]},
        {"title_cn": "八、阅读理解练习", "title_ru": "Chteniye s ponimaniyem",
         "content": [
             "Menya zovut Anna. Ya iz Moskvy. Ya studentka universiteta. Moy brat tozhe student. On iz Sankt-Peterburga.",
             "",
             ("header", "请回答："),
             "1. Kak eye zovut? (她叫什么？)",
             "2. Otkuda Anna? (她来自哪里？)",
             "3. Kto eye brat? (她弟弟是做什么的？)",
             ("note", "提示：studentka = 女大学生，universitet = 大学"),
         ]},
        {"title_cn": "九、翻译练习", "title_ru": "Perevod",
         "content": [
             ("header", "中译俄练习："),
             "1. 这是我朋友的书。",
             "2. 我来自中国。",
             "3. 她没有时间。",
             "",
             ("header", "参考答案："),
             "1. Eto kniga moego druga.",
             "2. Ya iz Kitaya.",
             "3. U neyo net vremeni.",
         ]},
        {"title_cn": "十、本课重点总结", "title_ru": "Itogi uroka",
         "content": [
             ("header", "第2格核心用法："),
             "① 表示所属：chyo? > ...的...",
             "② 表示来源：iz + 第2格 = 来自...",
             "③ 数字搭配：2-4用单数，5+用复数",
             "④ 否定结构：net + 第2格",
             ("header", "今日关键词汇："),
             "dom (房子) > doma",
             "kniga (书) > knigi",
             "student (学生) > studenta",
             ("note", "记住：名词第2格回答chya?"),
         ]},
    ]
    
    return make_lesson(1, lesson_num, title_cn, title_ru, content, output_dir)

output_dir = "/Users/mac/.openclaw/workspace/courses/俄语课件"
os.makedirs(output_dir, exist_ok=True)

path = gen_day1(output_dir)
print("\n俄语Day1课件生成完成！")
print("路径:", path)
