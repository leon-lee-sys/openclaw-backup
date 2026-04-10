from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

def add_slide_number(slide, num, total):
    """添加页码"""
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

# ===== 第1页：封面 =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar.line.fill.background()
txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第一课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

main_title = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main_title.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第1-3章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER

subtitle = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = subtitle.text_frame
p3 = tf3.paragraphs[0]
p3.text = "宇宙观 · 道之本体 · 无为而治"
p3.font.size = Pt(28)
p3.font.color.rgb = RGBColor(70, 130, 180)
p3.alignment = PP_ALIGN.CENTER

quote_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_q = quote_box.text_frame
tf_q.word_wrap = True
p_q1 = tf_q.paragraphs[0]
p_q1.text = "道可道，非常道。名可名，非常名。"
p_q1.font.size = Pt(22)
p_q1.font.italic = True
p_q1.font.color.rgb = RGBColor(30, 60, 114)
p_q1.alignment = PP_ALIGN.CENTER
p_q2 = tf_q.add_paragraph()
p_q2.text = "—— 《道德经》开篇"
p_q2.font.size = Pt(14)
p_q2.font.color.rgb = RGBColor(150, 150, 150)
p_q2.alignment = PP_ALIGN.CENTER

add_slide_number(slide1, 1, 10)

# ===== 第2页：第一章详解 =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar2.fill.solid()
title_bar2.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar2.line.fill.background()
txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf5 = txBox2.text_frame
p5 = tf5.paragraphs[0]
p5.text = "第一章 · 天地之始"
p5.font.size = Pt(32)
p5.font.bold = True
p5.font.color.rgb = RGBColor(255, 255, 255)

# 原文框
yuanwen_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
yuanwen_box.fill.solid()
yuanwen_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw = yuanwen_box.text_frame
tf_yw.word_wrap = True
p_yw = tf_yw.paragraphs[0]
p_yw.text = "【原文】"
p_yw.font.size = Pt(18)
p_yw.font.bold = True
p_yw.font.color.rgb = RGBColor(30, 60, 114)
for line in ["道可道，非常道；名可名，非常名。", "无名天地之始，有名万物之母。"]:
    p = tf_yw.add_paragraph()
    p.text = line
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

# 白话解读
content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(4))
tf_c = content_box.text_frame
tf_c.word_wrap = True
p_c1 = tf_c.paragraphs[0]
p_c1.text = "【白话解读】"
p_c1.font.size = Pt(18)
p_c1.font.bold = True
p_c1.font.color.rgb = RGBColor(30, 60, 114)
解读 = [
    ("道可道，非常道", "能用语言说出来的'道'，就不是永恒不变的道了。'道'是不可言说的。"),
    ("名可名，非常名", "能用名字命名的东西，都不是永恒的名字。名字只是代号，不是本质。"),
    ("无名天地之始", "'无'是天地的本原——宇宙最初是一片虚无。"),
    ("有名万物之母", "'有'是万物的母亲——天地万物从'有'中产生。")
]
for title, content in 解读:
    p = tf_c.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(8)

add_slide_number(slide2, 2, 10)

# ===== 第3页：第一章核心思想 =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar3.fill.solid()
title_bar3.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar3.line.fill.background()
txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf6 = txBox3.text_frame
p6 = tf6.paragraphs[0]
p6.text = "第一章 · 核心思想"
p6.font.size = Pt(32)
p6.font.bold = True
p6.font.color.rgb = RGBColor(255, 255, 255)

# 核心框
core_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
core_box.fill.solid()
core_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_core = core_box.text_frame
tf_core.word_wrap = True
p_core = tf_core.paragraphs[0]
p_core.text = "🌟 核心观点"
p_core.font.size = Pt(18)
p_core.font.bold = True
p_core.font.color.rgb = RGBColor(255, 255, 255)
p_core2 = tf_core.add_paragraph()
p_core2.text = "'道'是宇宙本源，看不见摸不着却支配万物运行。顺应规律比强行干预更有效。"
p_core2.font.size = Pt(16)
p_core2.font.color.rgb = RGBColor(255, 255, 255)

# 三个要点
points = [
    ("1. 道不可言说", "真正的'道'超出语言边界，就像禅宗'不立文字'——语言只是指月的手指，不是月亮本身。"),
    ("2. 名只是代号", "我们给事物起的名字，只是代号，不是本质。如'手机'这个名不代表手机真正的存在。"),
    ("3. 有无相生", "'无'中生'有'，宇宙从虚无中诞生。就像创意从空白中产生，灵感从无意识中涌现。")
]
y = 3.1
for title, content in points:
    pt_box = slide3.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.2))
    tf_pt = pt_box.text_frame
    tf_pt.word_wrap = True
    p_pt1 = tf_pt.paragraphs[0]
    p_pt1.text = title
    p_pt1.font.size = Pt(18)
    p_pt1.font.bold = True
    p_pt1.font.color.rgb = RGBColor(30, 60, 114)
    p_pt2 = tf_pt.add_paragraph()
    p_pt2.text = content
    p_pt2.font.size = Pt(14)
    p_pt2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.3

add_slide_number(slide3, 3, 10)

# ===== 第4页：第二章详解 =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar4.fill.solid()
title_bar4.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar4.line.fill.background()
txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf7 = txBox4.text_frame
p7 = tf7.paragraphs[0]
p7.text = "第二章 · 相对之美"
p7.font.size = Pt(32)
p7.font.bold = True
p7.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
yuanwen_box2.fill.solid()
yuanwen_box2.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw2 = yuanwen_box2.text_frame
tf_yw2.word_wrap = True
p_yw2 = tf_yw2.paragraphs[0]
p_yw2.text = "【原文】"
p_yw2.font.size = Pt(18)
p_yw2.font.bold = True
p_yw2.font.color.rgb = RGBColor(30, 60, 114)
p_yw2_2 = tf_yw2.add_paragraph()
p_yw2_2.text = "天下皆知美之为美，斯恶已；皆知善之为善，斯不善已。"
p_yw2_2.font.size = Pt(24)
p_yw2_2.font.color.rgb = RGBColor(51, 51, 51)

content_box2 = slide4.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(4.5))
tf_c2 = content_box2.text_frame
tf_c2.word_wrap = True
p_c4 = tf_c2.paragraphs[0]
p_c4.text = "【白话解读】"
p_c4.font.size = Pt(18)
p_c4.font.bold = True
p_c4.font.color.rgb = RGBColor(30, 60, 114)
解读2 = [
    ("美与丑对立", "天下人都知道美之所以为美，是因为有丑的存在作为对比——没有丑就无所谓美。"),
    ("善与恶对立", "都知道善之所以为善，是因为有不善（恶）存在——没有恶就无所谓善。"),
    ("对立才知标准", "美和善的标准，是通过对比才产生的。没有对比就没有判断的基准。"),
    ("核心：相对论", "美丑、善恶都是相对的、可转化的。坏事可以变好事，好事也可能变坏事。")
]
for title, content in 解读2:
    p = tf_c2.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide4, 4, 10)

# ===== 第5页：塞翁失马典故 =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar5.fill.solid()
title_bar5.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar5.line.fill.background()
txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf8 = txBox5.text_frame
p8 = tf8.paragraphs[0]
p8.text = "📜 中国典故：塞翁失马"
p8.font.size = Pt(32)
p8.font.bold = True
p8.font.color.rgb = RGBColor(255, 255, 255)

story_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.2))
story_box.fill.solid()
story_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_story = story_box.text_frame
tf_story.word_wrap = True
p_s1 = tf_story.paragraphs[0]
p_s1.text = "【故事原文】"
p_s1.font.size = Pt(16)
p_s1.font.bold = True
p_s1.font.color.rgb = RGBColor(30, 60, 114)
p_s2 = tf_story.add_paragraph()
p_s2.text = "近塞上之人有善术者，马无故亡而入胡。人皆吊之，其父曰：'此何遽不为福乎？'居数月，其马将胡骏马而归。"
p_s2.font.size = Pt(15)
p_s2.font.color.rgb = RGBColor(51, 51, 51)
p_s2.space_before = Pt(8)
p_s3 = tf_story.add_paragraph()
p_s3.text = "人皆贺之，其父曰：'此何遽不能为祸乎？'家富良马，其子好骑，坠而折其髀。人皆吊之，其父曰：'此何遽不为福乎？'"
p_s3.font.size = Pt(15)
p_s3.font.color.rgb = RGBColor(51, 51, 51)
p_s3.space_before = Pt(6)

# 转折
change_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(9), Inches(1))
change_box.fill.solid()
change_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_ch = change_box.text_frame
tf_ch.word_wrap = True
p_ch = tf_ch.paragraphs[0]
p_ch.text = "⚡ 故事的启示"
p_ch.font.size = Pt(18)
p_ch.font.bold = True
p_ch.font.color.rgb = RGBColor(255, 255, 255)
p_ch2 = tf_ch.add_paragraph()
p_ch2.text = "\"祸兮福所倚，福兮祸所伏\" —— 坏事可能变好事，好事也可能变坏事"
p_ch2.font.size = Pt(16)
p_ch2.font.color.rgb = RGBColor(255, 255, 255)

# 现代应用
app_box = slide5.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2.5))
tf_app = app_box.text_frame
tf_app.word_wrap = True
p_a1 = tf_app.paragraphs[0]
p_a1.text = "【现代应用】"
p_a1.font.size = Pt(16)
p_a1.font.bold = True
p_a1.font.color.rgb = RGBColor(30, 60, 114)
apps = [
    "• 职场：被裁员了？也许是转型新事业的机会",
    "• 投资：亏损了？也许是在教会你风险管理",
    "• 人际：被人伤害？也许是在帮你识人"
]
for app in apps:
    p = tf_app.add_paragraph()
    p.text = app
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide5, 5, 10)

# ===== 第6页：西方思想对照 =====
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar6.fill.solid()
title_bar6.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar6.line.fill.background()
txBox6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf9 = txBox6.text_frame
p9 = tf9.paragraphs[0]
p9.text = "🌍 西方思想对照"
p9.font.size = Pt(32)
p9.font.bold = True
p9.font.color.rgb = RGBColor(255, 255, 255)

philosophers = [
    ("�希腊 赫拉克利特", "\"对立统一\"", "没有黑暗就没有光明，没有战争就没有和平。宇宙是对立力量的和谐统一。", "约前500年"),
    ("🇩🇪 黑格尔", "\"辩证法\"", "正题→反题→合题，任何事物都在矛盾中发展转化。", "1800年代"),
    ("🇬🇧 牛顿", "\"相对运动\"", "运动是相对的，快慢取决于参照物。没有绝对标准。", "1687年"),
    ("🇺🇸 爱因斯坦", "\"相对论\"", "时间、空间都是相对的，取决于观察者的运动状态。", "1905年")
]
y = 1.4
for name, theory, desc, year in philosophers:
    phil_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.2))
    phil_box.fill.solid()
    phil_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_phil = phil_box.text_frame
    tf_phil.word_wrap = True
    p_n = tf_phil.paragraphs[0]
    p_n.text = f"{name} · {theory} · {year}"
    p_n.font.size = Pt(16)
    p_n.font.bold = True
    p_n.font.color.rgb = RGBColor(30, 60, 114)
    p_d = tf_phil.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.35

add_slide_number(slide6, 6, 10)

# ===== 第7页：第三章详解 =====
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar7.fill.solid()
title_bar7.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar7.line.fill.background()
txBox7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf10 = txBox7.text_frame
p10 = tf10.paragraphs[0]
p10.text = "第三章 · 无为而治"
p10.font.size = Pt(32)
p10.font.bold = True
p10.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box3 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.8))
yuanwen_box3.fill.solid()
yuanwen_box3.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw3 = yuanwen_box3.text_frame
tf_yw3.word_wrap = True
p_yw3 = tf_yw3.paragraphs[0]
p_yw3.text = "【原文】"
p_yw3.font.size = Pt(18)
p_yw3.font.bold = True
p_yw3.font.color.rgb = RGBColor(30, 60, 114)
for line in ["不尚贤，使民不争；", "不贵难得之货，使民不为盗；", "不见可欲，使民心不乱。"]:
    p = tf_yw3.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box3 = slide7.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(9), Inches(4))
tf_c3 = content_box3.text_frame
tf_c3.word_wrap = True
p_c7 = tf_c3.paragraphs[0]
p_c7.text = "【白话解读】"
p_c7.font.size = Pt(18)
p_c7.font.bold = True
p_c7.font.color.rgb = RGBColor(30, 60, 114)
解读3 = [
    ("不尚贤", "不推崇贤能之名利，人们就不会攀比争夺"),
    ("不贵货", "不珍视稀有财物，人们就不会起盗心"),
    ("不见欲", "不展示诱惑之物，人们心神就不会纷乱"),
    ("核心：无为", "不是什么都不做，而是不人为制造欲望和纷争")
]
for title, content in 解读3:
    p = tf_c3.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide7, 7, 10)

# ===== 第8页：文景之治 =====
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar8.fill.solid()
title_bar8.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar8.line.fill.background()
txBox8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf11 = txBox8.text_frame
p11 = tf11.paragraphs[0]
p11.text = "📜 历史验证：文景之治"
p11.font.size = Pt(32)
p11.font.bold = True
p11.font.color.rgb = RGBColor(255, 255, 255)

# 历史背景
history_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
history_box.fill.solid()
history_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_hist = history_box.text_frame
tf_hist.word_wrap = True
p_h1 = tf_hist.paragraphs[0]
p_h1.text = "【历史背景】（公元前180年-前141年）"
p_h1.font.size = Pt(16)
p_h1.font.bold = True
p_h1.font.color.rgb = RGBColor(30, 60, 114)
p_h2 = tf_hist.add_paragraph()
p_h2.text = "汉文帝、汉景帝以黄老学说治国，实行\"无为而治\"——减轻赋税、与民休息、不兴土木。40年后，国库充盈，百姓安居，创造了中国历史上第一个盛世。"
p_h2.font.size = Pt(15)
p_h2.font.color.rgb = RGBColor(51, 51, 51)

# 具体措施
measures_box = slide8.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(4))
tf_m = measures_box.text_frame
tf_m.word_wrap = True
p_m1 = tf_m.paragraphs[0]
p_m1.text = "【具体措施】"
p_m1.font.size = Pt(18)
p_m1.font.bold = True
p_m1.font.color.rgb = RGBColor(30, 60, 114)
measures = [
    ("田税降至三十税一", "农民负担大减，生产积极性提高"),
    ("废除酷刑", "废除连坐法等严刑峻法，社会更和谐"),
    ("不搞大型土木", "不征发百姓建宫殿，还利于民"),
    ("与匈奴和亲", "避免战争，集中精力发展经济")
]
for title, content in measures:
    p = tf_m.add_paragraph()
    p.text = f"✓ {title}：{content}"
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(8)

add_slide_number(slide8, 8, 10)

# ===== 第9页：生活应用 =====
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar9.fill.solid()
title_bar9.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar9.line.fill.background()
txBox9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf12 = txBox9.text_frame
p12 = tf12.paragraphs[0]
p12.text = "💡 生活中的智慧"
p12.font.size = Pt(32)
p12.font.bold = True
p12.font.color.rgb = RGBColor(255, 255, 255)

app_items = [
    ("🏢 职场管理", "不过度KPI考核，反而激发创造力\n少开会、少汇报，员工反而更主动\n案例：谷歌的'20%自由时间'"),
    ("👨‍👩‍👧 育儿教育", "不拿孩子和别人比较\n不强调'赢在起跑线'\n让孩子按自己的节奏成长"),
    ("🧘 心态调整", "好事来了别得意，可能藏着危机\n坏事来了别绝望，也许转机就在其中\n保持平常心，不被外物所累"),
    ("💰 投资理财", "不追求'高收益'，风险往往埋藏其中\n分散投资，不把鸡蛋放一个篮子\n记住'祸福相依'的辩证思维")
]
y = 1.4
for title, content in app_items:
    app_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.3))
    app_box.fill.solid()
    app_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_app = app_box.text_frame
    tf_app.word_wrap = True
    p_a = tf_app.paragraphs[0]
    p_a.text = title
    p_a.font.size = Pt(16)
    p_a.font.bold = True
    p_a.font.color.rgb = RGBColor(30, 60, 114)
    p_b = tf_app.add_paragraph()
    p_b.text = content
    p_b.font.size = Pt(13)
    p_b.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.45

add_slide_number(slide9, 9, 10)

# ===== 第10页：精华总结 =====
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar10.fill.solid()
title_bar10.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar10.line.fill.background()
txBox10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf13 = txBox10.text_frame
p13 = tf13.paragraphs[0]
p13.text = "💎 精华提炼 · 第一课总结"
p13.font.size = Pt(32)
p13.font.bold = True
p13.font.color.rgb = RGBColor(255, 255, 255)

summary_items = [
    ("第一章·道之本体", "\"道\"是宇宙本源，看不见却支配万物运行"),
    ("第二章·相对之美", "美丑、善恶对立统一，不执着于一边"),
    ("第三章·无为而治", "不过度干预，让事物按其本性自然发展")
]
y = 1.4
for title, content in summary_items:
    s_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.9))
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_s = s_box.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.text = f"{title}"
    p_s.font.size = Pt(16)
    p_s.font.bold = True
    p_s.font.color.rgb = RGBColor(30, 60, 114)
    p_s2 = tf_s.add_paragraph()
    p_s2.text = content
    p_s2.font.size = Pt(14)
    p_s2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.05

# 总结框
conclusion_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
conclusion_box.fill.solid()
conclusion_box.fill.fore_color.rgb = RGBColor(30, 60, 114)
tf_con = conclusion_box.text_frame
tf_con.word_wrap = True
p_con = tf_con.paragraphs[0]
p_con.text = "🌟 一句话总结"
p_con.font.size = Pt(18)
p_con.font.bold = True
p_con.font.color.rgb = RGBColor(255, 255, 255)
p_con2 = tf_con.add_paragraph()
p_con2.text = "道法自然，无为而治"
p_con2.font.size = Pt(24)
p_con2.font.bold = True
p_con2.font.color.rgb = RGBColor(255, 255, 255)
p_con2.alignment = PP_ALIGN.CENTER
p_con3 = tf_con.add_paragraph()
p_con3.text = "顺应规律，不加强制，让事物按其本性发展。"
p_con3.font.size = Pt(16)
p_con3.font.color.rgb = RGBColor(200, 200, 200)

add_slide_number(slide10, 10, 10)

# 保存
prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第1课.pptx')
print("PPT已生成：道德经_第1课.pptx (内容大幅丰富)")
