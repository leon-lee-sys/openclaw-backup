from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 配色：学术白 ──
COLORS = {
    'primary': RGBColor(0, 51, 102),     # 深蓝
    'secondary': RGBColor(102, 102, 102), # 灰
    'accent': RGBColor(204, 0, 0),       # 红色点缀
    'text': RGBColor(51, 51, 51),        # 正文
    'bg': RGBColor(255, 255, 255),       # 白色背景
    'light': RGBColor(240, 248, 255),    # 浅蓝灰
}

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['bg']

def title_bar(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(30); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def quote_box(slide, text, top=1.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(top), Inches(11.73), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['light']
    box.line.color.rgb = COLORS['primary']
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(20); p.font.name = 'KaiTi'
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    tf.margin_left = Pt(20); tf.margin_right = Pt(20)

def section_title(slide, text, top=1.4):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.73), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(22); p.font.bold = True
    p.font.color.rgb = COLORS['accent']

def body_text(slide, lines, top=2.1):
    left, top_, width, height = Inches(0.8), Inches(top), Inches(11.73), Inches(4.5)
    tb = slide.shapes.add_textbox(left, top_, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(17)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(10)

def insight_box(slide, title, points, top=5.3):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(top), Inches(11.73), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['primary']
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"💡 {title}"
    p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 220, 80)
    for pt in points:
        p2 = tf.add_paragraph()
        p2.text = f"  • {pt}"; p2.font.size = Pt(15)
        p2.font.color.rgb = RGBColor(220, 240, 255)

# ═══════════════════════════════════════
# Slide 1 — 封面
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
# 顶部色带
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(2.2))
bar.fill.solid(); bar.fill.fore_color.rgb = COLORS['primary']
bar.line.fill.background()
# 主标题
tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1.2))
p = tb.text_frame.paragraphs[0]
p.text = '《道德经》每日课件'; p.font.size = Pt(44); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER
# 副标题
tb2 = s.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.33), Inches(0.8))
p2 = tb2.text_frame.paragraphs[0]
p2.text = '第8天 · 第二十二章 至 第二十四章'; p2.font.size = Pt(26)
p2.font.color.rgb = RGBColor(200, 230, 255); p2.alignment = PP_ALIGN.CENTER
# 日期
tb3 = s.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.33), Inches(0.6))
p3 = tb3.text_frame.paragraphs[0]
p3.text = '2026年4月15日'; p3.font.size = Pt(20)
p3.font.color.rgb = COLORS['secondary']; p3.alignment = PP_ALIGN.CENTER
# 引言
tb4 = s.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.33), Inches(1.5))
p4 = tb4.text_frame.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
p4.text = '"曲则全，枉则直，洼则盈，敝则新"'; p4.font.size = Pt(22)
p4.font.name = 'KaiTi'; p4.font.color.rgb = COLORS['primary']
p4b = tb4.text_frame.add_paragraph(); p4b.alignment = PP_ALIGN.CENTER
p4b.text = '—— 委曲求全，以退为进'; p4b.font.size = Pt(16)
p4b.font.color.rgb = COLORS['secondary']

# ═══════════════════════════════════════
# Slide 2 — 第二十二章
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十二章｜曲则全')
quote_box(s, '曲则全，枉则直，洼则盈，敝则新，少则得，多则惑。\n是以圣人抱一为天下式。不自见，故明；不自是，故彰；\n不自伐，故有功；不自矜，故长。夫唯不争，故天下莫能与之争。')
section_title(s, '白话解读')
body_text(s, [
    '→ 弯曲才能保全，屈枉才能伸直，低洼才能盈满，破旧才能更新。',
    '→ 少取反而多得，贪多反而迷惑。圣人坚守"道"作为天下的法则。',
    '→ 不自我炫耀，反而更显眼；不自以为是，反而更彰显；',
    '→ 不自我夸耀，反而有功劳；不骄傲自满，反而能长久。',
    '→ 正因为不与人争，所以天下没有人能与他相争。'
], top=3.1)
insight_box(s, '核心智慧：柔弱胜刚强，以退为进', [
    '不争不是消极退缩，而是顺应规律、自然而然地成就',
])

# ═══════════════════════════════════════
# Slide 3 — 第二十二章 典故
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十二章｜典故案例')
section_title(s, '历史典故：孙膑与庞涓（减灶计）')
body_text(s, [
    '【背景】战国时期，孙膑与庞涓同门，庞涓忌惮孙膑之才，加害于他。',
    '【策略】孙膑逃至齐国。庞涓率魏军攻韩，齐国派孙膑援救。',
    '',
    '孙膑令齐军入魏境后逐日减少军灶：第一天十万，第二日五万，第三日三万。',
    '',
    '【结果】庞涓见灶减，误以为齐军怯懦逃散，率轻骑追击。',
    '孙膑在马陵设伏，庞涓兵败自刎。',
    '',
    '【感悟】孙膑"示弱"实为"曲则全"——减少灶数是主动示形，',
    '让敌人做出错误判断，以退为进，以弱胜强。这正是老子"不争而天下莫能与之争"',
    '的军事实践：不争一日之胜负，而争最终的全局胜利。'
], top=1.4)

# ═══════════════════════════════════════
# Slide 4 — 第二十三章
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十三章｜希言自然')
quote_box(s, '希言自然。故飘风不终朝，骤雨不终日。\n孰为此者？天地。天地尚不能久，而况于人乎？\n故从事于道者，道者同于道，德者同于德，失者同于失。\n同于道者，道亦乐得之；同于德者，德亦乐得之；同于失者，失亦乐得之。')
section_title(s, '白话解读')
body_text(s, [
    '→ 少说话才合乎自然。狂风刮不了一早晨，暴雨下不了一整天。',
    '→ 谁造成这种现象？是天地。天地尚且不能长久，何况人呢？',
    '→ 追求道的人，与道同行；追求德的人，与德同行；失去道德的人，则与失同行。',
    '→ 与道同行的人，道也乐于成就他；与德同行的人，德也乐于成就他；',
    '→ 与失同行的人，失去也会找上他。',
    '→ 诚信不足，就会有人不信任你。'
], top=3.1)
insight_box(s, '核心智慧：少言慎行，顺应自然节奏', [
    '狂风骤雨虽猛烈却短暂，合道而行虽缓慢却长久',
])

# ═══════════════════════════════════════
# Slide 5 — 第二十三章 典故
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十三章｜典故案例')
section_title(s, '历史典故：诸葛亮七擒孟获（攻心为上）')
body_text(s, [
    '【背景】蜀汉建兴三年，南方部落首领孟获起兵反抗。',
    '',
    '【策略】诸葛亮放弃"猛攻速胜"，采用"攻心为上"的柔性策略。',
    '对孟获只捉不杀，前后共擒获七次，每一次都释放。',
    '',
    '第一次擒获后，孟获不服，诸葛亮一笑放他回去整顿军马。',
    '第二次、第三次……每一次孟获都有理由不服，诸葛亮每一次都耐心再战。',
    '',
    '【结果】第七次擒获后，孟获终于心服："公，天威也！"南方从此安定。',
    '',
    '【感悟】诸葛亮深谙"飘风不终朝"之理——强力镇压如同飘风骤雨，来得快去得也快。',
    '而"七擒七纵"如同道之"希言"，以最少的暴力换取最长久的和平，',
    '正是本章"同于道者，道亦乐得之"的经典诠释。'
], top=1.4)

# ═══════════════════════════════════════
# Slide 6 — 第二十四章
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十四章｜企者不立')
quote_box(s, '企者不立，跨者不行，自见者不明，自是者不彰，\n自伐者无功，自矜者不长。\n其在道也，曰余食赘行，物或恶之，故有道者不处。')
section_title(s, '白话解读')
body_text(s, [
    '→ 踮起脚尖想要站得高，反而站不稳；大步快走想要走得快，反而走不远。',
    '→ 固执己见的人，反而不能明白事理；自以为是的人，反而不能彰显才华；',
    '→ 自我夸耀的人，反而没有功劳；骄傲自满的人，反而不能长久。',
    '→ 从道的角度来看，这些行为如同剩饭赘瘤，连鬼神都会厌恶。',
    '→ 所以有道之人决不这样做。'
], top=3.1)
insight_box(s, '核心智慧：戒除矜夸，顺其自然', [
    '过度用力适得其反，保持谦逊柔弱方为正道',
])

# ═══════════════════════════════════════
# Slide 7 — 第二十四章 典故
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十四章｜典故案例')
section_title(s, '历史典故：祢衡之死（才高而矜）')
body_text(s, [
    '【背景】祢衡是东汉末年名士，极有才华，但性格傲慢，矜才自负。',
    '',
    '【言行】曹操召见他，他裸身击鼓羞辱曹操；后至刘表处，仍言语不敬；',
    '黄祖性急，祢衡当众辱骂，黄祖忍无可忍，将其处斩，年仅二十五岁。',
    '',
    '【反思】祢衡之才不可否认，然"自矜者不长"——他的才华因过度自矜而反成祸端。',
    '正如老子所言，"余食赘行，物或恶之"，傲才视同残羹赘疣，令人厌恶。',
    '',
    '【对比】同为才士，荀攸"成就功业而若愚"，贾诩"以谦自守而全身"。',
    '真正有道之人，才华与谦逊并存，成就功业却不居功自傲。',
    '老子警告的"自见、自是、自伐、自矜"，是成就大事的四块绊脚石。'
], top=1.4)

# ═══════════════════════════════════════
# Slide 8 — 三章精华总结
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十二至二十四章 · 精华提炼')
# 三栏
for i, (ch, title, pts) in enumerate([
    ('第二十二章', '曲则全', ['柔弱胜刚强，不争而天下莫能争', '委曲求全，以退为进', '抱一为天下式']),
    ('第二十三章', '希言自然', ['少言慎行，顺应天地节奏', '狂风骤雨不能久，平和恒久', '信不足焉，有不信焉']),
    ('第二十四章', '企者不立', ['戒除矜夸，不自见自是', '过度用力适得其反', '余食赘行，物或恶之']),
]):
    left = Inches(0.5 + i * 4.2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             left, Inches(1.3), Inches(4.0), Inches(4.0))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['light']
    box.line.color.rgb = COLORS['primary']
    # 章节号
    tb = s.shapes.add_textbox(left + Inches(0.15), Inches(1.4), Inches(3.7), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = ch
    p.font.size = Pt(14); p.font.color.rgb = COLORS['secondary']
    # 标题
    tb2 = s.shapes.add_textbox(left + Inches(0.15), Inches(1.85), Inches(3.7), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = title
    p2.font.size = Pt(24); p2.font.bold = True; p2.font.color.rgb = COLORS['primary']
    # 要点
    tb3 = s.shapes.add_textbox(left + Inches(0.15), Inches(2.6), Inches(3.7), Inches(2.5))
    tf = tb3.text_frame; tf.word_wrap = True
    for j, pt in enumerate(pts):
        p3 = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p3.text = f'• {pt}'; p3.font.size = Pt(14)
        p3.font.color.rgb = COLORS['text']; p3.space_after = Pt(8)
# 底部总结
bot = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5))
bot.fill.solid(); bot.fill.fore_color.rgb = COLORS['primary']
bot.line.fill.background()
tbb = bot.text_frame; tbb.word_wrap = True
pb = tbb.paragraphs[0]; pb.text = '三章贯通：以柔克刚，顺其自然'
pb.font.size = Pt(22); pb.font.bold = True
pb.font.color.rgb = RGBColor(255, 220, 80); pb.alignment = PP_ALIGN.CENTER
pb2 = tbb.add_paragraph(); pb2.text = '核心心法：不自见、不自是、不自伐、不自矜 → 夫唯不争，故天下莫能与之争'
pb2.font.size = Pt(15); pb2.font.color.rgb = RGBColor(200, 230, 255); pb2.alignment = PP_ALIGN.CENTER

# ── 保存 ──
out = '/Users/mac/.openclaw/workspace/courses/道德经_第8天.pptx'
prs.save(out)
print(f'Saved: {out}')
