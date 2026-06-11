#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
俄语7天速成 Day 1 - Приветствие (问候与介绍)
"""
import os
import sys
sys.path.insert(0, '/opt/homebrew/lib/node_modules/openclaw/node_modules')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(30, 60, 114)
TEXT = RGBColor(51, 51, 51)
BLUE_LIGHT = RGBColor(70, 130, 180)
GOLD = RGBColor(180, 140, 50)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(60, 140, 90)

def make_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(200, 220, 255)

def add_footer(slide, day, theme):
    footer = slide.shapes.add_textbox(Inches(8.5), Inches(7.2), Inches(1.3), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "День %d/7" % day
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(160, 160, 160)
    p.alignment = PP_ALIGN.RIGHT

def add_box(slide, text, left, top, width, height, fill_color=None, font_size=16, bold=False, color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return box

def make_lesson():
    prs = Presentation()
    day = 1
    
    # Slide 1: Cover
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语7天速成班"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    main = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = main.text_frame
    p = tf.paragraphs[0]
    p.text = "День 1: Приветствие"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    sub = s.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "问候与介绍"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER
    
    add_footer(s, day, "Приветствие")
    
    # Slide 2: Core vocabulary
    s = make_slide(prs)
    set_title_bar(s, "核心词汇 · Слова", "Ключевые слова")
    
    words = [
        ("приветствие", "问候"),
        ("знакомство", "介绍/认识"),
        ("здравствуйте", "您好（正式）"),
        ("добрый день", "日安"),
        ("спасибо", "谢谢"),
        ("пожалуйста", "请/不客气"),
        ("извините", "对不起"),
        ("до свидания", "再见"),
    ]
    
    y = 1.6
    for ru, cn in words:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(4), Inches(0.5), fill_color=BLUE_LIGHT, font_size=18, color=WHITE)
        add_box(s, cn, Inches(4.7), Inches(y), Inches(4.8), Inches(0.5), fill_color=GREEN, font_size=18, color=WHITE)
        y += 0.6
    
    add_footer(s, day, "Приветствие")
    
    # Slide 3: Formal greetings
    s = make_slide(prs)
    set_title_bar(s, "正式问候 · Официальное приветствие", "Формальное")
    
    phrases = [
        ("Здравствуйте!", "您好！"),
        ("Доброе утро!", "早上好！"),
        ("Добрый день!", "日安！/下午好！"),
        ("Добрый вечер!", "晚上好！"),
        ("Как ваши дела?", "您好吗？"),
        ("Как поживаете?", "您过得怎么样？"),
        ("Очень приятно!", "很高兴认识您！"),
    ]
    
    y = 1.6
    for ru, cn in phrases:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(5.5), Inches(0.55), fill_color=PRIMARY, font_size=17, color=WHITE)
        add_box(s, cn, Inches(6.2), Inches(y), Inches(3.3), Inches(0.55), font_size=17, color=TEXT)
        y += 0.65
    
    add_footer(s, day, "Приветствие")
    
    # Slide 4: Informal greetings
    s = make_slide(prs)
    set_title_bar(s, "非正式问候 · Неформальное приветствие", "Разговорное")
    
    phrases2 = [
        ("Привет!", "你好！"),
        ("Как дела?", "怎么样？"),
        ("Что нового?", "有什么新鲜事？"),
        ("Как жизнь?", "日子过得怎么样？"),
        ("Пока!", "再见（朋友之间）"),
    ]
    
    y = 1.6
    for ru, cn in phrases2:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(5.5), Inches(0.55), fill_color=GREEN, font_size=17, color=WHITE)
        add_box(s, cn, Inches(6.2), Inches(y), Inches(3.3), Inches(0.55), font_size=17, color=TEXT)
        y += 0.7
    
    add_box(s, "正式：здравствуйте / 非正式：привет", Inches(0.5), Inches(5.5), Inches(9), Inches(0.5), fill_color=GOLD, font_size=14, color=WHITE)
    add_footer(s, day, "Приветствие")
    
    # Slide 5: Self introduction
    s = make_slide(prs)
    set_title_bar(s, "自我介绍 · Знакомство", "Представление себя")
    
    dialogues = [
        ("Меня зовут...", "我的名字是..."),
        ("Я из Шанхая.", "我来自上海。"),
        ("Я работаю менеджером.", "我是经理。"),
        ("Мне нравится читать книги.", "我喜欢读书。"),
    ]
    
    y = 1.6
    for ru, cn in dialogues:
        add_box(s, ru, Inches(0.5), Inches(y), Inches(5.5), Inches(0.55), fill_color=PRIMARY, font_size=16, color=WHITE)
        add_box(s, cn, Inches(6.2), Inches(y), Inches(3.3), Inches(0.55), font_size=16, color=TEXT)
        y += 0.65
    
    add_footer(s, day, "Приветствие")
    
    # Slide 6: Dialogue
    s = make_slide(prs)
    set_title_bar(s, "情景对话 · Диалог", "Ситуация")
    
    dlg = [
        "— Здравствуйте! Меня зовут Ли Мин. Очень приятно!",
        "— Здравствуйте! Я Анна. Тоже очень приятно!",
        "— Откуда вы?",
        "— Я из Китая. А вы?",
        "— Я из России. Вы работаете здесь?",
        "— Да, я менеджер в компании.",
    ]
    
    y = 1.6
    for line in dlg:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        y += 0.5
    
    add_footer(s, day, "Приветствие")
    
    # Slide 7: Grammar
    s = make_slide(prs)
    set_title_bar(s, "语法要点 · Грамматика", "Ключевые моменты")
    
    points = [
        "1. Меня зовут + 名字 = 自我介绍",
        "2. Я из + 地方 = 我来自...",
        "3. Я + 职业 = 我是...",
        "4. 也可以说：Моё имя + 名字",
    ]
    
    y = 1.6
    for point in points:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT
        y += 0.55
    
    add_footer(s, day, "Приветствие")
    
    # Slide 8: Practice
    s = make_slide(prs)
    set_title_bar(s, "课后练习 · Упражнения", "Практика")
    
    exercises = [
        "1. 用俄语说你好和再见",
        "2. 介绍自己：Меня зовут... Я из...",
        "3. 翻译：您好，很高兴认识您！",
        "4. 回答：Как ваши дела?",
    ]
    
    y = 1.7
    for ex in exercises:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = ex
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT
        y += 0.55
    
    add_footer(s, day, "Приветствие")
    
    # Slide 9: Summary
    s = make_slide(prs)
    set_title_bar(s, "今日总结 · Итоги", "Что мы выучили")
    
    summary = [
        "✅ 正式/非正式问候语",
        "✅ 自我介绍：名字、来自哪里、职业",
        "✅ 礼貌用语：спасибо, пожалуйста, извините",
        "✅ 简单对话练习",
    ]
    
    y = 1.7
    for item in summary:
        t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GREEN
        y += 0.6
    
    add_box(s, "Завтра: День 2 - На работе (工作职场)", Inches(1), Inches(5.5), Inches(8), Inches(0.6), fill_color=PRIMARY, font_size=16, color=WHITE)
    add_footer(s, day, "Приветствие")
    
    # Save
    output_dir = "/Users/mac/.openclaw/workspace/courses/俄语课件"
    os.makedirs(output_dir, exist_ok=True)
    path = os.path.join(output_dir, "俄语速成_Day1.pptx")
    prs.save(path)
    print("Done: 俄语速成_Day1.pptx")
    return path

if __name__ == "__main__":
    make_lesson()
