from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

COLORS = {
    'primary': RGBColor(0, 51, 102),
    'secondary': RGBColor(102, 102, 102),
    'accent': RGBColor(204, 0, 0),
    'text': RGBColor(51, 51, 51),
    'bg': RGBColor(255, 255, 255),
    'light': RGBColor(240, 248, 255),
}

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['bg']

def title_bar(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(30); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def quote_box(slide, text, top=1.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(top), Inches(11.73), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['light']
    box.line.color.rgb = COLORS['primary']
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(20); p.font.name = 'KaiTi'
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    tf.margin_left = Pt(20); tf.margin_right = Pt(20)

def section_title(slide, text, top=1.4):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.73), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(22); p.font.bold = True
    p.font.color.rgb = COLORS['accent']

def body_text(slide, lines, top=2.1):
    left, top_, width, height = Inches(0.8), Inches(top), Inches(11.73), Inches(4.5)
    tb = slide.shapes.add_textbox(left, top_, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(17)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(10)

def insight_box(slide, title, points, top=5.3):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(top), Inches(11.73), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['primary']
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"💡 {title}"
    p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 220, 80)
    for pt in points:
        p2 = tf.add_paragraph()
        p2.text = f"  • {pt}"; p2.font.size = Pt(15)
        p2.font.color.rgb = RGBColor(220, 240, 255)

# Slide 1 - Cover
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(2.2))
bar.fill.solid(); bar.fill.fore_color.rgb = COLORS['primary']
bar.line.fill.background()
tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1.2))
p = tb.text_frame.paragraphs[0]
p.text = '《道德经》每日课件'; p.font.size = Pt(44); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER
tb2 = s.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.33), Inches(0.8))
p2 = tb2.text_frame.paragraphs[0]
p2.text = '第10天 · 第三十一章 至 第三十三章'; p2.font.size = Pt(26)
p2.font.color.rgb = RGBColor(200, 230, 255); p2.alignment = PP_ALIGN.CENTER
tb3 = s.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.33), Inches(0.6))
p3 = tb3.text_frame.paragraphs[0]
p3.text = '2026年4月17日'; p3.font.size = Pt(20)
p3.font.color.rgb = COLORS['secondary']; p3.alignment = PP_ALIGN.CENTER
tb4 = s.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.33), Inches(1.5))
p4 = tb4.text_frame.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
p4.text = '「夫佳兵者，不祥之器……君子居则贵左，用兵则贵右」'; p4.font.size = Pt(22)
p4.font.name = 'KaiTi'; p4.font.color.rgb = COLORS['primary']
p4b = tb4.text_frame.add_paragraph(); p4b.alignment = PP_ALIGN.CENTER
p4b.text = '—— 胜而不美，以哀悲泣之'; p4b.font.size = Pt(16)
p4b.font.color.rgb = COLORS['secondary']

# Slide 2 - Ch31
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第三十一章｜夫佳兵者')
quote_box(s, '夫佳兵者，不祥之器。物或恶之，故有道者不处。\n君子居则贵左，用兵则贵右。吉事尚左，凶事尚右。\n偏将军居左，上将军居右，言以丧礼处之。')
section_title(s, '白话解读')
body_text(s, [
    '→ 精兵利器是不吉祥的东西，人们都厌恶它，所以有道之人不使用它。',
    '→ 君子平时以左边为尊，用兵时以右边为尊。吉庆以左为上，凶丧以右为上。',
    '→ 偏将军站在左边，上将军站在右边，用办丧事的礼仪对待战争。',
    '→ 杀人众多，要以悲哀的心情去看待；战胜了也要以丧礼来处理。'
], top=3.1)
insight_box(s, '核心智慧：兵者不祥，慎用而悲', [
    '战争是不得已的手段，应用悲伤之心对待，而非以胜利为荣',
])

# Slide 3 - Ch31 Story
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第三十一章｜典故案例')
section_title(s, '历史典故：商鞅之死与穷兵黩武之鉴')
body_text(s, [
    '【背景】商鞅辅佐秦孝公变法，使秦国崛起，却因严刑峻法得罪贵族。',
    '秦孝公死后，商鞅被公子虔诬告谋反，遭车裂之刑，全家灭族。',
    '',
    '【老子警示】老子说「佳兵者不祥」，商鞅恰是一个注解。',
    '他以法家铁腕立威，虽令秦一时强盛，却种下严刑峻法的祸根。',
    '秦以暴力得天下，亦以暴力失天下——二世而亡，非天命，乃人祸。',
    '',
    '【对比】汉初张良运筹帷幄，却功成身退，不贪军权，得以善终。',
    '而商鞅不知收敛，终致灭门。正合老子「有道者不处」之诫。',
    '',
    '【感悟】「言以丧礼处之」——战争是凶事，当以庄重、悲哀之心处置。',
    '以战为荣、以杀为乐者，必遭反噬。管理者当以此为鉴。'
], top=1.4)

# Slide 4 - Ch32
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第三十二章｜知止不殆')
quote_box(s, '道常无名，朴虽小，天下莫能臣也。\n侯王若能守之，万物将自宾。\n始制有名，名亦既有，夫亦将知止。知止所以不殆。')
section_title(s, '白话解读')
body_text(s, [
    '→ 道永远处于无名而质朴的状态，就像未经雕琢的木材，',
    '→ 虽然渺小，但天下没有人能让它臣服。侯王若能守住它，万物将自然归附。',
    '→ 天地之气相融合，就会降下甘露，无需人们命令就自然均匀。',
    '→ 制度建立后就有了名称，有了名称之后就要知道适可而止。',
    '→ 知道适可而止，才能避免危险。'
], top=3.1)
insight_box(s, '核心智慧：知止不殆，适可而止', [
    '朴虽小而天下莫能臣，知止是管理者最高明的修养',
])

# Slide 5 - Ch32 Story
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第三十二章｜典故案例')
section_title(s, '历史典故：乾隆十全武功与曾国藩知止之道')
body_text(s, [
    '【乾隆】乾隆在位60年，自诩有「十全武功」——平准噶尔、定回部、',
    '扫金川、靖台湾、降缅甸、定安南、受廓尔喀，武功显赫。',
    '',
    '【隐忧】然而连年战争耗尽国库，官员腐败横行，民变频发。',
    '乾隆晚年宠信和珅，吏治败坏，白莲教起义爆发，',
    '清王朝由盛转衰的伏笔已经埋下。',
    '',
    '【老子之戒】「知止所以不殆」——乾隆不知止。',
    '曾国藩平定太平天国，功高震主，却主动裁撤湘军，功成身退，',
    '正是深谙「知止」之道的典范。其家书流传至今，「知止」之悟实源于此章。'
], top=1.4)

# Slide 6 - Ch33
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第三十三章｜知人者智')
quote_box(s, '知人者智，自知者明。\n胜人者有力，自胜者强。\n知足者富，强行者有志。\n不失其所者久，死而不亡者寿。')
section_title(s, '白话解读')
body_text(s, [
    '→ 能了解别人的人是智慧的，能认识自己的人才是真正清明通达的。',
    '→ 能战胜别人是有力气的，能战胜自己才是真正的强者。',
    '→ 知道满足的人才是富有的，坚持力行的人有志向。',
    '→ 不丧失根基的人才能长久，肉体已死但精神长存的人才是长寿。'
], top=3.1)
insight_box(s, '核心智慧：自知自胜，死而不亡', [
    '真正的强大是战胜自己，精神不朽才是真正的长寿',
])

# Slide 7 - Ch33 Story
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第三十三章｜典故案例')
section_title(s, '历史典故：范蠡与张良（自知者明）')
body_text(s, [
    '【范蠡】越王勾践卧薪尝胆，终于灭吴复国。',
    '功成名就之际，范蠡急流勇退，泛舟五湖，经商三致千金，',
    '被后世尊为「商圣」。而文种不知止，终被勾践赐死。',
    '范蠡之言：「飞鸟尽，良弓藏；狡兔死，走狗烹。」',
    '',
    '【张良】张良谋刺秦始皇未遂，后在圯上遇黄石公授书，',
    '遂成帝王之师。他深谙「自胜者强」，功成不居，',
    '刘邦称其「运筹帷幄之中，决胜千里之外」，却托言辟谷，得以善终。',
    '',
    '【感悟】范蠡与张良，皆知「知止不殆」之理。',
    '知人者智已不易，自知者明更难——知道何时该退，',
    '方能在乱世中全身而退，在历史中流芳百世。'
], top=1.4)

# Slide 8 - Summary
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第三十一至三十三章 · 精华提炼')
for i, (ch, title, pts) in enumerate([
    ('第三十一章', '夫佳兵者', ['兵者不祥，有道不处', '以丧礼处战，慎用武力', '杀人众多，以哀悲泣之']),
    ('第三十二章', '知止不殆', ['道常无名，朴虽小莫能臣', '万物自宾，知止不殆', '适可而止，避免危险']),
    ('第三十三章', '知人者智', ['自知者明，自胜者强', '知足者富，不失所者久', '死而不亡者寿，精神不朽']),
]):
    left = Inches(0.5 + i * 4.2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             left, Inches(1.3), Inches(4.0), Inches(4.0))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['light']
    box.line.color.rgb = COLORS['primary']
    tb = s.shapes.add_textbox(left + Inches(0.15), Inches(1.4), Inches(3.7), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = ch
    p.font.size = Pt(14); p.font.color.rgb = COLORS['secondary']
    tb2 = s.shapes.add_textbox(left + Inches(0.15), Inches(1.85), Inches(3.7), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = title
    p2.font.size = Pt(24); p2.font.bold = True; p2.font.color.rgb = COLORS['primary']
    tb3 = s.shapes.add_textbox(left + Inches(0.15), Inches(2.6), Inches(3.7), Inches(2.5))
    tf = tb3.text_frame; tf.word_wrap = True
    for j, pt in enumerate(pts):
        p3 = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p3.text = f'• {pt}'; p3.font.size = Pt(14)
        p3.font.color.rgb = COLORS['text']; p3.space_after = Pt(8)
bot = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5))
bot.fill.solid(); bot.fill.fore_color.rgb = COLORS['primary']
bot.line.fill.background()
tbb = bot.text_frame; tbb.word_wrap = True
pb = tbb.paragraphs[0]; pb.text = '三章贯通：慎武知止，自知自胜'
pb.font.size = Pt(22); pb.font.bold = True
pb.font.color.rgb = RGBColor(255, 220, 80); pb.alignment = PP_ALIGN.CENTER
pb2 = tbb.add_paragraph()
pb2.text = '核心心法：兵者不祥 → 朴虽小而莫能臣，知止不殆 → 自知者明，自胜者强'
pb2.font.size = Pt(15); pb2.font.color.rgb = RGBColor(200, 230, 255); pb2.alignment = PP_ALIGN.CENTER

out = '/Users/mac/.openclaw/workspace/courses/道德经_第10天.pptx'
prs.save(out)
print(f'Saved: {out}')
