#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
道德经 第13天（40-42章）课件生成
风格：学术白
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案 - 学术白
COLORS = {
    'primary': RGBColor(0, 51, 102),      # 深蓝
    'secondary': RGBColor(102, 102, 102),  # 灰色
    'accent': RGBColor(178, 34, 34),      # 深红（强调）
    'text': RGBColor(51, 51, 51),         # 正文色
    'bg': RGBColor(255, 255, 255),        # 白色背景
    'header_bg': RGBColor(0, 51, 102),    # 深蓝
}

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['bg']

def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['header_bg']
    bar.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_text_box(slide, left, top, width, height, font_size=16, color=None, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf, txBox

def add_paragraph(tf, text, font_size=16, bold=False, color=None, space_before=0, space_after=8, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph() if len(tf.paragraphs) > 0 else tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS['text']
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.alignment = align
    return p

# ===== 第1页：封面 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

# 主标题
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "《道德经》每日一课"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.alignment = PP_ALIGN.CENTER

# 副标题
p2 = tf.add_paragraph()
p2.text = "第13天 · 第40-42章"
p2.font.size = Pt(26)
p2.font.color.rgb = COLORS['secondary']
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(16)

# 主题
p3 = tf.add_paragraph()
p3.text = "反者道之动 · 阴阳相生"
p3.font.size = Pt(20)
p3.font.color.rgb = COLORS['accent']
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(20)

# 日期
p4 = tf.add_paragraph()
p4.text = "2026年4月20日"
p4.font.size = Pt(16)
p4.font.color.rgb = COLORS['secondary']
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(40)

# 底部引言
bottom = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
btf = bottom.text_frame
bp = btf.paragraphs[0]
bp.text = "「上善若水，水善利万物而不争。」"
bp.font.size = Pt(14)
bp.font.italic = True
bp.font.color.rgb = COLORS['secondary']
bp.alignment = PP_ALIGN.CENTER

# ===== 第2页：第四十章 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "第四十章 · 反者道之动")

# 原文框
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(9.2), Inches(1.6))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(240, 248, 255)
box1.line.color.rgb = COLORS['primary']

tf1 = box1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "【原文】"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.space_before = Pt(6)

p = tf1.add_paragraph()
p.text = "反者道之动，弱者道之用。天下万物生于有，有生于无。"
p.font.size = Pt(18)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(10)

# 白话解读
txBox = slide.shapes.add_textbox(Inches(0.4), Inches(3.1), Inches(9.2), Inches(2.0))
tf2 = txBox.text_frame
tf2.word_wrap = True

p = tf2.paragraphs[0]
p.text = "【白话解读】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']

p = tf2.add_paragraph()
p.text = "道的运动是循环往复的，道的作用是微妙柔弱的。天下万物产生于具体的有（有形的存在），而有又产生于虚无的道（无形的本原）。这是老子对宇宙生成论的核心表述——从「无」到「有」的演化过程。"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(8)
p.line_spacing = 1.5

# 精华提炼
txBox2 = slide.shapes.add_textbox(Inches(0.4), Inches(5.2), Inches(9.2), Inches(1.8))
tf3 = txBox2.text_frame
tf3.word_wrap = True

p = tf3.paragraphs[0]
p.text = "【精华提炼】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']

bullets = [
    "「反」：返回、循环，道的运动规律是周而复始的",
    "「弱」：柔弱胜刚强，道以柔弱为用",
    "「有无相生」：一切从无中来，又向无中去"
]
for bullet in bullets:
    p = tf3.add_paragraph()
    p.text = f"• {bullet}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text']
    p.space_before = Pt(6)

# ===== 第3页：第四十一章 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "第四十一章 · 上士闻道")

# 原文
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(9.2), Inches(1.8))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(240, 248, 255)
box2.line.color.rgb = COLORS['primary']

tf4 = box2.text_frame
tf4.word_wrap = True
p = tf4.paragraphs[0]
p.text = "【原文】"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.space_before = Pt(6)

p = tf4.add_paragraph()
p.text = "上士闻道，勤而行之；中士闻道，若存若亡；下士闻道，大笑之。不笑不足以为道。"
p.font.size = Pt(17)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(10)

p = tf4.add_paragraph()
p.text = "故建言有之：明道若昧，进道若退，夷道若纇，上德若谷，大白若辱，广德若不足，建德若偷，质真若渝，大方无隅，大器晚成，大音希声，大象无形，道隐无名。"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(12)

# 白话
txBox3 = slide.shapes.add_textbox(Inches(0.4), Inches(3.2), Inches(9.2), Inches(2.2))
tf5 = txBox3.text_frame
tf5.word_wrap = True

p = tf5.paragraphs[0]
p.text = "【白话解读】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']

p = tf5.add_paragraph()
p.text = "上等士人听闻道，就会努力去实践；中等士人将信将疑；下等士人则加以嘲笑。不被嘲笑，就不足以称为道了。因而古语说：光明的道好像暗昧；前进的道好像后退；平坦的道好像崎岖；崇高的德好像低洼的山谷……真正的道理往往与表面现象相反，需要用心体悟。"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(8)
p.line_spacing = 1.4

# 精华
txBox4 = slide.shapes.add_textbox(Inches(0.4), Inches(5.5), Inches(9.2), Inches(1.5))
tf6 = txBox4.text_frame
tf6.word_wrap = True

p = tf6.paragraphs[0]
p.text = "【精华提炼】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']

p = tf6.add_paragraph()
p.text = "• 真金不怕火炼：真正的道，往往超越常人认知"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(5)

p = tf6.add_paragraph()
p.text = "• 大器晚成：成就大事需要时间与耐心"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(5)

p = tf6.add_paragraph()
p.text = "• 大音希声：最宏大的声音反而听不见"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(5)

# ===== 第4页：第四十二章 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "第四十二章 · 道生一")

# 原文
box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(9.2), Inches(1.5))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(240, 248, 255)
box3.line.color.rgb = COLORS['primary']

tf7 = box3.text_frame
tf7.word_wrap = True
p = tf7.paragraphs[0]
p.text = "【原文】"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.space_before = Pt(6)

p = tf7.add_paragraph()
p.text = "道生一，一生二，二生三，三生万物。万物负阴而抱阳，冲气以为和。"
p.font.size = Pt(17)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(10)

p = tf7.add_paragraph()
p.text = "人之所恶，唯孤、寡、不谷，而王公以为称。故物或损之而益，或益之而损。人之所教，我亦教之：强梁者不得其死，吾将以为教父。"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(10)

# 白话
txBox5 = slide.shapes.add_textbox(Inches(0.4), Inches(3.0), Inches(9.2), Inches(2.0))
tf8 = txBox5.text_frame
tf8.word_wrap = True

p = tf8.paragraphs[0]
p.text = "【白话解读】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']

p = tf8.add_paragraph()
p.text = "道是独一无二的原初存在，一生二（产生阴阳二气），二生三（阴阳二气交合），三生万物。万物背阴而向阳，阴阳二气相互激荡而达到和谐。世人厌恶孤、寡、不谷，但王公却用这些词自称。事物有时减损反而有益，有时增益反而有害。"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(8)
p.line_spacing = 1.4

# 精华
txBox6 = slide.shapes.add_textbox(Inches(0.4), Inches(5.1), Inches(9.2), Inches(1.9))
tf9 = txBox6.text_frame
tf9.word_wrap = True

p = tf9.paragraphs[0]
p.text = "【精华提炼】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']

bullets2 = [
    "「道生万物」：从无到有，从一到万，宇宙生成的自然规律",
    "「负阴抱阳」：万事万物皆包含对立统一的阴阳两面",
    "「冲气为和」：和谐来自对立力量的平衡与调适",
    "「强梁者不得其死」：过于强硬反而招致灭亡"
]
for b in bullets2:
    p = tf9.add_paragraph()
    p.text = f"• {b}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLORS['text']
    p.space_before = Pt(4)

# ===== 第5页：典故案例 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "典故案例 · 塞翁失马")

# 引言
txBox7 = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.6))
tf10 = txBox7.text_frame
p = tf10.paragraphs[0]
p.text = "「祸兮福所倚，福兮祸所伏。」—— 祸福相依，阴阳互转"
p.font.size = Pt(15)
p.font.italic = True
p.font.color.rgb = COLORS['secondary']
p.alignment = PP_ALIGN.CENTER

# 故事框
box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(2.0), Inches(9.2), Inches(4.5))
box4.fill.solid()
box4.fill.fore_color.rgb = RGBColor(250, 250, 250)
box4.line.color.rgb = COLORS['secondary']

tf11 = box4.text_frame
tf11.word_wrap = True

p = tf11.paragraphs[0]
p.text = "《淮南子·人间训》"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['accent']
p.space_before = Pt(12)
p.alignment = PP_ALIGN.CENTER

p = tf11.add_paragraph()
p.text = "从前，边塞地区住着一个精通道术的老人。他的马无缘无故跑到了胡人的地方，大家都为此安慰他。父亲说：'这未必不是一件好事。'"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(14)
p.line_spacing = 1.6

p = tf11.add_paragraph()
p.text = "过了几个月，那匹马带着一群胡地的骏马回来了。大家都来祝贺。父亲说：'这未必不会成为一件坏事。'家里有好马，儿子喜欢骑马，结果从马上摔下来折断了大腿。"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(14)
p.line_spacing = 1.6

p = tf11.add_paragraph()
p.text = "过了一年，胡人大举入侵，边塞壮年男子都拿起弓箭去打仗。靠近边塞的人，十个有九个都战死了。只有他儿子因为腿瘸的缘故，得以保全性命。"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(14)
p.line_spacing = 1.6

p = tf11.add_paragraph()
p.text = "所以《易经》说：'吉凶祸福都是自己招致的。'"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(14)
p.line_spacing = 1.6

# 底部总结
p = tf11.add_paragraph()
p.text = "✅ 印证第四十章「反者道之动」：事物发展到极点，必然向反面转化"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.space_before = Pt(18)

# ===== 第6页：总结页 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "第40-42章 · 精华总结")

# 三章核心
items = [
    ("第四十章", "反者道之动", "道的运动是循环往复的，万物生于有，有生于无"),
    ("第四十一章", "上士闻道", "真正悟道者勤而行之，真理往往超越表面认知"),
    ("第四十二章", "道生一", "阴阳相生，刚柔并济，强梁者不得其死"),
]

for i, (chapter, theme, desc) in enumerate(items):
    top = 1.4 + i * 1.7
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(top), Inches(9.2), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 248, 255)
    box.line.color.rgb = COLORS['primary']
    
    tf_box = box.text_frame
    tf_box.word_wrap = True
    
    p = tf_box.paragraphs[0]
    p.text = f"{chapter}：{theme}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.space_before = Pt(10)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf_box.add_paragraph()
    p.text = desc
    p.font.size = Pt(15)
    p.font.color.rgb = COLORS['text']
    p.space_before = Pt(8)
    p.alignment = PP_ALIGN.CENTER

# 底部寄语
bottom_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.4), Inches(7), Inches(0.8))
bottom_box.fill.solid()
bottom_box.fill.fore_color.rgb = COLORS['primary']
bottom_box.line.fill.background()

tf_btm = bottom_box.text_frame
p = tf_btm.paragraphs[0]
p.text = "💡 今日功课：体会「反者道之动」的智慧，在生活中观察阴阳转化之道"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# ===== 保存 =====
output_dir = "/Users/mac/.openclaw/workspace/courses"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "道德经_第13天.pptx")
prs.save(output_path)
print(f"✅ 课件已保存：{output_path}")
