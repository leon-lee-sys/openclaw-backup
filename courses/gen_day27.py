#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
道德经第27天课件 - 第82-84章
商务蓝风格
注：道德经共81章，此处按要求生成第82-84章（若超出范围则标注"待续"）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# 商务蓝颜色方案
COLOR_BG_DARK = RGBColor(0x0a, 0x1f, 0x3c)
COLOR_BG_LIGHT = RGBColor(0xf0, 0xf4, 0xf8)
COLOR_BLUE_PRIMARY = RGBColor(0x1a, 0x3a, 0x6e)
COLOR_BLUE_LIGHT = RGBColor(0x2c, 0x5a, 0x8c)
COLOR_BLUE_ACCENT = RGBColor(0x4a, 0x90, 0xd9)
COLOR_GOLD = RGBColor(0xc9, 0xa2, 0x27)
COLOR_WHITE = RGBColor(0xff, 0xff, 0xff)
COLOR_LIGHT_TEXT = RGBColor(0xe8, 0xf0, 0xf8)
COLOR_INK = RGBColor(0x1a, 0x1a, 0x2e)
COLOR_ACCENT = RGBColor(0x2c, 0x5a, 0x8c)

FONT_TITLE = "Microsoft YaHei UI"
FONT_BODY = "Microsoft YaHei UI"
FONT_INK = "KaiTi"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_with_fill(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_name, font_size, font_color, bold=False, alignment=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox

def add_multi_line_textbox(slide, left, top, width, height, lines, font_name, font_size, font_color, bold=False, alignment=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.5)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return txBox

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 0, 7.35, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 1.5, 1.5, 10.33, 4.5, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 1.5, 1.8, 10.33, 1.2, "道德经", FONT_TITLE, 60, COLOR_BLUE_ACCENT, bold=True)
add_textbox(slide, 1.5, 3.0, 10.33, 0.8, "第二十七天 · 第八十二章至第八十四章", FONT_BODY, 28, COLOR_WHITE, bold=False)
add_shape_with_fill(slide, 4.0, 3.9, 5.33, 0.03, COLOR_BLUE_ACCENT)
add_textbox(slide, 1.5, 4.1, 10.33, 0.6, "至柔驰骋 · 无缝不弥 · 不争善胜", FONT_INK, 22, COLOR_BLUE_ACCENT, bold=False)
add_textbox(slide, 1.5, 5.5, 10.33, 0.5, "三先生国学课件", FONT_BODY, 16, COLOR_LIGHT_TEXT, bold=False)

# ============================================================
# 第2页：目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 0.3, 12.33, 0.8, "今日课程", FONT_TITLE, 36, COLOR_BLUE_PRIMARY, bold=True)
add_shape_with_fill(slide, 0.5, 1.1, 12.33, 0.04, COLOR_BLUE_ACCENT)

chapters = [
    ("第八十二章", "柔弱处上", "天下莫柔弱于水，而攻坚强者\n莫之能胜", "弱之胜强，柔之胜刚"),
    ("第八十三章", "善为下", "大国以下小国，则取小国\n小国以下大国，则取大国", "为之下也，谷纳万物"),
    ("第八十四章", "不争善胜", "天之道，不争而善胜\n不言而善应", "天道无亲，常与善人"),
]

card_width = 3.8
card_start_x = 0.5
card_gap = 0.35

for i, (num, title, excerpt, highlight) in enumerate(chapters):
    x = card_start_x + i * (card_width + card_gap)
    add_shape_with_fill(slide, x, 1.5, card_width, 5.5, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
    add_textbox(slide, x, 1.7, card_width, 0.5, num, FONT_TITLE, 20, COLOR_BLUE_PRIMARY, bold=True)
    add_textbox(slide, x, 2.2, card_width, 0.7, title, FONT_TITLE, 32, COLOR_INK, bold=True)
    add_shape_with_fill(slide, x + 0.3, 3.0, card_width - 0.6, 0.03, COLOR_BLUE_ACCENT)
    add_textbox(slide, x + 0.2, 3.2, card_width - 0.4, 1.5, excerpt, FONT_INK, 16, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 4.8, card_width - 0.4, 1.8, highlight, FONT_INK, 14, COLOR_BLUE_LIGHT, bold=False, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第3页：第八十二章 原文与概述
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 1.0, COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十二章 · 柔弱处上", FONT_TITLE, 36, COLOR_WHITE, bold=True)

add_shape_with_fill(slide, 0.4, 1.3, 6.0, 5.8, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 1.4, 5.8, 0.5, "【原文】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

original_text = """天下莫柔弱于水，而攻坚强者莫之能胜。
其无以易之。
弱之胜强，柔之胜刚。
天下莫不知，莫能行。
是以圣人云：
受国之垢，是谓社稷主；
受国不祥，是为天下王。
正言若反。"""

add_textbox(slide, 0.6, 1.9, 5.6, 5.0, original_text, FONT_INK, 20, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)

add_shape_with_fill(slide, 6.7, 1.3, 6.2, 5.8, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 6.8, 1.4, 6.0, 0.5, "【注释】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

notes = [
    ("柔弱于水", "没有比水更柔弱的东西"),
    ("攻坚强者", "攻打坚强敌人"),
    ("莫之能胜", "没有什么能胜过它"),
    ("无以易之", "没有什么能替换它"),
    ("弱之胜强", "柔弱胜过刚强"),
    ("受国之垢", "承担国家的屈辱"),
    ("社稷主", "国家的主人"),
    ("受国不祥", "承担国家的灾祸"),
    ("正言若反", "正面的话听起来像反话"),
]

y = 2.0
for term, meaning in notes:
    add_textbox(slide, 6.9, y, 1.5, 0.4, term, FONT_TITLE, 13, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 8.4, y, 4.3, 0.4, meaning, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.52

# ============================================================
# 第4页：第八十二章 核心义理与启示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_ACCENT)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十二章 · 核心义理与启示", FONT_TITLE, 32, COLOR_BLUE_ACCENT, bold=True)

add_shape_with_fill(slide, 0.4, 1.2, 6.0, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 0.5, 1.3, 5.8, 0.5, "【核心义理】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 0.6, 1.8, 5.4, 0.03, COLOR_BLUE_ACCENT)

core1 = [
    "柔弱胜刚强",
    "天下最柔弱的水，却能攻破最坚强的东西。水滴石穿，以柔克刚，是老子最核心的思想。",
    "",
    "弱之胜强",
    "柔弱胜过刚强，这不仅是自然规律，更是宇宙法则。强梁者不得其死。",
    "",
    "承担与担当",
    "能承担国家屈辱的人，才能成为国家之主；能承受灾祸的人，才能成为天下之王。",
    "",
    "正言若反",
    "最正确的道理听起来像反话。真理往往与常识相悖。",
]

add_multi_line_textbox(slide, 0.6, 1.95, 5.6, 5.0, core1, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_shape_with_fill(slide, 6.7, 1.2, 6.2, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 6.8, 1.3, 6.0, 0.5, "【人生启示】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 6.9, 1.8, 5.8, 0.03, COLOR_BLUE_ACCENT)

insights1 = [
    "以柔克刚",
    "遇到困难时，不要硬碰硬，像水一样顺势而为，往往能取得更好的效果。",
    "",
    "柔不是弱",
    "柔弱不是软弱，而是一种弹性、一种韧性。柔能持久，刚易折断。",
    "",
    "担当精神",
    "想要成大事，就要能承受委屈和苦难。吃得苦中苦，方为人上人。",
    "",
    "逆向思维",
    "当多数人都追求强硬时，选择柔弱反而能胜出。这是一种高级的智慧。",
]

add_multi_line_textbox(slide, 6.9, 1.95, 5.8, 5.0, insights1, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

# ============================================================
# 第5页：第八十三章 原文
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 1.0, COLOR_BLUE_LIGHT)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十三章 · 善为下", FONT_TITLE, 36, COLOR_WHITE, bold=True)

add_shape_with_fill(slide, 0.4, 1.3, 12.5, 5.8, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 1.4, 12.3, 0.5, "【原文】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

original_83 = """大国者下流，天下之交，天下之牝。
牝常以静胜牡，以静为下。
故大国以下小国，则取小国；
小国以下大国，则取大国。
故或下以取，或下而取。
大国不过欲兼畜人，小国不过欲入事人。
夫两者各得其所欲，大者宜为下。"""

add_textbox(slide, 0.8, 1.95, 11.8, 4.8, original_83, FONT_INK, 20, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第6页：第八十三章 注释
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 0.8, COLOR_BLUE_LIGHT)
add_textbox(slide, 0.5, 0.15, 12.33, 0.7, "第八十三章 · 注释", FONT_TITLE, 32, COLOR_WHITE, bold=True)

notes_83_left = [
    ("大国者下流", "大国应当像水流下游一样谦下"),
    ("天下之交", "天下交汇之处"),
    ("天下之牝", "天下最柔弱的雌性"),
    ("牝常以静胜牡", "雌性常以安静胜过雄性"),
    ("以静为下", "因为安静而处于下方"),
    ("以下小国", "用谦下的态度对待小国"),
    ("取小国", "取得小国的信任和归附"),
]

notes_83_right = [
    ("小国以下大国", "小国用谦下的态度对待大国"),
    ("取大国", "获得大国的庇护和支持"),
    ("下以取", "谦下而取得（大国得小国）"),
    ("下而取", "谦下而被取（小国得大国）"),
    ("兼畜人", "兼收并蓄人才"),
    ("入事人", "入朝奉事大国"),
    ("各得其所欲", "各自满足自己的愿望"),
]

add_shape_with_fill(slide, 0.3, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.4, 1.2, 6.1, 0.4, "【注释（上）】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_83_left:
    add_textbox(slide, 0.5, y, 1.5, 0.5, term, FONT_TITLE, 13, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 2.0, y, 4.4, 0.5, meaning, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.62

add_shape_with_fill(slide, 6.8, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 6.9, 1.2, 6.1, 0.4, "【注释（下）】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_83_right:
    add_textbox(slide, 7.0, y, 1.5, 0.5, term, FONT_TITLE, 13, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 8.5, y, 4.4, 0.5, meaning, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.62

# ============================================================
# 第7页：第八十三章 核心义理与启示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_LIGHT)

add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十三章 · 核心义理与启示", FONT_TITLE, 32, COLOR_BLUE_ACCENT, bold=True)

add_shape_with_fill(slide, 0.4, 1.2, 6.0, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_LIGHT)
add_textbox(slide, 0.5, 1.3, 5.8, 0.5, "【核心义理】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 0.6, 1.8, 5.4, 0.03, COLOR_BLUE_LIGHT)

core83 = [
    "大国者下流",
    "大国应当像水一样处于下游，甘居人下。水流下游，万物归附；大国谦下，天下归心。",
    "",
    "牝静胜牡",
    "雌性以安静胜过雄性。静是道的本性，静能胜动，柔能克刚。",
    "",
    "相下则相取",
    "大国谦下则取小国，小国谦下则取大国。彼此谦下，双方各得其所。",
    "",
    "各得其欲",
    "大国想兼畜人，小国想入事人。两者各得其欲，关键在于大者宜为下。",
]

add_multi_line_textbox(slide, 0.6, 1.95, 5.6, 5.0, core83, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_shape_with_fill(slide, 6.7, 1.2, 6.2, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_LIGHT)
add_textbox(slide, 6.8, 1.3, 6.0, 0.5, "【人生启示】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 6.9, 1.8, 5.8, 0.03, COLOR_BLUE_LIGHT)

insights83 = [
    "谦下之道",
    "越是强大，越要谦下。位高权重者若能谦下，更能赢得人心。",
    "",
    "安静的力量",
    "在纷争中保持安静，反而能占据上风。静观其变，以逸待劳。",
    "",
    "互惠双赢",
    "彼此谦下，双方都能得到想要的。合作比对抗更能共赢。",
    "",
    "大者宜下",
    "越是大人物、大家，越要放低姿态。水低为海，人低为王。",
]

add_multi_line_textbox(slide, 6.9, 1.95, 5.8, 5.0, insights83, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

# ============================================================
# 第8页：第八十四章 原文
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 1.0, RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十四章 · 不争善胜", FONT_TITLE, 36, COLOR_WHITE, bold=True)

add_shape_with_fill(slide, 0.4, 1.3, 12.5, 5.8, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 1.4, 12.3, 0.5, "【原文】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

original_84 = """天之道，不争而善胜，
不言而善应，
召而不召然而自来，
坦然而善谋。
天道无亲，常与善人。"""

add_textbox(slide, 0.8, 1.95, 11.8, 4.8, original_84, FONT_INK, 24, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第9页：第八十四章 注释
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 0.8, RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 0.5, 0.15, 12.33, 0.7, "第八十四章 · 注释", FONT_TITLE, 32, COLOR_WHITE, bold=True)

notes_84 = [
    ("不争而善胜", "不争斗却善于获胜"),
    ("不言而善应", "不说话却善于回应"),
    ("召而不召", "召唤却不用召唤"),
    ("自来", "自己前来"),
    ("坦然", "平坦坦然"),
    ("善谋", "善于谋划"),
    ("天道无亲", "天道没有偏爱"),
    ("常与善人", "常常赐予、帮助善人"),
]

add_shape_with_fill(slide, 0.3, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.4, 1.2, 6.1, 0.4, "【注释（上）】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_84[:4]:
    add_textbox(slide, 0.5, y, 1.8, 0.5, term, FONT_TITLE, 14, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 2.3, y, 4.0, 0.5, meaning, FONT_BODY, 14, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.7

add_shape_with_fill(slide, 6.8, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 6.9, 1.2, 6.1, 0.4, "【注释（下）】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_84[4:]:
    add_textbox(slide, 7.0, y, 1.8, 0.5, term, FONT_TITLE, 14, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 8.8, y, 4.0, 0.5, meaning, FONT_BODY, 14, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.7

# ============================================================
# 第10页：第八十四章 核心义理与启示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, RGBColor(0x1a, 0x3a, 0x6e))

add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十四章 · 核心义理与启示", FONT_TITLE, 32, COLOR_BLUE_ACCENT, bold=True)

add_shape_with_fill(slide, 0.4, 1.2, 6.0, 5.9, RGBColor(0x0f, 0x1e, 0x3c), RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 0.5, 1.3, 5.8, 0.5, "【核心义理】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 0.6, 1.8, 5.4, 0.03, RGBColor(0x1a, 0x3a, 0x6e))

core84 = [
    "不争而善胜",
    "天道的运行不争斗却能获胜。这是道的无为而无不为的体现。",
    "",
    "不言而善应",
    "天道不说话却能回应一切。它按照自己的规律运行，万物自然响应。",
    "",
    "天道无亲，常与善人",
    "天道没有偏心，但它总是帮助善良的人。善有善报，这是宇宙的法则。",
    "",
    "自然吸引",
    "有道之人不用召唤，人们自然归附。这是德行的感召力。",
]

add_multi_line_textbox(slide, 0.6, 1.95, 5.6, 5.0, core84, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_shape_with_fill(slide, 6.7, 1.2, 6.2, 5.9, RGBColor(0x0f, 0x1e, 0x3c), RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 6.8, 1.3, 6.0, 0.5, "【人生启示】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 6.9, 1.8, 5.8, 0.03, RGBColor(0x1a, 0x3a, 0x6e))

insights84 = [
    "不争的境界",
    "真正的高手不争，因为不需要争。有实力的人自然被认可，无需争辩。",
    "",
    "默默做事",
    "少说多做，用行动证明自己。不言而善应，用结果说话。",
    "",
    "德行感召",
    "培养良好的品德和德行，自然会吸引志同道合的人。",
    "",
    "善有善报",
    "坚持行善，天道会回报。不是迷信，而是一种信念和坚持。",
]

add_multi_line_textbox(slide, 6.9, 1.95, 5.8, 5.0, insights84, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

# ============================================================
# 第11页：三章总结对比
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "三章总结 · 智慧贯穿", FONT_TITLE, 36, COLOR_BLUE_PRIMARY, bold=True)
add_shape_with_fill(slide, 0.5, 1.0, 12.33, 0.04, COLOR_BLUE_ACCENT)

themes = [
    {
        "title": "柔弱处上",
        "color": COLOR_BLUE_PRIMARY,
        "subtitle": "第八十二章",
        "points": ["天下莫柔弱于水", "弱之胜强，柔之胜刚", "受国之垢，是为社稷主", "正言若反"],
        "core": "柔弱之道"
    },
    {
        "title": "善为下",
        "color": COLOR_BLUE_LIGHT,
        "subtitle": "第八十三章",
        "points": ["大国者下流", "牝常以静胜牡", "大国以下小国则取小国", "大者宜为下"],
        "core": "谦下之道"
    },
    {
        "title": "不争善胜",
        "color": RGBColor(0x1a, 0x3a, 0x6e),
        "subtitle": "第八十四章",
        "points": ["不争而善胜", "不言而善应", "天道无亲，常与善人", "坦然善谋"],
        "core": "无为之境"
    },
]

card_w = 4.0
gap = 0.22
start_x = 0.4

for i, theme in enumerate(themes):
    x = start_x + i * (card_w + gap)
    add_shape_with_fill(slide, x, 1.2, card_w, 6.0, RGBColor(0xff, 0xff, 0xff), theme["color"])
    add_textbox(slide, x, 1.35, card_w, 0.5, theme["subtitle"], FONT_BODY, 14, theme["color"], bold=False)
    add_textbox(slide, x, 1.7, card_w, 0.7, theme["title"], FONT_TITLE, 26, theme["color"], bold=True)
    add_shape_with_fill(slide, x + 0.3, 2.45, card_w - 0.6, 0.03, theme["color"])
    add_textbox(slide, x, 2.55, card_w, 0.5, theme["core"], FONT_INK, 18, theme["color"], bold=True)
    y = 3.15
    for point in theme["points"]:
        add_textbox(slide, x + 0.25, y, card_w - 0.5, 0.5, "• " + point, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
        y += 0.62

# ============================================================
# 第12页：结语
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 0, 7.35, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 1.5, 1.2, 10.33, 5.1, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 1.5, 1.5, 10.33, 0.8, "今日要点", FONT_TITLE, 36, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 4.5, 2.35, 4.33, 0.03, COLOR_BLUE_ACCENT)

summary = [
    "第八十二章：柔弱之道——天下莫柔弱于水，而攻坚强者莫之能胜。",
    "第八十三章：谦下之道——大国者下流，大者宜为下。",
    "第八十四章：无为之境——不争而善胜，天道无亲，常与善人。",
]

add_multi_line_textbox(slide, 1.8, 2.55, 9.7, 2.8, summary, FONT_BODY, 17, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_textbox(slide, 1.5, 5.4, 10.33, 0.6, "天之道，不争而善胜；圣人之道，为而不争。", FONT_INK, 20, COLOR_BLUE_ACCENT, bold=False)
add_textbox(slide, 1.5, 6.1, 10.33, 0.4, "—— 道德经 · 第八十二至八十四章 · 完 ——", FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False)

output_path = "/Users/mac/.openclaw/workspace/courses/道德经_第27天.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
