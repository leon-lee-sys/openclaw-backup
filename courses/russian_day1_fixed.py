#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
俄语课件Day1 - 名词第1格和第2格
使用正确俄语字母（кириллица）
"""
import os
import sys
sys.path.insert(0, '/opt/homebrew/lib/node_modules/openclaw/node_modules')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(30, 60, 114)
TEXT = RGBColor(51, 51, 51)
BLUE_LIGHT = RGBColor(70, 130, 180)
GOLD = RGBColor(180, 140, 50)
WHITE = RGBColor(255, 255, 255)

def make_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(200, 220, 255)

def add_footer(slide, day_num):
    footer = slide.shapes.add_textbox(Inches(8.5), Inches(7.2), Inches(1.3), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "День %d/84" % day_num
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(160, 160, 160)
    p.alignment = PP_ALIGN.RIGHT

def add_box(slide, text, left, top, width, height, fill_color=None, font_size=16, color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if color:
        p.font.color.rgb = color
    return box

def make_lesson(day_num, title_cn, slides_data, output_dir):
    prs = Presentation()
    
    # Cover
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    main = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = main.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语基础语法 · 第%d天" % day_num
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    sub = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    tf2 = sub.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title_cn
    p2.font.size = Pt(36)
    p2.font.color.rgb = TEXT
    p2.alignment = PP_ALIGN.CENTER
    
    add_footer(s, day_num)
    
    # Content slides
    for slide_data in slides_data:
        s = make_slide(prs)
        set_title_bar(s, slide_data["title"], slide_data.get("subtitle"))
        
        y = 1.6
        for item in slide_data["content"]:
            if isinstance(item, tuple):
                if item[0] == "header":
                    h = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
                    tf = h.text_frame
                    p = tf.paragraphs[0]
                    p.text = item[1]
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = PRIMARY
                    y += 0.4
                elif item[0] == "example":
                    add_box(s, item[1], Inches(0.5), Inches(y), Inches(9), Inches(0.55), 
                            fill_color=BLUE_LIGHT, font_size=15, color=WHITE)
                    y += 0.7
                elif item[0] == "note":
                    add_box(s, item[1], Inches(0.5), Inches(y), Inches(9), Inches(0.5),
                            fill_color=GOLD, font_size=14, color=WHITE)
                    y += 0.65
            else:
                t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
                tf = t.text_frame
                p = tf.paragraphs[0]
                p.text = item
                p.font.size = Pt(15)
                p.font.color.rgb = TEXT
                y += 0.42
        
        add_footer(s, day_num)
    
    filename = "俄语_Day%d_L1.pptx" % day_num
    path = os.path.join(output_dir, filename)
    prs.save(path)
    print("Done: %s" % filename)
    return path

output_dir = "/Users/mac/.openclaw/workspace/courses/俄语课件"
os.makedirs(output_dir, exist_ok=True)

# Day 1 slides with correct Russian (Cyrillic)
day1_slides = [
    {"title": "一、什么是第1格？", "subtitle": "Именительный падеж",
     "content": [
         ("header", "第1格 = 主格，表示«谁/什么»"),
         "第1格是名词的原形，是所有名词的基础形式。",
         ("example", "Студент (学生) — Этот студент учится. (这个学生在学习)"),
         ("example", "Книга (书) — Эта книга интересная. (这本书很有趣)"),
         ("note", "记忆方法：第1格就是词典里单词的原形！"),
     ]},
    {"title": "二、第2格的基本用法", "subtitle": "Родительный падеж",
     "content": [
         ("header", "第2格主要表示«所属»关系，相当于中文的«的»："),
         "Это книга студента. = 这是学生的书。",
         "Это дом брата. = 这是弟弟的房子。",
         ("header", "第2格还表示«来自»："),
         "Я из Китая. = 我来自中国。",
         "Она из Москвы. = 她来自莫斯科。",
         ("note", "关键语法：из + 第2格 = 来自某地"),
     ]},
    {"title": "三、名词第2格单数变化规则", "subtitle": "Склонение существительных",
     "content": [
         ("header", "第一变格法（阴性名词）："),
         "以 -а / -я 结尾的阴性名词：",
         ("example", "книга → книги (书)"),
         ("example", "девушка → девушки (姑娘)"),
         ("header", "第二变格法（阳性/中性名词）："),
         "以硬辅音结尾的阳性名词：",
         ("example", "дом → дома (房子)"),
         ("example", "студент → студента (学生)"),
     ]},
    {"title": "四、名词第2格变化示例", "subtitle": "Примеры склонения",
     "content": [
         ("header", "阳性名词："),
         "студдент → студента (学生的)",
         "друг → друга (朋友的)",
         "врач → врача (医生的)",
         ("header", "阴性名词："),
         "книга → книги (书的)",
         "девушка → девушки (姑娘的)",
         ("note", "规律：通常在词尾加 -а 或 -и"),
     ]},
    {"title": "五、数量词后的第2格", "subtitle": "Родительный после числительных",
     "content": [
         ("header", "数字2-4后名词用单数第2格："),
         "два студента (两个学生)",
         "три книги (三本书)",
         "четыре девушки (四个姑娘)",
         ("header", "数字5以上名词用复数第2格："),
         "пять студентов (五个学生)",
         ("note", "关键：1之后单数，2-4单数第2格，5+复数第2格"),
     ]},
    {"title": "六、第2格否定用法", "subtitle": "Отрицательный родительный",
     "content": [
         ("header", "否定«没有/不是»时使用第2格："),
         "У меня нет времени. = 我没有时间。",
         "Это не моя книга. = 这不是我的书。",
         ("header", "нет = 没有（现在时）"),
         "не было = 没有（过去时）",
         "не будет = 没有（将来时）",
         ("note", "结构：У + 人称 + нет + 物"),
     ]},
    {"title": "七、阅读理解练习", "subtitle": "Чтение с пониманием",
     "content": [
         "Меня зовут Анна. Я из Москвы. Я студентка университета. Мой брат тоже студент. Он из Санкт-Петербурга.",
         "",
         ("header", "请回答："),
         "1. Как её зовут? (她叫什么？)",
         "2. Откуда Анна? (她来自哪里？)",
         "3. Кто её брат? (她弟弟是做什么的？)",
         ("note", "提示：студентка = 女大学生，университет = 大学"),
     ]},
    {"title": "八、翻译练习", "subtitle": "Перевод",
     "content": [
         ("header", "中译俄练习："),
         "1. 这是我朋友的书。",
         "2. 我来自中国。",
         "3. 她没有时间。",
         "",
         ("header", "参考答案："),
         "1. Это книга моего друга.",
         "2. Я из Китая.",
         "3. У неё нет времени.",
     ]},
    {"title": "九、本课重点总结", "subtitle": "Итоги урока",
     "content": [
         ("header", "第2格核心用法："),
         "① 表示所属：чьё? → ...的...",
         "② 表示来源：из + 第2格 = 来自...",
         "③ 数字搭配：2-4用单数，5+用复数",
         "④ 否定结构：нет + 第2格",
         ("header", "今日关键词汇："),
         "дом (房子) → дома",
         "книга (书) → книги",
         "студент (学生) → студента",
         ("note", "记住：名词第2格回答чья?"),
     ]},
]

make_lesson(1, "名词的第1格和第2格", day1_slides, output_dir)

print("\nDay1课件（正确俄语）生成完成！")
