# -*- coding: utf-8 -*-
"""
道德经第21天 课件生成器
章节：第64-66章
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_yuanwen_box(slide, lines, y=1.4, height=2.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(20)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=BLUE_LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(255, 255, 255)

def add_points(slide, points, y_start=3.1):
    y = y_start
    for title, content in points:
        pt_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT
        y += 1.1
    return y

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.5))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第21课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第64-66章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "为之于未有 · 千里之行 · 出生入死"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "合抱之木，生于毫末。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = PRIMARY
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第六十四章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = SECONDARY
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, 12)

# ===== 第2页：第六十四章详解 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s2, "第六十四章 · 为之于未有")
add_yuanwen_box(s2, [
    "其安易持，其未兆易谋，",
    "其脆易泮，其微易散。",
    "为之于未有，治之于未乱。",
    "",
    "合抱之木，生于毫末；",
    "九层之台，起于累土；",
    "千里之行，始于足下。",
    "",
    "为者败之，执者失之。",
    "是以圣人无为，故无败；",
    "无执，故无失。",
    "",
    "民之从事，常于几成而败之。",
    "慎终如始，则无败事。"
], y=1.4, height=5.0)
content2 = s2.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(2))
tf_c2 = content2.text_frame
tf_c2.word_wrap = True
pc2 = tf_c2.paragraphs[0]
pc2.text = "【白话解读】安稳时容易维持，未成形时容易谋划。在事情没有发生之前就处理好，在混乱没有形成之前就治理好。合抱的大树从嫩芽生长，九层高台从泥土堆起，千里远行从脚下开始。"
pc2.font.size = Pt(14)
pc2.font.color.rgb = TEXT
add_slide_number(s2, 2, 12)

# ===== 第3页：第六十四章核心 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s3, "第六十四章 · 核心思想")
add_core_box(s3, "🌟 核心观点", "为之于未有，治之于未乱。千里之行始于足下，慎终如始则无败事。")
add_points(s3, [
    ("1. 为之于未有", "在事情没有发生前就处理好——防患于未然。"),
    ("2. 千里之行始于足下", "再大的事业也是从第一步开始——不积跬步无以至千里。"),
    ("3. 慎终如始", "快到终点时像刚开始一样谨慎——坚持到底才能成功。"),
    ("4. 无为无败", "不妄为就不失败，不强求就不失去——顺其自然才是正道。")
])
add_slide_number(s3, 3, 12)

# ===== 第4页：第六十五章详解 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s4, "第六十五章 · 玄德深矣")
add_yuanwen_box(s4, [
    "古之善为道者，",
    "非以明民，将以愚之。",
    "",
    "民之难治，以其智多。",
    "故以智治国，国之贼；",
    "不以智治国，国之福。",
    "",
    "知此两者，亦稽式。",
    "常知稽式，是谓玄德。",
    "玄德深矣，远矣，",
    "与物反矣，",
    "然后乃至大顺。"
], y=1.4, height=4.4)
content4 = s4.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(9), Inches(2))
tf_c4 = content4.text_frame
tf_c4.word_wrap = True
pc4 = tf_c4.paragraphs[0]
pc4.text = "【白话解读】古代善于行道的人，不是让人民聪明巧智，而是让人民淳朴敦厚。人民难治理是因为巧智太多。用智巧治国是国家的祸害，不用智巧治国是国家的幸福。"
pc4.font.size = Pt(14)
pc4.font.color.rgb = TEXT
add_slide_number(s4, 4, 12)

# ===== 第5页：第六十五章核心 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s5, "第六十五章 · 核心思想")
add_core_box(s5, "🌟 核心观点", "以智治国，国之贼；不以智治国，国之福。玄德深远，与物反矣。")
add_points(s5, [
    ("1. 将以愚之", "让人民淳朴敦厚而非投机取巧——返璞归真的智慧。"),
    ("2. 以智治国国之贼", "用权谋智巧治理国家，反而是祸害——巧智带来混乱。"),
    ("3. 玄德", "深远的德行——不是小聪明，是大智若愚的境界。"),
    ("4. 与物反矣", "与社会流行的价值观相反——道往往与世俗背道而驰。")
])
add_slide_number(s5, 5, 12)

# ===== 第6页：现代反思 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s6, "第六十四-六十五章 · 现代反思")
items6 = [
    ("🏗️ 风险管理", "为之于未有 → 好的风险管理是在风险发生前就预防，而不是事后补救。"),
    ("👔 领导力", "慎终如始 → 创业者最怕的就是成功后失去初心，慎终如始才能持续成功。"),
    ("📚 教育", "将以愚之 → 教育的本质不是教人变得精明，而是让人淳朴厚道。"),
    ("🎯 专注", "千里之行始于足下 → 不要想着走捷径，每一步都认真走，就是最快的路。")
]
y6 = 1.4
for icon_title, explanation in items6:
    box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y6), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf6 = box.text_frame
    tf6.word_wrap = True
    p6_1 = tf6.paragraphs[0]
    p6_1.text = icon_title
    p6_1.font.size = Pt(16)
    p6_1.font.bold = True
    p6_1.font.color.rgb = PRIMARY
    p6_2 = tf6.add_paragraph()
    p6_2.text = explanation
    p6_2.font.size = Pt(14)
    p6_2.font.color.rgb = TEXT
    y6 += 1.4
add_slide_number(s6, 6, 12)

# ===== 第7页：第六十六章详解 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s7, "第六十六章 · 天下皆谓我大")
add_yuanwen_box(s7, [
    "天下皆谓我大，似不肖。",
    "夫唯大，故似不肖。",
    "若肖，久矣其细也夫！",
    "",
    "我有三宝，持而保之：",
    "一曰慈，二曰俭，三曰不敢为天下先。",
    "慈故能勇，",
    "俭故能广，",
    "不敢为天下先故能成器长。",
    "",
    "今舍慈且勇，",
    "舍俭且广，",
    "舍后且先，死矣！"
], y=1.4, height=4.4)
content7 = s7.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(9), Inches(2))
tf_c7 = content7.text_frame
tf_c7.word_wrap = True
pc7 = tf_c7.paragraphs[0]
pc7.text = "【白话解读】天下人都说我伟大，不像任何具体的东西。正因为伟大，所以不像。如果像，早就渺小了！我有三件宝贝：一是慈爱，二是俭朴，三是不敢走在天下人前面。"
pc7.font.size = Pt(14)
pc7.font.color.rgb = TEXT
add_slide_number(s7, 7, 12)

# ===== 第8页：第六十六章核心 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s8, "第六十六章 · 核心思想")
add_core_box(s8, "🌟 核心观点", "我有三宝：慈、俭、不敢为天下先。慈故能勇，俭故能广，不敢为天下先故能成器长。")
add_points(s8, [
    ("1. 慈", "慈爱是根本——慈爱才能真正勇敢，不是匹夫之勇。"),
    ("2. 俭", "俭朴才能宽广——不挥霍才能长久，俭才能有余。"),
    ("3. 不敢为天下先", "不争先反而能成为领导者——甘居人下才能领导人心。"),
    ("4. 舍之死矣", "舍弃这三宝而追求勇、广、先，就走向死亡了。")
])
add_slide_number(s8, 8, 12)

# ===== 第9页：西方思想 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s9, "🌍 西方思想对照")
phils = [
    ("亚里士多德 · 中庸之道", "美德处于两个极端之间——勇敢在卤莽和怯懦之间，俭朴在奢侈和吝啬之间。约前350年"),
    ("爱比克泰德 · 斯多葛", "真正的自由是控制自己的欲望，不被外物支配。约100年"),
    ("塞缪尔·斯迈尔斯 · 自助", "自助者天助之——自己的成功靠自己的努力。1859年"),
    ("乔布斯 · 专注", "专注意味着对成百上千的好主意说不。2011年")
]
y9 = 1.4
for name_theory, desc_year in phils:
    box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y9), Inches(9), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf9 = box.text_frame
    tf9.word_wrap = True
    p9_1 = tf9.paragraphs[0]
    p9_1.text = name_theory
    p9_1.font.size = Pt(16)
    p9_1.font.bold = True
    p9_1.font.color.rgb = PRIMARY
    p9_2 = tf9.add_paragraph()
    p9_2.text = desc_year
    p9_2.font.size = Pt(14)
    p9_2.font.color.rgb = TEXT
    y9 += 1.3
add_slide_number(s9, 9, 12)

# ===== 第10页：现代应用 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s10, "📌 现代应用")
add_points(s10, [
    ("自我管理", "慈故能勇 → 心中有爱才能有真正的勇气，不是逞强斗狠。"),
    ("理财观念", "俭故能广 → 收入再高，不挥霍才能积累财富，才能宽广从容。"),
    ("职场发展", "不敢为天下先 → 不急于出风头，先倾听再发言，反而更容易被尊重。"),
    ("创业哲学", "慎终如始 → 成功时保持创业初期的心态，不忘初心。")
])
add_slide_number(s10, 10, 12)

# ===== 第11页：中国典故 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s11, "📜 中国典故：范蠡三散财富")
box11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.5))
box11.fill.solid()
box11.fill.fore_color.rgb = BG
tf11 = box11.text_frame
tf11.word_wrap = True
p11_1 = tf11.paragraphs[0]
p11_1.text = "【史记·越王勾践世家】"
p11_1.font.size = Pt(16)
p11_1.font.bold = True
p11_1.font.color.rgb = PRIMARY
p11_2 = tf11.add_paragraph()
p11_2.text = "范蠡辅佐越王勾践灭吴复国后，主动辞职经商，三次散尽家财又三次重新富有。他居住于定陶，自称陶朱公，认为赚取财富是用来救助困厄之人，而非个人享受。他的儿子们却舍不得钱财，最终招来灾祸。"
p11_2.font.size = Pt(15)
p11_2.font.color.rgb = TEXT
p11_2.space_before = Pt(10)
insight11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.1), Inches(9), Inches(1))
insight11.fill.solid()
insight11.fill.fore_color.rgb = BLUE_LIGHT
tfi11 = insight11.text_frame
tfi11.word_wrap = True
pi11 = tfi11.paragraphs[0]
pi11.text = "💡 启示"
pi11.font.size = Pt(18)
pi11.font.bold = True
pi11.font.color.rgb = RGBColor(255, 255, 255)
pi11_2 = tfi11.add_paragraph()
pi11_2.text = "范蠡\"俭故能广\"——不把财富当私产，俭朴不奢，才能真正宽广从容。三散家财反而让他更有境界。\"俭\"不仅是节俭，更是心中有度。\""
pi11_2.font.size = Pt(14)
pi11_2.font.color.rgb = RGBColor(255, 255, 255)
app11 = s11.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2))
tfa11 = app11.text_frame
tfa11.word_wrap = True
pa11 = tfa11.paragraphs[0]
pa11.text = "【现代应用】财富：赚得多不如花得值，俭朴是持续的智慧。教育：慈爱是给孩子最好的礼物，不是物质。职场：不敢为天下先，先倾听再行动，反而更容易成功。"
pa11.font.size = Pt(14)
pa11.font.color.rgb = TEXT
add_slide_number(s11, 11, 12)

# ===== 第12页：总结 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s12, "第21课 · 三章精华总结")
sums = [
    ("第六十四章", "为之于未有", "合抱之木生于毫末，千里之行始于足下，慎终如始则无败事。"),
    ("第六十五章", "玄德深矣", "以智治国国之贼，返璞归真是真正的智慧。"),
    ("第六十六章", "我有三宝", "慈故能勇，俭故能广，不敢为天下先故能成器长。")
]
y12 = 1.4
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    y12 += 1.6
final = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(9), Inches(1))
final.fill.solid()
final.fill.fore_color.rgb = PRIMARY
tff = final.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "我有三宝，持而保之。"
pf1.font.size = Pt(20)
pf1.font.bold = True
pf1.font.color.rgb = RGBColor(255, 255, 255)
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "—— 道德经第六十六章"
pf2.font.size = Pt(14)
pf2.font.color.rgb = RGBColor(200, 200, 200)
pf2.alignment = PP_ALIGN.CENTER
add_slide_number(s12, 12, 12)

prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第21天.pptx')
print("Done: 道德经_第21天.pptx (12 slides)")
