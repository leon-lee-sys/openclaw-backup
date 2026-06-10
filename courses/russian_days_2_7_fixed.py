#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
俄语课件Day2-7 - 批量生成（正确俄语字母）
"""
import os
import sys
sys.path.insert(0, '/opt/homebrew/lib/node_modules/openclaw/node_modules')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(30, 60, 114)
TEXT = RGBColor(51, 51, 51)
BLUE_LIGHT = RGBColor(70, 130, 180)
GOLD = RGBColor(180, 140, 50)
WHITE = RGBColor(255, 255, 255)

def make_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(200, 220, 255)

def add_footer(slide, day_num):
    footer = slide.shapes.add_textbox(Inches(8.5), Inches(7.2), Inches(1.3), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "День %d/84" % day_num
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(160, 160, 160)
    p.alignment = PP_ALIGN.RIGHT

def add_box(slide, text, left, top, width, height, fill_color=None, font_size=16, color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if color:
        p.font.color.rgb = color
    return box

def make_lesson(day_num, title_cn, slides_data, output_dir):
    prs = Presentation()
    
    # Cover
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    main = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = main.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语基础语法 · 第%d天" % day_num
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    sub = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    tf2 = sub.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title_cn
    p2.font.size = Pt(36)
    p2.font.color.rgb = TEXT
    p2.alignment = PP_ALIGN.CENTER
    
    add_footer(s, day_num)
    
    for slide_data in slides_data:
        s = make_slide(prs)
        set_title_bar(s, slide_data["title"], slide_data.get("subtitle"))
        
        y = 1.6
        for item in slide_data["content"]:
            if isinstance(item, tuple):
                if item[0] == "header":
                    h = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
                    tf = h.text_frame
                    p = tf.paragraphs[0]
                    p.text = item[1]
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = PRIMARY
                    y += 0.4
                elif item[0] == "example":
                    add_box(s, item[1], Inches(0.5), Inches(y), Inches(9), Inches(0.55), 
                            fill_color=BLUE_LIGHT, font_size=15, color=WHITE)
                    y += 0.7
                elif item[0] == "note":
                    add_box(s, item[1], Inches(0.5), Inches(y), Inches(9), Inches(0.5),
                            fill_color=GOLD, font_size=14, color=WHITE)
                    y += 0.65
            else:
                t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
                tf = t.text_frame
                p = tf.paragraphs[0]
                p.text = item
                p.font.size = Pt(15)
                p.font.color.rgb = TEXT
                y += 0.42
        
        add_footer(s, day_num)
    
    filename = "俄语_Day%d_L1.pptx" % day_num
    path = os.path.join(output_dir, filename)
    prs.save(path)
    print("Done: %s" % filename)
    return path

output_dir = "/Users/mac/.openclaw/workspace/courses/俄语课件"
os.makedirs(output_dir, exist_ok=True)

# Day 2: 名词第3格
day2_slides = [
    {"title": "一、第3格的意义", "subtitle": "Дательный падеж",
     "content": [
         ("header", "第3格 = 间接宾语，表示«给谁/为谁»"),
         "Дайте книгу студенту. = 把书给学生。",
         ("example", "Я помогаю другу. = 我帮助朋友。"),
         ("note", "第3格回答问题：кому? чему? (给谁？向什么？)"),
     ]},
    {"title": "二、第3格变化规则", "subtitle": "Склонение",
     "content": [
         ("header", "阳性名词："),
         "студдент → студенту (给学生)",
         "друг → другу (给朋友)",
         ("header", "阴性名词："),
         "книга → книге (给书)",
         "девушка → девушке (给姑娘)",
         ("note", "规律：词尾加 -у 或 -е"),
     ]},
    {"title": "三、人称代词第3格", "subtitle": "Местоимения",
     "content": [
         ("header", "я → мне, ты → тебе, он → ему"),
         ("header", "она → ей, мы → нам, вы → вам"),
         ("example", "Они помогают мне. = 他们帮助我。"),
         ("example", "Я даю книгу ей. = 我把书给她。"),
         ("note", "记忆：人称代词第3格需要单独记忆！"),
     ]},
    {"title": "四、课后练习", "subtitle": "Упражнения",
     "content": [
         "1. 将下列名词变为第3格：",
         "   дом, книга, друг, студдент, девушка",
         "",
         "2. 用人称代词填空：",
         "   Я даю книгу ___ (он).",
         "   Он помогает ___ (она).",
         "",
         "3. 中译俄：",
         "   我把书给朋友。",
         "   妈妈帮助女儿。",
     ]},
]

# Day 3: 名词第4格
day3_slides = [
    {"title": "一、第4格的意义", "subtitle": "Винительный падеж",
     "content": [
         ("header", "第4格 = 直接宾语，表示«做什么/把谁»"),
         "Я читаю книгу. = 我读书。",
         ("example", "Он видит друга. = 他看见朋友。"),
         ("note", "第4格回答问题：кого? что? (谁？什么？)"),
     ]},
    {"title": "二、第4格变化规则", "subtitle": "Склонение",
     "content": [
         ("header", "阳性名词（活动物）："),
         "студдент → студента (看见学生)",
         "друг → друга (看见朋友)",
         ("header", "阳性名词（非活动物）："),
         "дом → дом (看见房子) — 同第1格",
         ("header", "阴性名词："),
         "книга → книгу (看见书)",
         ("note", "重点：活动物阳性第4格同第2格！"),
     ]},
    {"title": "三、课后练习", "subtitle": "Упражнения",
     "content": [
         "1. 将下列名词变为第4格：",
         "   дом, книга, друг, студдент",
         "",
         "2. 中译俄：",
         "   我在读书。",
         "   他看见了朋友。",
         "   妈妈在做汤。",
     ]},
]

# Day 4: 名词第5格
day4_slides = [
    {"title": "一、第5格的意义", "subtitle": "Творительный падеж",
     "content": [
         ("header", "第5格 = 工具格，表示«用什么/和谁一起»"),
         "Я пишу ручкой. = 我用笔写字。",
         ("example", "Он идёт с другом. = 他和朋友一起走。"),
         ("note", "第5格回答问题：кем? чем? (用谁？用什么？)"),
     ]},
    {"title": "二、第5格变化规则", "subtitle": "Склонение",
     "content": [
         ("header", "阳性/中性名词："),
         "студдент → студентом (用学生)",
         "окно → окном (用窗户)",
         ("header", "阴性名词："),
         "книга → книгой (用书)",
         "девушка → девушкой (和姑娘一起)",
         ("note", "规律：加 -ом 或 -ой"),
     ]},
    {"title": "三、课后练习", "subtitle": "Упражнения",
     "content": [
         "1. 将下列名词变为第5格：",
         "   дом, книга, друг, окно",
         "",
         "2. 中译俄：",
         "   我用笔写字。",
         "   她和朋友一起走。",
         "   他站在窗边。",
     ]},
]

# Day 5: 名词第6格
day5_slides = [
    {"title": "一、第6格的意义", "subtitle": "Предложный падеж",
     "content": [
         ("header", "第6格 = 前置词格，必须与前置词连用"),
         ("header", "常用前置词：в / на (在...里/上)"),
         ("example", "Книга на столе. = 书在桌子上。"),
         ("example", "Я в университете. = 我在大学里。"),
         ("note", "第6格回答问题：о ком? о чём? (关于谁？关于什么？)"),
     ]},
    {"title": "二、第6格变化规则", "subtitle": "Склонение",
     "content": [
         ("header", "阳性/中性名词："),
         "студдент → в университете (在大学里)",
         "окно → на окне (在窗户上)",
         ("header", "阴性名词："),
         "книга → в книге (在书里)",
         "девушка → о девушке (关于姑娘)",
         ("note", "规律：加 -е 或 -и"),
     ]},
    {"title": "三、六格综合练习", "subtitle": "Упражнения",
     "content": [
         "1. 将下列名词变为各格（单数）：",
         "   дом, книга, студдент, окно",
         "",
         "2. 用适当的格填空：",
         "   Я иду ___ (дом).",
         "   Книга лежит ___ (стол).",
         "   Я помогаю ___ (друг).",
     ]},
]

# Day 6: 形容词性数配合（一）
day6_slides = [
    {"title": "一、形容词性数配合原则", "subtitle": "Согласование прилагательных",
     "content": [
         ("header", "形容词必须与名词在性、数上一致："),
         ("header", "阳性：большой дом (大房子)"),
         ("header", "阴性：большая книга (大书)"),
         ("header", "中性：большое окно (大窗户)"),
         ("note", "记住：性数一致是俄语形容词的核心规则！"),
     ]},
    {"title": "二、阳性形容词词尾", "subtitle": "Окончания",
     "content": [
         ("header", "以 -ый / -ой 结尾的形容词："),
         "новый (新的) → нового, новому...",
         "большой (大的) → большого, большому...",
         ("header", "以 -ний 结尾的形容词："),
         "хороший (好的) → хорошего, хорошему...",
         ("note", "规律：阳性形容词回答Какой? (什么样的？)"),
     ]},
    {"title": "三、阴性形容词词尾", "subtitle": "Женский род",
     "content": [
         ("header", "以 -ая / -яя 结尾的形容词："),
         "новый → новая книга (新书)",
         "хороший → хорошая девушка (好姑娘)",
         ("note", "规律：阴性形容词回答Какая? (什么样的？)"),
     ]},
    {"title": "四、课后练习", "subtitle": "Упражнения",
     "content": [
         "1. 将形容词变为适当形式：",
         "   большой ___ (дом)",
         "   новая ___ (книга)",
         "   хороший ___ (друг)",
         "",
         "2. 中译俄：",
         "   一本新书",
         "   一个好学生",
         "   一扇大窗户",
     ]},
]

# Day 7: 形容词性数配合（二）
day7_slides = [
    {"title": "一、中性形容词词尾", "subtitle": "Средний род",
     "content": [
         ("header", "以 -ое / -ее 结尾的形容词："),
         "новое окно (新窗户)",
         "большое здание (大建筑)",
         ("header", "以 -нее 结尾的形容词："),
         "хорошее настроение (好心情)",
         ("note", "规律：中性形容词回答Какое? (什么样的？)"),
     ]},
    {"title": "二、形容词短尾形式", "subtitle": "Краткая форма",
     "content": [
         ("header", "有些形容词有短尾形式（作谓语用）："),
         "Он красивый. → Он красив. (他帅)",
         "Она красивая. → Она красива. (她漂亮)",
         ("header", "短尾形容词用在句子中作表语："),
         ("example", "Этот дом большой. → Этот дом болен. (这房子大)"),
         ("note", "注意：短尾形容词只用作谓语！"),
     ]},
    {"title": "三、综合练习", "subtitle": "Упражнения",
     "content": [
         "1. 将形容词变为适当形式（性数配合）：",
         "   большой ___ (дом, окно, комната)",
         "   новый ___ (книга, письмо, студдент)",
         "",
         "2. 选用适当形容词填空：",
         "   ___ погода (хороший/хорошая/хорошее)",
         "   ___ студдент (умный/умная/умное)",
     ]},
]

# Generate all lessons
make_lesson(2, "名词第3格（给谁）", day2_slides, output_dir)
make_lesson(3, "名词第4格（做什么）", day3_slides, output_dir)
make_lesson(4, "名词第5格（用什么）", day4_slides, output_dir)
make_lesson(5, "名词第6格（关于什么）", day5_slides, output_dir)
make_lesson(6, "形容词性数配合（一）", day6_slides, output_dir)
make_lesson(7, "形容词性数配合（二）", day7_slides, output_dir)

print("\nDay2-7课件（正确俄语）生成完成！")
