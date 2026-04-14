from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

# ===== 第1页：封面 =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar.line.fill.background()
txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第六课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

main_title = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main_title.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第19-21章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER

subtitle = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = subtitle.text_frame
p3 = tf3.paragraphs[0]
p3.text = "绝圣弃智 · 归根复命 · 众里寻他"
p3.font.size = Pt(28)
p3.font.color.rgb = RGBColor(70, 130, 180)
p3.alignment = PP_ALIGN.CENTER

quote_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_q = quote_box.text_frame
tf_q.word_wrap = True
p_q1 = tf_q.paragraphs[0]
p_q1.text = "绝圣弃智，民利百倍。"
p_q1.font.size = Pt(22)
p_q1.font.italic = True
p_q1.font.color.rgb = RGBColor(30, 60, 114)
p_q1.alignment = PP_ALIGN.CENTER
p_q2 = tf_q.add_paragraph()
p_q2.text = "—— 《道德经》第十九章"
p_q2.font.size = Pt(14)
p_q2.font.color.rgb = RGBColor(150, 150, 150)
p_q2.alignment = PP_ALIGN.CENTER

add_slide_number(slide1, 1, 10)

# ===== 第2页：第十九章详解 =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar2.fill.solid()
title_bar2.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar2.line.fill.background()
txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf5 = txBox2.text_frame
p5 = tf5.paragraphs[0]
p5.text = "第十九章 · 绝圣弃智"
p5.font.size = Pt(32)
p5.font.bold = True
p5.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.4))
yuanwen_box.fill.solid()
yuanwen_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw = yuanwen_box.text_frame
tf_yw.word_wrap = True
p_yw = tf_yw.paragraphs[0]
p_yw.text = "【原文】"
p_yw.font.size = Pt(18)
p_yw.font.bold = True
p_yw.font.color.rgb = RGBColor(30, 60, 114)
for line in ["绝圣弃智，民利百倍；", "绝仁弃义，民复孝慈；", "绝巧弃利，盗贼无有。"]:
    p = tf_yw.add_paragraph()
    p.text = line
    p.font.size = Pt(19)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(3.2))
tf_c = content_box.text_frame
tf_c.word_wrap = True
p_c1 = tf_c.paragraphs[0]
p_c1.text = "【白话解读】"
p_c1.font.size = Pt(18)
p_c1.font.bold = True
p_c1.font.color.rgb = RGBColor(30, 60, 114)
解读 = [
    ("绝圣弃智", "抛弃圣贤的帽子和智巧的算计——减少人为折腾，百姓反而得到更大利益。"),
    ("绝仁弃义", "放下仁义的标榜和道德的说教——家庭自然恢复孝慈本色，不必刻意提倡。"),
    ("绝巧弃利", "断绝投机取巧和追逐私利——盗贼失去生存土壤，自然消失。"),
    ("核心：返璞归真", "三个\"绝\"字，都是让社会回归简单、自然、淳朴的状态。"),
]
for title, content in 解读:
    p = tf_c.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(5)

add_slide_number(slide2, 2, 10)

# ===== 第3页：第十九章生活智慧 =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar3.fill.solid()
title_bar3.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar3.line.fill.background()
txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf6 = txBox3.text_frame
p6 = tf6.paragraphs[0]
p6.text = "第十九章 · 生活智慧"
p6.font.size = Pt(32)
p6.font.bold = True
p6.font.color.rgb = RGBColor(255, 255, 255)

core_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.3))
core_box.fill.solid()
core_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_core = core_box.text_frame
tf_core.word_wrap = True
p_core = tf_core.paragraphs[0]
p_core.text = "🌟 核心智慧"
p_core.font.size = Pt(18)
p_core.font.bold = True
p_core.font.color.rgb = RGBColor(255, 255, 255)
p_core2 = tf_core.add_paragraph()
p_core2.text = "少则得，多则惑——做减法的人生，反而更丰盛。"
p_core2.font.size = Pt(16)
p_core2.font.color.rgb = RGBColor(255, 255, 255)

points = [
    ("1. 断舍离的哲学", "不是清贫，而是精简。减少杂物、应酬、信息的干扰，心灵才能呼吸。"),
    ("2. 教育孩子的启示", "不刻意\"鸡娃\"、不填鸭式教学——让孩子自然成长，反而更有创造力。"),
    ("3. 企业管理的借鉴", "减少KPI表格、减少审批流程——流程越简单，执行越高效。")
]
y = 2.9
for title, content in points:
    pt_box = slide3.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
    tf_pt = pt_box.text_frame
    tf_pt.word_wrap = True
    p_pt1 = tf_pt.paragraphs[0]
    p_pt1.text = title
    p_pt1.font.size = Pt(18)
    p_pt1.font.bold = True
    p_pt1.font.color.rgb = RGBColor(30, 60, 114)
    p_pt2 = tf_pt.add_paragraph()
    p_pt2.text = content
    p_pt2.font.size = Pt(14)
    p_pt2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.2

add_slide_number(slide3, 3, 10)

# ===== 第4页：第二十章详解 =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar4.fill.solid()
title_bar4.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar4.line.fill.background()
txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf7 = txBox4.text_frame
p7 = tf7.paragraphs[0]
p7.text = "第二十章 · 归真之旅"
p7.font.size = Pt(32)
p7.font.bold = True
p7.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.2))
yuanwen_box2.fill.solid()
yuanwen_box2.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw2 = yuanwen_box2.text_frame
tf_yw2.word_wrap = True
p_yw2 = tf_yw2.paragraphs[0]
p_yw2.text = "【原文】"
p_yw2.font.size = Pt(18)
p_yw2.font.bold = True
p_yw2.font.color.rgb = RGBColor(30, 60, 114)
for line in ["绝学无忧。唯之与阿，相去几何？", "众人熙熙，如享太牢，如春登台。", "我独泊兮，其未兆，如婴儿之未孩。"]:
    p = tf_yw2.add_paragraph()
    p.text = line
    p.font.size = Pt(19)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box2 = slide4.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(4))
tf_c2 = content_box2.text_frame
tf_c2.word_wrap = True
p_c4 = tf_c2.paragraphs[0]
p_c4.text = "【白话解读】"
p_c4.font.size = Pt(18)
p_c4.font.bold = True
p_c4.font.color.rgb = RGBColor(30, 60, 114)
解读2 = [
    ("绝学无忧", "抛弃世俗之学，就没有忧愁——脱离比较，内心平静。"),
    ("唯之与阿", "恭敬的\"唯\"与怠慢的\"阿\"，差别有多大？——世人的尊卑荣辱，皆是虚荣。"),
    ("众人熙熙", "众人都兴高采烈，如赴盛宴，如春日登台——外在的热闹，终究是过眼云烟。"),
    ("我独泊未兆", "唯独我淡泊宁静，没有显露任何征兆——像婴儿还不会笑，返璞归真。"),
]
for title, content in 解读2:
    p = tf_c2.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(5)

add_slide_number(slide4, 4, 10)

# ===== 第5页：典故：老子的孤独 =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar5.fill.solid()
title_bar5.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar5.line.fill.background()
txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf8 = txBox5.text_frame
p8 = tf8.paragraphs[0]
p8.text = "📜 典故：老子的出关之谜"
p8.font.size = Pt(32)
p8.font.bold = True
p8.font.color.rgb = RGBColor(255, 255, 255)

story_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
story_box.fill.solid()
story_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_story = story_box.text_frame
tf_story.word_wrap = True
p_s1 = tf_story.paragraphs[0]
p_s1.text = "【史记·老子韩非列传】"
p_s1.font.size = Pt(16)
p_s1.font.bold = True
p_s1.font.color.rgb = RGBColor(30, 60, 114)
p_s2 = tf_story.add_paragraph()
p_s2.text = "老子见周朝衰落，遂骑青牛出关，西游隐居。关令尹喜望见紫气东来，知有圣人至，强留老子著书，遂成《道德经》五千言。传说他最后不知所终，留下无尽神秘。"
p_s2.font.size = Pt(15)
p_s2.font.color.rgb = RGBColor(51, 51, 51)
p_s2.space_before = Pt(8)
p_s3 = tf_story.add_paragraph()
p_s3.text = "老子的出关，正是\"绝学无忧\"的实践——不求闻达，不为世累，独来独往。"
p_s3.font.size = Pt(15)
p_s3.font.color.rgb = RGBColor(51, 51, 51)

west_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.6), Inches(9), Inches(1))
west_box.fill.solid()
west_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_west = west_box.text_frame
tf_west.word_wrap = True
p_west = tf_west.paragraphs[0]
p_west.text = "🌍 西方对照"
p_west.font.size = Pt(18)
p_west.font.bold = True
p_west.font.color.rgb = RGBColor(255, 255, 255)
p_west2 = tf_west.add_paragraph()
p_west2.text = "卢梭：\"我独异于人，而贵食母。\"——与老子一样，追求与众不同的精神独立。"
p_west2.font.size = Pt(14)
p_west2.font.color.rgb = RGBColor(200, 200, 200)

app_box = slide5.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2.5))
tf_app = app_box.text_frame
tf_app.word_wrap = True
p_a1 = tf_app.paragraphs[0]
p_a1.text = "【人生启示】"
p_a1.font.size = Pt(16)
p_a1.font.bold = True
p_a1.font.color.rgb = RGBColor(30, 60, 114)
apps = [
    "• 不必活在他人的眼光里——\"相去几何\"，尊卑荣辱皆虚妄",
    "• 热闹是他们的，你有自己的节奏——独泊未兆，如婴儿之未孩",
    "• 真正的成熟，是内心的宁静——能在喧嚣中保持淡泊"
]
for app in apps:
    p = tf_app.add_paragraph()
    p.text = app
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide5, 5, 10)

# ===== 第6页：第二十一章详解 =====
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar6.fill.solid()
title_bar6.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar6.line.fill.background()
txBox6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf9 = txBox6.text_frame
p9 = tf9.paragraphs[0]
p9.text = "第二十一章 · 众里寻他"
p9.font.size = Pt(32)
p9.font.bold = True
p9.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box3 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
yuanwen_box3.fill.solid()
yuanwen_box3.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw3 = yuanwen_box3.text_frame
tf_yw3.word_wrap = True
p_yw3 = tf_yw3.paragraphs[0]
p_yw3.text = "【原文】"
p_yw3.font.size = Pt(18)
p_yw3.font.bold = True
p_yw3.font.color.rgb = RGBColor(30, 60, 114)
for line in ["孔德之容，惟道是从。", "道之为物，惟恍惟惚。", "惚兮恍兮，其中有象；恍兮惚兮，其中有物。"]:
    p = tf_yw3.add_paragraph()
    p.text = line
    p.font.size = Pt(19)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box3 = slide6.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(4))
tf_c3 = content_box3.text_frame
tf_c3.word_wrap = True
p_c7 = tf_c3.paragraphs[0]
p_c7.text = "【白话解读】"
p_c7.font.size = Pt(18)
p_c7.font.bold = True
p_c7.font.color.rgb = RGBColor(30, 60, 114)
解读3 = [
    ("孔德之容", "大德之人的举止——\"孔\"是大，德是得道后的品质。"),
    ("惟道是从", "完全追随道——最大的德行，就是顺从道的规律。"),
    ("道之为物", "道这个东西——它恍惚无形（无具体形象），却真实存在。"),
    ("其中有象/有物", "恍惚之中，有形象可见；有实体可触——道虽无形，却主宰一切。"),
]
for title, content in 解读3:
    p = tf_c3.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(5)

add_slide_number(slide6, 6, 10)

# ===== 第7页：第二十一章·孔德之容 =====
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar7.fill.solid()
title_bar7.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar7.line.fill.background()
txBox7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf10 = txBox7.text_frame
p10 = tf10.paragraphs[0]
p10.text = "📜 典故：道不远人"
p10.font.size = Pt(32)
p10.font.bold = True
p10.font.color.rgb = RGBColor(255, 255, 255)

story_box2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.2))
story_box2.fill.solid()
story_box2.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_st2 = story_box2.text_frame
tf_st2.word_wrap = True
p_st2a = tf_st2.paragraphs[0]
p_st2a.text = "【王阳明·龙场悟道】"
p_st2a.font.size = Pt(16)
p_st2a.font.bold = True
p_st2a.font.color.rgb = RGBColor(30, 60, 114)
p_st2b = tf_st2.add_paragraph()
p_st2b.text = "王阳明被贬贵州龙场，在困顿中悟道：\"圣人之道，吾性自足，向之求理于事物者误也。\"——道不在外物，而在内心。正如老子所说：\"孔德之容，惟道是从\"，真正的德行，是内心有道。"
p_st2b.font.size = Pt(15)
p_st2b.font.color.rgb = RGBColor(51, 51, 51)
p_st2b.space_before = Pt(8)

west_box2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(9), Inches(1))
west_box2.fill.solid()
west_box2.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_w2 = west_box2.text_frame
tf_w2.word_wrap = True
p_w2a = tf_w2.paragraphs[0]
p_w2a.text = "⚡ 现代启示"
p_w2a.font.size = Pt(18)
p_w2a.font.bold = True
p_w2a.font.color.rgb = RGBColor(255, 255, 255)
p_w2b = tf_w2.add_paragraph()
p_w2b.text = "向外求名求利，永远焦虑；向内求道求心，才是根本"
p_w2b.font.size = Pt(15)
p_w2b.font.color.rgb = RGBColor(200, 200, 200)

app_box2 = slide7.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2.5))
tf_ap2 = app_box2.text_frame
tf_ap2.word_wrap = True
p_ap2a = tf_ap2.paragraphs[0]
p_ap2a.text = "【知行合一】"
p_ap2a.font.size = Pt(16)
p_ap2a.font.bold = True
p_ap2a.font.color.rgb = RGBColor(30, 60, 114)
apps2 = [
    "• 道的规律：惟恍惟惚——宇宙本质难以完全把握，保持敬畏",
    "• 德的践行：惟道是从——追随规律，而非追随欲望",
    "• 知行合一：明白道理不难，难的是真正按道理去做"
]
for app in apps2:
    p = tf_ap2.add_paragraph()
    p.text = app
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide7, 7, 10)

# ===== 第8页：江海处下 =====
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar8.fill.solid()
title_bar8.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar8.line.fill.background()
txBox8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf11 = txBox8.text_frame
p11 = tf11.paragraphs[0]
p11.text = "🌊 附：江海处下的智慧"
p11.font.size = Pt(32)
p11.font.bold = True
p11.font.color.rgb = RGBColor(255, 255, 255)

story_box3 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
story_box3.fill.solid()
story_box3.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_st3 = story_box3.text_frame
tf_st3.word_wrap = True
p_st3a = tf_st3.paragraphs[0]
p_st3a.text = "【上善若水·第八章延伸】"
p_st3a.font.size = Pt(16)
p_st3a.font.bold = True
p_st3a.font.color.rgb = RGBColor(30, 60, 114)
p_st3b = tf_st3.add_paragraph()
p_st3b.text = "江海之所以能为百谷王者，以其善下之，故能为百谷王。是以圣人欲上民，必以言下之；欲先民，必以身后之。——江海居下，故能纳百川；圣人谦下，故能得民心。"
p_st3b.font.size = Pt(15)
p_st3b.font.color.rgb = RGBColor(51, 51, 51)
p_st3b.space_before = Pt(8)

west_box3 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.6), Inches(9), Inches(1))
west_box3.fill.solid()
west_box3.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_w3 = west_box3.text_frame
tf_w3.word_wrap = True
p_w3a = tf_w3.paragraphs[0]
p_w3a.text = "🇬🇧 丘吉尔"
p_w3a.font.size = Pt(18)
p_w3a.font.bold = True
p_w3a.font.color.rgb = RGBColor(255, 255, 255)
p_w3b = tf_w3.add_paragraph()
p_w3b.text = "\"勇气是站起来的能力，但跌倒后爬起来更需要勇气。\"——处下不争，方能长久。"
p_w3b.font.size = Pt(14)
p_w3b.font.color.rgb = RGBColor(200, 200, 200)

app_box3 = slide8.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2.5))
tf_ap3 = app_box3.text_frame
tf_ap3.word_wrap = True
p_ap3a = tf_ap3.paragraphs[0]
p_ap3a.text = "【处下智慧】"
p_ap3a.font.size = Pt(16)
p_ap3a.font.bold = True
p_ap3a.font.color.rgb = RGBColor(30, 60, 114)
apps3 = [
    "• 谦受益，满招损——自视甚高的人，往往失去成长的空间",
    "• 欲上先下：想领导，先服务——团队成员感受到尊重，自然愿意追随",
    "• 功成不居：做了事不显摆——功劳归于团队，威望反而更高"
]
for app in apps3:
    p = tf_ap3.add_paragraph()
    p.text = app
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide8, 8, 10)

# ===== 第9页：生活应用 =====
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar9.fill.solid()
title_bar9.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar9.line.fill.background()
txBox9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf12 = txBox9.text_frame
p12 = tf12.paragraphs[0]
p12.text = "💡 生活中的智慧"
p12.font.size = Pt(32)
p12.font.bold = True
p12.font.color.rgb = RGBColor(255, 255, 255)

app_items = [
    ("🧘 内心做减法", "每天减少一件事——少刷一个App，少参加一个应酬，少买一件不需要的东西", ""),
    ("💼 职场心态", "不为KPI焦虑，不与人争高低——惟道是从，做好自己认为对的事", ""),
    ("❤️ 家庭关系", "少说教，多陪伴——绝圣弃智，用行动而非言语表达爱", ""),
    ("🎯 个人修养", "在喧嚣中保持泊然——众人的热闹是他们的，我的宁静是我的", "")
]
y = 1.4
for title, content, _ in app_items:
    app_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.3))
    app_box.fill.solid()
    app_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_app = app_box.text_frame
    tf_app.word_wrap = True
    p_a = tf_app.paragraphs[0]
    p_a.text = title
    p_a.font.size = Pt(16)
    p_a.font.bold = True
    p_a.font.color.rgb = RGBColor(30, 60, 114)
    p_b = tf_app.add_paragraph()
    p_b.text = content
    p_b.font.size = Pt(13)
    p_b.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.45

add_slide_number(slide9, 9, 10)

# ===== 第10页：精华总结 =====
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar10.fill.solid()
title_bar10.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar10.line.fill.background()
txBox10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf13 = txBox10.text_frame
p13 = tf13.paragraphs[0]
p13.text = "💎 精华提炼 · 第六课总结"
p13.font.size = Pt(32)
p13.font.bold = True
p13.font.color.rgb = RGBColor(255, 255, 255)

summary_items = [
    ("第十九章·绝圣弃智", "绝圣弃智，民利百倍——少折腾、少标榜，回归简单，反而更好"),
    ("第二十章·归真之旅", "我独泊兮，其未兆——在众人熙熙中，保持淡泊宁静，不从众"),
    ("第二十一章·孔德之容", "惟道是从——最大的德，是顺从宇宙规律，而非追逐名利")
]
y = 1.4
for title, content in summary_items:
    s_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.9))
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_s = s_box.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.text = f"{title}"
    p_s.font.size = Pt(16)
    p_s.font.bold = True
    p_s.font.color.rgb = RGBColor(30, 60, 114)
    p_s2 = tf_s.add_paragraph()
    p_s2.text = content
    p_s2.font.size = Pt(14)
    p_s2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.05

conclusion_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
conclusion_box.fill.solid()
conclusion_box.fill.fore_color.rgb = RGBColor(30, 60, 114)
tf_con = conclusion_box.text_frame
tf_con.word_wrap = True
p_con = tf_con.paragraphs[0]
p_con.text = "🌟 一句话总结"
p_con.font.size = Pt(18)
p_con.font.bold = True
p_con.font.color.rgb = RGBColor(255, 255, 255)
p_con2 = tf_con.add_paragraph()
p_con2.text = "绝学无忧，惟道是从"
p_con2.font.size = Pt(24)
p_con2.font.bold = True
p_con2.font.color.rgb = RGBColor(255, 255, 255)
p_con2.alignment = PP_ALIGN.CENTER
p_con3 = tf_con.add_paragraph()
p_con3.text = "减少世俗的焦虑，追随道的规律，在纷扰中守住内心那份淡泊与宁静。"
p_con3.font.size = Pt(16)
p_con3.font.color.rgb = RGBColor(200, 200, 200)

add_slide_number(slide10, 10, 10)

# 保存
output_path = '/Users/mac/.openclaw/workspace/courses/道德经_第6天.pptx'
prs.save(output_path)
print(f"PPT已生成：{output_path}")
