from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

# ===== 第1页：封面 =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar.line.fill.background()
txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第五课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

main_title = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main_title.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第16-18章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER

subtitle = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = subtitle.text_frame
p3 = tf3.paragraphs[0]
p3.text = "致虚守静 · 功成不居 · 大道废有仁义"
p3.font.size = Pt(28)
p3.font.color.rgb = RGBColor(70, 130, 180)
p3.alignment = PP_ALIGN.CENTER

quote_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_q = quote_box.text_frame
tf_q.word_wrap = True
p_q1 = tf_q.paragraphs[0]
p_q1.text = "致虚极，守静笃。"
p_q1.font.size = Pt(22)
p_q1.font.italic = True
p_q1.font.color.rgb = RGBColor(30, 60, 114)
p_q1.alignment = PP_ALIGN.CENTER
p_q2 = tf_q.add_paragraph()
p_q2.text = "—— 《道德经》第十六章"
p_q2.font.size = Pt(14)
p_q2.font.color.rgb = RGBColor(150, 150, 150)
p_q2.alignment = PP_ALIGN.CENTER

add_slide_number(slide1, 1, 10)

# ===== 第2页：第十六章详解 =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar2.fill.solid()
title_bar2.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar2.line.fill.background()
txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf5 = txBox2.text_frame
p5 = tf5.paragraphs[0]
p5.text = "第十六章 · 致虚守静"
p5.font.size = Pt(32)
p5.font.bold = True
p5.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.4))
yuanwen_box.fill.solid()
yuanwen_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw = yuanwen_box.text_frame
tf_yw.word_wrap = True
p_yw = tf_yw.paragraphs[0]
p_yw.text = "【原文】"
p_yw.font.size = Pt(18)
p_yw.font.bold = True
p_yw.font.color.rgb = RGBColor(30, 60, 114)
for line in ["致虚极，守静笃。万物并作，吾以观复。", "夫物芸芸，各复归其根。归根曰静，是谓复命。", "复命曰常，知常曰明。不知常，妄作凶。"]:
    p = tf_yw.add_paragraph()
    p.text = line
    p.font.size = Pt(19)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(3.2))
tf_c = content_box.text_frame
tf_c.word_wrap = True
p_c1 = tf_c.paragraphs[0]
p_c1.text = "【白话解读】"
p_c1.font.size = Pt(18)
p_c1.font.bold = True
p_c1.font.color.rgb = RGBColor(30, 60, 114)
解读 = [
    ("致虚极", "让心灵达到极度虚空的状态——放下执念，空杯才能盛水。"),
    ("守静笃", "坚定地保持宁静——静才能看清事物的本质和规律。"),
    ("观复", "观察万物循环往复——四时轮转、生死交替，都是道在运行。"),
    ("复命曰常", "回归本源是永恒的规律——知常才能明达，不妄动、不迷失。"),
]
for title, content in 解读:
    p = tf_c.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(5)

add_slide_number(slide2, 2, 10)

# ===== 第3页：第十六章生活智慧 =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar3.fill.solid()
title_bar3.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar3.line.fill.background()
txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf6 = txBox3.text_frame
p6 = tf6.paragraphs[0]
p6.text = "第十六章 · 生活智慧"
p6.font.size = Pt(32)
p6.font.bold = True
p6.font.color.rgb = RGBColor(255, 255, 255)

core_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.3))
core_box.fill.solid()
core_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_core = core_box.text_frame
tf_core.word_wrap = True
p_core = tf_core.paragraphs[0]
p_core.text = "🌟 核心智慧"
p_core.font.size = Pt(18)
p_core.font.bold = True
p_core.font.color.rgb = RGBColor(255, 255, 255)
p_core2 = tf_core.add_paragraph()
p_core2.text = "静观自照，回归本源——在喧嚣中守住内心宁静，才能看清方向。"
p_core2.font.size = Pt(16)
p_core2.font.color.rgb = RGBColor(255, 255, 255)

points = [
    ("1. 决定前先静心", "慌乱时做的决定往往错误。给自己一段静默时间，让心沉下来，答案自会浮现。"),
    ("2. 接受周期性规律", "事业有高低，人生有起伏。知道\"复命曰常\"，就能在低谷时保持平静。"),
    ("3. 每日静坐片刻", "哪怕10分钟不看手机，感受呼吸，让纷乱的念头自然沉淀。")
]
y = 2.9
for title, content in points:
    pt_box = slide3.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
    tf_pt = pt_box.text_frame
    tf_pt.word_wrap = True
    p_pt1 = tf_pt.paragraphs[0]
    p_pt1.text = title
    p_pt1.font.size = Pt(18)
    p_pt1.font.bold = True
    p_pt1.font.color.rgb = RGBColor(30, 60, 114)
    p_pt2 = tf_pt.add_paragraph()
    p_pt2.text = content
    p_pt2.font.size = Pt(14)
    p_pt2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.2

add_slide_number(slide3, 3, 10)

# ===== 第4页：第十七章详解 =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar4.fill.solid()
title_bar4.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar4.line.fill.background()
txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf7 = txBox4.text_frame
p7 = tf7.paragraphs[0]
p7.text = "第十七章 · 功成不居"
p7.font.size = Pt(32)
p7.font.bold = True
p7.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.2))
yuanwen_box2.fill.solid()
yuanwen_box2.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw2 = yuanwen_box2.text_frame
tf_yw2.word_wrap = True
p_yw2 = tf_yw2.paragraphs[0]
p_yw2.text = "【原文】"
p_yw2.font.size = Pt(18)
p_yw2.font.bold = True
p_yw2.font.color.rgb = RGBColor(30, 60, 114)
for line in ["太上，下知有之；其次，亲而誉之；", "其次，畏之；其次，侮之。", "信不足焉，有不信焉。悠兮，其贵言。"]:
    p = tf_yw2.add_paragraph()
    p.text = line
    p.font.size = Pt(19)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box2 = slide4.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(4))
tf_c2 = content_box2.text_frame
tf_c2.word_wrap = True
p_c4 = tf_c2.paragraphs[0]
p_c4.text = "【白话解读】"
p_c4.font.size = Pt(18)
p_c4.font.bold = True
p_c4.font.color.rgb = RGBColor(30, 60, 114)
解读2 = [
    ("太上下知有之", "最高明的领导者，百姓只知道他的存在——制度完善，无为而治。"),
    ("其次亲而誉之", "次一等的，百姓亲近并称赞他——德治仁政，有为教化。"),
    ("其次畏之侮之", "再次的，百姓害怕他；最次的，百姓轻蔑他——威权暴政，强压之下必有反抗。"),
    ("核心：贵言", "最好的治理，是少发号施令——言多必失，令繁必乱。"),
]
for title, content in 解读2:
    p = tf_c2.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(5)

add_slide_number(slide4, 4, 10)

# ===== 第5页：管理哲学典故 =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar5.fill.solid()
title_bar5.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar5.line.fill.background()
txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf8 = txBox5.text_frame
p8 = tf8.paragraphs[0]
p8.text = "📜 典故：功成不居的管理智慧"
p8.font.size = Pt(32)
p8.font.bold = True
p8.font.color.rgb =RGBColor(255, 255, 255)

story_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
story_box.fill.solid()
story_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_story = story_box.text_frame
tf_story.word_wrap = True
p_s1 = tf_story.paragraphs[0]
p_s1.text = "【张良功成身退】"
p_s1.font.size = Pt(16)
p_s1.font.bold = True
p_s1.font.color.rgb = RGBColor(30, 60, 114)
p_s2 = tf_story.add_paragraph()
p_s2.text = "张良佐汉高祖定天下，封留侯。后从赤松子游，学辟谷之术。天下既定，他不贪功名，急流勇退。刘邦后杀韩信、彭越，而张良独善其身，正是\"功成不居\"的最好例证。"
p_s2.font.size = Pt(15)
p_s2.font.color.rgb = RGBColor(51, 51, 51)
p_s2.space_before = Pt(8)
p_s3 = tf_story.add_paragraph()
p_s3.text = "老子的\"太上下知有之\"——最好的领导，是让团队感受不到你的存在，却离不开你的制度。"
p_s3.font.size = Pt(15)
p_s3.font.color.rgb = RGBColor(51, 51, 51)

west_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.6), Inches(9), Inches(1))
west_box.fill.solid()
west_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_west = west_box.text_frame
tf_west.word_wrap = True
p_west = tf_west.paragraphs[0]
p_west.text = "🌍 西方管理对照"
p_west.font.size = Pt(18)
p_west.font.bold = True
p_west.font.color.rgb = RGBColor(255, 255, 255)
p_west2 = tf_west.add_paragraph()
p_west2.text = "德鲁克：\"管理的最高境界，是无需管理。\" 优秀的制度，让执行者自觉运转。"
p_west2.font.size = Pt(14)
p_west2.font.color.rgb = RGBColor(200, 200, 200)

app_box = slide5.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2.5))
tf_app = app_box.text_frame
tf_app.word_wrap = True
p_a1 = tf_app.paragraphs[0]
p_a1.text = "【管理启示】"
p_a1.font.size = Pt(16)
p_a1.font.bold = True
p_a1.font.color.rgb = RGBColor(30, 60, 114)
apps = [
    "• 建立制度，而非依赖个人权威——\"太上，下知有之\"",
    "• 授权到位，不过度干预——让团队自主运转",
    "• 功成不居：成绩归团队，风险我承担"
]
for app in apps:
    p = tf_app.add_paragraph()
    p.text = app
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide5, 5, 10)

# ===== 第6页：第十八章详解 =====
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar6.fill.solid()
title_bar6.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar6.line.fill.background()
txBox6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf9 = txBox6.text_frame
p9 = tf9.paragraphs[0]
p9.text = "第十八章 · 大道废有仁义"
p9.font.size = Pt(32)
p9.font.bold = True
p9.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box3 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
yuanwen_box3.fill.solid()
yuanwen_box3.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw3 = yuanwen_box3.text_frame
tf_yw3.word_wrap = True
p_yw3 = tf_yw3.paragraphs[0]
p_yw3.text = "【原文】"
p_yw3.font.size = Pt(18)
p_yw3.font.bold = True
p_yw3.font.color.rgb = RGBColor(30, 60, 114)
for line in ["大道废，有仁义；", "智慧出，有大伪；", "六亲不和，有孝慈；国家昏乱，有忠臣。"]:
    p = tf_yw3.add_paragraph()
    p.text = line
    p.font.size = Pt(19)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box3 = slide6.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(4))
tf_c3 = content_box3.text_frame
tf_c3.word_wrap = True
p_c7 = tf_c3.paragraphs[0]
p_c7.text = "【白话解读】"
p_c7.font.size = Pt(18)
p_c7.font.bold = True
p_c7.font.color.rgb = RGBColor(30, 60, 114)
解读3 = [
    ("大道废，有仁义", "天道被废弃，才需要倡导仁义——道德败坏才需要道德说教。"),
    ("智慧出，有大伪", "智巧出现，诈伪也就跟着来了——机关算尽，反误了卿卿性命。"),
    ("六亲不和，有孝慈", "家庭失和，才需要彰显孝慈——平日和睦，何须特别标榜？"),
    ("核心：返璞归真", "最好的状态是无需刻意——真正的和谐，不需要对比和彰显。"),
]
for title, content in 解读3:
    p = tf_c3.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(5)

add_slide_number(slide6, 6, 10)

# ===== 第7页：大道废典故 =====
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar7.fill.solid()
title_bar7.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar7.line.fill.background()
txBox7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf10 = txBox7.text_frame
p10 = tf10.paragraphs[0]
p10.text = "📜 典故：大道废有仁义"
p10.font.size = Pt(32)
p10.font.bold = True
p10.font.color.rgb = RGBColor(255, 255, 255)

story_box2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.2))
story_box2.fill.solid()
story_box2.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_st2 = story_box2.text_frame
tf_st2.word_wrap = True
p_st2a = tf_st2.paragraphs[0]
p_st2a.text = "【庄子·外篇·在宥】"
p_st2a.font.size = Pt(16)
p_st2a.font.bold = True
p_st2a.font.color.rgb = RGBColor(30, 60, 114)
p_st2b = tf_st2.add_paragraph()
p_st2b.text = "孔子西藏游，谒见老子。老子曰：\"先生想以仁义教导君主，是扰乱其心啊！\"仁义是乱世之药——当社会道德尚好时，仁义无需标榜；当需要大力提倡时，问题已经严重了。"
p_st2b.font.size = Pt(15)
p_st2b.font.color.rgb = RGBColor(51, 51, 51)
p_st2b.space_before = Pt(8)

west_box2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(9), Inches(1))
west_box2.fill.solid()
west_box2.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_w2 = west_box2.text_frame
tf_w2.word_wrap = True
p_w2a = tf_w2.paragraphs[0]
p_w2a.text = "⚡ 现代职场警示"
p_w2a.font.size = Pt(18)
p_w2a.font.bold = True
p_w2a.font.color.rgb = RGBColor(255, 255, 255)
p_w2b = tf_w2.add_paragraph()
p_w2b.text = "一家公司反复强调\"诚信文化\"——说明诚信问题已经很突出了"
p_w2b.font.size = Pt(15)
p_w2b.font.color.rgb = RGBColor(200, 200, 200)

app_box2 = slide7.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2.5))
tf_ap2 = app_box2.text_frame
tf_ap2.word_wrap = True
p_ap2a = tf_ap2.paragraphs[0]
p_ap2a.text = "【辨证智慧】"
p_ap2a.font.size = Pt(16)
p_ap2a.font.bold = True
p_ap2a.font.color.rgb = RGBColor(30, 60, 114)
apps2 = [
    "• 警惕\"道德标榜\"：天天讲道德，可能道德正在滑坡",
    "• 看实质而非形式：六亲和睦的家庭，无需\"孝子\"奖状",
    "• 问题出现后的补救，远不如预防——\"为之于未有\""
]
for app in apps2:
    p = tf_ap2.add_paragraph()
    p.text = app
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide7, 7, 10)

# ===== 第8页：西方哲学对照 =====
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar8.fill.solid()
title_bar8.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar8.line.fill.background()
txBox8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf11 = txBox8.text_frame
p11 = tf11.paragraphs[0]
p11.text = "🌍 西方思想对照"
p11.font.size = Pt(32)
p11.font.bold = True
p11.font.color.rgb = RGBColor(255, 255, 255)

philosophers = [
    ("🇫🇷 卢梭", "\"自然状态\"", "\"文明是疾病，良心是天使。\"——文明越复杂，人性越扭曲。", "1754年"),
    ("🇩🇪 马克思", "\"异化理论\"", "资本主义使人\"异化\"——人创造了机器，却成为机器的附属品。", "1844年"),
    ("🇺🇸 梭罗", "\"文明的反面\"", "\"文明创造了更体面的苦刑。\"——物质越丰富，束缚越深。", "1849年"),
    ("🇬🇧 斯密", "\"道德情感论\"", "同情心是人的本性，市场经济反而让人冷漠——资本主义的道德代价。", "1759年")
]
y = 1.4
for name, theory, desc, year in philosophers:
    phil_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.2))
    phil_box.fill.solid()
    phil_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_phil = phil_box.text_frame
    tf_phil.word_wrap = True
    p_n = tf_phil.paragraphs[0]
    p_n.text = f"{name} · {theory} · {year}"
    p_n.font.size = Pt(16)
    p_n.font.bold = True
    p_n.font.color.rgb = RGBColor(30, 60, 114)
    p_d = tf_phil.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.35

add_slide_number(slide8, 8, 10)

# ===== 第9页：生活应用 =====
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar9.fill.solid()
title_bar9.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar9.line.fill.background()
txBox9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf12 = txBox9.text_frame
p12 = tf12.paragraphs[0]
p12.text = "💡 生活中的智慧"
p12.font.size = Pt(32)
p12.font.bold = True
p12.font.color.rgb = RGBColor(255, 255, 255)

app_items = [
    ("🧘 每日静心", "晨起静坐10分钟，不刷手机，感受\"致虚极，守静笃\"的清凉", ""),
    ("💼 领导力修炼", "少开会、少发指令——建立好制度，让团队\"下知有之\"自动运转", ""),
    ("❤️ 家庭关系", "家庭和睦不需要\"孝子模范\"——日常的陪伴，胜过节日的仪式", ""),
    ("🎯 个人成长", "追求内在充实，而非外在标签——真正的成熟，不需要向人证明", "")
]
y = 1.4
for title, content, _ in app_items:
    app_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.3))
    app_box.fill.solid()
    app_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_app = app_box.text_frame
    tf_app.word_wrap = True
    p_a = tf_app.paragraphs[0]
    p_a.text = title
    p_a.font.size = Pt(16)
    p_a.font.bold = True
    p_a.font.color.rgb = RGBColor(30, 60, 114)
    p_b = tf_app.add_paragraph()
    p_b.text = content
    p_b.font.size = Pt(13)
    p_b.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.45

add_slide_number(slide9, 9, 10)

# ===== 第10页：精华总结 =====
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar10.fill.solid()
title_bar10.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar10.line.fill.background()
txBox10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf13 = txBox10.text_frame
p13 = tf13.paragraphs[0]
p13.text = "💎 精华提炼 · 第五课总结"
p13.font.size = Pt(32)
p13.font.bold = True
p13.font.color.rgb = RGBColor(255, 255, 255)

summary_items = [
    ("第十六章·致虚守静", "万物并作，吾以观复——静观本质，回归本源，知晓常理方能明达"),
    ("第十七章·功成不居", "太上下知有之——最好的领导，是让制度代替个人权威"),
    ("第十八章·大道废", "大道废，有仁义——社会问题已出现，才需要道德标榜")
]
y = 1.4
for title, content in summary_items:
    s_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.9))
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_s = s_box.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.text = f"{title}"
    p_s.font.size = Pt(16)
    p_s.font.bold = True
    p_s.font.color.rgb = RGBColor(30, 60, 114)
    p_s2 = tf_s.add_paragraph()
    p_s2.text = content
    p_s2.font.size = Pt(14)
    p_s2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.05

conclusion_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
conclusion_box.fill.solid()
conclusion_box.fill.fore_color.rgb = RGBColor(30, 60, 114)
tf_con = conclusion_box.text_frame
tf_con.word_wrap = True
p_con = tf_con.paragraphs[0]
p_con.text = "🌟 一句话总结"
p_con.font.size = Pt(18)
p_con.font.bold = True
p_con.font.color.rgb = RGBColor(255, 255, 255)
p_con2 = tf_con.add_paragraph()
p_con2.text = "致虚极，守静笃"
p_con2.font.size = Pt(24)
p_con2.font.bold = True
p_con2.font.color.rgb = RGBColor(255, 255, 255)
p_con2.alignment = PP_ALIGN.CENTER
p_con3 = tf_con.add_paragraph()
p_con3.text = "在喧嚣中守住宁静，在纷扰中回归本质，功成不居，自然自在。"
p_con3.font.size = Pt(16)
p_con3.font.color.rgb = RGBColor(200, 200, 200)

add_slide_number(slide10, 10, 10)

# 保存
prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第5天.pptx')
print("PPT已生成：道德经_第5天.pptx")
