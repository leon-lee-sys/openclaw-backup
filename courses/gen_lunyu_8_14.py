#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
论语 第8-14天课件生成 - 简化版
"""
import os
import sys
sys.path.insert(0, '/opt/homebrew/lib/node_modules/openclaw/node_modules')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
ACCENT = RGBColor(180, 60, 60)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.2), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "%d/%d" % (num, total)
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_text_box(slide, text, x, y, width, height, font_size=18, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box

def add_rounded_box(slide, text, x, y, width, height, bg_color=BG, text_color=TEXT, font_size=16):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    return box

CHAPTERS = [
    # Day 8: 泰伯第八
    [
        {"num": 8, "name": "泰伯第八", "yuanwen": ["子曰：泰伯其可谓至德也已矣，三以天下让，民无得而称焉。"], "translation": "泰伯是周太王的长子，他三次把天下让给季历，百姓找不到词语来称赞他。", "core": ["至德：泰伯的让国是最高德行的体现", "三以天下让：多次礼让，发自内心", "民无得而称：真正的高尚无法言说"]},
        {"num": 9, "name": "士不可不弘毅", "yuanwen": ["曾子曰：士不可不弘毅，任重而道远。仁以为己任，不亦重乎？死而后已，不亦远乎？"], "translation": "读书人必须有远大的志向和坚定的意志，因为负担重而路途远。把仁德作为自己的使命，这个担子不是很重吗？直到死才能放下，这不是很远吗？", "core": ["弘毅：志向远大、意志坚定", "任重道远：仁德是一生的使命", "死而后已：真正的坚持"]},
        {"num": 10, "name": "可以学艺", "yuanwen": ["子曰：诗三百，一言以蔽之，曰思无邪。"], "translation": "《诗经》三百首，用一句话概括，就是思想纯正。", "core": ["诗教功能：文学有教化人心的力量", "思无邪：真正美好的诗歌能净化心灵", "情感到理：从真挚情感到达道德境界"]}
    ],
    # Day 9: 乡党第十
    [
        {"num": 11, "name": "祭如在", "yuanwen": ["祭如在，祭神如神在。不诚无物。"], "translation": "祭祀祖先要好像祖先真的在场，祭祀神祇要好像神真的存在。不真诚的祭祀没有意义。", "core": ["祭祀要真诚：形式不是最重要，内心虔诚才是根本", "敬如在场：这是诚的体现", "不祭无义：形式主义的祭祀没有意义"]},
        {"num": 12, "name": "食不语", "yuanwen": ["食不语，寝不言。"], "translation": "吃饭时不交谈，睡觉时不说话。", "core": ["专心一致：吃饭就专心吃饭，睡觉就专心睡觉", "食不言：养生之道", "寝不言：睡觉不说话影响睡眠质量"]},
        {"num": 13, "name": "君子不急", "yuanwen": ["君子不重则不威，学则不固。主忠信，无友不如己者，过则勿惮改。"], "translation": "君子不庄重就没有威严，学习就不能牢固。以忠诚和信用为主，不要和不如自己的人交朋友，有了过错就不要害怕改正。", "core": ["不重不威：庄重是威严的基础", "学则不固：学习让人不断进步", "过则勿惮改：不怕改正错误"]}
    ],
    # Day 10: 先进第十一
    [
        {"num": 14, "name": "颜渊好学", "yuanwen": ["哀公问：弟子孰为好学？孔子对曰：有颜回者好学，不迁怒，不贰过。不幸短命死矣！今也则亡，未闻好学者也。"], "translation": "鲁哀公问：你的学生谁是最好学的？孔子说：颜回最好学，他不迁怒于人，不重复犯同样的过错。不幸短命死了！现在没有像他这样的人了！", "core": ["不迁怒：情绪管理，不向无关的人发泄", "不贰过：从失败中学习，不重复犯错", "好学的定义：不是学得多，而是人格完善"]},
        {"num": 15, "name": "中庸之为德", "yuanwen": ["子曰：中庸之为德也，其至矣乎！民鲜久矣。"], "translation": "孔子说：中庸作为德行，是最高的了！百姓很少能长期实行它。", "core": ["中庸之道：不偏不倚，恰到好处", "民鲜久矣：很少人能长期坚持", "日常体现：中庸就在百姓日用之中"]},
        {"num": 16, "name": "学而时习", "yuanwen": ["子曰：学而时习之，不亦说乎？有朋自远方来，不亦乐乎？人不知而不愠，不亦君子乎？"], "translation": "学习知识后时常复习不是很喜悦吗？有志同道合的朋友从远方来不是很快乐吗？别人不了解我，我却不恼怒，这不也是君子吗？", "core": ["学而时习：学习需要不断复习", "有朋远来：志同道合的朋友是人生财富", "人不知不愠：不被认同时保持内心平静"]}
    ],
    # Day 11: 颜渊第十二
    [
        {"num": 17, "name": "四维正", "yuanwen": ["孔子曰：君子之德风，小人之德草。草上之风必偃。"], "translation": "领导者的德行像风，百姓的德行像草。风吹过来，草必然会倒下。", "core": ["上行下效：上面的人影响下面", "德风：领导者的示范作用", "草必偃：被影响是必然的"]},
        {"num": 18, "name": "己所不欲", "yuanwen": ["仲弓问仁。子曰：出门如见大宾，使民如承大祭。己所不欲，勿施于人。"], "translation": "仲弓问什么是仁。孔子说：出门要像会见重要宾客，使唤百姓要像承当大祭典。自己不愿意的，不要强加给别人。", "core": ["恕道：推己及人", "己所不欲勿施于人：儒家的黄金法则", "出门如见大宾：尊重每一个人"]},
        {"num": 19, "name": "求卓复", "yuanwen": ["冉求曰：非不说子之道，力不足也。子曰：力不足者中道而废，今女画。"], "translation": "冉求说：我不是不喜欢您的学说，是才能不够。孔子说：才能不够的人是做一半就放弃，你是自己给自己画了界限。", "core": ["画地为限：自己限制自己", "力不足还是不想：很多时候是态度问题", "突破自我设限：相信能力可以培养"]}
    ],
    # Day 12: 子路第十三
    [
        {"num": 20, "name": "政者正也", "yuanwen": ["子路问政。子曰：先之劳之。请益。曰：无倦。"], "translation": "子路问如何治理政事。孔子说：自己带头做，百姓就会勤劳工作。子路请求多讲一些。孔子说：不要倦怠。", "core": ["先之：以身作则", "劳之：带动百姓勤劳", "无倦：坚持不懈"]},
        {"num": 21, "name": "其身正", "yuanwen": ["子曰：其身正，不令而行；其身不正，虽令不从。"], "translation": "孔子说：自身正了，不发布命令事情也能实行；自身不正，即使发布命令也没人听从。", "core": ["以身作则：行动比言语更重要", "身正令行：自身正了不需要强制", "身不正虽令不从：权力不如威信"]},
        {"num": 22, "name": "上有所好", "yuanwen": ["上好礼，则民莫敢不敬；上好义，则民莫敢不服；上好信，则民莫敢不用情。"], "translation": "上面的人喜好礼仪，百姓就不敢不尊敬；上面的人喜好道义，百姓就不敢不服从；上面的人喜好诚信，百姓就不敢不真诚。", "core": ["上行下效：上面的喜好影响下面", "上好礼民莫敢不敬：领导者的示范作用", "民的响应是必然的：不是拍马屁而是规律"]}
    ],
    # Day 13: 宪问第十四
    [
        {"num": 23, "name": "为政以德", "yuanwen": ["为政以德譬如北辰，居其所而众星共之。"], "translation": "用德行来治理国家，就像北极星——它在自己的位置上，而众星都环绕着它。", "core": ["德治优于力治：感化而非强迫", "北辰之喻：有德的领导者像北极星", "无为而治：德行所至，人心自服"]},
        {"num": 24, "name": "仁者无敌", "yuanwen": ["仁者不忧，知者不惑，勇者不惧。"], "translation": "有仁德的人不忧愁，有智慧的人不迷惑，勇敢的人不恐惧。", "core": ["仁者不忧：内心安定没有忧虑", "知者不惑：有智慧就能辨别是非", "勇者不惧：勇敢的人一往无前"]},
        {"num": 25, "name": "修身起践", "yuanwen": ["子路曰：何以为身？子曰：恭、宽、信、敏、惠。恭则不侮，宽则得众，信则人任焉，敏则有功，惠则足以使人。"], "translation": "子路问如何修养自身。孔子说：恭敬、宽厚、诚信、勤敏、慈惠。恭敬就不会被侮辱，宽厚就能得到众人，信实就会被人信任，勤敏就能有功绩，慈惠就能很好地使用人。", "core": ["五德：恭、宽、信、敏、惠", "恭则不侮：恭敬是立身之本", "惠则足以使人：施恩惠才能用人"]}
    ],
    # Day 14: 卫灵公第十五
    [
        {"num": 26, "name": "及人诸", "yuanwen": ["子贡问曰：有一言而可以终身行之者乎？子曰：其恕乎！己所不欲，勿施于人。"], "translation": "子贡问：有没有一个字可以终身奉行的？孔子说：大概是'恕'吧——自己不愿意的，不要强加给别人。", "core": ["一以贯之：孔子思想的核心是恕", "恕道：推己及人的方法论", "终身奉行：简单但不容易"]},
        {"num": 27, "name": "知进退", "yuanwen": ["子曰：可与共学，未可与适道；可与适道，未可与立；可与立，未可与权。"], "translation": "孔子说：有的人可以一起学习，未必可以一起得道；有的人可以一起得道，未必可以一起坚守；有的人可以一起坚守，未必可以一起通权达变。", "core": ["学→道→立→权：境界的递进", "知进退：知道什么时候该坚持什么时候该变通", "权道：通权达变是最高的境界"]},
        {"num": 28, "name": "求道", "yuanwen": ["子曰：求仁得仁，又何怨？"], "translation": "孔子说：追求仁德而得到仁德，还有什么可抱怨的呢？", "core": ["求仁得仁：追求什么就得到什么", "无怨：内心平静不抱怨", "知足：知道自己要什么"]}
    ]
]

def make_pptx(day, chapters, output_dir):
    prs = Presentation()
    
    # Cover slide
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "国学经典精读 · 第%d天" % day
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    main = s.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
    tf2 = main.text_frame
    p2 = tf2.paragraphs[0]
    titles = {8:"泰伯第八", 9:"乡党第十", 10:"先进第十一", 11:"颜渊第十二", 12:"子路第十三", 13:"宪问第十四", 14:"卫灵公第十五"}
    p2.text = "论语 · %s" % titles.get(day, "")
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.alignment = PP_ALIGN.CENTER
    
    add_slide_number(s, 1, 14)
    
    # Chapter slides
    slide_num = 2
    for ch in chapters:
        # Chapter intro
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_title_bar(s, "第%d章 · %s" % (ch["num"], ch["name"]))
        
        add_rounded_box(s, "[原文]\n" + "\n".join(ch["yuanwen"]), 0.5, 1.5, 9, 2.5, BG, TEXT, 18)
        add_rounded_box(s, "[译文]\n" + ch["translation"], 0.5, 4.2, 9, 1.5, BLUE_LIGHT, RGBColor(255,255,255), 16)
        
        y = 5.9
        for point in ch["core"]:
            add_text_box(s, "• " + point, 0.5, y, 9, 0.5, 15, False, PRIMARY)
            y += 0.5
        
        add_slide_number(s, slide_num, 14)
        slide_num += 1
    
    path = os.path.join(output_dir, "论语_第%d天.pptx" % day)
    prs.save(path)
    print("Done: 论语_第%d天.pptx" % day)
    return path

if __name__ == "__main__":
    output_dir = "/Users/mac/.openclaw/workspace/courses/论语"
    os.makedirs(output_dir, exist_ok=True)
    
    for i, chapters in enumerate(CHAPTERS):
        day = i + 8
        make_pptx(day, chapters, output_dir)
    
    print("\n第8-14天课件生成完成!")
