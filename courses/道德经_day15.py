# -*- coding: utf-8 -*-
"""
道德经第15天 课件生成器
章节：第49-51章
风格：精读深度讲解版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
ACCENT = RGBColor(204, 0, 0)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_yuanwen_box(slide, lines, y=1.4, height=2.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(20)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=BLUE_LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(255, 255, 255)

def add_points(slide, points, y_start=3.1):
    y = y_start
    for title, content in points:
        pt_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT
        y += 1.1
    return y

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.5))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第15课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第46-48章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "天下有道 · 不窥牖见天道 · 为学日益"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "天下有道，却走马以粪。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = PRIMARY
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第四十六章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = SECONDARY
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, 12)

# ===== 第2页：第四十六章详解 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s2, "第四十六章 · 天下有道")
add_yuanwen_box(s2, [
    "天下有道，却走马以粪；",
    "天下无道，戎马生于郊。",
    "",
    "罪莫大于可欲，",
    "祸莫大于不知足，",
    "咎莫大于欲得。",
    "",
    "故知足之足，常足矣。"
], y=1.4, height=3.8)
content2 = s2.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2.5))
tf_c2 = content2.text_frame
tf_c2.word_wrap = True
pc2 = tf_c2.paragraphs[0]
pc2.text = "【白话解读】天下有道时，战马用来耕田；天下无道时，战马在战场生育。最大的罪过是永无止境的欲望。知足才是真正的富足。"
pc2.font.size = Pt(14)
pc2.font.color.rgb = TEXT
add_slide_number(s2, 2, 12)

# ===== 第3页：第四十六章核心 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s3, "第四十六章 · 核心思想")
add_core_box(s3, "🌟 核心观点", "最大的罪恶是永无止境的欲望，知足才是真正的富足。")
add_points(s3, [
    ("1. 天下有道", "政治清明时，战马退役耕田，社会安宁。"),
    ("2. 天下无道", "政治黑暗时，战马用于战争，连怀胎的母马都要上战场。"),
    ("3. 知足之足", "知道满足的那种满足，才是真正的、永远的满足。"),
    ("4. 欲壑难填", "不知足是一切罪恶、灾祸、过失的根源。")
])
add_slide_number(s3, 3, 12)

# ===== 第4页：第四十七章详解 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s4, "第四十七章 · 不窥牖见天道")
add_yuanwen_box(s4, [
    "不出户，知天下；",
    "不窥牖，见天道。",
    "",
    "其出弥远，其知弥少。",
    "",
    "是以圣人不行而知，",
    "不见而名，",
    "不为而成。"
], y=1.4, height=3.8)
content4 = s4.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2.5))
tf_c4 = content4.text_frame
tf_c4.word_wrap = True
pc4 = tf_c4.paragraphs[0]
pc4.text = "【白话解读】不出门能知天下事，不望窗外能见天道。走得越远，知道得越少。所以圣人不远行却知道，不看见却明了，不妄为却成功。"
pc4.font.size = Pt(14)
pc4.font.color.rgb = TEXT
add_slide_number(s4, 4, 12)

# ===== 第5页：第四十七章核心 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s5, "第四十七章 · 核心思想")
add_core_box(s5, "🌟 核心观点", "认识天道不需要远行，真正的智慧在于内省和悟道，而非走马观花。")
add_points(s5, [
    ("1. 不出户知天下", "通过推理和内省，可以了解天下大事。"),
    ("2. 不窥牖见天道", "通过观察自然规律，可以理解天道运行。"),
    ("3. 出弥远知弥少", "跑得越远，见识越少——忙于外物，反而迷失内心。"),
    ("4. 圣人之道", "圣人不依赖亲身经历，而是通过道来推断一切。")
])
add_slide_number(s5, 5, 12)

# ===== 第6页：现代反思 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s6, "第四十六-四十七章 · 现代反思")
items6 = [
    ("🏃 信息过载", "其出弥远，其知弥少 → 每天刷手机10小时，获得的有效信息可能不如专注读一本书。"),
    ("💰 消费主义", "罪莫大于可欲 → 商家制造焦虑，让人觉得\"必须拥有\"，实际上是不知足的陷阱。"),
    ("📚 内省之道", "不出户，知天下 → 每天留出时间独处思考，比不断向外索取更能提升认知。"),
    ("🎯 专注力", "不见而名，不为而成 → 减少不必要的行动，专注最重要的事，反而更容易成功。")
]
y6 = 1.4
for icon_title, explanation in items6:
    box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y6), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf6 = box.text_frame
    tf6.word_wrap = True
    p6_1 = tf6.paragraphs[0]
    p6_1.text = icon_title
    p6_1.font.size = Pt(16)
    p6_1.font.bold = True
    p6_1.font.color.rgb = PRIMARY
    p6_2 = tf6.add_paragraph()
    p6_2.text = explanation
    p6_2.font.size = Pt(14)
    p6_2.font.color.rgb = TEXT
    y6 += 1.4
add_slide_number(s6, 6, 12)

# ===== 第7页：第四十八章详解 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s7, "第四十八章 · 为学日益")
add_yuanwen_box(s7, [
    "为学日益，",
    "为道日损，",
    "损之又损，",
    "以至于无为，",
    "无为而无不为。",
    "",
    "取天下常以无事，",
    "及其有事，",
    "不足以取天下。"
], y=1.4, height=4.0)
content7 = s7.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(2))
tf_c7 = content7.text_frame
tf_c7.word_wrap = True
pc7 = tf_c7.paragraphs[0]
pc7.text = "【白话解读】追求学问，知识每天增加；追求大道，欲望每天减少。减到极致，就能无为；无为，反而什么都能做成。"
pc7.font.size = Pt(14)
pc7.font.color.rgb = TEXT
add_slide_number(s7, 7, 12)

# ===== 第8页：第四十八章核心 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s8, "第四十八章 · 核心思想")
add_core_box(s8, "🌟 核心观点", "为学是做加法，为道是做减法。放下执念，才能顺道而行，无为而无不为。")
add_points(s8, [
    ("1. 为学日益", "追求知识和技能，需要不断积累——这是入世的态度。"),
    ("2. 为道日损", "追求大道，需要不断放下——放下执念、偏见、欲望。"),
    ("3. 无为而无不为", "不妄为反而什么都做成了——顺道而行，事半功倍。"),
    ("4. 取天下常以无事", "治理天下要顺其自然，不扰民，不折腾。")
])
add_slide_number(s8, 8, 12)

# ===== 第9页：西方思想 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s9, "🌍 西方思想对照")
phils = [
    ("苏格拉底 · 认识你自己", "知道自己不知道，是智慧的开始。追求知识要先承认无知。约前400年"),
    ("笛卡尔 · 怀疑一切", "我思故我在——通过彻底怀疑找到最可靠的真理。1641年"),
    ("奥卡姆 · 剃刀原理", "如无必要，勿增实体——最简单的解释往往是最正确的。约1320年"),
    ("爱因斯坦 · 想象力", "想象力比知识更重要。知识有限，而想象力无限。1952年")
]
y9 = 1.4
for name_theory, desc_year in phils:
    box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y9), Inches(9), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf9 = box.text_frame
    tf9.word_wrap = True
    p9_1 = tf9.paragraphs[0]
    p9_1.text = name_theory
    p9_1.font.size = Pt(16)
    p9_1.font.bold = True
    p9_1.font.color.rgb = PRIMARY
    p9_2 = tf9.add_paragraph()
    p9_2.text = desc_year
    p9_2.font.size = Pt(14)
    p9_2.font.color.rgb = TEXT
    y9 += 1.3
add_slide_number(s9, 9, 12)

# ===== 第10页：现代应用 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s10, "📌 现代应用")
add_points(s10, [
    ("个人成长", "为学日益——保持学习，每天进步一点点，积少成多。"),
    ("断舍离生活", "为道日损——减少不必要的物质和执念，轻装上阵。"),
    ("企业管理", "无为而无不为——给团队足够的空间，不事事微管理。"),
    ("投资理财", "知足之足，常足矣——设定合理的收益目标，不贪婪。")
])
add_slide_number(s10, 10, 12)

# ===== 第11页：中国典故 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s11, "📜 中国典故：杨震\"四知\"拒金")
box11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.5))
box11.fill.solid()
box11.fill.fore_color.rgb = BG
tf11 = box11.text_frame
tf11.word_wrap = True
p11_1 = tf11.paragraphs[0]
p11_1.text = "【后汉书·杨震传】"
p11_1.font.size = Pt(16)
p11_1.font.bold = True
p11_1.font.color.rgb = PRIMARY
p11_2 = tf11.add_paragraph()
p11_2.text = "杨震调任东莱太守时，路过昌邑县。县令王密曾受杨震提拔，深夜怀揣十斤金子要送给杨震。杨震拒绝。王密说：\"深夜无人知道。\"杨震说：\"天知、神知、我知、子知，何谓无知？\""
p11_2.font.size = Pt(16)
p11_2.font.color.rgb = TEXT
p11_2.space_before = Pt(10)
insight11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.1), Inches(9), Inches(1))
insight11.fill.solid()
insight11.fill.fore_color.rgb = BLUE_LIGHT
tfi11 = insight11.text_frame
tfi11.word_wrap = True
pi11 = tfi11.paragraphs[0]
pi11.text = "💡 启示"
pi11.font.size = Pt(18)
pi11.font.bold = True
pi11.font.color.rgb = RGBColor(255, 255, 255)
pi11_2 = tfi11.add_paragraph()
pi11_2.text = "杨震的\"四知\"体现了\"为道日损\"的精神——放下私欲，心怀坦荡，自然无愧。知足不贪，才能常足。"
pi11_2.font.size = Pt(14)
pi11_2.font.color.rgb = RGBColor(255, 255, 255)
app11 = s11.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2))
tfa11 = app11.text_frame
tfa11.word_wrap = True
pa11 = tfa11.paragraphs[0]
pa11.text = "【现代应用】职场拒绝不当利益，生活知足常乐，处事心怀坦荡。君子慎独，不因无人知晓而放松自我要求。"
pa11.font.size = Pt(14)
pa11.font.color.rgb = TEXT
add_slide_number(s11, 11, 12)

# ===== 第12页：总结 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s12, "第15课 · 三章精华总结")
sums = [
    ("第四十六章", "天下有道", "知足之足，常足矣。不知足是一切罪恶灾祸的根源。"),
    ("第四十七章", "不窥牖见天道", "不出户知天下，为道日损，圣人不行而成。"),
    ("第四十八章", "为学日益", "为学做加法，为道做减法，无为而无不为。")
]
y12 = 1.4
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    y12 += 1.6
final = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(9), Inches(1))
final.fill.solid()
final.fill.fore_color.rgb = PRIMARY
tff = final.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "为学日益，为道日损。"
pf1.font.size = Pt(20)
pf1.font.bold = True
pf1.font.color.rgb = RGBColor(255, 255, 255)
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "—— 道德经第四十八章"
pf2.font.size = Pt(14)
pf2.font.color.rgb = RGBColor(200, 200, 200)
pf2.alignment = PP_ALIGN.CENTER
add_slide_number(s12, 12, 12)

prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第15天.pptx')
print("Done: 道德经_第15天.pptx (12 slides)")
