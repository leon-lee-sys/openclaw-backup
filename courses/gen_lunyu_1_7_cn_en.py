#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
论语 第1-7天精品课件 - 中英双语版（无拼音）
"""
import os
import sys
sys.path.insert(0, '/opt/homebrew/lib/node_modules/openclaw/node_modules')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(30, 60, 114)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)
GOLD = RGBColor(180, 140, 50)
ACCENT = RGBColor(180, 60, 60)

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_footer(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.2), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "%d/%d" % (num, total)
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def make_pptx(day, chapters, output_dir):
    prs = Presentation()
    titles = {
        1:"学而第一", 2:"为政第二", 3:"八佾第三", 4:"里仁第四",
        5:"公冶长第五", 6:"雍也第六", 7:"述而第七"
    }
    
    # Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    tx = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "国学经典精读 · 第%d天" % day
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    main = s.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
    tf2 = main.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "论语 · " + titles.get(day, "")
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.alignment = PP_ALIGN.CENTER
    
    add_footer(s, 1, 14)
    
    # Chapter slides
    for idx, ch in enumerate(chapters):
        slide_num = idx + 2
        
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_title_bar(s, "第%d章 · %s" % (ch["num"], ch["name"]))
        
        # Chinese original - larger font
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = "[原文]"
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        for line in ch["yuanwen"]:
            pp = tf.add_paragraph()
            pp.text = line
            pp.font.size = Pt(24)
            pp.font.color.rgb = TEXT
        
        # English translation
        t_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.5), Inches(9), Inches(1.3))
        t_box.fill.solid()
        t_box.fill.fore_color.rgb = BLUE_LIGHT
        tf = t_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "[Translation]"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(200, 220, 255)
        p2 = tf.add_paragraph()
        p2.text = ch["translation"]
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        
        # Core points
        y = 5.0
        for title, content in ch["core"]:
            # Title
            pt_box = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
            tf = pt_box.text_frame
            p1 = tf.paragraphs[0]
            p1.text = title
            p1.font.size = Pt(16)
            p1.font.bold = True
            p1.font.color.rgb = PRIMARY
            y += 0.35
            # Content
            p2 = tf.add_paragraph()
            p2.text = content
            p2.font.size = Pt(14)
            p2.font.color.rgb = TEXT
            y += 0.5
        
        add_footer(s, slide_num, 14)
    
    path = os.path.join(output_dir, "论语_第%d天.pptx" % day)
    prs.save(path)
    print("Done: 论语_第%d天.pptx" % day)
    return path

output_dir = "/Users/mac/.openclaw/workspace/courses/论语"
os.makedirs(output_dir, exist_ok=True)

CHAPTERS = [
    # Day 1: 学而第一
    [
        {"num": 1, "name": "学而时习", "yuanwen": ["子曰：学而时习之，不亦说乎？", "有朋自远方来，不亦乐乎？", "人不知而不愠，不亦君子乎？"], "translation": "The Master said: 'Is it not pleasant to learn with constant practice? Is it not delightful to have friends coming from distant quarters? Is it not a gentleman who is never angry when others do not understand him?'", "core": [("1. Learning and Practice", "学而时习：学习知识后时常复习，是喜悦的"), ("2. Friendship", "有朋自远方来：志同道合的朋友是人生的财富"), ("3. Patience", "人不知而不愠：别人不了解我，我却不恼怒")]}
    ],
    # Day 2: 为政第二
    [
        {"num": 2, "name": "为政以德", "yuanwen": ["为政以德，譬如北辰，", "居其所而众星共之。"], "translation": "Govern with virtue, and the people will be like the North Star - it keeps its place while all the stars pay homage to it.", "core": [("1. Rule by Virtue", "为政以德：用德行治理国家"), ("2. Like the North Star", "譬如北辰：有德的领导者像北极星一样"), ("3. Natural Attraction", "众星共之：不需要强迫，人心自然归附")]}
    ],
    # Day 3: 八佾第三
    [
        {"num": 3, "name": "礼之用", "yuanwen": ["礼之用，和为贵。", "先王之道，斯为美，小大由之。"], "translation": "In the practice of ritual, harmony is most valuable. The ways of the ancient kings were excellent - both great and small matters followed this principle.", "core": [("1. Harmony in Ritual", "礼之用，和为贵：礼仪的最高境界是和谐"), ("2. Ancient Wisdom", "先王之道：古代圣王的治国之道"), ("3. Proportional Application", "小大由之：无论大事小事都依此原则")]}
    ],
    # Day 4: 里仁第四
    [
        {"num": 4, "name": "里仁为美", "yuanwen": ["里仁为美，择不处仁，", "焉得知？"], "translation": "To dwell in the neighborhood of benevolence is beautiful. If you do not choose to live among the benevolent, how can you be wise?", "core": [("1. Living with Benevolence", "里仁为美：居住在有仁德的地方是美好的"), ("2. Wise Choice", "择不处仁：选择住处要选择有仁德之处"), ("3. True Wisdom", "焉得知：这才是真正的智慧")]}
    ],
    # Day 5: 公冶长第五
    [
        {"num": 5, "name": "君子之道", "yuanwen": ["子谓南容：邦有道不废，邦无道免于刑戮。", "以其兄之子妻之。"], "translation": "Confucius said of Nanrong: 'In a state with good government, he would not be dismissed; in a state with bad government, he would escape punishment.' He married his niece to him.", "core": [("1. Adaptability", "邦有道不废：国家政治清明时不会被废弃"), ("2. Prudence", "邦无道免于刑戮：国家混乱时能免于刑罚"), ("3. Family Trust", "以其兄之子妻之：把侄女嫁给他，说明孔子对他的信任")]}
    ],
    # Day 6: 雍也第六
    [
        {"num": 6, "name": "君子周急", "yuanwen": ["子华使于齐，冉子为其母请粟。", "子曰：与之釜。请益，曰：与之庾。"], "translation": "Zihua was sent on a mission to Qi. Ran Zi asked for grain for his mother. The Master said: 'Give her one fu.' Ran Zi asked for more. The Master said: 'Give her one yu.'", "core": [("1. Appropriate Reward", "子华使齐：公西赤出使齐国"), ("2. Right Amount", "冉子请粟：冉有为公西赤母亲请求粮食"), ("3. Proportional Giving", "与之釜庾：适量给予，不过多也不过少")]}
    ],
    # Day 7: 述而第七
    [
        {"num": 7, "name": "三人行", "yuanwen": ["三人行，必有我师焉。", "择其善者而从之，", "其不善者而改之。"], "translation": "When I walk with two others, they may be my teachers. I will select their good qualities and follow them, and their bad qualities and avoid them.", "core": [("1. Learning from Others", "三人行必有我师：每个人身上都有值得学习的"), ("2. Follow the Good", "择其善者而从之：选择别人的优点来学习"), ("3. Self-Reflection", "其不善者而改之：以别人的缺点为镜子，改正自己")]}
    ]
]

for i, chapters in enumerate(CHAPTERS):
    day = i + 1
    make_pptx(day, chapters, output_dir)

print("\n第1-7天论语课件（中英双语版）生成完成！")
