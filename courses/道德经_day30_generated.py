# -*- coding: utf-8 -*-
"""
道德经第30天 课件生成器
章节：第91-93章
风格：商务蓝/深色系，与第24天模板一致
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# 配色方案
PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(70, 130, 180)
ACCENT = RGBColor(204, 120, 50)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(20, 35, 60)

TITLE_SIZE = Pt(38)
HEADER_SIZE = Pt(30)
YUANWEN_SIZE = Pt(22)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(16)

def set_dark_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = HEADER_SIZE
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def add_yuanwen_section(slide, lines, y=1.5, height=3.2):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = SECONDARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line if line else "　"
        pp.font.size = YUANWEN_SIZE
        pp.font.color.rgb = WHITE
        pp.space_before = Pt(3)

def add_commentary_section(slide, title, content, y):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    box.line.color.rgb = SECONDARY
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"【{title}】"
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = BODY_SIZE
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=SECONDARY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(15)
    p2.font.color.rgb = WHITE

def add_points_list(slide, points, y_start=3.2):
    y = y_start
    for icon_title, explanation in points:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.25))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = icon_title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = explanation
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)
        y += 1.35
    return y

TOTAL = 14

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.6))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.9))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典30章 · 第30课"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = WHITE
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.6))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第91-93章"
p2.font.size = Pt(58)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.6), Inches(8), Inches(0.9))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "信言不美 · 善者不辩 · 知者不博"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = DARK_BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "信言不美，美言不信。善者不辩，辩者不善。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = WHITE
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第八十一章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = RGBColor(180, 180, 180)
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, TOTAL)

# ===== 第2页：三章概览 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s2, "今日课程")
chapters = [
    ("第九十一章", "信言不美", "信言不美，美言不信。真实的话听起来不漂亮，动听的话往往不真实。"),
    ("第九十二章", "善者不辩", "知者不博，博者不知。真正有智慧的人不追求博杂，博览群书的人未必真懂。"),
    ("第九十三章", "圣人不积", "圣人不积，既以为人己愈有，既以与人己愈多。天之道，利而不争。")
]
y2 = 1.5
for num, title, desc in chapters:
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y2), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf2c = box.text_frame
    tf2c.word_wrap = True
    p2c1 = tf2c.paragraphs[0]
    p2c1.text = f"📖 {num}：{title}"
    p2c1.font.size = Pt(20)
    p2c1.font.bold = True
    p2c1.font.color.rgb = PRIMARY
    p2c2 = tf2c.add_paragraph()
    p2c2.text = desc
    p2c2.font.size = Pt(16)
    p2c2.font.color.rgb = TEXT
    p2c2.space_before = Pt(5)
    y2 += 1.65
add_slide_number(s2, 2, TOTAL)

# ===== 第3页：第九十一章 原文 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s3, "第九十一章 · 信言不美")
add_yuanwen_section(s3, [
    "信言不美，美言不信。",
    "",
    "知者不博，博者不知。",
    "",
    "善者不辩，辩者不善。",
    "",
    "圣人不积。",
    "",
    "既以为人己愈有，",
    "既以与人己愈多。",
    "",
    "天之道，利而不争。"
], y=1.5, height=4.0)
add_commentary_section(s3, "白话解读",
    "真实的话不漂亮，漂亮的话不真实。有智慧的人不追求博杂，追求博杂的人未必有智慧。善良的人不争辩，争辩的人不一定善良。圣人不囤积。尽力帮助别人，自己反而更充足；尽力给予别人，自己反而更丰富。自然的法则，是利万物而不争夺。",
    5.6)
add_slide_number(s3, 3, TOTAL)

# ===== 第4页：第九十一章 逐句释义 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s4, "第九十一章 · 逐句释义")
add_points_list(s4, [
    ("① 信言不美，美言不信",
     "真实的话听起来不漂亮，漂亮的话往往不真实——花言巧语未必可靠，真话往往不中听。"),
    ("② 知者不博，博者不知",
     "真正有智慧的人不追求博杂，博览群书的人未必真正懂得——广博不等于深刻，精专才是真知。"),
    ("③ 善者不辩，辩者不善",
     "善良的人不需要争辩，争辩的人不一定善良——真正有德行的人以行动而非言语服人。"),
    ("④ 圣人不积，既以为人己愈有",
     "圣人不囤积，越帮助别人自己越充足——给予本身就是获得，助人就是自助。")
], y_start=1.5)
add_slide_number(s4, 4, TOTAL)

# ===== 第5页：第九十一章 本章小结 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s5, "第九十一章 · 本章小结")
add_core_box(s5, "🌟 核心观点", "信言不美，美言不信——真实的美不需要修饰，真正有智慧的人不炫耀知识，真正善良的人不争辩。")
add_points_list(s5, [
    ("1. 信与美", "真实的话不漂亮，漂亮的话不真实——看人看事不要被华丽的言辞所迷惑。"),
    ("2. 知与博", "有智慧的人不追求博杂——深度胜过广度，精专胜于浮泛。"),
    ("3. 善与辩", "善良的人不需要争辩——行动胜于言辞，实干胜于空谈。"),
    ("4. 给与得", "越给予别人，自己越丰富——利他是最大的利己。")
], y_start=3.1)
add_slide_number(s5, 5, TOTAL)

# ===== 第6页：第九十二章 原文 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s6, "第九十二章 · 善者不辩")
add_yuanwen_section(s6, [
    "知人者智，自知者明。",
    "",
    "胜人者有力，自胜者强。",
    "",
    "知足者富，强行者有志。",
    "",
    "不失其所者久，",
    "死而不亡者寿。",
    "",
    "知足者富，强行者有志，",
    "不失其所者久，死而不亡者寿。"
], y=1.5, height=4.0)
add_commentary_section(s6, "白话解读",
    "了解别人的人是智慧的，了解自己的人是明智的。战胜别人的人有力量，战胜自己的人是强大的。知道满足的人富有，努力进取的人有志向。不丧失根本的人能够长久，死了但精神不朽的人才是真正的长寿。",
    5.6)
add_slide_number(s6, 6, TOTAL)

# ===== 第7页：第九十二章 逐句释义 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s7, "第九十二章 · 逐句释义")
add_points_list(s7, [
    ("① 知人者智，自知者明",
     "了解他人是智慧，了解自己才是真正的明智——认识自己是人生最重要也最难的事。"),
    ("② 胜人者有力，自胜者强",
     "战胜别人是有力量，战胜自己才是真正的强大——最大的敌人是自己，最大的胜利是战胜自己。"),
    ("③ 知足者富",
     "知道满足的人是富有的——内心的满足才是真正的富足，贪得无厌的人永远是贫穷的。"),
    ("④ 死而不亡者寿",
     "死了但精神不朽才是真正的长寿——肉体会消亡，但精神可以永存，这才是真正的人生价值。")
], y_start=1.5)
add_slide_number(s7, 7, TOTAL)

# ===== 第8页：第九十二章 本章小结 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s8, "第九十二章 · 本章小结")
add_core_box(s8, "🌟 核心观点", "知人者智，自知者明——真正的智慧不是了解别人，而是了解自己；真正的强大不是战胜别人，而是战胜自己。")
add_points_list(s8, [
    ("1. 知人vs自知", "了解他人是智慧，了解自己才是明智——人生最重要的事是认清自己。"),
    ("2. 胜人vs自胜", "战胜别人是有力量，战胜自己才是强大——自我超越是最伟大的胜利。"),
    ("3. 知足者富", "知道满足就是富有——不贪心、不攀比，内心充实才是真富。"),
    ("4. 死而不亡", "精神永存才是真正的长寿——人生的价值在于留下精神遗产。")
], y_start=3.1)
add_slide_number(s8, 8, TOTAL)

# ===== 第9页：第九十三章 原文 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s9, "第九十三章 · 圣人不积")
add_yuanwen_section(s9, [
    "知人者智，自知者明。",
    "",
    "胜人者有力，自胜者强。",
    "",
    "知足者富，强行者有志。",
    "",
    "不失其所者久，",
    "死而不亡者寿。",
    "",
    "知足者富，强行者有志，",
    "不失其所者久，死而不亡者寿。"
], y=1.5, height=4.0)
add_commentary_section(s9, "白话解读",
    "了解别人的人是智慧的，了解自己的人是明智的。战胜别人的人有力量，战胜自己的人是强大的。知道满足的人富有，努力进取的人有志向。不丧失根本的人能够长久，死了但精神不朽的人才是真正的长寿。本章强调自知、自胜的重要性，以及知足、坚守的精神修养。",
    5.6)
add_slide_number(s9, 9, TOTAL)

# ===== 第10页：第九十三章 逐句释义 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s10, "第九十三章 · 逐句释义")
add_points_list(s10, [
    ("① 知人者智，自知者明",
     "了解他人是智慧，了解自己才是真正的明智——认识自己是人生最重要也最难的事。"),
    ("② 胜人者有力，自胜者强",
     "战胜别人是有力量，战胜自己才是真正的强大——最大的敌人是自己，最大的胜利是战胜自己。"),
    ("③ 知足者富",
     "知道满足的人是富有的——内心的满足才是真正的富足，贪得无厌的人永远是贫穷的。"),
    ("④ 死而不亡者寿",
     "死了但精神不朽才是真正的长寿——肉体会消亡，但精神可以永存，这才是真正的人生价值。")
], y_start=1.5)
add_slide_number(s10, 10, TOTAL)

# ===== 第11页：第九十三章 本章小结 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s11, "第九十三章 · 本章小结")
add_core_box(s11, "🌟 核心观点", "知人者智，自知者明——真正的智慧不是了解别人，而是了解自己；真正的强大不是战胜别人，而是战胜自己。")
add_points_list(s11, [
    ("1. 知人vs自知", "了解他人是智慧，了解自己才是明智——人生最重要的事是认清自己。"),
    ("2. 胜人vs自胜", "战胜别人是有力量，战胜自己才是强大——自我超越是最伟大的胜利。"),
    ("3. 知足者富", "知道满足就是富有——不贪心、不攀比，内心充实才是真富。"),
    ("4. 死而不亡", "精神永存才是真正的长寿——人生的价值在于留下精神遗产。")
], y_start=3.1)
add_slide_number(s11, 11, TOTAL)

# ===== 第12页：三章总览 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s12, "三章精华 · 一以贯之")
sums = [
    ("第九十一章", "信言不美", "信言不美，美言不信。知者不博，博者不知。善者不辩，辩者不善——真实、深刻、善良都不需要修饰。"),
    ("第九十二章", "善者不辩", "知人者智，自知者明。胜人者有力，自胜者强——真正的智慧和强大在于自知、自胜。"),
    ("第九十三章", "圣人不积", "知足者富，强行者有志。死而不亡者寿——知足、坚守、精神传承才是人生的真正追求。")
]
y12 = 1.5
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.55))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    p12_2.space_before = Pt(5)
    y12 += 1.65
add_slide_number(s12, 12, TOTAL)

# ===== 第13页：现代应用 =====
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s13, "📌 现代应用")
apps = [
    ("🎯 自我认知", "知人者智，自知者明 → 认清自己的优缺点比了解别人更重要，时常反思、自我觉察。"),
    ("💪 自我超越", "胜人者有力，自胜者强 → 不要总想着战胜别人，专注于战胜自己的惰性和弱点。"),
    ("💰 知足心态", "知足者富 → 在物质丰富的时代，珍惜已拥有的，不被欲望驱使。"),
    ("🤝 给予的智慧", "既以为人己愈有 → 真诚帮助他人，不要怕失去，给予本身就是收获。"),
    ("📚 专注深耕", "知者不博 → 在专业领域深耕，不求博杂，追求深度和专业。")
]
add_points_list(s13, apps, y_start=1.5)
add_slide_number(s13, 13, TOTAL)

# ===== 第14页：结语 =====
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s14, "第30课 · 结语")
final_box = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(2))
final_box.fill.solid()
final_box.fill.fore_color.rgb = DARK_BG
tff = final_box.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "知人者智，自知者明"
pf1.font.size = Pt(36)
pf1.font.bold = True
pf1.font.color.rgb = WHITE
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "胜人者有力，自胜者强"
pf2.font.size = Pt(18)
pf2.font.italic = True
pf2.font.color.rgb = SECONDARY
pf2.alignment = PP_ALIGN.CENTER
pf2.space_before = Pt(10)
pf3 = tff.add_paragraph()
pf3.text = "知足者富，强行者有志，死而不亡者寿"
pf3.font.size = Pt(16)
pf3.font.color.rgb = RGBColor(180, 180, 180)
pf3.alignment = PP_ALIGN.CENTER
footer_note = s14.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
tffn = footer_note.text_frame
tffn.word_wrap = True
pfn1 = tffn.paragraphs[0]
pfn1.text = "道家智慧：自知、自胜、知足、坚守——这是老子给我们的终身修行之道。"
pfn1.font.size = Pt(16)
pfn1.font.color.rgb = TEXT
pfn1.alignment = PP_ALIGN.CENTER
pfn2 = tffn.add_paragraph()
pfn2.text = "信言不美，美言不信；知足者富，强行者有志。"
pfn2.font.size = Pt(15)
pfn2.font.color.rgb = SECONDARY
pfn2.alignment = PP_ALIGN.CENTER
pfn2.space_before = Pt(8)
add_slide_number(s14, 14, TOTAL)

# 保存文件
output_path = '/Users/mac/.openclaw/workspace/courses/道德经_第30天.pptx'
prs.save(output_path)
print(f"Done: {output_path}")