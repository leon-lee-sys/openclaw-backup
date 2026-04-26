#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
道德经 第14天（43-45章）课件生成
风格：学术白
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

COLORS = {
    'primary': RGBColor(0, 51, 102),
    'secondary': RGBColor(102, 102, 102),
    'accent': RGBColor(178, 34, 34),
    'text': RGBColor(51, 51, 51),
    'bg': RGBColor(255, 255, 255),
    'header_bg': RGBColor(0, 51, 102),
}

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['bg']

def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['header_bg']
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_box(slide, left, top, width, height, fill_color, border_color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = border_color
    return box

# ===== 第1页：封面 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "《道德经》每日一课"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "第14天 · 第43-45章"
p2.font.size = Pt(26)
p2.font.color.rgb = COLORS['secondary']
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(16)

p3 = tf.add_paragraph()
p3.text = "知者不言 · 为学日益 · 大成若缺"
p3.font.size = Pt(20)
p3.font.color.rgb = COLORS['accent']
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(20)

p4 = tf.add_paragraph()
p4.text = "2026年4月21日"
p4.font.size = Pt(16)
p4.font.color.rgb = COLORS['secondary']
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(40)

bottom = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
btf = bottom.text_frame
bp = btf.paragraphs[0]
bp.text = "「上善若水，水善利万物而不争。」"
bp.font.size = Pt(14)
bp.font.italic = True
bp.font.color.rgb = COLORS['secondary']
bp.alignment = PP_ALIGN.CENTER

# ===== 第2页：第四十三章 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "第四十三章 · 知者不言")

box1 = add_box(slide, 0.4, 1.3, 9.2, 1.6, RGBColor(240, 248, 255), COLORS['primary'])
tf1 = box1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "【原文】"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.space_before = Pt(6)
p = tf1.add_paragraph()
p.text = "天下皆知美之为美，斯恶已；皆知善之为善，斯不善已。故有无相生，难易相成，长短相较，高下相倾，音声相和，前后相随。"
p.font.size = Pt(17)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(10)

txBox2 = slide.shapes.add_textbox(Inches(0.4), Inches(3.1), Inches(9.2), Inches(2.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "【白话解读】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']
p = tf2.add_paragraph()
p.text = "天下人都知道美之所以为美，丑的观念就产生了；都知道善之所以为善，不善的观念也就产生了。所以有与无相互生成，难与易相互成就，长与短相互显现，高与下相互依存，音与声相互和谐，前与后相互跟随——这便是宇宙万象对立统一的根本规律。"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(8)
p.line_spacing = 1.5

txBox3 = slide.shapes.add_textbox(Inches(0.4), Inches(5.2), Inches(9.2), Inches(1.8))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "【精华提炼】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']
bullets = [
    "「相对而生」：一切价值都是相对比较而存在的",
    "「相生相成」：对立事物互相依存、互相成就",
    "「自然之美」：不刻意比较，顺其自然才是真"
]
for b in bullets:
    p = tf3.add_paragraph()
    p.text = f"• {b}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text']
    p.space_before = Pt(6)

# ===== 第3页：第四十四章 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "第四十四章 · 为学日益")

box2 = add_box(slide, 0.4, 1.3, 9.2, 1.8, RGBColor(240, 248, 255), COLORS['primary'])
tf4 = box2.text_frame
tf4.word_wrap = True
p = tf4.paragraphs[0]
p.text = "【原文】"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.space_before = Pt(6)
p = tf4.add_paragraph()
p.text = "为学日益，为道日损。损之又损，以至于无为，无为而无不为。"
p.font.size = Pt(17)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(10)
p = tf4.add_paragraph()
p.text = "取天下常以无事，及其有事，不足以取天下。"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(10)

txBox4 = slide.shapes.add_textbox(Inches(0.4), Inches(3.2), Inches(9.2), Inches(2.2))
tf5 = txBox4.text_frame
tf5.word_wrap = True
p = tf5.paragraphs[0]
p.text = "【白话解读】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']
p = tf5.add_paragraph()
p.text = "追求学问，知识每天都会增加；追求大道，欲望每天都会减少。减少再减少，一直到达自然无为的境界，无为反而能成就一切。治理天下要靠不妄为，如果政令繁苛，就不足以治理天下了。"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(8)
p.line_spacing = 1.4

txBox5 = slide.shapes.add_textbox(Inches(0.4), Inches(5.5), Inches(9.2), Inches(1.5))
tf6 = txBox5.text_frame
tf6.word_wrap = True
p = tf6.paragraphs[0]
p.text = "【精华提炼】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']
bullets2 = [
    "「为学日益」：知识技能需要不断积累",
    "「为道日损」：修心养性需要不断放下执念",
    "「无为而无不为」：不妄为反而什么都做成了"
]
for b in bullets2:
    p = tf6.add_paragraph()
    p.text = f"• {b}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text']
    p.space_before = Pt(5)

# ===== 第4页：第四十五章 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "第四十五章 · 大成若缺")

box3 = add_box(slide, 0.4, 1.3, 9.2, 1.5, RGBColor(240, 248, 255), COLORS['primary'])
tf7 = box3.text_frame
tf7.word_wrap = True
p = tf7.paragraphs[0]
p.text = "【原文】"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.space_before = Pt(6)
p = tf7.add_paragraph()
p.text = "大成若缺，其用不弊。大盈若冲，其用不穷。大直若屈，大巧若拙，大辩若讷。"
p.font.size = Pt(17)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(10)
p = tf7.add_paragraph()
p.text = "躁胜寒，静胜热。清静为天下正。"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(10)

txBox6 = slide.shapes.add_textbox(Inches(0.4), Inches(3.0), Inches(9.2), Inches(2.0))
tf8 = txBox6.text_frame
tf8.word_wrap = True
p = tf8.paragraphs[0]
p.text = "【白话解读】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']
p = tf8.add_paragraph()
p.text = "最完美的东西好像有欠缺，但它的作用不会衰竭。最充实的东西好像有空虚，但它的作用不会穷尽。最正直的东西好像有弯曲，最灵巧的东西好像很笨拙，最善辩的人好像言语迟钝。躁动能战胜寒冷，安静能战胜炎热。清静无为才是天下的正道。"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(8)
p.line_spacing = 1.4

txBox7 = slide.shapes.add_textbox(Inches(0.4), Inches(5.1), Inches(9.2), Inches(1.9))
tf9 = txBox7.text_frame
tf9.word_wrap = True
p = tf9.paragraphs[0]
p.text = "【精华提炼】"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLORS['accent']
bullets3 = [
    "「大成若缺」：真正完美的事物看似有缺陷",
    "「大巧若拙」：真正高明的人看似笨拙",
    "「清静为天下正」：内心清静是治天下的根本"
]
for b in bullets3:
    p = tf9.add_paragraph()
    p.text = f"• {b}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text']
    p.space_before = Pt(6)

# ===== 第5页：典故案例 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "典故案例 · 扁鹊的医术")

txBox8 = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.6))
tf10 = txBox8.text_frame
p = tf10.paragraphs[0]
p.text = "「大巧若拙」—— 真正的高手，往往不动声色"
p.font.size = Pt(15)
p.font.italic = True
p.font.color.rgb = COLORS['secondary']
p.alignment = PP_ALIGN.CENTER

box4 = add_box(slide, 0.4, 2.0, 9.2, 4.5, RGBColor(250, 250, 250), COLORS['secondary'])
tf11 = box4.text_frame
tf11.word_wrap = True

p = tf11.paragraphs[0]
p.text = "《鹖冠子》"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['accent']
p.space_before = Pt(12)
p.alignment = PP_ALIGN.CENTER

p = tf11.add_paragraph()
p.text = "魏文侯问扁鹊：'你们三兄弟都精于医术，谁的医术最高？'扁鹊说：'大哥最高，二哥次之，我最差。'"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(14)
p.line_spacing = 1.6

p = tf11.add_paragraph()
p.text = "文侯问：'为什么？'扁鹊答：'大哥治病，在病情发作之前就铲除病根，病人不知道自己有病，所以不以为功。我二哥治病，在病初起时症状尚轻，病人以为只能治小病，所以名气只及乡里。而我治病，往往在病情严重时，用针灸、刺血、敷药等大动作，病人和家属都以为我医术高超，所以名闻天下。'"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(14)
p.line_spacing = 1.6

p = tf11.add_paragraph()
p.text = "真正的高手，在事情还未恶化时已化解，却不为人知——正如「大成若缺，大巧若拙」。"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']
p.space_before = Pt(14)
p.line_spacing = 1.6

p = tf11.add_paragraph()
p.text = "✅ 印证第四十五章「清静为天下正」：最高明的治理，是消弭问题于无形"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['primary']
p.space_before = Pt(18)

# ===== 第6页：总结页 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "第43-45章 · 精华总结")

items = [
    ("第四十三章", "天下皆知", "美丑、善恶相对而生，有无、难易、长短、高下相互依存"),
    ("第四十四章", "为学日益", "求学知识日增，求道欲望日减，无为而无不为"),
    ("第四十五章", "大成若缺", "真正完美看似有缺，大巧若拙，清静为天下正"),
]

for i, (chapter, theme, desc) in enumerate(items):
    top = 1.4 + i * 1.7
    box = add_box(slide, 0.4, top, 9.2, 1.5, RGBColor(245, 248, 255), COLORS['primary'])
    tf_box = box.text_frame
    tf_box.word_wrap = True
    p = tf_box.paragraphs[0]
    p.text = f"{chapter}：{theme}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.space_before = Pt(10)
    p.alignment = PP_ALIGN.CENTER
    p = tf_box.add_paragraph()
    p.text = desc
    p.font.size = Pt(15)
    p.font.color.rgb = COLORS['text']
    p.space_before = Pt(8)
    p.alignment = PP_ALIGN.CENTER

bottom_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.4), Inches(7), Inches(0.8))
bottom_box.fill.solid()
bottom_box.fill.fore_color.rgb = COLORS['primary']
bottom_box.line.fill.background()
tf_btm = bottom_box.text_frame
p = tf_btm.paragraphs[0]
p.text = "💡 今日功课：体会「大成若缺」的智慧，放下完美执念，清静面对生活"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

output_dir = "/Users/mac/.openclaw/workspace/courses"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "道德经_第14天.pptx")
prs.save(output_path)
print(f"✅ 课件已保存：{output_path}")
