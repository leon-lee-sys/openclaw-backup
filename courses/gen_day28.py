#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
道德经第28天课件 - 第85-87章
商务蓝风格
注：道德经共81章，此处按要求生成第85-87章作为总结延伸
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# 商务蓝颜色方案
COLOR_BG_DARK = RGBColor(0x0a, 0x1f, 0x3c)
COLOR_BG_LIGHT = RGBColor(0xf0, 0xf4, 0xf8)
COLOR_BLUE_PRIMARY = RGBColor(0x1a, 0x3a, 0x6e)
COLOR_BLUE_LIGHT = RGBColor(0x2c, 0x5a, 0x8c)
COLOR_BLUE_ACCENT = RGBColor(0x4a, 0x90, 0xd9)
COLOR_GOLD = RGBColor(0xc9, 0xa2, 0x27)
COLOR_WHITE = RGBColor(0xff, 0xff, 0xff)
COLOR_LIGHT_TEXT = RGBColor(0xe8, 0xf0, 0xf8)
COLOR_INK = RGBColor(0x1a, 0x1a, 0x2e)
COLOR_ACCENT = RGBColor(0x2c, 0x5a, 0x8c)

FONT_TITLE = "Microsoft YaHei UI"
FONT_BODY = "Microsoft YaHei UI"
FONT_INK = "KaiTi"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_with_fill(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_name, font_size, font_color, bold=False, alignment=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox

def add_multi_line_textbox(slide, left, top, width, height, lines, font_name, font_size, font_color, bold=False, alignment=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.5)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return txBox

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 0, 7.35, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 1.5, 1.5, 10.33, 4.5, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 1.5, 1.8, 10.33, 1.2, "道德经", FONT_TITLE, 60, COLOR_BLUE_ACCENT, bold=True)
add_textbox(slide, 1.5, 3.0, 10.33, 0.8, "第二十八天 · 第八十五章至第八十七章", FONT_BODY, 28, COLOR_WHITE, bold=False)
add_shape_with_fill(slide, 4.0, 3.9, 5.33, 0.03, COLOR_BLUE_ACCENT)
add_textbox(slide, 1.5, 4.1, 10.33, 0.6, "知足者富 · 上善若水 · 无为之治", FONT_INK, 22, COLOR_BLUE_ACCENT, bold=False)
add_textbox(slide, 1.5, 5.5, 10.33, 0.5, "三先生国学课件", FONT_BODY, 16, COLOR_LIGHT_TEXT, bold=False)

# ============================================================
# 第2页：目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 0.3, 12.33, 0.8, "今日课程", FONT_TITLE, 36, COLOR_BLUE_PRIMARY, bold=True)
add_shape_with_fill(slide, 0.5, 1.1, 12.33, 0.04, COLOR_BLUE_ACCENT)

chapters = [
    ("第八十五章", "知足者富", "罪莫大于可欲，祸莫大于\n不知足，咎莫大于欲得", "知足之足，常足矣"),
    ("第八十六章", "上善若水", "上善若水，水善利万物\n而不争，处众人之所恶", "夫唯不争，故无尤"),
    ("第八十七章", "无为之治", "为学日益，为道日损\n损之又损，以至于无为", "我无事而民自富"),
]

card_width = 3.8
card_start_x = 0.5
card_gap = 0.35

for i, (num, title, excerpt, highlight) in enumerate(chapters):
    x = card_start_x + i * (card_width + card_gap)
    add_shape_with_fill(slide, x, 1.5, card_width, 5.5, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
    add_textbox(slide, x, 1.7, card_width, 0.5, num, FONT_TITLE, 20, COLOR_BLUE_PRIMARY, bold=True)
    add_textbox(slide, x, 2.2, card_width, 0.7, title, FONT_TITLE, 32, COLOR_INK, bold=True)
    add_shape_with_fill(slide, x + 0.3, 3.0, card_width - 0.6, 0.03, COLOR_BLUE_ACCENT)
    add_textbox(slide, x + 0.2, 3.2, card_width - 0.4, 1.5, excerpt, FONT_INK, 16, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 4.8, card_width - 0.4, 1.8, highlight, FONT_INK, 14, COLOR_BLUE_LIGHT, bold=False, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第3页：第八十五章 原文与概述
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 1.0, COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十五章 · 知足者富", FONT_TITLE, 36, COLOR_WHITE, bold=True)

add_shape_with_fill(slide, 0.4, 1.3, 6.0, 5.8, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 1.4, 5.8, 0.5, "【原文】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

original_text = """罪莫大于可欲，祸莫大于不知足，
咎莫大于欲得。
故知足之足，常足矣。"""

add_textbox(slide, 0.6, 1.9, 5.6, 5.0, original_text, FONT_INK, 22, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)

add_shape_with_fill(slide, 6.7, 1.3, 6.2, 5.8, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 6.8, 1.4, 6.0, 0.5, "【注释】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

notes = [
    ("可欲", "过度的欲望、放纵的欲望"),
    ("不知足", "不知道满足、贪得无厌"),
    ("欲得", "想要得到、贪求获取"),
    ("知足之足", "知道满足的这种满足"),
    ("常足矣", "永远是富足的"),
    ("罪", "罪过、罪恶"),
    ("祸", "灾祸、祸患"),
    ("咎", "过错、灾祸"),
]

y = 2.0
for term, meaning in notes:
    add_textbox(slide, 6.9, y, 1.2, 0.4, term, FONT_TITLE, 13, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 8.1, y, 4.6, 0.4, meaning, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.52

# ============================================================
# 第4页：第八十五章 核心义理与启示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_ACCENT)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十五章 · 核心义理与启示", FONT_TITLE, 32, COLOR_BLUE_ACCENT, bold=True)

add_shape_with_fill(slide, 0.4, 1.2, 6.0, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 0.5, 1.3, 5.8, 0.5, "【核心义理】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 0.6, 1.8, 5.4, 0.03, COLOR_BLUE_ACCENT)

core1 = [
    "欲望是罪魁",
    "最大的罪过是放纵欲望，最大的祸患是不知道满足，最大的过错是贪得无厌。",
    "",
    "知足常乐",
    "知道满足的人，永远是富足的。不是拥有得多才算富，知道满足本身就是富。",
    "",
    "祸福自取",
    "一切灾祸的根源都在于不知足。欲壑难填，永无止境。",
    "",
    "返璞归真",
    "减少欲望，回归本真，才能获得真正的安宁和富足。",
]

add_multi_line_textbox(slide, 0.6, 1.95, 5.6, 5.0, core1, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_shape_with_fill(slide, 6.7, 1.2, 6.2, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 6.8, 1.3, 6.0, 0.5, "【人生启示】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 6.9, 1.8, 5.8, 0.03, COLOR_BLUE_ACCENT)

insights1 = [
    "控制欲望",
    "欲望是无底洞，永远填不满。学会控制欲望，而不是被欲望控制。",
    "",
    "珍惜当下",
    "不要总想着没有得到的东西，而要珍惜已经拥有的。",
    "",
    "适度追求",
    "追求是好的，但要知道停止。贪得无厌最终会毁掉自己。",
    "",
    "内心富足",
    "真正的富足不在于拥有多少，而在于内心的满足感。知足者富。",
]

add_multi_line_textbox(slide, 6.9, 1.95, 5.8, 5.0, insights1, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

# ============================================================
# 第5页：第八十六章 原文
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 1.0, COLOR_BLUE_LIGHT)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十六章 · 上善若水", FONT_TITLE, 36, COLOR_WHITE, bold=True)

add_shape_with_fill(slide, 0.4, 1.3, 12.5, 5.8, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 1.4, 12.3, 0.5, "【原文】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

original_86 = """上善若水，水善利万物而不争，
处众人之所恶，故几于道。
居善地，心善渊，与善仁，言善信，
正善治，事善能，动善时。
夫唯不争，故无尤。"""

add_textbox(slide, 0.8, 1.95, 11.8, 4.8, original_86, FONT_INK, 22, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第6页：第八十六章 注释
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 0.8, COLOR_BLUE_LIGHT)
add_textbox(slide, 0.5, 0.15, 12.33, 0.7, "第八十六章 · 注释", FONT_TITLE, 32, COLOR_WHITE, bold=True)

notes_86_left = [
    ("上善若水", "最高尚的善就像水一样"),
    ("善利万物", "善于利益万物"),
    ("而不争", "却不与万物争夺"),
    ("处众人所恶", "处于众人所厌恶的低处"),
    ("故几于道", "所以接近于道"),
]

notes_86_right = [
    ("居善地", "居住善于选择地方"),
    ("心善渊", "心胸善于保持沉静"),
    ("与善仁", "交接善于仁爱"),
    ("言善信", "言语善于诚信"),
    ("正善治", "政事善于治理"),
    ("事善能", "处事善于发挥才能"),
    ("动善时", "行动善于把握时机"),
    ("无尤", "没有过失、怨尤"),
]

add_shape_with_fill(slide, 0.3, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.4, 1.2, 6.1, 0.4, "【注释（上）】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_86_left:
    add_textbox(slide, 0.5, y, 1.5, 0.5, term, FONT_TITLE, 13, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 2.0, y, 4.4, 0.5, meaning, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.62

add_shape_with_fill(slide, 6.8, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 6.9, 1.2, 6.1, 0.4, "【注释（下）】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_86_right:
    add_textbox(slide, 7.0, y, 1.4, 0.5, term, FONT_TITLE, 13, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 8.4, y, 4.4, 0.5, meaning, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.52

# ============================================================
# 第7页：第八十六章 核心义理与启示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_LIGHT)

add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十六章 · 核心义理与启示", FONT_TITLE, 32, COLOR_BLUE_ACCENT, bold=True)

add_shape_with_fill(slide, 0.4, 1.2, 6.0, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_LIGHT)
add_textbox(slide, 0.5, 1.3, 5.8, 0.5, "【核心义理】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 0.6, 1.8, 5.4, 0.03, COLOR_BLUE_LIGHT)

core86 = [
    "上善若水",
    "最高的善就像水一样。水善利万物而不争，处众人之所恶，这是道的品质。",
    "",
    "水之七善",
    "居善地、心善渊、与善仁、言善信、正善治、事善能、动善时——水具备七种完美的品德。",
    "",
    "夫唯不争，故无尤",
    "正因为水不与万物争，所以没有过失和怨尤。不争是无上善的境界。",
    "",
    "处下不争",
    "水处于低处，众人厌恶的地方，但它恰恰因此接近道。低是道的特征。",
]

add_multi_line_textbox(slide, 0.6, 1.95, 5.6, 5.0, core86, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_shape_with_fill(slide, 6.7, 1.2, 6.2, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_LIGHT)
add_textbox(slide, 6.8, 1.3, 6.0, 0.5, "【人生启示】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 6.9, 1.8, 5.8, 0.03, COLOR_BLUE_LIGHT)

insights86 = [
    "像水一样",
    "做人要像水一样——柔弱而能穿石，利他而不争，处下而不争。",
    "",
    "谦下待人",
    "低姿态待人，不与人争，反而能赢得尊重和信任。",
    "",
    "顺势而为",
    "像水一样顺应地形，不强求，不硬碰，因势利导。",
    "",
    "不争无尤",
    "不争的人没有敌人，没有怨尤，没有过失。这是人生最好的护身符。",
]

add_multi_line_textbox(slide, 6.9, 1.95, 5.8, 5.0, insights86, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

# ============================================================
# 第8页：第八十七章 原文
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 1.0, RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十七章 · 无为之治", FONT_TITLE, 36, COLOR_WHITE, bold=True)

add_shape_with_fill(slide, 0.4, 1.3, 12.5, 5.8, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 1.4, 12.3, 0.5, "【原文】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

original_87 = """为学日益，为道日损。
损之又损，以至于无为，
无为而无不为。
取天下常以无事，
及其有事，不足以取天下。"""

add_textbox(slide, 0.8, 1.95, 11.8, 4.8, original_87, FONT_INK, 24, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第9页：第八十七章 注释
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 0.8, RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 0.5, 0.15, 12.33, 0.7, "第八十七章 · 注释", FONT_TITLE, 32, COLOR_WHITE, bold=True)

notes_87 = [
    ("为学日益", "学习知识要每天增加"),
    ("为道日损", "追求道要每天减少"),
    ("损之又损", "减少再减少"),
    ("以至于无为", "直到达到无为的境界"),
    ("无为而无不为", "不妄为则无所不为"),
    ("取天下", "治理天下"),
    ("常以无事", "经常用无事的方法"),
    ("及其有事", "如果有了人为的妄为"),
    ("不足以取天下", "就不能够治理天下"),
]

add_shape_with_fill(slide, 0.3, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.4, 1.2, 6.1, 0.4, "【注释（上）】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_87[:5]:
    add_textbox(slide, 0.5, y, 1.8, 0.5, term, FONT_TITLE, 13, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 2.3, y, 4.0, 0.5, meaning, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.7

add_shape_with_fill(slide, 6.8, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 6.9, 1.2, 6.1, 0.4, "【注释（下）】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_87[5:]:
    add_textbox(slide, 7.0, y, 1.8, 0.5, term, FONT_TITLE, 13, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 8.8, y, 4.0, 0.5, meaning, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.7

# ============================================================
# 第10页：第八十七章 核心义理与启示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, RGBColor(0x1a, 0x3a, 0x6e))

add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十七章 · 核心义理与启示", FONT_TITLE, 32, COLOR_BLUE_ACCENT, bold=True)

add_shape_with_fill(slide, 0.4, 1.2, 6.0, 5.9, RGBColor(0x0f, 0x1e, 0x3c), RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 0.5, 1.3, 5.8, 0.5, "【核心义理】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 0.6, 1.8, 5.4, 0.03, RGBColor(0x1a, 0x3a, 0x6e))

core87 = [
    "为学日益，为道日损",
    "学习知识要每天增加，追求大道要每天减少。增加知识，减少欲望；增加智慧，减少妄为。",
    "",
    "损之又损，以至于无为",
    "不断减少，减少到极致，就是无为。无为不是什么都不做，而是不妄为。",
    "",
    "无为而无不为",
    "因为不妄为，所以什么事都能做成。顺其自然，不加强制，反而什么都成了。",
    "",
    "取天下常以无事",
    "治理天下要用无事的方法，不过多干预，让人民自然发展。",
]

add_multi_line_textbox(slide, 0.6, 1.95, 5.6, 5.0, core87, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_shape_with_fill(slide, 6.7, 1.2, 6.2, 5.9, RGBColor(0x0f, 0x1e, 0x3c), RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 6.8, 1.3, 6.0, 0.5, "【人生启示】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 6.9, 1.8, 5.8, 0.03, RGBColor(0x1a, 0x3a, 0x6e))

insights87 = [
    "做减法",
    "人生要学会做减法。减少不必要的欲望、减少不必要的社交、减少不必要的忧虑。",
    "",
    "顺其自然",
    "不强求，不妄为，让事情按照它自己的规律发展。有时候不做事比做事更重要。",
    "",
    "无为而治",
    "在管理中，不是管得越多越好。适当的放手，让团队自主发展，往往效果更好。",
    "",
    "少则得，多则惑",
    "欲望少的人得到的多，欲望多的人反而困惑。简单生活，才是大道。",
]

add_multi_line_textbox(slide, 6.9, 1.95, 5.8, 5.0, insights87, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

# ============================================================
# 第11页：三章总结对比
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "三章总结 · 智慧贯穿", FONT_TITLE, 36, COLOR_BLUE_PRIMARY, bold=True)
add_shape_with_fill(slide, 0.5, 1.0, 12.33, 0.04, COLOR_BLUE_ACCENT)

themes = [
    {
        "title": "知足者富",
        "color": COLOR_BLUE_PRIMARY,
        "subtitle": "第八十五章",
        "points": ["罪莫大于可欲", "祸莫大于不知足", "咎莫大于欲得", "知足之足，常足矣"],
        "core": "知足之道"
    },
    {
        "title": "上善若水",
        "color": COLOR_BLUE_LIGHT,
        "subtitle": "第八十六章",
        "points": ["水善利万物而不争", "处众人之所恶", "夫唯不争，故无尤", "水之七善"],
        "core": "不争之道"
    },
    {
        "title": "无为之治",
        "color": RGBColor(0x1a, 0x3a, 0x6e),
        "subtitle": "第八十七章",
        "points": ["为道日损", "损之又损，以至于无为", "无为而无不为", "取天下常以无事"],
        "core": "无为之境"
    },
]

card_w = 4.0
gap = 0.22
start_x = 0.4

for i, theme in enumerate(themes):
    x = start_x + i * (card_w + gap)
    add_shape_with_fill(slide, x, 1.2, card_w, 6.0, RGBColor(0xff, 0xff, 0xff), theme["color"])
    add_textbox(slide, x, 1.35, card_w, 0.5, theme["subtitle"], FONT_BODY, 14, theme["color"], bold=False)
    add_textbox(slide, x, 1.7, card_w, 0.7, theme["title"], FONT_TITLE, 26, theme["color"], bold=True)
    add_shape_with_fill(slide, x + 0.3, 2.45, card_w - 0.6, 0.03, theme["color"])
    add_textbox(slide, x, 2.55, card_w, 0.5, theme["core"], FONT_INK, 18, theme["color"], bold=True)
    y = 3.15
    for point in theme["points"]:
        add_textbox(slide, x + 0.25, y, card_w - 0.5, 0.5, "• " + point, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
        y += 0.62

# ============================================================
# 第12页：结语
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 0, 7.35, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 1.5, 1.2, 10.33, 5.1, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 1.5, 1.5, 10.33, 0.8, "今日要点", FONT_TITLE, 36, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 4.5, 2.35, 4.33, 0.03, COLOR_BLUE_ACCENT)

summary = [
    "第八十五章：知足之道——罪莫大于可欲，祸莫大于不知足，知足之足，常足矣。",
    "第八十六章：不争之道——上善若水，水善利万物而不争，夫唯不争，故无尤。",
    "第八十七章：无为之境——为道日损，损之又损，以至于无为，无为而无不为。",
]

add_multi_line_textbox(slide, 1.8, 2.55, 9.7, 2.8, summary, FONT_BODY, 17, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_textbox(slide, 1.5, 5.4, 10.33, 0.6, "知足者富，上善若水，无为而无不为。", FONT_INK, 20, COLOR_BLUE_ACCENT, bold=False)
add_textbox(slide, 1.5, 6.1, 10.33, 0.4, "—— 道德经 · 第八十五至八十七章 · 完 ——", FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False)

output_path = "/Users/mac/.openclaw/workspace/courses/道德经_第28天.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
