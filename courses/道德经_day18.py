# -*- coding: utf-8 -*-
"""
道德经第18天 课件生成器
章节：第58-60章
风格：精读深度讲解版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
ACCENT = RGBColor(204, 0, 0)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_yuanwen_box(slide, lines, y=1.4, height=2.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(20)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=BLUE_LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(255, 255, 255)

def add_points(slide, points, y_start=3.1):
    y = y_start
    for title, content in points:
        pt_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT
        y += 1.1
    return y

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.5))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第18课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第58-60章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "其政闷闷 · 治人事天 · 出生入死"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "其政闷闷，其民淳淳。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = PRIMARY
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第五十八章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = SECONDARY
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, 12)

# ===== 第2页：第五十八章详解 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s2, "第五十八章 · 其政闷闷")
add_yuanwen_box(s2, [
    "其政闷闷，其民淳淳；",
    "其政察察，其民缺缺。",
    "",
    "祸兮福之所倚，",
    "福兮祸之所伏。",
    "孰知其极？其无正。",
    "",
    "正复为奇，善复为妖。",
    "人之迷，其日固久。",
    "",
    "是以圣人方而不割，廉而不刿，",
    "直而不肆，光而不耀。"
], y=1.4, height=4.5)
content2 = s2.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(2))
tf_c2 = content2.text_frame
tf_c2.word_wrap = True
pc2 = tf_c2.paragraphs[0]
pc2.text = "【白话解读】政治宽厚，人民就淳朴；政治严苛，人民就狡诈。灾祸中藏着幸福的种子，幸福里也藏着灾祸。祸福相依，辩证看问题才是道的智慧。"
pc2.font.size = Pt(14)
pc2.font.color.rgb = TEXT
add_slide_number(s2, 2, 12)

# ===== 第3页：第五十八章核心 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s3, "第五十八章 · 核心思想")
add_core_box(s3, "🌟 核心观点", "灾祸中藏着幸福的种子，幸福里也藏着灾祸。祸福相依，辩证看问题才是道的智慧。")
add_points(s3, [
    ("1. 闷闷与察察", "政治宽厚则民淳朴，政治严苛则民狡诈——宽仁才能得人心。"),
    ("2. 祸福相依", "灾祸中藏着幸福的种子，幸福里也藏着灾祸——事物无不相互转化。"),
    ("3. 其无正", "没有绝对的正邪善恶——一切都在变化中，没有永恒的标准。"),
    ("4. 方而不割", "圣人方正而不伤人，有棱角而不刺人——外圆内方的智慧。")
])
add_slide_number(s3, 3, 12)

# ===== 第4页：第五十九章详解 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s4, "第五十九章 · 治人事天")
add_yuanwen_box(s4, [
    "治人事天，莫若啬。",
    "",
    "夫唯啬，是谓早服；",
    "早服谓之重积德；",
    "重积德则无不克；",
    "无不克则莫知其极；",
    "莫知其极，可以有国。",
    "",
    "有国之母，可以长久。",
    "是谓深根固柢、",
    "长生久视之道。"
], y=1.4, height=4.2)
content4 = s4.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(2))
tf_c4 = content4.text_frame
tf_c4.word_wrap = True
pc4 = tf_c4.paragraphs[0]
pc4.text = "【白话解读】治理人民、侍奉上天，没有比节俭更重要的了。节俭就是早做准备，早做准备就是不断积累德行。积累德行就没有什么不能克服的，没有什么不能克服就无法估计其力量——这样才能保有国家。"
pc4.font.size = Pt(14)
pc4.font.color.rgb = TEXT
add_slide_number(s4, 4, 12)

# ===== 第5页：第五十九章核心 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s5, "第五十九章 · 核心思想")
add_core_box(s5, "🌟 核心观点", "啬就是节俭、收敛。及早顺应天道，不断积累德行，才能深根固柢，长生久视。")
add_points(s5, [
    ("1. 啬", "啬就是节俭、收敛——不过度消耗，才能细水长流。"),
    ("2. 早服", "及早顺应天道——早做准备，防患于未然。"),
    ("3. 重积德", "不断积累德行——德行深厚才能拥有无穷的力量。"),
    ("4. 深根固柢", "根深叶茂，基业长青——养生和治国都是一个道理。")
])
add_slide_number(s5, 5, 12)

# ===== 第6页：现代反思 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s6, "第五十八-五十九章 · 现代反思")
items6 = [
    ("🏛️ 管理哲学", "其政闷闷 → 好的管理是无声的渗透，不是处处设卡。"),
    ("💰 理财智慧", "啬 → 赚钱重要，存钱更重要。细水长流才能真正富裕。"),
    ("🔄 祸福转化", "祸兮福所倚 → 失败中藏着成功的种子，顺境中也要警惕风险。"),
    ("🌱 长期主义", "深根固柢 → 无论是养生、理财还是事业，都需要长期积累，不追求速成。")
]
y6 = 1.4
for icon_title, explanation in items6:
    box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y6), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf6 = box.text_frame
    tf6.word_wrap = True
    p6_1 = tf6.paragraphs[0]
    p6_1.text = icon_title
    p6_1.font.size = Pt(16)
    p6_1.font.bold = True
    p6_1.font.color.rgb = PRIMARY
    p6_2 = tf6.add_paragraph()
    p6_2.text = explanation
    p6_2.font.size = Pt(14)
    p6_2.font.color.rgb = TEXT
    y6 += 1.4
add_slide_number(s6, 6, 12)

# ===== 第7页：第六十章详解 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s7, "第六十章 · 治大国若烹小鲜")
add_yuanwen_box(s7, [
    "治大国，若烹小鲜。",
    "",
    "以道莅天下，",
    "其鬼不神。",
    "非其鬼不神，",
    "其神不伤人；",
    "非其神不伤人，",
    "圣人亦不伤人。",
    "",
    "夫两不相伤，",
    "故德交归焉。"
], y=1.4, height=4.0)
content7 = s7.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(2))
tf_c7 = content7.text_frame
tf_c7.word_wrap = True
pc7 = tf_c7.paragraphs[0]
pc7.text = "【白话解读】治理大国，像煎小鱼一样——不能频繁翻动，否则会碎。以道来治理天下，鬼神都不灵了。不是鬼神不灵，而是它们不伤害人；不是它们不伤害人，而是圣人也根本不伤害人。双方互不伤害，德行就回归于道了。"
pc7.font.size = Pt(14)
pc7.font.color.rgb = TEXT
add_slide_number(s7, 7, 12)

# ===== 第8页：第六十章核心 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s8, "第六十章 · 核心思想")
add_core_box(s8, "🌟 核心观点", "治理大国若烹小鲜——不要折腾，顺其自然，无为而治，德行归道。")
add_points(s8, [
    ("1. 烹小鲜", "煎小鱼不能翻动太多次——治理国家不能朝令夕改，政策要稳定。"),
    ("2. 以道莅天下", "以道的原则治理天下——顺应规律，不妄为。"),
    ("3. 两不相伤", "圣人和百姓互不伤害——最好的治理是让所有人都能各安其位。"),
    ("4. 德交归焉", "德行回归于道——无为而治的终极目标。")
])
add_slide_number(s8, 8, 12)

# ===== 第9页：西方思想 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s9, "🌍 西方思想对照")
phils = [
    ("黑格尔 · 辩证法", "正反合——事物内部矛盾推动发展，正与反相互依存。1807年"),
    ("尼采 · 永恒轮回", "一切都在无限循环中，祸福交替是自然的永恒法则。1882年"),
    ("老子与赫拉克利特", "变才是唯一的不变——东西方哲人都发现了变化的本质。约前500年"),
    ("海森堡 · 测不准原理", "观测本身影响结果——人的介入会改变事物的走向。1927年")
]
y9 = 1.4
for name_theory, desc_year in phils:
    box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y9), Inches(9), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf9 = box.text_frame
    tf9.word_wrap = True
    p9_1 = tf9.paragraphs[0]
    p9_1.text = name_theory
    p9_1.font.size = Pt(16)
    p9_1.font.bold = True
    p9_1.font.color.rgb = PRIMARY
    p9_2 = tf9.add_paragraph()
    p9_2.text = desc_year
    p9_2.font.size = Pt(14)
    p9_2.font.color.rgb = TEXT
    y9 += 1.3
add_slide_number(s9, 9, 12)

# ===== 第10页：现代应用 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s10, "📌 现代应用")
add_points(s10, [
    ("企业管理", "烹小鲜 → 战略一旦确定，就要保持稳定，不要总变来变去。"),
    ("投资理财", "重积德 → 不追求一夜暴富，慢慢积累复利才是王道。"),
    ("人际交往", "方而不割 → 待人要圆融，有原则但不伤人。"),
    ("个人修养", "啬 → 精力有限，要啬己，不要过度消耗在无意义的事上。")
])
add_slide_number(s10, 10, 12)

# ===== 第11页：中国典故 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s11, "📜 中国典故：刘邦的\"约法三章\"")
box11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.5))
box11.fill.solid()
box11.fill.fore_color.rgb = BG
tf11 = box11.text_frame
tf11.word_wrap = True
p11_1 = tf11.paragraphs[0]
p11_1.text = "【史记·高祖本纪】"
p11_1.font.size = Pt(16)
p11_1.font.bold = True
p11_1.font.color.rgb = PRIMARY
p11_2 = tf11.add_paragraph()
p11_2.text = "刘邦率军攻入咸阳时，与秦民约法三章：'杀人者死，伤人及盗抵罪。'其余繁苛法令一律废除。这一简单的约定，让秦地百姓大喜，纷纷犒劳军队。刘邦凭借这简单的\"约法三章\"赢得了民心，为日后建立汉朝奠定了基础。"
p11_2.font.size = Pt(15)
p11_2.font.color.rgb = TEXT
p11_2.space_before = Pt(10)
insight11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.1), Inches(9), Inches(1))
insight11.fill.solid()
insight11.fill.fore_color.rgb = BLUE_LIGHT
tfi11 = insight11.text_frame
tfi11.word_wrap = True
pi11 = tfi11.paragraphs[0]
pi11.text = "💡 启示"
pi11.font.size = Pt(18)
pi11.font.bold = True
pi11.font.color.rgb = RGBColor(255, 255, 255)
pi11_2 = tfi11.add_paragraph()
pi11_2.text = "\"烹小鲜\"的智慧——简单稳定的政策比复杂多变的法令更能得人心。治理如此，待人也是如此。"
pi11_2.font.size = Pt(14)
pi11_2.font.color.rgb = RGBColor(255, 255, 255)
app11 = s11.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2))
tfa11 = app11.text_frame
tfa11.word_wrap = True
pa11 = tfa11.paragraphs[0]
pa11.text = "【现代应用】职场：简洁直接的管理方式比繁文缛节更有效。生活：简单的生活习惯更容易坚持。交友：真诚简单的交往更深得人心。"
pa11.font.size = Pt(14)
pa11.font.color.rgb = TEXT
add_slide_number(s11, 11, 12)

# ===== 第12页：总结 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s12, "第18课 · 三章精华总结")
sums = [
    ("第五十八章", "其政闷闷", "祸福相依，辩证看世界。方而不割，光而不耀。"),
    ("第五十九章", "治人事天", "啬为根本，重积德，深根固柢，长生久视。"),
    ("第六十章", "治大国若烹小鲜", "以道莅天下，无为而治，两不相伤，德交归焉。")
]
y12 = 1.4
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    y12 += 1.6
final = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(9), Inches(1))
final.fill.solid()
final.fill.fore_color.rgb = PRIMARY
tff = final.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "治大国，若烹小鲜。"
pf1.font.size = Pt(20)
pf1.font.bold = True
pf1.font.color.rgb = RGBColor(255, 255, 255)
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "—— 道德经第六十章"
pf2.font.size = Pt(14)
pf2.font.color.rgb = RGBColor(200, 200, 200)
pf2.alignment = PP_ALIGN.CENTER
add_slide_number(s12, 12, 12)

prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第18天.pptx')
print("Done: 道德经_第18天.pptx (12 slides)")
