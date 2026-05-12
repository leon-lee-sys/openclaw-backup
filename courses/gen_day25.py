#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
道德经第25天课件 - 第79-81章
商务蓝风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# 商务蓝颜色方案
COLOR_BG_DARK = RGBColor(0x0a, 0x1f, 0x3c)      # 深蓝背景
COLOR_BG_LIGHT = RGBColor(0xf0, 0xf4, 0xf8)     # 浅蓝灰背景
COLOR_BLUE_PRIMARY = RGBColor(0x1a, 0x3a, 0x6e) # 主蓝
COLOR_BLUE_LIGHT = RGBColor(0x2c, 0x5a, 0x8c)    # 中蓝
COLOR_BLUE_ACCENT = RGBColor(0x4a, 0x90, 0xd9)  # 亮蓝
COLOR_GOLD = RGBColor(0xc9, 0xa2, 0x27)         # 金色点缀
COLOR_WHITE = RGBColor(0xff, 0xff, 0xff)
COLOR_LIGHT_TEXT = RGBColor(0xe8, 0xf0, 0xf8)   # 浅色文字
COLOR_INK = RGBColor(0x1a, 0x1a, 0x2e)           # 墨色
COLOR_ACCENT = RGBColor(0x2c, 0x5a, 0x8c)       # 蓝色强调

# 字体
FONT_TITLE = "Microsoft YaHei UI"
FONT_BODY = "Microsoft YaHei UI"
FONT_INK = "KaiTi"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_with_fill(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_name, font_size, font_color, bold=False, alignment=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox

def add_multi_line_textbox(slide, left, top, width, height, lines, font_name, font_size, font_color, bold=False, alignment=PP_ALIGN.CENTER, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.5)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return txBox

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 0, 7.35, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 1.5, 1.5, 10.33, 4.5, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 1.5, 1.8, 10.33, 1.2, "道德经", FONT_TITLE, 60, COLOR_BLUE_ACCENT, bold=True)
add_textbox(slide, 1.5, 3.0, 10.33, 0.8, "第二十五天 · 第七十九章至第八十一章", FONT_BODY, 28, COLOR_WHITE, bold=False)
add_shape_with_fill(slide, 4.0, 3.9, 5.33, 0.03, COLOR_BLUE_ACCENT)
add_textbox(slide, 1.5, 4.1, 10.33, 0.6, "报怨以德 · 小国寡民 · 为而不争", FONT_INK, 22, COLOR_BLUE_ACCENT, bold=False)
add_textbox(slide, 1.5, 5.5, 10.33, 0.5, "三先生国学课件", FONT_BODY, 16, COLOR_LIGHT_TEXT, bold=False)

# ============================================================
# 第2页：目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 0.3, 12.33, 0.8, "今日课程", FONT_TITLE, 36, COLOR_BLUE_PRIMARY, bold=True)
add_shape_with_fill(slide, 0.5, 1.1, 12.33, 0.04, COLOR_BLUE_ACCENT)

chapters = [
    ("第七十九章", "报怨以德", "和大怨，必有余怨\n安可以为善？", "天道无亲，常与善人"),
    ("第八十章", "小国寡民", "小国寡民，使有什伯之器\n而不用", "邻国相望，鸡犬之声\n相闻"),
    ("第八十一章", "为而不争", "信言不美，美言不信\n善者不辩，辩者不善", "天之道，利而不害\n圣人之道，为而不争"),
]

card_width = 3.8
card_start_x = 0.5
card_gap = 0.35

for i, (num, title, excerpt, highlight) in enumerate(chapters):
    x = card_start_x + i * (card_width + card_gap)
    add_shape_with_fill(slide, x, 1.5, card_width, 5.5, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
    add_textbox(slide, x, 1.7, card_width, 0.5, num, FONT_TITLE, 20, COLOR_BLUE_PRIMARY, bold=True)
    add_textbox(slide, x, 2.2, card_width, 0.7, title, FONT_TITLE, 32, COLOR_INK, bold=True)
    add_shape_with_fill(slide, x + 0.3, 3.0, card_width - 0.6, 0.03, COLOR_BLUE_ACCENT)
    add_textbox(slide, x + 0.2, 3.2, card_width - 0.4, 1.5, excerpt, FONT_INK, 16, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 4.8, card_width - 0.4, 1.8, highlight, FONT_INK, 14, COLOR_BLUE_LIGHT, bold=False, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第3页：第七十九章 原文与概述
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 1.0, COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第七十九章 · 报怨以德", FONT_TITLE, 36, COLOR_WHITE, bold=True)

add_shape_with_fill(slide, 0.4, 1.3, 6.0, 5.8, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 1.4, 5.8, 0.5, "【原文】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

original_text = """和大怨，必有余怨，安可以为善？
是以圣人执左契，而不责于人。
有德司契，无德司彻。
天道无亲，常与善人。"""

add_textbox(slide, 0.6, 1.9, 5.6, 5.0, original_text, FONT_INK, 22, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)

add_shape_with_fill(slide, 6.7, 1.3, 6.2, 5.8, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 6.8, 1.4, 6.0, 0.5, "【注释】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

notes = [
    ("大怨", "深重的怨恨、仇怨"),
    ("余怨", "残余的怨恨"),
    ("执左契", "持有借据的存根。左契：债权人持有的凭证"),
    ("责于人", "向人索取、追究"),
    ("司契", "掌管契约的人，喻指有德者"),
    ("司彻", "负责收取田租的人，喻指斤斤计较"),
    ("无亲", "没有偏爱、不偏不倚"),
    ("与善人", "帮助、赐福于善人"),
]

y = 2.0
for term, meaning in notes:
    add_textbox(slide, 6.9, y, 1.2, 0.4, term, FONT_TITLE, 14, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 8.1, y, 4.6, 0.4, meaning, FONT_BODY, 14, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.52

# ============================================================
# 第4页：第七十九章 核心义理与启示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_ACCENT)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第七十九章 · 核心义理与启示", FONT_TITLE, 32, COLOR_BLUE_ACCENT, bold=True)

add_shape_with_fill(slide, 0.4, 1.2, 6.0, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 0.5, 1.3, 5.8, 0.5, "【核心义理】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 0.6, 1.8, 5.4, 0.03, COLOR_BLUE_ACCENT)

core1 = [
    "和解怨恨的智慧",
    "老子认为，怨恨即使被调和，也必有残余，无法真正达到善的境界。真正的智慧是不要结下怨恨。",
    "",
    "圣人不责于人",
    "圣人持有借据却不向人索取，展现的是一种宽容、仁慈的处世态度。",
    "",
    "有德与无德的对比",
    "有德者如掌管契约者，宽厚待人；无德者如掌管税收者，斤斤计较。天道公正，眷顾善人。",
]

add_multi_line_textbox(slide, 0.6, 1.95, 5.6, 5.0, core1, FONT_BODY, 15, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_shape_with_fill(slide, 6.7, 1.2, 6.2, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 6.8, 1.3, 6.0, 0.5, "【人生启示】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 6.9, 1.8, 5.8, 0.03, COLOR_BLUE_ACCENT)

insights1 = [
    "防怨胜于解怨",
    "与其费尽心思化解矛盾，不如从一开始就不要与人结怨。",
    "",
    "宽容是一种智慧",
    "当别人对不起我们时，选择宽容而非报复，这是一种内在的力量。",
    "",
    "积德行善",
    "天道无亲，常与善人。善行会积累福报，最终回馈自身。",
    "",
    "不斤斤计较",
    "有德之人待人宽厚，不在小事上与人争执计较。",
]

add_multi_line_textbox(slide, 6.9, 1.95, 5.8, 5.0, insights1, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

# ============================================================
# 第5页：第八十章 原文
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 1.0, COLOR_BLUE_LIGHT)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十章 · 小国寡民", FONT_TITLE, 36, COLOR_WHITE, bold=True)

add_shape_with_fill(slide, 0.4, 1.3, 12.5, 5.8, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 1.4, 12.3, 0.5, "【原文】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

original_80 = """小国寡民，使有什伯之器而不用，使民重死而不远徙。
虽有舟舆，无所乘之；虽有甲兵，无所陈之。
使民复结绳而用之。
甘其食，美其服，安其居，乐其俗。
邻国相望，鸡犬之声相闻，民至老死不相往来。"""

add_textbox(slide, 0.8, 1.95, 11.8, 4.8, original_80, FONT_INK, 22, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第6页：第八十章 注释
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 0.8, COLOR_BLUE_LIGHT)
add_textbox(slide, 0.5, 0.15, 12.33, 0.7, "第八十章 · 注释", FONT_TITLE, 32, COLOR_WHITE, bold=True)

notes_80_left = [
    ("小国寡民", "国家小，人民少"),
    ("什伯之器", "各种各样的器具"),
    ("不用", "不需要使用"),
    ("重死", "珍视生命"),
    ("不远徙", "不向远方迁移"),
    ("舟舆", "船只和车辆"),
    ("甲兵", "铠甲和兵器"),
    ("无所陈之", "没有地方陈列使用"),
]

notes_80_right = [
    ("结绳而用", "用绳子打结记事"),
    ("甘其食", "以自己的食物为甘美"),
    ("美其服", "以自己的服饰为美好"),
    ("安其居", "安于自己的居所"),
    ("乐其俗", "乐于自己的风俗"),
    ("邻国相望", "相邻的国家可以互相望见"),
    ("鸡犬之声", "鸡鸣狗吠的声音"),
    ("不相往来", "彼此不互相往来交流"),
]

add_shape_with_fill(slide, 0.3, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.4, 1.2, 6.1, 0.4, "【注释（上）】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_80_left:
    add_textbox(slide, 0.5, y, 1.5, 0.45, term, FONT_TITLE, 14, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 2.0, y, 4.4, 0.45, meaning, FONT_BODY, 14, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.62

add_shape_with_fill(slide, 6.8, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 6.9, 1.2, 6.1, 0.4, "【注释（下）】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_80_right:
    add_textbox(slide, 7.0, y, 1.5, 0.45, term, FONT_TITLE, 14, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 8.5, y, 4.4, 0.45, meaning, FONT_BODY, 14, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.62

# ============================================================
# 第7页：第八十章 核心义理与启示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_LIGHT)

add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十章 · 核心义理与启示", FONT_TITLE, 32, COLOR_BLUE_ACCENT, bold=True)

add_shape_with_fill(slide, 0.4, 1.2, 6.0, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_LIGHT)
add_textbox(slide, 0.5, 1.3, 5.8, 0.5, "【核心义理】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 0.6, 1.8, 5.4, 0.03, COLOR_BLUE_LIGHT)

core80 = [
    "小国寡民的理想",
    "不是反对发展，而是追求一种适度、简朴的社会状态。国小民少，资源充足，纷争减少。",
    "",
    "返璞归真",
    "使民复结绳而用之，不是倒退，而是在物质极大丰富后，选择回归本真，简约生活。",
    "",
    "知足常乐",
    "甘其食，美其服，安其居，乐其俗。不追求奢华，但求满足。内心的富足才是真正的富足。",
    "",
    "邻里的和谐",
    "鸡犬之声相闻，却不相往来——不是冷漠，而是不互相干扰，各自安好。",
]

add_multi_line_textbox(slide, 0.6, 1.95, 5.6, 5.0, core80, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_shape_with_fill(slide, 6.7, 1.2, 6.2, 5.9, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_LIGHT)
add_textbox(slide, 6.8, 1.3, 6.0, 0.5, "【人生启示】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 6.9, 1.8, 5.8, 0.03, COLOR_BLUE_LIGHT)

insights80 = [
    "知足者富",
    "不是拥有得多才富，而是知道满足、不贪求的人，才是真正的富足。",
    "",
    "简化生活",
    "减少不必要的物质追求和欲望，把精力放在真正重要的事情上。",
    "",
    "内心安宁",
    "不攀比、不焦虑，安于当下，享受已有的生活。",
    "",
    "和谐共处",
    "与邻为善，但不互相干涉。保持适度的距离，反而更能长久。",
]

add_multi_line_textbox(slide, 6.9, 1.95, 5.8, 5.0, insights80, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

# ============================================================
# 第8页：第八十一章 原文
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 1.0, RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十一章 · 为而不争", FONT_TITLE, 36, COLOR_WHITE, bold=True)

add_shape_with_fill(slide, 0.4, 1.3, 12.5, 5.8, RGBColor(0xff, 0xff, 0xff), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 1.4, 12.3, 0.5, "【原文】", FONT_TITLE, 18, COLOR_BLUE_PRIMARY, bold=True)

original_81 = """信言不美，美言不信。
善者不辩，辩者不善。
知者不博，博者不知。
圣人不积，既以为人己愈有，既以与人己愈多。
天之道，利而不害；
圣人之道，为而不争。"""

add_textbox(slide, 0.8, 1.95, 11.8, 4.8, original_81, FONT_INK, 24, COLOR_INK, bold=False, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第9页：第八十一章 注释
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_shape_with_fill(slide, 0, 0.1, 13.33, 0.8, RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 0.5, 0.15, 12.33, 0.7, "第八十一章 · 注释", FONT_TITLE, 32, COLOR_WHITE, bold=True)

notes_81 = [
    ("信言不美", "真实的话不华美"),
    ("美言不信", "华美的话不真实"),
    ("善者不辩", "真正善良的人不需要争辩"),
    ("辩者不善", "善于争辩的人不一定善良"),
    ("知者不博", "真正有智慧的人不一定博学"),
    ("博者不知", "博学的人不一定有真智慧"),
    ("圣人不积", "圣人不为自己积蓄"),
    ("为人", "帮助别人、给予别人"),
    ("愈有", "更加富有"),
    ("与人", "给予别人"),
    ("愈多", "更加丰富"),
    ("利而不害", "利益万物而不伤害万物"),
    ("为而不争", "有所作为但不与人争"),
]

add_shape_with_fill(slide, 0.3, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.4, 1.2, 6.1, 0.4, "【注释】", FONT_TITLE, 16, COLOR_BLUE_PRIMARY, bold=True)

y = 1.7
for term, meaning in notes_81[:7]:
    add_textbox(slide, 0.5, y, 1.6, 0.5, term, FONT_TITLE, 14, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 2.1, y, 4.3, 0.5, meaning, FONT_BODY, 14, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.63

add_shape_with_fill(slide, 6.8, 1.1, 6.3, 6.0, RGBColor(0xe8, 0xf0, 0xf8), COLOR_BLUE_PRIMARY)

y = 1.7
for term, meaning in notes_81[7:]:
    add_textbox(slide, 7.0, y, 1.6, 0.5, term, FONT_TITLE, 14, COLOR_BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 8.6, y, 4.3, 0.5, meaning, FONT_BODY, 14, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
    y += 0.63

# ============================================================
# 第10页：第八十一章 核心义理与启示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, RGBColor(0x1a, 0x3a, 0x6e))

add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "第八十一章 · 核心义理与启示", FONT_TITLE, 32, COLOR_BLUE_ACCENT, bold=True)

add_shape_with_fill(slide, 0.4, 1.2, 6.0, 5.9, RGBColor(0x0f, 0x1e, 0x3c), RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 0.5, 1.3, 5.8, 0.5, "【核心义理】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 0.6, 1.8, 5.4, 0.03, RGBColor(0x1a, 0x3a, 0x6e))

core81 = [
    "真假与美丑",
    "真实的话往往不漂亮，漂亮的话往往不真实。这是老子对表象与本质的深刻洞察。",
    "",
    "善与辩",
    "真正的善不需要自我辩解，真正有智慧的人不追求博学。内在的品德胜于外在的表现。",
    "",
    "给予即是获得",
    "越给予别人，自己越富有；越帮助别人，自己越充足。这是宇宙的法则。",
    "",
    "为而不争",
    "天之道是利益万物而不伤害，圣人之道是有所作为但不争竞。这是最妙的道。",
]

add_multi_line_textbox(slide, 0.6, 1.95, 5.6, 5.0, core81, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_shape_with_fill(slide, 6.7, 1.2, 6.2, 5.9, RGBColor(0x0f, 0x1e, 0x3c), RGBColor(0x1a, 0x3a, 0x6e))
add_textbox(slide, 6.8, 1.3, 6.0, 0.5, "【人生启示】", FONT_TITLE, 20, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 6.9, 1.8, 5.8, 0.03, RGBColor(0x1a, 0x3a, 0x6e))

insights81 = [
    "真诚为本",
    "说话做事以真诚为根本，不要为了好听而说假话。真诚的人最有力量。",
    "",
    "少说多做",
    "少一些无谓的争辩，多一些实际的行动。用结果说话，比用嘴巴说服更有力。",
    "",
    "助人即助己",
    "帮助他人不是失去，而是得到。给予本身就是最大的回报。",
    "",
    "不争的智慧",
    "不争，不是消极不作为，而是不计较得失、不与人争锋。夫唯不争，天下莫能与之争。",
]

add_multi_line_textbox(slide, 6.9, 1.95, 5.8, 5.0, insights81, FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

# ============================================================
# 第11页：三章总结对比
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_LIGHT)
add_shape_with_fill(slide, 0, 0, 13.33, 0.1, COLOR_BLUE_PRIMARY)
add_textbox(slide, 0.5, 0.2, 12.33, 0.8, "三章总结 · 智慧贯穿", FONT_TITLE, 36, COLOR_BLUE_PRIMARY, bold=True)
add_shape_with_fill(slide, 0.5, 1.0, 12.33, 0.04, COLOR_BLUE_ACCENT)

themes = [
    {
        "title": "报怨以德",
        "color": COLOR_BLUE_PRIMARY,
        "subtitle": "第七十九章",
        "points": ["和解怨恨的智慧", "宽容不责于人", "天道常与善人", "有德司契，无德司彻"],
        "core": "宽恕之道"
    },
    {
        "title": "小国寡民",
        "color": COLOR_BLUE_LIGHT,
        "subtitle": "第八十章",
        "points": ["知足常乐", "返璞归真", "甘其食，美其服", "安居乐俗，不相往来"],
        "core": "知足之道"
    },
    {
        "title": "为而不争",
        "color": RGBColor(0x1a, 0x3a, 0x6e),
        "subtitle": "第八十一章",
        "points": ["信言不美，美言不信", "善者不辩，辩者不善", "为人愈有，与人愈多", "天之道，利而不害"],
        "core": "无为之境"
    },
]

card_w = 4.0
gap = 0.22
start_x = 0.4

for i, theme in enumerate(themes):
    x = start_x + i * (card_w + gap)
    add_shape_with_fill(slide, x, 1.2, card_w, 6.0, RGBColor(0xff, 0xff, 0xff), theme["color"])
    add_textbox(slide, x, 1.35, card_w, 0.5, theme["subtitle"], FONT_BODY, 14, theme["color"], bold=False)
    add_textbox(slide, x, 1.7, card_w, 0.7, theme["title"], FONT_TITLE, 26, theme["color"], bold=True)
    add_shape_with_fill(slide, x + 0.3, 2.45, card_w - 0.6, 0.03, theme["color"])
    add_textbox(slide, x, 2.55, card_w, 0.5, theme["core"], FONT_INK, 18, theme["color"], bold=True)
    y = 3.15
    for point in theme["points"]:
        add_textbox(slide, x + 0.25, y, card_w - 0.5, 0.5, "• " + point, FONT_BODY, 13, COLOR_INK, bold=False, alignment=PP_ALIGN.LEFT)
        y += 0.62

# ============================================================
# 第12页：结语
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_BG_DARK)
add_shape_with_fill(slide, 0, 0, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 0, 7.35, 13.33, 0.15, COLOR_BLUE_ACCENT)
add_shape_with_fill(slide, 1.5, 1.2, 10.33, 5.1, RGBColor(0x0f, 0x2a, 0x4f), COLOR_BLUE_ACCENT)
add_textbox(slide, 1.5, 1.5, 10.33, 0.8, "今日要点", FONT_TITLE, 36, COLOR_BLUE_ACCENT, bold=True)
add_shape_with_fill(slide, 4.5, 2.35, 4.33, 0.03, COLOR_BLUE_ACCENT)

summary = [
    "第七十九章：宽容之道——报怨以德，不责于人，天道无亲，常与善人。",
    "第八十章：知足之道——甘其食，美其服，安其居，乐其俗，不相往来。",
    "第八十一章：无为之境——为而不争，利而不害，既以与人己愈多。",
]

add_multi_line_textbox(slide, 1.8, 2.55, 9.7, 2.8, summary, FONT_BODY, 17, COLOR_LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT)

add_textbox(slide, 1.5, 5.4, 10.33, 0.6, "天之道，利而不害；圣人之道，为而不争。", FONT_INK, 20, COLOR_BLUE_ACCENT, bold=False)
add_textbox(slide, 1.5, 6.1, 10.33, 0.4, "—— 道德经 · 第七十九至八十一章 · 完 ——", FONT_BODY, 14, COLOR_LIGHT_TEXT, bold=False)

output_path = "/Users/mac/.openclaw/workspace/courses/道德经_第25天.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
