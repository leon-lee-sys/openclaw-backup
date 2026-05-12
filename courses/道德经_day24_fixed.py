# -*- coding: utf-8 -*-
"""
道德经第24天 课件生成器
章节：第76-78章
风格：商务蓝/深色系，正文字体比标准模板大2号
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# 配色方案 - 商务蓝/深色系
PRIMARY = RGBColor(30, 60, 114)       # 深蓝
SECONDARY = RGBColor(70, 130, 180)    # 中蓝
ACCENT = RGBColor(204, 120, 50)       # 橙色强调
TEXT = RGBColor(51, 51, 51)           # 深灰文字
BG = RGBColor(245, 247, 250)          # 浅灰蓝背景
WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(20, 35, 60)        # 深蓝背景

# 字体加大2号的配置（原模板约16-18pt，加大2号后约18-20pt）
TITLE_SIZE = Pt(38)
HEADER_SIZE = Pt(30)
YUANWEN_SIZE = Pt(22)    # 原文加大2号
BODY_SIZE = Pt(18)       # 正文加大2号
SMALL_SIZE = Pt(16)      # 辅助说明加大2号

def set_dark_title_bar(slide, title):
    """深色标题栏"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = HEADER_SIZE
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def add_yuanwen_section(slide, lines, y=1.5, height=3.2):
    """原文区域"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = SECONDARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line if line else "　"
        pp.font.size = YUANWEN_SIZE
        pp.font.color.rgb = WHITE
        pp.space_before = Pt(3)

def add_commentary_section(slide, title, content, y):
    """逐句释义区域"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    box.line.color.rgb = SECONDARY
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"【{title}】"
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = BODY_SIZE
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=SECONDARY):
    """核心观点框"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(15)
    p2.font.color.rgb = WHITE

def add_points_list(slide, points, y_start=3.2):
    """要点列表"""
    y = y_start
    for icon_title, explanation in points:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.25))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = icon_title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = explanation
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(4)
        y += 1.35
    return y

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.6))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.9))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第24课"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = WHITE
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.6))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第76-78章"
p2.font.size = Pt(58)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.6), Inches(8), Inches(0.9))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "出生入死 · 柔弱处上 · 为而不争"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = DARK_BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "天下莫柔弱于水，而攻坚强者莫之能胜。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = WHITE
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第七十八章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = RGBColor(180, 180, 180)
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, 14)

# ===== 第2页：三章概览 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s2, "今日课程")
chapters = [
    ("第七十六章", "出生入死", "民之生也柔弱，死也坚强。万物草木之生也柔脆，死也枯槁。故坚强者死之徒，柔弱者生之徒。"),
    ("第七十七章", "柔弱处上", "天下之至柔，驰骋天下之至坚。无有入无间。柔弱胜刚强。"),
    ("第七十八章", "为而不争", "天下莫柔弱于水，而攻坚强者莫之能胜。以柔克刚，知其雄，守其雌。")
]
y2 = 1.5
for num, title, desc in chapters:
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y2), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf2c = box.text_frame
    tf2c.word_wrap = True
    p2c1 = tf2c.paragraphs[0]
    p2c1.text = f"📖 {num}：{title}"
    p2c1.font.size = Pt(20)
    p2c1.font.bold = True
    p2c1.font.color.rgb = PRIMARY
    p2c2 = tf2c.add_paragraph()
    p2c2.text = desc
    p2c2.font.size = Pt(16)
    p2c2.font.color.rgb = TEXT
    p2c2.space_before = Pt(5)
    y2 += 1.65
add_slide_number(s2, 2, 14)

# ===== 第3页：第七十六章 原文 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s3, "第七十六章 · 出生入死")
add_yuanwen_section(s3, [
    "人之生也柔弱，其死也坚强。",
    "万物草木之生也柔脆，其死也枯槁。",
    "",
    "故坚强者死之徒，柔弱者生之徒。",
    "",
    "兵强则不胜，木强则兵。",
    "",
    "强大处下，柔弱处上。"
], y=1.5, height=3.8)
add_commentary_section(s3, "白话解读",
    "人活着时身体柔软，死了则变得僵硬。万物草木活着时枝条柔脆，死了则干枯。所以坚强的东西属于死亡一类，柔弱的东西属于生存一类。军队过于强大反而不能取胜，树木过于强硬则会被砍伐。强大的处于下位，柔弱的反而在上。",
    5.4)
add_slide_number(s3, 3, 14)

# ===== 第4页：第七十六章 逐句释义 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s4, "第七十六章 · 逐句释义")
add_points_list(s4, [
    ("① 人之生也柔弱，其死也坚强",
     "人活着时身体柔软灵活，死了则僵硬——柔弱是生命力的象征，坚强是死亡的预兆。"),
    ("② 万物草木之生也柔脆，其死也枯槁",
     "草木活着时枝条柔韧，死了则干枯——自然万物无一例外，柔弱代表生机，枯槁代表死亡。"),
    ("③ 坚强者死之徒，柔弱者生之徒",
     "属于坚强一类的是死亡之旅，属于柔弱一类的是生存之旅——老子揭示了强弱转化的辩证法则。"),
    ("④ 强大处下，柔弱处上",
     "强大的反而处于下方，柔弱的反而处于上方——柔弱不是软弱，而是更具生命力的状态。")
], y_start=1.5)
add_slide_number(s4, 4, 14)

# ===== 第5页：第七十六章 本章小结 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s5, "第七十六章 · 本章小结")
add_core_box(s5, "🌟 核心观点", "柔弱胜刚强——生命的力量不在于强硬，而在于柔韧。真正强大的，是那些懂得示弱、保持柔顺的人。")
add_points_list(s5, [
    ("1. 柔弱代表生机", "人活着时柔弱，死了才僵硬——柔弱是生命力的体现，是活力的象征。"),
    ("2. 坚强属于死亡", "坚强的事物最终都会走向消亡——过于强硬反而会招致失败。"),
    ("3. 兵强则不胜", "军队过于强硬反而无法获胜——逞强不是力量，柔韧才是真正的强大。"),
    ("4. 强大处下", "真正强大的人懂得处下——处下不是软弱，而是智慧。")
], y_start=3.1)
add_slide_number(s5, 5, 14)

# ===== 第6页：第七十七章 原文 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s6, "第七十七章 · 柔弱处上")
add_yuanwen_section(s6, [
    "天下之至柔，驰骋天下之至坚。",
    "无有入无间。",
    "",
    "吾是以知无为之有益。",
    "",
    "天下莫柔弱于水，",
    "而攻坚强者莫之能胜。",
    "",
    "其无以易之。",
    "",
    "柔弱胜刚强。",
    "",
    "弱之胜强，柔之胜刚，",
    "天下莫不知，莫能行。"
], y=1.5, height=4.5)
add_commentary_section(s6, "白话解读",
    "天下最柔弱的东西，能够穿透最坚硬的东西。无形的存在能够进入没有间隙的地方。我因此知道无为是有益的。天下没有什么比水更柔弱，但攻坚克强却没有能胜过水的。因为没有什么可以替代它。柔弱能战胜刚强。弱能胜强，柔能胜刚，天下没有人不知道，但没有人能够实行。",
    6.1)
add_slide_number(s6, 6, 14)

# ===== 第7页：第七十七章 逐句释义 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s7, "第七十七章 · 逐句释义")
add_points_list(s7, [
    ("① 天下之至柔，驰骋天下之至坚",
     "天下最柔弱的东西，能驾驭最坚硬的东西——水能穿石，风能蚀木，柔弱之力超越强硬。"),
    ("② 无有入无间",
     "无形的力量能够穿透没有间隙的东西——"无"比"有"更有力量，无为比妄为更有效果。"),
    ("③ 天下莫柔弱于水，而攻坚强者莫之能胜",
     "天下没有什么比水更柔弱，但攻坚强者没有能胜过水的——柔能克刚，以柔克刚是自然法则。"),
    ("④ 天下莫不知，莫能行",
     "天下没有人不知道柔胜刚的道理，但没有人能够实行——知易行难，需要真正的智慧与修养。")
], y_start=1.5)
add_slide_number(s7, 7, 14)

# ===== 第8页：第七十七章 本章小结 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s8, "第七十七章 · 本章小结")
add_core_box(s8, "🌟 核心观点", "柔弱胜刚强——天下最柔的水，能攻克天下最刚强的东西。以柔克刚，是道的智慧。")
add_points_list(s8, [
    ("1. 至柔驰骋至坚", "最柔弱的东西能驾驭最坚硬的东西——柔弱不是无能，而是最具穿透力的力量。"),
    ("2. 无有入无间", "无形的力量能穿透无隙的东西——无为胜有为，无形胜有形。"),
    ("3. 水之柔弱", "水虽柔弱却能攻克坚强——柔弱胜刚强是最普遍的自然法则。"),
    ("4. 知而不行", "天下人莫不知柔胜刚，却莫能行——知道容易做到难，需要长期修养。")
], y_start=3.1)
add_slide_number(s8, 8, 14)

# ===== 第9页：第七十八章 原文 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s9, "第七十八章 · 为而不争")
add_yuanwen_section(s9, [
    "天下莫柔弱于水，",
    "而攻坚强者莫之能胜。",
    "",
    "以柔克刚，以弱胜强。",
    "",
    "知其雄，守其雌，",
    "此为溪常德。",
    "",
    "复归于婴儿，",
    "复归于朴。",
    "",
    "为而不争，",
    "天之道也。"
], y=1.5, height=4.0)
add_commentary_section(s9, "白话解读",
    "天下没有什么比水更柔弱，但攻坚克强却没有能胜过水的。以柔克刚，以弱胜强。知道什么是雄强，却甘愿处于雌柔，这就是接近永恒的德。再回归到婴儿般纯朴的状态，再回归到道的本源。有所作为但不与人争夺，这是自然的法则。",
    5.6)
add_slide_number(s9, 9, 14)

# ===== 第10页：第七十八章 逐句释义 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s10, "第七十八章 · 逐句释义")
add_points_list(s10, [
    ("① 天下莫柔弱于水，而攻坚强者莫之能胜",
     "天下最柔弱的是水，但攻坚克强没有能胜过水的——水是柔弱胜刚强的最好例证。"),
    ("② 知其雄，守其雌",
     "知道什么是雄强，却甘愿处于雌柔——明知自己有能力，却选择谦下，这是最高的德行。"),
    ("③ 复归于婴儿，复归于朴",
     "回归到婴儿般的纯真，回归到道的本源——返璞归真，复归自然，是人生的最高境界。"),
    ("④ 为而不争，天之道也",
     "有所作为但不与人争夺，这是自然的法则——做该做的事，但不争功、不争利。")
], y_start=1.5)
add_slide_number(s10, 10, 14)

# ===== 第11页：第七十八章 本章小结 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s11, "第七十八章 · 本章小结")
add_core_box(s11, "🌟 核心观点", "为而不争——做该做的事，但不强求结果，不与人争。这是天之道，也是最高明的处世之道。")
add_points_list(s11, [
    ("1. 柔弱胜刚强", "水虽柔弱却无坚不摧——柔弱不是无能，而是最强大的力量。"),
    ("2. 知雄守雌", "知道自己的优势，却甘愿谦下——这是一种深层的智慧和修养。"),
    ("3. 复归于朴", "回归本真、返璞归真——不被世俗所染，保持内心的纯朴。"),
    ("4. 为而不争", "做事情但不争抢——天之道是利万物而不争，人的最高境界也是为而不争。")
], y_start=3.1)
add_slide_number(s11, 11, 14)

# ===== 第12页：三章总览 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s12, "三章精华 · 一以贯之")
sums = [
    ("第七十六章", "出生入死", "柔弱生之徒，坚强死之徒。强大处下，柔弱处上。生命的力量在于柔韧而非强硬。"),
    ("第七十七章", "柔弱处上", "天下之至柔，驰骋天下之至坚。无有入无间。柔弱胜刚强，道之奥义。"),
    ("第七十八章", "为而不争", "知其雄，守其雌，复归于朴。为而不争，天之道也——最高明的处世之道。")
]
y12 = 1.5
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.55))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    p12_2.space_before = Pt(5)
    y12 += 1.65
add_slide_number(s12, 12, 14)

# ===== 第13页：现代应用 =====
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s13, "📌 现代应用")
apps = [
    ("🏋️ 健身与健康", "身体保持柔韧则健康，僵硬则衰亡——经常拉伸、瑜伽、太极，保持身体的柔韧性。"),
    ("💼 商业竞争", "柔弱胜刚强 → 竞争中不逞强，善用柔韧策略，以退为进，才能在商海中立于不败之地。"),
    ("🧘 内心修养", "复归于婴儿 → 减少成人世界的算计与焦虑，保持内心的纯真与好奇，活得更轻松自在。"),
    ("🤝 人际关系", "为而不争 → 与人相处多做少说，不争功不抢功，反而更容易获得尊重与信任。"),
    ("💰 投资理财", "弱之胜强 → 稳健投资不冒险，不追求高风险高回报，长期坚持反而收获更丰。")
]
add_points_list(s13, apps, y_start=1.5)
add_slide_number(s13, 13, 14)

# ===== 第14页：结语 =====
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_title_bar(s14, "第24课 · 结语")
final_box = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(2))
final_box.fill.solid()
final_box.fill.fore_color.rgb = DARK_BG
tff = final_box.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "柔弱胜刚强"
pf1.font.size = Pt(36)
pf1.font.bold = True
pf1.font.color.rgb = WHITE
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "天下莫柔弱于水，而攻坚强者莫之能胜"
pf2.font.size = Pt(18)
pf2.font.italic = True
pf2.font.color.rgb = SECONDARY
pf2.alignment = PP_ALIGN.CENTER
pf2.space_before = Pt(10)
pf3 = tff.add_paragraph()
pf3.text = "为而不争，天之道也"
pf3.font.size = Pt(16)
pf3.font.color.rgb = RGBColor(180, 180, 180)
pf3.alignment = PP_ALIGN.CENTER
footer_note = s14.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
tffn = footer_note.text_frame
tffn.word_wrap = True
pfn1 = tffn.paragraphs[0]
pfn1.text = "道家智慧：至柔驰骋至坚，无有入无间。柔弱不是软弱，而是最具生命力的力量。"
pfn1.font.size = Pt(16)
pfn1.font.color.rgb = TEXT
pfn1.alignment = PP_ALIGN.CENTER
pfn2 = tffn.add_paragraph()
pfn2.text = "知雄守雌，复归于朴——这是老子留给我们的处世良方。"
pfn2.font.size = Pt(15)
pfn2.font.color.rgb = SECONDARY
pfn2.alignment = PP_ALIGN.CENTER
pfn2.space_before = Pt(8)
add_slide_number(s14, 14, 14)

# 保存文件
output_path = '/Users/mac/.openclaw/workspace/courses/道德经_第24天.pptx'
prs.save(output_path)
print(f"Done: {output_path}")