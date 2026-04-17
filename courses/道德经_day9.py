from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

COLORS = {
    'primary': RGBColor(0, 51, 102),
    'secondary': RGBColor(102, 102, 102),
    'accent': RGBColor(204, 0, 0),
    'text': RGBColor(51, 51, 51),
    'bg': RGBColor(255, 255, 255),
    'light': RGBColor(240, 248, 255),
}

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['bg']

def title_bar(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(30); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def quote_box(slide, text, top=1.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(top), Inches(11.73), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['light']
    box.line.color.rgb = COLORS['primary']
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(20); p.font.name = 'KaiTi'
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    tf.margin_left = Pt(20); tf.margin_right = Pt(20)

def section_title(slide, text, top=1.4):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.73), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(22); p.font.bold = True
    p.font.color.rgb = COLORS['accent']

def body_text(slide, lines, top=2.1):
    left, top_, width, height = Inches(0.8), Inches(top), Inches(11.73), Inches(4.5)
    tb = slide.shapes.add_textbox(left, top_, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(17)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(10)

def insight_box(slide, title, points, top=5.3):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(top), Inches(11.73), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['primary']
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"💡 {title}"
    p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 220, 80)
    for pt in points:
        p2 = tf.add_paragraph()
        p2.text = f"  • {pt}"; p2.font.size = Pt(15)
        p2.font.color.rgb = RGBColor(220, 240, 255)

# ═══════════════════════════════════════
# Slide 1 — 封面
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(2.2))
bar.fill.solid(); bar.fill.fore_color.rgb = COLORS['primary']
bar.line.fill.background()
tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1.2))
p = tb.text_frame.paragraphs[0]
p.text = '《道德经》每日课件'; p.font.size = Pt(44); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER
tb2 = s.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.33), Inches(0.8))
p2 = tb2.text_frame.paragraphs[0]
p2.text = '第9天 · 第二十五章 至 第二十七章'; p2.font.size = Pt(26)
p2.font.color.rgb = RGBColor(200, 230, 255); p2.alignment = PP_ALIGN.CENTER
tb3 = s.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.33), Inches(0.6))
p3 = tb3.text_frame.paragraphs[0]
p3.text = '2026年4月16日'; p3.font.size = Pt(20)
p3.font.color.rgb = COLORS['secondary']; p3.alignment = PP_ALIGN.CENTER
tb4 = s.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.33), Inches(1.5))
p4 = tb4.text_frame.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
p4.text = '"有物混成，先天地生……独立不改，周行而不殆"'; p4.font.size = Pt(22)
p4.font.name = 'KaiTi'; p4.font.color.rgb = COLORS['primary']
p4b = tb4.text_frame.add_paragraph(); p4b.alignment = PP_ALIGN.CENTER
p4b.text = '—— 道法自然，万物之宗'; p4b.font.size = Pt(16)
p4b.font.color.rgb = COLORS['secondary']

# ═══════════════════════════════════════
# Slide 2 — 第二十五章
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十五章｜有物混成')
quote_box(s, '有物混成，先天地生。寂兮寥兮，独立不改，\n周行而不殆，可以为天下母。\n吾不知其名，字之曰道，强为之名曰大。')
section_title(s, '白话解读')
body_text(s, [
    '→ 有一个浑然一体的东西，在天地之前就已经存在。',
    '→ 它寂静空虚，不依靠任何外力而独立存在，循环运行从不停止，',
    '→ 可以作为天下的根本。我不知道它的名字，勉强把它叫做"道"，',
    '→ 再勉强给它起个名字叫"大"。',
    '→ 道大、天大、地大、人也大。宇宙中有四大，人居其中之一。',
    '→ 人效法地，地效法天，天效法道，道效法自然。'
], top=3.1)
insight_box(s, '核心智慧：道法自然，人应顺应四大法则', [
    '道是宇宙本源，人应效法天地，遵循自然规律而不妄为',
])

# ═══════════════════════════════════════
# Slide 3 — 第二十五章 典故
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十五章｜典故案例')
section_title(s, '历史典故：汉初"文景之治"（道法自然）')
body_text(s, [
    '【背景】秦末战乱后，汉初民生凋敝，社会急需休养生息。',
    '',
    '【策略】汉高祖刘邦入关后，约法三章，轻徭薄赋。',
    '至汉文帝、景帝，继续推行"黄老之术"——以道家思想治国，',
    '政府尽量少干预，让百姓自然恢复。',
    '',
    '【结果】农业税收降至三十税一，甚至一度全免；',
    '废除肉刑，废止诽谤罪；边境和平，对匈奴和亲。',
    '四十年间，百姓富足，粮仓堆积如山，史称"文景之治"。',
    '',
    '【感悟】"人法地，地法天，天法道，道法自然"——',
    '治国如同种田，不能揠苗助长。汉初帝王深谙此道，',
    '以"无为"之政，收"无不为"之效，正是老子"道法自然"的绝佳诠释。',
], top=1.4)

# ═══════════════════════════════════════
# Slide 4 — 第二十六章
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十六章｜重为轻根')
quote_box(s, '重为轻根，静为躁君。\n是以圣人终日行，不离辎重。\n虽有荣观，燕处超然。奈何万乘之主，而以身轻天下？\n轻则失本，躁则失君。')
section_title(s, '白话解读')
body_text(s, [
    '→ 稳重是轻率的根本，宁静是躁动的主宰。',
    '→ 因此圣人每天行走，始终不离开载重车辆。',
    '→ 虽然有奢华的观赏享乐，却能安然超脱，不被物欲所困。',
    '→ 为什么身为拥有万辆战车的大国君主，却以轻率的态度治理天下呢？',
    '→ 轻率就会失去根本，躁动就会失去主宰的地位。'
], top=3.1)
insight_box(s, '核心智慧：稳重宁静，守住根本', [
    '轻浮躁动是领导大忌，稳重沉静才能驾驭全局',
])

# ═══════════════════════════════════════
# Slide 5 — 第二十六章 典故
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十六章｜典故案例')
section_title(s, '历史典故：刘备"髀肉复生"（稳重之志）')
body_text(s, [
    '【背景】刘备在荆州依附刘表，一日如厕时发现自己大腿上长出了肥肉。',
    '',
    '【感慨】刘备不禁潸然泪下——自寄人篱下以来，髀肉复生，',
    '而功业未建，年华虚度。他对刘表说：',
    '"从前在大腿上常年起茧，如今光景倒好，大腿又长肥肉了。',
    '日月流逝，老将至矣，而功业不建，怎能不悲伤！"',
    '',
    '【对比】同为枭雄，曹操性格轻锐果决，善于把握时机，但有时过于自信；',
    '孙权沉稳持重，守成有方；刘备以"稳重"著称，髀肉复生之事足见其心志。',
    '',
    '【感悟】"重为轻根，静为躁君"——刘备的泪水正是对这一古训的最好注脚。',
    '稳重不是迟钝，而是对根本的坚守；宁静不是无为，而是对大业的执念。',
    '领导者的轻浮是国家倾覆的先兆，刘备深谙此理，所以能三分天下有其一。'
], top=1.4)

# ═══════════════════════════════════════
# Slide 6 — 第二十七章
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十七章｜善行无辙迹')
quote_box(s, '善行无辙迹，善言无瑕谪，善数不用筹策，\n善闭无关楗而不可开，善结无绳约而不可解。\n是以圣人常善救人，故无弃人；常善救物，故无弃物。')
section_title(s, '白话解读')
body_text(s, [
    '→ 善于行走的不留痕迹，善于说话的不留破绽，善于计数的不用筹码，',
    '→ 善于关闭的无须门闩却无人能开，善于捆绑的无须绳索却无人能解。',
    '→ 因此圣人总是善于帮助人，所以没有被遗弃的人；',
    '→ 总是善于利用万物，所以没有被遗弃的物。',
    '→ 这就是袭明——承续大道之明。'
], top=3.1)
insight_box(s, '核心智慧：善行无迹，物尽其用', [
    '最高明的帮助是润物无声，不留痕迹地成就他人',
])

# ═══════════════════════════════════════
# Slide 7 — 第二十七章 典故
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十七章｜典故案例')
section_title(s, '历史典故：子贡赎人与子路受牛（善行有迹 vs 善行无迹）')
body_text(s, [
    '【子贡赎人】鲁国有一条法律：赎出在外的鲁国人，可向国库报销。',
    '子贡赎人于外，却拒绝报销，自认为这是好事。',
    '孔子批评他：你这样做，以后谁还敢赎人？报销是应得的，不应收道德绑架。',
    '',
    '【子路受牛】子路救起一个落水者，对方送了一头牛感谢他，子路欣然接受。',
    '孔子高兴地说：以后鲁国救人的人会更多！',
    '',
    '【感悟】子贡"善行"却留了"辙迹"——他的高尚反而堵住了别人行善的路；',
    '子路"善行"无迹——接受谢礼让善行变得可复制、可推广。',
    '老子说"善行无辙迹"，最高明的善行不是让人望而生畏，',
    '而是春风化雨，让更多人在不知不觉中跟随。',
    '管理者当学子路：让善行成为习惯，而非高高在上的标杆。'
], top=1.4)

# ═══════════════════════════════════════
# Slide 8 — 三章精华总结
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十五至二十七章 · 精华提炼')
for i, (ch, title, pts) in enumerate([
    ('第二十五章', '有物混成', ['道为宇宙本源，先天地生', '人法地、地法天、天法道、道法自然', '四大之中，人居其一']),
    ('第二十六章', '重为轻根', ['稳重是轻率的根本', '静为躁君，宁静致远', '轻则失本，躁则失君']),
    ('第二十七章', '善行无辙迹', ['善于助人，无人被弃', '善于用物，无物可弃', '善行无迹，润物无声']),
]):
    left = Inches(0.5 + i * 4.2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             left, Inches(1.3), Inches(4.0), Inches(4.0))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['light']
    box.line.color.rgb = COLORS['primary']
    tb = s.shapes.add_textbox(left + Inches(0.15), Inches(1.4), Inches(3.7), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = ch
    p.font.size = Pt(14); p.font.color.rgb = COLORS['secondary']
    tb2 = s.shapes.add_textbox(left + Inches(0.15), Inches(1.85), Inches(3.7), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = title
    p2.font.size = Pt(24); p2.font.bold = True; p2.font.color.rgb = COLORS['primary']
    tb3 = s.shapes.add_textbox(left + Inches(0.15), Inches(2.6), Inches(3.7), Inches(2.5))
    tf = tb3.text_frame; tf.word_wrap = True
    for j, pt in enumerate(pts):
        p3 = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p3.text = f'• {pt}'; p3.font.size = Pt(14)
        p3.font.color.rgb = COLORS['text']; p3.space_after = Pt(8)
bot = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5))
bot.fill.solid(); bot.fill.fore_color.rgb = COLORS['primary']
bot.line.fill.background()
tbb = bot.text_frame; tbb.word_wrap = True
pb = tbb.paragraphs[0]; pb.text = '三章贯通：道法自然，稳重行事，善行无迹'
pb.font.size = Pt(22); pb.font.bold = True
pb.font.color.rgb = RGBColor(255, 220, 80); pb.alignment = PP_ALIGN.CENTER
pb2 = tbb.add_paragraph()
pb2.text = '核心心法：人法道，道法自然 → 重为轻根，静为躁君 → 善行无辙迹，常善救人'
pb2.font.size = Pt(15); pb2.font.color.rgb = RGBColor(200, 230, 255); pb2.alignment = PP_ALIGN.CENTER

out = '/Users/mac/.openclaw/workspace/courses/道德经_第9天.pptx'
prs.save(out)
print(f'Saved: {out}')
