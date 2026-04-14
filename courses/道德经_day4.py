from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

# ===== 第1页：封面 =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar.line.fill.background()
txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第四课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

main_title = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main_title.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第10-12章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER

subtitle = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = subtitle.text_frame
p3 = tf3.paragraphs[0]
p3.text = "专气致柔 · 虚无之用 · 五色目盲"
p3.font.size = Pt(28)
p3.font.color.rgb = RGBColor(70, 130, 180)
p3.alignment = PP_ALIGN.CENTER

quote_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_q = quote_box.text_frame
tf_q.word_wrap = True
p_q1 = tf_q.paragraphs[0]
p_q1.text = "专气致柔，能婴儿乎？"
p_q1.font.size = Pt(22)
p_q1.font.italic = True
p_q1.font.color.rgb = RGBColor(30, 60, 114)
p_q1.alignment = PP_ALIGN.CENTER
p_q2 = tf_q.add_paragraph()
p_q2.text = "—— 《道德经》第十章"
p_q2.font.size = Pt(14)
p_q2.font.color.rgb = RGBColor(150, 150, 150)
p_q2.alignment = PP_ALIGN.CENTER

add_slide_number(slide1, 1, 10)

# ===== 第2页：第十章详解 =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar2.fill.solid()
title_bar2.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar2.line.fill.background()
txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf5 = txBox2.text_frame
p5 = tf5.paragraphs[0]
p5.text = "第十章 · 专气致柔"
p5.font.size = Pt(32)
p5.font.bold = True
p5.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.2))
yuanwen_box.fill.solid()
yuanwen_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw = yuanwen_box.text_frame
tf_yw.word_wrap = True
p_yw = tf_yw.paragraphs[0]
p_yw.text = "【原文】"
p_yw.font.size = Pt(18)
p_yw.font.bold = True
p_yw.font.color.rgb = RGBColor(30, 60, 114)
for line in ["载营魄抱一，能无离乎？", "专气致柔，能婴儿乎？", "涤除玄览，能无疵乎？"]:
    p = tf_yw.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(3.5))
tf_c = content_box.text_frame
tf_c.word_wrap = True
p_c1 = tf_c.paragraphs[0]
p_c1.text = "【白话解读】"
p_c1.font.size = Pt(18)
p_c1.font.bold = True
p_c1.font.color.rgb = RGBColor(30, 60, 114)
解读 = [
    ("营魄抱一", "精神和形体合一，能否不分离？——身心统一是健康之本。"),
    ("专气致柔", "凝聚精气达到柔顺，能否像婴儿一样？——柔弱才是生命力的象征。"),
    ("涤除玄览", "清除杂念反观内心，能否没有瑕疵？——心灵要常常洗练。"),
]
for title, content in 解读:
    p = tf_c.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide2, 2, 10)

# ===== 第3页：第十章生活智慧 =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar3.fill.solid()
title_bar3.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar3.line.fill.background()
txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf6 = txBox3.text_frame
p6 = tf6.paragraphs[0]
p6.text = "第十章 · 生活智慧"
p6.font.size = Pt(32)
p6.font.bold = True
p6.font.color.rgb = RGBColor(255, 255, 255)

core_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.3))
core_box.fill.solid()
core_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_core = core_box.text_frame
tf_core.word_wrap = True
p_core = tf_core.paragraphs[0]
p_core.text = "🌟 核心智慧"
p_core.font.size = Pt(18)
p_core.font.bold = True
p_core.font.color.rgb = RGBColor(255, 255, 255)
p_core2 = tf_core.add_paragraph()
p_core2.text = "身心合一、柔顺自然、心灵澄澈——这是回归大道的三重修养。"
p_core2.font.size = Pt(16)
p_core2.font.color.rgb = RGBColor(255, 255, 255)

points = [
    ("1. 身心健康", "现代人身心分离——身体在工作，心却在焦虑。学会专注当下，身心合一。"),
    ("2. 柔顺之道", "强硬易折，柔弱长存。太极拳正是\"专气致柔\"的最好体现。"),
    ("3. 心灵洗练", "每天给自己一段静心时间，洗去杂念，恢复心灵的清明。")
]
y = 2.9
for title, content in points:
    pt_box = slide3.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
    tf_pt = pt_box.text_frame
    tf_pt.word_wrap = True
    p_pt1 = tf_pt.paragraphs[0]
    p_pt1.text = title
    p_pt1.font.size = Pt(18)
    p_pt1.font.bold = True
    p_pt1.font.color.rgb = RGBColor(30, 60, 114)
    p_pt2 = tf_pt.add_paragraph()
    p_pt2.text = content
    p_pt2.font.size = Pt(14)
    p_pt2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.2

add_slide_number(slide3, 3, 10)

# ===== 第4页：第十一章详解 =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar4.fill.solid()
title_bar4.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar4.line.fill.background()
txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf7 = txBox4.text_frame
p7 = tf7.paragraphs[0]
p7.text = "第十一章 · 虚无之用"
p7.font.size = Pt(32)
p7.font.bold = True
p7.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
yuanwen_box2.fill.solid()
yuanwen_box2.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw2 = yuanwen_box2.text_frame
tf_yw2.word_wrap = True
p_yw2 = tf_yw2.paragraphs[0]
p_yw2.text = "【原文】"
p_yw2.font.size = Pt(18)
p_yw2.font.bold = True
p_yw2.font.color.rgb = RGBColor(30, 60, 114)
for line in ["三十辐共一毂，当其无，有车之用。", "埏埴以为器，当其无，有器之用。"]:
    p = tf_yw2.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box2 = slide4.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(4))
tf_c2 = content_box2.text_frame
tf_c2.word_wrap = True
p_c4 = tf_c2.paragraphs[0]
p_c4.text = "【白话解读】"
p_c4.font.size = Pt(18)
p_c4.font.bold = True
p_c4.font.color.rgb = RGBColor(30, 60, 114)
解读2 = [
    ("车毂之喻", "三十根辐条汇聚车毂，正因为中间空虚，车才能行走——\"无\"才是真正的\"用\"。"),
    ("器皿之喻", "揉捏泥土做成器皿，正因为中间空虚，器皿才能盛物——虚空才是价值所在。"),
    ("核心：有无相生", "有与无相互依存，有之以为利，无之以为用——看不到的空间才是关键。"),
]
for title, content in 解读2:
    p = tf_c2.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide4, 4, 10)

# ===== 第5页：有无相生典故 =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar5.fill.solid()
title_bar5.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar5.line.fill.background()
txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf8 = txBox5.text_frame
p8 = tf8.paragraphs[0]
p8.text = "📜 典故：有无相生"
p8.font.size = Pt(32)
p8.font.bold = True
p8.font.color.rgb = RGBColor(255, 255, 255)

story_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
story_box.fill.solid()
story_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_story = story_box.text_frame
tf_story.word_wrap = True
p_s1 = tf_story.paragraphs[0]
p_s1.text = "【老子出关】"
p_s1.font.size = Pt(16)
p_s1.font.bold = True
p_s1.font.color.rgb = RGBColor(30, 60, 114)
p_s2 = tf_story.add_paragraph()
p_s2.text = "函谷关令尹喜见紫气东来，知有圣人至。果然，老子骑青牛过函谷。尹喜请老子著书，老子留下《道德经》五千言，后西出函谷，不知所终。紫气东来，是因为有无形之气，才有了有形之经典。"
p_s2.font.size = Pt(15)
p_s2.font.color.rgb = RGBColor(51, 51, 51)
p_s2.space_before = Pt(8)
p_s3 = tf_story.add_paragraph()
p_s3.text = "无形的\"道\"，生出了有形的《道德经》——有因无生，道法自然。"
p_s3.font.size = Pt(15)
p_s3.font.color.rgb = RGBColor(51, 51, 51)

west_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.6), Inches(9), Inches(1))
west_box.fill.solid()
west_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_west = west_box.text_frame
tf_west.word_wrap = True
p_west = tf_west.paragraphs[0]
p_west.text = "🌍 西方哲学对照：海德格尔"
p_west.font.size = Pt(18)
p_west.font.bold = True
p_west.font.color.rgb = RGBColor(255, 255, 255)
p_west2 = tf_west.add_paragraph()
p_west2.text = "\"我们只关注存在者，却遗忘了存在本身。\"——海德格尔（1927）《存在与时间》"
p_west2.font.size = Pt(14)
p_west2.font.color.rgb = RGBColor(200, 200, 200)

app_box = slide5.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2.5))
tf_app = app_box.text_frame
tf_app.word_wrap = True
p_a1 = tf_app.paragraphs[0]
p_a1.text = "【现代启示】"
p_a1.font.size = Pt(16)
p_a1.font.bold = True
p_a1.font.color.rgb = RGBColor(30, 60, 114)
apps = [
    "• 留白艺术：书法、绘画、音乐的留白，是\"无\"的妙用",
    "• 组织管理：团队中的\"授权空间\"，让成员自主发挥",
    "• 产品设计：App的\"负空间\"界面，反而更易用"
]
for app in apps:
    p = tf_app.add_paragraph()
    p.text = app
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide5, 5, 10)

# ===== 第6页：第十二章详解 =====
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar6.fill.solid()
title_bar6.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar6.line.fill.background()
txBox6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf9 = txBox6.text_frame
p9 = tf9.paragraphs[0]
p9.text = "第十二章 · 五色目盲"
p9.font.size = Pt(32)
p9.font.bold = True
p9.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box3 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2))
yuanwen_box3.fill.solid()
yuanwen_box3.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw3 = yuanwen_box3.text_frame
tf_yw3.word_wrap = True
p_yw3 = tf_yw3.paragraphs[0]
p_yw3.text = "【原文】"
p_yw3.font.size = Pt(18)
p_yw3.font.bold = True
p_yw3.font.color.rgb = RGBColor(30, 60, 114)
for line in ["五色令人目盲；五音令人耳聋；五味令人口爽；", "驰骋畋猎，令人心发狂；难得之货，令人行妨。"]:
    p = tf_yw3.add_paragraph()
    p.text = line
    p.font.size = Pt(19)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box3 = slide6.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(4))
tf_c3 = content_box3.text_frame
tf_c3.word_wrap = True
p_c7 = tf_c3.paragraphs[0]
p_c7.text = "【白话解读】"
p_c7.font.size = Pt(18)
p_c7.font.bold = True
p_c7.font.color.rgb = RGBColor(30, 60, 114)
解读3 = [
    ("五色目盲", "缤纷色彩让人眼花缭乱——过度刺激感官的，反而伤害感官。"),
    ("五音耳聋", "嘈杂声音让人耳聋——噪音让人失去对宁静的感知。"),
    ("五味口爽", "浓烈味道让味觉失灵——重口味反而降低对淡味的感受。"),
    ("核心：适度", "外物是为人服务的，人不应成为外物的奴隶")
]
for title, content in 解读3:
    p = tf_c3.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide6, 6, 10)

# ===== 第7页：五色目盲典故 =====
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar7.fill.solid()
title_bar7.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar7.line.fill.background()
txBox7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf10 = txBox7.text_frame
p10 = tf10.paragraphs[0]
p10.text = "📜 典故：五色目盲"
p10.font.size = Pt(32)
p10.font.bold = True
p10.font.color.rgb = RGBColor(255, 255, 255)

story_box2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.2))
story_box2.fill.solid()
story_box2.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_st2 = story_box2.text_frame
tf_st2.word_wrap = True
p_st2a = tf_st2.paragraphs[0]
p_st2a.text = "【庄子·天地】"
p_st2a.font.size = Pt(16)
p_st2a.font.bold = True
p_st2a.font.color.rgb = RGBColor(30, 60, 114)
p_st2b = tf_st2.add_paragraph()
p_st2b.text = "子贡南游，过汉阴，见一老人凿隧入井，抱瓮灌溉。子贡问：\"为何不用桔槔（汲水工具）？\"老人答：\"有机械必有机心——有机心，则纯白不备，神生不定。\"——追求便利与技巧，反而失去内心的纯粹与平静。"
p_st2b.font.size = Pt(15)
p_st2b.font.color.rgb = RGBColor(51, 51, 51)
p_st2b.space_before = Pt(8)

west_box2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(9), Inches(1))
west_box2.fill.solid()
west_box2.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_w2 = west_box2.text_frame
tf_w2.word_wrap = True
p_w2a = tf_w2.paragraphs[0]
p_w2a.text = "⚡ 数字时代警示"
p_w2a.font.size = Pt(18)
p_w2a.font.bold = True
p_w2a.font.color.rgb = RGBColor(255, 255, 255)
p_w2b = tf_w2.add_paragraph()
p_w2b.text = "五色令人目盲——刷短视频3小时，脑子\"中毒\"，专注力归零"
p_w2b.font.size = Pt(15)
p_w2b.font.color.rgb = RGBColor(200, 200, 200)

app_box2 = slide7.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2.5))
tf_ap2 = app_box2.text_frame
tf_ap2.word_wrap = True
p_ap2a = tf_ap2.paragraphs[0]
p_ap2a.text = "【解决方案】"
p_ap2a.font.size = Pt(16)
p_ap2a.font.bold = True
p_ap2a.font.color.rgb = RGBColor(30, 60, 114)
apps2 = [
    "• 设定每日屏幕时间上限（如2小时）",
    "• 培养静心习惯：阅读、书法、太极",
    "• 亲近自然：山水草木，让人回归本真"
]
for app in apps2:
    p = tf_ap2.add_paragraph()
    p.text = app
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide7, 7, 10)

# ===== 第8页：西方哲学对照 =====
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar8.fill.solid()
title_bar8.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar8.line.fill.background()
txBox8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf11 = txBox8.text_frame
p11 = tf11.paragraphs[0]
p11.text = "🌍 西方思想对照"
p11.font.size = Pt(32)
p11.font.bold = True
p11.font.color.rgb = RGBColor(255, 255, 255)

philosophers = [
    ("🇬🇷 第欧根尼", "\"极简主义\"", "住在一个木桶里，晒太阳吃面包——欲望越少，自由越多。", "前4世纪"),
    ("🇺🇸 梭罗", "\"瓦尔登湖\"", "\"大多数奢侈品，让生活变得复杂。\"简化生活，返璞归真。", "1854年"),
    ("🇩🇪 本雅明", "\"机械复制\"", "工业时代使艺术品失去\"光晕\"——过度复制使文化失去灵魂。", "1935年"),
    ("🇬🇧 密尔", "\"功利主义批判\"", "追求最大幸福不等于感官放纵——高质量的快乐更深刻持久。", "1863年")
]
y = 1.4
for name, theory, desc, year in philosophers:
    phil_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.2))
    phil_box.fill.solid()
    phil_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_phil = phil_box.text_frame
    tf_phil.word_wrap = True
    p_n = tf_phil.paragraphs[0]
    p_n.text = f"{name} · {theory} · {year}"
    p_n.font.size = Pt(16)
    p_n.font.bold = True
    p_n.font.color.rgb = RGBColor(30, 60, 114)
    p_d = tf_phil.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.35

add_slide_number(slide8, 8, 10)

# ===== 第9页：生活应用 =====
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar9.fill.solid()
title_bar9.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar9.line.fill.background()
txBox9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf12 = txBox9.text_frame
p12 = tf12.paragraphs[0]
p12.text = "💡 生活中的智慧"
p12.font.size = Pt(32)
p12.font.bold = True
p12.font.color.rgb = RGBColor(255, 255, 255)

app_items = [
    ("🏢 极简办公", "办公桌上堆满文件，反而找不到重要文件\n留白才有空间感——乔布斯也坚持极简工位", ""),
    ("🧘 心灵减法", "减少无效社交、减少碎片化信息输入\n给心灵留出空白，才能思考真正重要的事", ""),
    ("💼 投资哲学", "\"有\"是账面数字，\"无\"是抗风险能力\n留足安全边际，不All-in——这才是\"虚无之用\"", ""),
    ("👨‍👩‍👧 家庭教育", "不给孩子报满培训班\n留白才有成长空间——\"专气致柔\"的下一代", "")
]
y = 1.4
for title, content, _ in app_items:
    app_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.3))
    app_box.fill.solid()
    app_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_app = app_box.text_frame
    tf_app.word_wrap = True
    p_a = tf_app.paragraphs[0]
    p_a.text = title
    p_a.font.size = Pt(16)
    p_a.font.bold = True
    p_a.font.color.rgb = RGBColor(30, 60, 114)
    p_b = tf_app.add_paragraph()
    p_b.text = content
    p_b.font.size = Pt(13)
    p_b.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.45

add_slide_number(slide9, 9, 10)

# ===== 第10页：精华总结 =====
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar10.fill.solid()
title_bar10.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar10.line.fill.background()
txBox10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf13 = txBox10.text_frame
p13 = tf13.paragraphs[0]
p13.text = "💎 精华提炼 · 第四课总结"
p13.font.size = Pt(32)
p13.font.bold = True
p13.font.color.rgb = RGBColor(255, 255, 255)

summary_items = [
    ("第十章·专气致柔", "身心抱一，返璞归真——像婴儿一样纯粹，才能获得真正的力量"),
    ("第十一章·虚无之用", "有之以为利，无之以为用——虚空才是最大的用处"),
    ("第十二章·五色目盲", "为腹不为目，回归本真——不被感官欲望支配")
]
y = 1.4
for title, content in summary_items:
    s_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.9))
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_s = s_box.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.text = f"{title}"
    p_s.font.size = Pt(16)
    p_s.font.bold = True
    p_s.font.color.rgb = RGBColor(30, 60, 114)
    p_s2 = tf_s.add_paragraph()
    p_s2.text = content
    p_s2.font.size = Pt(14)
    p_s2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.05

conclusion_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
conclusion_box.fill.solid()
conclusion_box.fill.fore_color.rgb = RGBColor(30, 60, 114)
tf_con = conclusion_box.text_frame
tf_con.word_wrap = True
p_con = tf_con.paragraphs[0]
p_con.text = "🌟 一句话总结"
p_con.font.size = Pt(18)
p_con.font.bold = True
p_con.font.color.rgb = RGBColor(255, 255, 255)
p_con2 = tf_con.add_paragraph()
p_con2.text = "上善若水，为而不争"
p_con2.font.size = Pt(24)
p_con2.font.bold = True
p_con2.font.color.rgb = RGBColor(255, 255, 255)
p_con2.alignment = PP_ALIGN.CENTER
p_con3 = tf_con.add_paragraph()
p_con3.text = "守住内心的虚空与柔软，不被外物所役，才能获得真正的自在。"
p_con3.font.size = Pt(16)
p_con3.font.color.rgb = RGBColor(200, 200, 200)

add_slide_number(slide10, 10, 10)

# 保存
prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第4天.pptx')
print("PPT已生成：道德经_第4天.pptx")
