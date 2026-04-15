from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 配色：学术白
COLORS = {
    'primary': RGBColor(0, 51, 102),     # 深蓝
    'secondary': RGBColor(102, 102, 102), # 灰
    'accent': RGBColor(204, 0, 0),       # 红色点缀
    'text': RGBColor(51, 51, 51),        # 正文
    'bg': RGBColor(255, 255, 255),       # 白色背景
    'light': RGBColor(240, 248, 255),    # 浅蓝灰
}

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['bg']

def title_bar(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(30); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def quote_box(slide, text, top=1.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(top), Inches(11.73), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['light']
    box.line.color.rgb = COLORS['primary']
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(20); p.font.name = 'KaiTi'
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    tf.margin_left = Pt(20); tf.margin_right = Pt(20)

def section_title(slide, text, top=1.4):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.73), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(22); p.font.bold = True
    p.font.color.rgb = COLORS['accent']

def body_text(slide, lines, top=2.1):
    left, top_, width, height = Inches(0.8), Inches(top), Inches(11.73), Inches(4.5)
    tb = slide.shapes.add_textbox(left, top_, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(17)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(10)

def insight_box(slide, title, points, top=5.3):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(top), Inches(11.73), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['primary']
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"💡 {title}"
    p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 220, 80)
    for pt in points:
        p2 = tf.add_paragraph()
        p2.text = f"  • {pt}"; p2.font.size = Pt(15)
        p2.font.color.rgb = RGBColor(220, 240, 255)

# Slide 1 — 封面
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(2.2))
bar.fill.solid(); bar.fill.fore_color.rgb = COLORS['primary']
bar.line.fill.background()
tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1.2))
p = tb.text_frame.paragraphs[0]
p.text = '《道德经》每日课件'; p.font.size = Pt(44); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER
tb2 = s.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.33), Inches(0.8))
p2 = tb2.text_frame.paragraphs[0]
p2.text = '第8天 · 第二十五章 至 第二十七章'; p2.font.size = Pt(26)
p2.font.color.rgb = RGBColor(200, 230, 255); p2.alignment = PP_ALIGN.CENTER
tb3 = s.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.33), Inches(0.6))
p3 = tb3.text_frame.paragraphs[0]
p3.text = '2026年4月15日'; p3.font.size = Pt(20)
p3.font.color.rgb = COLORS['secondary']; p3.alignment = PP_ALIGN.CENTER
tb4 = s.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.33), Inches(1.5))
p4 = tb4.text_frame.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
p4.text = '"有物混成，先天地生"'; p4.font.size = Pt(22)
p4.font.name = 'KaiTi'; p4.font.color.rgb = COLORS['primary']
p4b = tb4.text_frame.add_paragraph(); p4b.alignment = PP_ALIGN.CENTER
p4b.text = '—— 道先天地而生，为天下母'; p4b.font.size = Pt(16)
p4b.font.color.rgb = COLORS['secondary']

# Slide 2 — 第二十五章
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十五章｜有物混成')
quote_box(s, '有物混成，先天地生。寂兮寥兮，独立不改，\n周行而不殆，可以为天下母。\n吾不知其名，字之曰道，强为之名曰大。\n大曰逝，逝曰远，远曰反。故道大，天大，地大，王亦大。')
section_title(s, '白话解读')
body_text(s, [
    '→ 有个东西在天地之前就存在，寂静空虚却永恒运行，是天下万物的本源。',
    '→ 我不知它的名字，称它为"道"，勉强叫它"大"。',
    '→ 道广大无边而运行不息，伸展遥远而返回本原。',
    '→ 道大、天大、地大、王大，宇宙四大，人居其一。',
], top=3.1)
insight_box(s, '核心智慧：道为宇宙本源', [
    '认识"道"的本体，理解万物运行的根本规律',
])

# Slide 3 — 第二十五章 典故
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十五章｜典故案例')
section_title(s, '历史典故：庖丁解牛')
body_text(s, [
    '【出处】《庄子·养生主》',
    '',
    '【故事】有一个名为丁的厨师替梁惠王宰牛，手所触碰的地方、肩所倚靠的地方、',
    '脚所踩踏的地方、膝所抵住的地方，都发出声响，进刀时发出嚐然声响，',
    '没有一处不符合节奏。梁惠王赞叹："好啊，你的技艺怎么达到这种地步？"',
    '',
    '【奥秘】庖丁说："臣所探究的是道，已经超过一般的技术。最初我宰牛时，',
    '看到的都是完整的牛；三年之后，就再也看不到整头的牛了。"',
    '',
    '【感悟】庖丁解牛19年仍刀刃如新，因为他"以神遇而不以目视"，',
    '直接与道相合。这正是老子所说的"有物混成，先天地生"——',
    '当技艺精进到极致，便与"道"合一，以无厚入有间，游刃有余。'
], top=1.4)

# Slide 4 — 第二十六章
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十六章｜致誉无誉')
quote_box(s, '重为轻根，静为躁君。\n是以圣人终日行不离辎重，虽有荣观，燕处超然。\n奈何万乘之主，而以身轻天下？\n轻则失本，躁则失君。')
section_title(s, '白话解读')
body_text(s, [
    '→ 稳重是轻率的根本，宁静是躁动的主宰。',
    '→ 圣人每天行走都不离开载重之车，虽有荣华富贵，却超然物外不以名利为念。',
    '→ 为什么万乘之国的君主，还要用轻率的态度来治理天下呢？',
    '→ 轻率会失去根本，躁动会失去主宰。',
], top=3.1)
insight_box(s, '核心智慧：静重为安身之本', [
    '内心宁静方能驾驭外在躁动，稳重行事方能立于不败',
])

# Slide 5 — 第二十六章 典故
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十六章｜典故案例')
section_title(s, '历史典故：刘备与曹操的"静重"之别')
body_text(s, [
    '【背景】三国时期，刘备与曹操同为枭雄，但性格截然不同。',
    '',
    '【曹操】性格刚猛果断，喜欢亲临前线，率兵冲锋。',
    '官渡之战时亲自带队夜袭乌巢，展现出极强的个人魄力与行动力。',
    '但也因此多次陷入险境，如宛城之战差点被张绣活捉。',
    '',
    '【刘备】虽起步晚、根基薄，但始终保持沉稳。',
    '三顾茅庐请诸葛亮出山，在位期间蜀汉政治相对稳定。',
    '尽管军事才能不及曹操，但因行事稳重而得人心。',
    '',
    '【感悟】刘备的成功在于"静为躁君"——不以一时胜负乱了方寸，',
    '始终坚持"复兴汉室"的政治理念，最终与魏、吴三分天下。',
    '而曹操虽能力出众，但"轻则失本"的教训也多次上演。',
    '对于管理者而言，团队的稳定根基在于领导者的静重与沉稳。'
], top=1.4)

# Slide 6 — 第二十七章
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十七章｜希言自然')
quote_box(s, '希言自然。\n飘风不终朝，骤雨不终日。\n孰为此者？天地。天地尚不能久，而况于人乎？\n故从事于道者，道者同于道，德者同于德，失者同于失。')
section_title(s, '白话解读')
body_text(s, [
    '→ 少言才合乎自然之道。狂风刮不了一早晨，暴雨下不了一整天。',
    '→ 谁造成这一切？是天地。天地尚且不能持久，何况人呢？',
    '→ 追求道的人，与道同行；追求德的人，与德同行；失去道德的人，则与失同行。',
    '→ 与道同行的人，道也乐于成就他；与德同行的人，德也乐于成就他。',
], top=3.1)
insight_box(s, '核心智慧：少言慎行，顺应自然节奏', [
    '暴风骤雨虽猛烈却短暂，合道而行虽缓慢却长久',
])

# Slide 7 — 第二十七章 典故
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十七章｜典故案例')
section_title(s, '历史典故：大禹治水与"飘风骤雨"')
body_text(s, [
    '【背景】上古时期，洪水泛滥，百姓流离失所。',
    '',
    '【鲧治水】鲧用"堵"的方法，筑堤防水，结果九年不成，洪水反而更凶。',
    '【大禹治水】大禹改用"疏导"的方法，顺应水往低处流的自然规律，',
    '开凿河道，引水入海，终于平息了水患。',
    '',
    '【对比】鲧的"猛攻"如同飘风骤雨，短期内看似有力，实则违背自然规律，注定失败。',
    '大禹的"疏导"如同"希言自然"——不多言、不妄为，只是顺应水的本性加以引导。',
    '',
    '【成果】大禹在外治水十三年，三过家门而不入，最终成功。',
    '',
    '【感悟】大禹的成功在于"道法自然"——不强为、不妄作，顺应事物本身的规律，',
    '以最小的干预实现最大的效果。这正是"飘风不终朝，骤雨不终日"给我们的启示：',
    '强力为之的事情不会长久，只有合乎道的作为才能持久。'
], top=1.4)

# Slide 8 — 三章精华总结
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, '第二十五至二十七章 · 精华提炼')
for i, (ch, title, pts) in enumerate([
    ('第二十五章', '有物混成', ['道先天地而生，为天下母', '道大、天大、地大、王大', '大曰逝，逝曰远，远曰反']),
    ('第二十六章', '致誉无誉', ['重为轻根，静为躁君', '稳重行事方能立于不败', '轻则失本，躁则失君']),
    ('第二十七章', '希言自然', ['少言慎行顺应自然', '飘风骤雨不能长久', '同于道者，道亦乐得之']),
]):
    left = Inches(0.5 + i * 4.2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             left, Inches(1.3), Inches(4.0), Inches(4.0))
    box.fill.solid(); box.fill.fore_color.rgb = COLORS['light']
    box.line.color.rgb = COLORS['primary']
    tb = s.shapes.add_textbox(left + Inches(0.15), Inches(1.4), Inches(3.7), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = ch
    p.font.size = Pt(14); p.font.color.rgb = COLORS['secondary']
    tb2 = s.shapes.add_textbox(left + Inches(0.15), Inches(1.85), Inches(3.7), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = title
    p2.font.size = Pt(24); p2.font.bold = True; p2.font.color.rgb = COLORS['primary']
    tb3 = s.shapes.add_textbox(left + Inches(0.15), Inches(2.6), Inches(3.7), Inches(2.5))
    tf = tb3.text_frame; tf.word_wrap = True
    for j, pt in enumerate(pts):
        p3 = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p3.text = f'• {pt}'; p3.font.size = Pt(14)
        p3.font.color.rgb = COLORS['text']; p3.space_after = Pt(8)
bot = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5))
bot.fill.solid(); bot.fill.fore_color.rgb = COLORS['primary']
bot.line.fill.background()
tbb = bot.text_frame; tbb.word_wrap = True
pb = tbb.paragraphs[0]; pb.text = '三章贯通：认识道的本体，保持静重之心，顺应自然之道'
pb.font.size = Pt(20); pb.font.bold = True
pb.font.color.rgb = RGBColor(255, 220, 80); pb.alignment = PP_ALIGN.CENTER
pb2 = tbb.add_paragraph(); pb2.text = '上善若水，水善利万物而不争。'
pb2.font.size = Pt(16); pb2.font.color.rgb = RGBColor(200, 230, 255); pb2.alignment = PP_ALIGN.CENTER

# 保存
out = '/Users/mac/.openclaw/workspace/courses/道德经_第8天_正式版.pptx'
prs.save(out)
print(f'Saved: {out}')
