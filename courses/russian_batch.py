#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
俄语课件批量生成 - 第2-7天
名词第3格、第4格、第5格、第6格 + 形容词性数配合
"""
import os
import sys
sys.path.insert(0, '/opt/homebrew/lib/node_modules/openclaw/node_modules')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(30, 60, 114)
TEXT = RGBColor(51, 51, 51)
BLUE_LIGHT = RGBColor(70, 130, 180)
GOLD = RGBColor(180, 140, 50)
WHITE = RGBColor(255, 255, 255)

def make_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(200, 220, 255)

def add_footer(slide, day_num):
    footer = slide.shapes.add_textbox(Inches(8.5), Inches(7.2), Inches(1.3), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Day%d/84" % day_num
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(160, 160, 160)
    p.alignment = PP_ALIGN.RIGHT

def add_box(slide, text, left, top, width, height, fill_color=None, font_size=16, color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if color:
        p.font.color.rgb = color
    return box

def make_lesson(day_num, title_cn, title_ru, slides_content, output_dir):
    prs = Presentation()
    
    # Cover
    s = make_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    main = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = main.text_frame
    p = tf.paragraphs[0]
    p.text = "俄语基础语法 · 第%d天" % day_num
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    sub = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    tf2 = sub.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title_cn
    p2.font.size = Pt(36)
    p2.font.color.rgb = TEXT
    p2.alignment = PP_ALIGN.CENTER
    
    ru_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5), Inches(6), Inches(0.8))
    ru_box.fill.solid()
    ru_box.fill.fore_color.rgb = BLUE_LIGHT
    ru_box.line.fill.background()
    tf = ru_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_ru
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    add_footer(s, day_num)
    
    # Content slides
    for slide_data in slides_content:
        s = make_slide(prs)
        set_title_bar(s, slide_data["title"], slide_data.get("subtitle", ""))
        
        y = 1.6
        for item in slide_data["content"]:
            if isinstance(item, tuple):
                if item[0] == "header":
                    h = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
                    tf = h.text_frame
                    p = tf.paragraphs[0]
                    p.text = item[1]
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = PRIMARY
                    y += 0.4
                elif item[0] == "example":
                    add_box(s, item[1], Inches(0.5), Inches(y), Inches(9), Inches(0.55), 
                            fill_color=BLUE_LIGHT, font_size=15, color=WHITE)
                    y += 0.7
                elif item[0] == "note":
                    add_box(s, item[1], Inches(0.5), Inches(y), Inches(9), Inches(0.5),
                            fill_color=GOLD, font_size=14, color=WHITE)
                    y += 0.65
            else:
                t = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
                tf = t.text_frame
                p = tf.paragraphs[0]
                p.text = item
                p.font.size = Pt(15)
                p.font.color.rgb = TEXT
                y += 0.42
        
        add_footer(s, day_num)
    
    filename = "俄语_Day%d_L1.pptx" % day_num
    path = os.path.join(output_dir, filename)
    prs.save(path)
    print("Done: %s" % filename)
    return path

output_dir = "/Users/mac/.openclaw/workspace/courses/俄语课件"
os.makedirs(output_dir, exist_ok=True)

# =====================
# Day 2: 名词第3格
# =====================
day2_slides = [
    {"title": "一、第3格的意义", "subtitle": "Datelny padesh",
     "content": [
         ("header", "第3格 = 间接宾语，表示'给谁/为谁'"),
         "Daite knigu studentu. = 把书给学生。",
         ("example", "Ya pomogayu drugu. = 我帮助朋友。"),
         ("note", "第3格回答问题：komu? chemu? (给谁？向谁？)"),
     ]},
    {"title": "二、第3格变化规则", "subtitle": "Sklonenie",
     "content": [
         ("header", "阳性名词："),
         "student > studentu (给学生)",
         "drug > drugu (给朋友)",
         ("header", "阴性名词："),
         "kniga > knige (给书)",
         "devushka > devushke (给姑娘)",
         ("note", "规律：词尾加 -u 或 -e"),
     ]},
    {"title": "三、人称代词第3格", "subtitle": "Mestoimeniya",
     "content": [
         ("header", "我 > мне, 你 > тебе, 他 > ему"),
         ("header", "她 > ей, 我们 > нам, 你们 > вам"),
         ("example", "Oni pomogayut mne. = 他们帮助我。"),
         ("example", "Ya dayu knigu yey. = 我把书给她。"),
         ("note", "记忆：人称代词第3格需要单独记忆！"),
     ]},
    {"title": "四、第3格与动词搭配", "subtitle": "Glagoly s datelnym",
     "content": [
         ("header", "常用动词："),
         "pomogat (帮助) + 第3格",
         "davat (给) + 第3格",
         "rasskazyvat (告诉) + 第3格",
         ("example", "Uchitel pomogayet studentam. = 老师帮助学生们。"),
     ]},
    {"title": "五、课后练习", "subtitle": "Uprazhneniya",
     "content": [
         "1. 将下列名词变为第3格：",
         "   dom, kniga, drug, student, devushka",
         "",
         "2. 用人称代词填空：",
         "   Ya dayu knigu ___ (on).",
         "   On pomogayet ___ (ona).",
         "",
         "3. 中译俄：",
         "   我把书给朋友。",
         "   妈妈帮助女儿。",
     ]},
]

# =====================
# Day 3: 名词第4格
# =====================
day3_slides = [
    {"title": "一、第4格的意义", "subtitle": "Vinitelny padesh",
     "content": [
         ("header", "第4格 = 直接宾语，表示'做什么/把谁'"),
         "Ya chitayu knigu. = 我读书。",
         ("example", "On vidit drugа. = 他看见朋友。"),
         ("note", "第4格回答问题：kogo? chto? (谁？什么？)"),
     ]},
    {"title": "二、第4格变化规则", "subtitle": "Sklonenie",
     "content": [
         ("header", "阳性名词（活动物）："),
         "student > studenta (看见学生)",
         "drug > druga (看见朋友)",
         ("header", "阳性名词（非活动物）："),
         "dom > dom (看见房子)",
         ("header", "阴性名词："),
         "kniga > knigu (看见书)",
         ("note", "重点：活动物阳性第4格同第2格！"),
     ]},
    {"title": "三、第四格与完成体", "subtitle": "Sovershennyy vid",
     "content": [
         ("header", "未完成体：Ya chitayu knigu. = 我在读书。"),
         ("header", "完成体：Ya prochitayu knigu. = 我读完书了。"),
         ("example", "Ya-delayu zadachu. = 我在做作业。"),
         ("example", "Ya sdayu zadachu. = 我提交作业。"),
         ("note", "完成体强调结果，未完成体强调过程"),
     ]},
    {"title": "四、课后练习", "subtitle": "Uprazhneniya",
     "content": [
         "1. 将下列名词变为第4格：",
         "   dom, kniga, drug, student",
         "",
         "2. 中译俄：",
         "   我在读书。",
         "   他看见了朋友。",
         "   妈妈在做汤。",
     ]},
]

# =====================
# Day 4: 名词第5格
# =====================
day4_slides = [
    {"title": "一、第5格的意义", "subtitle": "Tvoritelny padesh",
     "content": [
         ("header", "第5格 = 工具格，表示'用什么/和谁一起'"),
         "Ya pishu ruchkoy. = 我用笔写字。",
         ("example", "On idyot s drugom. = 他和朋友一起走。"),
         ("note", "第5格回答问题：kem? chem? (用谁？用什么？)"),
     ]},
    {"title": "二、第5格变化规则", "subtitle": "Sklonenie",
     "content": [
         ("header", "阳性/中性名词："),
         "student > studentom (用学生)",
         "okno > oknom (用窗户)",
         ("header", "阴性名词："),
         "kniga > knigoy (用书)",
         "devushka > devushkoy (和姑娘一起)",
         ("note", "规律：加 -om 或 -oy"),
     ]},
    {"title": "三、第5格与动词搭配", "subtitle": "Glagoly",
     "content": [
         ("header", "常用动词："),
         "stat (站在) + 第5格",
         "saditsya (坐在) + 第5格",
         "gorit (燃烧) + 第5格",
         ("example", "Kniga lezhit na stole. = 书在桌子上。"),
         ("note", "注意：на + 第6格 = 在...上面"),
     ]},
    {"title": "四、课后练习", "subtitle": "Uprazhneniya",
     "content": [
         "1. 将下列名词变为第5格：",
         "   dom, kniga, drug, okno",
         "",
         "2. 中译俄：",
         "   我用笔写字。",
         "   她和朋友一起走。",
         "   他站在窗边。",
     ]},
]

# =====================
# Day 5: 名词第6格
# =====================
day5_slides = [
    {"title": "一、第6格的意义", "subtitle": "Predlozhny padesh",
     "content": [
         ("header", "第6格 = 前置词格，必须与前置词连用"),
         ("header", "常用前置词：в/на (在...里/上)"),
         ("example", "Kniga na stole. = 书在桌子上。"),
         ("example", "Ya v universitete. = 我在大学里。"),
         ("note", "第6格回答问题：o kom? o chem? (关于谁？关于什么？)"),
     ]},
    {"title": "二、第6格变化规则", "subtitle": "Sklonenie",
     "content": [
         ("header", "阳性/中性名词："),
         "student > v universitete (在大学里)",
         "okno > na okne (在窗户上)",
         ("header", "阴性名词："),
         "kniga > v knige (在书里)",
         "devushka > o devushke (关于姑娘)",
         ("note", "规律：加 -e 或 -i"),
     ]},
    {"title": "三、六格综合练习", "subtitle": "Uprazhneniya",
     "content": [
         "1. 将下列名词变为各格（单数）：",
         "   dom, kniga, student, okno",
         "",
         "2. 用适当的格填空：",
         "   Ya idu ___ (дом).",
         "   Kniga lezhit ___ (стол).",
         "   Ya pomogayu ___ (друг).",
     ]},
]

# =====================
# Day 6: 形容词性数配合（一）
# =====================
day6_slides = [
    {"title": "一、形容词性数配合原则", "subtitle": "Soglasovaniye prilagatelnykh",
     "content": [
         ("header", "形容词必须与名词在性、数上一致："),
         ("header", "阳性：большой дом (大房子)"),
         ("header", "阴性：большая kniga (大书)"),
         ("header", "中性：большoye okno (大窗户)"),
         ("note", "记住：性数一致是俄语形容词的核心规则！"),
     ]},
    {"title": "二、阳性形容词词尾", "subtitle": "Okonchaniya",
     "content": [
         ("header", "以 -ый / -ой 结尾的形容词："),
         "новый (新的) > нового, новому...",
         "большой (大的) > большого, большому...",
         ("header", "以 -ний 结尾的形容词："),
         "хороший (好的) > хорошего, хорошему...",
         ("note", "规律：阳性形容词回答Какой? (什么样的？)"),
     ]},
    {"title": "三、阴性形容词词尾", "subtitle": "Zhenskiy rod",
     "content": [
         ("header", "以 -ая / -яя 结尾的形容词："),
         "новый > новая kniga (新书)",
         "хороший > хорошая devushka (好姑娘)",
         ("header", "特殊："),
         "большой > большая komnata (大房间)",
         ("note", "规律：阴性形容词回答Какая? (什么样的？)"),
     ]},
    {"title": "四、课后练习", "subtitle": "Uprazhneniya",
     "content": [
         "1. 将形容词变为适当形式：",
         "   большой ___ (дом)",
         "   новая ___ (книга)",
         "   хороший ___ (друг)",
         "",
         "2. 中译俄：",
         "   一本新书",
         "   一个好学生",
         "   一扇大窗户",
     ]},
]

# =====================
# Day 7: 形容词性数配合（二）
# =====================
day7_slides = [
    {"title": "一、中性形容词词尾", "subtitle": "Sredniy rod",
     "content": [
         ("header", "以 -ое / -ee 结尾的形容词："),
         "новое окно (新窗户)",
         "большое здание (大建筑)",
         ("header", "以 -нее 结尾的形容词："),
         "хорошее настроение (好心情)",
         ("note", "规律：中性形容词回答Какое? (什么样的？)"),
     ]},
    {"title": "二、形容词短尾形式", "subtitle": "Kratkaya forma",
     "content": [
         ("header", "有些形容词有短尾形式（作谓语用）："),
         "Он красивый. > Он красив. (他帅)",
         "Она красивая. > Она красива. (她漂亮)",
         ("header", "短尾形容词用在句子中作表语："),
         ("example", "Этот дом bolshoy. > Этот дом bolen. (这房子大)"),
         ("note", "注意：短尾形容词只用作谓语！"),
     ]},
    {"title": "三、综合练习", "subtitle": "Uprazhneniya",
     "content": [
         "1. 将形容词变为适当形式（性数配合）：",
         "   большой ___ (дом, окно, комната)",
         "   новый ___ (книга, письмо, студент)",
         "",
         "2. 选用适当形容词填空：",
         "   ___ погода (хороший/хорошая/хорошее)",
         "   ___ студент (умный/умная/умное)",
     ]},
]

# 生成所有课件
all_days = [
    (2, "名词第3格", "Datelny padesh", day2_slides),
    (3, "名词第4格", "Vinitelny padesh", day3_slides),
    (4, "名词第5格", "Tvoritelny padesh", day4_slides),
    (5, "名词第6格", "Predlozhny padesh", day5_slides),
    (6, "形容词性数配合（一）", "Soglasovaniye 1", day6_slides),
    (7, "形容词性数配合（二）", "Soglasovaniye 2", day7_slides),
]

for day_num, title_cn, title_ru, slides in all_days:
    make_lesson(day_num, title_cn, title_ru, slides, output_dir)

print("\n第2-7天课件生成完成！")
