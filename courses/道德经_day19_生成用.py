# -*- coding: utf-8 -*-
"""
道德经第19天 课件生成器
章节：第61-63章
风格：精读深度讲解版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

PRIMARY = RGBColor(30, 60, 114)
SECONDARY = RGBColor(102, 102, 102)
ACCENT = RGBColor(204, 0, 0)
TEXT = RGBColor(51, 51, 51)
BG = RGBColor(245, 245, 240)
BLUE_LIGHT = RGBColor(70, 130, 180)

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def set_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_yuanwen_box(slide, lines, y=1.4, height=2.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "【原文】"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    for line in lines:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(20)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)

def add_core_box(slide, title, subtitle, color=BLUE_LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(255, 255, 255)

def add_points(slide, points, y_start=3.1):
    y = y_start
    for title, content in points:
        pt_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.1))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT
        y += 1.1
    return y

# ===== 第1页：封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.5))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第19课"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
main = s1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第61-63章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.alignment = PP_ALIGN.CENTER
sub = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = sub.text_frame
p3 = tf3.paragraphs[0]
p3.text = "大国者下流 · 道者万物之奥 · 为无为"
p3.font.size = Pt(28)
p3.font.color.rgb = SECONDARY
p3.alignment = PP_ALIGN.CENTER
quote = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = BG
tfq = quote.text_frame
tfq.word_wrap = True
pq1 = tfq.paragraphs[0]
pq1.text = "大国者下流，天下之交。"
pq1.font.size = Pt(24)
pq1.font.italic = True
pq1.font.color.rgb = PRIMARY
pq1.alignment = PP_ALIGN.CENTER
pq2 = tfq.add_paragraph()
pq2.text = "—— 第六十一章"
pq2.font.size = Pt(14)
pq2.font.color.rgb = SECONDARY
pq2.alignment = PP_ALIGN.CENTER
add_slide_number(s1, 1, 12)

# ===== 第2页：第六十一章详解 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s2, "第六十一章 · 大国者下流")
add_yuanwen_box(s2, [
    "大国者下流，天下之交。",
    "天下之牝，牝常以静胜牡，",
    "以静为下。",
    "",
    "故大国以下小国，",
    "则取小国；",
    "小国以下大国，",
    "则取大国。",
    "",
    "故或下以取，",
    "或下而取。",
    "大国不过欲兼畜人，",
    "小国不过欲入事人。",
    "",
    "夫两者各得其所欲，",
    "大者宜为下。"
], y=1.4, height=4.8)
content2 = s2.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(2))
tf_c2 = content2.text_frame
tf_c2.word_wrap = True
pc2 = tf_c2.paragraphs[0]
pc2.text = "【白话解读】大国应当像江河的下游，处于天下交汇之地。雌柔常以安静战胜雄强，因为安静处于下方。所以大国用谦下的态度对待小国，就能取得小国的信任；小国用谦下的态度对待大国，就能被大国容纳。双方各取所需，大国尤其应该处下。"
pc2.font.size = Pt(13)
pc2.font.color.rgb = TEXT
add_slide_number(s2, 2, 12)

# ===== 第3页：第六十一章核心 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s3, "第六十一章 · 核心思想")
add_core_box(s3, "🌟 核心观点", "大国者下流——真正强大的国家懂得处下，以谦和的态度对待小国，才能赢得尊重与信任。")
add_points(s3, [
    ("1. 下流之道", "大国像江河下游，包容万流——强大的同时保持谦下，才是真正的强大。"),
    ("2. 牝静胜牡", "雌柔以静胜刚——安静、谦下才是真正的力量。"),
    ("3. 两者各得其欲", "大国小国各取所需——外交的本质是互利，不是单方索取。"),
    ("4. 大者宜为下", "越强大越要谦下——德高望重的领导者从不盛气凌人。")
])
add_slide_number(s3, 3, 12)

# ===== 第4页：第六十二章详解 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s4, "第六十二章 · 道者万物之奥")
add_yuanwen_box(s4, [
    "道者，万物之奥。",
    "善人之宝，不善人之所保。",
    "",
    "美言可以市，",
    "尊行可以加人。",
    "人之不善，何弃之有？",
    "",
    "故立天子，置三公，",
    "虽有拱璧以先驷马，",
    "不如坐进此道。",
    "",
    "古之所以贵此道者何？",
    "不曰：求以得，有罪以免与？",
    "故为天下贵。"
], y=1.4, height=4.2)
content4 = s4.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(2))
tf_c4 = content4.text_frame
tf_c4.word_wrap = True
pc4 = tf_c4.paragraphs[0]
pc4.text = "【白话解读】道是万物的庇护所。善良的人把它当作宝贝，不善的人也用它来保护自己。美好的言语可以换来尊敬，善良的行为可以提升人。对不善的人，道又怎会舍弃他们呢？所以即使贵为天子，设置三公，拥有美玉驷马，也不如安心求道。自古以来人们为何如此看重道？因为有求就能得到，有罪也能被赦免——道是天下最珍贵的东西。"
pc4.font.size = Pt(13)
pc4.font.color.rgb = TEXT
add_slide_number(s4, 4, 12)

# ===== 第5页：第六十二章核心 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s5, "第六十二章 · 核心思想")
add_core_box(s5, "🌟 核心观点", "道者万物之奥——道庇护万物，无论善与不善，有求就能得，有罪就能免。")
add_points(s5, [
    ("1. 万物之奥", "道是万物的庇护所——不论好人坏人都需要道的保护。"),
    ("2. 善人之宝", "善良的人把道当作珍宝——道让善人的善行更加坚定。"),
    ("3. 不善人之所保", "不善的人也能用道保护自己——道慈悲广大，不弃一人。"),
    ("4. 求以得，有罪以免", "有求就能得到，有罪可以被赦免——道的恩典是无条件的。")
])
add_slide_number(s5, 5, 12)

# ===== 第6页：第六十三章详解 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s6, "第六十三章 · 为无为")
add_yuanwen_box(s6, [
    "为无为，事无事，味无味。",
    "",
    "大小多少，报怨以德。",
    "",
    "图难于其易，",
    "为大于其细。",
    "天下难事，必作于易；",
    "天下大事，必作于细。",
    "",
    "是以圣人终不为大，",
    "故能成其大。",
    "",
    "夫轻诺必寡信，",
    "多易必多难。",
    "",
    "是以圣人犹难之，",
    "故终无难。"
], y=1.4, height=4.5)
content6 = s6.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(2))
tf_c6 = content6.text_frame
tf_c6.word_wrap = True
pc6 = tf_c6.paragraphs[0]
pc6.text = "【白话解读】以无为的态度去作为，以不搅扰的方式去行事，以恬淡无味为滋味。无论大事小事、大仇大怨，都以德相报。做难的事要从容易的地方入手，做大事要从细小的地方着眼。容易的事认真做，小事认真做，就是难事大事的解决之道。圣人始终不做大事，所以能成就大事。"
pc6.font.size = Pt(13)
pc6.font.color.rgb = TEXT
add_slide_number(s6, 6, 12)

# ===== 第7页：第六十三章核心 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s7, "第六十三章 · 核心思想")
add_core_box(s7, "🌟 核心观点", "为无为，事无事——最高明的作为是不妄为，从容易和小处着手，终能成大事。")
add_points(s7, [
    ("1. 为无为", "以无为的态度去作为——不折腾、不妄动、不多事，顺其自然。"),
    ("2. 图难于其易", "做难事先从容易的地方入手——防微杜渐，把问题扼杀在萌芽。"),
    ("3. 报怨以德", "用德行来回应怨恨——以德报怨，化解矛盾而非激化矛盾。"),
    ("4. 终无难", "圣人把容易的事当难事认真做——所以最终没有真正的难事。")
])
add_slide_number(s7, 7, 12)

# ===== 第8页：现代反思 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s8, "第六十一-六十三章 · 现代反思")
items8 = [
    ("🏛️ 外交哲学", "大国者下流 → 强大的国家更需要谦和，否则众叛亲离。"),
    ("💰 理财智慧", "图难于其易 → 理财从存钱开始，从小额投资开始，不要想着一夜暴富。"),
    ("🔄 为人处世", "报怨以德 → 与其记仇不如以德报怨，格局大的人不纠缠小事。"),
    ("🌱 长期主义", "为大于其细 → 成就大业的人，都懂得把每一件小事做到极致。")
]
y8 = 1.4
for icon_title, explanation in items8:
    box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y8), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf8 = box.text_frame
    tf8.word_wrap = True
    p8_1 = tf8.paragraphs[0]
    p8_1.text = icon_title
    p8_1.font.size = Pt(16)
    p8_1.font.bold = True
    p8_1.font.color.rgb = PRIMARY
    p8_2 = tf8.add_paragraph()
    p8_2.text = explanation
    p8_2.font.size = Pt(14)
    p8_2.font.color.rgb = TEXT
    y8 += 1.4
add_slide_number(s8, 8, 12)

# ===== 第9页：西方思想 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s9, "🌍 西方思想对照")
phils = [
    ("修昔底德 · 雅典与斯巴达", "新兴大国必然挑战现有霸权，但战争并非不可避免——大国博弈更需要智慧而非冲动。"),
    ("马基雅维利 · 君主论", "君主应以被畏惧而非被爱戴为目标——但老子认为谦下才是长久之道。约1513年"),
    ("老子的\"处下\"与林肯", "林肯：没有谁比我更渺小——谦逊的领导者赢得最深的尊重。1855年"),
    ("修昔底德陷阱", "格雷厄姆·艾利森：大国竞争不必走向战争——修昔底德陷阱可以被智慧规避。2012年")
]
y9 = 1.4
for name_theory, desc_year in phils:
    box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y9), Inches(9), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf9 = box.text_frame
    tf9.word_wrap = True
    p9_1 = tf9.paragraphs[0]
    p9_1.text = name_theory
    p9_1.font.size = Pt(15)
    p9_1.font.bold = True
    p9_1.font.color.rgb = PRIMARY
    p9_2 = tf9.add_paragraph()
    p9_2.text = desc_year
    p9_2.font.size = Pt(13)
    p9_2.font.color.rgb = TEXT
    y9 += 1.3
add_slide_number(s9, 9, 12)

# ===== 第10页：现代应用 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s10, "📌 现代应用")
add_points(s10, [
    ("企业管理", "大国者下流 → 越是大企业，越要谦和对待供应商和员工，才能长久。"),
    ("投资理财", "图难于其易 → 投资理财从稳健开始，不要被高风险产品迷惑。"),
    ("人际交往", "报怨以德 → 记仇伤身，以德报怨不是软弱而是格局。"),
    ("个人修养", "为无为 → 少折腾、少焦虑，顺其自然反而更容易成事。")
])
add_slide_number(s10, 10, 12)

# ===== 第11页：中国典故 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s11, "📜 中国典故：刘备的\"以德服人\"")
box11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.5))
box11.fill.solid()
box11.fill.fore_color.rgb = BG
tf11 = box11.text_frame
tf11.word_wrap = True
p11_1 = tf11.paragraphs[0]
p11_1.text = "【三国演义】"
p11_1.font.size = Pt(16)
p11_1.font.bold = True
p11_1.font.color.rgb = PRIMARY
p11_2 = tf11.add_paragraph()
p11_2.text = '刘备三顾茅庐请诸葛亮出山，面对诸葛亮的傲慢与考验毫不动怒，最终以真诚打动了卧龙先生。刘备的成功不是靠武力，而是靠德行和谦和——他懂得"大国者下流"的智慧，用谦下的态度赢得了天下英才的归心，建立了蜀汉政权。'
p11_2.font.size = Pt(15)
p11_2.font.color.rgb = TEXT
p11_2.space_before = Pt(10)
insight11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.1), Inches(9), Inches(1))
insight11.fill.solid()
insight11.fill.fore_color.rgb = BLUE_LIGHT
tfi11 = insight11.text_frame
tfi11.word_wrap = True
pi11 = tfi11.paragraphs[0]
pi11.text = "💡 启示"
pi11.font.size = Pt(18)
pi11.font.bold = True
pi11.font.color.rgb = RGBColor(255, 255, 255)
pi11_2 = tfi11.add_paragraph()
pi11_2.text = '真正的大国、大家、大人物，都懂得"处下"的智慧。谦和不是软弱，而是最强大的力量。'
pi11_2.font.size = Pt(14)
pi11_2.font.color.rgb = RGBColor(255, 255, 255)
app11 = s11.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2))
tfa11 = app11.text_frame
tfa11.word_wrap = True
pa11 = tfa11.paragraphs[0]
pa11.text = "【现代应用】职场：谦和的领导更能凝聚团队力量。商场：服务态度决定客户忠诚度。人生：越谦逊的人越能获得尊重。"
pa11.font.size = Pt(14)
pa11.font.color.rgb = TEXT
add_slide_number(s11, 11, 12)

# ===== 第12页：总结 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_title_bar(s12, "第19课 · 三章精华总结")
sums = [
    ("第六十一章", "大国者下流", "真正强大的国家懂得处下，以谦和赢得尊重与信任。"),
    ("第六十二章", "道者万物之奥", "道庇护万物，有求以得，有罪以免，是天下最贵之物。"),
    ("第六十三章", "为无为", "图难于其易，为大于其细——把容易和小事做好，就是难事大事的成功之道。")
]
y12 = 1.4
for ch, title, content in sums:
    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y12), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    tf12 = box.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = f"📖 {ch}：{title}"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = PRIMARY
    p12_2 = tf12.add_paragraph()
    p12_2.text = content
    p12_2.font.size = Pt(15)
    p12_2.font.color.rgb = TEXT
    y12 += 1.6
final = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(9), Inches(1))
final.fill.solid()
final.fill.fore_color.rgb = PRIMARY
tff = final.text_frame
tff.word_wrap = True
pf1 = tff.paragraphs[0]
pf1.text = "为无为，事无事，味无味。"
pf1.font.size = Pt(20)
pf1.font.bold = True
pf1.font.color.rgb = RGBColor(255, 255, 255)
pf1.alignment = PP_ALIGN.CENTER
pf2 = tff.add_paragraph()
pf2.text = "—— 道德经第六十三章"
pf2.font.size = Pt(14)
pf2.font.color.rgb = RGBColor(200, 200, 200)
pf2.alignment = PP_ALIGN.CENTER
add_slide_number(s12, 12, 12)

prs.save('/Users/mac/.openclaw/workspace/courses/道德经_第19天.pptx')
print("Done: 道德经_第19天.pptx (12 slides, 61-63章)")
