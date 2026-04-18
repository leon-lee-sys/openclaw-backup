#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title, color=RGBColor(0, 51, 102)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_footer(slide, text):
    ft = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.4))
    tf = ft.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(250, 250, 250)
PRIMARY = RGBColor(0, 51, 102)
ACCENT = RGBColor(200, 50, 50)

# ==================== Slide 1: 封面 ====================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1, WHITE)
add_title_bar(s1, "国学经典24章 · 第十一课", PRIMARY)

# 主标题
mt = s1.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
tf = mt.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "《道德经》第34-36章"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

# 副标题
st = s1.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(1.2))
tf = st.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "大道泛兮 · 执大象 · 将欲取之"
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(100, 100, 100)
p.alignment = PP_ALIGN.CENTER

# 分隔线
line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.6), Inches(3), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = PRIMARY
line.line.fill.background()

# 日期和天数
dt = s1.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.6))
tf = dt.text_frame
p = tf.paragraphs[0]
p.text = "2026年4月18日 · 第11天"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(130, 130, 130)
p.alignment = PP_ALIGN.CENTER

add_footer(s1, "小燕子出品 🐦")

# ==================== Slide 2: 第三十四章 原文+解读 ====================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, LIGHT_BG)
add_title_bar(s2, "第三十四章 · 大道泛兮")
add_footer(s2, "2/10")

# 原文框
yb = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(9.2), Inches(1.6))
yb.fill.solid()
yb.fill.fore_color.rgb = RGBColor(240, 245, 250)
yb.line.color.rgb = PRIMARY

yt = s2.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(1.4))
tf = yt.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "【原文】"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = ACCENT
p2 = tf.add_paragraph()
p2.text = "大道泛兮，其可左右。万物恃之而生而不辞，功成不名有。衣养万物而不为主，常无欲，可名于小。万物归焉而不为主，可名为大。以其终不自为大，故能成其大。"
p2.font.size = Pt(14)
p2.font.color.rgb = RGBColor(50, 50, 50)
p2.space_before = Pt(6)

# 解读
jd = s2.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(3.8))
tf = jd.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "【白话解读】"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY

items = [
    ("大道泛兮，其可左右", "道广泛流行，无所不在，如同水一样流淌万物之中"),
    ("万物恃之而生而不辞", "万物依靠道生存，道从不推辞，从不夸耀"),
    ("功成不名有", "成就了功业却不居功、不占有"),
    ("衣养万物而不为主", "养育万物却不做主人，没有控制欲"),
    ("常无欲，可名于小", "永远没有私欲，看似渺小却最伟大"),
    ("以其终不自为大，故能成其大", "正因为始终不自以为大，所以能成就其伟大"),
]
for k, v in items:
    p = tf.add_paragraph()
    p.text = f"• {k}：{v}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_before = Pt(4)

# ==================== Slide 3: 第三十四章 生活智慧 ====================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, LIGHT_BG)
add_title_bar(s3, "第三十四章 · 生活智慧")
add_footer(s3, "3/10")

# 核心智慧框
zh = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.9))
zh.fill.solid()
zh.fill.fore_color.rgb = PRIMARY

zht = s3.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(8.8), Inches(0.6))
tf = zht.text_frame
p = tf.paragraphs[0]
p.text = "🌟 核心智慧：无为之德 — 成就万物而不居功"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 四条建议
tips = [
    ("成就他人，不居功", "真正的领导力是成就团队成员，让他们感受到是自己的功劳，而非你在控制。道的特点是'衣养万物而不为主'——真正的强者成就万物但不占为己有。"),
    ("放下控制欲", "越想控制，越会失去。管理者若总想牢牢把控权力，反而会让团队失去活力。道'常无欲'，所以能成其大。"),
    ("谦卑处下", "自以为大者，反而难成大事。那些功成名就却谦逊的人，正是'终不自为大，故能成其大'的践行者。"),
    ("顺势而为", "道在万物之中，不强为、不硬推。做事顺应规律，如同水流自然，无需刻意用力。"),
]

y = 2.4
for title, desc in tips:
    box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(9.2), Inches(1.05))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(200, 200, 200)
    
    tt = s3.shapes.add_textbox(Inches(0.6), Inches(y + 0.08), Inches(8.8), Inches(0.4))
    tf = tt.text_frame
    p = tf.paragraphs[0]
    p.text = f"• {title}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    td = s3.shapes.add_textbox(Inches(0.6), Inches(y + 0.48), Inches(8.8), Inches(0.5))
    tf = td.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(80, 80, 80)
    
    y += 1.15

# ==================== Slide 4: 第三十五章 原文+解读 ====================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, LIGHT_BG)
add_title_bar(s4, "第三十五章 · 执大象")
add_footer(s4, "4/10")

yb4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(9.2), Inches(1.4))
yb4.fill.solid()
yb4.fill.fore_color.rgb = RGBColor(240, 245, 250)
yb4.line.color.rgb = PRIMARY

yt4 = s4.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(1.2))
tf = yt4.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "【原文】"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = ACCENT
p2 = tf.add_paragraph()
p2.text = "执大象，天下往。往而不害，安平泰。乐与饵，过客止。道之出口，淡乎其无味，视之不足见，听之不足闻，用之不足既。"
p2.font.size = Pt(14)
p2.font.color.rgb = RGBColor(50, 50, 50)
p2.space_before = Pt(6)

jd4 = s4.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(3.8))
tf = jd4.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "【白话解读】"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY

items4 = [
    ("执大象，天下往", "掌握了大道，天下人都来归附"),
    ("往而不害，安平泰", "归附而不相互伤害，大家都安宁平和"),
    ("乐与饵，过客止", "美妙的音乐和美食，会让过客停下脚步"),
    ("道之出口，淡乎其无味", "但道说起来，却淡得没有味道"),
    ("视之不足见，听之不足闻", "看它看不见，听它听不到"),
    ("用之不足既", "用它却用不完"),
]
for k, v in items4:
    p = tf.add_paragraph()
    p.text = f"• {k}：{v}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_before = Pt(4)

# ==================== Slide 5: 第三十五章 典故+启示 ====================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5, LIGHT_BG)
add_title_bar(s5, "第三十五章 · 典故与启示")
add_footer(s5, "5/10")

# 典故
dg_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(4.4), Inches(2.8))
dg_box.fill.solid()
dg_box.fill.fore_color.rgb = RGBColor(245, 240, 250)
dg_box.line.color.rgb = RGBColor(150, 100, 180)

dg_t = s5.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(4.0), Inches(2.6))
tf = dg_t.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📜 典故"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 50, 150)

p = tf.add_paragraph()
p.text = "\n【孔子的感叹】"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(80, 80, 80)
p = tf.add_paragraph()
p.text = "\n孔子适周，问礼于老子。归来后感叹："
p.font.size = Pt(11)
p.font.color.rgb = RGBColor(100, 100, 100)
p = tf.add_paragraph()
p.text = "\n'吾所见老子，其犹龙乎？'"
p.font.size = Pt(12)
p.font.italic = True
p.font.color.rgb = PRIMARY
p = tf.add_paragraph()
p.text = "\n道之出口淡而无味，却用之不竭——真正的大道，从不张扬却影响深远。"
p.font.size = Pt(11)
p.font.color.rgb = RGBColor(100, 100, 100)

# 启示
qs_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.4), Inches(2.8))
qs_box.fill.solid()
qs_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
qs_box.line.color.rgb = PRIMARY

qs_t = s5.shapes.add_textbox(Inches(5.4), Inches(1.4), Inches(4.0), Inches(2.6))
tf = qs_t.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🌍 管理启示"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY

items5 = [
    "执大象：抓住本质，不被表象迷惑",
    "安平泰：利他才能安人，害人者必自害",
    "淡乎无味：真正的管理是润物细无声",
]
for item in items5:
    p = tf.add_paragraph()
    p.text = f"\n• {item}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(60, 60, 60)

# 底部总结
btm = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.3), Inches(9.2), Inches(2.4))
btm.fill.solid()
btm.fill.fore_color.rgb = WHITE
btm.line.color.rgb = RGBColor(180, 180, 180)

btm_t = s5.shapes.add_textbox(Inches(0.6), Inches(4.4), Inches(8.8), Inches(2.2))
tf = btm_t.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "💡 现代启示"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY

tips5 = [
    "企业价值观：短期吸引靠产品，长期留住靠文化（'乐与饵，过客止'）",
    "领导力真谛：真正的影响力不在于喊口号，而在于润物无声（'道之出口，淡乎其无味'）",
    "品牌建设：追求的不是知名度，而是让人'用之不足既'的价值",
]
for tip in tips5:
    p = tf.add_paragraph()
    p.text = f"\n• {tip}"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(70, 70, 70)

# ==================== Slide 6: 第三十六章 原文+解读 ====================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6, LIGHT_BG)
add_title_bar(s6, "第三十六章 · 将欲取之")
add_footer(s6, "6/10")

yb6 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(9.2), Inches(1.8))
yb6.fill.solid()
yb6.fill.fore_color.rgb = RGBColor(240, 245, 250)
yb6.line.color.rgb = PRIMARY

yt6 = s6.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(1.6))
tf = yt6.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "【原文】"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = ACCENT
p2 = tf.add_paragraph()
p2.text = "将欲取天下而为之，吾见其不得已。天下神器，不可为也。为者败之，执者失之。是以圣人无为，故无败；无执，故无失。物，或行或随，或呴或吹，或强或羸，或挫或隳。是以圣人去甚，去奢，去泰。"
p2.font.size = Pt(13)
p2.font.color.rgb = RGBColor(50, 50, 50)
p2.space_before = Pt(6)

jd6 = s6.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(3.5))
tf = jd6.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "【白话解读】"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY

items6 = [
    ("将欲取天下而为之，吾见其不得已", "想靠强力来治理天下，我看他达不到目的"),
    ("天下神器，不可为也", "天下是神圣的器物，不能靠强行作为"),
    ("为者败之，执者失之", "强行作为必然失败，紧紧把控必然失去"),
    ("圣人无为，故无败；无执，故无失", "圣人不妄为，所以不失败；不强执，所以不失去"),
    ("去甚，去奢，去泰", "去掉极端、奢侈、过度"),
]
for k, v in items6:
    p = tf.add_paragraph()
    p.text = f"• {k}：{v}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_before = Pt(4)

# ==================== Slide 7: 第三十六章 典故+启示 ====================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7, LIGHT_BG)
add_title_bar(s7, "第三十六章 · 典故与生活智慧")
add_footer(s7, "7/10")

# 典故
gd7 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(4.4), Inches(2.6))
gd7.fill.solid()
gd7.fill.fore_color.rgb = RGBColor(245, 248, 255)
gd7.line.color.rgb = PRIMARY

gd7_t = s7.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(4.0), Inches(2.4))
tf = gd7_t.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📜 历史典故"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY

p = tf.add_paragraph()
p.text = "\n【秦朝灭亡的教训】"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(80, 80, 80)
p = tf.add_paragraph()
p.text = "\n秦始皇统一六国后，实行严刑峻法、焚书坑儒，拼命想'执'住天下。结果二世而亡，正是'为者败之，执者失之'的典型例证。"
p.font.size = Pt(11)
p.font.color.rgb = RGBColor(100, 100, 100)
p = tf.add_paragraph()
p.text = "\n【汉初的黄老之治】"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(80, 80, 80)
p = tf.add_paragraph()
p.text = "\n汉高祖刘邦之后，文景二帝采用'无为而治'，轻徭薄赋，与民休息。成就了中国第一个盛世——'文景之治'。"
p.font.size = Pt(11)
p.font.color.rgb = RGBColor(100, 100, 100)

# 智慧
zh7 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.4), Inches(2.6))
zh7.fill.solid()
zh7.fill.fore_color.rgb = RGBColor(255, 248, 240)
zh7.line.color.rgb = ACCENT

zh7_t = s7.shapes.add_textbox(Inches(5.4), Inches(1.4), Inches(4.0), Inches(2.4))
tf = zh7_t.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🌟 生活智慧"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT

tips7 = [
    "不去甚：不做过头事，说话留余地",
    "不去奢：生活简单，精神更富足",
    "不去泰：得意时不张狂，失意时不颓废",
]
for tip in tips7:
    p = tf.add_paragraph()
    p.text = f"\n• {tip}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(80, 80, 80)

# 底部四象限
quads = [
    ("🧘 心态", "去除极端心态\n保持平和中正", Inches(0.4), Inches(4.1), RGBColor(230, 245, 230)),
    ("💼 做事", "去除强取豪夺\n顺势而为", Inches(2.55), Inches(4.1), RGBColor(230, 238, 250)),
    ("❤️ 交往", "去除占有欲\n给予而非索取", Inches(4.7), Inches(4.1), RGBColor(250, 240, 245)),
    ("🎯 成长", "去除急功近利\n脚踏实地", Inches(6.85), Inches(4.1), RGBColor(255, 248, 230)),
]

for title, desc, lx, ty, color in quads:
    box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, ty, Inches(2.0), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = RGBColor(200, 200, 200)
    
    tt = s7.shapes.add_textbox(lx + Inches(0.1), ty + Inches(0.1), Inches(1.8), Inches(0.5))
    tf = tt.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    dt = s7.shapes.add_textbox(lx + Inches(0.1), ty + Inches(0.6), Inches(1.8), Inches(1.8))
    tf = dt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

# ==================== Slide 8: 西方思想对照 ====================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8, LIGHT_BG)
add_title_bar(s8, "西方思想对照 · 四位思想家名言")
add_footer(s8, "8/10")

philosophers = [
    ("🇫🇷 伏尔泰", "自然法", "1750年", "'顺应自然，方能获得自由。'——道法自然，欧洲启蒙思想的核心。"),
    ("🇩🇪 黑格尔", "辩证法", "1807年", "'矛盾推动发展，对立统一。'——如同道家有无相生之理。"),
    ("🇺🇸 梭罗", "瓦尔登湖", "1849年", "'简化生活，回归本真。'——与道家返璞归真不谋而合。"),
    ("🇬🇧 斯密", "看不见的手", "1776年", "'市场自发秩序，比刻意调控更有效。'——无为而治的经济学表达。"),
]

y8 = 1.35
for name, concept, year, quote in philosophers:
    box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y8), Inches(9.2), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(200, 200, 200)
    
    nt = s8.shapes.add_textbox(Inches(0.6), Inches(y8 + 0.1), Inches(2.5), Inches(0.4))
    tf = nt.text_frame
    p = tf.paragraphs[0]
    p.text = f"{name} · {concept}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    yt = s8.shapes.add_textbox(Inches(3.2), Inches(y8 + 0.12), Inches(1.5), Inches(0.35))
    tf = yt.text_frame
    p = tf.paragraphs[0]
    p.text = year
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    
    qt = s8.shapes.add_textbox(Inches(0.6), Inches(y8 + 0.55), Inches(8.8), Inches(0.65))
    tf = qt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(80, 80, 80)
    
    y8 += 1.42

# ==================== Slide 9: 生活中的智慧 ====================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9, LIGHT_BG)
add_title_bar(s9, "生活中的智慧 · 四场景应用")
add_footer(s9, "9/10")

scenes = [
    ("🧘 每日自省", PRIMARY, "每天留10分钟独处，问自己：今天有没有'去甚、去奢、去泰'？顺境时检视是否傲慢，逆境时检视是否怨天尤人。"),
    ("💼 团队管理", RGBColor(0, 100, 80), "少用控制，多用成就。好的管理者是'衣养万物而不为主'——提供平台，让团队成员自己成长。"),
    ("❤️ 家庭关系", RGBColor(150, 50, 50), "家人之间，去除占有欲和控制欲。不是'执'住对方，而是支持对方成为更好的自己。"),
    ("🎯 个人成长", RGBColor(100, 100, 0), "不急功近利，不追求'立竿见影'。真正的成长是'大道泛兮'——慢慢渗透，功到自然成。"),
]

y9 = 1.3
for title, color, desc in scenes:
    box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y9), Inches(9.2), Inches(1.35))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    
    # 左侧色块
    lb = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(y9), Inches(0.12), Inches(1.35))
    lb.fill.solid()
    lb.fill.fore_color.rgb = color
    lb.line.fill.background()
    
    tt = s9.shapes.add_textbox(Inches(0.7), Inches(y9 + 0.1), Inches(8.7), Inches(0.4))
    tf = tt.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    
    dt = s9.shapes.add_textbox(Inches(0.7), Inches(y9 + 0.55), Inches(8.7), Inches(0.7))
    tf = dt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(80, 80, 80)
    
    y9 += 1.45

# ==================== Slide 10: 总结 ====================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10, WHITE)
add_title_bar(s10, "精华提炼 · 第十一课总结")
add_footer(s10, "10/10")

# 三章核心
chapters = [
    ("第三十四章 · 大道泛兮", "无为之德", "成就万物而不居功，'以其终不自为大，故能成其大'", PRIMARY),
    ("第三十五章 · 执大象", "道法自然", "掌握本质，润物无声，'道之出口，淡乎其无味，用之不足既'", RGBColor(0, 100, 80)),
    ("第三十六章 · 将欲取之", "去甚去泰", "不为不执，去除极端、奢侈、过度，方能'无败无失'", RGBColor(150, 50, 50)),
]

y10 = 1.3
for title, theme, desc, color in chapters:
    box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y10), Inches(9.2), Inches(1.35))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(248, 248, 255)
    box.line.color.rgb = color
    
    # 左侧色块
    lb = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(y10), Inches(0.15), Inches(1.35))
    lb.fill.solid()
    lb.fill.fore_color.rgb = color
    lb.line.fill.background()
    
    tt = s10.shapes.add_textbox(Inches(0.7), Inches(y10 + 0.08), Inches(8.7), Inches(0.4))
    tf = tt.text_frame
    p = tf.paragraphs[0]
    p.text = f"{title}｜{theme}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = color
    
    dt = s10.shapes.add_textbox(Inches(0.7), Inches(y10 + 0.5), Inches(8.7), Inches(0.75))
    tf = dt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(80, 80, 80)
    
    y10 += 1.45

# 一句话总结
summary_box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.65), Inches(9.2), Inches(1.2))
summary_box.fill.solid()
summary_box.fill.fore_color.rgb = PRIMARY

st10 = s10.shapes.add_textbox(Inches(0.6), Inches(5.75), Inches(8.8), Inches(1.0))
tf = st10.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🌟 一句话总结"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(200, 220, 255)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "大道泛兮，道法自然，无为之德成就万物；将欲取之，去甚去泰，不争之争方能无败。"
p2.font.size = Pt(14)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER

# 保存
output = "/Users/mac/.openclaw/workspace/courses/道德经_第11天.pptx"
prs.save(output)
print(f"已保存: {output}")
