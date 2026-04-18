#!/usr/bin/env python3
"""
OpenClaw精通指南 - 60页PPT生成器
基于Mac Mini的OpenClaw全面指南
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os
from datetime import datetime

# ========== 配色方案 ==========
STYLES = {
    "封面": {"primary": RGBColor(30, 58, 95), "secondary": RGBColor(52, 152, 219), "bg": RGBColor(236, 240, 241), "text": RGBColor(255, 255, 255)},
    "章节页": {"primary": RGBColor(30, 58, 95), "secondary": RGBColor(52, 152, 219), "bg": RGBColor(44, 62, 80), "text": RGBColor(255, 255, 255)},
    "内容页": {"primary": RGBColor(30, 58, 95), "secondary": RGBColor(52, 152, 219), "bg": RGBColor(255, 255, 255), "text": RGBColor(44, 62, 80)},
    "代码页": {"primary": RGBColor(44, 62, 80), "secondary": RGBColor(46, 204, 113), "bg": RGBColor(44, 62, 80), "text": RGBColor(46, 204, 113)},
}

# ========== PPT内容结构（60页）==========
SLIDES = [
    # ========== 第一部分：封面与目录 ==========
    {"type": "封面", "title": "OpenClaw 精通指南", "subtitle": "基于Mac Mini的AI助手网关完整教程", "footer": "2026年版"},
    {"type": "目录", "title": "目录", "items": ["第一章：OpenClaw基本介绍", "第二章：安装指南", "第三章：主要配置方法", "第四章：主要Skill介绍", "第五章：常见问题与处理", "第六章：飞书连接教程", "第七章：常用指令集"]},
    
    # ========== 第一章：基本介绍 ==========
    {"type": "章节", "title": "第一章", "subtitle": "OpenClaw基本介绍", "page": 1},
    {"type": "内容", "title": "什么是OpenClaw？", "bullets": ["OpenClaw是一个多通道AI网关", "连接聊天应用到AI编程助手", "支持WhatsApp、Telegram、Discord、iMessage等多种渠道", "运行在您自己的硬件上，数据完全私密"], "page": 1},
    {"type": "内容", "title": "核心特点", "bullets": ["🖥️ 自托管：运行在您的Mac Mini上", "📱 多通道：一个网关同时服务多个聊天应用", "🤖 Agent原生：内置编程Agent、记忆、会话管理", "🔓 开源：MIT许可证，社区驱动"], "page": 2},
    {"type": "内容", "title": "系统架构", "bullets": ["Chat Apps + Plugins → Gateway → Pi Agent/CLI/Control UI", "Gateway是会话、路由和频道连接的单一可信源", "支持浏览器Dashboard、CLI命令行、飞书等接入"], "page": 3},
    {"type": "内容", "title": "系统要求", "bullets": ["硬件：Mac Mini（推荐Apple Silicon M1/M2/M3）", "系统：macOS 12 Monterey 或更高版本", "软件：Node.js 22 LTS 或 Node.js 24（推荐）", "网络：稳定的互联网连接"], "page": 4},
    {"type": "内容", "title": "Dashboard界面", "bullets": ["本地默认地址：http://127.0.0.1:18789/", "支持浏览器Dashboard进行配置、会话管理", "远程访问支持Tailscale组网方案"], "page": 5},
    
    # ========== 第二章：安装指南 ==========
    {"type": "章节", "title": "第二章", "subtitle": "安装指南", "page": 1},
    {"type": "内容", "title": "安装前准备", "bullets": ["确认Node.js版本：node --version", "建议使用Homebrew安装Node 22 LTS", "准备好API Key（MiniMax、阿里云等）", "确认Mac Mini有足够的磁盘空间（建议50GB+）"], "page": 1},
    {"type": "内容", "title": "安装步骤", "bullets": ["方式一：npm全局安装：npm install -g openclaw@latest", "方式二：Homebrew安装：brew install openclaw", "安装完成后验证：openclaw --version"], "page": 2},
    {"type": "代码", "title": "安装命令", "code": "# Node.js检查\nnode --version\n\n# npm全局安装OpenClaw\nnpm install -g openclaw@latest\n\n# 验证安装\nopenclaw --version\n\n# 查看帮助\nopenclaw --help"},
    {"type": "内容", "title": "初始化配置", "bullets": ["运行初始化向导：openclaw setup", "配置工作区目录（默认~/.openclaw/workspace）", "选择安装模式：本地或远程", "配置API Key（MiniMax推荐）"], "page": 3},
    {"type": "代码", "title": "初始化命令", "code": "# 启动初始化向导\nopenclaw setup\n\n# 非交互式快速初始化\nopenclaw setup --non-interactive\n\n# 指定工作区\nopenclaw setup --workspace ~/.openclaw/workspace\n\n# 本地模式\nopenclaw setup --mode local"},
    {"type": "内容", "title": "启动网关服务", "bullets": ["方式一：前台运行 openclaw gateway run", "方式二：安装为后台服务 openclaw gateway install", "方式三：使用launchd（macOS推荐）", "启动后访问 http://127.0.0.1:18789/"], "page": 4},
    {"type": "代码", "title": "服务管理命令", "code": "# 前台运行\nopenclaw gateway run\n\n# 安装为系统服务\nopenclaw gateway install\n\n# 启动服务\nopenclaw gateway start\n\n# 查看服务状态\nopenclaw gateway status\n\n# 重启服务\nopenclaw gateway restart"},
    
    # ========== 第三章：主要配置方法 ==========
    {"type": "章节", "title": "第三章", "subtitle": "主要配置方法", "page": 1},
    {"type": "内容", "title": "配置文件位置与结构", "bullets": ["配置文件：~/.openclaw/openclaw.json", "工作区目录：~/.openclaw/workspace/", "认证信息：~/.openclaw/auth-profiles.json", "日志文件：~/.openclaw/logs/"], "page": 1},
    {"type": "内容", "title": "API Key配置", "bullets": ["MiniMax（推荐）：sk-cp-开头", "阿里云DashScope：sk-开头", "Tavily搜索API：tvly-开头", "配置路径：openclaw config set"], "page": 2},
    {"type": "代码", "title": "配置API Key", "code": "# 配置MiniMax\nopenclaw config set auth.providers.minimax.key YOUR_MINIMAX_KEY\n\n# 配置阿里云DashScope\nopenclaw config set auth.providers.modelstudio.key YOUR_DASHSCOPE_KEY\n\n# 配置Tavily\nopenclaw config set secrets.provider.tavily tvly-YOUR-KEY\n\n# 验证配置\nopenclaw models status"},
    {"type": "内容", "title": "飞书频道配置", "bullets": ["配置飞书机器人App ID和App Secret", "配置Webhook地址用于接收消息", "设置消息路由规则", "配置消息加密（如需要）"], "page": 3},
    {"type": "代码", "title": "飞书配置示例", "code": "# 配置飞书渠道\nopenclaw channels add --channel feishu\n\n# 配置飞书App ID\nopenclaw config set channels.feishu.appId CLI_XXXX\n\n# 配置飞书App Secret\nopenclaw config set channels.feishu.appSecret YOUR_SECRET\n\n# 查看渠道状态\nopenclaw channels status"},
    {"type": "内容", "title": "Agent配置", "bullets": ["默认Agent：Pi Agent", "可配置多个隔离Agent", "设置Agent的系统提示词（SOUL.md）", "配置Agent的记忆和上下文策略"], "page": 4},
    {"type": "代码", "title": "Agent管理命令", "code": "# 列出所有Agent\nopenclaw agents list\n\n# 添加新Agent\nopenclaw agents add my-agent --workspace ~/.openclaw/agents/my-agent\n\n# 绑定渠道到Agent\nopenclaw agents bind --agent my-agent --bind feishu:default\n\n# 设置默认Agent\nopenclaw config set agents.defaults.model minimax/MiniMax-M2.7"},
    {"type": "内容", "title": "Cron定时任务配置", "bullets": ["支持cron表达式定时执行", "可执行消息发送、Agent调用等", "常用场景：新闻早报、学习提醒、数据备份", "配置路径：~/.openclaw/workspace/"], "page": 5},
    {"type": "代码", "title": "Cron任务管理", "code": "# 创建Cron任务\nopenclaw cron add --name \"新闻早报\" --cron \"0 7 * * *\" \\\n  --message \"生成今日新闻早报\" --session isolated --announce\n\n# 列出所有Cron任务\nopenclaw cron list\n\n# 查看任务状态\nopenclaw cron runs --id TASK_ID\n\n# 删除任务\nopenclaw cron rm TASK_ID"},
    
    # ========== 第四章：主要Skill介绍 ==========
    {"type": "章节", "title": "第四章", "subtitle": "主要Skill介绍", "page": 1},
    {"type": "内容", "title": "Skill系统概述", "bullets": ["Skill是OpenClaw的功能扩展模块", "安装路径：~/.openclaw/workspace/skills/", "可通过ClawHub安装：npx clawhub@latest install", "已安装54+个Skills"], "page": 1},
    {"type": "内容", "title": "PPT智能生成（ppt-smart-generator）", "bullets": ["支持8种PPT类型：工作汇报、数据分析、融资路演等", "智能生成专业大纲", "数据图表自动生成与整合", "网络资料深度搜集"], "page": 2},
    {"type": "内容", "title": "Excel自动化（automate-excel）", "bullets": ["数据清洗与完整性校验", "多维度透视分析", "可视化图表生成", "支持HR、财务、资产、项目等数据类型"], "page": 3},
    {"type": "内容", "title": "PPT生成器（pptx-generator）", "bullets": ["专业PPT生成器", "支持多种风格：商务蓝、科技深、学术白、创意紫", "可编辑PPTX格式", "图表支持：柱状图、折线图、饼图"], "page": 4},
    {"type": "内容", "title": "记忆备份（memory-backup）", "bullets": ["自动备份对话历史", "每日定时备份到GitHub", "重要对话即刻记录到MEMORY.md", "支持导出和恢复"], "page": 5},
    {"type": "内容", "title": "飞书日历（lark-calendar）", "bullets": ["读取飞书日历事件", "创建飞书日历日程", "设置日程提醒", "每日晚8点主动同步"], "page": 6},
    {"type": "内容", "title": "学习类Skills", "bullets": ["learning-cards：学习卡片系统", "flash-forge：AI闪卡生成器", "spaced-repetition-teaching：间隔重复记忆教学", "ai-meeting-notes：会议纪要（含行动项）"], "page": 7},
    {"type": "内容", "title": "调研类Skills", "bullets": ["research-assistant：深度调研助手", "tavily-search：网络搜索", "brave-web-search：Brave搜索备份", "seo-competitor-analysis：SEO竞品分析"], "page": 8},
    {"type": "代码", "title": "Skill管理命令", "code": "# 列出已安装Skills\nopenclaw skills list\n\n# 搜索ClawHub上的Skills\nopenclaw skills search ppt\n\n# 安装Skill\nnpx clawhub@latest install ppt-smart-generator\n\n# 查看Skill详情\nopenclaw skills info ppt-smart-generator\n\n# 检查Skill状态\nopenclaw skills check"},
    
    # ========== 第五章：常见问题与处理 ==========
    {"type": "章节", "title": "第五章", "subtitle": "常见问题与处理方法", "page": 1},
    {"type": "内容", "title": "安装问题", "bullets": ["问题：npm安装失败", "解决：检查Node.js版本，更新到22 LTS或24", "问题：权限错误（EACCES）", "解决：使用sudo或配置npm prefix"], "page": 1},
    {"type": "内容", "title": "网关连接问题", "bullets": ["问题：无法访问127.0.0.1:18789", "解决：检查网关是否运行，尝试重启服务", "问题：Gateway超时", "解决：检查端口占用，查看日志 openclaw logs"], "page": 2},
    {"type": "内容", "title": "API Key问题", "bullets": ["问题：All models failed", "解决：检查API Key是否有效，更新过期Key", "问题：认证失败（HTTP 401）", "解决：重新配置API Key，使用 openclaw models auth"], "page": 3},
    {"type": "代码", "title": "诊断命令", "code": "# 运行诊断工具\nopenclaw doctor\n\n# 检查网关健康\nopenclaw gateway health\n\n# 查看详细状态\nopenclaw status --deep\n\n# 检查模型状态\nopenclaw models status --probe\n\n# 查看最近日志\nopenclaw logs --limit 100"},
    {"type": "内容", "title": "飞书连接问题", "bullets": ["问题：飞书机器人无响应", "解决：检查App ID和App Secret配置", "问题：消息发送失败", "解决：检查飞书机器人的权限设置"], "page": 4},
    {"type": "content", "title": "Cron任务问题", "bullets": ["问题：Cron任务失败", "解决：检查任务超时设置，查看任务日志", "问题：任务未执行", "解决：检查cron表达式是否正确，确认网关运行中"], "page": 5},
    {"type": "content", "title": "Skill问题", "bullets": ["问题：Skill无法加载", "解决：检查依赖是否安装，查看 openclaw skills check", "问题：Skill执行报错", "解决：检查Skill文档，查看脚本路径"], "page": 6},
    
    # ========== 第六章：飞书连接教程 ==========
    {"type": "章节", "title": "第六章", "subtitle": "飞书连接教程", "page": 1},
    {"type": "内容", "title": "飞书机器人准备", "bullets": ["在飞书开放平台创建企业自建应用", "配置应用名称和图标", "获取App ID（CLI开头）和App Secret", "添加机器人能力"], "page": 1},
    {"type": "内容", "title": "配置飞书应用权限", "bullets": ["添加消息权限：im:message", "添加通讯录权限：contact:user.employee_id:readonly", "配置应用可用范围", "发布应用版本"], "page": 2},
    {"type": "内容", "title": "配置Webhook", "bullets": ["在飞书应用配置中启用Webhook", "设置请求地址指向OpenClaw网关", "配置消息加密（可选，推荐使用", "测试Webhook连接"], "page": 3},
    {"type": "代码", "title": "飞书配置步骤", "code": "# 1. 添加飞书渠道\nopenclaw channels add --channel feishu\n\n# 2. 配置App ID\nopenclaw config set channels.feishu.appId CLI_XXXXXXXXXXXXXX\n\n# 3. 配置App Secret\nopenclaw config set channels.feishu.appSecret YOUR_APP_SECRET\n\n# 4. 重启网关\nopenclaw gateway restart\n\n# 5. 验证连接\nopenclaw channels status"},
    {"type": "内容", "title": "配置三先生信息", "bullets": ["获取三先生的飞书Open ID", "配置路由规则将消息路由到正确的Agent", "设置消息接收和发送的账号绑定", "测试双向通讯"], "page": 4},
    {"type": "内容", "title": "飞书消息格式", "bullets": ["支持文本消息", "支持图片、文件、语音", "支持@提及和回复", "支持卡片消息（Interactive）"], "page": 5},
    {"type": "代码", "title": "发送飞书消息", "code": "# 发送文本消息\nopenclaw message send --channel feishu \\\n  --target user:ou_XXXXXXXXX \\\n  --message \"测试消息\"\n\n# 发送文件\nopenclaw message send --channel feishu \\\n  --target user:ou_XXXXXXXXX \\\n  --media /path/to/file.pptx"},
    
    # ========== 第七章：常用指令集 ==========
    {"type": "章节", "title": "第七章", "subtitle": "OpenClaw常用指令集", "page": 1},
    {"type": "内容", "title": "网关管理命令", "bullets": ["openclaw gateway run：前台运行", "openclaw gateway start：启动服务", "openclaw gateway stop：停止服务", "openclaw gateway restart：重启服务"], "page": 1},
    {"type": "代码", "title": "网关管理", "code": "# 运行网关\nopenclaw gateway run\n\n# 守护进程模式\nopenclaw gateway install\nopenclaw gateway start\n\n# 检查状态\nopenclaw gateway status\n\n# 查看日志\nopenclaw logs --follow"},
    {"type": "内容", "title": "配置管理命令", "bullets": ["openclaw config get <path>：获取配置值", "openclaw config set <path> <value>：设置配置", "openclaw config file：显示配置文件路径", "openclaw config validate：验证配置"], "page": 2},
    {"type": "代码", "title": "配置管理", "code": "# 查看配置\nopenclaw config get agents.defaults.model\n\n# 设置配置\nopenclaw config set agents.defaults.model minimax/MiniMax-M2.7\n\n# 添加渠道\nopenclaw channels add --channel feishu\n\n# 验证配置\nopenclaw config validate"},
    {"type": "内容", "title": "渠道管理命令", "bullets": ["openclaw channels list：列出渠道", "openclaw channels status：查看状态", "openclaw channels add：添加渠道", "openclaw channels logs：查看日志"], "page": 3},
    {"type": "代码", "title": "渠道管理", "code": "# 列出所有渠道\nopenclaw channels list\n\n# 检查渠道健康\nopenclaw channels status --probe\n\n# 添加新渠道\nopenclaw channels add --channel telegram --token YOUR_BOT_TOKEN\n\n# 查看渠道日志\nopenclaw channels logs --channel feishu"},
    {"type": "内容", "title": "Cron定时任务", "bullets": ["openclaw cron list：列出任务", "openclaw cron add：创建任务", "openclaw cron rm <id>：删除任务", "openclaw cron runs --id <id>：查看执行记录"], "page": 4},
    {"type": "代码", "title": "Cron命令", "code": "# 创建每日早报任务\nopenclaw cron add \\\n  --name \"每日新闻早报\" \\\n  --cron \"0 7 * * *\" \\\n  --message \"生成今日要闻早报\" \\\n  --session isolated \\\n  --announce \\\n  --channel feishu\n\n# 查看所有任务\nopenclaw cron list\n\n# 手动触发任务\nopenclaw cron run TASK_ID"},
    {"type": "内容", "title": "会话管理命令", "bullets": ["openclaw sessions list：列出会话", "openclaw sessions cleanup：清理会话", "openclaw status：查看状态", "openclaw health：健康检查"], "page": 5},
    {"type": "代码", "title": "会话管理", "code": "# 查看活跃会话\nopenclaw sessions list\n\n# 深度状态检查\nopenclaw status --deep\n\n# 健康检查\nopenclaw health\n\n# 清理过期会话\nopenclaw sessions cleanup"},
    {"type": "内容", "title": "Skill命令", "bullets": ["openclaw skills list：列出Skills", "openclaw skills check：检查状态", "openclaw skills info <name>：查看详情", "npx clawhub@latest install <slug>：安装Skill"], "page": 6},
    {"type": "代码", "title": "Skill命令", "code": "# 列出已安装Skills\nopenclaw skills list\n\n# 检查Skills健康状态\nopenclaw skills check\n\n# 搜索Skills\nopenclaw skills search ppt\n\n# 安装Skill\nnpx clawhub@latest install ppt-smart-generator"},
    {"type": "内容", "title": "模型管理命令", "bullets": ["openclaw models list：列出可用模型", "openclaw models status：查看模型状态", "openclaw models set <model>：设置默认模型", "openclaw models auth：配置认证信息"], "page": 7},
    {"type": "代码", "title": "模型管理", "code": "# 列出可用模型\nopenclaw models list\n\n# 检查模型认证状态\nopenclaw models status --probe\n\n# 设置默认模型\nopenclaw models set minimax/MiniMax-M2.7\n\n# 配置API Key\nopenclaw models auth add --provider minimax"},
    {"type": "内容", "title": "诊断与修复", "bullets": ["openclaw doctor：运行诊断", "openclaw doctor --fix：自动修复", "openclaw backup create：创建备份", "openclaw reset：重置配置"], "page": 8},
    {"type": "代码", "title": "诊断修复", "code": "# 运行完整诊断\nopenclaw doctor\n\n# 自动修复问题\nopenclaw doctor --fix\n\n# 创建备份\nopenclaw backup create\n\n# 重置（谨慎使用）\nopenclaw reset --scope config+creds+sessions --yes"},
    
    # ========== 总结页 ==========
    {"type": "内容", "title": "快速参考卡", "bullets": ["安装：npm install -g openclaw@latest", "启动：openclaw gateway run", "配置：openclaw config set", "渠道：openclaw channels add", "任务：openclaw cron add", "Skills：npx clawhub@latest install"], "page": 1},
    {"type": "内容", "title": "资源链接", "bullets": ["官方文档：https://docs.openclaw.ai", "ClawHub：https://clawhub.ai", "GitHub：https://github.com/openclaw/openclaw", "社区Discord：https://discord.gg/clawd", "本地Dashboard：http://127.0.0.1:18789/"], "page": 2},
    {"type": "封面", "title": "感谢观看", "subtitle": "OpenClaw 精通指南", "footer": "小燕子出品 · 2026"},
]

# ========== 生成PPT ==========
def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))

def add_title_slide(prs, slide_data):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = STYLES["封面"]["primary"]
    background.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = STYLES["封面"]["text"]
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if "subtitle" in slide_data:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data["subtitle"]
        p.font.size = Pt(28)
        p.font.color.rgb = STYLES["封面"]["secondary"]
        p.alignment = PP_ALIGN.CENTER
    
    # 页脚
    if "footer" in slide_data:
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data["footer"]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.CENTER

def add_chapter_slide(prs, slide_data):
    """添加章节分隔幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = STYLES["章节页"]["bg"]
    background.line.fill.background()
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = STYLES["章节页"]["text"]
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if "subtitle" in slide_data:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data["subtitle"]
        p.font.size = Pt(32)
        p.font.color.rgb = STYLES["章节页"]["secondary"]
        p.alignment = PP_ALIGN.CENTER
    
    # 页码
    if "page" in slide_data:
        page_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"- {slide_data['page']} -"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, slide_data):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部色条
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = STYLES["内容页"]["primary"]
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = STYLES["内容页"]["text"]
    
    # 内容
    if "bullets" in slide_data:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(slide_data["bullets"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = STYLES["内容页"]["text"]
            p.space_after = Pt(12)
    
    # 页码
    if "page" in slide_data:
        page_box = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.3))
        p = page_box.text_frame.paragraphs[0]
        p.text = str(slide_data["page"])
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(150, 150, 150)

def add_code_slide(prs, slide_data):
    """添加代码幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = STYLES["代码页"]["bg"]
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = STYLES["代码页"]["text"]
    
    # 代码框
    code_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(12.733), Inches(6))
    tf = code_box.text_frame
    tf.word_wrap = True
    
    lines = slide_data.get("code", "").strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Courier New"
        p.font.color.rgb = STYLES["代码页"]["text"]
        p.space_after = Pt(4)

def add_toc_slide(prs, slide_data):
    """添加目录幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(250, 250, 250)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = STYLES["内容页"]["primary"]
    
    # 目录项
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5))
    tf = content_box.text_frame
    
    items = slide_data.get("items", [])
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = STYLES["内容页"]["text"]
        p.space_after = Pt(16)

def generate_ppt(output_path):
    """生成完整PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in SLIDES:
        slide_type = slide_data.get("type", "内容")
        
        if slide_type == "封面":
            add_title_slide(prs, slide_data)
        elif slide_type == "章节":
            add_chapter_slide(prs, slide_data)
        elif slide_type == "内容":
            add_content_slide(prs, slide_data)
        elif slide_type == "代码":
            add_code_slide(prs, slide_data)
        elif slide_type == "目录":
            add_toc_slide(prs, slide_data)
    
    prs.save(output_path)
    print(f"✅ PPT已生成：{output_path}")
    print(f"共 {len(SLIDES)} 页")
    return output_path

if __name__ == "__main__":
    output_dir = os.path.expanduser("~/openclaw-data/ppt")
    os.makedirs(output_dir, exist_ok=True)
    output_path = f"{output_dir}/OpenClaw精通指南_{datetime.now().strftime('%Y%m%d')}.pptx"
    generate_ppt(output_path)
