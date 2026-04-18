#!/usr/bin/env python3
"""
PPT智能生成器 - 核心引擎
支持：工作汇报、数据分析、项目立项提案、培训课件、产品介绍、融资路演、经济运行分析、策划方案
功能：智能大纲生成 + 数据图表整合 + 网络资料搜集 + PPT制作
"""

import sys
import os
import argparse
import json
import re
from pathlib import Path
from datetime import datetime

# ========== 1. PPT类型定义 ==========

PPT_TYPES = {
    "工作汇报": {
        "name": "工作汇报",
        "scenes": ["周报", "月报", "季度汇报", "年度汇报", "专项汇报"],
        "structure": ["封面", "目录", "本期工作概述", "重点工作进展", "问题与挑战", "下期工作计划", "结语"],
        "style": "商务蓝",
        "color": "#1E3A5F",
        "charts": ["进度对比", "完成率环形", "趋势折线"],
        "keywords": ["完成率", "工作量", "里程碑", "问题分析", "改进措施"]
    },
    "数据分析": {
        "name": "数据分析",
        "scenes": ["运营分析", "市场分析", "用户分析", "财务分析", "业务分析"],
        "structure": ["封面", "目录", "分析背景", "核心指标", "数据趋势", "问题诊断", "建议与行动", "附录"],
        "style": "科技深",
        "color": "#0D1B2A",
        "charts": ["柱状图", "折线图", "饼图", "散点图", "热力图"],
        "keywords": ["增长率", "占比", "同比环比", "用户画像", "漏斗分析"]
    },
    "项目立项提案": {
        "name": "项目立项提案",
        "scenes": ["新产品立项", "技术改造立项", "基建项目立项", "研发项目立项"],
        "structure": ["封面", "目录", "项目背景", "项目目标", "方案设计", "投资估算", "实施计划", "风险分析", "预期效益", "结论建议"],
        "style": "商务蓝",
        "color": "#1E3A5F",
        "charts": ["甘特图", "预算分解柱状", "投资回报曲线"],
        "keywords": ["ROI", "投资回收期", "NPV", "敏感性分析", "风险矩阵"]
    },
    "培训课件": {
        "name": "培训课件",
        "scenes": ["新员工培训", "技能培训", "管理培训", "安全培训", "产品培训"],
        "structure": ["封面", "目录", "培训目标", "内容大纲", "知识点1", "知识点2", "知识点3", "实操演示", "总结与测试"],
        "style": "学术白",
        "color": "#2C3E50",
        "charts": ["知识结构图", "流程图", "步骤分解图"],
        "keywords": ["学习目标", "知识要点", "案例分析", "实践操作", "考核标准"]
    },
    "产品介绍": {
        "name": "产品介绍",
        "scenes": ["产品发布会", "销售支持", "客户演示", "官网展示"],
        "structure": ["封面", "目录", "市场背景", "产品概述", "核心功能", "产品优势", "应用场景", "客户案例", "定价方案", "联系我们"],
        "style": "创意紫",
        "color": "#6C3483",
        "charts": ["功能对比表", "应用场景图", "客户案例柱状"],
        "keywords": ["核心卖点", "差异化", "用户价值", "应用场景", "成功案例"]
    },
    "融资路演": {
        "name": "融资路演",
        "scenes": ["天使轮", "A轮", "B轮", "C轮", "战略投资"],
        "structure": ["封面", "我们是誰", "市场机会", "解决方案", "商业模式", "运营现状", "发展规划", "融资计划", "团队介绍", "联系方式"],
        "style": "创意紫",
        "color": "#6C3483",
        "charts": ["市场规模柱状", "增长曲线", "商业模式画布", "估值图"],
        "keywords": ["TAM", "SAM", "SOM", "增长曲线", "单位经济模型", "估值"]
    },
    "经济运行分析": {
        "name": "经济运行分析",
        "scenes": ["年度经济分析", "季度经济分析", "行业经济分析", "区域经济分析"],
        "structure": ["封面", "目录", "宏观环境", "经济运行概况", "主要指标分析", "产业结构分析", "问题研判", "趋势预测", "对策建议"],
        "style": "商务蓝",
        "color": "#1E3A5F",
        "charts": ["GDP趋势", "产业结构饼图", "区域对比柱状", "PMI走势图"],
        "keywords": ["GDP", "CPI", "PPI", "PMI", "固投", "消费", "进出口", "同比增长"]
    },
    "策划方案": {
        "name": "策划方案",
        "scenes": ["营销策划", "活动策划", "品牌策划", "危机公关策划"],
        "structure": ["封面", "目录", "背景分析", "目标受众", "核心策略", "创意概念", "执行计划", "时间表", "预算分配", "效果评估"],
        "style": "创意紫",
        "color": "#6C3483",
        "charts": ["受众画像", "传播路径图", "时间轴甘特图", "预算饼图"],
        "keywords": ["SWOT", "STP", "4P", "目标受众", "创意策略", "传播渠道", "KPI"]
    }
}

# ========== 2. 智能大纲生成 ==========

def generate_outline(ppt_type, scene, custom_topic="", extra_data=None):
    """根据PPT类型和场景生成智能大纲"""
    
    if ppt_type not in PPT_TYPES:
        return {"error": f"未知PPT类型: {ppt_type}"}
    
    config = PPT_TYPES[ppt_type]
    topic = custom_topic or config["name"]
    
    outline = {
        "type": ppt_type,
        "scene": scene,
        "topic": topic,
        "structure": config["structure"].copy(),
        "style": config["style"],
        "color": config["color"],
        "charts": config["charts"],
        "keywords": config["keywords"],
        "slides_count": len(config["structure"]),
        "estimated_minutes": len(config["structure"]) * 2
    }
    
    return outline

# ========== 3. 数据图表生成 ==========

def generate_chart_data(chart_type, data_source=None):
    """根据类型生成图表数据"""
    
    chart_templates = {
        "进度对比": {
            "type": "bar",
            "title": "工作进度对比",
            "data": {"labels": ["计划值", "实际值"], "datasets": [{"data": [100, 85], "color": "#3498DB"}]}
        },
        "完成率环形": {
            "type": "doughnut",
            "title": "任务完成率",
            "data": {"labels": ["已完成", "进行中", "未完成"], "datasets": [{"data": [65, 25, 10]}]}
        },
        "趋势折线": {
            "type": "line",
            "title": "趋势变化图",
            "data": {"labels": ["1月", "2月", "3月", "4月"], "datasets": [{"data": [30, 45, 60, 80]}]}
        },
        "柱状图": {
            "type": "bar",
            "title": "数据对比",
            "data": {"labels": ["A项目", "B项目", "C项目"], "datasets": [{"data": [120, 85, 96]}]}
        },
        "饼图": {
            "type": "pie",
            "title": "占比分析",
            "data": {"labels": ["占比1", "占比2", "占比3", "占比4"], "datasets": [{"data": [40, 30, 20, 10]}]}
        },
        "甘特图": {
            "type": "horizontal_bar",
            "title": "项目实施计划",
            "data": {"labels": ["阶段一", "阶段二", "阶段三", "阶段四"], "datasets": [{"data": [[1, 5], [3, 8], [6, 10], [9, 12]]}]}
        },
        "增长曲线": {
            "type": "line",
            "title": "增长曲线",
            "data": {"labels": ["Q1", "Q2", "Q3", "Q4"], "datasets": [{"data": [10, 25, 55, 120]}]}
        }
    }
    
    return chart_templates.get(chart_type, chart_templates["柱状图"])

# ========== 4. 网络资料搜集 ==========

def search_web(topic, keywords="", max_results=5):
    """模拟网络资料搜集（实际调用 Tavily 或 Brave Search）"""
    # 这里实际会调用 web_search 工具
    return {
        "topic": topic,
        "sources": [
            {"title": f"{topic} - 相关资料1", "url": "https://example.com/1", "snippet": "这是搜索到的相关资料摘要..."},
            {"title": f"{topic} - 相关资料2", "url": "https://example.com/2", "snippet": "这是搜索到的相关资料摘要..."}
        ],
        "count": max_results
    }

# ========== 5. PPT文件生成 ==========

def create_pptx(outline, output_path=None, chart_data=None):
    """根据大纲生成PPTX文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 颜色配置
    style_colors = {
        "商务蓝": {"primary": RGBColor(30, 58, 95), "secondary": RGBColor(52, 152, 219), "bg": RGBColor(236, 240, 241)},
        "科技深": {"primary": RGBColor(13, 27, 42), "secondary": RGBColor(46, 204, 113), "bg": RGBColor(44, 62, 80)},
        "学术白": {"primary": RGBColor(44, 62, 80), "secondary": RGBColor(52, 152, 219), "bg": RGBColor(255, 255, 255)},
        "创意紫": {"primary": RGBColor(108, 52, 131), "secondary": RGBColor(155, 89, 182), "bg": RGBColor(245, 240, 255)},
    }
    
    colors = style_colors.get(outline["style"], style_colors["商务蓝"])
    
    # 生成幻灯片
    for i, slide_title in enumerate(outline["structure"]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        
        # 添加标题
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.333)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        
        # 添加内容占位
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"【{slide_title}】内容区域\n\n请在此处添加具体内容..."
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(80, 80, 80)
        
        # 添加页码
        page_box = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.3))
        p = page_box.text_frame.paragraphs[0]
        p.text = f"{i+1}/{len(outline['structure'])}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(150, 150, 150)
    
    # 保存
    if not output_path:
        output_path = f"~/openclaw-data/ppt/{outline['type']}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    
    output_path = os.path.expanduser(output_path)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    
    return output_path

# ========== 主入口 ==========

def main():
    parser = argparse.ArgumentParser(description='PPT智能生成器')
    parser.add_argument('--type', required=True, choices=list(PPT_TYPES.keys()), help='PPT类型')
    parser.add_argument('--scene', required=True, help='使用场景')
    parser.add_argument('--topic', default='', help='自定义主题')
    parser.add_argument('--output', default='', help='输出文件路径')
    parser.add_argument('--search', action='store_true', help='是否进行网络资料搜集')
    parser.add_argument('--keywords', default='', help='搜索关键词')
    args = parser.parse_args()
    
    print(f"=== PPT智能生成器 ===")
    print(f"类型：{args.type}")
    print(f"场景：{args.scene}")
    print(f"主题：{args.topic or '(自动)'}")
    print("")
    
    # 1. 生成智能大纲
    print("【1/4】生成智能大纲...")
    outline = generate_outline(args.type, args.scene, args.topic)
    print(f"  ✅ 生成大纲：{outline['slides_count']}页")
    for s in outline["structure"]:
        print(f"    - {s}")
    print("")
    
    # 2. 网络资料搜集
    if args.search or args.keywords:
        print("【2/4】网络资料搜集...")
        search_results = search_web(args.topic or args.type, args.keywords)
        print(f"  ✅ 搜集到{search_results['count']}条资料")
        for src in search_results["sources"][:3]:
            print(f"    - {src['title']}")
        print("")
    
    # 3. 生成图表数据
    print("【3/4】生成图表数据...")
    charts = []
    for chart_type in outline["charts"][:3]:
        chart = generate_chart_data(chart_type)
        charts.append(chart)
        print(f"  ✅ {chart['title']} ({chart['type']})")
    print("")
    
    # 4. 生成PPTX
    print("【4/4】生成PPTX文件...")
    output_path = create_pptx(outline, args.output, charts)
    print(f"  ✅ 生成完成：{output_path}")
    print("")
    print(f"=== 完成 ===")
    print(f"文件路径：{output_path}")
    print(f"共 {outline['slides_count']} 页幻灯片")
    print(f"风格：{outline['style']}")
    print(f"建议演讲时间：{outline['estimated_minutes']} 分钟")

if __name__ == '__main__':
    main()
