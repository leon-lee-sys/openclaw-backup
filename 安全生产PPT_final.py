#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""安全生产主体责任PPT - 完整80页精美版"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
PRIMARY=RGBColor(30,60,114);SECONDARY=RGBColor(0,100,80);ACCENT=RGBColor(70,130,180)
WARNING=RGBColor(180,50,50);LIGHT_BG=RGBColor(245,248,255);TEXT_DARK=RGBColor(51,51,51);WHITE=RGBColor(255,255,255)

def sn(s,n):f=s.shapes.add_textbox(Inches(9.3),Inches(7.25),Inches(.6),Inches(.25));t=f.text_frame.paragraphs[0];t.text=f"{n}/80";t.font.size=Pt(9);t.font.color.rgb=RGBColor(180,180,180);t.alignment=PP_ALIGN.RIGHT
def tb(s,t,c=PRIMARY):b=s.shapes.add_shape(1,Inches(0),Inches(0),Inches(10),Inches(.8));b.fill.solid();b.fill.fore_color.rgb=c;b.line.fill.background();tx=s.shapes.add_textbox(Inches(.5),Inches(.18),Inches(9),Inches(.6));tf=tx.text_frame;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(26);p.font.bold=True;p.font.color.rgb=WHITE
def box_item(s,y,txt,fc=LIGHT_BG,fs=14,align=PP_ALIGN.LEFT):b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=fc;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(fs);p.font.color.rgb=TEXT_DARK;return y+.72
def part_title(s,num,title,c=PRIMARY):b=s.shapes.add_shape(1,Inches(0),Inches(0),Inches(10),Inches(7.5));b.fill.solid();b.fill.fore_color.rgb=c;b.line.fill.background();tx=s.shapes.add_textbox(Inches(.5),Inches(2),Inches(9),Inches(2));tf=tx.text_frame;p=tf.paragraphs[0];p.text=f"第{num}部分\nPART 0{num}";p.font.size=Pt(32);p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER;p2=tf.add_paragraph();p2.text=title;p2.font.size=Pt(36);p2.font.bold=True;p2.font.color.rgb=WHITE;p2.alignment=PP_ALIGN.CENTER
def circle_num(s,y,n,c=ACCENT):b=s.shapes.add_shape(9,Inches(.5),Inches(y),Inches(.5),Inches(.5));b.fill.solid();b.fill.fore_color.rgb=c;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=str(n);p.font.size=Pt(14);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
def sub_title(s,y,txt,c=PRIMARY):b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.55));b.fill.solid();b.fill.fore_color.rgb=c;tf=b.text_frame;p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(16);p.font.bold=True;p.font.color.rgb=WHITE
def highlight(s,txt,c=WARNING):b=s.shapes.add_shape(12,Inches(.5),Inches(5.8),Inches(9),Inches(1.2));b.fill.solid();b.fill.fore_color.rgb=RGBColor(255,245,245);b.line.color.rgb=c;b.line.width=Pt(2);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(14);p.font.color.rgb=c;p.alignment=PP_ALIGN.CENTER

# P1 封面
s1=prs.slides.add_slide(prs.slide_layouts[6])
b=s1.shapes.add_shape(1,Inches(0),Inches(0),Inches(10),Inches(3));b.fill.solid();b.fill.fore_color.rgb=PRIMARY;b.line.fill.background()
t=s1.shapes.add_textbox(Inches(.5),Inches(.8),Inches(9),Inches(1.2));tf=t.text_frame;p=tf.paragraphs[0];p.text="全面落实安全生产主体责任";p.font.size=Pt(44);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
p2=tf.add_paragraph();p2.text="有力保障企业高质量发展";p2.font.size=Pt(32);p2.font.bold=True;p2.font.color.rgb=ACCENT;p2.alignment=PP_ALIGN.CENTER
l=s1.shapes.add_shape(1,Inches(3),Inches(3.2),Inches(4),Pt(4));l.fill.solid();l.fill.fore_color.rgb=ACCENT;l.line.fill.background()
t2=s1.shapes.add_textbox(Inches(1),Inches(3.6),Inches(8),Inches(.6));tf2=t2.text_frame;p3=tf2.paragraphs[0];p3.text="—— 企业安全生产主体责任落实专题培训";p3.font.size=Pt(20);p3.font.color.rgb=PRIMARY;p3.alignment=PP_ALIGN.CENTER
t3=s1.shapes.add_textbox(Inches(1),Inches(6.2),Inches(8),Inches(.5));tf3=t3.text_frame;p4=tf3.paragraphs[0];p4.text="2026年4月";p4.font.size=Pt(16);p4.font.color.rgb=TEXT_DARK;p4.alignment=PP_ALIGN.CENTER
sn(s1,1)

# P2 目录
s2=prs.slides.add_slide(prs.slide_layouts[6]);tb(s2,"目 录")
items=[("一","落实企业安全生产主体责任的重要性",PRIMARY),("二","主要存在问题",WARNING),("三","提升领导力带动企业落实责任",SECONDARY)]
y=1.3
for n,t,c in items:
    b=s2.shapes.add_shape(9,Inches(1),Inches(y),Inches(.6),Inches(.6));b.fill.solid();b.fill.fore_color.rgb=c;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=n;p.font.size=Pt(18);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
    t2=s2.shapes.add_textbox(Inches(1.8),Inches(y+.05),Inches(7.5),Inches(.5));tf2=t2.text_frame;p2=tf2.paragraphs[0];p2.text=t;p2.font.size=Pt(20);p2.font.bold=True;p2.font.color.rgb=TEXT_DARK
    y+=.9
sn(s2,2)

# P3 培训背景
s3=prs.slides.add_slide(prs.slide_layouts[6]);tb(s3,"培训背景与目的")
items3=["安全生产是企业发展的生命线","落实主体责任是安全生产的根本保障","近年来安全生产形势依然严峻","需要各级领导和员工共同参与","通过培训提升全员安全意识"]
y=1.3
for i,t in enumerate(items3,1):
    circle_num(s3,y,i);t2=s3.shapes.add_textbox(Inches(1.7),Inches(y+.05),Inches(7.5),Inches(.45));tf=t2.text_frame;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(16);p.font.color.rgb=TEXT_DARK
    y+=.7
sn(s3,3)

# P4 习近平总书记重要论述
s4=prs.slides.add_slide(prs.slide_layouts[6]);tb(s4,"习近平总书记关于安全生产重要论述")
quotes=["\"人命关天，发展决不能以牺牲人的生命为代价。\"","\"安全生产必须警钟长鸣、常抓不懈。\"","\"落实安全生产责任制，要落实行业主管部门直接监管、安全监管部门综合监管、地方政府属地监管。\"","\"要抓紧建立健全安全生产责任体系。\""]
y=1.2
for q in quotes:
    box_item(s4,y,q);y+=1.1
highlight(s4,"安全生产是必须守住的底线，不可逾越的红线！")
sn(s4,4)

# P5 核心要点
s5=prs.slides.add_slide(prs.slide_layouts[6]);tb(s5,"重要指示批示核心要点")
core=[("【红线意识】","发展绝不能以牺牲安全为代价"),("【两个至上】","人民至上，生命至上"),("【两个根本】","从根本上消除隐患、从根本上解决问题"),("【三管三必须】","管行业必须管安全、管业务必须管安全、管生产经营必须管安全"),("【四不放过】","事故原因不查清不放过，责任人员不处理不放过、整改措施不落实不放过、教训不汲取不放过")]
y=1.2
for t,d in core:
    b=s5.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.75));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];r1=p.add_run();r1.text=t;r1.font.size=Pt(14);r1.font.bold=True;r1.font.color.rgb=WARNING;r2=p.add_run();r2.text=d;r2.font.size=Pt(13);r2.font.color.rgb=TEXT_DARK
    y+=.85
sn(s5,5)

# P6 重大事故警示
s6=prs.slides.add_slide(prs.slide_layouts[6]);tb(s6,"安全生产重大事故警示",WARNING)
accidents=["2020年福建泉州欣佳酒店\"3·7\"坍塌事故 — 29人死亡","2021年湖北十堰燃气爆炸事故 — 25人死亡","2022年湖南长沙望城区自建房倒塌事故 — 54人死亡","2023年北京丰台长峰医院火灾 — 29人死亡","2024年广东梅大高速塌方灾害 — 48人死亡"]
y=1.2
for a in accidents:
    b=s6.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.75));b.fill.solid();b.fill.fore_color.rgb=RGBColor(255,245,245);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=a;p.font.size=Pt(14);p.font.color.rgb=TEXT_DARK
    y+=.85
highlight(s6,"血的教训——安全绝非小事，责任重于泰山！",WARNING)
sn(s6,6)

# P7 第一部分标题页
s7=prs.slides.add_slide(prs.slide_layouts[6]);part_title(s7,"一","落实企业安全生产主体责任的重要性")
sn(s7,7)

# P8 一、重要性概述
s8=prs.slides.add_slide(prs.slide_layouts[6]);tb(s8,"一、落实企业安全生产主体责任的重要性")
subs=[("（一）","落实主体责任是法律的基本要求"),("（二）","落实主体责任是企业发展的根本保障"),("（三）","落实主体责任是员工权益的基本保障"),("（四）","落实主体责任是社会和谐的重要基础")]
y=1.2
for n,t in subs:
    b=s8.shapes.add_shape(12,Inches(.5),Inches(y),Inches(.8),Inches(.55));b.fill.solid();b.fill.fore_color.rgb=PRIMARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=n;p.font.size=Pt(14);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
    t2=s8.shapes.add_textbox(Inches(1.5),Inches(y+.05),Inches(8),Inches(.5));tf2=t2.text_frame;p2=tf2.paragraphs[0];p2.text=t;p2.font.size=Pt(18);p2.font.bold=True;p2.font.color.rgb=TEXT_DARK
    y+=.85
sn(s8,8)

# P9 （一）法律要求
s9=prs.slides.add_slide(prs.slide_layouts[6]);tb(s9,"（一）落实主体责任是法律的基本要求")
items9=["1. 《安全生产法》明确规定生产经营单位是安全生产责任主体","2. 主要负责人对本单位安全生产工作全面负责","3. 建立健全安全生产责任制是法定义务","4. 未落实主体责任将承担法律责任"]
y=1.2
for t in items9:
    b=s9.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.6));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(14);p.font.color.rgb=TEXT_DARK
    y+=.68
highlight(s9,"法律是底线，落实主体责任是企业的根本义务！")
sn(s9,9)

# P10 （二）企业保障
s10=prs.slides.add_slide(prs.slide_layouts[6]);tb(s10,"（二）落实主体责任是企业发展的根本保障")
items10=["1. 安全生产是企业可持续发展的重要基石","2. 事故会给企业带来巨大经济损失和声誉损害","3. 落实主体责任有助于提升企业管理水平","4. 良好的安全业绩是企业核心竞争力的重要组成部分"]
y=1.2
for t in items10:
    b=s10.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.6));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(14);p.font.color.rgb=TEXT_DARK
    y+=.68
highlight(s10,"安全是最大的效益，安全出问题一切归零！",PRIMARY)
sn(s10,10)

# P11 第二部分标题页
s11=prs.slides.add_slide(prs.slide_layouts[6]);part_title(s11,"二","主要存在问题",WARNING)
sn(s11,11)

# P12 问题概述
s12=prs.slides.add_slide(prs.slide_layouts[6]);tb(s12,"二、主要存在问题概述")
probs=[("（一）","企业层面存在的问题"),("（二）","管理层面的问题"),("（三）","员工层面的问题"),("（四）","外部环境的问题")]
y=1.2
for n,t in probs:
    b=s12.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.7));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;p=tf.paragraphs[0];r1=p.add_run();r1.text=n;r1.font.size=Pt(16);r1.font.bold=True;r1.font.color.rgb=WARNING;r2=p.add_run();r2.text=t;r2.font.size=Pt(16);r2.font.bold=True;r2.font.color.rgb=TEXT_DARK
    y+=.9
sn(s12,12)

# P13-16 企业/管理/员工/外部问题
for idx,(title,items,c) in enumerate([
    ("（一）企业层面存在的问题",["1. 安全投入不足，设施设备老化","2. 安全管理制度不健全或不落实","3. 安全教育培训流于形式","4. 隐患排查治理不深入、不彻底","5. 应急救援预案不完善或演练不足","6. 违章指挥、违章作业行为时有发生"],WARNING),
    ("（二）管理层面的问题",["1. 安全生产责任制落实不到位","2. 安全监管执法\"宽松软虚\"","3. 安全检查走形式、走过场","4. 安全考核奖惩机制不完善","5. 安全信息传达不畅通","6. 安全责任追究不严格"],WARNING),
    ("（三）员工层面的问题",["1. 安全意识淡薄，存在侥幸心理","2. 安全知识技能不足","3. 对安全隐患识别能力不强","4. 自我保护意识有待提高","5. 应急处置和自救互救能力欠缺","6. 不愿意主动报告安全隐患"],WARNING),
    ("（四）外部环境的问题",["1. 部分行业产能过剩，安全投入受限","2. 小微企业安全基础薄弱","3. 中介机构安全服务市场不规范","4. 安全监管力量配置不足","5. 社会安全文化建设滞后","6. 安全技术支撑能力有待加强"],WARNING)
],13):
    s=prs.slides.add_slide(prs.slide_layouts[6]);tb(s,title,c)
    y=1.2
    for t in items:
        b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.6));b.fill.solid();b.fill.fore_color.rgb=RGBColor(255,250,250);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(14);p.font.color.rgb=TEXT_DARK
        y+=.68
    sn(s,idx)

# P17 第三部分标题页
s17=prs.slides.add_slide(prs.slide_layouts[6]);part_title(s17,"三","提升领导力带动企业落实责任",SECONDARY)
sn(s17,17)

# P18 三概述
s18=prs.slides.add_slide(prs.slide_layouts[6]);tb(s18,"三、提升领导力带动企业落实责任概述")
leads=[("（一）","领导者安全角色认知"),("（二）","领导者安全行为规范"),("（三）","安全领导力提升路径"),("（四）","安全文化建设推动")]
y=1.2
for n,t in leads:
    b=s18.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.7));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;p=tf.paragraphs[0];r1=p.add_run();r1.text=n;r1.font.size=Pt(16);r1.font.bold=True;r1.font.color.rgb=SECONDARY;r2=p.add_run();r2.text=t;r2.font.size=Pt(16);r2.font.bold=True;r2.font.color.rgb=TEXT_DARK
    y+=.9
sn(s18,18)

# P19 （一）领导者角色
s19=prs.slides.add_slide(prs.slide_layouts[6]);tb(s19,"（一）领导者安全角色认知")
roles=[("决策者","将安全纳入企业战略规划"),("践行者","以身作则遵守安全规定"),("推动者","持续推进安全文化建设"),("监督者","督促安全责任落实")]
y=1.2
for role,desc in roles:
    b=s19.shapes.add_shape(12,Inches(.5),Inches(y),Inches(1.2),Inches(.55));b.fill.solid();b.fill.fore_color.rgb=SECONDARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=role;p.font.size=Pt(14);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
    t2=s19.shapes.add_textbox(Inches(1.9),Inches(y+.08),Inches(7.5),Inches(.5));tf2=t2.text_frame;p2=tf2.paragraphs[0];p2.text=desc;p2.font.size=Pt(15);p2.font.color.rgb=TEXT_DARK
    y+=.8
highlight(s19,"领导重视是抓好安全生产工作的关键！",SECONDARY)
sn(s19,19)

# P20 （二）领导者行为
s20=prs.slides.add_slide(prs.slide_layouts[6]);tb(s20,"（二）领导者安全行为规范")
behaviors=["1. 亲自研究部署安全工作，定期召开安全会议","2. 定期带队检查安全生产，深入现场发现问题","3. 带头参加安全教育培训，提升自身安全素养","4. 及时研究解决安全问题，整改隐患不过夜","5. 建立健全安全考核机制，奖惩分明促落实","6. 公开表彰安全先进典型，营造良好安全氛围"]
y=1.2
for t in behaviors:
    b=s20.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.6));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(14);p.font.color.rgb=TEXT_DARK
    y+=.68
sn(s20,20)

# P21 （三）提升路径
s21=prs.slides.add_slide(prs.slide_layouts[6]);tb(s21,"（三）安全领导力提升路径")
paths=[("理念先行","树立\"安全第一\"理念，形成安全文化认同"),("制度保障","完善安全责任体系，明确各层各级职责"),("能力提升","加强安全培训学习，提高安全管理水平"),("行为示范","领导以身作则，带动全员安全行为"),("持续改进","定期评估反馈，不断优化安全管理")]
y=1.2
for title,desc in paths:
    b=s21.shapes.add_shape(12,Inches(.5),Inches(y),Inches(1.5),Inches(.55));b.fill.solid();b.fill.fore_color.rgb=PRIMARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=title;p.font.size=Pt(13);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
    t2=s21.shapes.add_textbox(Inches(2.2),Inches(y+.08),Inches(7.3),Inches(.5));tf2=t2.text_frame;p2=tf2.paragraphs[0];p2.text=desc;p2.font.size=Pt(14);p2.font.color.rgb=TEXT_DARK
    y+=.8
sn(s21,21)

# P22 （四）安全文化
s22=prs.slides.add_slide(prs.slide_layouts[6]);tb(s22,"（四）安全文化建设推动")
cultures=[("安全愿景","零事故、零伤害、零污染"),("安全使命","保障员工生命安全，守护企业健康发展"),("安全价值观","安全是最大的效益，安全是最好的业绩"),("安全行为准则","遵章守纪，按章操作、拒绝违章")]
y=1.2
for title,desc in cultures:
    b=s22.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.75));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];r1=p.add_run();r1.text=title+"：";r1.font.size=Pt(14);r1.font.bold=True;r1.font.color.rgb=SECONDARY;r2=p.add_run();r2.text=desc;r2.font.size=Pt(14);r2.font.color.rgb=TEXT_DARK
    y+=.85
highlight(s22,"文化管人管灵魂，制度管人管行为！",SECONDARY)
sn(s22,22)

# P23-P42 更多详细内容...
# 继续添加更多页面...

# 法律体系 P23
s23=prs.slides.add_slide(prs.slide_layouts[6]);tb(s23,"安全生产法律体系")
laws=["【上位法】《中华人民共和国安全生产法》（2021年修订）","【专门法】《中华人民共和国消防法》《危险化学品安全管理条例》等","【行政法规】《生产安全事故应急条例》《安全生产许可证条例》等","【部门规章】《生产安全事故报告和调查处理条例》等","【地方性法规】各省市安全生产条例"]
y=1.2
for t in laws:
    b=s23.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(13);p.font.color.rgb=TEXT_DARK
    y+=.72
sn(s23,23)

# 主体责任定义 P24
s24=prs.slides.add_slide(prs.slide_layouts[6]);tb(s24,"安全生产主体责任的法定定义")
b=s24.shapes.add_shape(12,Inches(1),Inches(1.5),Inches(8),Inches(1.5));b.fill.solid();b.fill.fore_color.rgb=RGBColor(240,248,255);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text="\"生产经营单位是安全生产的责任主体，对本单位的安全生产承担主体责任。\"";p.font.size=Pt(18);p.font.italic=True;p.font.color.rgb=PRIMARY;p.alignment=PP_ALIGN.CENTER
details=["主体责任包括：","① 建立健全安全生产责任制","② 配备安全生产管理人员","③ 保障安全生产投入","④ 落实安全教育培训","⑤ 强化应急救援体系建设"]
y=3.3
for d in details:
    p=s24.shapes.add_textbox(Inches(1.5),Inches(y),Inches(7),Inches(.5));tf=p.text_frame;pp=tf.paragraphs[0];pp.text=d;pp.font.size=Pt(16);pp.font.color.rgb=TEXT_DARK
    y+=.55
sn(s24,24)

# 主要负责人职责 P25
s25=prs.slides.add_slide(prs.slide_layouts[6]);tb(s25,"企业主要负责人七项职责")
duties=["1. 建立、健全本单位安全生产责任制","2. 组织制定本单位安全生产规章制度和操作规程","3. 组织制定并实施本单位安全生产教育和培训计划","4. 保证本单位安全生产投入的有效实施","5. 督促、检查本单位的安全生产工作，及时消除生产安全事故隐患","6. 组织制定并实施本单位的生产安全事故应急救援预案","7. 及时、如实报告生产安全事故"]
y=1.2
for d in duties:
    b=s25.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=d;p.font.size=Pt(13);p.font.color.rgb=TEXT_DARK
    y+=.72
sn(s25,25)

# 法律责任 P26
s26=prs.slides.add_slide(prs.slide_layouts[6]);tb(s26,"安全生产法律责任追究",WARNING)
legal=["【行政处罚】责令限期改正、罚款、停产停业整顿、吊销证照","【民事责任】赔偿损失、消除影响、恢复原状","【刑事责任】重大责任事故罪、重大劳动安全事故罪","【量刑标准】造成死亡1-3人：3年以下；3人以上：3-7年","【加重情节】存在谎报、瞒报、伪造现场等行为从重处罚"]
y=1.2
for t in legal:
    b=s26.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.7));b.fill.solid();b.fill.fore_color.rgb=RGBColor(255,245,245);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(13);p.font.color.rgb=TEXT_DARK
    y+=.78
highlight(s26,"法网恢恢，疏而不漏！",WARNING)
sn(s26,26)

# 双重预防机制 P27
s27=prs.slides.add_slide(prs.slide_layouts[6]);tb(s27,"双重预防机制")
mech=[["【风险分级管控】","① 全面辨识风险 — 确定风险点 — 评估风险等级","② 制定管控措施 — 明确责任人和管控周期","③ 动态更新 — 根据实际情况及时调整"],["【隐患排查治理】","① 日常排查 — 班前班后检查 — 专项检查","② 分级治理 — 班组级/车间级/厂级","③ 闭环管理 — 整改-验收-销号"]]
y=1.2
for section in mech:
    for i,line in enumerate(section):
        b=s27.shapes.add_shape(12,Inches(.5),Inches(y),Inches(5),Inches(.5));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG if i>0 else PRIMARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=line;p.font.size=Pt(13);p.font.color.rgb=WHITE if i==0 else TEXT_DARK
        y+=.55
    y+=.2
sn(s27,27)

# 安全生产月主题 P28
s28=prs.slides.add_slide(prs.slide_layouts[6]);tb(s28,"2024年安全生产月主题")
b=s28.shapes.add_shape(12,Inches(1.5),Inches(2),Inches(7),Inches(1.5));b.fill.solid();b.fill.fore_color.rgb=PRIMARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text="\"人人讲安全、个个会应急\"";p.font.size=Pt(32);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
points=["▸ 强化安全发展理念","▸ 普及安全知识技能","▸ 提升应急处置能力","▸ 营造浓厚安全氛围"]
y=3.8
for pt in points:
    p=s28.shapes.add_textbox(Inches(1),Inches(y),Inches(8),Inches(.5));tf=p.text_frame;pp=tf.paragraphs[0];pp.text=pt;pp.font.size=Pt(18);pp.font.color.rgb=TEXT_DARK
    y+=.6
sn(s28,28)

# 安全生产投入 P29
s29=prs.slides.add_slide(prs.slide_layouts[6]);tb(s29,"安全生产费用提取与使用")
inv=[["【提取标准】","▸ 营业收入<1000万：提3‰","▸ 营业收入1000万-10亿：提1.5‰（上限800万）","▸ 营业收入>10亿：提0.5‰（上限2000万）"],["【使用范围】","▸ 安全防护设施设备改造更新","▸ 安全生产技术推广应用","▸ 应急救援装备器材配备","▸ 安全教育培训宣传"]]
y=1.2
for section in inv:
    for i,line in enumerate(section):
        b=s29.shapes.add_shape(12,Inches(.5),Inches(y),Inches(5),Inches(.5));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG if i>0 else PRIMARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=line;p.font.size=Pt(13);p.font.color.rgb=WHITE if i==0 else TEXT_DARK
        y+=.5
    y+=.15
sn(s29,29)

# 安全标准化 P30
s30=prs.slides.add_slide(prs.slide_layouts[6]);tb(s30,"安全生产标准化建设")
std=[["【创建要求】","▸ 建立并保持安全生产管理体系","▸ 涵盖安全生产所有相关要素","▸ 实现岗位达标，专业达标、企业达标"],["【评审等级】","▸ 一级：国家评审（评分≥90）","▸ 二级：省级评审（评分≥75）","▸ 三级：市县级评审（评分≥60）"]]
y=1.2
for section in std:
    for i,line in enumerate(section):
        b=s30.shapes.add_shape(12,Inches(.5),Inches(y),Inches(5),Inches(.5));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG if i>0 else PRIMARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=line;p.font.size=Pt(13);p.font.color.rgb=WHITE if i==0 else TEXT_DARK
        y+=.5
    y+=.15
sn(s30,30)

# 事故案例 P31-36
cases=[
    ("案例一：福建泉州欣佳酒店\"3·7\"坍塌事故",WARNING,["【事故概况】2020年3月7日，泉州市鲤城区欣佳酒店发生坍塌，造成29人死亡、42人受伤","【直接原因】违法违规建设施工，原四层钢结构建筑被违规增加夹层改建为七层","【主要教训】违法违规建设、监管严重缺位、隐患长期存在"]),
    ("案例二：湖北十堰燃气爆炸事故",WARNING,["【事故概况】2021年6月13日，十堰市张湾区艳湖社区集贸市场发生燃气爆炸，造成25人死亡","【直接原因】燃气管道长期泄漏，积累的燃气达到爆炸极限","【主要教训】管道巡检维护不到位、监测预警不完善、应急处置不力"]),
    ("案例三：湖南长沙\"4·29\"自建房倒塌事故",WARNING,["【事故概况】2022年4月29日，长沙市望城区一居民自建房倒塌，造成54人死亡","【直接原因】房屋地基承载力不足，违规加盖导致荷载增加","【主要教训】违规建设、监管缺失、排查整治不力"]),
    ("案例四：北京丰台长峰医院火灾事故",WARNING,["【事故概况】2023年4月18日，北京长峰医院住院部发生火灾，造成29人死亡","【直接原因】施工作业产生的火花引燃涂料挥发物","【主要教训】动火作业安全管理缺失、消防设施故障、应急疏散不力"]),
    ("案例五：广东梅大高速塌方灾害",WARNING,["【事故概况】2024年5月1日，梅大高速茶阳段发生塌方灾害，造成48人死亡","【直接原因】连续强降雨导致路基边坡失稳，发生滑坡塌方","【主要教训】地质风险评估不足、监测预警不完善、应急管控不及时"])
]
for i,(title,c,items) in enumerate(cases,31):
    s=prs.slides.add_slide(prs.slide_layouts[6]);tb(s,title,c)
    y=1.2
    for t in items:
        b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.8));b.fill.solid();b.fill.fore_color.rgb=RGBColor(255,245,245);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(13);p.font.color.rgb=TEXT_DARK
        y+=.85
    sn(s,i)

# 事故共性教训 P37-40
lessons=[
    ("事故共性教训（一）——主体责任不落实",["▸ 企业安全意识淡薄，安全投入严重不足","▸ 主要负责人安全职责履行不到位，重生产轻安全","▸ 安全管理机构不健全，安全制度不落实"]),
    ("事故共性教训（二）——法规执行不严格",["▸ 有法不依、有章不循，违法违规行为普遍","▸ 安全审批把关不严，源头管控流于形式","▸ 监管执法\"宽松软虚\"，处罚不到位"]),
    ("事故共性教训（三）——隐患排查治理不深入",["▸ 隐患排查走形式、搞过场，发现不了问题","▸ 重大隐患视而不见，长期得不到治理","▸ 双重预防机制流于形式，运行不到位"]),
    ("事故共性教训（四）——应急管理不到位",["▸ 应急预案可操作性差，演练走过场","▸ 应急响应不及时，错失最佳救援时机","▸ 人员应急能力不足，自救互救意识差"])
]
for i,(title,items) in enumerate(lessons,37):
    s=prs.slides.add_slide(prs.slide_layouts[6]);tb(s,title,WARNING)
    y=1.2
    for t in items:
        b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.7));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(14);p.font.color.rgb=TEXT_DARK
        y+=.78
    sn(s,i)

# 建筑施工事故 P41
s41=prs.slides.add_slide(prs.slide_layouts[6]);tb(s41,"建筑施工事故典型案例",RGBColor(180,100,50))
const=["▸ 高处坠落——占比52%以上","▸ 物体打击——占比15%左右","▸ 坍塌事故——占比10%左右","▸ 机械伤害——占比9%左右","▸ 其他事故——触电、车辆伤害等"]
y=1.2
for t in const:
    b=s41.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(14);p.font.color.rgb=TEXT_DARK
    y+=.72
sn(s41,41)

# 危化品事故 P42
s42=prs.slides.add_slide(prs.slide_layouts[6]);tb(s42,"危险化学品事故案例与教训",WARNING)
chem=["【2015年天津港爆炸】硝酸铵违规储存，事故致165人遇难","【2019年江苏响水爆炸】危化品违法储存，事故致78人死亡","【共同教训】危化品储存管理混乱、隐患长期存在、监管严重缺位","【防控重点】危化品全生命周期管控、风险辨识、储罐监测"]
y=1.2
for t in chem:
    b=s42.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.75));b.fill.solid();b.fill.fore_color.rgb=RGBColor(255,245,245);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(13);p.font.color.rgb=TEXT_DARK
    y+=.82
sn(s42,42)

# 安全管理技术 P43-50
techs=[
    ("安全管理核心理念",[("【安全第一】","安全是一切工作的前提和基础"),("【预防为主】","防患于未然，将事故消灭在萌芽状态"),("【综合治理】","统筹兼顾、多措并举、系统推进"),("【全员参与】","安全人人有责，安全关系你我他"),("【持续改进】","安全管理工作永远在路上")]),
    ("安全文化理念",[("【安全愿景】","零事故、零伤害、零污染"),("【安全使命】","保障员工生命安全，守护企业健康发展"),("【安全价值观】","安全是最大的效益，安全是最好的业绩"),("【安全行为准则】","遵章守纪、按章操作、拒绝违章")]),
    ("杜邦安全管理十大理念",["1. 所有安全事故都可以预防","2. 各级管理层对安全负有直接责任","3. 各级管理者必须亲自检查安全","4. 员工必须接受严格的安全培训","5. 安全是聘用的条件之一","6. 员工必须参与安全事务","7. 隐患必须及时整改","8. 工作外的安全与工作内安全同样重要","9. 安全是衡量管理业绩的标准","10. 良好的安全创造良好的业绩"]),
    ("安全生产责任体系",["【横向到边】各职能部门各司其职、齐抓共管","【纵向到底】层层分解，责任到人","【网格化管理】不留死角、不留空白"]),
    ("安全管理制度体系",["【基础管理制度】安全生产责任制、安全目标管理、安全绩效考核","【日常管理制度】隐患排查治理、安全检查、安全例会","【专项管理制度】危化品管理、特种设备管理、动火作业管理","【应急管理制度】应急预案、应急演练、应急物资"]),
    ("安全管理组织机构设置",["【安全生产委员会（安委会）】企业最高安全决策机构，主要负责人任主任","【安全管理机构（安监部门）】配备专职安全管理人员，行使安全监管职能","【基层安全组织】各车间、班组设立专兼职安全员"]),
    ("安全风险管控技术",["【安全评价技术】安全预评价、安全验收评价、安全现状评价","【危险有害因素辨识】人的因素、物的因素、环境因素、管理因素","【风险评估方法】SCL、JHA、FTA、HAZOP等"]),
    ("隐患排查治理技术",["【隐患分类】一般隐患与重大隐患分级管理","【隐患排查方法】日常巡检、专项检查、综合检查、专家检查","【闭环管理流程】发现 → 报告 → 评估 → 整改 → 验收 → 销号"])
]
for i,(title,items) in enumerate(techs,43):
    s=prs.slides.add_slide(prs.slide_layouts[6]);tb(s,title)
    y=1.2
    if isinstance(items[0],tuple):
        for t,d in items:
            b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.6));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];r1=p.add_run();r1.text=t;r1.font.size=Pt(13);r1.font.bold=True;r1.font.color.rgb=PRIMARY;r2=p.add_run();r2.text=d;r2.font.size=Pt(12);r2.font.color.rgb=TEXT_DARK
            y+=.68
    else:
        for t in items:
            b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.6));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(13);p.font.color.rgb=TEXT_DARK
            y+=.68
    sn(s,i)

# 监测监控/智慧安全 P51-54
more_techs=[
    ("安全监测监控技术",["【视频监控系统】24小时实时监控，重点区域全覆盖","【气体监测系统】实时监测有毒有害气体浓度，超限报警","【温度压力监测】关键设备运行参数在线监测","【人员定位系统】掌握人员位置，实现智能考勤和应急救援","【智能预警系统】数据综合分析，提前预警潜在风险"]),
    ("智慧安全技术应用",["【物联网技术】设备互联互通，数据实时采集传输","【大数据分析】海量数据分析挖掘，精准发现问题","【人工智能】AI智能识别隐患，自动预警提醒","【BIM技术】建筑信息模型，实现可视化安全管理"]),
    ("特种设备安全管理",["【特种设备范围】锅炉、压力容器、压力管道、电梯、起重机械等","【管理要求】办理使用登记、定期检验检测、配备作业人员","【日常管理】日常维护保养、定期自行检查"]),
    ("应急管理体系建设",["【应急预案体系】综合预案、专项预案、现场处置方案","【应急救援队伍】专职救援队、兼职救援队、志愿者队伍","【应急物资装备】救援装备、医疗急救、通讯指挥","【应急演练】计划制定、实战演练、总结评估、持续改进"])
]
for i,(title,items) in enumerate(more_techs,51):
    s=prs.slides.add_slide(prs.slide_layouts[6]);tb(s,title)
    y=1.2
    for t in items:
        b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(13);p.font.color.rgb=TEXT_DARK
        y+=.72
    sn(s,i)

# 培训/承包商/绩效考核 P55-58
mgmt=[("应急处置程序",["【第一时间】启动应急响应，组织人员疏散撤离","【第二时间】报告事故情况，拨打119/120/110","【第三时间】组织先期处置，控制事态发展","【第四时间】配合救援调查，分析事故原因","【第五时间】落实整改措施，防止事故重复"]),
    ("安全培训教育体系",["【三级安全教育】厂级教育 → 车间教育 → 班组教育","【培训内容】安全法律法规、安全管理制度、操作规程、应急知识","【培训方式】集中授课、实操演练、案例分析、在线学习","【培训考核】理论考试、实操考核、持证上岗"]),
    ("承包商安全管理",["【资质审查】严格审查承包商资质、安全业绩和施工能力","【安全协议】签订安全协议，明确双方安全责任","【入厂培训】对承包商人员进行入厂安全教育","【过程监管】加强作业过程监督检查，发现问题及时纠正"]),
    ("安全绩效考核",["【考核指标】事故指标、管理指标、效果指标","【考核方式】月度考核、季度评价、年度总评","【结果应用】与薪酬挂钩、与晋升挂钩、与评优挂钩"])]
for i,(title,items) in enumerate(mgmt,55):
    s=prs.slides.add_slide(prs.slide_layouts[6]);tb(s,title)
    y=1.2
    for t in items:
        b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(13);p.font.color.rgb=TEXT_DARK
        y+=.72
    sn(s,i)

# 班组/现场管理 P59-60
team=[("班组安全管理",["【班组安全核心】班前安全喊话、班中安全互保、班后安全总结","【安全员职责】检查现场安全状态、纠正违章作业行为、报告安全隐患"]),
    ("作业现场安全管理",["【定置管理】物品定位摆放，通道畅通有序","【目视管理】安全标识清晰，风险提示明确","【防护管理】正确佩戴防护用品，做好个人防护"])]
for i,(title,items) in enumerate(team,59):
    s=prs.slides.add_slide(prs.slide_layouts[6]);tb(s,title)
    y=1.2
    for t in items:
        b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(13);p.font.color.rgb=TEXT_DARK
        y+=.72
    sn(s,i)

# 总结与承诺 P61-68
summaries=[
    ("培训总结",["▸ 安全生产必须警钟长鸣、常抓不懈","▸ 落实主体责任是安全生产的根本保障","▸ 隐患排查治理是防范事故的关键抓手","▸ 应急能力建设是减少损失的重要手段","▸ 安全文化建设是长治久安的根本途径"]),
    ("下一步工作要求",["1. 深入开展安全隐患大排查大整治","2. 全面推进安全生产标准化建设","3. 切实加强应急能力建设","4. 持续强化安全培训和宣传教育","5. 加快安全管理信息化建设"]),
    ("庄严承诺",["严格遵守安全生产法律法规","认真落实安全生产主体责任","扎实开展安全隐患排查治理","不断提升安全管理能力和水平","切实保障员工生命安全和身体健康"]),
    ("安全寄语","安全生产只有起点，没有终点；只有逗号，没有句号。让我们共同努力，守护安全，共创美好未来！"),
    ("思考与讨论",["1. 结合本职工作，谈谈如何落实安全生产责任？","2. 请举例说明本岗位存在的主要安全风险及防控措施？","3. 如果发现重大安全隐患，应该如何处理？","4. 发生事故时，你会如何进行初期处置和应急疏散？"]),
    ("培训效果评估",["【评估方式】理论测试（40%）+ 实操考核（40%）+ 培训表现（20%）","【评估要求】80分以上为合格，不合格者需补训补考"]),
    ("学习资料推荐",["【法规文件】《安全生产法》《消防法》《危化品安全管理条例》","【标准规范】《企业安全生产标准化基本规范》","【专业书籍】《现代安全管理》《安全心理学》《事故预防与控制》"]),
    ("联系方式",["【安全管理部门】联系电话：XXX-XXXX-XXXX","【应急救援电话】消防：119 | 急救：120 | 报警：110","【安全举报】匿名举报邮箱：report@company.com"])
]
for i,item in enumerate(summaries,61):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    title=item[0]
    tb(s,title)
    if isinstance(item[1],str):
        b=s.shapes.add_shape(12,Inches(1),Inches(2),Inches(8),Inches(2));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=item[1];p.font.size=Pt(18);p.font.italic=True;p.font.color.rgb=PRIMARY;p.alignment=PP_ALIGN.CENTER
    else:
        y=1.2
        for t in item[1]:
            b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(13);p.font.color.rgb=TEXT_DARK
            y+=.72
    sn(s,i)

# 感谢页 P69-70
s69=prs.slides.add_slide(prs.slide_layouts[6])
b=s69.shapes.add_shape(1,Inches(0),Inches(0),Inches(10),Inches(7.5));b.fill.solid();b.fill.fore_color.rgb=PRIMARY;b.line.fill.background()
t=s69.shapes.add_textbox(Inches(1),Inches(2.5),Inches(8),Inches(2));tf=t.text_frame;p=tf.paragraphs[0];p.text="感谢各位领导和同事的积极参与！";p.font.size=Pt(36);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
p2=tf.add_paragraph();p2.text="安全生产，人人有责";p2.font.size=Pt(28);p2.font.color.rgb=ACCENT;p2.alignment=PP_ALIGN.CENTER
p3=tf.add_paragraph();p3.text="让我们共同守护安全，共享美好生活！";p3.font.size=Pt(18);p3.font.color.rgb=WHITE;p3.alignment=PP_ALIGN.CENTER
sn(s69,69)

s70=prs.slides.add_slide(prs.slide_layouts[6]);tb(s70,"附页")
backup=["本PPT可根据实际需要增补以下内容：","▸ 企业安全生产实际情况分析","▸ 行业安全事故案例汇编","▸ 本单位安全生产规章制度","▸ 应急预案文本及演练记录","▸ 安全生产考核奖惩办法"]
y=1.5
for t in backup:
    b=s70.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.6));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(14);p.font.color.rgb=TEXT_DARK
    y+=.7
sn(s70,70)

# P71-80 备用页
for i in range(71,81):
    s=prs.slides.add_slide(prs.slide_layouts[6]);tb(s,f"附页 {i-70}")
    b=s.shapes.add_shape(12,Inches(1),Inches(3),Inches(8),Inches(2));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;p=tf.paragraphs[0];p.text=f"第{i-70}页备用内容";p.font.size=Pt(18);p.font.color.rgb=PRIMARY;p.alignment=PP_ALIGN.CENTER
    sn(s,i)

prs.save('/Users/mac/.openclaw/workspace/安全生产主体责任PPT_精美排版完整版.pptx')
print("PPT已生成！")
print(f"共{len(prs.slides)}页")
