#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""道德经 第22-24章 课件（正确版）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
PRIMARY=RGBColor(30,60,114);SECONDARY=RGBColor(0,100,80);ACCENT=RGBColor(70,130,180)
LIGHT_BG=RGBColor(245,248,255);TEXT_DARK=RGBColor(51,51,51);WHITE=RGBColor(255,255,255)

def sn(s,n,t=7):
    f=s.shapes.add_textbox(Inches(9.3),Inches(7.25),Inches(.6),Inches(.25))
    tf=f.text_frame.paragraphs[0];tf.text=f"{n}/{t}";tf.font.size=Pt(9);tf.font.color.rgb=RGBColor(180,180,180);tf.alignment=PP_ALIGN.RIGHT

def tb(s,t,c=PRIMARY):
    b=s.shapes.add_shape(1,Inches(0),Inches(0),Inches(10),Inches(.8));b.fill.solid();b.fill.fore_color.rgb=c;b.line.fill.background()
    tx=s.shapes.add_textbox(Inches(.5),Inches(.18),Inches(9),Inches(.6));tf=tx.text_frame;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(24);p.font.bold=True;p.font.color.rgb=WHITE

def box_item(s,y,txt,c=LIGHT_BG,fs=14):
    b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=c;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(fs);p.font.color.rgb=TEXT_DARK;return y+.72

def quote_box(s,y,txt):
    b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(1.2));b.fill.solid();b.fill.fore_color.rgb=RGBColor(240,248,255);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(16);p.font.italic=True;p.font.color.rgb=PRIMARY;p.alignment=PP_ALIGN.CENTER

# 封面
s1=prs.slides.add_slide(prs.slide_layouts[6])
b=s1.shapes.add_shape(1,Inches(0),Inches(0),Inches(10),Inches(3.2));b.fill.solid();b.fill.fore_color.rgb=PRIMARY;b.line.fill.background()
t=s1.shapes.add_textbox(Inches(.5),Inches(.8),Inches(9),Inches(1.2));tf=t.text_frame;p=tf.paragraphs[0];p.text="道德经";p.font.size=Pt(52);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
p2=tf.add_paragraph();p2.text="第22-24章";p2.font.size=Pt(36);p2.font.color.rgb=ACCENT;p2.alignment=PP_ALIGN.CENTER
l=s1.shapes.add_shape(1,Inches(3),Inches(3.4),Inches(4),Pt(4));l.fill.solid();l.fill.fore_color.rgb=ACCENT;l.line.fill.background()
t2=s1.shapes.add_textbox(Inches(1),Inches(3.8),Inches(8),Inches(.6));tf2=t2.text_frame;p3=tf2.paragraphs[0];p3.text="第二十二章 不自是 | 第二十三章 希言自然\n第二十四章 自知者明";p3.font.size=Pt(16);p3.font.color.rgb=PRIMARY;p3.alignment=PP_ALIGN.CENTER
t3=s1.shapes.add_textbox(Inches(1),Inches(5.5),Inches(8),Inches(.5));tf3=t3.text_frame;p4=tf3.paragraphs[0];p4.text="2026年4月14日 · 第7天";p4.font.size=Pt(16);p4.font.color.rgb=TEXT_DARK;p4.alignment=PP_ALIGN.CENTER
sn(s1,1)

# 第22章
s2=prs.slides.add_slide(prs.slide_layouts[6]);tb(s2,"第二十二章 不自是")
quote_box(s2,1.2,"曲则全，枉则直，洼则盈，\n敝则新，少则得，多则惑。\n是以圣人抱一为天下式。\n不自是，故彰；不自见，故明；\n不自伐，故有功；不自矜，故长。")
y=2.7
box_item(s2,y,"【解读】委曲反而能保全，弯曲反而能伸直，低洼反而能充盈。");y+=.7
box_item(s2,y,"【要点】");y+=.55
for t in ["1. 曲则全：受得住委屈才能保全","2. 不自是故彰：不自以为是，反而彰显","3. 不自伐故有功：不自我夸耀，反而有功劳","4. 夫唯不争，天下莫能与之争：正因为不争，天下没人能争得过"]:
    box_item(s2,y,t,fs=12);y+=.55
sn(s2,2)

# 第23章
s3=prs.slides.add_slide(prs.slide_layouts[6]);tb(s3,"第二十三章 希言自然")
quote_box(s3,1.2,"希言自然。\n飘风不终朝，骤雨不终日。\n孰为此者？天地。\n天地尚不能久，而况于人乎？")
y=2.7
box_item(s3,y,"【解读】狂风刮不了一早晨，暴雨下不了一整天。谁造成的？天地。");y+=.7
box_item(s3,y,"【要点】");y+=.55
for t in ["1. 少说话才合乎自然之道","2. 飘风骤雨都是短暂的，不持久","3. 暴风骤雨比喻强为之事不可长久","4. 故从事于道者，道者同于道：追随道的人，道也与他同在"]:
    box_item(s3,y,t,fs=12);y+=.55
sn(s3,3)

# 第24章
s4=prs.slides.add_slide(prs.slide_layouts[6]);tb(s4,"第二十四章 自知者明")
quote_box(s4,1.2,"企者不立，跨者不行，\n自见者不明，自是者不彰，\n自伐者无功，自矜者不长。\n其在道也，曰余食赘行，\n物或恶之，故有道者不处。")
y=2.7
box_item(s4,y,"【解读】踮起脚尖站不稳，迈大步走不远。自以为是反而不明。");y+=.7
box_item(s4,y,"【要点】");y+=.55
for t in ["1. 企者不立：踮脚站不稳","2. 跨者不行：大步走不远","3. 自见者不明：固执己见反而不明白","4. 其在道也，余食赘行：从道的角度看，这些就像剩饭赘肉"]:
    box_item(s4,y,t,fs=12);y+=.55
sn(s4,4)

# 总结
s5=prs.slides.add_slide(prs.slide_layouts[6]);tb(s5,"第22-24章 总结")
y=1.2
for t in ["第二十二章：曲则全，不自是故彰","第二十三章：希言自然，飘风骤雨不终日","第二十四章：企者不立，自知者明"]:
    box_item(s5,y,t);y+=.7
b=s5.shapes.add_shape(12,Inches(1),Inches(4),Inches(8),Inches(1.5));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text="【核心要义】\n不争、谦下、自知之明。\n道家强调柔弱胜刚强，不自以为是方能有所成就。";p.font.size=Pt(14);p.font.color.rgb=PRIMARY;p.alignment=PP_ALIGN.CENTER
sn(s5,5)

# 进度表
s6=prs.slides.add_slide(prs.slide_layouts[6]);tb(s6,"道德经学习进度表")
y=1.2
for t in ["第1-3天 | 第1-9章（道可道，非常道）","第4-6天 | 第10-21章（上善若水等）","第7天（今天）| 第22-24章","第8天 | 第25-27章","第9天 | 第28-30章","第10天 | 第31-33章","..."]:
    box_item(s6,y,t);y+=.65
sn(s6,6)

# 附：正确课件说明
s7=prs.slides.add_slide(prs.slide_layouts[6]);tb(s7,"说明")
b=s7.shapes.add_shape(12,Inches(1),Inches(2),Inches(8),Inches(3));b.fill.solid();b.fill.fore_color.rgb=RGBColor(255,245,245);b.line.color.rgb=RGBColor(180,50,50);b.line.width=Pt(2);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text="【重要更正】\n\n本课件为第22-24章内容。\n\n之前错误发送的第31-35章课件作废，请以本版本为准。\n\n每日学习3章，从第1章依次递进。";p.font.size=Pt(14);p.font.color.rgb=RGBColor(180,50,50);p.alignment=PP_ALIGN.CENTER
sn(s7,7)

prs.save('/Users/mac/.openclaw/workspace/道德经_第7天_正确版.pptx')
print("道德经第7天正确版课件已生成：第22-24章")
