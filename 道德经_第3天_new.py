from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

def add_slide_number(slide, num, total):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

# ===== 第1页：封面 =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar.line.fill.background()
txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "国学经典24章 · 第三天"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

main_title = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf2 = main_title.text_frame
p2 = tf2.paragraphs[0]
p2.text = "《道德经》第7-9章"
p2.font.size = Pt(56)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER

subtitle = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf3 = subtitle.text_frame
p3 = tf3.paragraphs[0]
p3.text = "天长地久 · 上善若水 · 功遂身退"
p3.font.size = Pt(28)
p3.font.color.rgb = RGBColor(70, 130, 180)
p3.alignment = PP_ALIGN.CENTER

quote_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(6), Inches(1.2))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_q = quote_box.text_frame
tf_q.word_wrap = True
p_q1 = tf_q.paragraphs[0]
p_q1.text = "「上善若水，水善利万物而不争。」"
p_q1.font.size = Pt(22)
p_q1.font.italic = True
p_q1.font.color.rgb = RGBColor(30, 60, 114)
p_q1.alignment = PP_ALIGN.CENTER
p_q2 = tf_q.add_paragraph()
p_q2.text = "—— 《道德经》第八章"
p_q2.font.size = Pt(14)
p_q2.font.color.rgb = RGBColor(150, 150, 150)
p_q2.alignment = PP_ALIGN.CENTER

add_slide_number(slide1, 1, 12)

# ===== 第2页：第七章详解 =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar2.fill.solid()
title_bar2.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar2.line.fill.background()
txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf5 = txBox2.text_frame
p5 = tf5.paragraphs[0]
p5.text = "第七章 · 天长地久"
p5.font.size = Pt(32)
p5.font.bold = True
p5.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.2))
yuanwen_box.fill.solid()
yuanwen_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw = yuanwen_box.text_frame
tf_yw.word_wrap = True
p_yw = tf_yw.paragraphs[0]
p_yw.text = "【原文】"
p_yw.font.size = Pt(18)
p_yw.font.bold = True
p_yw.font.color.rgb = RGBColor(30, 60, 114)
for line in ["天长地久。天地所以能长且久者，以其不自生，故能长生。", "是以圣人后其身而身先，外其身而身存。", "非以其无私邪？故能成其私。"]:
    p = tf_yw.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(3.5))
tf_c = content_box.text_frame
tf_c.word_wrap = True
p_c1 = tf_c.paragraphs[0]
p_c1.text = "【白话解读】"
p_c1.font.size = Pt(18)
p_c1.font.bold = True
p_c1.font.color.rgb = RGBColor(30, 60, 114)
解读 = [
    ("天地之道", "天地之所以永恒存在，是因为它们不为自己而运作，故能长生。"),
    ("圣人之道", "圣人把自己放在最后，反而居于前列；置身度外，反而保全自身。"),
    ("核心智慧", "正因为无私，反而成就了自己——看似矛盾，实为天道。"),
]
for title, content in 解读:
    p = tf_c.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide2, 2, 12)

# ===== 第3页：第七章典故 =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar3.fill.solid()
title_bar3.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar3.line.fill.background()
txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf6 = txBox3.text_frame
p6 = tf6.paragraphs[0]
p6.text = "第七章 · 典故与案例"
p6.font.size = Pt(32)
p6.font.bold = True
p6.font.color.rgb = RGBColor(255, 255, 255)

stories = [
    ("【中国典故】范蠡三散财", "春秋范蠡辅佐越王勾践灭吴后，主动隐退，散尽家财，泛舟五湖。后又白手起家成为巨富，被后人尊为'商圣'。正因不争功名，反而名垂千古。"),
    ("【历史镜鉴】张良的智慧", "张良运筹帷幄，助刘邦定天下功成名就后，选择辟谷修道，不恋权位。功成身退，保全自身与家族安稳。"),
    ("【国外案例】华盛顿急流勇退", "美国国父华盛顿在连任两届总统后，主动放弃权力，开创了和平交接的先河。他放弃的是权力，得到的是永恒的尊敬。"),
]
y = 1.4
for title, content in stories:
    box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_b = box.text_frame
    tf_b.word_wrap = True
    p_b1 = tf_b.paragraphs[0]
    p_b1.text = title
    p_b1.font.size = Pt(16)
    p_b1.font.bold = True
    p_b1.font.color.rgb = RGBColor(30, 60, 114)
    p_b2 = tf_b.add_paragraph()
    p_b2.text = content
    p_b2.font.size = Pt(13)
    p_b2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.6

add_slide_number(slide3, 3, 12)

# ===== 第4页：第七章总结 =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar4.fill.solid()
title_bar4.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar4.line.fill.background()
txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf7 = txBox4.text_frame
p7 = tf7.paragraphs[0]
p7.text = "第七章 · 精华提炼"
p7.font.size = Pt(32)
p7.font.bold = True
p7.font.color.rgb = RGBColor(255, 255, 255)

points = [
    ("核心要义", "天地不为自己，故能长生；圣人不私，故能成其私。"),
    ("关键词", "无私、后其身、外其身、功成身退。"),
    ("生活践行", "竞争中不争一时之名，得理时让人三分。"),
    ("管理启示", "成事之后不居功，主动分享荣誉给团队。"),
]
y = 1.4
for title, content in points:
    box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_b = box.text_frame
    tf_b.word_wrap = True
    p_b1 = tf_b.paragraphs[0]
    p_b1.text = f"• {title}：{content}"
    p_b1.font.size = Pt(16)
    p_b1.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.1

conclusion_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(9), Inches(1))
conclusion_box.fill.solid()
conclusion_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_con = conclusion_box.text_frame
tf_con.word_wrap = True
p_con = tf_con.paragraphs[0]
p_con.text = "一句话总结：不为自己，反而能成就自己；懂得退让，才是真正的领先。"
p_con.font.size = Pt(16)
p_con.font.color.rgb = RGBColor(255, 255, 255)
p_con.alignment = PP_ALIGN.CENTER

add_slide_number(slide4, 4, 12)

# ===== 第5页：第八章详解 =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar5.fill.solid()
title_bar5.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar5.line.fill.background()
txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf8 = txBox5.text_frame
p8 = tf8.paragraphs[0]
p8.text = "第八章 · 上善若水"
p8.font.size = Pt(32)
p8.font.bold = True
p8.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.4))
yuanwen_box2.fill.solid()
yuanwen_box2.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw2 = yuanwen_box2.text_frame
tf_yw2.word_wrap = True
p_yw2 = tf_yw2.paragraphs[0]
p_yw2.text = "【原文】"
p_yw2.font.size = Pt(18)
p_yw2.font.bold = True
p_yw2.font.color.rgb = RGBColor(30, 60, 114)
for line in ["上善若水。水善利万物而不争，处众人之所恶，故几于道。", "居善地，心善渊，与善仁，言善信，正善治，事善能，动善时。", "夫唯不争，故无尤。"]:
    p = tf_yw2.add_paragraph()
    p.text = line
    p.font.size = Pt(17)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box2 = slide5.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(3.5))
tf_c2 = content_box2.text_frame
tf_c2.word_wrap = True
p_c2 = tf_c2.paragraphs[0]
p_c2.text = "【白话解读】"
p_c2.font.size = Pt(18)
p_c2.font.bold = True
p_c2.font.color.rgb = RGBColor(30, 60, 114)
解读2 = [
    ("最高善行", "最高善行就像水：水善于帮助万物，却不与万物争夺。"),
    ("接近于道", "水甘居众人厌恶的低洼之地，所以最接近于「道」。"),
    ("人生七善", "居地、心渊、与仁、言信、政治，事能、动时——水的七大美德。"),
    ("不争无尤", "正因为不争，所以无过——这便是水智慧的完美诠释。"),
]
for title, content in 解读2:
    p = tf_c2.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(5)

add_slide_number(slide5, 5, 12)

# ===== 第6页：第八章典故 =====
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar6.fill.solid()
title_bar6.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar6.line.fill.background()
txBox6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf9 = txBox6.text_frame
p9 = tf9.paragraphs[0]
p9.text = "第八章 · 典故与案例"
p9.font.size = Pt(32)
p9.font.bold = True
p9.font.color.rgb = RGBColor(255, 255, 255)

stories6 = [
    ("【中国典故】大禹治水", "不像父亲鲧那样用堵截，而是顺水之性，疏通九河，让水自行流入大海。不争之争，是最高明的治理。"),
    ("【历史人物】曹参的'萧规曹随'", "汉相曹参继萧何之后，不折腾、不创新，沿用既有制度。不争功，不炫耀，反而百姓安宁。"),
    ("【国外案例】圣雄甘地", "甘地以'非暴力不合作'抗争，看似柔软如水，却最终瓦解了大英帝国。柔弱胜刚强，不争而天下莫能与之争。"),
]
y = 1.4
for title, content in stories6:
    box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_b = box.text_frame
    tf_b.word_wrap = True
    p_b1 = tf_b.paragraphs[0]
    p_b1.text = title
    p_b1.font.size = Pt(16)
    p_b1.font.bold = True
    p_b1.font.color.rgb = RGBColor(30, 60, 114)
    p_b2 = tf_b.add_paragraph()
    p_b2.text = content
    p_b2.font.size = Pt(13)
    p_b2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.6

add_slide_number(slide6, 6, 12)

# ===== 第7页：第八章总结 =====
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar7.fill.solid()
title_bar7.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar7.line.fill.background()
txBox7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf10 = txBox7.text_frame
p10 = tf10.paragraphs[0]
p10.text = "第八章 · 精华提炼"
p10.font.size = Pt(32)
p10.font.bold = True
p10.font.color.rgb = RGBColor(255, 255, 255)

points7 = [
    ("核心要义", "最高善行如水，不争而利万物，故几于道。"),
    ("关键词", "上善若水、不争、利万物、七善（居，心，与，言，正，事，动）。"),
    ("生活践行", "遇事不硬顶，顺势而为；待人不争长短，平和相处。"),
    ("商业启示", "好的产品如水，润物无声，不强迫用户，自然受欢迎。"),
]
y = 1.4
for title, content in points7:
    box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_b = box.text_frame
    tf_b.word_wrap = True
    p_b1 = tf_b.paragraphs[0]
    p_b1.text = f"• {title}：{content}"
    p_b1.font.size = Pt(16)
    p_b1.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.1

conclusion_box7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(9), Inches(1))
conclusion_box7.fill.solid()
conclusion_box7.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_con7 = conclusion_box7.text_frame
tf_con7.word_wrap = True
p_con7 = tf_con7.paragraphs[0]
p_con7.text = "一句话总结：柔弱胜刚强，不争而无尤——水是道最形象的化身。"
p_con7.font.size = Pt(16)
p_con7.font.color.rgb = RGBColor(255, 255, 255)
p_con7.alignment = PP_ALIGN.CENTER

add_slide_number(slide7, 7, 12)

# ===== 第8页：第九章详解 =====
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar8.fill.solid()
title_bar8.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar8.line.fill.background()
txBox8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf11 = txBox8.text_frame
p11 = tf11.paragraphs[0]
p11.text = "第九章 · 功遂身退"
p11.font.size = Pt(32)
p11.font.bold = True
p11.font.color.rgb = RGBColor(255, 255, 255)

yuanwen_box3 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(2.2))
yuanwen_box3.fill.solid()
yuanwen_box3.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_yw3 = yuanwen_box3.text_frame
tf_yw3.word_wrap = True
p_yw3 = tf_yw3.paragraphs[0]
p_yw3.text = "【原文】"
p_yw3.font.size = Pt(18)
p_yw3.font.bold = True
p_yw3.font.color.rgb = RGBColor(30, 60, 114)
for line in ["持而盈之，不如其已。揣而锐之，不可长保。", "金玉满堂，莫之能守。富贵而骄，自遗其咎。", "功遂身退，天之道。"]:
    p = tf_yw3.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(4)

content_box3 = slide8.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(3.5))
tf_c3 = content_box3.text_frame
tf_c3.word_wrap = True
p_c3 = tf_c3.paragraphs[0]
p_c3.text = "【白话解读】"
p_c3.font.size = Pt(18)
p_c3.font.bold = True
p_c3.font.color.rgb = RGBColor(30, 60, 114)
解读3 = [
    ("持盈揣锐", "执持盈满，不如及时停止；磨砺锋芒，难以长久保持。"),
    ("金玉满堂", "金玉满堂，无法守得住；富贵而骄，必自留祸患。"),
    ("功遂身退", "功成业就之后，全身而退——这是自然天道运行的法则。"),
    ("核心警示", "贪多求全、锋芒毕露、富贵而骄——三者皆为取祸之道。"),
]
for title, content in 解读3:
    p = tf_c3.add_paragraph()
    p.text = f"• {title}：{content}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_before = Pt(6)

add_slide_number(slide8, 8, 12)

# ===== 第9页：第九章典故 =====
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar9.fill.solid()
title_bar9.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar9.line.fill.background()
txBox9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf12 = txBox9.text_frame
p12 = tf12.paragraphs[0]
p12.text = "第九章 · 典故与案例"
p12.font.size = Pt(32)
p12.font.bold = True
p12.font.color.rgb = RGBColor(255, 255, 255)

stories9 = [
    ("【中国典故】范蠡功成身退", "越王灭吴后，范蠡深知'蜚鸟尽，良弓藏'之理，主动泛舟五湖。而文种不听劝，终被勾践赐死。功成身退，保全性命；贪恋权位，祸及其身。"),
    ("【曾国藩的智慧", "曾国藩平定太平天国后，主动削减湘军，裁撤万余兵力，向朝廷表明无政治野心。正是这份'功遂身退'的清醒，保全了自身与家族。"),
    ("【国外案例】诺贝尔的明智选择", "诺贝尔在世时将大部分财产设立和平奖，功成之后不将财富留给子孙，而是回馈社会。他的名字因此永垂不朽。"),
]
y = 1.4
for title, content in stories9:
    box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_b = box.text_frame
    tf_b.word_wrap = True
    p_b1 = tf_b.paragraphs[0]
    p_b1.text = title
    p_b1.font.size = Pt(16)
    p_b1.font.bold = True
    p_b1.font.color.rgb = RGBColor(30, 60, 114)
    p_b2 = tf_b.add_paragraph()
    p_b2.text = content
    p_b2.font.size = Pt(13)
    p_b2.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.6

add_slide_number(slide9, 9, 12)

# ===== 第10页：第九章总结 =====
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar10.fill.solid()
title_bar10.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar10.line.fill.background()
txBox10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf13 = txBox10.text_frame
p13 = tf13.paragraphs[0]
p13.text = "第九章 · 精华提炼"
p13.font.size = Pt(32)
p13.font.bold = True
p13.font.color.rgb = RGBColor(255, 255, 255)

points10 = [
    ("核心要义", "功成身退是自然天道，贪恋富贵、锋芒毕露必招祸患。"),
    ("关键词", "持盈、揣锐、富贵而骄、功遂身退。"),
    ("生活践行", "见好就收，不追求完美结局；得意时想到退路。"),
    ("领导力启示", "功高震主时，主动让功于团队，不独占荣誉。"),
]
y = 1.4
for title, content in points10:
    box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_b = box.text_frame
    tf_b.word_wrap = True
    p_b1 = tf_b.paragraphs[0]
    p_b1.text = f"• {title}：{content}"
    p_b1.font.size = Pt(16)
    p_b1.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.1

conclusion_box10 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(9), Inches(1))
conclusion_box10.fill.solid()
conclusion_box10.fill.fore_color.rgb = RGBColor(70, 130, 180)
tf_con10 = conclusion_box10.text_frame
tf_con10.word_wrap = True
p_con10 = tf_con10.paragraphs[0]
p_con10.text = "一句话总结：月圆则亏，水满则溢——功成之日，即是退隐之时。"
p_con10.font.size = Pt(16)
p_con10.font.color.rgb = RGBColor(255, 255, 255)
p_con10.alignment = PP_ALIGN.CENTER

add_slide_number(slide10, 10, 12)

# ===== 第11页：三章总览 =====
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar11 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar11.fill.solid()
title_bar11.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar11.line.fill.background()
txBox11 = slide11.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf14 = txBox11.text_frame
p14 = tf14.paragraphs[0]
p14.text = "三天精华 · 第7-9章总览"
p14.font.size = Pt(32)
p14.font.bold = True
p14.font.color.rgb = RGBColor(255, 255, 255)

summary_items = [
    ("第七章 · 天长地久", "无私故能成其私", "后其身而身先，外其身而身存"),
    ("第八章 · 上善若水", "不争故无尤", "水善利万物而不争"),
    ("第九章 · 功遂身退", "功成身退天之道", "持而盈之不如其已"),
]
y = 1.4
for title, core, quote in summary_items:
    box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf_b = box.text_frame
    tf_b.word_wrap = True
    p_b1 = tf_b.paragraphs[0]
    p_b1.text = title
    p_b1.font.size = Pt(18)
    p_b1.font.bold = True
    p_b1.font.color.rgb = RGBColor(30, 60, 114)
    p_b2 = tf_b.add_paragraph()
    p_b2.text = f"核心：{core}"
    p_b2.font.size = Pt(14)
    p_b2.font.color.rgb = RGBColor(51, 51, 51)
    p_b3 = tf_b.add_paragraph()
    p_b3.text = f"原文：{quote}"
    p_b3.font.size = Pt(13)
    p_b3.font.italic = True
    p_b3.font.color.rgb = RGBColor(70, 130, 180)
    y += 1.6

add_slide_number(slide11, 11, 12)

# ===== 第12页：结束页 =====
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
title_bar12.fill.solid()
title_bar12.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar12.line.fill.background()
txBox12 = slide12.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
tf15 = txBox12.text_frame
p15 = tf15.paragraphs[0]
p15.text = "第7-9章 · 学习完成"
p15.font.size = Pt(32)
p15.font.bold = True
p15.font.color.rgb = RGBColor(255, 255, 255)

quote_box12 = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.5), Inches(7), Inches(2))
quote_box12.fill.solid()
quote_box12.fill.fore_color.rgb = RGBColor(245, 245, 240)
tf_q12 = quote_box12.text_frame
tf_q12.word_wrap = True
p_q12 = tf_q12.paragraphs[0]
p_q12.text = "「知足者富，强行者有志。"
p_q12.font.size = Pt(22)
p_q12.font.italic = True
p_q12.font.color.rgb = RGBColor(30, 60, 114)
p_q12.alignment = PP_ALIGN.CENTER
p_q12b = tf_q12.add_paragraph()
p_q12b.text = "不失其所者久，死而不亡者寿。」"
p_q12b.font.size = Pt(22)
p_q12b.font.italic = True
p_q12b.font.color.rgb = RGBColor(30, 60, 114)
p_q12b.alignment = PP_ALIGN.CENTER

themes = slide12.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
tf_t = themes.text_frame
p_t = tf_t.paragraphs[0]
p_t.text = "天长地久 · 上善若水 · 功遂身退"
p_t.font.size = Pt(24)
p_t.font.color.rgb = RGBColor(70, 130, 180)
p_t.alignment = PP_ALIGN.CENTER

add_slide_number(slide12, 12, 12)

# 保存
prs.save('/Users/mac/.openclaw/workspace/道德经_第3天.pptx')
print("PPT已生成：道德经_第3天.pptx")
print("共12页")
