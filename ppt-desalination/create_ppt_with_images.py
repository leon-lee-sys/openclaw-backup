#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
三维界面冷蒸发海水淡化 PPT 生成脚本 - 包含图片版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
TITLE_BLUE = RgbColor(0, 51, 102)
ACCENT_BLUE = RgbColor(0, 102, 153)
LIGHT_BLUE = RgbColor(204, 229, 255)
DARK_GRAY = RgbColor(64, 64, 64)
WHITE = RgbColor(255, 255, 255)

# 图片路径
IMG_DIR = "/Users/mac/.openclaw/workspace/ppt-desalination/images"

def add_title_slide(prs, title, subtitle, author, institution, date):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.5))
    tf = info_box.text_frame
    for text in [author, institution, date]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        if tf.paragraphs[0].text:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(20) if author in text else Pt(18) if institution in text else Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_toc_slide(prs, items):
    """添加目录页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "目  录"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    y_pos = 2.0
    for i, item in enumerate(items, 1):
        num_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(0.8), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i:02d}"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        text_box = slide.shapes.add_textbox(Inches(3), Inches(y_pos), Inches(8), Inches(0.6))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        
        y_pos += 0.9
    
    return slide

def add_part_slide(prs, part_num, part_title):
    """添加章节过渡页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(13.333), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(12.333), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Part {part_num}"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = part_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, sub_bullets=None):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_BLUE
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
        
        if sub_bullets and i < len(sub_bullets) and sub_bullets[i]:
            for sub in sub_bullets[i]:
                p = tf.add_paragraph()
                p.text = f"  - {sub}"
                p.font.size = Pt(18)
                p.font.color.rgb = RgbColor(100, 100, 100)
                p.space_after = Pt(8)
    
    return slide

def add_image_slide(prs, title, img_path, caption=""):
    """添加图片页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_BLUE
    
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(11.333))
    
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(128, 128, 128)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """添加表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_BLUE
    
    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5)).table
    
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    for row_idx, row in enumerate(rows, 1):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_flow_slide(prs, title, steps):
    """添加流程图页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_BLUE
    
    box_width = Inches(2.2)
    box_height = Inches(1.2)
    start_x = Inches(0.8)
    y_pos = Inches(3)
    
    for i, step in enumerate(steps):
        x_pos = start_x + i * (box_width + Inches(0.3))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE
        shape.line.color.rgb = ACCENT_BLUE
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_BLUE
        p.alignment = PP_ALIGN.CENTER
        
        if i < len(steps) - 1:
            arrow_x = x_pos + box_width + Inches(0.05)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, y_pos + Inches(0.5), Inches(0.2), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BLUE
            arrow.line.fill.background()
    
    return slide

# ========== 开始生成PPT ==========

# 第1页：封面
add_title_slide(prs, 
    "三维界面冷蒸发海水淡化",
    "新型多孔钛三维界面技术在海水淡化中的应用研究",
    "汇报人：李响",
    "上海海事大学",
    "2026年4月")

# 第2页：目录
add_toc_slide(prs, [
    "项目简介和研究目的",
    "研究内容",
    "研究方法和技术路线",
    "创新点及拟达到目标",
    "进度安排"
])

# 第3页：Part 01
add_part_slide(prs, "01", "项目简介和研究目的")

# 第4页：项目简介
add_content_slide(prs, "1.1 研究背景", [
    "全球水资源危机日益严峻，约20亿人面临水资源短缺",
    "海水淡化是解决水资源短缺的重要途径",
    "传统海水淡化技术存在能耗高、成本高等问题",
    "亟需开发新型低碳、高效的海水淡化技术"
])

# 第5页：技术痛点
add_table_slide(prs, "1.2 传统技术痛点分析", 
    ["技术类型", "能耗", "主要问题"],
    [
        ["多级闪蒸（MSF）", "8-10 kWh/m³", "能耗高、设备庞大"],
        ["反渗透（RO）", "3-5 kWh/m³", "膜污染严重、预处理复杂"],
        ["低温多效蒸馏（MED）", "6-8 kWh/m³", "结垢问题突出"],
        ["传统蒸发技术", "5-8 kWh/m³", "热效率低、操作复杂"]
    ])

# 第6页：研究目的
add_content_slide(prs, "1.3 研究目的", [
    "核心目标：揭示耦合机理与抗结垢机制，为海水淡化提供理论支撑",
    "具体目标一：构建多孔钛三维界面蒸发海水淡化系统",
    "具体目标二：深入揭示传质传热耦合机理，探究孔隙结构对蒸发性能的影响",
    "具体目标三：揭示结垢规律并提出抗结垢策略，为高效海水淡化提供理论与工程指导"
])

# 第7页：Part 02
add_part_slide(prs, "02", "研究内容")

# 第8页：研究内容概述
add_content_slide(prs, "2.1 研究内容概述", [
    "多孔钛三维界面材料设计：材料制备与表征、孔隙结构优化",
    "传质传热机理研究：耦合机理分析、理论模型构建",
    "抗结垢机制研究：结垢规律探究、抗垢策略开发"
])

# 第9页：多孔钛技术
add_content_slide(prs, "2.2 多孔钛三维界面技术", [
    "多孔钛材料特性：高比表面积、优异导热性、良好耐腐蚀性、可调控孔隙结构",
    "三维界面设计原理：毛细作用增强、蒸发面积最大化、热传导优化、蒸汽逃逸通道",
    "技术优势：能耗降低50%以上、常温（20-30°C）蒸发、结垢显著减轻、效率提高30%以上"
])

# 第10页：系统示意图
add_image_slide(prs, "2.3 系统组成示意图", 
    f"{IMG_DIR}/system_diagram.png",
    "图1：多孔钛三维界面冷蒸发海水淡化系统示意图")

# 第11页：原理对比图
add_image_slide(prs, "2.4 传统平面蒸发 vs 三维界面蒸发",
    f"{IMG_DIR}/principle_diagram.png",
    "图2：传统平面蒸发与三维界面蒸发原理对比")

# 第12页：技术对比
add_table_slide(prs, "2.5 技术性能对比", 
    ["指标", "传统技术", "本项目技术"],
    [
        ["蒸发温度", "60-80°C", "常温（20-30°C）"],
        ["能耗", "高", "降低50%以上"],
        ["结垢程度", "严重", "显著减轻"],
        ["蒸发效率", "中等", "提高30%以上"],
        ["设备复杂度", "复杂", "简化"]
    ])

# 第13页：Part 03
add_part_slide(prs, "03", "研究方法和技术路线")

# 第14页：技术路线
add_image_slide(prs, "3.1 总体技术路线",
    f"{IMG_DIR}/tech_route.png",
    "图3：项目总体技术路线图")

# 第15页：实验方法
add_content_slide(prs, "3.2 材料制备与表征技术", [
    "材料制备方法：粉末冶金法、溶胶-凝胶法、3D打印法",
    "微观结构表征：SEM形貌观察、BET孔隙分析、XRD晶相鉴定、XPS表面化学分析"
])

# 第16页：表征技术
add_table_slide(prs, "3.3 先进表征技术", 
    ["技术", "英文缩写", "主要用途", "关键参数"],
    [
        ["扫描电镜", "SEM", "微观形貌观察", "放大倍数、表面元素分布"],
        ["比表面积分析", "BET", "孔隙结构测定", "比表面积、孔径分布"],
        ["X射线衍射", "XRD", "晶相鉴定分析", "晶粒大小、晶体结构"],
        ["X射线光电子能谱", "XPS", "表面化学分析", "元素组成、化学态"]
    ])

# 第17页：计算模拟
add_content_slide(prs, "3.4 计算模拟方法", [
    "分子动力学模拟（MD）：研究水分子与材料表面相互作用、计算蒸发能垒、分析离子分布规律",
    "有限元模拟（FEM）：温度场分布模拟、速度场分析、浓度场模拟"
])

# 第18页：理论模型
add_content_slide(prs, "3.5 理论模型构建", [
    "传热传质模型：基于能量守恒的能量方程、基于质量守恒的质量方程、考虑毛细作用的渗流模型",
    "蒸发速率预测模型：结合孔结构和操作条件、预测蒸发速率、指导结构优化"
])

# 第19页：性能评价
add_table_slide(prs, "3.6 性能评价指标体系", 
    ["指标类别", "具体指标", "单位", "测试方法"],
    [
        ["蒸发性能", "蒸发速率", "kg/m²·h", "质量法"],
        ["蒸发性能", "能量效率", "%", "热平衡法"],
        ["水质指标", "产水盐度", "mg/L", "电导率法"],
        ["稳定性", "运行周期", "天", "长期测试"],
        ["耐久性", "材料寿命", "年", "加速老化试验"]
    ])

# 第20页：Part 04
add_part_slide(prs, "04", "创新点及拟达到目标")

# 第21页：创新点总览图
add_image_slide(prs, "4.1 四大核心创新点",
    f"{IMG_DIR}/innovation_summary.png",
    "图4：四大核心创新点概览")

# 第22页：创新点详解
add_content_slide(prs, "4.2 创新点详解", [
    "创新一 - 三维界面设计：提出三维界面协同增效设计理念，实现毛细作用、热传导、蒸发扩散的协同优化",
    "创新二 - 多孔钛材料：首创多孔钛三维骨架作为蒸发基底，高导热+高比表面积+优异耐腐蚀",
    "创新三 - 耦合机理：首次系统揭示传热-传质-相变多物理场耦合机理，建立定量预测模型",
    "创新四 - 抗垢策略：提出疏水-亲水图案化抗垢新策略，结合外场辅助实现双重抗垢"
])

# 第23页：材料结构图
add_image_slide(prs, "4.3 多孔钛材料微观结构",
    f"{IMG_DIR}/material_structure.png",
    "图5：多孔钛材料微观结构示意")

# 第24页：结垢机理图
add_image_slide(prs, "4.4 结垢机理与抗垢策略",
    f"{IMG_DIR}/fouling_mechanism.png",
    "图6：结垢形成机理及抗结垢策略")

# 第25页：关键技术指标
add_image_slide(prs, "4.5 性能指标对比",
    f"{IMG_DIR}/performance_comparison.png",
    "图7：传统技术与本项目目标性能对比")

# 第26页：预期成果
add_content_slide(prs, "4.6 预期研究成果", [
    "学术论文：SCI一区论文3-4篇，核心期刊2-3篇",
    "专利申请：发明专利3-5项，实用新型2-3项",
    "技术标准：形成1套技术标准或规范",
    "人才培养：培养研究生3-5名"
])

# 第27页：应用前景
add_content_slide(prs, "4.7 应用前景与市场", [
    "应用领域：海水淡化、苦咸水处理、废水回收利用、工业废水零排放",
    "市场规模：全球百亿美元级，中国千亿人民币级",
    "社会效益：缓解水资源短缺，助力碳中和目标，带动材料装备全产业链发展"
])

# 第28页：Part 05
add_part_slide(prs, "05", "进度安排")

# 第29页：总体进度
add_flow_slide(prs, "5.1 项目总体进度（4年）", [
    "第1年\n材料设计", "第2年\n实验研究", "第3年\n理论分析", "第4年\n成果总结"
])

# 第30页：详细安排
add_table_slide(prs, "5.2 详细进度安排", 
    ["阶段", "时间", "主要任务", "预期成果"],
    [
        ["第一阶段", "第1年", "材料设计、制备、表征", "材料样品、基本性能数据"],
        ["第二阶段", "第2年", "性能测试、机理研究", "性能优化方案、机理报告"],
        ["第三阶段", "第3年", "理论分析、模型构建", "理论模型、学术论文"],
        ["第四阶段", "第4年", "优化改进、工程验证", "示范工程、结题验收"]
    ])

# 第31页：阶段一
add_content_slide(prs, "5.3 阶段一：方案设计与材料制备（第1年）", [
    "Q1-Q2：文献调研与方案设计：国内外研究现状调研、技术路线确定、实验方案设计",
    "Q3-Q4：多孔钛材料制备与表征：多孔钛材料制备、微观结构表征、初步性能测试"
])

# 第32页：阶段二
add_content_slide(prs, "5.4 阶段二：实验研究（第2年）", [
    "Q1-Q2：系统性能测试与优化：不同孔结构对比、操作参数优化、长期稳定性测试",
    "Q3-Q4：传热传质与结垢机理研究：传热传质机理实验、结垢机理研究、抗垢策略验证"
])

# 第33页：阶段三四
add_content_slide(prs, "5.5 阶段三/四：理论分析与成果总结（第3-4年）", [
    "第3年：理论模型构建与论文发表：理论模型建立、数值模拟验证、高水平论文发表",
    "第4年：优化改进与工程验证：工程应用优化、成本效益分析、项目结题验收"
])

# 第34页：风险管理
add_content_slide(prs, "5.6 项目风险管理", [
    "技术风险：材料制备工艺波动 → 预案：多方比较制备工艺",
    "进度风险：实验周期较长 → 预案：并行开展多项实验",
    "应用风险：成本控制挑战 → 预案：分阶段成本核算"
])

# 第35页：团队与条件
add_content_slide(prs, "5.7 研究团队与支撑条件", [
    "研究团队：教授2名、副教授2名、研究生5-8名",
    "实验平台：材料制备实验室、性能测试平台、微观表征中心",
    "计算资源：高性能计算集群、分子动力学软件",
    "合作交流：国内外学术会议、联合研究计划"
])

# 第36页：结语
add_content_slide(prs, "5.8 项目预期贡献", [
    "理论贡献：丰富界面蒸发理论体系",
    "技术贡献：突破海水淡化能效瓶颈",
    "应用贡献：为沿海地区水资源保障提供新方案",
    "感谢各位专家批评指正！"
])

# 第37页：致谢
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2), Inches(13.333), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "感谢各位专家的指导和建议！"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "欢迎批评指正"
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 保存文件
output_dir = "/Users/mac/.openclaw/workspace/ppt-desalination"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "三维界面冷蒸发海水淡化_PPT_图文版.pptx")
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print(f"共 {len(prs.slides)} 页")
