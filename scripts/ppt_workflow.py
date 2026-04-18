#!/usr/bin/env python3
"""
PPT智能生成工作流
支持：工作汇报、数据分析、项目提案、培训课件、产品介绍、融资路演、策划方案
"""

import os
import sys

# 添加项目路径
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ========== PPT类型定义 ==========
PPT_TYPES = {
    "工作汇报": {
        "emoji": "📊",
        "description": "向领导或团队汇报工作进展、成果、问题与计划",
        "outline": [
            "封面：标题 + 汇报人 + 日期",
            "目录：汇报大纲",
            "一、工作概述：背景、目标、范围",
            "二、当前进展：完成情况、关键里程碑",
            "三、数据成果：核心指标、对比分析",
            "四、存在问题：风险点、挑战",
            "五、下步计划：改进措施、下阶段目标",
            "六、需要支持：资源需求、协调事项",
            "总结页：核心要点回顾",
            "Q&A致谢页"
        ],
        "style": "商务蓝",
        "color": RGBColor(0, 51, 102)
    },
    "数据分析": {
        "emoji": "📈",
        "description": "基于数据的深度分析报告，展示洞察与结论",
        "outline": [
            "封面：报告标题 + 数据周期 + 汇报人",
            "目录：分析框架",
            "一、分析背景：业务问题、数据来源",
            "二、数据概览：数据量、维度、质量",
            "三、核心发现：3-5个关键洞察",
            "四、详细分析：图表 + 解读",
            "五、趋势预测：未来走势研判",
            "六、建议行动：基于数据的决策建议",
            "附录：数据说明、方法论",
            "致谢页"
        ],
        "style": "科技深",
        "color": RGBColor(25, 55, 95)
    },
    "项目提案": {
        "emoji": "💼",
        "description": "向决策者展示项目价值，争取立项支持",
        "outline": [
            "封面：项目名称 + 提案人 + 日期",
            "目录：提案结构",
            "一、项目背景：市场需求、问题痛点",
            "二、解决方案：核心思路、技术路线",
            "三、价值收益：预期效果、ROI测算",
            "四、实施计划：阶段划分、时间节点",
            "五、资源需求：预算、人力、技术",
            "六、风险评估：风险点、应对策略",
            "七、团队介绍：核心成员、资质背书",
            "总结与请求：期待支持",
            "联系方式"
        ],
        "style": "商务蓝",
        "color": RGBColor(0, 82, 139)
    },
    "培训课件": {
        "emoji": "📚",
        "description": "教学或培训用，内容系统、表达清晰",
        "outline": [
            "封面：课程名称 + 讲师 + 日期",
            "课程目标：学习收获、技能掌握",
            "目录：课程大纲",
            "一、导入：为什么学这个",
            "二、概念定义：核心知识点",
            "三、原理讲解：底层逻辑",
            "四、案例分析：实际应用场景",
            "五、实践操作：步骤演示",
            "六、要点回顾：核心总结",
            "七、作业/练习：巩固所学",
            "参考资料"
        ],
        "style": "学术白",
        "color": RGBColor(70, 70, 70)
    },
    "产品介绍": {
        "emoji": "🚀",
        "description": "介绍产品特点、优势、应用场景",
        "outline": [
            "封面：产品名称 + 品牌 + 日期",
            "目录：介绍大纲",
            "一、市场背景：行业趋势、用户痛点",
            "二、产品概述：核心定位、功能全景",
            "三、核心功能：重点功能详解",
            "四、产品优势：差异化卖点、竞争力",
            "五、应用场景：典型案例、客户证言",
            "六、技术特色：架构、安全、扩展性",
            "七、合作方式：定价、方案",
            "成功案例",
            "联系我们"
        ],
        "style": "创意紫",
        "color": RGBColor(100, 50, 150)
    },
    "融资路演": {
        "emoji": "💰",
        "description": "向投资人展示项目潜力，争取融资",
        "outline": [
            "封面：项目名称 + 融资轮次 + 日期",
            "一句话介绍：电梯演讲",
            "目录：路演大纲",
            "一、痛点与市场：问题够痛、市场够大",
            "二、解决方案：产品/服务核心",
            "三、商业模式：如何赚钱",
            "四、运营数据：增长势头、核心指标",
            "五、竞争分析：护城河",
            "六、团队介绍：创始人、核心成员",
            "七、融资计划：融资金额、用途、预期",
            "愿景与使命",
            "联系方式"
        ],
        "style": "科技深",
        "color": RGBColor(30, 60, 100)
    },
    "策划方案": {
        "emoji": "🎯",
        "description": "活动、营销、品牌策划方案",
        "outline": [
            "封面：方案名称 + 策划方 + 日期",
            "目录：方案结构",
            "一、背景分析：市场环境、目标受众",
            "二、策划目标：KPI设定",
            "三、核心创意：主题概念、亮点",
            "四、内容规划：具体安排、时间线",
            "五、资源支持：预算、渠道、合作伙伴",
            "六、执行计划：甘特图、责任分工",
            "七、风险预案：备选方案",
            "八、效果评估：复盘指标",
            "附录：详细预算表"
        ],
        "style": "创意紫",
        "color": RGBColor(150, 50, 80)
    }
}

# ========== 样式定义 ==========
STYLES = {
    "商务蓝": {
        "primary": RGBColor(0, 51, 102),
        "secondary": RGBColor(68, 114, 196),
        "accent": RGBColor(0, 176, 240),
        "bg": RGBColor(255, 255, 255),
        "text": RGBColor(50, 50, 50),
        "light": RGBColor(240, 245, 250)
    },
    "学术白": {
        "primary": RGBColor(70, 70, 70),
        "secondary": RGBColor(100, 100, 100),
        "accent": RGBColor(0, 112, 192),
        "bg": RGBColor(255, 255, 255),
        "text": RGBColor(50, 50, 50),
        "light": RGBColor(245, 245, 245)
    },
    "创意紫": {
        "primary": RGBColor(100, 50, 150),
        "secondary": RGBColor(150, 80, 180),
        "accent": RGBColor(200, 100, 200),
        "bg": RGBColor(255, 255, 255),
        "text": RGBColor(50, 50, 50),
        "light": RGBColor(248, 240, 255)
    },
    "科技深": {
        "primary": RGBColor(25, 55, 95),
        "secondary": RGBColor(50, 100, 150),
        "accent": RGBColor(0, 200, 200),
        "bg": RGBColor(240, 245, 250),
        "text": RGBColor(230, 230, 230),
        "light": RGBColor(30, 60, 100)
    },
    "极简灰": {
        "primary": RGBColor(80, 80, 80),
        "secondary": RGBColor(120, 120, 120),
        "accent": RGBColor(0, 0, 0),
        "bg": RGBColor(255, 255, 255),
        "text": RGBColor(60, 60, 60),
        "light": RGBColor(245, 245, 245)
    }
}

# ========== PPT生成核心类 ==========
class PPTTemplate:
    """PPT模板基类"""
    def __init__(self, ppt_type, custom_title=None, custom_author=None):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        
        self.ppt_type = ppt_type
        self.type_info = PPT_TYPES[ppt_type]
        self.title = custom_title or ppt_type
        self.author = custom_author or "小燕子"
        self.style = STYLES[self.type_info["style"]]
        
    def add_slide(self, layout_idx=6):
        """添加空白幻灯片"""
        layout = self.prs.slide_layouts[layout_idx]
        return self.prs.slides.add_slide(layout)
    
    def add_title_slide(self):
        """添加标题幻灯片"""
        slide = self.add_slide(6)
        # 背景
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.style["primary"]
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.type_info['emoji']} {self.type_info['description']}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER
        
        # 日期和作者
        info_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5))
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        from datetime import datetime
        p.text = f"{datetime.now().strftime('%Y年%m月%d日')}  |  {self.author}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(180, 180, 180)
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_toc_slide(self):
        """添加目录幻灯片"""
        slide = self.add_slide(6)
        
        # 标题栏
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.style["primary"]
        title_bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "目 录"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 目录内容
        outline = self.type_info["outline"]
        y = 1.8
        for i, item in enumerate(outline):
            # 条目
            box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(10), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{i+1}. {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.style["text"]
            y += 0.5
        
        return slide
    
    def add_content_slide(self, title, content_items, has_sub_items=False):
        """添加内容幻灯片"""
        slide = self.add_slide(6)
        
        # 标题栏
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.style["primary"]
        title_bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 内容
        y = 1.6
        for item in content_items:
            if isinstance(item, tuple):
                # 主标题 + 子项
                main, subs = item
                box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.5))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = f"▶ {main}"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = self.style["primary"]
                y += 0.45
                
                for sub in subs:
                    box = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(11), Inches(0.4))
                    tf = box.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"  • {sub}"
                    p.font.size = Pt(16)
                    p.font.color.rgb = self.style["text"]
                    y += 0.35
            else:
                box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.5))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = f"▶ {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = self.style["text"]
                y += 0.45
        
        return slide
    
    def add_summary_slide(self, summary_points):
        """添加总结幻灯片"""
        slide = self.add_slide(6)
        
        # 标题栏
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.style["primary"]
        title_bar.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "总 结"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 总结要点
        y = 1.8
        for point in summary_points:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y), Inches(11.3), Inches(0.7))
            box.fill.solid()
            box.fill.fore_color.rgb = self.style["light"]
            box.line.color.rgb = self.style["primary"]
            
            t = box.text_frame
            p = t.paragraphs[0]
            p.text = f"✓ {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.style["primary"]
            p.alignment = PP_ALIGN.CENTER
            y += 0.9
        
        return slide
    
    def add_thanks_slide(self):
        """添加致谢幻灯片"""
        slide = self.add_slide(6)
        
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.style["primary"]
        
        # 感谢
        thanks_box = slide.shapes.add_textbox(Inches(0), Inches(3), Inches(13.33), Inches(1))
        tf = thanks_box.text_frame
        p = tf.paragraphs[0]
        p.text = "感谢聆听"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.3), Inches(13.33), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = "THANK YOU"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER
        
        # 联系信息
        contact_box = slide.shapes.add_textbox(Inches(0), Inches(5.5), Inches(13.33), Inches(0.5))
        tf = contact_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"🐦 小燕子PPT智能生成 | {self.author}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(180, 180, 180)
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def save(self, path):
        """保存PPT"""
        self.prs.save(path)
        print(f"✅ PPT已生成: {path}")


# ========== 快捷生成函数 ==========
def generate_ppt(ppt_type, title=None, author=None, output_path=None):
    """
    快速生成PPT
    
    Args:
        ppt_type: PPT类型（工作汇报/数据分析/项目提案/培训课件/产品介绍/融资路演/策划方案）
        title: 自定义标题
        author: 作者
        output_path: 输出路径
    """
    if ppt_type not in PPT_TYPES:
        raise ValueError(f"不支持的类型: {ppt_type}，支持的类型: {list(PPT_TYPES.keys())}")
    
    template = PPTTemplate(ppt_type, title, author)
    outline = template.type_info["outline"]
    
    # 生成标题页
    template.add_title_slide()
    
    # 生成目录页
    template.add_toc_slide()
    
    # 生成内容页
    for item in outline[2:-1]:  # 跳过封面、目录、总结
        # 简化内容页，实际可根据item生成详细页
        template.add_content_slide(item, [
            "• 核心要点1",
            "• 核心要点2", 
            "• 核心要点3"
        ])
    
    # 生成总结页
    template.add_summary_slide([
        "核心要点回顾1",
        "核心要点回顾2",
        "核心要点回顾3"
    ])
    
    # 生成致谢页
    template.add_thanks_slide()
    
    # 保存
    if not output_path:
        output_dir = os.path.expanduser("~/Desktop/小燕子成果文件库/")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"{ppt_type}_智能生成.pptx")
    
    template.save(output_path)
    return output_path


# ========== 列出支持的类型 ==========
def list_ppt_types():
    """列出所有支持的PPT类型"""
    print("\n📊 PPT智能生成工作流 - 支持的类型：")
    print("=" * 50)
    for ptype, info in PPT_TYPES.items():
        print(f"{info['emoji']} {ptype}: {info['description']}")
    print("=" * 50)
    return PPT_TYPES


if __name__ == "__main__":
    list_ppt_types()
    print("\n使用方法：")
    print("  from ppt_workflow import generate_ppt")
    print("  generate_ppt('工作汇报', title='2026年Q1工作汇报', author='张三')")
