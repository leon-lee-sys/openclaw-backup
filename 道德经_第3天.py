from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# 配色：古典中国风
C_DARK   = RGBColor(80, 40, 20)      # 深褐（主色）
C_MID    = RGBColor(180, 90, 40)     # 赭石
C_GOLD   = RGBColor(200, 160, 60)    # 金色
C_CREAM  = RGBColor(252, 248, 238)   # 米白背景
C_TEXT   = RGBColor(50, 30, 15)      # 深棕正文
C_WHITE  = RGBColor(255, 255, 255)

# 字体大小规范（确保清晰可读）
FONT_SIZES = {
    'title_main': 28,      # 主标题
    'title_sub': 14,       # 副标题/章节
    'body': 16,            # 正文
    'body_small': 14,       # 小字正文
    'quote': 18,           # 引用
    'footer': 11,          # 页脚
}

def set_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_CREAM

def title_bar(slide, text, chapter=""):
    """标题栏：深色背景，单行标题，简洁清晰"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_DARK; bar.line.fill.background()
    
    # 主标题 - 居中显示在标题栏内
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # 章节标识放在标题下方（正文区域顶部）
    if chapter:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
        tf = sub.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = chapter
        p.font.size = Pt(12)
        p.font.color.rgb = C_GOLD
        p.alignment = PP_ALIGN.LEFT

def footer(slide, txt="《道德经》打卡学习 · 小燕子"):
    ft = slide.shapes.add_textbox(Inches(0), Inches(7.1), Inches(10), Inches(0.4))
    tf = ft.text_frame
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(FONT_SIZES['footer'])
    p.font.color.rgb = C_MID
    p.alignment = PP_ALIGN.CENTER

def add_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    
    # 装饰横线
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.0), Inches(7), Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_GOLD; ln.line.fill.background()
    
    # 主标题
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "道德经 · 第三天"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    p2 = tf.add_paragraph()
    p2.text = "第7章 · 第8章 · 第9章"
    p2.font.size = Pt(24)
    p2.font.color.rgb = C_MID
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)
    
    # 装饰横线2
    ln2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2), Inches(7), Inches(0.05))
    ln2.fill.solid(); ln2.fill.fore_color.rgb = C_GOLD; ln2.line.fill.background()
    
    # 引言
    quote = s.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1.2))
    tf2 = quote.text_frame; tf2.word_wrap = True
    pq = tf2.paragraphs[0]
    pq.text = "「上善若水，水善利万物而不争。」"
    pq.font.size = Pt(20)
    pq.font.italic = True
    pq.font.color.rgb = C_MID
    pq.alignment = PP_ALIGN.CENTER
    
    footer(s)

# ── 第7章 ──────────────────────────────────────────────────
def add_ch7():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "第7章 · 天长地久", "Chapter 7")
    footer(s)

    # 原文框
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(9), Inches(1.8))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(240, 232, 218); box.line.color.rgb = C_GOLD
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(8.4), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "天长地久。天地所以能长且久者，以其不自生，故能长生。是以圣人后其身而身先，外其身而身存。非以其无私邪？故能成其私。"
    p.font.size = Pt(FONT_SIZES['body'])
    p.font.color.rgb = C_DARK

    # 解读
    items = [
        "天地之所以永恒存在，是因为它们不为自己而运作，故能长生。",
        "圣人把自己放在最后，反而居于前列；置身度外，反而保全自身。",
        "正因为无私，反而成就了自己——看似矛盾，实为天道。",
    ]
    for i, txt in enumerate(items):
        y = 3.7 + i * 0.9
        p = s.shapes.add_textbox(Inches(0.8), y, Inches(8.4), Inches(0.8))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = f"• {txt}"
        pp.font.size = Pt(FONT_SIZES['body_small'])
        pp.font.color.rgb = C_TEXT

def add_ch7_case():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "第7章 · 典故与案例", "Chapter 7 · Stories")
    footer(s)

    left_items = [
        "【中国典故】范蠡三散财",
        "春秋范蠡辅佐越王勾践灭吴后，主动隐退，散尽家财，泛舟五湖。后又白手起家成为巨富，被后人尊为'商圣'。正因不争功名，反而名垂千古。",
        "",
        "【历史镜鉴】张良的智慧",
        "张良运筹帷幄，助刘邦定天下功成名就后，选择辟谷修道，不恋权位。功成身退，保全自身与家族安稳。",
    ]
    right_items = [
        "【国外案例】华盛顿急流勇退",
        "美国国父华盛顿在连任两届总统后，主动放弃权力，开创了和平交接的先河。他放弃的是权力，得到的是永恒的尊敬。",
        "",
        "【管理启示】",
        "真正卓越的领导者懂得'后其身而身先'——不争功、不揽权，反而赢得团队的追随与信任。",
    ]

    for i, txt in enumerate(left_items):
        p = s.shapes.add_textbox(Inches(0.5), 1.7 + i*0.88, Inches(4.4), Inches(0.85))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = txt
        pp.font.size = Pt(13)
        pp.font.color.rgb = C_TEXT

    for i, txt in enumerate(right_items):
        p = s.shapes.add_textbox(Inches(5.1), 1.7 + i*0.88, Inches(4.4), Inches(0.85))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = txt
        pp.font.size = Pt(13)
        pp.font.color.rgb = C_TEXT

def add_ch7_summary():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "第7章 · 精华提炼", "Chapter 7 · Insights")
    footer(s)

    points = [
        "核心要义：天地不为自己，故能长生；圣人不私，故能成其私。",
        "关键词：无私、后其身、外其身、功成身退。",
        "生活践行：竞争中不争一时之名，得理时让人三分。",
        "管理启示：成事之后不居功，主动分享荣誉给团队。",
    ]
    for i, pt in enumerate(points):
        p = s.shapes.add_textbox(Inches(0.8), 1.8 + i*0.95, Inches(8.4), Inches(0.85))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = f"• {pt}"
        pp.font.size = Pt(FONT_SIZES['body'])
        pp.font.color.rgb = C_TEXT

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(8.4), Inches(1.2))
    box.fill.solid(); box.fill.fore_color.rgb = C_MID; box.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(8.0), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "一句话总结：不为自己，反而能成就自己；懂得退让，才是真正的领先。"
    p.font.size = Pt(FONT_SIZES['body'])
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

# ── 第8章 ──────────────────────────────────────────────────
def add_ch8():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "第8章 · 上善若水", "Chapter 8")
    footer(s)

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(9), Inches(2.0))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(240, 232, 218); box.line.color.rgb = C_GOLD
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(8.4), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "上善若水。水善利万物而不争，处众人之所恶，故几于道。居善地，心善渊，与善仁，言善信，正善治，事善能，动善时。夫唯不争，故无尤。"
    p.font.size = Pt(FONT_SIZES['body'])
    p.font.color.rgb = C_DARK

    items = [
        "最高善行就像水：水善于帮助万物，却不与万物争夺。",
        "水甘居众人厌恶的低洼之地，所以最接近于「道」。",
        "人生七善：居地、心渊、与仁、言信、政治，事能、动时。",
        "正因为不争，所以无过——这便是水智慧的完美诠释。",
    ]
    for i, txt in enumerate(items):
        y = 3.9 + i * 0.85
        p = s.shapes.add_textbox(Inches(0.8), y, Inches(8.4), Inches(0.78))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = f"• {txt}"
        pp.font.size = Pt(FONT_SIZES['body_small'])
        pp.font.color.rgb = C_TEXT
    box.fill.solid(); box.fill.fore_color.rgb = C_MID; box.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(8.0), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "一句话总结：不为自己，反而能成就自己；懂得退让，才是真正的领先。"
    p.font.size = Pt(FONT_SIZES['body'])
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

# ── 第8章 ──────────────────────────────────────────────────
def add_ch8():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "第8章 · 上善若水", "Chapter 8")
    footer(s)

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(9), Inches(2.0))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(240, 232, 218); box.line.color.rgb = C_GOLD
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(8.4), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "上善若水。水善利万物而不争，处众人之所恶，故几于道。居善地，心善渊，与善仁，言善信，正善治，事善能，动善时。夫唯不争，故无尤。"
    p.font.size = Pt(FONT_SIZES['body'])
    p.font.color.rgb = C_DARK

    items = [
        "最高善行就像水：水善于帮助万物，却不与万物争夺。",
        "水甘居众人厌恶的低洼之地，所以最接近于「道」。",
        "人生七善：居地、心渊、与仁、言信、政治，事能、动时。",
        "正因为不争，所以无过——这便是水智慧的完美诠释。",
    ]
    for i, txt in enumerate(items):
        y = 3.9 + i * 0.85
        p = s.shapes.add_textbox(Inches(0.8), y, Inches(8.4), Inches(0.78))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = f"• {txt}"
        pp.font.size = Pt(FONT_SIZES['body_small'])
        pp.font.color.rgb = C_TEXT

def add_ch8_case():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "第8章 · 典故与案例", "Chapter 8 · Stories")
    footer(s)

    left_items = [
        "【中国典故】大禹治水",
        "不像父亲鲧那样用堵截，而是顺水之性，疏通九河，让水自行流入大海。不争之争，是最高明的治理。",
        "",
        "【历史人物】曹参的'萧规曹随'",
        "汉相曹参继萧何之后，不折腾、不创新，沿用既有制度。不争功，不炫耀，反而百姓安宁。",
    ]
    right_items = [
        "【国外案例】圣雄甘地",
        "甘地以'非暴力不合作'抗争，看似柔软如水，却最终瓦解了大英帝国。柔弱胜刚强，不争而天下莫能与之争。",
        "",
        "【日常应用】",
        "遇冲突时，先退一步——不是软弱，而是'居善地'；说话留余地——是'言善信'；顺势而为——是'动善时'。",
    ]
    for i, txt in enumerate(left_items):
        p = s.shapes.add_textbox(Inches(0.5), 1.7 + i*0.88, Inches(4.4), Inches(0.85))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = txt
        pp.font.size = Pt(13)
        pp.font.color.rgb = C_TEXT

    for i, txt in enumerate(right_items):
        p = s.shapes.add_textbox(Inches(5.1), 1.7 + i*0.88, Inches(4.4), Inches(0.85))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = txt
        pp.font.size = Pt(13)
        pp.font.color.rgb = C_TEXT

def add_ch8_summary():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "第8章 · 精华提炼", "Chapter 8 · Insights")
    footer(s)

    points = [
        "核心要义：最高善行如水，不争而利万物，故几于道。",
        "关键词：上善若水、不争、利万物、七善（居，心，与，言，正，事，动）。",
        "生活践行：遇事不硬顶，顺势而为；待人不争长短，平和相处。",
        "商业启示：好的产品如水，润物无声，不强迫用户，自然受欢迎。",
    ]
    for i, pt in enumerate(points):
        p = s.shapes.add_textbox(Inches(0.8), 1.8 + i*0.95, Inches(8.4), Inches(0.85))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = f"• {pt}"
        pp.font.size = Pt(FONT_SIZES['body'])
        pp.font.color.rgb = C_TEXT

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(8.4), Inches(1.2))
    box.fill.solid(); box.fill.fore_color.rgb = C_MID; box.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(8.0), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "一句话总结：柔弱胜刚强，不争而无尤——水是道最形象的化身。"
    p.font.size = Pt(FONT_SIZES['body'])
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

# ── 第9章 ──────────────────────────────────────────────────
def add_ch9():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "第9章 · 功遂身退", "Chapter 9")
    footer(s)

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(9), Inches(1.8))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(240, 232, 218); box.line.color.rgb = C_GOLD
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(8.4), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "持而盈之，不如其已。揣而锐之，不可长保。金玉满堂，莫之能守。富贵而骄，自遗其咎。功遂身退，天之道。"
    p.font.size = Pt(FONT_SIZES['body'])
    p.font.color.rgb = C_DARK

    items = [
        "执持盈满，不如及时停止；磨砺锋芒，难以长久保持。",
        "金玉满堂，无法守得住；富贵而骄，必自留祸患。",
        "功成业就之后，全身而退——这是自然天道运行的法则。",
        "核心警示：贪多求全、锋芒毕露、富贵而骄——三者皆为取祸之道。",
    ]
    for i, txt in enumerate(items):
        y = 3.7 + i * 0.9
        p = s.shapes.add_textbox(Inches(0.8), y, Inches(8.4), Inches(0.82))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = f"• {txt}"
        pp.font.size = Pt(FONT_SIZES['body_small'])
        pp.font.color.rgb = C_TEXT

def add_ch9_case():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "第9章 · 典故与案例", "Chapter 9 · Stories")
    footer(s)

    left_items = [
        "【中国典故】范蠡功成身退",
        "越王灭吴后，范蠡深知'蜚鸟尽，良弓藏'之理，主动泛舟五湖。而文种不听劝，终被勾践赐死。功成身退，保全性命；贪恋权位，祸及其身。",
        "",
        "【曾国藩的智慧】",
        "曾国藩平定太平天国后，主动削减湘军，裁撤万余兵力，向朝廷表明无政治野心。正是这份'功遂身退'的清醒，保全了自身与家族。",
    ]
    right_items = [
        "【国外案例】诺贝尔的明智选择",
        "诺贝尔在世时将大部分财产设立和平奖，功成之后不将财富留给子孙，而是回馈社会。他的名字因此永垂不朽。",
        "",
        "【投资理财启示】",
        "'金玉满堂，莫之能守'——单一资产过于集中，风险极高。分散配置、适时止盈，是现代投资对这一古老智慧的最佳注解。",
    ]
    for i, txt in enumerate(left_items):
        p = s.shapes.add_textbox(Inches(0.5), 1.7 + i*0.88, Inches(4.4), Inches(0.85))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = txt
        pp.font.size = Pt(13)
        pp.font.color.rgb = C_TEXT

    for i, txt in enumerate(right_items):
        p = s.shapes.add_textbox(Inches(5.1), 1.7 + i*0.88, Inches(4.4), Inches(0.85))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = txt
        pp.font.size = Pt(13)
        pp.font.color.rgb = C_TEXT

def add_ch9_summary():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "第9章 · 精华提炼", "Chapter 9 · Insights")
    footer(s)

    points = [
        "核心要义：功成身退是自然天道，贪恋富贵、锋芒毕露必招祸患。",
        "关键词：持盈、揣锐、富贵而骄、功遂身退。",
        "生活践行：见好就收，不追求完美结局；得意时想到退路。",
        "领导力启示：功高震主时，主动让功于团队，不独占荣誉。",
    ]
    for i, pt in enumerate(points):
        p = s.shapes.add_textbox(Inches(0.8), 1.8 + i*0.95, Inches(8.4), Inches(0.85))
        tf = p.text_frame; tf.word_wrap = True
        pp = tf.paragraphs[0]
        pp.text = f"• {pt}"
        pp.font.size = Pt(FONT_SIZES['body'])
        pp.font.color.rgb = C_TEXT

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(8.4), Inches(1.2))
    box.fill.solid(); box.fill.fore_color.rgb = C_MID; box.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(8.0), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "一句话总结：月圆则亏，水满则溢——功成之日，即是退隐之时。"
    p.font.size = Pt(FONT_SIZES['body'])
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

# ── 三章总结 ──────────────────────────────────────────────
def add_3ch_summary():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    title_bar(s, "三天精华 · 第7-9章总览", "Review: Ch 7-9")
    footer(s)

    data = [
        ("第7章", "天长地久", "无私故能成其私", "后其身而身先"),
        ("第8章", "上善若水", "不争故无尤", "水善利万物而不争"),
        ("第9章", "功遂身退", "功成身退天之道", "持而盈之不如其已"),
    ]
    for i, (ch, title, core, quote) in enumerate(data):
        y = 1.7 + i * 1.7
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(9), Inches(1.5))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(240, 232, 218); box.line.color.rgb = C_GOLD
        
        ch_tb = s.shapes.add_textbox(Inches(0.7), y+0.12, Inches(2), Inches(0.45))
        p = ch_tb.text_frame.paragraphs[0]
        p.text = ch
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_MID
        
        ttl = s.shapes.add_textbox(Inches(2.5), y+0.12, Inches(6.8), Inches(0.5))
        p = ttl.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        
        cr = s.shapes.add_textbox(Inches(2.5), y+0.65, Inches(6.8), Inches(0.4))
        p = cr.text_frame.paragraphs[0]
        p.text = f"核心：{core}"
        p.font.size = Pt(14)
        p.font.color.rgb = C_TEXT
        
        qt = s.shapes.add_textbox(Inches(2.5), y+1.05, Inches(6.8), Inches(0.4))
        p = qt.text_frame.paragraphs[0]
        p.text = f"原文：{quote}"
        p.font.size = Pt(13)
        p.font.italic = True
        p.font.color.rgb = C_MID

def add_end():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.6), Inches(7), Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_GOLD; ln.line.fill.background()
    
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "第7-9章 · 学习完成"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "天长地久 · 上善若水 · 功遂身退"
    p2.font.size = Pt(22)
    p2.font.color.rgb = C_MID
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)
    
    ln2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.5), Inches(7), Inches(0.05))
    ln2.fill.solid(); ln2.fill.fore_color.rgb = C_GOLD; ln2.line.fill.background()
    
    qt = s.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(1.2))
    tf2 = qt.text_frame; tf2.word_wrap = True
    pq = tf2.paragraphs[0]
    pq.text = "「知足者富，强行者有志。不失其所者久，死而不亡者寿。」"
    pq.font.size = Pt(17)
    pq.font.italic = True
    pq.font.color.rgb = C_MID
    pq.alignment = PP_ALIGN.CENTER
    
    footer(s)

# ── 生成 ───────────────────────────────────────────────────
add_cover()
add_ch7(); add_ch7_case(); add_ch7_summary()
add_ch8(); add_ch8_case(); add_ch8_summary()
add_ch9(); add_ch9_case(); add_ch9_summary()
add_3ch_summary()
add_end()

out = "/Users/mac/.openclaw/workspace/道德经_第3天.pptx"
prs.save(out)
print(f"已保存：{out}")
print(f"共 {len(prs.slides)} 页")
