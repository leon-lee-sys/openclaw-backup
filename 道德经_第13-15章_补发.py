#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""道德经 第13-15章 课件（补发版）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
PRIMARY=RGBColor(30,60,114);SECONDARY=RGBColor(0,100,80);ACCENT=RGBColor(70,130,180)
LIGHT_BG=RGBColor(245,248,255);TEXT_DARK=RGBColor(51,51,51);WHITE=RGBColor(255,255,255)

def sn(s,n,t=7):
    f=s.shapes.add_textbox(Inches(9.3),Inches(7.25),Inches(.6),Inches(.25))
    tf=f.text_frame.paragraphs[0];tf.text=f"{n}/{t}";tf.font.size=Pt(9);tf.font.color.rgb=RGBColor(180,180,180);tf.alignment=PP_ALIGN.RIGHT

def tb(s,t,c=PRIMARY):
    b=s.shapes.add_shape(1,Inches(0),Inches(0),Inches(10),Inches(.8));b.fill.solid();b.fill.fore_color.rgb=c;b.line.fill.background()
    tx=s.shapes.add_textbox(Inches(.5),Inches(.18),Inches(9),Inches(.6));tf=tx.text_frame;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(24);p.font.bold=True;p.font.color.rgb=WHITE

def box_item(s,y,txt,c=LIGHT_BG,fs=14):
    b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=c;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(fs);p.font.color.rgb=TEXT_DARK;return y+.72

def quote_box(s,y,txt):
    b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(1.2));b.fill.solid();b.fill.fore_color.rgb=RGBColor(240,248,255);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(16);p.font.italic=True;p.font.color.rgb=PRIMARY;p.alignment=PP_ALIGN.CENTER

# 封面
s1=prs.slides.add_slide(prs.slide_layouts[6])
b=s1.shapes.add_shape(1,Inches(0),Inches(0),Inches(10),Inches(3.2));b.fill.solid();b.fill.fore_color.rgb=PRIMARY;b.line.fill.background()
t=s1.shapes.add_textbox(Inches(.5),Inches(.8),Inches(9),Inches(1.2));tf=t.text_frame;p=tf.paragraphs[0];p.text="道德经";p.font.size=Pt(52);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
p2=tf.add_paragraph();p2.text="第13-15章（补发）";p2.font.size=Pt(36);p2.font.color.rgb=ACCENT;p2.alignment=PP_ALIGN.CENTER
l=s1.shapes.add_shape(1,Inches(3),Inches(3.4),Inches(4),Pt(4));l.fill.solid();l.fill.fore_color.rgb=ACCENT;l.line.fill.background()
t2=s1.shapes.add_textbox(Inches(1),Inches(3.8),Inches(8),Inches(.6));tf2=t2.text_frame;p3=tf2.paragraphs[0];p3.text="第十三章 身与天下 | 第十四章 视之不见\n第十五章 微妙玄通";p3.font.size=Pt(16);p3.font.color.rgb=PRIMARY;p3.alignment=PP_ALIGN.CENTER
t3=s1.shapes.add_textbox(Inches(1),Inches(5.5),Inches(8),Inches(.5));tf3=t3.text_frame;p4=tf3.paragraphs[0];p4.text="2026年4月（补发）";p4.font.size=Pt(16);p4.font.color.rgb=TEXT_DARK;p4.alignment=PP_ALIGN.CENTER
sn(s1,1)

# 第13章
s2=prs.slides.add_slide(prs.slide_layouts[6]);tb(s2,"第十三章 身与天下")
quote_box(s2,1.2,"宠辱若惊，贵大患若身。\n何谓宠辱若惊？\n宠为下，得之若惊，失之若惊，\n是谓宠辱若惊。\n何谓贵大患若身？\n吾所以有大患者，为吾有身，\n及吾无身，吾有何患？")
y=2.7
box_item(s2,y,"【解读】受宠或受辱都像受惊一样，害怕大的祸患是因为看重自身。");y+=.7
box_item(s2,y,"【要点】");y+=.55
for t in ["1. 宠辱若惊：受宠和受辱都令人惊惶","2. 贵大患若身：重视祸患如同重视身体","3. 无身则无患：如果没有自身，还有什么祸患","4. 故贵以身为天下：把天下看得比自身更重"]:
    box_item(s2,y,t,fs=12);y+=.55
sn(s2,2)

# 第14章
s3=prs.slides.add_slide(prs.slide_layouts[6]);tb(s3,"第十四章 视之不见")
quote_box(s3,1.2,"视之不见名曰夷，听之不闻名曰希，\n搏之不得名曰微。\n此三者不可致诘，故混而为一。\n其上不皦，其下不昧，\n绳绳不可名，复归于无物。")
y=2.7
box_item(s3,y,"【解读】道看不见、听不到、摸不着，是超越感官的存在。");y+=.7
box_item(s3,y,"【要点】");y+=.55
for t in ["1. 视之不见曰夷：看不见叫做夷","2. 听之不闻名曰希：听不到叫做希","3. 搏之不得名曰微：摸不着叫做微","4. 道之为物，惟恍惟惚：道这个东西，若有若无"]:
    box_item(s3,y,t,fs=12);y+=.55
sn(s3,3)

# 第15章
s4=prs.slides.add_slide(prs.slide_layouts[6]);tb(s4,"第十五章 微妙玄通")
quote_box(s4,1.2,"古之善为士者，微妙玄通，深不可识。\n夫唯不可识，故强为之容：\n豫兮若冬涉川，犹兮若畏四邻，\n俨兮其若客，涣兮若冰之将释，\n敦兮其若朴，旷兮其若谷。")
y=2.7
box_item(s4,y,"【解读】古代得道之人，精微奥妙而深远通达，深不可识。");y+=.7
box_item(s4,y,"【要点】");y+=.55
for t in ["1. 豫兮若冬涉川：谨慎啊，如冬天涉水过河","2. 犹兮若畏四邻：警觉啊，如害怕四邻","3. 俨兮其若客：庄重啊，如做客一样","4. 涣兮若冰之释：柔和啊，如冰雪消融"]:
    box_item(s4,y,t,fs=12);y+=.55
sn(s4,4)

# 总结
s5=prs.slides.add_slide(prs.slide_layouts[6]);tb(s5,"第13-15章 总结")
y=1.2
for t in ["第十三章：宠辱若惊，贵身若贵天下","第十四章：视之不见，听之不闻","第十五章：微妙玄通，深不可识"]:
    box_item(s5,y,t);y+=.7
b=s5.shapes.add_shape(12,Inches(1),Inches(4),Inches(8),Inches(1.5));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text="【核心要义】\n道超越感官，无法用言语形容。\n得道之人谨慎谦下，微妙玄通。";p.font.size=Pt(14);p.font.color.rgb=PRIMARY;p.alignment=PP_ALIGN.CENTER
sn(s5,5)

# 修正后的进度
s6=prs.slides.add_slide(prs.slide_layouts[6]);tb(s6,"修正后的学习进度")
y=1.2
for t in ["第1-4天 | 第1-12章 ✓","第5天 | 第13-15章 ← 补发（今天）","第6天 | 第16-18章（之前已发）","第7天 | 第19-21章（之前已发）","第8天 | 第22-24章（明天）","第9天 | 第25-27章","..."]:
    box_item(s6,y,t);y+=.65
sn(s6,6)

# 总结说明
s7=prs.slides.add_slide(prs.slide_layouts[6]);tb(s7,"总结说明")
b=s7.shapes.add_shape(12,Inches(1),Inches(2),Inches(8),Inches(3));b.fill.solid();b.fill.fore_color.rgb=RGBColor(255,245,245);b.line.color.rgb=RGBColor(180,50,50);b.line.width=Pt(2);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text="【问题原因】\n\n从第5天开始，课件制作就跳过了第13-15章。\n\n第5天直接发了16-18章，第6天发了19-21章，第7天发了22-24章。\n\n本次补发第13-15章，之后按正确顺序继续。";p.font.size=Pt(14);p.font.color.rgb=RGBColor(180,50,50);p.alignment=PP_ALIGN.CENTER
sn(s7,7)

prs.save('/Users/mac/.openclaw/workspace/道德经_第13-15章_补发版.pptx')
print("道德经第13-15章补发版课件已生成")
