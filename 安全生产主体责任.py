#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
安全生产主体责任PPT - 80页
题目：全面落实安全生产主体责任 有力保障企业高质量发展
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

def add_slide_number(slide, num, total=80):
    footer = slide.shapes.add_textbox(Inches(9), Inches(7.2), Inches(0.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def add_title_bar(slide, title_text, color=RGBColor(30, 60, 114)):
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = color
    title_bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_content_box(slide, content_list, start_y=1.5):
    """添加内容框"""
    for i, item in enumerate(content_list):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(start_y + i*0.85), Inches(9), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 240)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 51, 51)

# ==================== 第1-10页：封面与目录 ====================

# 第1页：封面
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(2.5))
title_bar1.fill.solid()
title_bar1.fill.fore_color.rgb = RGBColor(30, 60, 114)
title_bar1.line.fill.background()
main_title = slide1.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.5))
tf = main_title.text_frame
p = tf.paragraphs[0]
p.text = "全面落实安全生产主体责任"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "有力保障企业高质量发展"
p2.font.size = Pt(36)
p2.font.bold = True
p2.font.color.rgb = RGBColor(255, 255, 255)
p2.alignment = PP_ALIGN.CENTER

subtitle = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
tf2 = subtitle.text_frame
p3 = tf2.paragraphs[0]
p3.text = "—— 企业安全生产主体责任落实专题培训"
p3.font.size = Pt(24)
p3.font.color.rgb = RGBColor(70, 130, 180)
p3.alignment = PP_ALIGN.CENTER

add_slide_number(slide1, 1)

# 第2页：目录
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide2, "目 录")
contents = [
    "第一部分  安全生产重要指示批示",
    "第二部分  安全生产法律法规政策",
    "第三部分  国内外安全生产发展现状",
    "第四部分  重大典型安全生产事故案例分析",
    "第五部分  安全管理理念、技术和体系"
]
y = 1.5
for i, item in enumerate(contents):
    box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y), Inches(8), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{i+1}. {item}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 60, 114)
    y += 1.05
add_slide_number(slide2, 2)

# 第3页：培训背景
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide3, "培训背景与目的")
bg_content = [
    "▸ 安全生产是企业发展的生命线",
    "▸ 落实主体责任是安全生产的根本保障",
    "▸ 近年来安全生产形势依然严峻",
    "▸ 需要各级领导和员工共同参与",
    "▸ 通过培训提升全员安全意识"
]
add_content_box(slide3, bg_content, 1.5)
add_slide_number(slide3, 3)

# 第4页：习近平总书记重要论述
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide4, "习近平总书记关于安全生产重要论述")
xi_content = [
    "\"人命关天，发展决不能以牺牲人的生命为代价。\"",
    "\"安全生产必须警钟长鸣、常抓不懈。\"",
    "\"落实安全生产责任制，要落实行业主管部门直接监管、安全监管部门综合监管、地方政府属地监管。\"",
    "\"要抓紧建立健全安全生产责任体系。\""
]
y = 1.4
for item in xi_content:
    box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    y += 1.2
add_slide_number(slide4, 4)

# 第5页：安全生产重要指示批示（续）
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide5, "重要指示批示核心要点")
core_content = [
    "【红线意识】发展绝不能以牺牲安全为代价",
    "【两个至上】人民至上、生命至上",
    "【两个根本】从根本上消除隐患、从根本上解决问题",
    "【三管三必须】管行业必须管安全、管业务必须管安全、管生产经营必须管安全",
    "【四不放过】事故原因不查清不放过、责任人员不处理不放过、整改措施不落实不放过、教训不汲取不放过"
]
add_content_box(slide5, core_content, 1.5)
add_slide_number(slide5, 5)

# 第6页：安全生产重大事故警示
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide6, "安全生产重大事故警示", RGBColor(180, 50, 50))
warning_content = [
    "▸ 2020年福建泉州欣佳酒店\"3·7\"坍塌事故 — 29人死亡",
    "▸ 2021年湖北十堰燃气爆炸事故 — 25人死亡",
    "▸ 2022年湖南长沙望城区自建房倒塌事故 — 54人死亡",
    "▸ 2023年北京丰台长峰医院火灾 — 29人死亡",
    "▸ 2024年广东梅大高速塌方灾害 — 48人死亡"
]
add_content_box(slide6, warning_content, 1.5)
add_slide_number(slide6, 6)

# 第7页：安全生产重要指示批示-小结
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide7, "第一部分小结")
summary1 = [
    "✓ 习近平总书记关于安全生产的重要论述是做好安全生产工作的根本遵循",
    "✓ 必须坚持以人民为中心的发展思想",
    "✓ 牢牢守住安全生产底线",
    "✓ 切实保障人民群众生命财产安全"
]
add_content_box(slide7, summary1, 2)
add_slide_number(slide7, 7)

# 第8页：第二部分开场
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
part2 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(part2, "第二部分", RGBColor(0, 100, 80))
p = part2.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
tf = p.text_frame
p2 = tf.paragraphs[0]
p2.text = "安全生产法律法规政策"
p2.font.size = Pt(40)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER
add_slide_number(part2, 8)

# 第9页：安全生产法律体系
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide9, "我国安全生产法律体系")
law_content = [
    "【上位法】《中华人民共和国安全生产法》（2021年修订）",
    "【专门法】《中华人民共和国消防法》《危险化学品安全管理条例》等",
    "【行政法规】《生产安全事故应急条例》《安全生产许可证条例》等",
    "【部门规章】《生产安全事故报告和调查处理条例》等",
    "【地方性法规】各省市安全生产条例"
]
add_content_box(slide9, law_content, 1.5)
add_slide_number(slide9, 9)

# 第10页：安全生产法核心内容
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide10, "《安全生产法》核心内容")
law_detail = [
    "▸ 明确生产经营单位的主体责任",
    "▸ 规定主要负责人七项职责",
    "▸ 建立安全生产标准化建设制度",
    "▸ 完善安全风险分级管控和隐患排查治理双重预防机制",
    "▸ 强化法律责任追究"
]
add_content_box(slide10, law_detail, 1.5)
add_slide_number(slide10, 10)

# ==================== 第11-25页：法律政策 ====================

# 第11页：主体责任定义
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide11, "安全生产主体责任的法定定义")
definition = "\"生产经营单位是安全生产的责任主体，对本单位的安全生产承担主体责任。\""
box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(8), Inches(1.5))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(240, 248, 255)
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = definition
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = RGBColor(30, 60, 114)
p.alignment = PP_ALIGN.CENTER
details = ["主体责任包括：", "① 建立健全安全生产责任制", "② 配备安全生产管理人员", "③ 保障安全生产投入", "④ 落实安全教育培训", "⑤ 强化应急救援体系建设"]
y = 3.3
for d in details:
    p = slide11.shapes.add_textbox(Inches(1.5), Inches(y), Inches(7), Inches(0.5))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = d
    pp.font.size = Pt(16)
    pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.55
add_slide_number(slide11, 11)

# 第12页：企业主要负责人职责
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide12, "企业主要负责人七项职责")
duties = [
    "1. 建立、健全本单位安全生产责任制",
    "2. 组织制定本单位安全生产规章制度和操作规程",
    "3. 组织制定并实施本单位安全生产教育和培训计划",
    "4. 保证本单位安全生产投入的有效实施",
    "5. 督促、检查本单位的安全生产工作，及时消除生产安全事故隐患",
    "6. 组织制定并实施本单位的生产安全事故应急救援预案",
    "7. 及时、如实报告生产安全事故"
]
y = 1.4
for duty in duties:
    box = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = duty
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.75
add_slide_number(slide12, 12)

# 第13页：法律责任追究
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide13, "安全生产法律责任追究", RGBColor(180, 50, 50))
legal_content = [
    "【行政处罚】责令限期改正、罚款、停产停业整顿、吊销证照",
    "【民事责任】赔偿损失、消除影响、恢复原状",
    "【刑事责任】重大责任事故罪、重大劳动安全事故罪、强令违章冒险作业罪",
    "【量刑标准】造成死亡1-3人：3年以下；3人以上：3-7年",
    "【加重情节】存在谎报、瞒报、伪造现场等行为从重处罚"
]
add_content_box(slide13, legal_content, 1.5)
add_slide_number(slide13, 13)

# 第14页：2024年安全生产月主题
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide14, "2024年安全生产月主题")
theme = "\"人人讲安全、个个会应急\""
box = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2), Inches(7), Inches(1.5))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(30, 60, 114)
tf = box.text_frame
p = tf.paragraphs[0]
p.text = theme
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
points = ["▸ 强化安全发展理念", "▸ 普及安全知识技能", "▸ 提升应急处置能力", "▸ 营造浓厚安全氛围"]
y = 3.8
for pt in points:
    p = slide14.shapes.add_textbox(Inches(2), Inches(y), Inches(6), Inches(0.5))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = pt
    pp.font.size = Pt(18)
    pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.6
add_slide_number(slide14, 14)

# 第15页：政策文件体系
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide15, "安全生产重要政策文件")
policy = [
    "▸ 《中共中央 国务院关于推进安全生产领域改革发展的意见》",
    "▸ 《安全生产\"十四五\"规划》",
    "▸ 《生产安全事故应急条例》（国务院令第708号）",
    "▸ 《危险化学品安全综合治理方案》",
    "▸ 《工贸行业安全生产专项整治三年行动实施方案》",
    "▸ 《企业安全生产标准化基本规范（GB/T 33000-2016）》"
]
add_content_box(slide15, policy, 1.5)
add_slide_number(slide15, 15)

# 第16页：双重预防机制
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide16, "双重预防机制")
mechanism = [
    "【风险分级管控】",
    "① 全面辨识风险 — 确定风险点 — 评估风险等级",
    "② 制定管控措施 — 明确责任人和管控周期",
    "③ 动态更新 — 根据实际情况及时调整",
    "",
    "【隐患排查治理】",
    "① 日常排查 — 班前班后检查 — 专项检查",
    "② 分级治理 — 班组级/车间级/厂级",
    "③ 闭环管理 — 整改-验收-销号"
]
y = 1.4
for line in mechanism:
    p = slide16.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.5))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(13)
    if '【' in line:
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(30, 60, 114)
    else:
        pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.45
add_slide_number(slide16, 16)

# 第17页：安全生产投入保障
slide17 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide17, "安全生产费用提取与使用")
investment = [
    "【提取标准】",
    "▸ 营业收入<1000万：提3‰",
    "▸ 营业收入1000万-10亿：提1.5‰（上限800万）",
    "▸ 营业收入>10亿：提0.5‰（上限2000万）",
    "",
    "【使用范围】",
    "▸ 安全防护设施设备改造更新",
    "▸ 安全生产技术推广应用",
    "▸ 应急救援装备器材配备",
    "▸ 安全评价评估检测检验",
    "▸ 安全教育培训宣传"
]
y = 1.4
for line in investment:
    p = slide17.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.5))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(13)
    if '【' in line:
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(30, 60, 114)
    else:
        pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.42
add_slide_number(slide17, 17)

# 第18页：安全标准化建设
slide18 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide18, "安全生产标准化建设")
standard = [
    "【创建要求】",
    "▸ 建立并保持安全生产管理体系",
    "▸ 涵盖安全生产所有相关要素",
    "▸ 实现岗位达标、专业达标、企业达标",
    "",
    "【评审等级】",
    "▸ 一级：国家评审（评分≥90）",
    "▸ 二级：省级评审（评分≥75）",
    "▸ 三级：市县级评审（评分≥60）"
]
y = 1.4
for line in standard:
    p = slide18.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.5))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(13)
    if '【' in line:
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(30, 60, 114)
    else:
        pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.45
add_slide_number(slide18, 18)

# 第19页：法律政策小结
slide19 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide19, "第二部分小结")
summary2 = [
    "✓ 安全生产法律体系日趋完善",
    "✓ 主体责任法定化、清单化",
    "✓ 双重预防机制成为核心抓手",
    "✓ 安全生产投入有法定标准",
    "✓ 标准化建设持续推进"
]
add_content_box(slide19, summary2, 2)
add_slide_number(slide19, 19)

# 第20页：第三部分开场
slide20 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide20, "第三部分", RGBColor(0, 100, 80))
p = slide20.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
tf = p.text_frame
p2 = tf.paragraphs[0]
p2.text = "国内外安全生产发展现状"
p2.font.size = Pt(40)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER
add_slide_number(slide20, 20)

# 第21页：国际安全生产形势
slide21 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide21, "国际安全生产形势")
intl = [
    "【总体趋势】全球安全生产形势持续改善，但事故总量仍居高不下",
    "【发达国家】美、德、日等国安全生产水平较高，万人死亡率达0.02-0.04",
    "【发展中国家】工业化进程中事故高发，煤矿、建筑等行业风险突出",
    "【新兴领域】新能源、人工智能等带来新型安全风险",
    "【国际合作】ISO45001等国际标准推动全球职业健康安全管理"
]
add_content_box(slide21, intl, 1.5)
add_slide_number(slide21, 21)

# 第22页：我国安全生产现状
slide22 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide22, "我国安全生产现状")
china = [
    "【总体平稳】事故总量持续下降，2023年各类生产安全事故起数同比下降6.1%",
    "【结构优化】工矿商贸事故下降明显，较大以上事故得到有效遏制",
    "【基础薄弱】中小企业安全投入不足，安全管理水平参差不齐",
    "【风险突出】危化品、矿山、建筑施工等高危行业风险依然较高",
    "【新挑战】新产业新业态新模式带来新风险"
]
add_content_box(slide22, china, 1.5)
add_slide_number(slide22, 22)

# 第23页：国际先进经验
slide23 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide23, "国际先进经验借鉴")
experience = [
    "【德国】工伤保险与安全监管深度融合，浮动费率激励企业主动改善安全",
    "【美国】OSHA监管执法力度强，违规处罚严厉，安全文化深入人心",
    "【日本】精细化管理，\"零灾害\"运动，安全行为养成教育体系完善",
    "【北欧】高福利保障下的安全治理，员工参与度高，安全成为自觉",
    "【新加坡】安全距离管理、安全准入制度、严格的安全绩效考核"
]
add_content_box(slide23, experience, 1.5)
add_slide_number(slide23, 23)

# 第24页：差距与不足
slide24 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide24, "我国与先进水平的差距")
gap = [
    "【安全文化】安全意识有待提高，安全文化氛围不够浓厚",
    "【本质安全】安全技术装备水平与发达国家仍有差距",
    "【监管能力】基层监管力量薄弱，监管手段有待创新",
    "【应急能力】应急救援体系需进一步完善，专业化水平待提升",
    "【信息化水平】安全信息化建设起步较晚，大数据应用有待深化"
]
add_content_box(slide24, gap, 1.5)
add_slide_number(slide24, 24)

# 第25页：发展现状小结
slide25 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide25, "第三部分小结")
summary3 = [
    "✓ 国际安全生产管理经验值得借鉴",
    "✓ 我国安全生产形势持续好转",
    "✓ 但与发达国家仍有差距",
    "✓ 需要持续推进安全生产治理体系和能力现代化"
]
add_content_box(slide25, summary3, 2)
add_slide_number(slide25, 25)

# ==================== 第26-50页：事故案例分析 ====================

# 第26页：第四部分开场
slide26 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide26, "第四部分", RGBColor(0, 100, 80))
p = slide26.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
tf = p.text_frame
p2 = tf.paragraphs[0]
p2.text = "重大典型安全生产事故案例分析"
p2.font.size = Pt(36)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER
add_slide_number(slide26, 26)

# 第27页：案例分析的意义
slide27 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide27, "事故案例分析的意义")
meaning = [
    "▸ 血的教训——历史事故是用生命换来的警示",
    "▸ 警醒教育——让全员深刻认识到事故的严重后果",
    "▸ 举一反三——对照检查，排查自身存在的类似隐患",
    "▸ 完善制度——从事故中汲取教训，完善安全管理",
    "▸ 压实责任——强化各级人员的安全责任意识"
]
add_content_box(slide27, meaning, 2)
add_slide_number(slide27, 27)

# 第28页：2020年福建泉州欣佳酒店坍塌事故
slide28 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide28, "案例一：福建泉州欣佳酒店\"3·7\"坍塌事故", RGBColor(180, 50, 50))
case1 = [
    "【事故概况】2020年3月7日，泉州市鲤城区欣佳酒店发生坍塌，造成29人死亡、42人受伤",
    "【直接原因】违法违规建设施工，原四层钢结构建筑被违规增加夹层改建为七层，导致结构承载力不足",
    "【主要教训】",
    "  ① 违法违规建设：未办理规划、施工许可，违规发包给无资质施工队",
    "  ② 监管严重缺位：相关部门未发现和制止违法行为",
    "  ③ 隐患长期存在：违规改建持续数年，未被有效查处"
]
y = 1.4
for line in case1:
    p = slide28.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(12)
    pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.5
add_slide_number(slide28, 28)

# 第29页：2021年湖北十堰燃气爆炸事故
slide29 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide29, "案例二：湖北十堰燃气爆炸事故", RGBColor(180, 50, 50))
case2 = [
    "【事故概况】2021年6月13日，十堰市张湾区艳湖社区集贸市场发生燃气爆炸，造成25人死亡、138人受伤",
    "【直接原因】燃气管道长期泄漏，积累的燃气达到爆炸极限，遇明火发生爆炸",
    "【主要教训】",
    "  ① 管道日常巡检维护不到位，隐患排查不彻底",
    "  ② 安全监测预警系统不完善，未能及时发现泄漏",
    "  ③ 应急处置不及时，疏散管控不力"
]
y = 1.4
for line in case2:
    p = slide29.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(12)
    pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.5
add_slide_number(slide29, 29)

# 第30页：2022年湖南长沙自建房倒塌事故
slide30 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide30, "案例三：湖南长沙\"4·29\"自建房倒塌事故", RGBColor(180, 50, 50))
case3 = [
    "【事故概况】2022年4月29日，长沙市望城区一居民自建房倒塌，造成54人死亡、9人受伤",
    "【直接原因】房屋地基承载力不足，结构安全储备不足，违规加盖导致荷载增加",
    "【主要教训】",
    "  ① 违规建设：未办理建设审批，违法违规加盖",
    "  ② 监管缺失：违建长期存在，无人查处",
    "  ③ 排查整治不力：未发现和整治重大隐患"
]
y = 1.4
for line in case3:
    p = slide30.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(12)
    pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.5
add_slide_number(slide30, 30)

# 第31页：2023年北京丰台长峰医院火灾
slide31 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide31, "案例四：北京丰台长峰医院火灾事故", RGBColor(180, 50, 50))
case4 = [
    "【事故概况】2023年4月18日，北京长峰医院住院部发生火灾，造成29人死亡",
    "【直接原因】施工作业产生的火花引燃涂料挥发物，与空气混合后爆燃",
    "【主要教训】",
    "  ① 动火作业安全管理缺失，未清理可燃物",
    "  ② 消防设施故障，未能及时发现和处置",
    "  ③ 应急疏散不力，延误逃生时机"
]
y = 1.4
for line in case4:
    p = slide31.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(12)
    pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.5
add_slide_number(slide31, 31)

# 第32页：2024年广东梅大高速塌方
slide32 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide32, "案例五：广东梅大高速塌方灾害", RGBColor(180, 50, 50))
case5 = [
    "【事故概况】2024年5月1日，梅大高速茶阳段发生塌方灾害，造成48人死亡",
    "【直接原因】连续强降雨导致路基边坡失稳，发生滑坡塌方",
    "【主要教训】",
    "  ① 地质风险评估不足，对特殊地质路段重视不够",
    "  ② 监测预警体系不完善，未能提前发现险情",
    "  ③ 应急管控不及时，未能有效阻止车辆进入危险路段"
]
y = 1.4
for line in case5:
    p = slide32.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(12)
    pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.5
add_slide_number(slide32, 32)

# 第33-37页：事故共性教训
slide33 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide33, "事故共性教训分析（一）")
lessons1 = [
    "【主体责任不落实】",
    "▸ 企业安全意识淡薄，安全投入严重不足",
    "▸ 主要负责人安全职责履行不到位，重生产轻安全",
    "▸ 安全管理机构不健全，安全制度不落实"
]
y = 1.4
for line in lessons1:
    p = slide33.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(14)
    if '【' in line:
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(180, 50, 50)
    else:
        pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.55
add_slide_number(slide33, 33)

slide34 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide34, "事故共性教训分析（二）")
lessons2 = [
    "【法规执行不严格】",
    "▸ 有法不依、有章不循，违法违规行为普遍",
    "▸ 安全审批把关不严，源头管控流于形式",
    "▸ 监管执法\"宽松软虚\"，处罚不到位"
]
y = 1.4
for line in lessons2:
    p = slide34.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(14)
    if '【' in line:
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(180, 50, 50)
    else:
        pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.55
add_slide_number(slide34, 34)

slide35 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide35, "事故共性教训分析（三）")
lessons3 = [
    "【隐患排查治理不深入】",
    "▸ 隐患排查走形式、搞过场，发现不了问题",
    "▸ 重大隐患视而不见，长期得不到治理",
    "▸ 双重预防机制流于形式，运行不到位"
]
y = 1.4
for line in lessons3:
    p = slide35.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(14)
    if '【' in line:
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(180, 50, 50)
    else:
        pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.55
add_slide_number(slide35, 35)

slide36 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide36, "事故共性教训分析（四）")
lessons4 = [
    "【应急管理不到位】",
    "▸ 应急预案可操作性差，演练走过场",
    "▸ 应急响应不及时，错失最佳救援时机",
    "▸ 人员应急能力不足，自救互救意识差"
]
y = 1.4
for line in lessons4:
    p = slide36.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(14)
    if '【' in line:
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(180, 50, 50)
    else:
        pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.55
add_slide_number(slide36, 36)

slide37 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide37, "事故警示与反思")
reflection = [
    "【红线意识】安全是必须守住的底线，不可逾越的红线",
    "【问题导向】从事故教训中深刻反思，举一反三",
    "【系统治理】坚持综合施策、标本兼治",
    "【压实责任】确保安全生产责任落实到每个环节"
]
add_content_box(slide37, reflection, 2)
add_slide_number(slide37, 37)

# 第38-42页：建筑业事故案例
slide38 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide38, "建筑施工事故典型案例", RGBColor(180, 100, 50))
case_const = [
    "【事故类型分布】",
    "▸ 高处坠落——占比52%以上",
    "▸ 物体打击——占比15%左右",
    "▸ 坍塌事故——占比10%左右",
    "▸ 机械伤害——占比9%左右",
    "▸ 其他事故——触电、车辆伤害等"
]
y = 1.4
for line in case_const:
    p = slide38.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(14)
    if '【' in line:
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(30, 60, 114)
    else:
        pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.5
add_slide_number(slide38, 38)

slide39 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide39, "建筑施工事故预防措施")
prevent_const = [
    "▸ 严格落实施工安全措施费",
    "▸ 加强特种作业人员管理，持证上岗",
    "▸ 强化高处作业安全防护",
    "▸ 规范脚手架搭设与使用",
    "▸ 做好基坑支护和监测",
    "▸ 落实防高坠各项措施"
]
add_content_box(slide39, prevent_const, 1.5)
add_slide_number(slide39, 39)

# 第40页：危险化学品事故
slide40 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide40, "危险化学品事故案例与教训", RGBColor(180, 50, 50))
chem_cases = [
    "【2015年天津港爆炸】硝酸铵违规储存，事故致165人遇难",
    "【2019年江苏响水爆炸】危化品违法储存，事故致78人死亡",
    "【共同教训】危化品储存管理混乱、隐患长期存在、监管严重缺位",
    "【防控重点】危化品全生命周期管控、风险辨识、储罐监测"
]
add_content_box(slide40, chem_cases, 1.5)
add_slide_number(slide40, 40)

# 第41页：矿山事故
slide41 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide41, "矿山事故案例与教训", RGBColor(180, 50, 50))
mine_cases = [
    "【2023年阿拉善煤矿坍塌】露天煤矿坍塌，53人失联",
    "【2024年山西吕梁火灾】居民楼火灾，26人不幸遇难",
    "【共同教训】风险管控缺失、现场管理混乱、应急处置不力",
    "【防控重点】强化瓦斯水患治理、机电运输安全管理"
]
add_content_box(slide41, mine_cases, 1.5)
add_slide_number(slide41, 41)

# 第42页：第四部分小结
slide42 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide42, "第四部分小结")
summary4 = [
    "✓ 事故是血的教训，必须深刻汲取",
    "✓ 主体责任不落实是一切事故的根源",
    "✓ 隐患排查治理必须深入彻底",
    "✓ 应急管理必须常备不懈"
]
add_content_box(slide42, summary4, 2)
add_slide_number(slide42, 42)

# ==================== 第43-70页：安全管理理念、技术和体系 ====================

# 第43页：第五部分开场
slide43 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide43, "第五部分", RGBColor(0, 100, 80))
p = slide43.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
tf = p.text_frame
p2 = tf.paragraphs[0]
p2.text = "安全管理理念、技术和体系"
p2.font.size = Pt(40)
p2.font.bold = True
p2.font.color.rgb = RGBColor(30, 60, 114)
p2.alignment = PP_ALIGN.CENTER
add_slide_number(slide43, 43)

# 第44页：安全管理核心理念
slide44 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide44, "安全管理核心理念")
philosophy = [
    "【安全第一】安全是一切工作的前提和基础",
    "【预防为主】防患于未然，将事故消灭在萌芽状态",
    "【综合治理】统筹兼顾、多措并举、系统推进",
    "【全员参与】安全人人有责，安全关系你我他",
    "【持续改进】安全管理工作永远在路上"
]
y = 1.4
for line in philosophy:
    p = slide44.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(15)
    if '【' in line:
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(30, 60, 114)
    else:
        pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.6
add_slide_number(slide44, 44)

# 第45页：安全文化理念
slide45 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide45, "安全文化理念")
culture = [
    "【安全愿景】零事故、零伤害、零污染",
    "【安全使命】保障员工生命安全，守护企业健康发展",
    "【安全价值观】安全是最大的效益，安全是最好的业绩",
    "【安全行为准则】遵章守纪、按章操作、拒绝违章",
    "【安全承诺】高层做起、全员参与、持续改进"
]
add_content_box(slide45, culture, 1.5)
add_slide_number(slide45, 45)

# 第46页：杜邦安全管理理念
slide46 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide46, "杜邦安全管理十大理念")
dupont = [
    "1. 所有安全事故都可以预防",
    "2. 各级管理层对安全负有直接责任",
    "3. 各级管理者必须亲自检查安全",
    "4. 员工必须接受严格的安全培训",
    "5. 安全是聘用的条件之一",
    "6. 员工必须参与安全事务",
    "7. 隐患必须及时整改",
    "8. 工作外的安全与工作内安全同样重要",
    "9. 安全是衡量管理业绩的标准",
    "10. 良好的安全创造良好的业绩"
]
y = 1.4
for line in dupont:
    p = slide46.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.5))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(13)
    pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.48
add_slide_number(slide46, 46)

# 第47页：安全生产责任体系
slide47 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide47, "安全生产责任体系")
responsibility = [
    "【横向到边】各职能部门各司其职、齐抓共管",
    "  生产部门抓生产安全 | 技术部门抓技术安全 | 人力部门抓培训安全",
    "【纵向到底】层层分解、责任到人",
    "  企业主要负责人 → 分管领导 → 部门负责人 → 班组 → 岗位",
    "【网格化管理】不留死角、不留空白",
    "  划片包干、定人定责、网格联动"
]
y = 1.4
for line in responsibility:
    p = slide47.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.6))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(13)
    if '【' in line or line.startswith('  '):
        pp.font.color.rgb = RGBColor(51, 51, 51)
        if '【' in line:
            pp.font.bold = True
            pp.font.color.rgb = RGBColor(30, 60, 114)
    y += 0.5
add_slide_number(slide47, 47)

# 第48页：安全管理制度体系
slide48 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide48, "安全管理制度体系")
system = [
    "【基础管理制度】",
    "▸ 安全生产责任制 | 安全目标管理 | 安全绩效考核",
    "【日常管理制度】",
    "▸ 隐患排查治理 | 安全检查 | 安全例会 | 值班值守",
    "【专项管理制度】",
    "▸ 危化品管理 | 特种设备管理 | 动火作业管理 | 有限空间作业管理",
    "【应急管理制度】",
    "▸ 应急预案 | 应急演练 | 应急物资 | 应急培训"
]
y = 1.4
for line in system:
    p = slide48.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.55))
    tf = p.text_frame
    pp = tf.paragraphs[0]
    pp.text = line
    pp.font.size = Pt(12)
    if '【' in line:
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(30, 60, 114)
    else:
        pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.5
add_slide_number(slide48, 48)

# 第49页：安全管理组织机构
slide49 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide49, "安全管理组织机构设置")
org = [
    "【安全生产委员会（安委会）】",
    "▸ 企业最高安全决策机构，主要负责人任主任",
    "▸ 研究部署重大安全事项，协调解决突出问题",
    "【安全管理机构（安监部门）】",
    "▸ 配备专职安全管理人员，行使安全监管职能",
    "▸ 负责日常安全管理和监督检查工作",
    "【基层安全组织】",
    "▸ 各车间、班组设立专兼职安全员"
]
add_content_box(slide49, org, 1.4)
add_slide_number(slide49, 49)

# 第50页：安全风险管控技术
slide50 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide50, "安全风险管控技术（一）")
tech1 = [
    "【安全评价技术】",
    "▸ 安全预评价 | 安全验收评价 | 安全现状评价 | 专项评价",
    "【危险有害因素辨识】",
    "▸ 人的因素、物的因素、环境因素、管理因素",
    "【风险评估方法】",
    "▸ SCL（检查表法）| JHA（作业危害分析法）| FTA（故障树分析）| HAZOP（危险与可操作性分析）"
]
add_content_box(slide50, tech1, 1.4)
add_slide_number(slide50, 50)

# 第51页：隐患排查治理技术
slide51 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide51, "安全风险管控技术（二）——隐患排查治理")
hazard = [
    "【隐患分类】",
    "▸ 一般隐患：危害较小，整改难度不大",
    "▸ 重大隐患：严重危害，整改难度较大，需挂牌督办",
    "【隐患排查方法】",
    "▸ 日常巡检 | 专项检查 | 综合检查 | 专家检查",
    "【闭环管理流程】",
    "▸ 发现 → 报告 → 评估 → 整改 → 验收 → 销号"
]
add_content_box(slide51, hazard, 1.4)
add_slide_number(slide51, 51)

# 第52页：安全监测监控技术
slide52 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide52, "安全监测监控技术")
monitor = [
    "【视频监控系统】24小时实时监控，实现重点区域全覆盖",
    "【气体监测系统】实时监测有毒有害气体浓度，超限报警",
    "【温度压力监测】关键设备运行参数在线监测",
    "【人员定位系统】掌握人员位置，实现智能考勤和应急救援",
    "【智能预警系统】数据综合分析，提前预警潜在风险"
]
add_content_box(slide52, monitor, 1.5)
add_slide_number(slide52, 52)

# 第53页：智慧安全技术
slide53 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide53, "智慧安全技术应用")
smart = [
    "【物联网技术】设备互联互通，数据实时采集传输",
    "【大数据分析】海量数据分析挖掘，精准发现问题",
    "【人工智能】AI智能识别隐患，自动预警提醒",
    "【BIM技术】建筑信息模型，实现可视化安全管理",
    "【数字孪生】虚拟映射实体，构建智能化安全管控平台"
]
add_content_box(slide53, smart, 1.5)
add_slide_number(slide53, 53)

# 第54页：特种设备安全管理
slide54 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide54, "特种设备安全管理")
special = [
    "【特种设备范围】",
    "▸ 锅炉、压力容器、压力管道、电梯、起重机械、场（厂）内专用机动车辆等",
    "【管理要求】",
    "▸ 办理使用登记 | 定期检验检测 | 配备作业人员 | 建立技术档案",
    "【日常管理】",
    "▸ 日常维护保养 | 定期自行检查 | 应急演练演练"
]
add_content_box(slide54, special, 1.4)
add_slide_number(slide54, 54)

# 第55页：危化品安全管理
slide55 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide55, "危险化学品安全管理")
chemical = [
    "【危化品管理重点】",
    "▸ 实施目录管理 | 规范储存管理 | 严格作业审批",
    "【储存管理】",
    "▸ 分类储存 | 防火防爆 | 温湿度控制 | 应急物资配备",
    "【作业管理】",
    "▸ 动火作业、有限空间作业等必须执行作业许可",
    "【运输管理】",
    "▸ 危货运输资质 | 车辆实时监控 | 驾驶员安全管理"
]
add_content_box(slide55, chemical, 1.4)
add_slide_number(slide55, 55)

# 第56页：应急管理体系
slide56 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide56, "应急管理体系建设")
emergency = [
    "【应急预案体系】",
    "▸ 综合预案 | 专项预案 | 现场处置方案",
    "【应急救援队伍】",
    "▸ 专职救援队 | 兼职救援队 | 志愿者队伍",
    "【应急物资装备】",
    "▸ 救援装备 | 医疗急救 | 通讯指挥 | 后勤保障",
    "【应急演练】",
    "▸ 计划制定 | 实战演练 | 总结评估 | 持续改进"
]
add_content_box(slide56, emergency, 1.4)
add_slide_number(slide56, 56)

# 第57页：应急处置程序
slide57 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide57, "事故应急处置程序")
procedure = [
    "【第一时间】启动应急响应，组织人员疏散撤离",
    "【第二时间】报告事故情况，拨打119/120/110",
    "【第三时间】组织先期处置，控制事态发展",
    "【第四时间】配合救援调查，分析事故原因",
    "【第五时间】落实整改措施，防止事故重复"
]
add_content_box(slide57, procedure, 1.8)
add_slide_number(slide57, 57)

# 第58页：安全培训教育体系
slide58 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide58, "安全培训教育体系")
training = [
    "【三级安全教育】",
    "▸ 厂级教育 → 车间教育 → 班组教育",
    "【培训内容】",
    "▸ 安全法律法规 | 安全管理制度 | 操作规程 | 应急知识",
    "【培训方式】",
    "▸ 集中授课 | 实操演练 | 案例分析 | 在线学习",
    "【培训考核】",
    "▸ 理论考试 | 实操考核 | 持证上岗"
]
add_content_box(slide58, training, 1.4)
add_slide_number(slide58, 58)

# 第59页：承包商安全管理
slide59 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide59, "承包商安全管理")
contractor = [
    "【资质审查】严格审查承包商资质、安全业绩和施工能力",
    "【安全协议】签订安全协议，明确双方安全责任",
    "【入厂培训】对承包商人员进行入厂安全教育",
    "【过程监管】加强作业过程监督检查，发现问题及时纠正",
    "【考核评价】定期对承包商安全绩效进行评价考核"
]
add_content_box(slide59, contractor, 1.8)
add_slide_number(slide59, 59)

# 第60页：安全绩效考核
slide60 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide60, "安全绩效考核")
assessment = [
    "【考核指标】",
    "▸ 事故指标（工伤事故率、重大隐患数）",
    "▸ 管理指标（培训完成率、检查覆盖率）",
    "▸ 效果指标（隐患整改率、安全投入使用率）",
    "【考核方式】",
    "▸ 月度考核 | 季度评价 | 年度总评",
    "【结果应用】",
    "▸ 与薪酬挂钩 | 与晋升挂钩 | 与评优挂钩"
]
add_content_box(slide60, assessment, 1.4)
add_slide_number(slide60, 60)

# 第61-65页：安全领导力
slide61 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide61, "安全领导力提升（一）")
leadership1 = [
    "【领导者安全角色】",
    "▸ 决策者——将安全纳入企业战略",
    "▸ 践行者——以身作则遵守安全规定",
    "▸ 推动者——持续推进安全文化建设",
    "▸ 监督者——督促安全责任落实"
]
add_content_box(slide61, leadership1, 1.6)
add_slide_number(slide61, 61)

slide62 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide62, "安全领导力提升（二）")
leadership2 = [
    "【领导者安全行为】",
    "▸ 亲自研究部署安全工作",
    "▸ 定期带队检查安全生产",
    "▸ 参加安全教育培训",
    "▸ 主持召开安全会议",
    "▸ 及时研究解决安全问题"
]
add_content_box(slide62, leadership2, 1.6)
add_slide_number(slide62, 62)

slide63 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide63, "安全领导力提升（三）——走动管理")
walking = [
    "【走动管理要点】",
    "▸ 每天至少一次深入现场检查",
    "▸ 关注人员状态、设备设施、作业环境",
    "▸ 发现问题及时记录和整改",
    "▸ 与员工沟通交流安全话题",
    "▸ 营造高层重视安全的氛围"
]
add_content_box(slide63, walking, 1.6)
add_slide_number(slide63, 63)

slide64 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide64, "安全领导力提升（四）——安全沟通")
communication = [
    "【安全沟通技巧】",
    "▸ 倾听员工安全诉求",
    "▸ 鼓励员工报告安全隐患",
    "▸ 及时反馈处理结果",
    "▸ 表扬安全行为",
    "▸ 公开表彰安全先进"
]
add_content_box(slide64, communication, 1.6)
add_slide_number(slide64, 64)

slide65 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide65, "安全领导力提升（五）")
leadership5 = [
    "【安全决策原则】",
    "▸ 安全优先原则——安全与效率冲突时，安全优先",
    "▸ 预防优先原则——事前投入大于事后代价",
    "▸ 综合治理原则——多措并举、系统推进",
    "【决策要求】",
    "▸ 重大决策必须安全论证",
    "▸ 新项目必须安全审查",
    "▸ 新技术必须安全评估"
]
add_content_box(slide65, leadership5, 1.4)
add_slide_number(slide65, 65)

# 第66页：班组安全管理
slide66 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide66, "班组安全管理")
team = [
    "【班组安全核心】",
    "▸ 班前安全喊话——布置任务同时布置安全",
    "▸ 班中安全互保——互相监督、互相提醒",
    "▸ 班后安全总结——总结经验、改进不足",
    "【安全员职责】",
    "▸ 检查现场安全状态",
    "▸ 纠正违章作业行为",
    "▸ 报告安全隐患"
]
add_content_box(slide66, team, 1.4)
add_slide_number(slide66, 66)

# 第67页：作业现场安全管理
slide67 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide67, "作业现场安全管理")
site = [
    "【定置管理】物品定位摆放，通道畅通有序",
    "【目视管理】安全标识清晰，风险提示明确",
    "【三违管理】违章指挥、违章作业、违反劳动纪律",
    "【防护管理】正确佩戴防护用品，做好个人防护",
    "【环境管理】照明通风良好，温度湿度适宜"
]
add_content_box(slide67, site, 1.8)
add_slide_number(slide67, 67)

# 第68页：安全技术进步趋势
slide68 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide68, "安全技术发展趋势")
trend = [
    "【智能化】AI识别隐患，无人机巡检，机器人作业",
    "【信息化】大数据分析风险，云平台管理，移动互联",
    "【无人化】自动化替代危险岗位，远程控制操作",
    "【集成化】系统集成联动，综合预警指挥",
    "【标准化】国际标准对接，标准化建设深化"
]
add_content_box(slide68, trend, 1.8)
add_slide_number(slide68, 68)

# 第69页：第五部分小结
slide69 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide69, "第五部分小结")
summary5 = [
    "✓ 安全管理理念是安全生产的灵魂",
    "✓ 安全技术是安全生产的支撑",
    "✓ 安全体系是安全生产的保障",
    "✓ 安全领导力是落实责任的关键",
    "✓ 班组是安全生产的基础"
]
add_content_box(slide69, summary5, 2)
add_slide_number(slide69, 69)

# 第70页：落实主体责任的核心
slide70 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide70, "落实主体责任的核心要素")
core = [
    "【责任到位】横到边、纵到底，责任全覆盖",
    "【投入到位】安全投入有保障，设施设备完好",
    "【培训到位】全员培训，提高安全意识技能",
    "【管理到位】制度健全，执行有力，考核严格",
    "【应急到位】预案完善，演练经常，处置有效"
]
add_content_box(slide70, core, 1.8)
add_slide_number(slide70, 70)

# ==================== 第71-80页：总结与承诺 ====================

# 第71页：总结
slide71 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide71, "培训总结")
conclusion = [
    "▸ 安全生产必须警钟长鸣、常抓不懈",
    "▸ 落实主体责任是安全生产的根本保障",
    "▸ 隐患排查治理是防范事故的关键抓手",
    "▸ 应急能力建设是减少损失的重要手段",
    "▸ 安全文化建设是长治久安的根本途径"
]
add_content_box(slide71, conclusion, 2)
add_slide_number(slide71, 71)

# 第72页：工作要求
slide72 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide72, "下一步工作要求")
requirements = [
    "1. 深入开展安全隐患大排查大整治",
    "2. 全面推进安全生产标准化建设",
    "3. 切实加强应急能力建设",
    "4. 持续强化安全培训和宣传教育",
    "5. 加快安全管理信息化建设"
]
y = 1.6
for line in requirements:
    box = slide72.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(8.4), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 240)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.9
add_slide_number(slide72, 72)

# 第73页：庄严承诺
slide73 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide73, "庄严承诺")
promise_box = slide73.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(8), Inches(4))
promise_box.fill.solid()
promise_box.fill.fore_color.rgb = RGBColor(30, 60, 114)
tf = promise_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "我郑重承诺："
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
promises = [
    "严格遵守安全生产法律法规",
    "认真落实安全生产主体责任",
    "扎实开展安全隐患排查治理",
    "不断提升安全管理能力和水平",
    "切实保障员工生命安全和身体健康"
]
for pr in promises:
    pp = tf.add_paragraph()
    pp.text = "✓ " + pr
    pp.font.size = Pt(16)
    pp.font.color.rgb = RGBColor(255, 255, 255)
add_slide_number(slide73, 73)

# 第74页：安全寄语
slide74 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide74, "安全寄语")
quote = "\"安全生产只有起点，没有终点；只有逗号，没有句号。\""
quote_box = slide74.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(2))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(245, 248, 255)
tf = quote_box.text_frame
p = tf.paragraphs[0]
p.text = quote
p.font.size = Pt(22)
p.font.italic = True
p.font.color.rgb = RGBColor(30, 60, 114)
p.alignment = PP_ALIGN.CENTER
msg = slide74.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf2 = msg.text_frame
p2 = tf2.paragraphs[0]
p2.text = "让我们共同努力，守护安全，共创美好未来！"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(51, 51, 51)
p2.alignment = PP_ALIGN.CENTER
add_slide_number(slide74, 74)

# 第75-78页：思考题
slide75 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide75, "思考与讨论")
questions = [
    "1. 结合本职工作，谈谈如何落实安全生产责任？",
    "2. 请举例说明本岗位存在的主要安全风险及防控措施？",
    "3. 如果发现重大安全隐患，应该如何处理？",
    "4. 发生事故时，你会如何进行初期处置和应急疏散？",
    "5. 对照先进企业经验，我们还有哪些差距和不足？"
]
y = 1.6
for q in questions:
    p = slide75.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.8))
    tf = p.text_frame
    tf.word_wrap = True
    pp = tf.paragraphs[0]
    pp.text = q
    pp.font.size = Pt(15)
    pp.font.color.rgb = RGBColor(51, 51, 51)
    y += 0.8
add_slide_number(slide75, 75)

# 第76页：培训评估
slide76 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide76, "培训效果评估")
evaluation = [
    "【评估方式】",
    "▸ 理论测试（占40%）",
    "▸ 实操考核（占40%）",
    "▸ 培训表现（占20%）",
    "【评估要求】",
    "▸ 80分以上为合格",
    "▸ 不合格者需补训补考"
]
add_content_box(slide76, evaluation, 1.6)
add_slide_number(slide76, 76)

# 第77页：学习资料
slide77 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide77, "学习资料推荐")
materials = [
    "【法规文件】",
    "▸ 《安全生产法》《消防法》《危化品安全管理条例》",
    "【标准规范】",
    "▸ 《企业安全生产标准化基本规范》",
    "【专业书籍】",
    "▸ 《现代安全管理》《安全心理学》《事故预防与控制》"
]
add_content_box(slide77, materials, 1.6)
add_slide_number(slide77, 77)

# 第78页：联系方式
slide78 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide78, "联系方式")
contact = [
    "【安全管理部门】",
    "▸ 联系电话：XXX-XXXX-XXXX",
    "▸ 邮箱：safety@company.com",
    "【应急救援电话】",
    "▸ 消防：119 | 急救：120 | 报警：110",
    "【安全举报】",
    "▸ 匿名举报邮箱：report@company.com"
]
add_content_box(slide78, contact, 1.8)
add_slide_number(slide78, 78)

# 第79页：感谢页
slide79 = prs.slides.add_slide(prs.slide_layouts[6])
thanks = slide79.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
tf = thanks.text_frame
p = tf.paragraphs[0]
p.text = "感谢各位领导和同事的积极参与！"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(30, 60, 114)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "安全生产，人人有责"
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(70, 130, 180)
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "让我们共同守护安全，共享美好生活！"
p3.font.size = Pt(20)
p3.font.color.rgb = RGBColor(51, 51, 51)
p3.alignment = PP_ALIGN.CENTER
add_slide_number(slide79, 79)

# 第80页：备用页
slide80 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide80, "附页")
backup = [
    "本PPT可根据实际需要增补以下内容：",
    "▸ 企业安全生产实际情况分析",
    "▸ 行业安全事故案例汇编",
    "▸ 本单位安全生产规章制度",
    "▸ 应急预案文本及演练记录",
    "▸ 安全生产考核奖惩办法"
]
add_content_box(slide80, backup, 2)
add_slide_number(slide80, 80)

# 保存
prs.save('/Users/mac/.openclaw/workspace/安全生产主体责任PPT_80页.pptx')
print("PPT已生成：安全生产主体责任PPT_80页.pptx")
print(f"共80页")
