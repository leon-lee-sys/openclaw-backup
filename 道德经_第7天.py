#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""道德经 第31-35章 课件"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
PRIMARY=RGBColor(30,60,114);SECONDARY=RGBColor(0,100,80);ACCENT=RGBColor(70,130,180)
LIGHT_BG=RGBColor(245,248,255);TEXT_DARK=RGBColor(51,51,51);WHITE=RGBColor(255,255,255)

def sn(s,n,t=7):
    f=s.shapes.add_textbox(Inches(9.3),Inches(7.25),Inches(.6),Inches(.25))
    tf=f.text_frame.paragraphs[0];tf.text=f"{n}/{t}";tf.font.size=Pt(9);tf.font.color.rgb=RGBColor(180,180,180);tf.alignment=PP_ALIGN.RIGHT

def tb(s,t,c=PRIMARY):
    b=s.shapes.add_shape(1,Inches(0),Inches(0),Inches(10),Inches(.8));b.fill.solid();b.fill.fore_color.rgb=c;b.line.fill.background()
    tx=s.shapes.add_textbox(Inches(.5),Inches(.18),Inches(9),Inches(.6));tf=tx.text_frame;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(24);p.font.bold=True;p.font.color.rgb=WHITE

def box_item(s,y,txt,c=LIGHT_BG,fs=14,ic=PRIMARY):
    b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(.65));b.fill.solid();b.fill.fore_color.rgb=c;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(fs);p.font.color.rgb=TEXT_DARK
    return y+.72

def quote_box(s,y,txt):
    b=s.shapes.add_shape(12,Inches(.5),Inches(y),Inches(9),Inches(1.2));b.fill.solid();b.fill.fore_color.rgb=RGBColor(240,248,255);tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(16);p.font.italic=True;p.font.color.rgb=PRIMARY;p.alignment=PP_ALIGN.CENTER

# 封面
s1=prs.slides.add_slide(prs.slide_layouts[6])
b=s1.shapes.add_shape(1,Inches(0),Inches(0),Inches(10),Inches(3.2));b.fill.solid();b.fill.fore_color.rgb=PRIMARY;b.line.fill.background()
t=s1.shapes.add_textbox(Inches(.5),Inches(.8),Inches(9),Inches(1.2));tf=t.text_frame;p=tf.paragraphs[0];p.text="道德经";p.font.size=Pt(52);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
p2=tf.add_paragraph();p2.text="第31-35章";p2.font.size=Pt(36);p2.font.color.rgb=ACCENT;p2.alignment=PP_ALIGN.CENTER
l=s1.shapes.add_shape(1,Inches(3),Inches(3.4),Inches(4),Pt(4));l.fill.solid();l.fill.fore_color.rgb=ACCENT;l.line.fill.background()
t2=s1.shapes.add_textbox(Inches(1),Inches(3.8),Inches(8),Inches(.6));tf2=t2.text_frame;p3=tf2.paragraphs[0];p3.text="第三十一章 佳兵不祥 | 第三十二章 知止不殆\n第三十三章 知人者智 | 第三十四章 不自为大\n第三十五章 往而不害";p3.font.size=Pt(16);p3.font.color.rgb=PRIMARY;p3.alignment=PP_ALIGN.CENTER
t3=s1.shapes.add_textbox(Inches(1),Inches(5.5),Inches(8),Inches(.5));tf3=t3.text_frame;p4=tf3.paragraphs[0];p4.text="2026年4月14日 · 第7天";p4.font.size=Pt(16);p4.font.color.rgb=TEXT_DARK;p4.alignment=PP_ALIGN.CENTER
sn(s1,1)

# 第31章
s2=prs.slides.add_slide(prs.slide_layouts[6]);tb(s2,"第三十一章 佳兵不祥")
quote_box(s2,1.2,"夫佳兵者，不祥之器。\n物或恶之，故有道者不处。\n君子居则贵左，用兵则贵右。\n兵者不祥之器，非君子之器，\n不得已而用之，恬淡为上。")
y=2.7
for t in ["【解读】", "本章论战争之道。兵革为不祥之物，君子应淡然处之。", "【要点】", "1. 兵器为不祥之器，圣人不得已而用之", "2. 用兵应以恬淡为上，不可逞强好战", "3. 战争应以防御为主，不可轻动"]:
    if '【' in t:
        b=s2.shapes.add_shape(12,Inches(.5),Inches(y),Inches(1.5),Inches(.5));b.fill.solid();b.fill.fore_color.rgb=SECONDARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(12);p.font.bold=True;p.font.color.rgb=WHITE
    else:
        box_item(s2,y,t,fs=12);y-=.02
    y+=.55
sn(s2,2)

# 第32章
s3=prs.slides.add_slide(prs.slide_layouts[6]);tb(s3,"第三十二章 知止不殆")
quote_box(s3,1.2,"道常无名，朴。\n虽小，天下莫能臣。\n侯王若能守之，万物将自宾。\n天地相合，以降甘露，\n民莫之令而自均。")
y=2.7
for t in ["【解读】", "道如同未经雕琢的朴玉，虽小却为天下所尊贵。", "【要点】", "1. 道常无名，朴素的品质最为珍贵", "2. 侯王若能守住道，万物自然归顺", "3. 知止不殆，知道何时停止方可免于危险"]:
    if '【' in t:
        b=s3.shapes.add_shape(12,Inches(.5),Inches(y),Inches(1.5),Inches(.5));b.fill.solid();b.fill.fore_color.rgb=SECONDARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(12);p.font.bold=True;p.font.color.rgb=WHITE
    else:
        box_item(s3,y,t,fs=12);y-=.02
    y+=.55
sn(s3,3)

# 第33章
s4=prs.slides.add_slide(prs.slide_layouts[6]);tb(s4,"第三十三章 知人者智")
quote_box(s4,1.2,"知人者智，自知者明。\n胜人者有力，自胜者强。\n知足者富，强行者有志。\n不失其所者久，死而不亡者寿。")
y=2.7
for t in ["【解读】", "本章论修养境界。知人、自知、胜人、自胜，境界递进。", "【要点】", "1. 知人者智：了解他人是有智慧", "2. 自知者明：认识自己是真正的明智", "3. 自胜者强：战胜自己才是真正的强者", "4. 死而不亡：精神永存方为长寿"]:
    if '【' in t:
        b=s4.shapes.add_shape(12,Inches(.5),Inches(y),Inches(1.5),Inches(.5));b.fill.solid();b.fill.fore_color.rgb=SECONDARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(12);p.font.bold=True;p.font.color.rgb=WHITE
    else:
        box_item(s4,y,t,fs=12);y-=.02
    y+=.55
sn(s4,4)

# 第34章
s5=prs.slides.add_slide(prs.slide_layouts[6]);tb(s5,"第三十四章 不自为大")
quote_box(s5,1.2,"大道泛兮，其可左右。\n万物恃之而生而不辞，\n功成不名有。\n衣养万物而不为主，常无欲，\n可名于小；万物归焉而不为主，\n可名为大。\n以其终不自为大，故能成其大。")
y=2.7
for t in ["【解读】", "大道广泛流行，无所不至。因其不自以为大，故能成就其大。", "【要点】", "1. 大道泛兮：道如河水泛滥，无处不在", "2. 功成不名有：成就功业却不自以为有功", "3. 衣养万物而不为主：滋养万物却不做主人", "4. 不自为大故能成其大：因谦下而成就伟大"]:
    if '【' in t:
        b=s5.shapes.add_shape(12,Inches(.5),Inches(y),Inches(1.5),Inches(.5));b.fill.solid();b.fill.fore_color.rgb=SECONDARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(12);p.font.bold=True;p.font.color.rgb=WHITE
    else:
        box_item(s5,y,t,fs=12);y-=.02
    y+=.55
sn(s5,5)

# 第35章
s6=prs.slides.add_slide(prs.slide_layouts[6]);tb(s6,"第三十五章 往而不害")
quote_box(s6,1.2,"执大象，天下往。\n往而不害，安平泰。\n乐与饵，过客止。\n道之出口，淡乎其无味，\n视之不足见，听之不足闻，\n用之不足既。")
y=2.7
for t in ["【解读】", "得道者天下往，平和安泰。道虽平淡无味，用之不尽。", "【要点】", "1. 执大象，天下往：把握大道者，天下归心", "2. 往而不害，安平泰：归往而不相害，则平安泰和", "3. 乐与饵，过客止：美食音乐只能留住过客", "4. 道之用不足既：道的作用无穷无尽"]:
    if '【' in t:
        b=s6.shapes.add_shape(12,Inches(.5),Inches(y),Inches(1.5),Inches(.5));b.fill.solid();b.fill.fore_color.rgb=SECONDARY;b.line.fill.background();tf=b.text_frame;p=tf.paragraphs[0];p.text=t;p.font.size=Pt(12);p.font.bold=True;p.font.color.rgb=WHITE
    else:
        box_item(s6,y,t,fs=12);y-=.02
    y+=.55
sn(s6,6)

# 总结
s7=prs.slides.add_slide(prs.slide_layouts[6]);tb(s7,"第31-35章 总结")
y=1.2
for t in ["第三十一章：兵者不祥之器，恬淡为上", "第三十二章：道常无名，知止不殆", "第三十三章：知人自知，自胜者强", "第三十四章：不自为大，故能成其大", "第三十五章：执大象，天下往"]:
    box_item(s7,y,t);y+=.65
b=s7.shapes.add_shape(12,Inches(1),Inches(4.5),Inches(8),Inches(1.5));b.fill.solid();b.fill.fore_color.rgb=LIGHT_BG;tf=b.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text="【核心要义】\n道家以柔克刚、以退为进的思想贯穿始终。\n知止、自知、不自为大，方能成就大业。";p.font.size=Pt(14);p.font.color.rgb=PRIMARY;p.alignment=PP_ALIGN.CENTER
sn(s7,7)

prs.save('/Users/mac/.openclaw/workspace/道德经_第7天.pptx')
print("道德经第7天课件已生成：道德经_第7天.pptx")
